"""Generate an editable base PPTX template for report template v1 (Stage K1).

The v1 renderer builds slides programmatically and does NOT require this file.
This script produces an editable base deck (brand colours on the slide master +
a styled cover layout) that a designer can refine; if the file exists next to
the renderer the v1 builder will use it as the presentation base.

Usage:
    python build_template.py [output_path]
Default output: renderer/templates/report-template-v1.pptx
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

NAVY = RGBColor(0x0E, 0x1F, 0x3A)
ACCENT = RGBColor(0x1C, 0x6F, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(6858000)


def build_template(out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # A single guidance slide so the file is a valid, openable, editable deck.
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    box = slide.shapes.add_textbox(Emu(457200), Emu(2743200), Emu(8229600), Emu(1828800))
    tf = box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = "Digital Profile Audit — Report Template v1"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p = tf.add_paragraph()
    r2 = p.add_run()
    r2.text = "Editable base deck. The renderer generates content slides programmatically."
    r2.font.size = Pt(14)
    r2.font.color.rgb = ACCENT

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"Template written to {out_path}")


if __name__ == "__main__":
    default = os.path.join(os.path.dirname(__file__), "templates", "report-template-v1.pptx")
    build_template(sys.argv[1] if len(sys.argv) > 1 else default)
