"""Public render_orion_golden entrypoint."""

from __future__ import annotations

import base64
import json
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu

try:
    from client_text_contract import resolve_contract
except ImportError:  # pragma: no cover
    from renderer.client_text_contract import resolve_contract  # type: ignore

from .common import (
    SLIDE_H,
    SLIDE_W,
    _Ctx,
    _asset_map,
    _resolve_image_bytes,
    assert_render_font_family,
    get_bullet_measure,
    get_layout_telemetry,
    reset_bullet_measure,
    reset_layout_telemetry,
)
from .export import _export_png_pages, _write_pdf_fallback
from .slides import _render_slide

#: Версия формы вердикта мерного прогона: меняется форма — меняется строка.
BULLET_MEASURE_VERSION = "orion-bullet-measure-v1"


def _draw_deck(
    payload: dict[str, Any], *, log_assets: bool
) -> tuple[Presentation, list[str], str, list[dict[str, Any]], dict[str, dict[str, Any]]]:
    """Нарисовать деку в память: презентация, предупреждения, субъект, слайды, ассеты.

    Общая фаза настоящего рендера и мерного прогона: расходиться им негде,
    потому что это один и тот же код рисования. Всё, что происходит дальше —
    сохранение, конвертация, растр — к вопросу «сколько влезло» отношения не
    имеет и мерному прогону не нужно.
    """
    assert_render_font_family()
    reset_layout_telemetry()
    reset_bullet_measure()
    deck = payload.get("deckManifest") or {}
    report_spec = payload.get("reportSpec") or {}
    slides = list(deck.get("finalSlides") or [])
    if not slides:
        raise ValueError("deckManifest.finalSlides is empty")

    assets = _asset_map(payload)
    if log_assets:
        # Diagnostics for blank SERP slides (lengths only — never log base64).
        asset_diag = []
        for ref, asset in list(assets.items())[:20]:
            raw = _resolve_image_bytes(asset)
            asset_diag.append(
                {
                    "assetRef": ref,
                    "kind": asset.get("kind"),
                    "hasImageData": bool(asset.get("imageData")),
                    "imageDataChars": len(str(asset.get("imageData") or "")),
                    "hasStorageKey": bool(asset.get("storageKey")),
                    "resolvedBytes": len(raw) if raw else 0,
                }
            )
        print(
            "[orion-golden-render] assets",
            json.dumps(
                {
                    "assetCount": len(assets),
                    "serpSlides": sum(
                        1
                        for s in slides
                        if str(s.get("template") or "") == "orion_golden_serp_screenshot"
                    ),
                    "sample": asset_diag[:8],
                },
                ensure_ascii=False,
            ),
            flush=True,
        )
    subject = (report_spec.get("subject") or {}).get("displayName") or "Цифровой профиль"
    total = len(slides)
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    client_text_contract = resolve_contract(payload.get("clientTextContract"))

    warnings: list[str] = [f"client-text-contract:{client_text_contract.get('version')}"]
    for idx, slide in enumerate(slides, start=1):
        ctx = _Ctx(
            prs,
            idx,
            total,
            client_text_contract=client_text_contract,
            slide_key=str(slide.get("slideKey") or ""),
        )
        _render_slide(ctx, slide, assets)
        ctx.footer()
        warnings.extend(ctx.warnings)
    return prs, warnings, str(subject), slides, assets


def measure_orion_golden(payload: dict[str, Any]) -> dict[str, Any]:
    """Мерный прогон: нарисовать деку в память и отчитаться о пути буллетов.

    Ни PPTX, ни PDF, ни PNG, ни файла телеметрии — мера не оставляет следов,
    которые судят ворота выпуска. Отвечает тем же кодом рисования, поэтому
    разойтись с настоящим рендером ей негде: «сколько влезает на лист» —
    вопрос с одним ответом.
    """
    _draw_deck(payload, log_assets=False)
    # Потери сайдбара — той же мерой: телеметрия того же прохода, отобранная
    # по роли. Раньше вердикт знал только буллеты, и страница 62 прогона
    # DPA-2026-0053 дошла до выпуска с сайдбаром, потерявшим два блока.
    sidebars = [
        {
            "page": int(row.get("page") or 0),
            "field": str(row.get("name") or "")
            .replace(f"_p{row.get('page')}", "")
            .replace("orion_sidebar_", ""),
            "droppedLines": int(row.get("droppedLines") or 0),
            "requiredHeight": int(row.get("requiredHeight") or 0),
            "availableHeight": int(row.get("availableHeight") or 0),
        }
        for row in get_layout_telemetry()
        if str(row.get("role") or "") == "sidebar" and int(row.get("droppedLines") or 0) > 0
    ]
    return {"version": BULLET_MEASURE_VERSION, "pages": get_bullet_measure(), "sidebars": sidebars}


def render_orion_golden(payload: dict[str, Any]) -> dict[str, Any]:
    prs, warnings, subject, slides, assets = _draw_deck(payload, log_assets=True)

    with tempfile.TemporaryDirectory(prefix="orion-golden-") as tmp:
        tmp_path = Path(tmp)
        pptx_path = tmp_path / "report.pptx"
        prs.save(str(pptx_path))
        pdf_path = tmp_path / "report.pdf"
        pdf_ok = False
        pdf_mode = "fitz-fallback"
        try:
            from convert_pdf import convert_to_pdf

            convert_to_pdf(str(pptx_path), str(pdf_path))
            pdf_ok = pdf_path.exists() and pdf_path.stat().st_size > 0
            if pdf_ok:
                pdf_mode = "libreoffice"
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"libreoffice-failed:{exc}")

        if not pdf_ok:
            _write_pdf_fallback(slides, pdf_path, str(subject), assets)
            pdf_mode = "fitz-fallback"

        print(f"[orion-golden-render] pdfExportMode={pdf_mode} warnings={warnings}", flush=True)
        pages = _export_png_pages(pdf_path)
        telemetry = get_layout_telemetry()
        return {
            "slideCount": len(prs.slides),
            "pptxBase64": base64.b64encode(pptx_path.read_bytes()).decode("ascii"),
            "pdfBase64": base64.b64encode(pdf_path.read_bytes()).decode("ascii") if pdf_path.exists() else "",
            "pages": pages,
            "pdfExportMode": pdf_mode,
            "warnings": warnings,
            "layoutTelemetry": {
                "version": "orion-layout-telemetry-v1",
                "entries": telemetry,
            },
        }


if __name__ == "__main__":
    import sys

    data = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    out = render_orion_golden(data)
    Path(sys.argv[2]).write_bytes(base64.b64decode(out["pptxBase64"]))
    if out.get("pdfBase64"):
        Path(sys.argv[3]).write_bytes(base64.b64decode(out["pdfBase64"]))
    pages_dir = Path(sys.argv[4])
    pages_dir.mkdir(parents=True, exist_ok=True)
    for page in out.get("pages") or []:
        Path(pages_dir / f"page-{page['pageNumber']:02d}.png").write_bytes(
            base64.b64decode(page["contentBase64"])
        )
    meta = {
        "slideCount": out["slideCount"],
        "pages": len(out.get("pages") or []),
        "pdfExportMode": out.get("pdfExportMode"),
        "warnings": out.get("warnings") or [],
    }
    Path(pages_dir.parent / "golden-render-meta.json").write_text(json.dumps(meta), encoding="utf-8")
    telemetry = out.get("layoutTelemetry")
    if telemetry:
        Path(pages_dir.parent / "layout-telemetry.json").write_text(
            json.dumps(telemetry, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    print(json.dumps(meta))

