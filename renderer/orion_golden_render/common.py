"""ORION Golden Report renderer — shared theme, fonts, layout, and canvas helpers.
Split from orion_golden_renderer.py (REMEDIATION §9.5) — mechanical move only.
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
import tempfile
from functools import lru_cache
from pathlib import Path
from typing import Any, Callable

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

try:
    from client_text_contract import (
        renderer_strip_re,
        resolve_contract,
        sidebar_check_failures,
    )
except ImportError:  # pragma: no cover — package-style import inside container
    from renderer.client_text_contract import (  # type: ignore
        renderer_strip_re,
        resolve_contract,
        sidebar_check_failures,
    )

# Гарнитура отчёта. Inter — SIL OFL 1.1: встраивание в PPTX/PDF, отдаваемые
# клиенту, разрешено лицензией и не требует покупки. Кириллица нарисована
# авторами гарнитуры, а не подставлена запасным шрифтом.
#
# До этого стоял DejaVu Sans — системный шрифт Linux по умолчанию. Он не
# нейтрален: широкие пропорции и характерный рисунок читаются как «документ,
# свёрстанный чем попало», и для отчёта, который показывают состоятельному
# клиенту, это видно с первого взгляда.
#
# Побочно Inter у́же: та же строка занимает 387px против 418px у DejaVu (−7%),
# а жирное шире обычного на 3,6% вместо ~10%. Текста на странице помещается
# больше, а разрыв между замером и выводом стал меньше.
FONT = "Inter"
FS_TITLE = 26
FS_SECTION = 22
FS_SUBTITLE = 13
FS_BODY = 11  # min readable body ≥ 10.5pt
FS_CAPTION = 9  # footer/provenance 8.5–9pt
FS_COVER = 36
FS_LEAD = 20
FS_CARD_TITLE = 16

#: Типографическая шкала: единственный набор допустимых кеглей (ADR-0008).
#:
#: Замер отрисованной деки до этого шага: **шестнадцать** разных кеглей, до
#: восьми на одной странице, десять ступеней между 8,5 и 14 pt. Такой набор
#: глаз не читает как иерархию — он читает его как шум, и все блоки выглядят
#: одинаково важными.
#:
#: Шкала была объявлена и раньше (`FS_*`), но не соблюдалась: размеры писались
#: литералами (`Pt(14)`, `Pt(10.5)`, `Pt(8.5)`), а ступенчатое уменьшение
#: шрифта порождало 9,5 и 10,5 арифметикой. Один вопрос — два ответа.
TYPE_SCALE_PT: tuple[float, ...] = (
    FS_CAPTION,      # 9  — подпись, колонтитул, происхождение
    FS_BODY,         # 11 — основной текст
    FS_SUBTITLE,     # 13 — лид, заголовок карточки риска
    FS_CARD_TITLE,   # 16 — заголовок блока
    FS_LEAD,         # 20 — заголовок страницы
    FS_SECTION,      # 22 — заголовок раздела
    FS_TITLE,        # 26 — крупный заголовок
    FS_COVER,        # 36 — титул
)


def _scale_steps_below(value: float) -> list[float]:
    """Ступени шкалы ниже текущей, от крупной к мелкой."""
    return [s for s in reversed(TYPE_SCALE_PT) if s < value]


def scale_pt(value: float) -> float:
    """Ближайшая ступень шкалы, не крупнее запрошенного размера.

    Уменьшать кегль «на полпункта», чтобы текст влез, шкала не позволяет: по
    ADR-0008 не влезающий текст сокращается или переносится на следующую
    страницу, а не ужимается до нечитаемого. Поэтому округление — вниз, к
    ближайшей объявленной ступени, и ниже подписи опускаться некуда.
    """
    smaller = [s for s in TYPE_SCALE_PT if s <= value]
    return smaller[-1] if smaller else TYPE_SCALE_PT[0]


def next_smaller_scale_pt(value: float) -> float | None:
    """Следующая ступень вниз — для тех мест, где кегль всё же снижают."""
    smaller = [s for s in TYPE_SCALE_PT if s < value]
    return smaller[-1] if smaller else None

# Master slide 16:10 (12.8" × 8.0") — matches ORION reference aspect.
SLIDE_W = 11_704_320
SLIDE_H = 7_315_200
MARGIN_X = 480_000
CONTENT_W = SLIDE_W - 2 * MARGIN_X
FOOTER_Y = SLIDE_H - 440_000
# PDF-46 I.2 — hard clearance above confidential footer (measure underestimates
# multi-line theme cards; keep a wide safety band).
CONTENT_BOTTOM = SLIDE_H - 1_100_000

# Визуальная система cleeq (https://cleeq.ru) — только краска.
#
# Прежняя палитра (тёмно-синий + золото) досталась от «ORION Golden design
# system v2». Смена палитры не трогает ни сбор, ни клиентский текст: это
# именно те токены, которыми красят фон, карточки и акценты.
#
# Основной зелёный #24D875, чернила #101510, фиолетовый #AE7AFF,
# мятный лист #F6F8F4.
COVER_BG = RGBColor(0x10, 0x15, 0x10)
PAGE_BG = RGBColor(0xF6, 0xF8, 0xF4)
NAVY = RGBColor(0x10, 0x15, 0x10)
TITLE_COLOR = RGBColor(0xF7, 0xF9, 0xF5)
BODY_COLOR = RGBColor(0x10, 0x15, 0x10)
MUTED_COLOR = RGBColor(0x5B, 0x66, 0x5E)
ACCENT = RGBColor(0x24, 0xD8, 0x75)
VIOLET = RGBColor(0xAE, 0x7A, 0xFF)
CYAN = RGBColor(0x5B, 0xC8, 0xFF)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xD8, 0xE3, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_SOFT = RGBColor(0xE0, 0xFB, 0xEC)
VIOLET_SOFT = RGBColor(0xF0, 0xEE, 0xFF)
WARN_BG = RGBColor(0xFC, 0xF4, 0xF1)
RISK_BG = RGBColor(0xFF, 0xF1, 0xF0)
GOOD_BG = RGBColor(0xE0, 0xFB, 0xEC)
TONE_RISK = RGBColor(0xE1, 0x3D, 0x3D)
TONE_WARN = RGBColor(0xE3, 0x8A, 0x24)
TONE_GOOD = RGBColor(0x19, 0x64, 0x48)
METRIC_ACCENT = RGBColor(0x24, 0xD8, 0x75)
#: Акцент заголовка. Имя осталось прежним, чтобы не переписывать два десятка
#: мест ради переименования краски: теперь это зелёный cleeq, а не золото.
GOLD = ACCENT
DARK_RULE = RGBColor(0x24, 0x2C, 0x24)
COVER_SUBTITLE = RGBColor(0xC8, 0xD4, 0xCA)
STAGE_SHADOW = RGBColor(0xE4, 0xEB, 0xE4)

#: Маркер пункта — квадрат cleeq вместо точки. Глиф есть в Inter (U+25AA), и
#: ширина у него та же, что у «•», поэтому замер строки не меняется.
BULLET_GLYPH = "▪"

# REMEDIATION §6.1 — pattern sourced from client_text_contract.json
FORBIDDEN = renderer_strip_re()

# EMU helpers: 914400 EMU = 1 inch; 72 pt = 1 inch
EMU_PER_PT = 12_700
EMU_PER_INCH = 914_400

# Renderer-side text layout telemetry (page, font, box, measured vs available).
_LAYOUT_TELEMETRY: list[dict[str, Any]] = []


def reset_layout_telemetry() -> None:
    _LAYOUT_TELEMETRY.clear()


def get_layout_telemetry() -> list[dict[str, Any]]:
    return list(_LAYOUT_TELEMETRY)


# Мерные записи пути буллетов: доступная под список высота, высота каждого
# поданного элемента, сколько удержано. Живут отдельно от телеметрии разметки
# намеренно: телеметрию судят ворота выпуска, и запись мерного прогона в ней
# означала бы, что ворота осудили не тот документ.
_BULLET_MEASURE: list[dict[str, Any]] = []


def reset_bullet_measure() -> None:
    _BULLET_MEASURE.clear()


def get_bullet_measure() -> list[dict[str, Any]]:
    return list(_BULLET_MEASURE)


def record_bullet_measure(
    *,
    slide_key: str,
    page: int,
    available_height: int,
    max_items: int,
    item_heights: list[int],
    kept_items: int,
    dropped_bullets: int,
    dropped_lines: int,
) -> None:
    """Записать меру одной страницы пути буллетов.

    Высота элемента считается той же функцией, которой блок будет нарисован, —
    построитель складывает эти числа и сравнивает с `availableHeight`. Сумма
    отдельных высот всегда не меньше высоты того же набора одним куском
    (постоянная добавка входит в каждое слагаемое), поэтому укладка по мере
    консервативнее отрисовки — в ту сторону, где потери невозможны.
    """
    _BULLET_MEASURE.append(
        {
            "slideKey": slide_key,
            "page": page,
            "availableHeight": available_height,
            "maxItems": max_items,
            "itemHeights": item_heights,
            "keptItems": kept_items,
            "droppedBullets": dropped_bullets,
            "droppedLines": dropped_lines,
        }
    )


def record_text_layout(
    *,
    page: int,
    name: str,
    role: str,
    font_family: str,
    font_size_pt: float,
    box_width: int,
    box_height: int,
    available_height: int,
    required_height: int,
    measured_lines: int,
    text_length: int,
    clipped: bool,
    measurement_uncertain: bool = False,
    dropped_bullets: int = 0,
    dropped_lines: int = 0,
    drawn_blocks: int | None = None,
) -> None:
    entry: dict[str, Any] = {
        "page": page,
        "name": name,
        "role": role,
        "fontFamily": font_family,
        "fontSizePt": font_size_pt,
        "boxWidth": box_width,
        "boxHeight": box_height,
        "availableHeight": available_height,
        "requiredHeight": required_height,
        "measuredLines": measured_lines,
        "textLength": text_length,
        "clipped": clipped,
        "measurementUncertain": measurement_uncertain,
    }
    # Блоков на странице бывает больше одного (карточная сетка), и это не то
    # же самое, что число строк: пока карточки писались в `measuredLines`,
    # поле схемы «сколько строк намерено» отвечало на чужой вопрос.
    if drawn_blocks is not None:
        entry["drawnBlocks"] = drawn_blocks
    # Потеря содержимого — не то же самое, что вылезший за рамку текст, и
    # называется отдельно (шаг 16, 07.6).
    if dropped_bullets or dropped_lines:
        entry["droppedBullets"] = dropped_bullets
        entry["droppedLines"] = dropped_lines
    _LAYOUT_TELEMETRY.append(entry)


def _count_measured_lines(
    text: str,
    width_emu: int,
    font_size_pt: float,
    bold: bool = False,
) -> tuple[int, bool]:
    """Return (line_count, measurement_uncertain).

    Второй флаг остаётся ради телеметрии разметки, но «неточно» здесь больше не
    бывает: замер либо выполнен настоящими метриками, либо не выполнен вовсе и
    поднял ошибку. Приблизительный расчёт по числу символов убран: он отвечал
    «примерно столько строк» и расходился с отрисовкой на длинных словах —
    блок мерился одним, а рисовался другим.
    """
    return _wrapped_line_count(text, width_emu, font_size_pt, bold), False


#: Начертания, которыми рендерер и рисует, и меряет. Ключ — то, что в слайде
#: задаётся как `run.font.bold`; значение — файл, которым это же начертание
#: замеряется. Одна таблица на оба вопроса: пока их было два, они разошлись.
_FONT_FILES = {
    False: "Inter-Regular.otf",
    True: "Inter-Bold.otf",
}


def _font_path(bold: bool = False) -> str | None:
    """Файл того начертания, которым текст будет нарисован.

    Прежде функция возвращала единственный `DejaVuSans.ttf` на все случаи, а
    жирность в слайдах ставится через python-pptx (`run.font.bold = True`) в
    двадцати трёх местах. Ширина считалась по обычному начертанию, рисовалось
    жирное — оно шире, и строка не влезала в отведённое место. Замер по
    образцам отчёта: жирное шире обычного на 9,6–14,4%. Это и есть «текст
    выезжает за блоки» — заголовки, подписи и первые строки пунктов, то есть
    ровно те места, где стоит жирность.
    """
    filename = _FONT_FILES[bool(bold)]
    here = Path(__file__).resolve().parent
    override = os.environ.get("ORION_RENDER_FONT_BOLD" if bold else "ORION_RENDER_FONT")
    candidates = [
        override,
        str(here / "fonts" / filename),
        str(here.parent / "fonts" / filename),
        f"/usr/share/fonts/opentype/inter/{filename}",
        f"/usr/share/fonts/truetype/inter/{filename}",
        rf"C:\Windows\Fonts\{filename}",
    ]
    for path in candidates:
        if path and Path(path).is_file():
            return path
    return None


def _load_measure_font(font_size_pt: float, bold: bool):
    """Шрифт для замера — или внятный отказ.

    Приблизительный расчёт «все символы одной ширины» отсюда убран намеренно.
    Он давал неверную ширину всегда (шрифт пропорциональный), а для строки из
    «ш», «ж», «м» ошибка кратная — и при этом молчал, то есть выдавал вёрстку,
    про которую никто не знал, что она посчитана наугад. Тихая неточность здесь
    дороже явного падения: по тому же правилу в этом проекте
    `CONTENT_DROPPED_BY_RENDERER` стал CRITICAL.
    """
    if Image is None:
        raise RuntimeError(
            "ORION render: Pillow недоступен, замер текста невозможен. "
            "Pillow приходит зависимостью python-pptx — проверьте образ рендерера."
        )
    from PIL import ImageFont  # type: ignore

    fp = _font_path(bold)
    if not fp:
        weight = "жирное" if bold else "обычное"
        raise RuntimeError(
            f"ORION render font missing: нет файла для начертания «{weight}» "
            f"({_FONT_FILES[bool(bold)]}) — искали в renderer/fonts и системных каталогах"
        )
    return ImageFont.truetype(fp, size=max(8, int(round(font_size_pt * 96 / 72))))


#: Кегль, на котором снимается натуральная высота строки начертания.
#:
#: PIL округляет ascent/descent до целых пикселей, а рабочие кегли мелкие: на
#: 15 px (это 11 pt) отношение выходит 1,267 вместо 1,211 — лишние 4,6 % на
#: каждой строке. От кегля отношение не зависит, поэтому снимается один раз на
#: крупном размере: 750 pt — это ровно 1000 px, где округление незаметно.
_LINE_METRIC_PT = 750.0


@lru_cache(maxsize=4)
def _natural_line_ratio(bold: bool) -> float:
    """Натуральная высота строки начертания в долях кегля (у Inter ~1,21)."""
    font = _load_measure_font(_LINE_METRIC_PT, bold)
    ascent, descent = font.getmetrics()
    return (ascent + descent) / (_LINE_METRIC_PT * 96 / 72)


def font_line_step_emu(font_size_pt: float, line_spacing: float, bold: bool = False) -> int:
    """Шаг строки: натуральная строка начертания × выставленный межстрочный.

    Межстрочный, который ставится абзацу (`p.line_spacing`), — множитель к
    натуральной строке шрифта, а не к кеглю. Замер по растру эталона: тело
    11 pt с межстрочным 1,12 идёт шагом 190 500 EMU, а «кегль × 1,12» дало бы
    156 464 — на 18 % меньше нарисованного. Мера, посчитанная от кегля, обещает
    больше строк, чем помещается, и текст уезжает за блок.
    """
    return int(font_size_pt * EMU_PER_PT * _natural_line_ratio(bold) * line_spacing)


def text_width_px(text: str, font_size_pt: float, bold: bool = False) -> int:
    """Ширина строки тем начертанием, которым она будет нарисована.

    Единственный ответ на вопрос «сколько занимает эта строка»: до этого один и
    тот же цикл переноса жил в `_count_measured_lines` и в
    `measure_text_height` двумя копиями.
    """
    font = _load_measure_font(font_size_pt, bold)
    bbox = font.getbbox(text)
    return int(bbox[2] - bbox[0])


def _wrapped_line_count(
    text: str,
    width_emu: int,
    font_size_pt: float,
    bold: bool = False,
) -> int:
    """Сколько строк займёт текст в рамке заданной ширины."""
    raw = _safe_preserve_breaks(text)
    if not raw:
        return 1
    # Slightly narrower than box so PPTX wrap is not underestimated.
    width_px = max(40, int(width_emu / EMU_PER_INCH * 96 * 0.90))
    paragraphs = [p.strip() for p in re.split(r"\n+", raw) if p.strip()] or [""]
    font = _load_measure_font(font_size_pt, bold)
    total_lines = 0
    for para in paragraphs:
        words = para.split(" ")
        if not words:
            total_lines += 1
            continue
        line = ""
        lines = 1
        for word in words:
            word_box = font.getbbox(word)
            word_px = int(word_box[2] - word_box[0])
            if word_px > width_px:
                # Токен без пробелов шире рамки (домен, транслитерация, адрес)
                # PPTX ломает по знакам. Прежний счёт объявлял его одной
                # строкой любой длины, и мера занижалась во столько раз, во
                # сколько он шире рамки. Хвост токена считается занявшим свою
                # строку целиком: где именно он оборвётся, знает только вывод.
                if line:
                    lines += 1
                lines += -(-word_px // width_px) - 1
                line = ""
                continue
            trial = word if not line else f"{line} {word}"
            bbox = font.getbbox(trial)
            if int(bbox[2] - bbox[0]) <= width_px or not line:
                line = trial
            else:
                lines += 1
                line = word
        total_lines += max(1, lines)
    return max(1, total_lines)


def assert_render_font_family() -> str:
    """Старт/QA: шрифт замера совпадает со шрифтом вывода, оба начертания есть.

    Прежде проверка требовала именно DejaVu и падала на всём остальном — то
    есть заодно делала типографику неизменяемой, хотя следила совсем за другим.
    Требование теперь осмысленное и не мешает сменить гарнитуру (ADR-0008):
    **чем меряем, тем и рисуем, и оба начертания доступны**. Несовпадение
    по-прежнему валит старт — молча разъехавшиеся метрики и есть тот дефект,
    ради которого проверка писалась.
    """
    regular = _font_path(False)
    if not regular:
        raise RuntimeError(
            f"ORION render font missing: нет обычного начертания ({_FONT_FILES[False]}) "
            "в renderer/fonts или системных каталогах"
        )
    bold = _font_path(True)
    if not bold:
        raise RuntimeError(
            f"ORION render font missing: нет жирного начертания ({_FONT_FILES[True]}). "
            "Жирный текст мерился бы обычным начертанием и выезжал за рамку."
        )
    # Рисует python-pptx по имени семейства (FONT), меряет PIL по файлу. Если
    # это разные гарнитуры, замер верен для одной, а на слайде другая.
    family_token = FONT.replace(" ", "").lower()
    for path in (regular, bold):
        if family_token not in Path(path).name.replace(" ", "").lower():
            raise RuntimeError(
                f"ORION render font mismatch: рисуем «{FONT}», меряем {path} — "
                "замер относится не к тому шрифту, которым выводится текст"
            )
    return regular


def measure_text_height(
    text: str,
    width_emu: int,
    font_size_pt: float,
    line_spacing: float = 1.2,
    paragraph_spacing_pt: float = 6.0,
    bold: bool = False,
) -> int:
    """Высота текста в EMU, замеренная тем начертанием, которым он рисуется.

    `bold` обязан совпадать с тем, что ставится тексту в слайде
    (`run.font.bold`): жирное шире обычного на 9,6–14,4% по образцам отчёта, и
    расхождение здесь выдавливает строку за рамку.
    """
    # Переводы строк обязаны дожить до замера. `_safe` схлопывает любые пробелы,
    # включая `\n`, поэтому многострочная карточка мерилась как один абзац:
    # строки «упаковывались» плотнее, чем рисуются, высота выходила заниженной,
    # и текст вылезал за карточку, перекрываясь следующей (шаг 13, D1).
    raw = _safe_preserve_breaks(text)
    if not raw:
        return int(font_size_pt * EMU_PER_PT * line_spacing)
    paragraphs = [p.strip() for p in re.split(r"\n+", raw) if p.strip()] or [""]
    total_lines = _wrapped_line_count(text, width_emu, font_size_pt, bold)

    line_h = font_size_pt * EMU_PER_PT * line_spacing
    para_extra = max(0, len(paragraphs) - 1) * paragraph_spacing_pt * EMU_PER_PT
    # Safety margin: PPTX wraps more aggressively than PIL metrics.
    return int((total_lines * line_h + para_extra) * 1.18)


def _fit_text_to_height(
    text: str,
    width_emu: int,
    font_size_pt: float,
    max_h: int,
    *,
    line_spacing: float = 1.2,
    bold: bool = False,
) -> str:
    """Clip text so measured height fits max_h. Prefer complete sentences; never bare ellipsis stubs."""
    raw = _safe(text)
    if not raw:
        return ""
    if measure_text_height(raw, width_emu, font_size_pt, line_spacing=line_spacing, bold=bold) <= max_h:
        return raw
    # Prefer longest prefix that ends on a sentence boundary.
    sentences = re.split(r"(?<=[.!?…])\s+", raw)
    kept: list[str] = []
    for sent in sentences:
        trial = " ".join(kept + [sent]).strip()
        if measure_text_height(trial, width_emu, font_size_pt, line_spacing=line_spacing, bold=bold) <= max_h:
            kept.append(sent)
        else:
            break
    if kept:
        return " ".join(kept).strip()
    # Fall back to word clip without ellipsis; stop before dangling prepositions.
    words = raw.split()
    lo, hi = 1, len(words)
    best = words[0]
    while lo <= hi:
        mid = (lo + hi) // 2
        trial = " ".join(words[:mid])
        if measure_text_height(trial, width_emu, font_size_pt, line_spacing=line_spacing, bold=bold) <= max_h:
            best = trial
            lo = mid + 1
        else:
            hi = mid - 1
    return _trim_dangling_tail(best)


_DANGLING_TAIL = re.compile(
    r"(\s+(?:и|а|или|по|на|в|с|из|для|о|об|к|ко|у|от|до|при|без|над|под|про|через|как|что|чтобы|который|которая|которые|которых|баз|записям|доменом?|контекстом?|криминальным|компрометирующим|санкционным|нежелательный|ручной|негативным|т\.?\s*ч\.?))\s*$|"
    r"(\s+в\s+т\.?\s*ч\.?\s*$)|(\s+с\s+[А-ЯA-Z]\.?\s*$)|(\s+[А-ЯA-Z]\.?\s*$)",
    re.I,
)


def _trim_dangling_tail(text: str) -> str:
    val = text.strip()
    end = ""
    # Многоточие в конце — маркер обрезки, поставленный поисковой системой:
    # «Как Павел Дуров создал «ВКонтакте», Telegram и стал...». Прежний разбор
    # брал последнюю точку как знак конца, а остальные срезал вместе с
    # пунктуацией, и обрубок превращался в законченную с виду мысль
    # («…и стал.»). Читателю это врёт о полноте заголовка (шаг 15, E3).
    if val.endswith("…"):
        end = "…"
        val = val[:-1].rstrip(",;: ")
    elif val.endswith("..."):
        end = "…"
        val = val[:-3].rstrip(",;: ")
    elif val and val[-1] in ".!?":
        end = val[-1]
        val = val[:-1].rstrip(".,;: ")
    for _ in range(4):
        nxt = _DANGLING_TAIL.sub("", val).rstrip(".,;: ")
        if nxt == val:
            break
        val = nxt
    if not val:
        return ""
    return val + end


def _safe(text: object) -> str:
    val = re.sub(r"\s+", " ", str(text or "")).strip()
    val = FORBIDDEN.sub("", val)
    # Humanize residual enum-like tokens that may appear in summaries
    val = re.sub(r"\bWRONG[_\s-]?SUBJECT\b", "другой субъект", val, flags=re.I)
    val = re.sub(r"\bPENDING\b", "требует проверки", val, flags=re.I)
    val = re.sub(r"\bGPT\b", "модельный анализ", val)
    val = re.sub(r"\bVISUAL_ASSET_UNAVAILABLE\b", "визуальный экспорт недоступен", val, flags=re.I)
    val = re.sub(r"\bLEXISNEXIS_SIGNAL\b", "сигнал LexisNexis", val, flags=re.I)
    val = re.sub(r"\bPotential match\b", "Потенциальное совпадение", val, flags=re.I)
    return val.strip()


def _fit_lines_to_height(text: str, budget_emu: int, measure: Callable[[str], int]) -> str:
    """Оставляет столько целых строк, сколько помещается по высоте.

    Строка отбрасывается целиком: обрезка посередине даёт обрубок, который
    читается как сбой. Если не помещается даже первая, она остаётся — пустая
    карточка хуже плотной.

    Мера приходит от вызывающего и обязана быть той же, которой посчитан
    бюджет. Пока внутри был зашит межстрочный 1,2, страховка решала «что
    влезло» не тем прибором, которым мерился блок, — и выбрасывала строки,
    которые на листе помещались.

    Уборка повисшего ввода (`_close_dangling_lead_in`) осталась вызывающему:
    здесь она снималась всегда, а не только после реза, и разница «строк до /
    строк после» объявляла потерей чистку, при которой ничего не пропало.
    Считать потерю надо до гигиены — так же, как её считает
    `_clip_structured_bullet`, зовущий уборку только на ветке реза.
    """
    lines = [ln for ln in str(text or "").split("\n") if ln.strip()]
    if not lines:
        return ""
    kept: list[str] = []
    for ln in lines:
        trial = "\n".join(kept + [ln])
        if measure(trial) > budget_emu and kept:
            break
        kept.append(ln)
    return "\n".join(kept)


def _close_dangling_lead_in(text: str) -> str:
    """Двоеточие обещает продолжение — если его нет, обещание надо снять.

    Строка вида «…установлено: … В выборке присутствует материал «X» (Y):»
    получает двоеточие как ввод к списку цитат, которые идут следующими
    строками. Карточка с ограниченной высотой эти строки отбрасывает — и на
    главном слайде клиента оставалось предложение, обрывающееся двоеточием
    (шаг 13, C1).

    Чистый ввод («Найдены публикации:») без продолжения выбрасывается целиком:
    он не несёт содержания. Строка, в которой содержание есть, сохраняется и
    заканчивается точкой — мысль в ней уже высказана.

    Признак — содержание, а не длина. Порог в 80 знаков делил один дефект
    надвое: «Что делать: Проверить статусы дел, …направлениям:» сохранялся, а
    «Что делать: Проверить следующее:» исчезал целиком, унося текст аналитика.
    Так же терялась короткая строка с цитатой и адресом — «В выборке
    присутствует материал «Заголовок» (издание.ru):», — хотя в ней и заголовок,
    и источник (пункт BB бэклога).
    """
    val = (text or "").rstrip()
    if not val.endswith(":"):
        return text
    lines = val.split("\n")
    last = lines[-1].strip()
    if len(lines) > 1 and not _has_own_content(last):
        return "\n".join(lines[:-1]).rstrip()
    return val[:-1].rstrip(" ,;—-") + "."


def _has_own_content(line: str) -> bool:
    """Несёт ли строка что-то конкретное, кроме обещания списка.

    Конкретное — это цитата, адрес, число или собственный ярлык с текстом
    после него: во всех четырёх случаях строка уже что-то сказала, и снимать
    её значило бы терять сказанное. Ввод без единого такого признака
    («Найдены публикации:», «Примеры:») говорит только о том, что дальше будет
    список, — а списка после реза нет.
    """
    body = line.rstrip()
    if body.endswith(":"):
        body = body[:-1]
    if ":" in body:  # свой ярлык, и после него текст
        return True
    if any(ch in body for ch in "«»\"“”"):
        return True
    if any(ch.isdigit() for ch in body):
        return True
    return bool(_DOMAIN_IN_TEXT.search(body))


#: Адрес внутри фразы: «издание.ru», «example.com/page».
_DOMAIN_IN_TEXT = re.compile(r"[\w-]+\.[a-z]{2,}(?:/|\b)", re.IGNORECASE)


def _safe_preserve_breaks(text: object) -> str:
    """Like _safe, but keeps intentional newlines for structured theme cards."""
    raw = str(text or "").replace("\r\n", "\n")
    parts = []
    for ln in raw.split("\n"):
        cleaned = _safe(ln)
        if cleaned:
            parts.append(cleaned)
    return "\n".join(parts)


def plural_ru(n: int, one: str, few: str, many: str) -> str:
    """Russian plural agreement: 1 сигнал / 2–4 сигнала / 5+ сигналов."""
    abs_n = abs(int(n)) % 100
    last = abs_n % 10
    if 10 < abs_n < 20:
        return many
    if last == 1:
        return one
    if 2 <= last <= 4:
        return few
    return many


def _clip_words(text: str, max_chars: int) -> str:
    """Emergency brake only (PDF-36 D.4): drop the incomplete last sentence
    whole whenever a complete one fits — a shorter finished thought beats a
    longer broken one. Word-boundary cut is the last resort."""
    val = _safe(text)
    if len(val) <= max_chars:
        return _trim_dangling_tail(val)
    slice_ = val[:max_chars]
    punct = max(slice_.rfind(". "), slice_.rfind("! "), slice_.rfind("? "), slice_.rfind("; "))
    # Any complete sentence ≥20% of the budget wins over a broken tail.
    if punct > max_chars * 0.2:
        return slice_[: punct + 1].rstrip()
    sp = max(slice_.rfind(" "), slice_.rfind("\u00a0"))
    if sp > max_chars * 0.4:
        return _trim_dangling_tail(slice_[:sp])
    soft = re.sub(r"[^\s]{1,12}$", "", slice_).rstrip()
    if len(soft) > max_chars * 0.35:
        return _trim_dangling_tail(soft)
    return _trim_dangling_tail(slice_)


_META_LINE_RE = re.compile(
    r"^(Источники(?:\s+в\s+регионе)?|Примеры(?:\s+заголовков)?|Где видно|В корпусе|Всего по теме|Пример)\s*:",
    re.I,
)
# Allow one level of nested guillemets: «Экс-владелец «Главстроя»».
_QUOTE_BODY = r"(?:[^«»]|«[^»]*»)+"
_QUOTE_SOURCE_RE = re.compile(rf"^«{_QUOTE_BODY}»\s*—\s*источник\b", re.I)
_THEME_LINE_RE = re.compile(r"^«[^»]{2,80}»\s*$")


_QUOTE_SOURCE_INLINE_RE = re.compile(
    rf"«{_QUOTE_BODY}»\s*—\s*источник\s+[A-Za-z0-9][A-Za-z0-9.-]*",
    re.I,
)
_TAIL_SPLIT_RE = re.compile(
    r"(?=(?:Всего по теме:|В корпусе:|Где видно:|Что делать:|Для банка|Банки |Это усиливает|Риск в том|Деловой фон))"
)


def _reflow_g2b_bullet(raw: str) -> list[str] | None:
    """PDF-43 — restore theme / framing / quotes / scale when GPT flattened G.2b."""
    flat = re.sub(r"\s+", " ", raw.replace("\n", " ")).strip()
    if not flat:
        return None
    quotes = _QUOTE_SOURCE_INLINE_RE.findall(flat)
    if not quotes:
        return None
    theme = ""
    rest = flat
    theme_m = re.match(r"^(«[^»]{2,80}»)\s+(.*)$", flat)
    if (
        theme_m
        and "источник" not in theme_m.group(1).lower()
        and re.match(r"^(Найдены|Есть публикации|В открытой)", theme_m.group(2))
    ):
        theme = theme_m.group(1).strip()
        rest = theme_m.group(2).strip()
        quotes = _QUOTE_SOURCE_INLINE_RE.findall(rest)
    first = _QUOTE_SOURCE_INLINE_RE.search(rest)
    if not first:
        return None
    framing = rest[: first.start()].strip()
    if framing and re.search(r"Найдены|публик|материал", framing, re.I) and not framing.endswith(":"):
        framing = re.sub(r"[.:]\s*$", "", framing) + ":"
    last_end = 0
    for m in _QUOTE_SOURCE_INLINE_RE.finditer(rest):
        last_end = m.end()
    tail = rest[last_end:].strip()
    out: list[str] = []
    if theme:
        out.append(theme)
    if framing:
        out.append(framing)
    out.extend(q.strip() for q in quotes)
    if tail:
        out.extend(p.strip() for p in _TAIL_SPLIT_RE.split(tail) if p.strip())
    return out if len(out) >= 3 else None


def _normalize_nested_guillemets(text: str) -> str:
    """«Экс-владелец «Главстроя»» → «Экс-владелец \"Главстроя\"» for stable quote parse."""
    prev = None
    out = text
    while prev != out:
        prev = out
        out = re.sub(r"«([^«»]*)«([^»]+)»([^»]*)»", r'«\1"\2"\3»', out)
    return out


def _split_structured_bullet(text: str) -> list[str]:
    """PDF-38 F.1 / PDF-43 — normalize a theme bullet into theme / quotes / meta lines."""
    # PDF-47 — never flatten newlines here (_safe would glue G.2b into one line).
    raw = _normalize_nested_guillemets(_safe_preserve_breaks(text)).strip()
    if not raw:
        return []
    # Strip trailing finding marker before structural split (kept out of draw lines).
    raw_core = re.sub(r"\s*\[finding-[^\]]+\]\s*$", "", raw).strip() or raw
    lines_in = [ln.strip() for ln in raw_core.split("\n") if ln.strip()]
    needs_reflow = False
    for ln in lines_in:
        n = len(_QUOTE_SOURCE_INLINE_RE.findall(ln))
        if n > 1:
            needs_reflow = True
            break
        if n == 1 and re.search(
            r"(?:Всего по теме:|В корпусе:|Где видно:|Для банка|Банки |Риск в том|Что делать:)",
            ln,
        ):
            needs_reflow = True
            break
        if re.match(r"^«[^»]{2,80}»\s+(?:Найдены|Есть публикации|В открытой)", ln):
            needs_reflow = True
            break
    if needs_reflow or ("\n" not in raw_core and _QUOTE_SOURCE_INLINE_RE.search(raw_core)):
        reflowed = _reflow_g2b_bullet(raw_core)
        if reflowed:
            return reflowed
    if "\n" in raw_core:
        return lines_in
    # One-line «Theme» — body with optional Sources/Examples tails.
    m = re.match(r"^(«[^»]+»)\s*[—\-–:]?\s*(.+)$", raw_core, re.S)
    if m:
        theme, rest = m.group(1).strip(), m.group(2).strip()
        sources = re.search(
            r"(Источники(?:\s+в\s+регионе)?:\s*.+?)(?=\s*Примеры|\s*$)", rest, re.I
        )
        examples = re.search(r"((?:Примеры(?:\s+заголовков)?):\s*.+)$", rest, re.I)
        stats = rest
        if sources:
            stats = stats.replace(sources.group(0), "").strip()
        if examples:
            stats = stats.replace(examples.group(0), "").strip()
        stats = re.sub(r"\s+", " ", stats).strip(" ;,.")
        if stats and not re.search(r"[.!?…]$", stats):
            stats = f"{stats}."
        lines = [theme]
        if stats:
            lines.append(stats)
        if sources:
            lines.append(re.sub(r"\s+", " ", sources.group(1)).strip())
        if examples:
            lines.append(
                re.sub(r"\s+", " ", examples.group(1)).replace("Примеры заголовков:", "Примеры:").strip()
            )
        return lines
    return [raw_core]


def _structured_bullet_body(text: str) -> str:
    """Текст буллета без обрубков модели — то, что ножницы будут резать бюджетом.

    Выделено из `_clip_structured_bullet`, потому что на этот же вопрос отвечает
    учёт потерь в `ctx.bullets`: чистка заглушек («Что делать:.») — не потеря
    содержимого, и мерить её как срез значит поднимать CRITICAL там, где ничего
    не пропало. Одна функция на оба ответа.
    """
    return "\n".join(_cleaned_structured_lines(text))


def _cleaned_structured_lines(text: str) -> list[str]:
    lines = _split_structured_bullet(text)
    if not lines:
        return []
    # Drop empty / broken GPT stubs («Что делать:.», «контекстом —.», «Где видно:.»).
    cleaned: list[str] = []
    for ln in lines:
        if re.match(r"^(Что делать|Всего по теме|В корпусе|Где видно)\s*:\s*\.?$", ln, re.I):
            continue
        if re.search(r",\s*с негативным контекстом\s+[—–-]\s*\.?$", ln, re.I):
            continue
        # PDF-49 — dangling token must follow « or whitespace (not «…Дерипаски»).
        if re.search(
            r"«.*(?:«|\s)(?:из-за|и|в|во|на|по|с|со|о|об|and|or|of|the|to|for|with|from|by|over)\s*»",
            ln,
            re.I | re.S,
        ):
            continue
        # PDF-48 — drop quotes that end on a comma/colon stub («…Дерипаски,»).
        if re.search(r"«[^»]*[,;:]\s*»", ln):
            continue
        cleaned.append(ln)
    return cleaned


def _clip_structured_bullet(text: str, max_chars: int) -> str:
    """Keep whole structural lines only — never mid-cut a quote/phrase (PDF-45/46)."""
    lines = _cleaned_structured_lines(text)
    if not lines:
        return ""
    if sum(len(ln) for ln in lines) + max(0, len(lines) - 1) <= max_chars:
        return "\n".join(lines)
    kept: list[str] = []
    used = 0
    for ln in lines:
        sep = 1 if kept else 0
        room = max_chars - used - sep
        if room < 12:
            break
        if len(ln) <= room:
            kept.append(ln)
            used += sep + len(ln)
            continue
        # Quote / theme / meta: skip whole line rather than publish «…visa over».
        if _QUOTE_SOURCE_RE.match(ln) or ln.startswith("«") or _META_LINE_RE.match(ln):
            break
        # Non-quote prose: keep the head that fits.
        #
        # Граница предложения ищется **внутри** бюджета, и это единственное
        # правило реза прозы — то же, которым режет `_clip_words`. Прежде
        # `ln.rfind(". ")` искал границу по всей строке: у нарратива в 1387
        # знаков последняя точка лежала за 900-м, условие «граница помещается»
        # не выполнялось, и функция возвращала пустую строку. Страница 30
        # живого отчёта пришла с двумя буллетами из трёх, а исчезнувший нёс
        # метод проверки, дату и ссылку на статью.
        kept.append(_clip_words(ln, room))
        break
    # Never fall back to a mid-word / mid-phrase slice of the first line.
    return _close_dangling_lead_in("\n".join(kept))


def _bullet_line_style(line: str, *, is_first: bool) -> tuple[bool, RGBColor, float]:
    """Return (bold, color, size_pt) for one line inside a structured bullet."""
    if _META_LINE_RE.match(line):
        return False, MUTED_COLOR, FS_CAPTION
    # Concrete evidence quotes are body text, not theme headers.
    if _QUOTE_SOURCE_RE.match(line):
        return False, BODY_COLOR, float(FS_BODY)
    if is_first and (_THEME_LINE_RE.match(line) or (line.endswith(":") and len(line) <= 80)):
        return True, NAVY, FS_BODY
    if is_first and line.startswith("«") and "»" in line[:90] and "источник" not in line.lower():
        return True, NAVY, FS_BODY
    return False, BODY_COLOR, float(FS_BODY)


# REMEDIATION §6.2 — replace QA-violating sidebar fields; never fail the whole render.

SIDEBAR_SAFE_FALLBACK = "См. таблицу результатов на этой странице."


class _Ctx:
    def __init__(
        self,
        prs: Presentation,
        page: int,
        total: int,
        *,
        client_text_contract: dict[str, Any] | None = None,
        slide_key: str = "",
    ):
        self.prs = prs
        self.page = page
        self.total = total
        # Ключ страницы деки: по нему построитель находит свою страницу в
        # вердикте меры. Номер страницы для этого не годится — при перекладке
        # он и меняется.
        self.slide_key = slide_key
        self.client_text_contract = resolve_contract(client_text_contract)
        self.warnings: list[str] = []
        self.dark = False
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        self.slide = prs.slides.add_slide(layout)

    def footer(self) -> None:
        # Hairline rule + brand line left, page counter right (design v2).
        rule = self.slide.shapes.add_shape(
            1, Emu(MARGIN_X), Emu(FOOTER_Y - 40_000), Emu(CONTENT_W), Emu(12_700)
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = DARK_RULE if self.dark else CARD_BORDER
        rule.line.fill.background()
        brand = self.slide.shapes.add_textbox(
            Emu(MARGIN_X), Emu(FOOTER_Y), Emu(CONTENT_W // 2), Emu(250000)
        )
        bp = brand.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = "cleeq · Конфиденциально"
        br.font.name = FONT
        br.font.size = Pt(FS_CAPTION)
        br.font.color.rgb = MUTED_COLOR
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(FOOTER_Y), Emu(CONTENT_W), Emu(250000))
        try:
            box.name = f"orion_footer_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{self.page} / {self.total}"
        r.font.name = FONT
        r.font.size = Pt(FS_CAPTION)
        r.font.color.rgb = MUTED_COLOR

    def dark_bg(self) -> None:
        self.dark = True
        fill = self.slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COVER_BG

    def light_bg(self) -> None:
        self.dark = False
        fill = self.slide.background.fill
        fill.solid()
        # Мятный лист cleeq вместо белого: белые карточки на нём читаются как
        # отдельные плоскости, а не как продолжение фона.
        fill.fore_color.rgb = PAGE_BG

    def title(
        self,
        text: str,
        y: int = 280000,
        color: RGBColor = TITLE_COLOR,
        size: int = FS_TITLE,
        *,
        accent: bool = True,
    ) -> int:
        # Зелёная засечка cleeq у первой строки заголовка.
        #
        # Отступ заголовка от неё увеличен со 165 000 до 200 000 EMU, а сама
        # засечка стала шире и со скруглением. Высота полосы заголовка при этом
        # не меняется (`return y + 950_000` ниже): ёмкость страницы в строках
        # объявлена в TS-шаблоне и откалибрована по замеру рендерера, поэтому
        # сдвигать содержимое вниз ради воздуха нельзя — это выдавило бы блок
        # за нижнюю границу листа.
        text_x = MARGIN_X
        text_w = CONTENT_W
        if accent:
            bar_h = int(size * EMU_PER_PT * 1.15)
            bar = self.slide.shapes.add_shape(
                5, Emu(MARGIN_X), Emu(y + 40_000), Emu(70_000), Emu(bar_h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            try:
                bar.adjustments[0] = 0.5
            except Exception:  # noqa: BLE001
                pass
            text_x = MARGIN_X + 200_000
            text_w = CONTENT_W - 200_000
        box = self.slide.shapes.add_textbox(Emu(text_x), Emu(y), Emu(text_w), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(text)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = color
        return y + 950000

    def card(
        self,
        y: int,
        h: int = 2_200_000,
        *,
        x: int | None = None,
        w: int | None = None,
        fill: RGBColor = CARD_BG,
        border: RGBColor | None = CARD_BORDER,
        radius: float = 0.055,
    ):
        """Рисует карточку и возвращает фигуру — вызывающему может понадобиться
        погасить у неё тень, которую LibreOffice подставляет по умолчанию."""
        left = MARGIN_X if x is None else x
        width = CONTENT_W if w is None else w
        avail = max(200_000, min(h, CONTENT_BOTTOM - y))
        # 5 = rounded rectangle: мягкие углы cleeq на всех карточках.
        shape = self.slide.shapes.add_shape(5, Emu(left), Emu(y), Emu(width), Emu(avail))
        try:
            shape.adjustments[0] = max(radius, 0.08)
        except Exception:  # noqa: BLE001
            pass
        try:
            shape.name = f"orion_card_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if border is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = border
            shape.line.width = Pt(0.75)
        return shape

    def body(
        self,
        text: str,
        y: int,
        max_h: int = 900000,
        color: RGBColor = BODY_COLOR,
        *,
        x: int | None = None,
        w: int | None = None,
        font_size: int = FS_BODY,
        bold: bool = False,
    ) -> int:
        """Render body text and return actual bottom Y from measured content height.

        `bold` обязан совпадать с начертанием, которым текст будет нарисован:
        он и рисует, и меряет одной таблицей начертаний. Лид на «сцене» cleeq — жирный, и
        мерить его обычным начертанием значит занизить высоту абзаца.
        """
        left = MARGIN_X if x is None else x
        width = CONTENT_W if w is None else w
        avail = max(200000, min(max_h, CONTENT_BOTTOM - y))
        chunks = [c.strip() for c in re.split(r"\n+", _safe(text)) if c.strip()]
        if not chunks:
            return y
        # Prefer height-fit over crude char starvation (was ~200 chars at 420k emu).
        joined_raw = "\n".join(chunks[:8])
        # PDF-36 D.3 — before dropping sentences, shrink the font 1–2 pt
        # (min 9pt): full text at 10pt beats a cut paragraph at 11pt.
        if measure_text_height(joined_raw, width, font_size, line_spacing=1.2, bold=bold) > avail:
            # Кегль снижается по объявленной шкале, а не арифметикой:
            # `- 1` / `- 2` порождали 9,5 и 10,5, которых в шкале нет.
            for candidate in _scale_steps_below(font_size):
                if candidate < 9:
                    break
                if measure_text_height(joined_raw, width, candidate, line_spacing=1.2, bold=bold) <= avail:
                    font_size = candidate
                    break
        fitted = _fit_text_to_height(joined_raw, width, font_size, avail, line_spacing=1.2, bold=bold)
        fitted = _trim_dangling_tail(fitted)
        dangling = re.compile(
            r"(?:\bв\s+т\.?\s*ч\.?|\bс\s+[А-ЯA-Z]\.?|\b(?:как|что|чтобы|и|а|или|по|на|в|с)\b|,|;|—|–|-)\s*$",
            re.I,
        )
        if dangling.search(fitted) or (fitted and fitted[-1] not in ".!?…»)"):
            sentences = re.split(r"(?<=[.!?…])\s+", _safe(joined_raw))
            kept_s: list[str] = []
            for sent in sentences:
                trial = " ".join(kept_s + [sent]).strip()
                if measure_text_height(trial, width, font_size, line_spacing=1.2, bold=bold) <= avail:
                    kept_s.append(sent)
                else:
                    break
            # Drop a trailing incomplete clause (e.g. ends with «как»).
            while kept_s and (
                dangling.search(kept_s[-1]) or kept_s[-1][-1] not in ".!?…»)"
            ):
                kept_s.pop()
            if kept_s:
                fitted = " ".join(kept_s).strip()
            else:
                # Last resort: first sentence trimmed of dangling tail.
                first = _trim_dangling_tail(sentences[0] if sentences else fitted)
                fitted = first if first and not dangling.search(first) else _trim_dangling_tail(fitted)
        kept = [c for c in fitted.split("\n") if c.strip()] or ([fitted] if fitted else [])
        if not kept:
            return y
        joined = "\n".join(kept)
        needed = measure_text_height(joined, width, font_size, line_spacing=1.2, paragraph_spacing_pt=8, bold=bold)
        box_h = min(avail, max(needed + 40_000, int(font_size * EMU_PER_PT)))
        measured_lines, uncertain = _count_measured_lines(joined, width, font_size, bold)
        # Clipping = placed text does not fit the box. Fitting/truncating source is not layout overflow.
        clipped = needed > avail
        record_text_layout(
            page=self.page,
            name=f"orion_text_body_p{self.page}",
            role="text",
            font_family=FONT,
            font_size_pt=font_size,
            box_width=width,
            box_height=box_h,
            available_height=avail,
            required_height=needed,
            measured_lines=measured_lines,
            text_length=len(joined),
            clipped=clipped,
            measurement_uncertain=uncertain,
        )
        box = self.slide.shapes.add_textbox(Emu(left), Emu(y), Emu(width), Emu(box_h))
        try:
            box.name = f"orion_text_body_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for chunk in kept:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = chunk
            r.font.name = FONT
            r.font.bold = bold
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
        return y + box_h

    def content_card(
        self,
        *,
        title: str | None,
        text: str,
        x: int,
        y: int,
        width: int,
        min_h: int = 320_000,
        max_h: int = 2_400_000,
        padding: int = 100_000,
        tone: str = "neutral",
        title_size: int = 10,
        body_size: int = 11,
        skip_if_stub: bool = False,
    ) -> int:
        """Draw a content-sized card; clip text to fit; return actual bottom Y.

        `skip_if_stub`: не рисовать карточку, если от текста осталась только
        обрубленная часть. Карточка «Действие» со словом «Проверить» вместо
        рекомендации занимает полосу во всю ширину и не сообщает ничего —
        отсутствие карточки честнее (шаг 13, D3).
        """
        fill = {
            "accent": ACCENT_SOFT,
            "warn": WARN_BG,
            "risk": RISK_BG,
            "good": GOOD_BG,
        }.get(tone, CARD_BG)
        # Design v2: card titles pick up the tone colour instead of flat navy.
        title_color = {
            "accent": TONE_GOOD,
            "warn": TONE_WARN,
            "risk": TONE_RISK,
            "good": TONE_GOOD,
        }.get(tone, NAVY)
        title_s = _safe(title or "")
        # PDF-45 — keep structured theme newlines (quotes / scale) inside cards.
        body_s = _safe_preserve_breaks(text) if "\n" in str(text or "") else _safe(text)
        budget = max(200_000, min(max_h, CONTENT_BOTTOM - y))
        # Shrink padding on short cards so body text is not starved to a stub line.
        pad = padding
        if budget < 420_000:
            pad = min(padding, 60_000)
        elif budget < 560_000:
            pad = min(padding, 80_000)
        inner_w = max(120_000, width - 2 * pad)
        title_h = 0
        if title_s:
            # Заголовок карточки рисуется жирным (`r.font.bold = True` ниже),
            # поэтому и меряется жирным. Иначе высота карточки занижается, и
            # текст выходит за её нижний край.
            title_h = (
                measure_text_height(title_s, inner_w, title_size, line_spacing=1.15, bold=True)
                + 30_000
            )
        full_body_h = measure_text_height(body_s, inner_w, body_size, line_spacing=1.2) if body_s else 0
        needed = 2 * pad + title_h + full_body_h + 30_000
        # Short client phrases must stay complete; height estimator is conservative and
        # otherwise collapses actions like «Исключить из digital profile…» to one word.
        short_phrase = bool(body_s) and len(body_s) <= 110 and len(body_s.split()) <= 16
        if body_s and needed > budget and not short_phrase:
            body_budget = max(100_000, budget - 2 * pad - title_h)
            # PDF-36 D.3 — step the font down (−1/−2 pt, min 9pt) before any
            # text is dropped: +20–30% capacity, imperceptible to the eye.
            for candidate in _scale_steps_below(body_size):
                if candidate < 9:
                    break
                if measure_text_height(body_s, inner_w, candidate, line_spacing=1.2) <= body_budget:
                    body_size = candidate
                    break
            else:
                body_size = next_smaller_scale_pt(body_size) or FS_CAPTION
            body_h = measure_text_height(body_s, inner_w, body_size, line_spacing=1.2)
            if body_h > body_budget:
                # Prefer dropping whole structural lines over mid-phrase clips.
                if "\n" in body_s:
                    kept_lines: list[str] = []
                    for ln in body_s.split("\n"):
                        trial = "\n".join(kept_lines + [ln]) if kept_lines else ln
                        if measure_text_height(trial, inner_w, body_size, line_spacing=1.2) <= body_budget:
                            kept_lines.append(ln)
                        else:
                            break
                    body_s = "\n".join(kept_lines) if kept_lines else body_s.split("\n")[0]
                else:
                    body_s = _fit_text_to_height(body_s, inner_w, body_size, body_budget)
                # Строки могли быть отброшены — снять обещание, если продолжения нет.
                body_s = _close_dangling_lead_in(body_s)
                if skip_if_stub:
                    original = _safe(text)
                    fitted = _safe(body_s)
                    # Обрубок: от текста осталась малая часть и меньше строки.
                    if len(fitted) < min(60, int(len(original) * 0.5)):
                        return y
                body_h = measure_text_height(body_s, inner_w, body_size, line_spacing=1.2) if body_s else 0
        else:
            body_h = full_body_h
        h = max(min_h, min(budget, 2 * pad + title_h + body_h + 30_000))
        self.card(y, h=h, x=x, w=width, fill=fill)
        cy = y + pad
        if title_s:
            box = self.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(inner_w), Emu(max(title_h, 160_000)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title_s
            r.font.name = FONT
            r.font.bold = True
            r.font.size = Pt(title_size)
            r.font.color.rgb = title_color
            cy += title_h
        if body_s:
            rem = max(100_000, y + h - cy - pad)
            box = self.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(inner_w), Emu(rem))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body_s
            r.font.name = FONT
            r.font.size = Pt(body_size)
            r.font.color.rgb = BODY_COLOR
        return y + h

    def metric_chips(self, metrics: list[dict[str, Any]], x: int, y: int, width: int) -> int:
        """Render 2-column metric chips; return bottom Y."""
        items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:4]
        if not items:
            return y
        cols = 2
        gap = 70_000
        chip_w = (width - gap) // cols
        chip_h = 640_000
        row_y = y
        for idx, m in enumerate(items):
            col = idx % cols
            if idx > 0 and col == 0:
                row_y += chip_h + gap
            cx = x + col * (chip_w + gap)
            tone = str(m.get("tone") or "neutral")
            fill = {"risk": RISK_BG, "warn": WARN_BG, "good": GOOD_BG}.get(tone, CARD_BG)
            value_color = {
                "risk": TONE_RISK,
                "warn": TONE_WARN,
                "good": TONE_GOOD,
                "neutral": METRIC_ACCENT,
            }.get(tone, METRIC_ACCENT)
            value = _clip_words(_safe(m.get("value")), 36)
            label = _clip_words(_safe(m.get("label")), 28)
            self.card(row_y, h=chip_h, x=cx, w=chip_w, fill=fill)
            box = self.slide.shapes.add_textbox(
                Emu(cx + 50_000), Emu(row_y + 70_000), Emu(chip_w - 100_000), Emu(chip_h - 140_000)
            )
            tf = box.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = value
            r0.font.name = FONT
            r0.font.bold = True
            r0.font.size = Pt(14 if len(value) > 12 else 16)
            r0.font.color.rgb = value_color
            p1 = tf.add_paragraph()
            p1.space_before = Pt(4)
            r1 = p1.add_run()
            r1.text = label
            r1.font.name = FONT
            r1.font.size = Pt(9)
            r1.font.color.rgb = MUTED_COLOR
        rows = (len(items) + cols - 1) // cols
        return y + rows * chip_h + max(0, rows - 1) * gap

    def bullets(
        self,
        items: list[str],
        y: int,
        color: RGBColor = BODY_COLOR,
        max_items: int = 8,
        max_chars: int = 900,
        bottom: int | None = None,
        x: int | None = None,
        width: int | None = None,
    ) -> int:
        """Нарисовать список буллетов; вернуть нижнюю Y.

        `x`/`width` — колонка списка; по умолчанию вся полоса содержимого.
        Страница AI-ответов печатает тела под панелью в левой колонке, а
        справа стоит сайдбар полной высоты: без колонки список залезал бы
        под него, а с сайдбаром в треть высоты тот терял блоки (0053, стр. 62).

        `bottom` — граница, ниже которой рисовать нельзя: по умолчанию низ
        полосы содержимого, но на странице, где под списком стоит ещё карточка,
        её место обязан вычесть тот, кто эту карточку рисует. Иначе список
        занимает лист целиком, а карточке остаётся ноль.
        """
        col_x = MARGIN_X if x is None else x
        col_w = CONTENT_W if width is None else width
        dangling = re.compile(
            r"(?:\bв\s+т\.?\s*ч\.?|\bс\s+[А-ЯA-Z]\.?|\b[А-ЯA-Z]\.?|,|;|—|–|-|\()\s*$",
            re.I,
        )
        # Потеря содержимого на пути к листу считается здесь и только здесь:
        # опустевший после ножниц буллет молчал, и страница 30 живого отчёта
        # пришла без нарратива при чистой телеметрии.
        #
        # Чистятся **все** поданные элементы, а не первые `max_items`: раньше
        # хвост списка не доходил ни до листа, ни до счётчика потерь — подали
        # десять, нарисовали восемь, и узнать об этом было неоткуда. Мере тоже
        # нужны высоты всех поданных: по ним построитель и решает, что унести на
        # следующий лист.
        prepared: list[str] = []
        submitted = 0
        clipped_chars = 0
        for idx, b in enumerate(items):
            raw = _safe(b)
            if not raw:
                prepared.append("")
                continue
            submitted += 1
            clipped = _clip_structured_bullet(raw, max_chars)
            # Считается ровно то, что срезал бюджет. Сравнение идёт с телом
            # буллета без обрубков модели: чистка заглушек и снятие оборванного
            # хвоста ниже убирают уже сломанное, и мерить их как потерю значило
            # бы красить приёмку живого прогона там, где ничего не пропало.
            if clipped:
                clipped_chars += max(0, len(_structured_bullet_body(raw)) - len(clipped))
            # Drop incomplete trailing parentheticals from mid-clip
            # (e.g. «…спорах (в.» from «(в т.ч. материалы на …)»).
            clipped = re.sub(r"\s*\([^)\n]*$", "", clipped).rstrip(" :;")
            flat_tail = clipped.replace("\n", " ").rstrip()
            if dangling.search(flat_tail):
                punct = max(clipped.rfind(". "), clipped.rfind("! "), clipped.rfind("? "))
                if punct > 40:
                    clipped = clipped[: punct + 1].strip()
                else:
                    clipped = _trim_dangling_tail(clipped)
                clipped = re.sub(r"\s*\([^)\n]*$", "", clipped).rstrip(" :;")
            flat_tail = clipped.replace("\n", " ").rstrip()
            if dangling.search(flat_tail) or flat_tail.endswith(("—", "–", "-", "(")):
                # Prefer dropping the broken last structural line over failing render.
                lines = [ln.strip() for ln in clipped.split("\n") if ln.strip()]
                while lines and (
                    dangling.search(lines[-1])
                    or lines[-1].endswith(("—", "–", "-", "("))
                    or "(" in lines[-1] and ")" not in lines[-1]
                ):
                    lines.pop()
                clipped = "\n".join(lines).strip()
                flat_tail = clipped.replace("\n", " ").rstrip()
            if not clipped:
                prepared.append("")
                continue
            # Обрубок останавливает рендер только там, где его напечатают:
            # элемент за пределом `max_items` до листа не доходит, и падать на
            # нём значило бы ронять документ из-за того, чего в нём нет.
            if idx < max_items and (
                dangling.search(flat_tail) or flat_tail.endswith(("—", "–", "-", "("))
            ):
                raise RuntimeError(f"ORION bullet dangling on p{self.page}: {flat_tail[-48:]}")
            prepared.append(clipped)
        kept = [text for text in prepared[:max_items] if text]
        # PDF-46 I.2 — whole bullets only; never paint past CONTENT_BOTTOM.
        #
        # Мерка приведена к тому, что рисуется (шаг 16, 07.6). Прежняя мерила
        # весь блок одним куском на FS_BODY с межстрочным 1.2 и отбивкой 6 pt
        # между **всеми** строками, а рисуются они иначе: строки «Где видно…» и
        # «Всего по теме…» — кеглем подписи, внутренние отбивки 1 pt,
        # межстрочный 1.12. Сверху лежал ещё запас 1.18 «на всякий случай».
        #
        # Замер на финальном прогоне: оценка превышала факт в 1.7–1.8 раза, и
        # страница-продолжение, вмещающая три тематических блока, принимала два.
        # Отсюда пять страниц «Россия — резюме аудита», заполненных на 18–62 %.
        page_avail = max(0, (CONTENT_BOTTOM if bottom is None else bottom) - y - 120_000)
        measure_slack = 1.08

        def _bullet_block_height(blocks: list[str]) -> int:
            total = 0
            for b in blocks:
                parts = _split_structured_bullet(b) or [b]
                for li, line in enumerate(parts):
                    text = f"{BULLET_GLYPH} {line}" if li == 0 else f"   {line}"
                    # Начертание берётся то же, которым строка будет нарисована
                    # (ниже, в самом выводе — тот же `_bullet_line_style`).
                    # Здесь признак жирности выбрасывался: `_, _, size_pt`, — и
                    # заголовок темы, который рисуется жирным, мерился обычным.
                    # Жирное шире на 9,6–14,4%, поэтому блок буллетов считался
                    # ниже, чем печатается, и уезжал за нижнюю границу листа.
                    bold_line, _, size_pt = _bullet_line_style(line, is_first=(li == 0))
                    total += measure_text_height(
                        text,
                        col_w,
                        size_pt,
                        line_spacing=1.12,
                        paragraph_spacing_pt=0,
                        bold=bold_line,
                    )
                    space_before = 6 if li == 0 else 1
                    space_after = 6 if li == len(parts) - 1 else 1
                    total += int((space_before + space_after) * EMU_PER_PT)
            return int(total * measure_slack) + 60_000

        dropped_lines = 0
        while kept and _bullet_block_height(kept) > page_avail:
            if len(kept) == 1:
                # Drop whole structural lines from the sole bullet until it fits.
                parts = _split_structured_bullet(kept[0]) or [kept[0]]
                while len(parts) > 1 and _bullet_block_height(["\n".join(parts)]) > page_avail:
                    parts.pop()
                    dropped_lines += 1
                kept = ["\n".join(parts)] if parts and _bullet_block_height(["\n".join(parts)]) <= page_avail else []
                break
            kept.pop()
        # Один ответ на «сколько потеряно»: подано минус нарисовано. Прежде
        # слагаемых было три (опустевшие, срезанные по высоте, невлезшие), а
        # четвёртое — элементы за пределом `max_items` — не считалось вовсе.
        dropped_bullets = max(0, submitted - len(kept))
        # Выброшенное содержимое не молчит. Прежде лишние блоки исчезали из
        # отчёта без следа: подали четыре — нарисовалось два, и узнать об этом
        # было неоткуда. Пагинация обязана резать по страницам сама, а этот
        # цикл — последний рубеж, о срабатывании которого надо знать.
        #
        # Укороченный ножницами буллет — такая же потеря: до читателя не дошёл
        # хвост. Он записывается тем же входом, но не считается выброшенным
        # блоком: инспектор различает «блока нет вовсе» и «блок обрезан».
        if dropped_bullets or dropped_lines or clipped_chars:
            record_text_layout(
                page=self.page,
                name=(
                    f"orion_bullets_dropped_p{self.page}"
                    if dropped_bullets or dropped_lines
                    else f"orion_bullets_clipped_p{self.page}"
                ),
                role="bullets",
                font_family=FONT,
                font_size_pt=FS_BODY,
                box_width=col_w,
                box_height=page_avail,
                available_height=page_avail,
                required_height=_bullet_block_height([t for t in prepared if t]),
                measured_lines=len(kept),
                text_length=sum(len(_safe(b)) for b in items),
                clipped=True,
                measurement_uncertain=False,
                dropped_bullets=dropped_bullets,
                dropped_lines=dropped_lines,
            )
        # Мера пишется всегда, а не только при потере: построитель раскладывает
        # по ней и чистые страницы тоже — иначе он не знал бы, сколько на них
        # осталось места.
        record_bullet_measure(
            slide_key=self.slide_key,
            page=self.page,
            available_height=page_avail,
            max_items=max_items,
            item_heights=[_bullet_block_height([t]) if t else 0 for t in prepared],
            kept_items=len(kept),
            dropped_bullets=dropped_bullets,
            dropped_lines=dropped_lines,
        )
        if not kept:
            return y
        text_lines: list[str] = []
        for b in kept:
            parts = _split_structured_bullet(b) or [b]
            text_lines.append(f"{BULLET_GLYPH} {parts[0]}")
            text_lines.extend(f"   {p}" for p in parts[1:])
        text = "\n".join(text_lines)
        needed = _bullet_block_height(kept)
        avail = max(200_000, min(needed, page_avail))
        # Hard cap: textbox bottom must stay at/above CONTENT_BOTTOM.
        if y + avail > CONTENT_BOTTOM:
            avail = max(0, CONTENT_BOTTOM - y)
        if avail < 200_000:
            return y
        box = self.slide.shapes.add_textbox(Emu(col_x), Emu(y), Emu(col_w), Emu(avail))
        tf = box.text_frame
        tf.word_wrap = True
        first_para = True
        for bi, bullet in enumerate(kept):
            lines = _split_structured_bullet(bullet) or [bullet]
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if first_para else tf.add_paragraph()
                first_para = False
                p.space_before = Pt(6 if li == 0 else 1)
                p.space_after = Pt(6 if li == len(lines) - 1 else 1)
                p.line_spacing = 1.12
                bold, line_color, size_pt = _bullet_line_style(line, is_first=(li == 0))
                # Fallback color arg only for flat single-line bullets.
                if len(lines) == 1 and color != BODY_COLOR:
                    line_color = color
                if li == 0:
                    # Маркер — отдельный прогон, чтобы квадрат был цветным, а
                    # текст пункта — чернилами. Замер этого не касается: строка
                    # меряется целиком, вместе с маркером и пробелом.
                    marker = p.add_run()
                    marker.text = f"{BULLET_GLYPH} "
                    marker.font.name = FONT
                    marker.font.bold = bold
                    marker.font.size = Pt(size_pt)
                    marker.font.color.rgb = ACCENT if bi % 2 == 0 else VIOLET
                r = p.add_run()
                r.text = line if li == 0 else f"   {line}"
                r.font.name = FONT
                r.font.bold = bold
                r.font.size = Pt(size_pt)
                r.font.color.rgb = line_color
        return y + min(avail, needed + 60_000)


def disable_shape_shadow(shape) -> None:
    """Погасить мягкую тень, которую LibreOffice рисует фигуре по умолчанию.

    Одного пустого `<a:effectLst/>` мало: python-pptx кладёт каждой фигуре
    ссылку на эффект темы (`<a:effectRef idx="2">`), и конвертер выполняет
    именно её. Тень размывается примерно на десять точек за границей фигуры —
    для карточки посреди листа это незаметно, а для сцены, доходящей до нижней
    границы полосы контента, это чернила ниже границы, и растровая проверка
    видит их (ADR-0007).
    """
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass
    style = shape._element.find(qn("p:style"))
    if style is None:
        return
    ref = style.find(qn("a:effectRef"))
    if ref is not None:
        ref.set("idx", "0")


def _asset_map(payload: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {str(a.get("assetRef")): a for a in payload.get("assets") or []}


def _resolve_image_bytes(asset: dict[str, Any] | None) -> bytes | None:
    """Load PNG/JPEG bytes from inline imageData or DATA_ROOT storageKey."""
    if not asset:
        return None
    img_data = asset.get("imageData")
    if img_data:
        try:
            raw = base64.b64decode(str(img_data))
            if len(raw) > 500:
                return raw
        except Exception:  # noqa: BLE001
            pass
    storage_key = str(asset.get("storageKey") or "").strip().lstrip("/")
    if storage_key:
        data_root = Path(os.environ.get("DATA_ROOT", "/data"))
        # storage keys are relative to digital-profile root; DATA_ROOT usually mounts that root.
        candidates = [
            data_root / storage_key,
            data_root / "digital-profile" / storage_key,
            Path(storage_key),
        ]
        for path in candidates:
            try:
                if path.is_file() and path.stat().st_size > 500:
                    return path.read_bytes()
            except OSError:
                continue
    return None


def _embed_image_contain(ctx: _Ctx, asset: dict[str, Any] | None, y: int, h: int = 4800000) -> bool:
    """Place image inside (MARGIN_X, y, CONTENT_W, h) preserving aspect ratio."""
    if not asset:
        ctx.body("Визуальный материал недоступен для данного раздела.", y)
        return False
    raw = _resolve_image_bytes(asset)
    if not raw:
        title = _safe(asset.get("title") or "Источник")
        domain = _safe(asset.get("caption") or asset.get("storageKey") or "")
        ctx.card(y, h)
        ctx.body(
            f"{title}\n{domain}\nИзображение недоступно — показаны источник и описание.",
            y + 120000,
            max_h=h - 200000,
        )
        return False
    box_w, box_h = CONTENT_W, h
    iw, ih = box_w, box_h
    if Image is not None:
        try:
            with Image.open(io.BytesIO(raw)) as im:
                iw, ih = im.size
        except Exception:  # noqa: BLE001
            pass
    scale = min(box_w / max(iw, 1), box_h / max(ih, 1))
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    left = MARGIN_X + (box_w - draw_w) // 2
    top = y + (box_h - draw_h) // 2
    stream = io.BytesIO(raw)
    ctx.slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(draw_w), height=Emu(draw_h))
    return True


def _embed_image(ctx: _Ctx, asset: dict[str, Any] | None, y: int, h: int = 4800000) -> None:
    _embed_image_contain(ctx, asset, y, h)


def _first_visual_asset(refs: list[Any], assets: dict[str, dict[str, Any]]) -> dict[str, Any] | None:
    for ref in refs:
        asset = assets.get(str(ref))
        if asset and _resolve_image_bytes(asset):
            return asset
    return None




def _rounded_portrait_png(
    raw: bytes, size_px: int = 900, radius_ratio: float = 0.14
) -> bytes | None:
    """Квадратное фото со скруглёнными углами — портрет на обложке.

    Широкая сетка изображений (~1200×680) — это коллаж; из него берётся область
    первой карточки, иначе на обложку попал бы кусок сетки, а не лицо.
    """
    if Image is None:
        return None
    try:
        from PIL import ImageDraw
    except ImportError:  # pragma: no cover
        return None
    try:
        with Image.open(io.BytesIO(raw)) as src:
            im = src.convert("RGBA")
        w, h = im.size
        if w >= 900 and h >= 500 and (w / max(h, 1)) > 1.35:
            sx, sy = w / 1200.0, h / 680.0
            box = (
                int(48 * sx),
                int(88 * sy),
                int((48 + 344) * sx),
                int((88 + 144) * sy),
            )
            im = im.crop(box)
            w, h = im.size
        side = min(w, h)
        if side < 32:
            return None
        left = (w - side) // 2
        top = (h - side) // 2
        im = im.crop((left, top, left + side, top + side)).resize(
            (size_px, size_px), Image.Resampling.LANCZOS
        )
        radius = max(8, int(size_px * radius_ratio))
        mask = Image.new("L", (size_px, size_px), 0)
        ImageDraw.Draw(mask).rounded_rectangle(
            (0, 0, size_px - 1, size_px - 1), radius=radius, fill=255
        )
        out = Image.new("RGBA", (size_px, size_px), (0, 0, 0, 0))
        out.paste(im, (0, 0), mask)
        buf = io.BytesIO()
        out.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:  # noqa: BLE001
        return None


def _embed_cover_portrait(ctx: _Ctx, asset: dict[str, Any] | None) -> bool:
    """Портрет субъекта на обложке: квадрат в зелёной и фиолетовой рамке."""
    raw = _resolve_image_bytes(asset)
    if not raw:
        return False
    png = _rounded_portrait_png(raw)
    if not png:
        return False
    side = 4_000_000
    left = 7_050_000
    top = 1_450_000
    for pad, color in ((220_000, VIOLET), (110_000, ACCENT)):
        d = side + pad * 2
        frame = ctx.slide.shapes.add_shape(5, Emu(left - pad), Emu(top - pad), Emu(d), Emu(d))
        frame.fill.background()
        frame.line.color.rgb = color
        frame.line.width = Pt(2.0)
        try:
            frame.adjustments[0] = 0.14
        except Exception:  # noqa: BLE001
            pass
    stream = io.BytesIO(png)
    ctx.slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(side), height=Emu(side))
    return True
