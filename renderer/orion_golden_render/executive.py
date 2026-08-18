"""Executive dashboard, risk matrix grid, and profile overview templates."""

from __future__ import annotations

import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from .common import (
    FS_SUBTITLE,
    ACCENT,
    ACCENT_SOFT,
    BODY_COLOR,
    CARD_BG,
    CARD_BORDER,
    CONTENT_BOTTOM,
    CONTENT_W,
    EMU_PER_PT,
    FONT,
    FS_BODY,
    FS_CAPTION,
    FS_SECTION,
    GOOD_BG,
    MARGIN_X,
    MUTED_COLOR,
    NAVY,
    RISK_BG,
    TITLE_COLOR,
    WARN_BG,
    WHITE,
    _Ctx,
    _bullet_line_style,
    _clip_structured_bullet,
    _clip_words,
    _close_dangling_lead_in,
    _safe,
    _safe_preserve_breaks,
    _fit_lines_to_height,
    _split_structured_bullet,
    _structured_bullet_body,
    _wrapped_line_count,
    font_line_step_emu,
    measure_text_height,
    record_text_layout,
)
from .layout_cleeq import (
    bars_color,
    content_stage,
    draw_level_bars,
    render_action_block,
    render_hero_metrics_row,
    stage_heading,
    tone_to_bars,
)
from .visual import (
    _render_kpi_cards,
    _render_status_badge,
    _title_line_estimate,
    _tone_fill,
    _tone_value_color,
)

def _render_executive_dashboard(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title, 240000, NAVY, FS_SECTION)
    raw_narrative = str(slide.get("narrative") or "")
    paras = [p.strip() for p in re.split(r"\n+", raw_narrative) if p.strip()][:3]
    if not paras and raw_narrative:
        # Split long single paragraph into ~3 chunks on sentence boundaries.
        parts = re.split(r"(?<=[.!?])\s+", _safe(raw_narrative))
        paras = []
        buf = ""
        for part in parts:
            trial = f"{buf} {part}".strip() if buf else part
            if len(trial) > 280 and buf:
                paras.append(buf)
                buf = part
            else:
                buf = trial
            if len(paras) >= 3:
                break
        if buf and len(paras) < 3:
            paras.append(buf)
        paras = paras[:3] or [_clip_words(_safe(raw_narrative), 420)]
    # Резюме в языке cleeq: ключевые цифры ведут страницу, всё остальное лежит
    # на одной белой сцене — лид, темы, следующий шаг. Прежняя раскладка (текст
    # слева, метрики справа, две карточки внизу) держала одно и то же
    # содержимое в четырёх рамках разного тона: рамки читались как иерархия,
    # которой между ними нет.
    metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
    metrics_bottom = y
    if metrics:
        metrics_bottom = render_hero_metrics_row(
            ctx, metrics[:4], MARGIN_X, y, CONTENT_W, tone_value_color=_tone_value_color
        )
        y = metrics_bottom + 140_000
    content_stage(ctx, y, top=metrics_bottom + 40_000 if metrics else None, corner_marks=True)
    for idx, para in enumerate(paras):
        # PDF-36 D.3 — высоту считает мерка, а не предварительная обрезка по
        # знакам: 420 знаков морили строку §7.2 и выводы лида.
        y = ctx.body(_safe(para), y, max_h=1_100_000, bold=idx == 0) + 30_000
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)][:2]
    actions = [a for a in (slide.get("actions") or []) if isinstance(a, dict)][:1]
    theme_bullets: list[str] = []
    for finding in findings:
        detail = _safe(finding.get("detail") or "")
        headline = _safe(finding.get("headline") or "")
        # Prefer detail; avoid duplicating headline when detail already starts with it.
        if detail and headline and detail.lower().startswith(headline.lower()[:24].lower()):
            text = detail
        elif detail:
            text = detail
        else:
            text = headline
        # Keep structured newlines; ctx.bullets drops whole lines if tight.
        text = _clip_structured_bullet(text, 900) or text
        if text:
            theme_bullets.append(text)
    if theme_bullets:
        y = stage_heading(ctx, "Коротко по итогам аудита", y)
        y = ctx.bullets(theme_bullets, y, max_items=len(theme_bullets), max_chars=900) + 40_000
    if actions:
        label = _safe(actions[0].get("label"))
        # Полная фраза или ничего: оборванная рекомендация не рекомендация.
        if len(label) > 280:
            punct = max(label.rfind(". "), label.rfind("! "), label.rfind("? "), label.rfind("; "))
            if punct > 60:
                label = label[: punct + 1].strip()
        render_action_block(ctx, label, y, max_h=1_000_000, heading="Следующий шаг")


#: Межстрочный, который выставляется абзацам тела карточки при отрисовке.
_CARD_LINE_SPACING = 1.12
#: Единственный запас карточной меры.
#:
#: Общая `measure_text_height` держит свой запас ×1,18 «на всякий случай», а
#: карточка брала сверху ещё ×1,12 — вместе с межстрочным 1,2 вместо
#: нарисованного 1,12 это давало треть лишней высоты: на растре эталона под
#: последней строкой карточки оставалось 41–42 px пустоты при паддинге 16 px.
#: Здесь запас один и назван: перенос PPTX чуть агрессивнее замера PIL, и
#: несколько процентов на это нужны — но именно несколько.
_CARD_MEASURE_SLACK = 1.06


def _card_paragraph_spacing(index: int) -> tuple[int, int]:
    """Отбивки абзаца строки тела карточки в пунктах: до и после.

    Объявлены отдельно, потому что их читают двое — замер и вывод; разойдутся
    они, и высота карточки перестанет совпадать с нарисованной.
    """
    return (0 if index == 0 else 1), 1


def _card_line_style(line: str, index: int, detail_font: float) -> tuple[bool, RGBColor, float]:
    """Начертание, цвет и кегль строки тела карточки.

    Один ответ на замер и на вывод. Разъедутся стороны хоть на одной ветке —
    и высота поедет на конкретных карточках: мета-строка («Всего по теме:»,
    «Где видно:») рисуется кеглем подписи, а мерилась в составе тела кеглем
    основного текста.
    """
    bold, line_color, size_pt = _bullet_line_style(line, is_first=(index == 0))
    if index == 0 and line.startswith("«"):
        return False, BODY_COLOR, float(detail_font)
    if not bold:
        return False, line_color, float(detail_font) if line_color == BODY_COLOR else FS_CAPTION
    return bold, line_color, size_pt


def _card_body_height(detail: str, text_w: int, detail_font: float) -> int:
    """Высота тела карточки — по тем строкам, кеглям и отбивкам, какими оно рисуется.

    Общая `measure_text_height` сюда не годится: она меряет тело одним куском
    одним кеглем, межстрочным 1,2 и отбивкой абзаца 6 pt, а рисуется оно
    построчно — своим кеглем на строку, межстрочным 1,12 и отбивками 0/1 и
    1 pt. Умножать её результат на ещё один запас тем более незачем: разница
    была третью высоты карточки, и «лишняя» карточка уезжала на следующую
    страницу с пустого листа.

    Умолчания `measure_text_height` при этом не трогаются: у неё десятки
    потребителей на всех шаблонах, и смена умолчаний — пересчёт всей деки.
    Путь буллетов решён так же: локальной мерой «как рисуется».
    """
    lines = _split_structured_bullet(detail) or ([detail] if detail else [])
    total = 0
    for index, line in enumerate(lines):
        bold, _color, size_pt = _card_line_style(line, index, detail_font)
        total += _wrapped_line_count(line, text_w, size_pt, bold) * font_line_step_emu(
            size_pt, _CARD_LINE_SPACING, bold
        )
        total += sum(_card_paragraph_spacing(index)) * EMU_PER_PT
    return int(total * _CARD_MEASURE_SLACK)


def _dropped_line_count(before: str, after: str) -> int:
    """Сколько строк тела не доехало до листа при резе `before` → `after`.

    `before` — то, что обязано было доехать: тело уже без обрубков модели и без
    повисшего ввода. Уборка этих заглушек — не потеря содержимого, и посчитать
    её потерей значит поднять CRITICAL там, где ничего не пропало: обоим
    инспекторам достаточно ненулевого `droppedLines`, чтобы остановить выдачу
    отчёта. Тем же правилом живёт учёт потерь в `ctx.bullets`.

    Укороченная строка считается снятой целиком: её хвост до клиента не дошёл
    так же, как целая строка, а «полстроки» телеметрии выразить нечем.
    """
    before_lines = [ln for ln in before.split("\n") if ln.strip()]
    after_lines = [ln for ln in after.split("\n") if ln.strip()]
    dropped = max(0, len(before_lines) - len(after_lines))
    if not dropped and len("".join(after_lines)) < len("".join(before_lines)):
        return 1
    return dropped


def _render_risk_matrix_grid(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title or "Матрица рисков", 240000, NAVY, FS_SECTION)
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)]
    if not findings:
        bullets = [_safe(b) for b in slide.get("bullets") or [] if _safe(b)]
        findings = [{"headline": b.split("—")[0].strip()[:60], "detail": b, "tone": "warn"} for b in bullets]
    badge_w = 1_900_000
    # Сколько карточек на листе — решает пагинация построителя (реестр
    # шаблонов). Рендерер отвечает только за то, влезает ли карточка целиком;
    # свой срез списка был вторым ответом на чужой вопрос.
    drawn = 0
    drawn_lines = 0
    dropped_lines = 0
    stepped_down = False
    grid_top = y
    for finding in findings:
        tone = str(finding.get("tone") or "warn")
        pill = _safe(finding.get("status") or finding.get("severity") or "")
        if len(pill) > 28:
            pill = _clip_words(pill, 28)
        headline = _safe(finding.get("headline") or "Тема")
        if len(headline) > 72:
            headline = _clip_words(headline, 72)
        # Переносы — структура карточки, а не оформление: по ним отделяется
        # строка «Что делать» и по ним же меряется высота. `_safe` схлопывал их
        # первым делом, вся строчная логика ниже становилась мёртвой, карточка
        # раздувалась и последняя на странице не рисовалась вовсе.
        raw_detail = _safe_preserve_breaks(finding.get("detail") or "")
        detail = _clip_structured_bullet(raw_detail, 900)
        # Бюджет знаков рендерера — такой же рез, как страховка по высоте, и
        # молчать о нём нельзя. Живой построитель в него не упирается (его
        # бюджет — 320 знаков), а замороженный пак старого прогона и чужой пак
        # упираются: с июльских карточек уезжала треть текста при чистой
        # телеметрии.
        dropped_lines += _dropped_line_count(_structured_bullet_body(raw_detail), detail)
        # PDF-40 G.1b — drop a leading «Theme» line that duplicates the headline.
        detail_lines = _split_structured_bullet(detail) or ([detail] if detail else [])
        if detail_lines and detail_lines[0].startswith("«"):
            theme_core = detail_lines[0].strip("«»").split("—")[0].strip().lower()
            head_core = headline.split("—")[0].strip().lower()
            if theme_core and (theme_core in head_core or head_core in theme_core):
                detail_lines = detail_lines[1:]
                detail = "\n".join(detail_lines)
        # Мерить надо ровно то, что будет нарисовано. При отрисовке текст
        # прогоняется через _split_structured_bullet, и тот может раздробить
        # его сильнее, чем было при замере: строк становится больше, высота —
        # выше, и текст вылезал за карточку, перекрываясь следующей
        # (шаг 13, D1). Повторное дробление уже разбитого текста идемпотентно.
        if detail:
            detail = "\n".join(_split_structured_bullet(detail) or [detail])
        # Prefer complete source text; only sentence-fit if height is tight.
        marker = _safe(finding.get("manualReview") or "")
        # Embed "requires review" into status instead of a cramped footer.
        if marker and re.search(r"требует", pill or "", re.I) is None and "проверк" in marker.lower():
            if pill and "проверк" not in pill.lower():
                pill = pill  # keep level; marker dropped from footer
        text_w = CONTENT_W - badge_w - 220_000 if pill else int(CONTENT_W * 0.92)
        left = MARGIN_X + 100_000
        pad_y = 100_000
        # Заголовок карточки риска рисуется жирным (`r.font.bold = True` ниже),
        # поэтому и меряется жирным: от `headline_h` зависит и высота карточки,
        # и решение «не рисовать карточку, которая не влезает над колонтитулом».
        headline_h = measure_text_height(headline, text_w, 13, line_spacing=1.15, bold=True)
        # PDF-36 E.2 — a card that cannot fit whole above the footer is not
        # drawn at all (TS pagination carries the rest to the continuation
        # page); a half-card bleeding into the footer reads as a defect.
        remaining = CONTENT_BOTTOM - y - 100_000
        if remaining < max(520_000, pad_y + headline_h + pad_y):
            break
        detail_font = FS_BODY
        if detail:
            # Grow card to fit complete detail when possible.
            detail_h = _card_body_height(detail, text_w, detail_font)
            needed = pad_y + headline_h + 40_000 + detail_h + pad_y
            max_h = remaining
            if needed > max_h:
                # PDF-36 D.3 — font step-down before dropping lines.
                for candidate in (FS_CAPTION,):
                    cand_h = _card_body_height(detail, text_w, candidate)
                    trial = pad_y + headline_h + 40_000 + cand_h + pad_y
                    if trial <= max_h:
                        detail_font = candidate
                        detail_h = cand_h
                        needed = trial
                        # Защита сработала: содержимое доехало целиком, но
                        # клиент видит карточку мельче соседних. Это и есть
                        # `clipped` по определению §8 — и до сих пор оно было
                        # немым, потому что признак выставлялся по факту
                        # потери, то есть означал прямо противоположное.
                        stepped_down = True
                        break
            # Карточка либо влезает целиком, либо не рисуется и уходит
            # пагинации (E.2). Прежде здесь стоял третий ответ на вопрос «что
            # в карточке»: цикл выбрасывал мета-строки и попами с конца снимал
            # строки, пока карточка не влезет, — молча. На живой странице он
            # съел «Что делать», а телеметрия ничего об этом не сказала.
            if needed > max_h:
                break
            h = min(max_h, needed)
        else:
            h = max(420_000, pad_y + headline_h + pad_y)
        h = max(420_000, min(h, remaining))
        if y + h > CONTENT_BOTTOM:
            break
        ctx.card(y, h=h, fill=_tone_fill(tone))
        # Счёт нарисованного стоит здесь, в точке рисования: посчитанное в
        # другом месте разойдётся с нарисованным, и телеметрия снова перестанет
        # быть свидетельством.
        drawn += 1
        # Headline
        box = ctx.slide.shapes.add_textbox(Emu(left), Emu(y + pad_y), Emu(text_w), Emu(max(headline_h + 20_000, 160_000)))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = headline
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_SUBTITLE)
        r.font.color.rgb = NAVY
        text_y = y + pad_y + headline_h + 30_000
        if detail:
            rem = max(120_000, y + h - text_y - pad_y)
            # Страховка на случай остаточной погрешности замера: короб текста
            # не может быть выше карточки, поэтому лишние строки отбрасываются
            # целиком, а не обрезаются посередине. Мера у неё та же, которой
            # посчитана карточка, — иначе она снимала бы строки, которые на
            # листе помещались.
            fitted = _fit_lines_to_height(
                detail, rem, lambda text: _card_body_height(text, text_w, detail_font)
            )
            # Снятая здесь строка до клиента не дошла — считается как потеря и
            # блокирует выдачу. Пустое тело при этом не пропускает карточку:
            # она уже нарисована и заняла свою высоту, а `continue` уводил от
            # бейджа и от приращения y, и следующая карточка ложилась поверх.
            dropped_lines += _dropped_line_count(detail, fitted)
            # Уборка повисшего ввода — после счёта. Двоеточие без продолжения
            # снимается и тогда, когда по высоте влезло всё: посчитанная
            # потерей, эта уборка останавливала выдачу здорового отчёта.
            detail = _close_dangling_lead_in(fitted)
        # Условие повторяется: на фиттере тело могло опустеть целиком.
        if detail:
            box = ctx.slide.shapes.add_textbox(Emu(left), Emu(text_y), Emu(text_w), Emu(rem))
            tf = box.text_frame
            tf.word_wrap = True
            # PDF-38 F.1 / G.1 — hierarchical detail inside the card bounds.
            draw_lines = _split_structured_bullet(detail) or [detail]
            first = True
            for li, line in enumerate(draw_lines):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                space_before, space_after = _card_paragraph_spacing(li)
                p.space_before = Pt(space_before)
                p.space_after = Pt(space_after)
                p.line_spacing = _CARD_LINE_SPACING
                bold, line_color, size_pt = _card_line_style(line, li, detail_font)
                # Строки переноса, а не структурные: поле схемы называется
                # «сколько строк намерено», и все остальные его производители
                # пишут туда именно перенос.
                drawn_lines += _wrapped_line_count(line, text_w, size_pt, bold)
                r = p.add_run()
                r.text = line
                r.font.name = FONT
                r.font.bold = bold
                r.font.size = Pt(size_pt)
                r.font.color.rgb = line_color
        if pill:
            # E.1 — badge height from conservative wrap estimate; long Russian
            # labels like «Требует подтверждения» always get ≥2 lines of room
            # so the plate border never bisects the text (PDF 36 p5–7).
            bx = MARGIN_X + CONTENT_W - badge_w - 80_000
            by = y + pad_y
            inner_w = badge_w - 140_000
            pill_size = FS_CAPTION
            lines = _title_line_estimate(pill, inner_w, pill_size, max_lines=3)
            if len(pill) > 14 or " " in pill:
                lines = max(lines, 2)
            line_h = int(pill_size * 12_700 * 1.35)
            filled = tone_to_bars(tone, pill)
            # Шкале степени нужна своя полоса под словом уровня: делений пять,
            # высота 90 000 плюс отбивка. Без этого запаса шкала ложилась прямо
            # на слово и читалась как зачёркивание.
            bars_h = 130_000
            need = lines * line_h + 60_000 + bars_h
            bh = min(max(420_000, need + 140_000), max(420_000, h - 2 * pad_y))
            # Soft tone fill, no hard border — a stroked plate was reading as
            # a mid-text strike-through after LibreOffice PDF conversion.
            ctx.card(by, h=bh, x=bx, w=badge_w, fill=_tone_fill(tone), border=None)
            b = ctx.slide.shapes.add_textbox(
                Emu(bx + 70_000), Emu(by + 70_000), Emu(inner_w), Emu(bh - 140_000)
            )
            btf = b.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = pill
            br.font.name = FONT
            br.font.bold = True
            br.font.size = Pt(pill_size)
            # Цвет слова — от самой степени, а не от тона карточки: тон у темы
            # часто не задан, и общий цвет метрик (зелёный) на слове
            # «Критический» читался бы как «всё в порядке».
            br.font.color.rgb = bars_color(filled)
            # Шкала степени cleeq: пять делений под словом уровня. Слово
            # остаётся — шкала показывает то же самое, но с одного взгляда и
            # в сравнении с соседними темами.
            bars_w = 5 * 120_000 + 4 * 30_000
            bars_y = by + 70_000 + lines * line_h + 40_000
            if bars_y + 90_000 <= by + bh - 40_000:
                draw_level_bars(
                    ctx,
                    bx + (badge_w - bars_w) // 2,
                    bars_y,
                    filled=filled,
                    hot=bars_color(filled),
                )
        y += h + 50_000
        if y > CONTENT_BOTTOM - 360_000:
            break
    # Карточка, не влезшая целиком, не рисуется — это решение рендерера, и о
    # нём обязан узнать контур: запись читает инспектор геометрии и превращает
    # ненулевую потерю в CONTENT_DROPPED_BY_RENDERER. Пока записи не было,
    # страница теряла тему при зелёной приёмке. Запись пишется всегда, а не
    # только при потере: её отсутствие приёмка считает отказом, иначе «нет
    # записей» неотличимо от «потерь нет».
    used = max(0, y - grid_top)
    record_text_layout(
        page=ctx.page,
        name=f"orion_risk_matrix_p{ctx.page}",
        role="cards",
        font_family=FONT,
        font_size_pt=FS_BODY,
        box_width=CONTENT_W,
        box_height=used,
        available_height=max(0, CONTENT_BOTTOM - grid_top),
        required_height=used,
        measured_lines=drawn_lines,
        text_length=sum(len(_safe(f.get("detail") or "")) for f in findings),
        # `clipped` и `dropped_*` отвечают на разные вопросы. Первый — «защита
        # сработала»: карточке понижен кегль, но содержимое доехало целиком.
        # Второй — «содержимое потеряно». Пока клип выставлялся по факту
        # потери, он значил ровно противоположное записанному в §8, а
        # понижение кегля не было слышно вовсе.
        clipped=stepped_down,
        dropped_bullets=len(findings) - drawn,
        dropped_lines=dropped_lines,
        drawn_blocks=drawn,
    )


def _render_profile_overview(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title, 240000, NAVY, FS_SECTION)
    y = _render_status_badge(
        ctx, slide.get("statusBadge") if isinstance(slide.get("statusBadge"), dict) else None, MARGIN_X, y, CONTENT_W
    )
    y += 80_000
    metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
    ru = [m for m in metrics if "RU" in _safe(m.get("label")).upper() or "Росс" in _safe(m.get("label"))]
    uae = [m for m in metrics if "UAE" in _safe(m.get("label")).upper() or "ОАЭ" in _safe(m.get("label"))]
    other = [m for m in metrics if m not in ru and m not in uae]
    half = (CONTENT_W - 100_000) // 2
    if ru or uae:
        left_metrics = (ru or metrics[:4])[:4]
        right_metrics = (uae or metrics[4:8])[:4]
        hdr_h = 320_000
        for label, xx in (("Россия", MARGIN_X), ("ОАЭ", MARGIN_X + half + 100_000)):
            ctx.card(y, h=hdr_h, x=xx, w=half, fill=ACCENT_SOFT)
            box = ctx.slide.shapes.add_textbox(Emu(xx + 90_000), Emu(y + 90_000), Emu(half - 180_000), Emu(160_000))
            rr = box.text_frame.paragraphs[0].add_run()
            rr.text = label
            rr.font.name = FONT
            rr.font.bold = True
            rr.font.size = Pt(FS_SUBTITLE)
            rr.font.color.rgb = NAVY
        y += hdr_h + 80_000
        left_bottom = _render_kpi_cards(ctx, left_metrics, MARGIN_X, y, half, cols=2)
        right_bottom = _render_kpi_cards(ctx, right_metrics, MARGIN_X + half + 100_000, y, half, cols=2)
        y = max(left_bottom, right_bottom) + 100_000
    else:
        y = _render_kpi_cards(ctx, metrics[:8], MARGIN_X, y, CONTENT_W, cols=4) + 100_000
    if other:
        y = _render_kpi_cards(ctx, other[:4], MARGIN_X, y, CONTENT_W, cols=4) + 80_000
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)][:3]
    for finding in findings:
        if y > CONTENT_BOTTOM - 360_000:
            break
        y = ctx.content_card(
            title=_clip_words(_safe(finding.get("headline")), 48),
            text=_safe(finding.get("detail") or ""),
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=280_000,
            max_h=700_000,
            tone=str(finding.get("tone") or "neutral"),
            title_size=11,
            body_size=11,
        )
        y += 60_000



