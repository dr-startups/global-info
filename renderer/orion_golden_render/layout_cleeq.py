"""Компоновки cleeq: плавающая сцена, квадратные маркеры, шкала уровня.

Перенесено из ветки `feature/checkpoint` (1a949bd — «cleeq-стиль отчёта»).
Там же взяты пропорции: страница мятная, содержимое лежит на белой «сцене» с
мягкой тенью, маркер — квадрат, а не точка, степень риска показана шкалой из
пяти делений.

Перенесён визуальный язык, а не его геометрия. В исходной ветке новые
компоновки резали текст (`_clip_words` до влезания, `bullets[:5]`,
фиксированная высота строки таблицы) и писали кегли литералами — 40, 28, 12,5.
На отрисовке эталонной деки той ветки это даёт `overflow: 5`, обрывы фраз на
полуслове и шестнадцать кеглей вместо восьми. Здесь то же самое выглядит так:

- сцена рисуется **под** обычным потоком текста и выступает за колонку наружу,
  поэтому ширина строки и ёмкость страницы не меняются ни на EMU. Ёмкость листа
  меряет сам код отрисовки (`ctx.bullets`), и построитель спрашивает её мерным
  прогоном — подвинуть содержимое значит потерять его;
- кегли берутся из объявленной шкалы (`TYPE_SCALE_PT`): литералы в коде и дают
  шестнадцать кеглей на деку вместо восьми;
- текст меряется тем начертанием, которым рисуется: жирный лид, померенный
  обычным, занижает высоту абзаца и роняет содержимое ниже границы контента.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from .common import (
    ACCENT,
    CARD_BORDER,
    CONTENT_BOTTOM,
    CONTENT_W,
    FONT,
    FS_BODY,
    FS_CAPTION,
    FS_SECTION,
    FS_TITLE,
    MARGIN_X,
    MUTED_COLOR,
    STAGE_SHADOW,
    TONE_RISK,
    TONE_WARN,
    VIOLET,
    WHITE,
    _Ctx,
    _clip_words,
    disable_shape_shadow,
    _fit_text_to_height,
    _safe,
    measure_text_height,
)

#: Насколько сцена выступает за текстовую колонку. Внутренние поля карточки
#: получаются из этого выступа, а колонка остаётся прежней ширины.
#:
#: Больше брать нельзя: тень сцены — это чернила на растре, и растровая
#: проверка считает дефектом всё, что заходит за 0,6 бокового поля
#: (288 000 EMU от края). При выступе 100 000 правый край тени приходится на
#: 11 359 320 при пороге 11 416 320 — запас 57 000 EMU.
STAGE_BLEED = 100_000

#: Смещение тени. Тень — единственное, чем белая сцена отличается от мятного
#: листа настолько, чтобы её было видно: белое по #F6F8F4 глаз почти не ловит.
STAGE_SHADOW_DX = 35_000
STAGE_SHADOW_DY = 45_000

#: Высота ряда метрик. Совпадает с `_render_kpi_cards`: hero-плитка отличается
#: кеглем цифры и шириной, но не высотой ряда — иначе содержимое ниже уезжает
#: за нижнюю границу листа.
METRIC_ROW_H = 780_000


def draw_stage(ctx: _Ctx, x: int, y: int, w: int, h: int) -> None:
    """Белая сцена с мягкой тенью — плоскость, на которой лежит содержимое.

    Тень у сцены ровно одна: сдвинутый серый прямоугольник, чью геометрию мы
    знаем. Собственную мягкую тень фигуры LibreOffice рисует по умолчанию, и
    размывается она примерно на десять точек **за** границей фигуры — то есть
    ниже границы контентной области, где растровая проверка справедливо видит
    чернила. Гасим её здесь, а не поднимаем сцену: подниматься ей некуда, текст
    доходит до самого низа полосы.
    """
    shadow = ctx.slide.shapes.add_shape(
        5, Emu(x + STAGE_SHADOW_DX), Emu(y + STAGE_SHADOW_DY), Emu(w), Emu(h)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = STAGE_SHADOW
    shadow.line.fill.background()
    disable_shape_shadow(shadow)
    try:
        shadow.adjustments[0] = 0.06
    except Exception:  # noqa: BLE001
        pass
    card = ctx.card(y, h=h, x=x, w=w, fill=WHITE, border=None, radius=0.06)
    disable_shape_shadow(card)


def content_stage(
    ctx: _Ctx,
    y: int,
    *,
    bottom: int | None = None,
    top: int | None = None,
    corner_marks: bool = False,
) -> int:
    """Сцена вокруг обычной текстовой колонки; возвращает её низ.

    Колонка (`MARGIN_X`..`+CONTENT_W`) не сдвигается и не сужается: сцена
    выступает наружу на `STAGE_BLEED`, и этот выступ и есть её внутреннее поле.
    Так страница получает белую плоскость, не теряя ни строки содержимого.

    `top` задаётся, когда над сценой уже что-то нарисовано: сцена кладётся
    поверх более ранних фигур, и выступ вверх срезал бы нижний край плиток
    метрик вместе с подписями.
    """
    limit = CONTENT_BOTTOM if bottom is None else min(bottom, CONTENT_BOTTOM)
    top = max(0, y - STAGE_BLEED) if top is None else max(0, top)
    # Тень обязана остаться выше границы контентной области: растровая проверка
    # считает дефектом любые чернила ниже неё, а тень — это чернила.
    height = limit - top - STAGE_SHADOW_DY - 60_000
    if height < 400_000:
        return y
    left = MARGIN_X - STAGE_BLEED
    width = CONTENT_W + 2 * STAGE_BLEED
    draw_stage(ctx, left, top, width, height)
    if corner_marks:
        draw_corner_marks(ctx, left, top, width, height)
    return top + height


def draw_corner_marks(ctx: _Ctx, x: int, y: int, w: int, h: int) -> None:
    """Зелёные уголки по краям сцены выводов."""
    for cx, cy in (
        (x + 40_000, y + 40_000),
        (x + w - 120_000, y + 40_000),
        (x + 40_000, y + h - 120_000),
        (x + w - 120_000, y + h - 120_000),
    ):
        mark = ctx.slide.shapes.add_shape(1, Emu(cx), Emu(cy), Emu(80_000), Emu(14_000))
        mark.fill.solid()
        mark.fill.fore_color.rgb = ACCENT
        mark.line.fill.background()


def draw_level_bars(
    ctx: _Ctx,
    x: int,
    y: int,
    *,
    filled: int,
    total: int = 5,
    seg_w: int = 120_000,
    seg_h: int = 90_000,
    gap: int = 30_000,
    hot: RGBColor = TONE_RISK,
    cold: RGBColor = CARD_BORDER,
) -> None:
    """Шкала степени: пять делений, закрашено столько, какова степень."""
    for i in range(total):
        sx = x + i * (seg_w + gap)
        shape = ctx.slide.shapes.add_shape(5, Emu(sx), Emu(y), Emu(seg_w), Emu(seg_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hot if i < filled else cold
        shape.line.fill.background()
        try:
            shape.adjustments[0] = 0.35
        except Exception:  # noqa: BLE001
            pass


def tone_to_bars(tone: str, pill: str) -> int:
    """Сколько делений шкалы закрашено. Слово степени важнее тона карточки."""
    blob = f"{tone} {pill}".lower()
    if any(k in blob for k in ("крит", "danger", "крайне")):
        return 5
    if any(k in blob for k in ("высок", "high", "risk")):
        return 4
    if any(k in blob for k in ("сред", "medium", "warn")):
        return 3
    if any(k in blob for k in ("низк", "low")):
        return 2
    return 3


def bars_color(filled: int) -> RGBColor:
    return TONE_RISK if filled >= 4 else TONE_WARN


def render_hero_metrics_row(
    ctx: _Ctx,
    metrics: list[dict[str, Any]],
    x: int,
    y: int,
    width: int,
    *,
    tone_value_color,
) -> int:
    """Ряд метрик, где первая — герой страницы: крупная цифра во всю плитку.

    Высота ряда — та же, что у обычного ряда плиток (`METRIC_ROW_H`). В
    исходной ветке hero-полоса была высотой 1 350 000 плюс отбивка, то есть
    забирала у страницы 570 000 EMU — примерно три с половиной строки текста,
    которые дальше нечем было вернуть, кроме обрезки.

    Кегли: герой — `FS_TITLE`, остальные — `FS_SECTION`, длинные значения
    опускаются до `FS_BODY`. Ступень `FS_LEAD` намеренно не используется: на
    странице и так есть заголовок (`FS_SECTION`), подписи (`FS_BODY`) и
    колонтитул (`FS_CAPTION`), а шкала разрешает четыре ступени на страницу.
    """
    items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:4]
    if not items:
        return y
    hero, rest = items[0], items[1:]
    gap = 80_000
    hero_w = int(width * 0.34) if rest else width
    _metric_tile(
        ctx,
        hero,
        x,
        y,
        hero_w,
        METRIC_ROW_H,
        value_size=FS_TITLE,
        tone_value_color=tone_value_color,
    )
    if rest:
        rest_x = x + hero_w + gap
        rest_w = width - hero_w - gap
        tile_w = (rest_w - gap * (len(rest) - 1)) // len(rest)
        for i, m in enumerate(rest):
            _metric_tile(
                ctx,
                m,
                rest_x + i * (tile_w + gap),
                y,
                tile_w,
                METRIC_ROW_H,
                value_size=FS_SECTION,
                tone_value_color=tone_value_color,
            )
    return y + METRIC_ROW_H


def render_metric_tiles(
    ctx: _Ctx,
    metrics: list[dict[str, Any]],
    x: int,
    y: int,
    width: int,
    *,
    tone_value_color,
    cols: int = 3,
) -> int:
    """Ряд равных плиток — продолжение ряда метрик под hero-строкой.

    Кегль тот же, что у плиток hero-строки: страница уже держит четыре ступени
    шкалы, и пятая (`FS_LEAD`) на ней запрещена.
    """
    items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:cols]
    if not items:
        return y
    gap = 80_000
    tile_w = (width - gap * (len(items) - 1)) // len(items)
    for i, m in enumerate(items):
        _metric_tile(
            ctx,
            m,
            x + i * (tile_w + gap),
            y,
            tile_w,
            METRIC_ROW_H,
            value_size=FS_SECTION,
            tone_value_color=tone_value_color,
        )
    return y + METRIC_ROW_H


def _metric_tile(
    ctx: _Ctx,
    metric: dict[str, Any],
    x: int,
    y: int,
    w: int,
    h: int,
    *,
    value_size: float,
    tone_value_color,
) -> None:
    tone = str(metric.get("tone") or "neutral")
    value = _clip_words(_safe(metric.get("value")), 36)
    label = _clip_words(_safe(metric.get("label")), 40)
    ctx.card(y, h=h, x=x, w=w, fill=WHITE, border=None, radius=0.1)
    # Длинное значение — это не цифра, а фраза вроде «Данные не собраны»:
    # крупным кеглем она не помещается и распадается на три строки.
    size = value_size if len(value) <= 10 else FS_BODY
    inner_w = w - 280_000
    # Плотная интерлиньяж и малая отбивка — не украшение: подпись метрики вроде
    # «Тем повышенного внимания» встаёт в две строки, и при обычном интерлиньяже
    # вторая строка выходит за нижний край плитки. Текстовая рамка в PPTX не
    # обрезает, поэтому подпись просто ложилась поверх того, что ниже.
    box = ctx.slide.shapes.add_textbox(
        Emu(x + 140_000), Emu(y + 90_000), Emu(inner_w), Emu(h - 150_000)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.line_spacing = 1.0
    r0 = p0.add_run()
    r0.text = value
    r0.font.name = FONT
    r0.font.bold = True
    r0.font.size = Pt(size)
    r0.font.color.rgb = tone_value_color(tone)
    # Подписи длиннее двух строк в плитку не помещаются ни при какой отбивке:
    # ступень ниже (`FS_CAPTION`) на странице уже есть — это колонтитул, так что
    # новой ступени шкалы здесь не заводится.
    label_size = FS_BODY
    two_lines = 2 * measure_text_height("x", inner_w, FS_BODY, line_spacing=1.0, paragraph_spacing_pt=0)
    if measure_text_height(label, inner_w, FS_BODY, line_spacing=1.0, paragraph_spacing_pt=0) > two_lines:
        label_size = FS_CAPTION
    p1 = tf.add_paragraph()
    p1.space_before = Pt(2)
    p1.line_spacing = 1.0
    r1 = p1.add_run()
    r1.text = label
    r1.font.name = FONT
    r1.font.size = Pt(label_size)
    r1.font.color.rgb = MUTED_COLOR


def stage_heading(ctx: _Ctx, text: str, y: int, *, x: int | None = None, w: int | None = None) -> int:
    """Зелёный подзаголовок внутри сцены («Действие», «Выводы»).

    Кегль основного текста: на странице уже заняты четыре ступени шкалы, и
    иерархию здесь несут цвет и начертание, а не пятый размер.
    """
    left = MARGIN_X if x is None else x
    width = CONTENT_W if w is None else w
    box = ctx.slide.shapes.add_textbox(Emu(left), Emu(y), Emu(width), Emu(200_000))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = _safe(text)
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(FS_BODY)
    r.font.color.rgb = ACCENT
    return y + 200_000


def render_action_block(
    ctx: _Ctx, label: str, y: int, *, max_h: int, heading: str = "Действие"
) -> int:
    """«Действие» на сцене: зелёный подзаголовок и текст рекомендации.

    Обрубок не печатается. Карточка «Действие» со словом «Проверить» вместо
    рекомендации занимает полосу во всю ширину и не сообщает ничего — правило
    из `content_card(skip_if_stub=True)`, которое здесь пришлось бы потерять
    вместе с карточкой.
    """
    text = _safe(label)
    if not text:
        return y
    avail = min(max_h, CONTENT_BOTTOM - y)
    if avail < 420_000:
        return y
    body_avail = avail - 260_000
    fitted = _fit_text_to_height(text, CONTENT_W, FS_BODY, body_avail)
    if len(fitted) < min(60, int(len(text) * 0.5)):
        return y
    y = stage_heading(ctx, heading, y)
    return ctx.body(text, y, max_h=body_avail) + 40_000


def alternating_color(index: int) -> RGBColor:
    """Чередование зелёного и фиолетового у маркеров — ритм страницы cleeq."""
    return ACCENT if index % 2 == 0 else VIOLET
