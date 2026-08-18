"""Per-template slide dispatch (_render_slide)."""

from __future__ import annotations

import io
import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

from .common import (
    ACCENT,
    BULLET_GLYPH,
    FS_CARD_TITLE,
    FS_COVER,
    FS_LEAD,
    FS_SUBTITLE,
    BODY_COLOR,
    CARD_BG,
    CARD_BORDER,
    CONTENT_BOTTOM,
    CONTENT_W,
    COVER_BG,
    COVER_SUBTITLE,
    CYAN,
    FONT,
    FS_BODY,
    FS_CAPTION,
    FS_SECTION,
    FS_TITLE,
    MARGIN_X,
    MUTED_COLOR,
    NAVY,
    SLIDE_H,
    TONE_RISK,
    VIOLET,
    WHITE,
    _Ctx,
    _clip_words,
    _embed_cover_portrait,
    _embed_image,
    _first_visual_asset,
    _resolve_image_bytes,
    _safe,
    _safe_preserve_breaks,
    _trim_dangling_tail,
    record_text_layout,
)
from .executive import (
    _render_executive_dashboard,
    _render_profile_overview,
    _render_risk_matrix_grid,
)
from .layout_cleeq import (
    alternating_color,
    content_stage,
    render_action_block,
    render_hero_metrics_row,
    render_metric_tiles,
)
from .visual import (
    _add_search_table,
    _render_kpi_cards,
    _render_status_badge,
    _render_visual_with_sidebar,
    _tone_value_color,
)


def _draw_cleeq_cover_art(ctx: _Ctx) -> None:
    """Абстрактные полосы бренда справа — обложка без портрета субъекта.

    Полосы держатся в границах листа и внутри боковых полей. В исходной ветке
    они уходили за правый край (навылет), и на эталонной деке это давало пять
    `out-of-bounds` от инспектора геометрии, две «пустые панели» от ручной
    визуальной проверки и два дефекта растровой — то есть обложка одна валила
    три приёмочных ворот из девятнадцати. Обрез — приём хороший, но он должен
    быть решением, объявленным проверке, а не её обходом.
    """
    right = 11_350_000
    bands = [
        (7_100_000, 520_000, right - 7_100_000, 2_300_000, ACCENT),
        (7_900_000, 1_450_000, right - 7_900_000, 2_400_000, VIOLET),
        (8_700_000, 2_500_000, right - 8_700_000, 1_700_000, CYAN),
        (7_450_000, 3_800_000, right - 7_450_000, 1_150_000, ACCENT),
        (8_900_000, 4_600_000, right - 8_900_000, 1_450_000, VIOLET),
    ]
    for x, y, w, h, color in bands:
        shape = ctx.slide.shapes.add_shape(5, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        try:
            shape.adjustments[0] = 0.5
        except Exception:  # noqa: BLE001
            pass


#: Сколько места держится под рекомендацией и футнотом, когда над ними стоит
#: список строк. Числа те же, которыми ограничена карточка «Что проверить» и
#: подпись под ней: список не вправе занять их место — иначе рекомендация
#: молча ужимается до обрубка, а методология не печатается вовсе.
#:
#: `FOOTNOTE_RESERVE` — замер: подпись 9 pt в 300 знаков (потолок клипа) занимает
#: 363 697 EMU. Этой же величиной проверяется, есть ли место под футнот: пока
#: резерв и порог были разными числами (360 000 и 400 000), футнот мог не
#: напечататься на полном листе — и молча.
ACTION_CARD_MAX_H = 1_100_000
CARD_GAP = 110_000
FOOTNOTE_RESERVE = 380_000


def _render_status_cards(
    ctx: _Ctx,
    slide: dict[str, Any],
    title: str,
    narrative: str,
    bullets: list[str],
    *,
    status_title: str,
    bullets_as_card: bool,
) -> None:
    """Карточная страница: статус, содержимое, рекомендация, футнот методологии.

    Один макет на два шаблона — пустое состояние поверхности и страницу
    фактической проверки Википедии. Разница ровно одна: буллеты печатаются
    карточкой «Что это означает» (пустое состояние объясняет, чем плохо
    отсутствие материалов) или списком строк выдачи. Копии геометрии не
    заводится намеренно: методология и рекомендация обязаны печататься на обеих
    страницах, а два одинаковых макета расходятся с первой же правкой.
    """
    ctx.light_bg()
    y = ctx.title(title, 320000, NAVY)
    # Карточка с одним заголовком — это пустой озаглавленный блок, и приёмка
    # считает его дефектом. На странице-продолжении нарратива нет по построению
    # (он принадлежит первой странице блока), поэтому карточка там не рисуется.
    if narrative:
        # PDF-36 D.3 — cards height-fit with font step-down; no char starvation.
        y = ctx.content_card(
            title=status_title,
            text=narrative,
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=380_000,
            max_h=1_700_000,
            tone="accent",
            title_size=11,
            body_size=FS_BODY,
        )
        y += CARD_GAP
    actions = [a for a in (slide.get("actions") or []) if isinstance(a, dict)]
    methodology = _safe(slide.get("methodologyNote") or "")
    source_note = _safe(slide.get("sourceNote") or "")
    footnote = " ".join(x for x in (methodology, source_note) if x)
    if bullets and bullets_as_card:
        y = ctx.content_card(
            title="Что это означает",
            text="\n".join(bullets[:4]),
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=360_000,
            max_h=2_200_000,
            tone="neutral",
            title_size=11,
            body_size=11,
        )
        y += CARD_GAP
    elif bullets:
        reserved = (ACTION_CARD_MAX_H + CARD_GAP if actions else 0) + (
            FOOTNOTE_RESERVE if footnote else 0
        )
        y = ctx.bullets(bullets, y, max_items=8, bottom=CONTENT_BOTTOM - reserved)
        y += CARD_GAP
    if actions:
        y = ctx.content_card(
            title="Что проверить",
            text=_safe(actions[0].get("label")),
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=320_000,
            max_h=ACTION_CARD_MAX_H,
            tone="warn",
            title_size=11,
            body_size=11,
        )
        y += CARD_GAP
    if not footnote:
        return
    if y <= CONTENT_BOTTOM - FOOTNOTE_RESERVE:
        ctx.body(
            _clip_words(footnote, 300),
            y,
            max_h=CONTENT_BOTTOM - y - 60_000,
            color=MUTED_COLOR,
            font_size=9,
        )
        return
    # Ненапечатанная методология — потеря содержимого, а не мелочь вёрстки: на
    # странице проверки она и есть предмет страницы. На ветке со списком место
    # под неё держит резерв, и сюда попасть нельзя; на карточной ветке карточки
    # могут съесть лист — тогда об этом обязано быть слышно.
    record_text_layout(
        page=ctx.page,
        name=f"orion_footnote_dropped_p{ctx.page}",
        role="footnote",
        font_family=FONT,
        font_size_pt=9,
        box_width=CONTENT_W,
        box_height=0,
        available_height=max(0, CONTENT_BOTTOM - y),
        required_height=FOOTNOTE_RESERVE,
        measured_lines=0,
        text_length=len(footnote),
        clipped=True,
        measurement_uncertain=False,
        dropped_lines=1,
    )


def _render_slide(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    template = str(slide.get("template") or "")
    # Level 2.5 — named pre-built layout variant picked by the GPT composer.
    # Validated upstream (TS registry); unknown values fall back to default.
    variant = str(slide.get("layoutVariant") or "")
    title = _safe(slide.get("title") or "ORION")
    narrative = _safe(slide.get("narrative") or "")
    # PDF-47 — keep structured theme newlines; _safe() collapses them and then
    # nested «…«…»» quotes cannot reflow → theme-only stubs («Офшоры»).
    bullets = [
        _safe_preserve_breaks(b)
        for b in slide.get("bullets") or []
        if _safe_preserve_breaks(b)
    ]
    refs = slide.get("assetRefs") or []
    primary = assets.get(str(refs[0])) if refs else None

    if template == "orion_golden_cover":
        # Обложка cleeq: чернильный лист, словесный знак, зелёный надзаголовок,
        # имя субъекта крупно, метаданные и чип. Справа — портрет субъекта из
        # выдачи, а если превью нет, абстрактные полосы бренда.
        ctx.dark_bg()
        portrait = _first_visual_asset(list(refs), assets) or primary
        if _embed_cover_portrait(ctx, portrait):
            # Полотно слева, чтобы подпись читалась поверх плитки портрета.
            veil = ctx.slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(6_200_000), Emu(SLIDE_H))
            veil.fill.solid()
            veil.fill.fore_color.rgb = COVER_BG
            veil.line.fill.background()
        else:
            _draw_cleeq_cover_art(ctx)
        brand = ctx.slide.shapes.add_textbox(
            Emu(MARGIN_X), Emu(420_000), Emu(4_800_000), Emu(320_000)
        )
        br = brand.text_frame.paragraphs[0].add_run()
        br.text = "cleeq"
        br.font.name = FONT
        br.font.bold = True
        br.font.size = Pt(FS_SUBTITLE)
        br.font.color.rgb = WHITE
        hero = (title or "Цифровой профиль").strip()
        for sep in (" — ", " – ", " - "):
            if sep in hero:
                hero = hero.split(sep, 1)[-1].strip() or hero
                break
        kicker = ctx.slide.shapes.add_textbox(
            Emu(MARGIN_X), Emu(3_250_000), Emu(5_800_000), Emu(280_000)
        )
        kr = kicker.text_frame.paragraphs[0].add_run()
        kr.text = "Цифровой профиль"
        kr.font.name = FONT
        kr.font.bold = True
        kr.font.size = Pt(FS_BODY)
        kr.font.color.rgb = ACCENT
        name_box = ctx.slide.shapes.add_textbox(
            Emu(MARGIN_X), Emu(3_600_000), Emu(5_900_000), Emu(1_400_000)
        )
        ntf = name_box.text_frame
        ntf.word_wrap = True
        nr = ntf.paragraphs[0].add_run()
        nr.text = _safe(hero)
        nr.font.name = FONT
        nr.font.bold = True
        nr.font.size = Pt(FS_COVER)
        nr.font.color.rgb = WHITE
        ctx.body(
            narrative or "Конфиденциально. Подготовлено для внутреннего использования клиента.",
            5_050_000,
            max_h=600_000,
            color=COVER_SUBTITLE,
            w=5_800_000,
            font_size=FS_SUBTITLE,
        )
        # Чип держится над границей контентной области: ниже неё чернил быть не
        # должно, и растровая проверка ловит это по отрисованной странице.
        chip = ctx.slide.shapes.add_shape(
            5, Emu(MARGIN_X), Emu(5_760_000), Emu(3_400_000), Emu(320_000)
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = RGBColor(0x1A, 0x24, 0x1A)
        chip.line.fill.background()
        try:
            chip.adjustments[0] = 0.5
        except Exception:  # noqa: BLE001
            pass
        chip_t = ctx.slide.shapes.add_textbox(
            Emu(MARGIN_X + 120_000), Emu(5_810_000), Emu(3_100_000), Emu(240_000)
        )
        ctr = chip_t.text_frame.paragraphs[0].add_run()
        ctr.text = "Аудит · стратегия · конфиденциально"
        ctr.font.name = FONT
        ctr.font.size = Pt(FS_BODY)
        ctr.font.color.rgb = ACCENT
        return

    if template == "orion_golden_toc":
        # Содержание cleeq: мятный лист, белая сцена, зелёные и фиолетовые
        # номера разделов, названия чернилами по волосяным линейкам.
        ctx.light_bg()
        y = ctx.title("Содержание отчёта", 320_000, NAVY, FS_TITLE)
        entries = [
            _clip_words(b, 110)
            for b in (bullets or ["Резюме", "Россия", "ОАЭ", "Compliance", "LexisNexis", "Рекомендации"])
        ][:10]
        stage_top = y + 80_000
        stage_bottom = content_stage(ctx, stage_top)
        row_h = min(560_000, max(380_000, (stage_bottom - stage_top - 280_000) // max(1, len(entries))))
        ry = stage_top + 120_000
        for i, entry in enumerate(entries, start=1):
            num = ctx.slide.shapes.add_textbox(
                Emu(MARGIN_X + 60_000), Emu(ry), Emu(700_000), Emu(row_h)
            )
            nr = num.text_frame.paragraphs[0].add_run()
            nr.text = f"{i:02d}"
            nr.font.name = FONT
            nr.font.bold = True
            nr.font.size = Pt(FS_LEAD)
            nr.font.color.rgb = alternating_color(i - 1)
            box = ctx.slide.shapes.add_textbox(
                Emu(MARGIN_X + 820_000),
                Emu(ry + 40_000),
                Emu(CONTENT_W - 900_000),
                Emu(row_h - 60_000),
            )
            tf = box.text_frame
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = entry
            r.font.name = FONT
            r.font.bold = True
            r.font.size = Pt(FS_SUBTITLE)
            r.font.color.rgb = NAVY
            rule = ctx.slide.shapes.add_shape(
                1,
                Emu(MARGIN_X + 60_000),
                Emu(ry + row_h - 40_000),
                Emu(CONTENT_W - 120_000),
                Emu(9_000),
            )
            rule.fill.solid()
            rule.fill.fore_color.rgb = CARD_BORDER
            rule.line.fill.background()
            ry += row_h
        return

    if template == "orion_golden_executive_dashboard":
        _render_executive_dashboard(ctx, slide, title)
        return

    if template == "orion_golden_risk_matrix_grid":
        _render_risk_matrix_grid(ctx, slide, title)
        return

    if template == "orion_golden_profile_overview":
        _render_profile_overview(ctx, slide, title)
        return

    if template == "orion_golden_executive_card":
        # Страница-документ cleeq: заголовок снаружи, всё содержимое — на одной
        # белой сцене. Карточек внутри сцены нет: белое на белом не читается как
        # отдельный блок, и вместо иерархии выходит рябь из рамок.
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
        if metrics:
            metrics_bottom = render_hero_metrics_row(
                ctx, metrics[:4], MARGIN_X, y, CONTENT_W, tone_value_color=_tone_value_color
            )
            y = metrics_bottom + 140_000
        content_stage(ctx, y, top=metrics_bottom + 40_000 if metrics else None)
        narr = narrative.strip()
        if variant == "accent-headline" and narr:
            # Accent-headline: the lead sentence leads the page, the rest of the
            # narrative follows, then the detail bullets.
            sentences = re.split(r"(?<=[.!?…])\s+", narr)
            lead = sentences[0] if sentences else narr
            rest = " ".join(sentences[1:]).strip()
            y = ctx.body(lead, y, max_h=1_100_000, bold=True) + 60_000
            if rest:
                y = ctx.body(
                    rest,
                    y,
                    max_h=min(2_400_000, CONTENT_BOTTOM - y - (1_200_000 if bullets else 100_000)),
                ) + 60_000
            if bullets:
                ctx.bullets(bullets, y, max_items=6, max_chars=900)
            return
        if narr and not bullets:
            ctx.body(narr, y, max_h=min(3_200_000, CONTENT_BOTTOM - y - 100_000), bold=True)
            return
        if narr:
            y = ctx.body(
                narr,
                y,
                max_h=min(2_000_000, CONTENT_BOTTOM - y - (1_200_000 if bullets else 100_000)),
                bold=True,
            ) + 60_000
        if bullets:
            ctx.bullets(bullets, y, max_items=6, max_chars=900)
        return

    if template == "orion_golden_risk_matrix":
        _render_risk_matrix_grid(ctx, slide, title or "Матрица рисков")
        return

    if template == "orion_golden_region_divider":
        ctx.dark_bg()
        if variant == "hero":
            # Разделитель cleeq: зелёный столб, фиолетовая засечка, крупный титул.
            bar = ctx.slide.shapes.add_shape(
                5, Emu(MARGIN_X), Emu(2_200_000), Emu(140_000), Emu(2_200_000)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            try:
                bar.adjustments[0] = 0.5
            except Exception:  # noqa: BLE001
                pass
            accent = ctx.slide.shapes.add_shape(
                5, Emu(MARGIN_X + 220_000), Emu(2_200_000), Emu(90_000), Emu(700_000)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = VIOLET
            accent.line.fill.background()
            text_x = MARGIN_X + 420_000
            text_w = CONTENT_W - 420_000
            box = ctx.slide.shapes.add_textbox(Emu(text_x), Emu(2_250_000), Emu(text_w), Emu(1_100_000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = _safe(title)
            r.font.name = FONT
            r.font.bold = True
            r.font.size = Pt(FS_COVER)
            r.font.color.rgb = WHITE
            if narrative:
                ctx.body(
                    narrative,
                    3_500_000,
                    max_h=1_500_000,
                    color=COVER_SUBTITLE,
                    font_size=FS_SUBTITLE,
                    x=text_x,
                    w=text_w,
                )
            return
        ctx.title(title, 2800000, WHITE, FS_COVER)
        return

    if template == "orion_golden_metrics_dashboard":
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        badge = slide.get("statusBadge") if isinstance(slide.get("statusBadge"), dict) else None
        if badge:
            y = _render_status_badge(ctx, badge, MARGIN_X, y, CONTENT_W) + 80_000
        # Страница метрик cleeq: ключевая цифра ведёт страницу, остальное лежит
        # на одной сцене. Вариант `kpi-first` отдельной веткой больше не нужен —
        # цифры ведут страницу всегда, и это ровно то, ради чего он заводился.
        #
        # Порядок содержимого прежний: метрики → нарратив → действие → темы.
        # Ёмкость тоже прежняя: ряд метрик той же высоты, что раньше, а сцена
        # рисуется под текстом и его не двигает.
        metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
        if metrics:
            y = render_hero_metrics_row(
                ctx, metrics[:4], MARGIN_X, y, CONTENT_W, tone_value_color=_tone_value_color
            )
            if len(metrics) > 4:
                y = render_metric_tiles(
                    ctx,
                    metrics[4:7],
                    MARGIN_X,
                    y + 80_000,
                    CONTENT_W,
                    tone_value_color=_tone_value_color,
                )
            metrics_bottom = y
            y += 140_000
        # Уголки — метка «сцены выводов» в языке cleeq.
        content_stage(
            ctx,
            y,
            top=metrics_bottom + 40_000 if metrics else None,
            corner_marks=True,
        )
        if narrative:
            y = ctx.body(
                narrative,
                y,
                max_h=900_000 if metrics else 1_100_000,
                bold=True,
            ) + 70_000
        actions = [a for a in (slide.get("actions") or []) if isinstance(a, dict)]
        if actions and (not bullets or (CONTENT_BOTTOM - y) > 1_800_000):
            y = render_action_block(
                ctx, _safe(actions[0].get("label")), y, max_h=1_000_000
            )
        if bullets:
            # Потолок читаемости, а не ёмкости: сколько блоков влезает, решает
            # мерка высоты, приведённая к тому, что рисуется (шаг 16, 07.6).
            ctx.bullets(bullets, y, max_items=6, max_chars=900)
        return

    if template == "orion_golden_serp_screenshot":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_image_grid":
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        if len(refs) == 1:
            primary_grid = assets.get(str(refs[0])) if refs else None
            if primary_grid and _resolve_image_bytes(primary_grid):
                _embed_image(ctx, primary_grid, y + 60000, h=5_200_000)
                cap = _safe(primary_grid.get("caption") or "")
                if cap:
                    ctx.body(cap, CONTENT_BOTTOM - 380000, max_h=320000, color=MUTED_COLOR)
                return
        cols = 3
        cell_w = CONTENT_W // 3 - 80_000
        cell_h = 1_600_000
        gap = 120000
        max_rows = max(1, int((CONTENT_BOTTOM - y + gap) // (cell_h + gap)))
        max_cells = max_rows * cols
        for idx, ref in enumerate(refs):
            if idx >= max_cells:
                break
            row = idx // cols
            col = idx % cols
            cx = MARGIN_X + col * (cell_w + gap)
            cy = y + row * (cell_h + gap)
            asset = assets.get(str(ref))
            raw = _resolve_image_bytes(asset) if asset else None
            if raw:
                stream = io.BytesIO(raw)
                iw, ih = cell_w, cell_h
                if Image is not None:
                    try:
                        with Image.open(io.BytesIO(raw)) as im:
                            iw, ih = im.size
                    except Exception:
                        pass
                scale = min(cell_w / max(iw, 1), cell_h / max(ih, 1))
                dw, dh = int(iw * scale), int(ih * scale)
                left = cx + (cell_w - dw) // 2
                top = cy + (cell_h - dh) // 2
                ctx.slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(dw), height=Emu(dh))
            else:
                shape = ctx.slide.shapes.add_shape(5, Emu(cx), Emu(cy), Emu(cell_w), Emu(cell_h))
                try:
                    shape.adjustments[0] = 0.05
                except Exception:  # noqa: BLE001
                    pass
                shape.fill.solid()
                shape.fill.fore_color.rgb = CARD_BG
                shape.line.color.rgb = CARD_BORDER
                shape.line.width = Pt(0.75)
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = _safe((asset or {}).get("title") or "Недоступно")
                r.font.size = Pt(FS_CAPTION)
        return

    if template == "orion_golden_video_cards":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_knowledge_panel":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_surface_panel":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_lexis_visual_page":
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        _embed_image(ctx, primary, y + 60000, h=5_200_000)
        return

    if template == "orion_golden_compliance_visual_page":
        # Dow Jones / World-Check approved screenshots — same layout as Lexis visual.
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        _embed_image(ctx, primary, y + 60000, h=5_200_000)
        return

    if template == "orion_golden_search_table":
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        # Сцена под таблицей: строки таблицы лежат на белой плоскости, а не
        # висят на мятном фоне.
        content_stage(ctx, y)
        if narrative:
            # Keep 1–2 complete sentences above the table; never end on «как/и/с».
            intro = _safe(narrative)
            sentences = [s.strip() for s in re.split(r"(?<=[.!?…])\s+", intro) if s.strip()]
            complete = [
                s
                for s in sentences
                if s.endswith((".", "!", "?", "…"))
                and not re.search(r"(?:\bкак|\bи|\bс|\bв|\bпо|,|;|—)\s*$", s, re.I)
            ]
            if complete:
                intro = " ".join(complete[:2])
            else:
                intro = "Таблица фиксирует сохранённые позиции поисковой выдачи."
            intro = _trim_dangling_tail(intro)
            y = ctx.body(intro, y, max_h=900000, color=MUTED_COLOR)
            y = y + 40000
        table = slide.get("table") if isinstance(slide.get("table"), dict) else None
        headers = list((table or {}).get("headers") or [])
        rows = list((table or {}).get("rows") or [])
        groups = list((table or {}).get("groups") or [])
        if not rows and bullets:
            # Fallback: parse bullet lines into a compact table
            headers = ["Поз.", "Домен", "Заголовок", "Риск"]
            parsed: list[list[str]] = []
            for bullet in bullets[:10]:
                raw = _safe(bullet)
                m = re.match(
                    r"^(?:\[([Н·N.])\]\s*)?#?\s*(\d+)\s+([^\s—\-]+)\s*[—\-–]\s*(.+)$",
                    raw,
                )
                if m:
                    mark = "Н" if m.group(1) in ("Н", "N") else "·"
                    parsed.append([m.group(2), m.group(3), _clip_words(m.group(4), 70), mark])
                else:
                    parsed.append(["—", "—", _clip_words(raw, 80), "·"])
            rows = parsed
        if headers and rows:
            # Render every row the (paginated) slide carries — no hidden cap.
            # Keep up to 5 headers so Запрос can be stripped inside the helper
            # without also dropping Статус.
            _add_search_table(ctx, y, headers[:5], rows, groups)
        elif bullets:
            avail = max(400000, CONTENT_BOTTOM - y)
            box = ctx.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(avail))
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for bullet in bullets[:18]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(2)
                p.space_after = Pt(5)
                r = p.add_run()
                clipped = _clip_words(bullet, 240)
                r.text = f"{BULLET_GLYPH} {clipped}"
                r.font.name = FONT
                r.font.size = Pt(FS_BODY)
                r.font.color.rgb = TONE_RISK if clipped.startswith("[Н]") else BODY_COLOR
        return

    if template == "orion_golden_no_data_compact":
        # C.2 — honest empty state as a structured page: status card,
        # "what it means" card, recommendation card, methodology footnote.
        _render_status_cards(
            ctx,
            slide,
            title,
            narrative or "Для данного раздела недостаточно подтверждённых данных.",
            bullets,
            status_title="Статус сбора",
            bullets_as_card=True,
        )
        return

    if template == "orion_golden_wikipedia_check":
        # Страница фактической проверки: тот же карточный макет, что и у
        # пустого состояния, но буллеты — строки поисковой выдачи, и печатаются
        # они списком между результатом проверки и рекомендацией.
        _render_status_cards(
            ctx,
            slide,
            title,
            narrative,
            bullets,
            status_title="Результат проверки",
            bullets_as_card=False,
        )
        return

    if template == "orion_golden_audit_dashboard":
        # ORION regional résumé: themes left-ish via bullets top, KPI counters below.
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        if narrative:
            # PDF-36 D.3 — ctx.body height-fits with font step-down; feed it
            # the full text instead of a pre-starved 520-char slice.
            y = ctx.body(_clip_words(narrative, 760), y, max_h=1200000)
            y = y + 80000
        if bullets:
            ctx.bullets(bullets, y, max_items=14, max_chars=340)
        return

    # orion_golden_prose (continuation themes) + default section / appendix
    ctx.light_bg()
    y = ctx.title(title, 280000, NAVY, FS_SECTION)
    content_stage(ctx, y)
    # PDF-36 D.3 / PDF-47 — page bottom is the real budget; theme cards need
    # the full 900-char structured budget (520 starved nested-quote bullets).
    short_narrative = _clip_words(narrative, 900) if narrative else ""
    if short_narrative and not bullets:
        ctx.body(short_narrative, y, max_h=CONTENT_BOTTOM - y - 100000, bold=True)
        return
    if short_narrative:
        y = ctx.body(short_narrative, y, max_h=1100000, bold=True)
        y = y + 80000
    if bullets:
        ctx.bullets(bullets, y, max_items=9, max_chars=900)


