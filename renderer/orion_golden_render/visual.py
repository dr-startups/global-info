"""Sidebar, KPI cards, visual+sidebar layout, and search tables."""

from __future__ import annotations

import io
import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Emu, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

from .common import (
    TYPE_SCALE_PT,
    FS_LEAD,
    FS_SUBTITLE,
    ACCENT,
    ACCENT_SOFT,
    BODY_COLOR,
    CARD_BG,
    CARD_BORDER,
    CONTENT_BOTTOM,
    CONTENT_W,
    EMU_PER_PT,
    FONT,
    FORBIDDEN,
    FS_BODY,
    FS_CAPTION,
    GOOD_BG,
    MARGIN_X,
    METRIC_ACCENT,
    MUTED_COLOR,
    NAVY,
    RISK_BG,
    SIDEBAR_SAFE_FALLBACK,
    TONE_GOOD,
    TONE_RISK,
    TONE_WARN,
    WARN_BG,
    WHITE,
    _Ctx,
    _clip_words,
    _embed_image,
    _embed_image_contain,
    _first_visual_asset,
    _fit_text_to_height,
    _resolve_image_bytes,
    _safe,
    _wrapped_line_count,
    measure_text_height,
    plural_ru,
    record_bullet_measure,
    record_text_layout,
)
from .layout_cleeq import level_step

#: Поля ячейки python-pptx: по 0.1″ с каждой стороны. Перенос случается по
#: полезной ширине, а не по ширине колонки.
CELL_MARGINS_EMU = 2 * 91_440

#: Кегль бейджа статуса — на полступени крупнее подписи; тем же и меряется.
BADGE_PT = 9.5

#: Суффикс ключа страницы у мерной записи таблицы.
#:
#: Мера таблицы и мера пути буллетов лежат в одном списке вердикта, и ключом в
#: нём служит ключ страницы. Без суффикса перекладка буллетов нашла бы запись
#: таблицы по имени страницы и прочитала бы числа про строки таблицы как числа
#: про блоки списка. Суффикс разводит два вопроса, не трогая ни форму записи,
#: ни версию вердикта (`orion-bullet-measure-v1`): рендерер прошлой версии
#: просто не пишет этих строк, и приложение остаётся с раскладкой сида.
#:
#: Знака «#» в ключах страниц деки нет, поэтому с настоящим ключом суффикс
#: столкнуться не может.
TABLE_MEASURE_KEY_SUFFIX = "#table"

try:
    from client_text_contract import sidebar_check_failures
except ImportError:  # pragma: no cover
    from renderer.client_text_contract import sidebar_check_failures  # type: ignore

def _sidebar_word_budget(text: str, max_words: int = 70) -> str:
    """Keep complete sentences within a soft word budget; never emit one-word stubs."""
    raw = _safe(text)
    if not raw:
        return ""
    words = raw.split()
    if len(words) <= max_words:
        return raw
    sentences = re.split(r"(?<=[.!?…])\s+", raw)
    kept: list[str] = []
    count = 0
    for sent in sentences:
        w = len(sent.split())
        if kept and count + w > max_words:
            break
        if not kept and w > max_words:
            return _clip_words(sent, max(80, max_words * 6))
        kept.append(sent)
        count += w
        if count >= max_words:
            break
    return " ".join(kept).strip() or _clip_words(raw, max(80, max_words * 6))


def _qa_preview(text: str, match_index: int = 0) -> str:
    """Short, safe preview around a matched token for QA diagnostics."""
    raw = _safe(text)
    start = max(0, match_index - 20)
    snippet = raw[start:match_index + 40].strip()
    snippet = re.sub(r"\s+", " ", snippet)
    return (snippet[:60] + "…") if len(snippet) > 60 else snippet


def _sidebar_sanitize_field(ctx: _Ctx, field: str, text: str) -> str:
    """REMEDIATION §6.2 — bad field → neutral fallback + warning (no raise)."""
    value = _safe(text)
    if not value:
        return value
    failures = sidebar_check_failures(value, field, ctx.client_text_contract)
    if not failures:
        return value
    for msg in failures:
        preview = _qa_preview(value)
        ctx.warnings.append(f"sidebar-qa:p{ctx.page}:{field}:{msg} preview=\"{preview}\"")
    return SIDEBAR_SAFE_FALLBACK


def _sidebar_loss(
    ctx: _Ctx,
    field: str,
    kind: str,
    avail: int,
    needed: int,
    size: float,
    w: int,
    lost: str,
) -> None:
    """Потеря текста на панели — вслух: предупреждением и записью разметки.

    Обе ветки укладки панели молчали. Необязательный блок, у которого не влезло
    даже первое предложение, выбрасывался целиком, а блок, из которого осталась
    часть предложений, укорачивался — и ни то ни другое не оставляло следа: на
    прогоне 91 (стр. 46) страница подсказок не напечатала «Что сделать» и вывод
    о собранном наборе, а в `layout-telemetry.json` не было ни одной записи об
    этой странице при 62 записях по деке.

    `dropped_bullets`/`dropped_lines` намеренно не трогаются: они поднимают
    `CONTENT_DROPPED_BY_RENDERER`, а это блокер выдачи, и живой прогон вставал
    бы на законном выходе модели. Клип-код блокером не является и уезжает
    оператору строкой предупреждения.

    `text_length` меряет **потерянное**, а не весь блок: на обрезке часть
    предложений напечатана, и запись про «700 знаков» там, где до клиента не
    доехало 300, врёт оператору в ту же сторону, что и молчание.
    """
    ctx.warnings.append(f"sidebar-qa:p{ctx.page}:{field}:{kind}")
    record_text_layout(
        page=ctx.page,
        name=f"orion_sidebar_{field}_p{ctx.page}",
        role="sidebar",
        font_family=FONT,
        font_size_pt=size,
        box_width=w,
        box_height=max(0, avail),
        available_height=max(0, avail),
        required_height=needed,
        measured_lines=0,
        text_length=len(lost),
        clipped=True,
    )


def _sidebar_analysis(ctx: _Ctx, slide: dict[str, Any], x: int, y: int, w: int, h: int) -> None:
    """Unified client sidebar panel (v57): one column, no stacked framed cards."""
    analysis = slide.get("visualAnalysis") or {}
    if not isinstance(analysis, dict):
        analysis = {}

    mode = str(analysis.get("sidebarMode") or "")
    headline = _safe(analysis.get("headlineConclusion") or slide.get("clientTakeaway") or "Вывод")
    meaning = _safe(analysis.get("clientMeaning") or analysis.get("whyItMatters") or "")
    visible = _safe(analysis.get("whatIsVisible") or "")
    explanations = analysis.get("highlightExplanations") or []
    if not isinstance(explanations, list):
        explanations = []
    actions = analysis.get("recommendedActions") or []
    # Рисуются **все** рекомендации, а не первая: остальные молча пропадали, а
    # ворот приёмки проверяет каждую — на нагрузке с двумя рекомендациями он
    # краснел бы на ветке рендерера, а не на дефекте.
    action = (
        " ".join(x for x in (_safe(a) for a in actions) if x) if isinstance(actions, list) else ""
    )
    provenance = _safe(analysis.get("provenanceLabel") or "")
    more_n = int(analysis.get("moreSignalsCount") or 0)

    has_frames = any(
        isinstance(ex, dict) and str(ex.get("frameTone") or "") in {"red", "amber"} for ex in explanations
    )
    if mode == "adverse_explanation" or has_frames:
        mid_title = "Почему выделено"
        mid_bits = []
        for ex in explanations[:2]:
            if not isinstance(ex, dict):
                continue
            reason = _safe(ex.get("clientReason") or "")
            if reason:
                mid_bits.append(reason)
        if more_n > 0:
            noun = plural_ru(more_n, "похожий сигнал", "похожих сигнала", "похожих сигналов")
            mid_bits.append(f"Ещё {more_n} {noun}.")
        mid_body = " ".join(mid_bits) if mid_bits else visible
    else:
        mid_title = "Что показывает экран"
        mid_body = visible or meaning

    # Soft-degrade contract violations per field (§6.2); render continues.
    headline = _sidebar_sanitize_field(ctx, "headlineConclusion", headline)
    mid_body = _sidebar_sanitize_field(ctx, "whatIsVisible", mid_body)
    meaning = _sidebar_sanitize_field(ctx, "clientMeaning", meaning)
    action = _sidebar_sanitize_field(ctx, "recommendedActions", action)

    # Draw one outer panel
    pad = 70_000
    gap = 55_000
    cy = y + pad
    max_bottom = min(y + h, CONTENT_BOTTOM) - pad
    ctx.card(y, h=min(h, max_bottom - y + pad), x=x, w=w, fill=CARD_BG)

    def write_block(
        title: str | None,
        body: str,
        *,
        field: str,
        size: float = 11,
        bold_title: bool = True,
        required: bool = False,
    ) -> None:
        nonlocal cy
        if not body:
            return
        title_h = 200_000 if title else 0
        # Prefer complete text; do not ellipsis-clip sidebar
        fitted = body
        needed = measure_text_height(fitted, w - 2 * pad, size, line_spacing=1.2)
        avail = max_bottom - cy - 160_000 - title_h
        # PDF-36 D.3 — shrink the font 1–1.5 pt before dropping sentences.
        if needed > avail:
            # Только ступени шкалы: «минус полтора пункта» её нарушало.
            for candidate in [x for x in reversed(TYPE_SCALE_PT) if x < size]:
                if candidate < 9.5:
                    break
                cand_h = measure_text_height(fitted, w - 2 * pad, candidate, line_spacing=1.2)
                if cand_h <= avail:
                    size = candidate
                    needed = cand_h
                    break
        if needed > avail:
            # Keep complete sentences only (no ellipsis). When even the first
            # sentence cannot fit: required headline → safe fallback (§6.2);
            # optional blocks are dropped whole.
            sentences = re.split(r"(?<=[.!?…])\s+", fitted)
            kept: list[str] = []
            for sent in sentences:
                trial = " ".join(kept + [sent]).strip()
                if measure_text_height(trial, w - 2 * pad, size, line_spacing=1.2) <= avail:
                    kept.append(sent)
                else:
                    break
            if not kept:
                if required:
                    ctx.warnings.append(
                        f"sidebar-qa:p{ctx.page}:{field}:overflow without complete sentence"
                    )
                    fitted = SIDEBAR_SAFE_FALLBACK
                else:
                    _sidebar_loss(ctx, field, "dropped", avail, needed, size, w, body)
                    return
            else:
                kept_text = " ".join(kept)
                _sidebar_loss(
                    ctx, field, "truncated", avail, needed, size, w, body[len(kept_text) :]
                )
                fitted = kept_text
        if title:
            box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(w - 2 * pad), Emu(220_000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.name = FONT
            r.font.bold = bold_title
            r.font.size = Pt(FS_BODY)
            r.font.color.rgb = NAVY
            cy += 200_000
        bh = measure_text_height(fitted, w - 2 * pad, size, line_spacing=1.2)
        box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(w - 2 * pad), Emu(max(bh, 120_000)))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = fitted
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = BODY_COLOR
        cy += bh + gap

    write_block(None, headline, field="headlineConclusion", size=FS_BODY, required=True)
    write_block(mid_title, mid_body, field="whatIsVisible", size=11)
    if meaning and meaning != mid_body and meaning != headline:
        write_block("Что это значит", meaning, field="clientMeaning", size=11)
    if action:
        write_block("Что сделать", action, field="recommendedActions", size=11)
    if provenance:
        # Fine print, no frame
        if cy >= max_bottom - 80_000:
            # Ветка достижима только когда обязательный вывод не поместился и
            # заменён запасной фразой: у прочих блоков в запасе остаётся
            # 160 000 EMU, и подпись после них влезает всегда.
            _sidebar_loss(
                ctx,
                "provenanceLabel",
                "dropped",
                max_bottom - cy,
                measure_text_height(provenance, w - 2 * pad, FS_CAPTION, line_spacing=1.2),
                FS_CAPTION,
                w,
                provenance,
            )
        if cy < max_bottom - 80_000:
            box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(min(cy, max_bottom - 120_000)), Emu(w - 2 * pad), Emu(140_000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = provenance
            r.font.name = FONT
            r.font.size = Pt(FS_CAPTION)
            r.font.color.rgb = MUTED_COLOR



def _tone_fill(tone: str) -> RGBColor:
    """Подложка карточки по тону.

    Словарей тонов в проекте два, и это осознанно: плитки метрик говорят
    `risk|warn|good|accent`, а клиентская шкала риска — `danger|warn|neutral`
    (`orion-golden/client/risk-scale.ts`). Знать надо оба: `danger` был
    рендереру неизвестен и красился дефолтом, то есть **белым**, — и карточка
    «Высокий» выходила единственной без тревожного фона, при пяти красных
    делениях шкалы (пункт BX бэклога).

    `neutral` красится белым намеренно: нижняя ступень тревоги не несёт.
    """
    return {
        "risk": RISK_BG,
        "danger": RISK_BG,
        "warn": WARN_BG,
        "good": GOOD_BG,
        "accent": ACCENT_SOFT,
        "neutral": CARD_BG,
    }.get(tone, CARD_BG)


def _tone_value_color(tone: str) -> RGBColor:
    # Нейтральная метрика окрашена брендовым зелёным: цифра — то, ради чего
    # плитку смотрят, и она обязана быть заметнее подписи под ней.
    return {
        "risk": TONE_RISK,
        "warn": TONE_WARN,
        "good": TONE_GOOD,
        "neutral": METRIC_ACCENT,
    }.get(tone, METRIC_ACCENT)


def _render_kpi_cards(ctx: _Ctx, metrics: list[dict[str, Any]], x: int, y: int, width: int, cols: int = 2) -> int:
    items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:6]
    if not items:
        return y
    # Геометрия плитки осталась прежней намеренно. В cleeq-варианте плитка выше
    # (900 000) и с большим зазором (100 000); замер показал цену: на стр.30
    # рендерер выбросил две строки (1 494 449 при 1 294 368), потому что ёмкость
    # страницы в строках объявлена в TS-шаблоне и откалибрована по нынешней
    # высоте обвязки. Воздух вокруг цифр не стоит потерянного содержимого.
    gap = 80_000
    card_w = (width - gap * (cols - 1)) // cols
    card_h = 780_000
    row_y = y
    for idx, m in enumerate(items):
        col = idx % cols
        if idx > 0 and col == 0:
            row_y += card_h + gap
        cx = x + col * (card_w + gap)
        tone = str(m.get("tone") or "neutral")
        # Keep room for Russian status phrases like «Данные не собраны» / «0 / 10».
        value = _clip_words(_safe(m.get("value")), 36)
        label = _clip_words(_safe(m.get("label")), 28)
        # Белая плитка на мятном листе + скруглённая полоса тона: плитка
        # читается как отдельная плоскость, а не как заливка тоном.
        ctx.card(row_y, h=card_h, x=cx, w=card_w, fill=WHITE, border=None, radius=0.1)
        stripe = ctx.slide.shapes.add_shape(
            5, Emu(cx + 55_000), Emu(row_y + 120_000), Emu(60_000), Emu(card_h - 240_000)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = _tone_value_color(tone)
        stripe.line.fill.background()
        try:
            stripe.adjustments[0] = 0.5
        except Exception:  # noqa: BLE001
            pass
        box = ctx.slide.shapes.add_textbox(
            Emu(cx + 180_000), Emu(row_y + 100_000), Emu(card_w - 250_000), Emu(card_h - 180_000)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = value
        r0.font.name = FONT
        r0.font.bold = True
        r0.font.size = Pt(FS_LEAD if len(value) <= 10 else FS_SUBTITLE if len(value) <= 22 else FS_BODY)
        r0.font.color.rgb = _tone_value_color(tone)
        p1 = tf.add_paragraph()
        p1.space_before = Pt(5)
        r1 = p1.add_run()
        r1.text = label
        r1.font.name = FONT
        r1.font.size = Pt(FS_BODY)
        r1.font.color.rgb = MUTED_COLOR
    rows = (len(items) + cols - 1) // cols
    return y + rows * card_h + max(0, rows - 1) * gap



def _render_status_badge(ctx: _Ctx, badge: dict[str, Any] | None, x: int, y: int, width: int) -> int:
    if not isinstance(badge, dict) or not _safe(badge.get("label")):
        return y
    tone = str(badge.get("tone") or "neutral")
    h = 360_000
    ctx.card(y, h=h, x=x, w=width, fill=_tone_fill(tone))
    box = ctx.slide.shapes.add_textbox(Emu(x + 90_000), Emu(y + 90_000), Emu(width - 180_000), Emu(200_000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _clip_words(_safe(badge.get("label")), 48)
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(FS_SUBTITLE)
    r.font.color.rgb = _tone_value_color(tone)
    return y + h



def _render_visual_with_sidebar(
    ctx: _Ctx,
    slide: dict[str, Any],
    assets: dict[str, dict[str, Any]],
    title: str,
) -> None:
    """Title + left visual (contain) + right analytical sidebar.
    PDF-36 E.4 — pages without a visual switch to a full-width card layout
    instead of wasting 62% of the page on «материал недоступен»."""
    ctx.light_bg()
    y = ctx.title(title, 280000, NAVY)
    refs = slide.get("assetRefs") or []
    visual = _first_visual_asset(refs, assets)
    has_sidebar = bool(slide.get("visualAnalysis") or slide.get("clientTakeaway") or slide.get("bullets"))
    if not visual:
        _render_analysis_cards_full_width(ctx, slide, y)
        return
    img_w = int(CONTENT_W * 0.62) if has_sidebar else CONTENT_W
    side_w = CONTENT_W - img_w - 120000
    img_h = CONTENT_BOTTOM - y - 80000
    # Temporarily embed into a narrower box by using contain math inline
    raw = _resolve_image_bytes(visual)
    if raw:
        iw, ih = img_w, img_h
        if Image is not None:
            try:
                with Image.open(io.BytesIO(raw)) as im:
                    iw, ih = im.size
            except Exception:  # noqa: BLE001
                pass
        # Top-align visual with sidebar (do not vertically center in full column).
        scale = min(img_w / max(iw, 1), (img_h - 60000) / max(ih, 1))
        dw, dh = int(iw * scale), int(ih * scale)
        left = MARGIN_X + (img_w - dw) // 2
        top = y + 60000
        ctx.slide.shapes.add_picture(io.BytesIO(raw), Emu(left), Emu(top), width=Emu(dw), height=Emu(dh))
    else:
        ctx.body("Визуальный материал недоступен.", y + 80000, max_h=600000, color=MUTED_COLOR)
    if has_sidebar and side_w > 400000:
        _sidebar_analysis(ctx, slide, MARGIN_X + img_w + 120000, y + 60000, side_w, img_h - 60000)


def _render_analysis_cards_full_width(ctx: _Ctx, slide: dict[str, Any], y: int) -> None:
    """No-visual layout: headline accent card, two analysis cards side by
    side, an action card and fine-print provenance — the sidebar content
    spread over the whole page instead of a starved right column."""
    analysis = slide.get("visualAnalysis") or {}
    if not isinstance(analysis, dict):
        analysis = {}
    headline = _sidebar_sanitize_field(
        ctx,
        "headlineConclusion",
        _safe(analysis.get("headlineConclusion") or slide.get("clientTakeaway") or slide.get("narrative") or ""),
    )
    visible = _sidebar_sanitize_field(
        ctx, "whatIsVisible", _safe(analysis.get("whatIsVisible") or "")
    )
    meaning = _sidebar_sanitize_field(
        ctx, "clientMeaning", _safe(analysis.get("clientMeaning") or analysis.get("whyItMatters") or "")
    )
    actions = analysis.get("recommendedActions") or []
    action = _sidebar_sanitize_field(
        ctx, "recommendedActions", _safe(actions[0]) if isinstance(actions, list) and actions else ""
    )
    provenance = _safe(analysis.get("provenanceLabel") or "")
    explanations = analysis.get("highlightExplanations") or []
    reasons = [
        _safe(ex.get("clientReason"))
        for ex in explanations
        if isinstance(ex, dict) and _safe(ex.get("clientReason"))
    ][:2]

    if not any((headline, visible, meaning, action)):
        reason = _safe(slide.get("blockedReason") or "Материалы по данному разделу недоступны.")
        ctx.body(reason, y + 80_000, max_h=800_000, color=MUTED_COLOR)
        return

    if headline:
        y = ctx.content_card(
            title=None,
            text=headline,
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=340_000,
            max_h=1_200_000,
            tone="accent",
            body_size=FS_BODY,
        )
        y += 110_000
    gap = 110_000
    col_w = (CONTENT_W - gap) // 2
    left_body = visible or (reasons[0] if reasons else "")
    right_body = meaning or (reasons[1] if len(reasons) > 1 else "")
    if left_body and right_body:
        max_col_h = max(600_000, CONTENT_BOTTOM - y - (620_000 if action else 200_000))
        lb = ctx.content_card(
            title="Что показывает раздел",
            text=left_body,
            x=MARGIN_X,
            y=y,
            width=col_w,
            min_h=420_000,
            max_h=max_col_h,
            tone="neutral",
            title_size=11,
            body_size=11,
        )
        rb = ctx.content_card(
            title="Что это значит",
            text=right_body,
            x=MARGIN_X + col_w + gap,
            y=y,
            width=col_w,
            min_h=420_000,
            max_h=max_col_h,
            tone="neutral",
            title_size=11,
            body_size=11,
        )
        y = max(lb, rb) + gap
    elif left_body or right_body:
        y = ctx.content_card(
            title="Что показывает раздел" if left_body else "Что это значит",
            text=left_body or right_body,
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=420_000,
            max_h=max(600_000, CONTENT_BOTTOM - y - (620_000 if action else 200_000)),
            tone="neutral",
            title_size=11,
            body_size=11,
        )
        y += gap
    if action and y < CONTENT_BOTTOM - 420_000:
        y = ctx.content_card(
            title="Что сделать",
            text=action,
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=340_000,
            max_h=900_000,
            tone="warn",
            title_size=11,
            body_size=11,
        )
        y += 90_000
    if provenance and y < CONTENT_BOTTOM - 200_000:
        ctx.body(provenance, y, max_h=300_000, color=MUTED_COLOR, font_size=FS_CAPTION)


def _title_line_estimate(text: str, col_width_emu: int, font_pt: float, max_lines: int = 2) -> int:
    """Грубая оценка числа строк с потолком `max_lines`.

    Единственный потребитель — плашка темы в карточке матрицы рисков
    (`executive.py`), где потолок и есть смысл: плашка не вправе расти выше
    трёх строк. Высоты таблиц считает не она, а `_wrapped_line_count` — тем же
    переносом и тем же шрифтом, которыми рисуют; здешняя модель занижала строку
    (потолок «две строки» при 10 pt на ячейке, которую красят 9-м), и таблица
    уезжала ниже поля при «чистой» объявленной геометрии.
    """
    text = (text or "").strip()
    if not text:
        return 1
    char_w = font_pt * EMU_PER_PT * 0.52
    max_chars = max(1, int(col_width_emu / char_w))
    words = text.split()
    lines = 1
    cur = 0
    for w in words:
        add = len(w) if cur == 0 else cur + 1 + len(w)
        if add <= max_chars:
            cur = add
        else:
            lines += 1
            cur = min(len(w), max_chars)
            if lines >= max_lines:
                return max_lines
    return min(lines, max_lines)


#: Цвет точки по ступени: danger / warn / neutral клиентской шкалы.
_STEP_DOT_COLORS = {
    "high": RGBColor(0xB9, 0x1C, 0x1C),
    "medium": RGBColor(0xC2, 0x41, 0x0C),
    "low": RGBColor(0x64, 0x74, 0x8B),
}


def _status_tone(status: str) -> tuple[str, "RGBColor"]:
    s = (status or "").strip().lower()
    # C.4 — negative/sanction values must never carry a green marker.
    if "нежелат" in s or "негатив" in s or "санкц" in s or "критич" in s:
        return "●", RGBColor(0xB9, 0x1C, 0x1C)
    # LIKELY_SUBJECT (§2.1) and manual-review statuses — amber, not green.
    # «Не подтверждено», «статус не зафиксирован», «Не проверено» — незакрытый
    # вопрос, а не риск и не «всё в порядке»: янтарный. Проверяется до ветки
    # «подтверждено» ниже, иначе отрицание попало бы в неё по подстроке.
    #
    # «Не проверено» — оценка строки выдачи, чью страницу не открывали. Без
    # своей ветки она красилась серым — тем же цветом, что «Нейтральный», и
    # разница между «проверили, чисто» и «не заходили» пропадала на листе.
    if (
        "вероятн" in s
        or "проверк" in s
        or "не провер" in s
        or "требует" in s
        or "pep" in s
        or "не подтвержд" in s
        or "не зафиксирован" in s
    ):
        return "●", RGBColor(0xC2, 0x41, 0x0C)
    # Подтверждённое аналитиком совпадение — подтверждённый комплаенс-риск.
    # Зелёная точка здесь читалась бы как одобрение (то же правило, что для
    # «санкц»), а незнакомое слово в этой функции зелёное по умолчанию.
    if "подтвержд" in s:
        return "●", RGBColor(0xB9, 0x1C, 0x1C)
    # Материал о другом лице занимает своё место в выдаче, но оценкой субъекта
    # не является: полый маркер и серый цвет отличают его и от негатива, и от
    # зелёного «всё в порядке». Зелёный здесь читался бы как одобрение
    # однофамильца.
    if "друго" in s:
        return "○", RGBColor(0x94, 0xA3, 0xB8)
    # Слова клиентской шкалы. Ступень узнаёт то же место, что и шкала делений,
    # — второго словаря ступеней в рендерере нет.
    step = level_step(s)
    if step:
        return "●", _STEP_DOT_COLORS[step]
    # Зелёный достаётся по имени. Единственный его владелец — «Позитивный»:
    # так называется маркер в легенде таблицы. Пока зелёный был умолчанием,
    # легенда обещала маркер, которого эта функция не знала вовсе.
    if "позитив" in s:
        return "●", RGBColor(0x04, 0x78, 0x57)
    # E.6 — neutral verdicts read gray, green stays for explicit positives.
    if "нейтрал" in s or s in {"·", "—", "-", ""}:
        return "●", RGBColor(0x64, 0x74, 0x8B)
    # Незнакомый статус — серый, а не зелёный.
    #
    # Зелёный был умолчанием, и любое слово, которого функция не знает, читалось
    # в отчёте для банка как «всё в порядке». Опасность была записана в
    # комментариях дважды и один раз уже стоила ветки: «Подтверждено
    # аналитиком» красилось зелёным, пока для него не завели своё правило.
    #
    # Замер на живом прогоне 22.08: все шесть встреченных значений попадают в
    # названные ветки — то есть зелёный достижим только незнакомым словом.
    # Неизвестность не благополучие (пункт V).
    return "●", RGBColor(0x64, 0x74, 0x8B)


def _add_search_table(
    ctx: _Ctx,
    y: int,
    headers: list[str],
    rows: list[list[str]],
    groups: list[dict[str, Any]] | None = None,
    *,
    bottom: int | None = None,
    declared_top: int | None = None,
) -> None:
    """
    Grouped SERP position table. Renders EVERY row the slide carries (no cap) —
    TS pagination already guaranteed geometric fit. Query is shown as a compact
    group-header band (spec §4), status as a colored badge (spec §5).

    Полосы адреса под строкой больше нет: адрес вернулся в колонку «Ссылка», и
    печатать его вторым способом значило бы печатать один факт дважды. Ветка
    полосы снята вместе с параметром `row_addresses` — её единственным входом
    были страницы выдачи, а они теперь пятиколоночные с адресом в ячейке.
    Оставить её «на всякий случай» — завести мёртвый путь, который следующий
    читатель примет за живой.

    `bottom` — низ бюджета страницы (низ белой сцены). Превышение пишется
    событием разметки уровня CRITICAL: таблица, нарисованная ниже поля, — это
    та же тихая потеря содержимого, что и невлезший буллет.

    `declared_top` — **объявленный** верх таблицы, от которого считается бюджет
    меры. Он не равен `y`: фактический верх зависит от длины вводного абзаца, а
    его переписывает стадия 2 уже **после** того, как мера снята. Ёмкость,
    выведенная из факта, была бы верна для черновика и завышена для готовой
    страницы — то есть таблица уехала бы мимо поля. Верх не объявлен (или не
    объявлен низ) — меры нет вовсе, и построитель остаётся с раскладкой сида.
    """
    # Body layout is 4 cols: Позиция | Домен | Заголовок | Статус.
    # If TS sends a leading «Запрос» column, drop it — query lives in group bands.
    hdr = [str(h) for h in headers]
    data_rows = [list(r) for r in rows]
    if len(hdr) >= 5 and re.search(r"запрос|query", hdr[0], re.I):
        hdr = hdr[1:]
        data_rows = [r[1:] if len(r) > 1 else r for r in data_rows]
    cols = max(1, min(5, len(hdr)))
    headers = hdr
    groups = groups or []

    # Row plan: header + interleaved group bands + data rows.
    plan: list[tuple[str, Any]] = [("header", headers)]
    if groups:
        for g in groups:
            start = int(g.get("rowStart", 0))
            count = int(g.get("rowCount", 0))
            label = str(g.get("queryDisplay") or "")
            qtag = g.get("qTag")
            band = f"Запрос: {label}" if not qtag else f"{qtag} — {label}"
            plan.append(("group", band))
            for r in data_rows[start : start + count]:
                plan.append(("data", r))
    else:
        for r in data_rows:
            plan.append(("data", r))

    # Ширины колонок выбираются по смыслу заголовков, а не по их длине.
    # Прежний признак `len(headers[0]) > 3` был прокси вопроса «первая колонка
    # номерная?», и прокси ошибался в обе стороны: «Тема» (4 буквы) уходила в
    # номерную ветку и получала 14 % ширины при двух счётчиках на 86 %, а
    # «Поз.» — в текстовую и получала 14 % под двузначное число.
    if cols == 2:
        prop = [0.24, 0.76]
    elif cols == 5 and re.search(r"^\s*№", str(headers[0])):
        # Первая таблица выдачи: № | Ссылка | Заголовок | Тип источника | Оценка.
        #
        # Признак ветки — колонка «№», и он выбран не случайно: **у второй
        # таблицы выдачи колонки позиции нет вовсе**, это решение владельца и
        # оно закреплено юнитом. Значит признак не сломается от переименования
        # любой другой колонки — в отличие от признака по «Найдено по запросу»,
        # который жил бы ровно до первой правки формулировки.
        #
        # Прежние доли `[0.05, 0.22, 0.44, 0.15, 0.14]` — те самые, при которых
        # адрес не открывался: 22 % это 229 px полезных, куда входит 62 знака.
        # Новые померены `_wrapped_line_count` на корпусе прогона 72 и на
        # предельных значениях построителя, а не подобраны:
        #
        #   0.05 «№»            — двузначное число;
        #   0.34 «Ссылка»       — 328 px полезных: адрес корпуса ложится в
        #                         1…3 строки у 45 строк из 46, предел 165
        #                         знаков самым широким знаком — 7 строк, из
        #                         которых и выведена ёмкость листа;
        #   0.27 «Заголовок»    — 257 px: предел 95 знаков не больше 5 строк;
        #   0.20 «Тип источника» — 186 px: «Официальный сайт / госресурс» в одну
        #                         строку (0.18 даёт две);
        #   0.14 «Оценка»       — 125 px: «● Нежелательный» в одну строку (0.13
        #                         даёт две).
        prop = [0.05, 0.34, 0.27, 0.20, 0.14]
    elif cols == 5:
        # Вторая таблица выдачи: Ссылка | Заголовок | Найдено по запросу | Тип
        # источника | Оценка. Номера строк здесь нет.
        #
        # Доли **уравнивают три широкие колонки**: при пределах построителей
        # (`SERP_ADDRESS_MAX_CHARS` 165, `SERP_TITLE_MAX_CHARS` 95,
        # `SERP_FOUND_BY_MAX_CHARS` 80 — все три режут в самом построителе) и
        # самом широком знаке 9 pt каждая из них даёт ровно семь нарисованных
        # строк, то есть худшая законная строка второй таблицы равна худшей
        # строке первой и ёмкость листа у обеих одна.
        #
        #   0.30 «Ссылка»            — 287 px полезных: тот же предел адреса 165
        #                              знаков, что и у первой таблицы (он один на
        #                              обе и берётся по этой, узкой, колонке);
        #   0.20 «Заголовок»         — 186 px: 95 знаков в 7 строк;
        #   0.16 «Найдено по запросу» — 145 px: 80 знаков в 7 строк;
        #   0.20 «Тип источника»     — 186 px: «Официальный сайт / госресурс» в одну;
        #   0.14 «Оценка»            — 125 px: «● Нежелательный» в одну.
        #
        # Замер отвергает доли `[0.30, 0.22, 0.14, 0.20, 0.14]`: при 0.14 запрос
        # предельной длины даёт **восемь** строк (1 173 480 EMU), ёмкость падает
        # до 2, и число листов второй таблицы удваивается.
        prop = [0.30, 0.20, 0.16, 0.20, 0.14]
    elif cols == 3:
        # Текст плюс счётчики («Тема | Публикаций | Из них нежелательных»):
        # ведёт текстовая колонка. 0.60 держит тему предельной длины (120
        # символов, 868 px при 10pt) в двух строках, 0.20 вмещает самый длинный
        # заголовок счётчика в одну.
        prop = [0.60, 0.20, 0.20]
    elif cols == 4 and re.search(r"баз[аы]\s+данных", str(headers[0]), re.I):
        # Комплаенс-сводка: «База данных | Тип совпадения | Совпадение по имени
        # | Статус проверки». Общая четырёхколоночная ветка отдавала 42 %
        # ширины третьей колонке и 18 % статусу, где стоит самая длинная
        # законная строка отчёта — «Не подтверждено (статус в артефактах
        # прогона не зафиксирован)»: она не влезала в две строки, LibreOffice
        # тянул строку по содержимому и таблица уезжала вниз. Доли выверены
        # настоящими метриками шрифта, а не моделью `_title_line_estimate`
        # (у той потолок — две строки, и на длинной ячейке она молчит).
        #
        # Ветка узнаётся по первой колонке, а не по словам третьей: третья
        # печатала «Оценку совпадения», теперь печатает имя записи, и признак,
        # завязанный на неё, молча отправил бы таблицу в общую ветку —
        # заголовки прежние, ширины чужие. Имени нужна ширина под трёхчастное
        # ФИО заглавными («КИРИЛЛ СЕРГЕЕВИЧ КУЛЕБАКИН») в одну строку.
        prop = [0.14, 0.26, 0.26, 0.34]
    elif headers and re.search(r"^\s*(№|поз)", str(headers[0]), re.I):
        prop = [0.07, 0.22, 0.53, 0.18][:cols]
    else:
        prop = [0.14, 0.26, 0.42, 0.18][:cols]
    widths = [max(500_000, int(CONTENT_W * p)) for p in prop]
    leftover = CONTENT_W - sum(widths)
    if leftover != 0 and widths:
        widths[2 if cols > 2 else len(widths) - 1] += leftover

    # C.4 — the colored status badge belongs only to a genuine status column;
    # generic value columns («Значение», «Комментарий») stay plain text so
    # negative categories never receive a misleading green marker.
    last_header = str(headers[cols - 1]) if cols - 1 < len(headers) else ""
    # E.6 — «Оценка» is the SERP verdict column and must carry the badge too.
    badge_last_col = bool(re.search(r"статус|риск|провер|оценк", last_header, re.I))

    def _is_badge(col: int) -> bool:
        return col == cols - 1 and badge_last_col

    def _cell_font_pt(col: int) -> float:
        """Кегль ячейки — тот, которым её и красят: у бейджа он свой."""
        return BADGE_PT if _is_badge(col) else float(FS_CAPTION)

    def _painted(col: int, value: str) -> str:
        """Текст ячейки в том виде, в каком он будет нарисован.

        У бейджа впереди точка и пробел — два знака, которых мера не видела.
        Щель та же, что была у кегля: меряем одно, рисуем другое.
        """
        return f"{_status_tone(value)[0]} {value}" if _is_badge(col) else value

    def _cell_height(text: str, width_emu: int, pt: float) -> int:
        """Сколько займёт ячейка: настоящий перенос по полезной ширине."""
        lines = _wrapped_line_count(text, max(1, width_emu - CELL_MARGINS_EMU), pt)
        return lines * int(pt * EMU_PER_PT * 1.2)

    # Per-row heights.
    #
    # Высота считается тем же переносом и тем же шрифтом, которыми ячейку
    # рисуют (`_wrapped_line_count`), и **тем кеглем, которым она красится**.
    # Прежде здесь стояла своя модель с потолком в две строки при 10 pt, а
    # красили при 9: строка в три нарисованных строки объявлялась двумя,
    # LibreOffice тянул её по содержимому, таблица уезжала вниз — и объявленная
    # геометрия при этом оставалась «чистой».
    #
    # Поля ячейки python-pptx (0.1″ с каждой стороны) вычитаются: перенос
    # случается по полезной ширине, а не по ширине колонки.
    pad = int(6 * EMU_PER_PT)
    header_h = int(26 * EMU_PER_PT)
    group_h = int(18 * EMU_PER_PT)
    heights: list[int] = []
    for kind, payload in plan:
        if kind == "header":
            heights.append(header_h)
        elif kind == "group":
            heights.append(group_h)
        else:
            # Высота строки — по самой высокой ячейке, по всем колонкам сразу.
            # Списка «измеряемых колонок» здесь нет намеренно: он расходился с
            # пропорциями (мерили колонку 2 — а текст жил в колонке 0) и строка
            # объявлялась в разы ниже своего содержимого.
            heights.append(
                max(
                    _cell_height(
                        _painted(c, str(payload[c]) if c < len(payload) else ""),
                        widths[c],
                        _cell_font_pt(c),
                    )
                    for c in range(cols)
                )
                + pad
            )

    table_rows = len(plan)
    table_h = sum(heights)
    # Бюджет листа — низ белой сцены, а не низ слайда.
    #
    # Прежде отрисовщик бюджета не знал вовсе: `slides.py` звал `content_stage`
    # и выбрасывал возвращаемое значение, а таблица рисовала ту высоту, которая
    # получилась. Нарисованное мимо страницы — та же тихая потеря содержимого,
    # что и невлезший буллет, поэтому превышение объявляется событием CRITICAL
    # (`TABLE_ROW_PARTIALLY_VISIBLE`), а не остаётся молча на растре.
    # Мера таблицы пишется **всегда**, а не только при переполнении: раскрой по
    # ней получают и листы, которые влезли, — иначе построитель не знает,
    # сколько на них осталось места, и режет по худшему законному случаю. До
    # этой записи о таблицах не сообщалось ничего: в телеметрии прогона 91 было
    # 62 записи при 39 листах с таблицами и ни одной с ролью «table».
    #
    # Потерей вердикта переполнение таблицы **не** объявляется. Строки таблицы
    # циклу перекладки буллетов не подвластны: увидев `droppedLines > 0` на
    # странице, которую он двинуть не может, цикл объявил бы несходимость и
    # остановил бы оплаченный прогон из-за одной высокой строки. О переполнении
    # говорит запись разметки ниже — она и есть громкое событие.
    if bottom is not None and declared_top is not None:
        record_bullet_measure(
            slide_key=f"{ctx.slide_key}{TABLE_MEASURE_KEY_SUFFIX}",
            page=ctx.page,
            available_height=max(0, bottom - declared_top),
            max_items=len(plan),
            item_heights=heights,
            kept_items=len(plan),
            dropped_bullets=0,
            dropped_lines=0,
        )
    if bottom is not None and y + table_h > bottom:
        record_text_layout(
            page=ctx.page,
            name=f"orion_search_table_p{ctx.page}",
            role="table",
            font_family=FONT,
            font_size_pt=FS_CAPTION,
            box_width=CONTENT_W,
            box_height=table_h,
            available_height=max(0, bottom - y),
            required_height=table_h,
            measured_lines=len(plan),
            text_length=sum(len(str(payload)) for _kind, payload in plan),
            clipped=True,
        )
    shape = ctx.slide.shapes.add_table(table_rows, cols, Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(table_h))
    tbl = shape.table
    for i, w in enumerate(widths):
        tbl.columns[i].width = Emu(w)
    for i, h in enumerate(heights):
        tbl.rows[i].height = Emu(h)

    def paint(cell: Any, text: str, *, bold: bool = False, color: Any = BODY_COLOR, bg: Any = WHITE, size: float = FS_CAPTION, clip: bool = True) -> None:
        # Status badges ("● Нежелательный") are complete labels, not clipped
        # prose — the dangling-tail trimmer would strip the word after the dot.
        cell.text = _clip_words(text, 200) if clip else _safe(text)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        cell.text_frame.word_wrap = True
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg

    for r_idx, (kind, payload) in enumerate(plan):
        if kind == "header":
            for c in range(cols):
                label = str(payload[c]) if c < len(payload) else ""
                # Шапка таблицы — зелёная полоса cleeq. Текст на ней чернилами,
                # а не белым: белое по #24D875 даёт контраст 1,8:1, то есть
                # заголовок столбца пришлось бы угадывать.
                paint(tbl.cell(r_idx, c), label, bold=True, color=NAVY, bg=ACCENT, size=FS_CAPTION)
        elif kind == "group":
            merged = tbl.cell(r_idx, 0)
            merged.merge(tbl.cell(r_idx, cols - 1))
            paint(merged, str(payload), bold=True, color=NAVY, bg=ACCENT_SOFT, size=FS_CAPTION)
        else:
            row = payload
            status = str(row[cols - 1] if len(row) >= cols else "").strip()
            status_l = status.lower()
            # Подсветка строки — от статусной колонки, и только от неё. В
            # карточке записи последняя колонка называется «Значение», и по
            # слову «требует» янтарь доставался строке «Категория: Требует
            # ручной классификации» — то есть окрашивалась не та строка. Правило
            # то же, что у бейджа ниже: обобщённая колонка значений статусом не
            # является.
            adverse = badge_last_col and "нежелат" in status_l
            likely = badge_last_col and (
                "вероятн" in status_l or "проверк" in status_l or "требует" in status_l
            )
            if adverse:
                row_bg = RGBColor(0xFE, 0xF2, 0xF2)
            elif likely:
                row_bg = RGBColor(0xFF, 0xF7, 0xED)  # soft amber for LIKELY
            else:
                row_bg = WHITE
            for c in range(cols):
                val = str(row[c]) if c < len(row) else ""
                if c == cols - 1 and badge_last_col:
                    dot, tone = _status_tone(val)
                    paint(tbl.cell(r_idx, c), f"{dot} {val}", color=tone, bg=row_bg, size=BADGE_PT, clip=False)
                else:
                    paint(tbl.cell(r_idx, c), val, bg=row_bg, size=FS_CAPTION)


