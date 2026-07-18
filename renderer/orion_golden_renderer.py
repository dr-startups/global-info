"""ORION Golden Report renderer — R10 deterministic PPTX/PDF from ReportSpec + deck manifest."""

from __future__ import annotations

import base64
import io
import json
import os
import re
import tempfile
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

FONT = "DejaVu Sans"
FS_TITLE = 26
FS_SECTION = 22
FS_SUBTITLE = 13
FS_BODY = 11  # min readable body ≥ 10.5pt
FS_CAPTION = 9  # footer/provenance 8.5–9pt

# Master slide 16:10 (12.8" × 8.0") — matches ORION reference aspect.
SLIDE_W = 11_704_320
SLIDE_H = 7_315_200
MARGIN_X = 480_000
CONTENT_W = SLIDE_W - 2 * MARGIN_X
FOOTER_Y = SLIDE_H - 440_000
CONTENT_BOTTOM = SLIDE_H - 700_000

NAVY = RGBColor(0x0B, 0x1A, 0x33)
TITLE_COLOR = RGBColor(0xF8, 0xFA, 0xFC)
BODY_COLOR = RGBColor(0x33, 0x41, 0x55)
MUTED_COLOR = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_SOFT = RGBColor(0xEF, 0xF6, 0xFF)
WARN_BG = RGBColor(0xFF, 0xF7, 0xED)
RISK_BG = RGBColor(0xFE, 0xF2, 0xF2)
GOOD_BG = RGBColor(0xEC, 0xFD, 0xF5)
TONE_RISK = RGBColor(0xB9, 0x1C, 0x1C)
TONE_WARN = RGBColor(0xC2, 0x41, 0x0C)
TONE_GOOD = RGBColor(0x04, 0x78, 0x57)

FORBIDDEN = re.compile(
    r"(storage/|C:\\\\|openai[_-]?api[_-]?key|cmr[a-z0-9]{10,}|adverse_media|requires_review)",
    re.I,
)

# EMU helpers: 914400 EMU = 1 inch; 72 pt = 1 inch
EMU_PER_PT = 12_700
EMU_PER_INCH = 914_400

# Renderer-side text layout telemetry (page, font, box, measured vs available).
_LAYOUT_TELEMETRY: list[dict[str, Any]] = []


def reset_layout_telemetry() -> None:
    _LAYOUT_TELEMETRY.clear()


def get_layout_telemetry() -> list[dict[str, Any]]:
    return list(_LAYOUT_TELEMETRY)


def record_text_layout(
    *,
    page: int,
    name: str,
    role: str,
    font_family: str,
    font_size_pt: float,
    box_width: int,
    box_height: int,
    available_height: int,
    required_height: int,
    measured_lines: int,
    text_length: int,
    clipped: bool,
    measurement_uncertain: bool = False,
) -> None:
    _LAYOUT_TELEMETRY.append(
        {
            "page": page,
            "name": name,
            "role": role,
            "fontFamily": font_family,
            "fontSizePt": font_size_pt,
            "boxWidth": box_width,
            "boxHeight": box_height,
            "availableHeight": available_height,
            "requiredHeight": required_height,
            "measuredLines": measured_lines,
            "textLength": text_length,
            "clipped": clipped,
            "measurementUncertain": measurement_uncertain,
        }
    )


def _count_measured_lines(
    text: str,
    width_emu: int,
    font_size_pt: float,
) -> tuple[int, bool]:
    """Return (line_count, measurement_uncertain)."""
    raw = _safe(text)
    if not raw:
        return 1, False
    width_px = max(40, int(width_emu / EMU_PER_INCH * 96 * 0.90))
    paragraphs = [p.strip() for p in re.split(r"\n+", raw) if p.strip()] or [""]
    total_lines = 0
    uncertain = False
    font = None
    if Image is not None:
        try:
            from PIL import ImageFont  # type: ignore

            fp = _font_path()
            if fp:
                font = ImageFont.truetype(fp, size=max(8, int(round(font_size_pt * 96 / 72))))
            else:
                uncertain = True
        except Exception:  # noqa: BLE001
            uncertain = True
            font = None
    else:
        uncertain = True

    for para in paragraphs:
        words = para.split(" ")
        if not words:
            total_lines += 1
            continue
        line = ""
        lines = 1
        for word in words:
            trial = word if not line else f"{line} {word}"
            if font is not None:
                try:
                    bbox = font.getbbox(trial)
                    tw = bbox[2] - bbox[0]
                except Exception:  # noqa: BLE001
                    tw = int(len(trial) * font_size_pt * 0.58 * 96 / 72)
                    uncertain = True
            else:
                tw = int(len(trial) * font_size_pt * 0.58 * 96 / 72)
            if tw <= width_px or not line:
                line = trial
            else:
                lines += 1
                line = word
        total_lines += max(1, lines)
    return max(1, total_lines), uncertain


def _font_path() -> str | None:
    """Return DejaVu Sans file used for both measurement and PPTX family."""
    here = Path(__file__).resolve().parent
    candidates = [
        os.environ.get("ORION_RENDER_FONT"),
        str(here / "fonts" / "DejaVuSans.ttf"),
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        r"C:\Windows\Fonts\DejaVuSans.ttf",
    ]
    for path in candidates:
        if path and Path(path).is_file():
            return path
    return None


def assert_render_font_family() -> str:
    """Startup/QA: measurement font must be DejaVu Sans family."""
    fp = _font_path()
    if not fp:
        raise RuntimeError("ORION render font missing: expected DejaVuSans.ttf under renderer/fonts")
    name = Path(fp).name.lower()
    if "dejavu" not in name:
        raise RuntimeError(f"ORION render font mismatch: expected DejaVu Sans, got {fp}")
    # Optional Linux check
    try:
        import subprocess

        out = subprocess.check_output(["fc-match", "DejaVu Sans"], text=True, stderr=subprocess.DEVNULL)
        if "DejaVu" not in out and "dejavu" not in out.lower():
            raise RuntimeError(f"fc-match DejaVu Sans unexpected: {out.strip()}")
    except FileNotFoundError:
        pass
    except subprocess.CalledProcessError:
        pass
    return fp


def measure_text_height(
    text: str,
    width_emu: int,
    font_size_pt: float,
    line_spacing: float = 1.2,
    paragraph_spacing_pt: float = 6.0,
) -> int:
    """Measure wrapped text height in EMU using real font metrics when available."""
    raw = _safe(text)
    if not raw:
        return int(font_size_pt * EMU_PER_PT * line_spacing)
    # Slightly narrower than box so PPTX wrap is not underestimated.
    width_px = max(40, int(width_emu / EMU_PER_INCH * 96 * 0.90))
    paragraphs = [p.strip() for p in re.split(r"\n+", raw) if p.strip()] or [""]
    total_lines = 0
    font = None
    if Image is not None:
        try:
            from PIL import ImageFont  # type: ignore

            fp = _font_path()
            if fp:
                font = ImageFont.truetype(fp, size=max(8, int(round(font_size_pt * 96 / 72))))
        except Exception:  # noqa: BLE001
            font = None

    for para in paragraphs:
        words = para.split(" ")
        if not words:
            total_lines += 1
            continue
        line = ""
        lines = 1
        for word in words:
            trial = word if not line else f"{line} {word}"
            if font is not None:
                try:
                    bbox = font.getbbox(trial)
                    tw = bbox[2] - bbox[0]
                except Exception:  # noqa: BLE001
                    tw = int(len(trial) * font_size_pt * 0.58 * 96 / 72)
            else:
                tw = int(len(trial) * font_size_pt * 0.58 * 96 / 72)
            if tw <= width_px or not line:
                line = trial
            else:
                lines += 1
                line = word
        total_lines += max(1, lines)

    line_h = font_size_pt * EMU_PER_PT * line_spacing
    para_extra = max(0, len(paragraphs) - 1) * paragraph_spacing_pt * EMU_PER_PT
    # Safety margin: PPTX wraps more aggressively than PIL metrics.
    return int((total_lines * line_h + para_extra) * 1.18)


def _fit_text_to_height(
    text: str,
    width_emu: int,
    font_size_pt: float,
    max_h: int,
    *,
    line_spacing: float = 1.2,
) -> str:
    """Clip text so measured height fits max_h. Prefer complete sentences; never bare ellipsis stubs."""
    raw = _safe(text)
    if not raw:
        return ""
    if measure_text_height(raw, width_emu, font_size_pt, line_spacing=line_spacing) <= max_h:
        return raw
    # Prefer longest prefix that ends on a sentence boundary.
    sentences = re.split(r"(?<=[.!?…])\s+", raw)
    kept: list[str] = []
    for sent in sentences:
        trial = " ".join(kept + [sent]).strip()
        if measure_text_height(trial, width_emu, font_size_pt, line_spacing=line_spacing) <= max_h:
            kept.append(sent)
        else:
            break
    if kept:
        return " ".join(kept).strip()
    # Fall back to word clip without ellipsis; stop before dangling prepositions.
    words = raw.split()
    lo, hi = 1, len(words)
    best = words[0]
    while lo <= hi:
        mid = (lo + hi) // 2
        trial = " ".join(words[:mid])
        if measure_text_height(trial, width_emu, font_size_pt, line_spacing=line_spacing) <= max_h:
            best = trial
            lo = mid + 1
        else:
            hi = mid - 1
    return _trim_dangling_tail(best)


_DANGLING_TAIL = re.compile(
    r"(\s+(?:и|а|или|по|на|в|с|из|для|о|об|к|ко|у|от|до|при|без|над|под|про|через|как|что|чтобы|который|которая|которые|которых|баз|записям|доменом?|контекстом?|криминальным|компрометирующим|санкционным|нежелательный|ручной|негативным|т\.?\s*ч\.?))\s*$|"
    r"(\s+в\s+т\.?\s*ч\.?\s*$)|(\s+с\s+[А-ЯA-Z]\.?\s*$)|(\s+[А-ЯA-Z]\.?\s*$)",
    re.I,
)


def _trim_dangling_tail(text: str) -> str:
    val = text.strip()
    end = ""
    if val and val[-1] in ".!?…":
        end = val[-1]
        val = val[:-1].rstrip(".,;: ")
    for _ in range(4):
        nxt = _DANGLING_TAIL.sub("", val).rstrip(".,;: ")
        if nxt == val:
            break
        val = nxt
    if not val:
        return ""
    return val + end


def _safe(text: object) -> str:
    val = re.sub(r"\s+", " ", str(text or "")).strip()
    val = FORBIDDEN.sub("", val)
    # Humanize residual enum-like tokens that may appear in summaries
    val = re.sub(r"\bWRONG[_\s-]?SUBJECT\b", "другой субъект", val, flags=re.I)
    val = re.sub(r"\bPENDING\b", "требует проверки", val, flags=re.I)
    val = re.sub(r"\bGPT\b", "модельный анализ", val)
    return val.strip()


def _clip_words(text: str, max_chars: int) -> str:
    """Clip on sentence/word boundary; avoid mid-thought stubs and dangling prepositions."""
    val = _safe(text)
    if len(val) <= max_chars:
        return _trim_dangling_tail(val)
    slice_ = val[:max_chars]
    punct = max(slice_.rfind(". "), slice_.rfind("! "), slice_.rfind("? "), slice_.rfind("; "))
    if punct > max_chars * 0.45:
        return slice_[: punct + 1].rstrip()
    sp = max(slice_.rfind(" "), slice_.rfind("\u00a0"))
    if sp > max_chars * 0.4:
        return _trim_dangling_tail(slice_[:sp])
    soft = re.sub(r"[^\s]{1,12}$", "", slice_).rstrip()
    if len(soft) > max_chars * 0.35:
        return _trim_dangling_tail(soft)
    return _trim_dangling_tail(slice_)


class _Ctx:
    def __init__(self, prs: Presentation, page: int, total: int):
        self.prs = prs
        self.page = page
        self.total = total
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        self.slide = prs.slides.add_slide(layout)

    def footer(self) -> None:
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(FOOTER_Y), Emu(CONTENT_W), Emu(250000))
        try:
            box.name = f"orion_footer_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{self.page} / {self.total}"
        r.font.name = FONT
        r.font.size = Pt(FS_CAPTION)
        r.font.color.rgb = MUTED_COLOR

    def dark_bg(self) -> None:
        fill = self.slide.background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    def light_bg(self) -> None:
        fill = self.slide.background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

    def title(self, text: str, y: int = 280000, color: RGBColor = TITLE_COLOR, size: int = FS_TITLE) -> int:
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(text)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = color
        return y + 950000

    def card(
        self,
        y: int,
        h: int = 2_200_000,
        *,
        x: int | None = None,
        w: int | None = None,
        fill: RGBColor = CARD_BG,
    ) -> None:
        left = MARGIN_X if x is None else x
        width = CONTENT_W if w is None else w
        avail = max(200_000, min(h, CONTENT_BOTTOM - y))
        shape = self.slide.shapes.add_shape(1, Emu(left), Emu(y), Emu(width), Emu(avail))
        try:
            shape.name = f"orion_card_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = CARD_BORDER

    def body(
        self,
        text: str,
        y: int,
        max_h: int = 900000,
        color: RGBColor = BODY_COLOR,
        *,
        x: int | None = None,
        w: int | None = None,
        font_size: int = FS_BODY,
    ) -> int:
        """Render body text and return actual bottom Y from measured content height."""
        left = MARGIN_X if x is None else x
        width = CONTENT_W if w is None else w
        avail = max(200000, min(max_h, CONTENT_BOTTOM - y))
        chunks = [c.strip() for c in re.split(r"\n+", _safe(text)) if c.strip()]
        if not chunks:
            return y
        # Prefer height-fit over crude char starvation (was ~200 chars at 420k emu).
        joined_raw = "\n".join(chunks[:8])
        fitted = _fit_text_to_height(joined_raw, width, font_size, avail, line_spacing=1.2)
        fitted = _trim_dangling_tail(fitted)
        dangling = re.compile(
            r"(?:\bв\s+т\.?\s*ч\.?|\bс\s+[А-ЯA-Z]\.?|\b(?:как|что|чтобы|и|а|или|по|на|в|с)\b|,|;|—|–|-)\s*$",
            re.I,
        )
        if dangling.search(fitted) or (fitted and fitted[-1] not in ".!?…»)"):
            sentences = re.split(r"(?<=[.!?…])\s+", _safe(joined_raw))
            kept_s: list[str] = []
            for sent in sentences:
                trial = " ".join(kept_s + [sent]).strip()
                if measure_text_height(trial, width, font_size, line_spacing=1.2) <= avail:
                    kept_s.append(sent)
                else:
                    break
            # Drop a trailing incomplete clause (e.g. ends with «как»).
            while kept_s and (
                dangling.search(kept_s[-1]) or kept_s[-1][-1] not in ".!?…»)"
            ):
                kept_s.pop()
            if kept_s:
                fitted = " ".join(kept_s).strip()
            else:
                # Last resort: first sentence trimmed of dangling tail.
                first = _trim_dangling_tail(sentences[0] if sentences else fitted)
                fitted = first if first and not dangling.search(first) else _trim_dangling_tail(fitted)
        kept = [c for c in fitted.split("\n") if c.strip()] or ([fitted] if fitted else [])
        if not kept:
            return y
        joined = "\n".join(kept)
        needed = measure_text_height(joined, width, font_size, line_spacing=1.2, paragraph_spacing_pt=8)
        box_h = min(avail, max(needed + 40_000, int(font_size * EMU_PER_PT)))
        measured_lines, uncertain = _count_measured_lines(joined, width, font_size)
        # Clipping = placed text does not fit the box. Fitting/truncating source is not layout overflow.
        clipped = needed > avail
        record_text_layout(
            page=self.page,
            name=f"orion_text_body_p{self.page}",
            role="text",
            font_family=FONT,
            font_size_pt=font_size,
            box_width=width,
            box_height=box_h,
            available_height=avail,
            required_height=needed,
            measured_lines=measured_lines,
            text_length=len(joined),
            clipped=clipped,
            measurement_uncertain=uncertain,
        )
        box = self.slide.shapes.add_textbox(Emu(left), Emu(y), Emu(width), Emu(box_h))
        try:
            box.name = f"orion_text_body_p{self.page}"
        except Exception:  # noqa: BLE001
            pass
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for chunk in kept:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(8)
            r = p.add_run()
            r.text = chunk
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
        return y + box_h

    def content_card(
        self,
        *,
        title: str | None,
        text: str,
        x: int,
        y: int,
        width: int,
        min_h: int = 320_000,
        max_h: int = 2_400_000,
        padding: int = 100_000,
        tone: str = "neutral",
        title_size: int = 10,
        body_size: int = 11,
    ) -> int:
        """Draw a content-sized card; clip text to fit; return actual bottom Y."""
        fill = {
            "accent": ACCENT_SOFT,
            "warn": WARN_BG,
            "risk": RISK_BG,
            "good": GOOD_BG,
        }.get(tone, CARD_BG)
        title_s = _safe(title or "")
        body_s = _safe(text)
        budget = max(200_000, min(max_h, CONTENT_BOTTOM - y))
        # Shrink padding on short cards so body text is not starved to a stub line.
        pad = padding
        if budget < 420_000:
            pad = min(padding, 60_000)
        elif budget < 560_000:
            pad = min(padding, 80_000)
        inner_w = max(120_000, width - 2 * pad)
        title_h = 0
        if title_s:
            title_h = measure_text_height(title_s, inner_w, title_size, line_spacing=1.15) + 30_000
        full_body_h = measure_text_height(body_s, inner_w, body_size, line_spacing=1.2) if body_s else 0
        needed = 2 * pad + title_h + full_body_h + 30_000
        # Short client phrases must stay complete; height estimator is conservative and
        # otherwise collapses actions like «Исключить из digital profile…» to one word.
        short_phrase = bool(body_s) and len(body_s) <= 110 and len(body_s.split()) <= 16
        if body_s and needed > budget and not short_phrase:
            body_budget = max(100_000, budget - 2 * pad - title_h)
            body_s = _fit_text_to_height(body_s, inner_w, body_size, body_budget)
            body_h = measure_text_height(body_s, inner_w, body_size, line_spacing=1.2) if body_s else 0
        else:
            body_h = full_body_h
        h = max(min_h, min(budget, 2 * pad + title_h + body_h + 30_000))
        self.card(y, h=h, x=x, w=width, fill=fill)
        cy = y + pad
        if title_s:
            box = self.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(inner_w), Emu(max(title_h, 160_000)))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title_s
            r.font.name = FONT
            r.font.bold = True
            r.font.size = Pt(title_size)
            r.font.color.rgb = NAVY
            cy += title_h
        if body_s:
            rem = max(100_000, y + h - cy - pad)
            box = self.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(inner_w), Emu(rem))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = body_s
            r.font.name = FONT
            r.font.size = Pt(body_size)
            r.font.color.rgb = BODY_COLOR
        return y + h

    def metric_chips(self, metrics: list[dict[str, Any]], x: int, y: int, width: int) -> int:
        """Render 2-column metric chips; return bottom Y."""
        items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:4]
        if not items:
            return y
        cols = 2
        gap = 70_000
        chip_w = (width - gap) // cols
        chip_h = 640_000
        row_y = y
        for idx, m in enumerate(items):
            col = idx % cols
            if idx > 0 and col == 0:
                row_y += chip_h + gap
            cx = x + col * (chip_w + gap)
            tone = str(m.get("tone") or "neutral")
            fill = {"risk": RISK_BG, "warn": WARN_BG, "good": GOOD_BG}.get(tone, CARD_BG)
            value_color = {"risk": TONE_RISK, "warn": TONE_WARN, "good": TONE_GOOD}.get(tone, NAVY)
            value = _clip_words(_safe(m.get("value")), 36)
            label = _clip_words(_safe(m.get("label")), 28)
            self.card(row_y, h=chip_h, x=cx, w=chip_w, fill=fill)
            box = self.slide.shapes.add_textbox(
                Emu(cx + 50_000), Emu(row_y + 70_000), Emu(chip_w - 100_000), Emu(chip_h - 140_000)
            )
            tf = box.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            r0 = p0.add_run()
            r0.text = value
            r0.font.name = FONT
            r0.font.bold = True
            r0.font.size = Pt(14 if len(value) > 12 else 16)
            r0.font.color.rgb = value_color
            p1 = tf.add_paragraph()
            p1.space_before = Pt(4)
            r1 = p1.add_run()
            r1.text = label
            r1.font.name = FONT
            r1.font.size = Pt(9)
            r1.font.color.rgb = MUTED_COLOR
        rows = (len(items) + cols - 1) // cols
        return y + rows * chip_h + max(0, rows - 1) * gap

    def bullets(self, items: list[str], y: int, color: RGBColor = BODY_COLOR, max_items: int = 8, max_chars: int = 320) -> int:
        dangling = re.compile(
            r"(?:\bв\s+т\.?\s*ч\.?|\bс\s+[А-ЯA-Z]\.?|\b[А-ЯA-Z]\.?|,|;|—|–|-)\s*$",
            re.I,
        )
        kept: list[str] = []
        for b in items[:max_items]:
            raw = _safe(b)
            if not raw:
                continue
            clipped = _clip_words(raw, max_chars)
            if dangling.search(clipped):
                # Prefer previous sentence boundary inside the clip window.
                punct = max(clipped.rfind(". "), clipped.rfind("! "), clipped.rfind("? "))
                if punct > 40:
                    clipped = clipped[: punct + 1].strip()
                else:
                    clipped = _trim_dangling_tail(clipped)
            if dangling.search(clipped) or clipped.endswith(("—", "–", "-")):
                raise RuntimeError(f"ORION bullet dangling on p{self.page}: {clipped[-48:]}")
            kept.append(clipped)
        if not kept:
            return y
        text = "\n".join(f"• {b}" for b in kept)
        needed = measure_text_height(text, CONTENT_W, FS_BODY, line_spacing=1.2, paragraph_spacing_pt=6)
        avail = max(300000, min(needed + 80_000, CONTENT_BOTTOM - y))
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(avail))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for bullet in kept:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(4)
            p.space_after = Pt(8)
            p.line_spacing = 1.15
            r = p.add_run()
            r.text = f"• {bullet}"
            r.font.name = FONT
            r.font.size = Pt(FS_BODY)
            r.font.color.rgb = color
        return y + min(avail, needed + 60_000)


def _asset_map(payload: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {str(a.get("assetRef")): a for a in payload.get("assets") or []}


def _resolve_image_bytes(asset: dict[str, Any] | None) -> bytes | None:
    """Load PNG/JPEG bytes from inline imageData or DATA_ROOT storageKey."""
    if not asset:
        return None
    img_data = asset.get("imageData")
    if img_data:
        try:
            raw = base64.b64decode(str(img_data))
            if len(raw) > 500:
                return raw
        except Exception:  # noqa: BLE001
            pass
    storage_key = str(asset.get("storageKey") or "").strip().lstrip("/")
    if storage_key:
        data_root = Path(os.environ.get("DATA_ROOT", "/data"))
        # storage keys are relative to digital-profile root; DATA_ROOT usually mounts that root.
        candidates = [
            data_root / storage_key,
            data_root / "digital-profile" / storage_key,
            Path(storage_key),
        ]
        for path in candidates:
            try:
                if path.is_file() and path.stat().st_size > 500:
                    return path.read_bytes()
            except OSError:
                continue
    return None


def _embed_image_contain(ctx: _Ctx, asset: dict[str, Any] | None, y: int, h: int = 4800000) -> bool:
    """Place image inside (MARGIN_X, y, CONTENT_W, h) preserving aspect ratio."""
    if not asset:
        ctx.body("Визуальный материал недоступен для данного раздела.", y)
        return False
    raw = _resolve_image_bytes(asset)
    if not raw:
        title = _safe(asset.get("title") or "Источник")
        domain = _safe(asset.get("caption") or asset.get("storageKey") or "")
        ctx.card(y, h)
        ctx.body(
            f"{title}\n{domain}\nИзображение недоступно — показаны источник и описание.",
            y + 120000,
            max_h=h - 200000,
        )
        return False
    box_w, box_h = CONTENT_W, h
    iw, ih = box_w, box_h
    if Image is not None:
        try:
            with Image.open(io.BytesIO(raw)) as im:
                iw, ih = im.size
        except Exception:  # noqa: BLE001
            pass
    scale = min(box_w / max(iw, 1), box_h / max(ih, 1))
    draw_w = int(iw * scale)
    draw_h = int(ih * scale)
    left = MARGIN_X + (box_w - draw_w) // 2
    top = y + (box_h - draw_h) // 2
    stream = io.BytesIO(raw)
    ctx.slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(draw_w), height=Emu(draw_h))
    return True


def _embed_image(ctx: _Ctx, asset: dict[str, Any] | None, y: int, h: int = 4800000) -> None:
    _embed_image_contain(ctx, asset, y, h)


def _first_visual_asset(refs: list[Any], assets: dict[str, dict[str, Any]]) -> dict[str, Any] | None:
    for ref in refs:
        asset = assets.get(str(ref))
        if asset and _resolve_image_bytes(asset):
            return asset
    return None


def _sidebar_word_budget(text: str, max_words: int = 70) -> str:
    """Keep complete sentences within a soft word budget; never emit one-word stubs."""
    raw = _safe(text)
    if not raw:
        return ""
    words = raw.split()
    if len(words) <= max_words:
        return raw
    sentences = re.split(r"(?<=[.!?…])\s+", raw)
    kept: list[str] = []
    count = 0
    for sent in sentences:
        w = len(sent.split())
        if kept and count + w > max_words:
            break
        if not kept and w > max_words:
            return _clip_words(sent, max(80, max_words * 6))
        kept.append(sent)
        count += w
        if count >= max_words:
            break
    return " ".join(kept).strip() or _clip_words(raw, max(80, max_words * 6))


def _qa_preview(text: str, match_index: int = 0) -> str:
    """Short, safe preview around a matched token for QA diagnostics."""
    raw = _safe(text)
    start = max(0, match_index - 20)
    snippet = raw[start:match_index + 40].strip()
    snippet = re.sub(r"\s+", " ", snippet)
    return (snippet[:60] + "…") if len(snippet) > 60 else snippet


def _sidebar_analysis(ctx: _Ctx, slide: dict[str, Any], x: int, y: int, w: int, h: int) -> None:
    """Unified client sidebar panel (v57): one column, no stacked framed cards."""
    analysis = slide.get("visualAnalysis") or {}
    if not isinstance(analysis, dict):
        analysis = {}
    qa_failures: list[str] = []

    def fail(msg: str) -> None:
        qa_failures.append(f"p{ctx.page}:{msg}")

    mode = str(analysis.get("sidebarMode") or "")
    headline = _safe(analysis.get("headlineConclusion") or slide.get("clientTakeaway") or "Вывод")
    meaning = _safe(analysis.get("clientMeaning") or analysis.get("whyItMatters") or "")
    visible = _safe(analysis.get("whatIsVisible") or "")
    explanations = analysis.get("highlightExplanations") or []
    if not isinstance(explanations, list):
        explanations = []
    actions = analysis.get("recommendedActions") or []
    action = _safe(actions[0]) if isinstance(actions, list) and actions else ""
    provenance = _safe(analysis.get("provenanceLabel") or "")
    more_n = int(analysis.get("moreSignalsCount") or 0)

    has_frames = any(
        isinstance(ex, dict) and str(ex.get("frameTone") or "") in {"red", "amber"} for ex in explanations
    )
    if mode == "adverse_explanation" or has_frames:
        mid_title = "Почему выделено"
        mid_bits = []
        for ex in explanations[:2]:
            if not isinstance(ex, dict):
                continue
            reason = _safe(ex.get("clientReason") or "")
            if reason:
                mid_bits.append(reason)
        if more_n > 0:
            mid_bits.append(f"Ещё {more_n} похожих сигналов.")
        mid_body = " ".join(mid_bits) if mid_bits else visible
    else:
        mid_title = "Что показывает экран"
        mid_body = visible or meaning

    # Hard client-safe bans in sidebar
    banned = re.compile(
        r"(\[DEMO\]|\.example\b|\bAPI\b|\bSUGGESTION\b|knowledge-строк|не\s+live|\bprovider\b|\bmanifest\b|\bsynthetic\b|\breconstruction\b|\bдвижок\b)",
        re.I,
    )
    field_names = {
        "headline": "headlineConclusion",
        "mid": "whatIsVisible",
        "meaning": "clientMeaning",
        "action": "recommendedActions",
    }
    for label, text in (("headline", headline), ("mid", mid_body), ("meaning", meaning), ("action", action)):
        field = field_names.get(label, label)
        if "…" in text or "..." in text:
            fail(f'sidebar ellipsis in {field}: "{_qa_preview(text)}"')
        m = banned.search(text)
        if m:
            # Diagnostic (not a bypass): surface field, matched token and a safe preview.
            fail(f'sidebar forbidden token "{m.group(0)}" in {field}: "{_qa_preview(text, m.start())}"')

    # Draw one outer panel
    pad = 70_000
    gap = 55_000
    cy = y + pad
    max_bottom = min(y + h, CONTENT_BOTTOM) - pad
    ctx.card(y, h=min(h, max_bottom - y + pad), x=x, w=w, fill=CARD_BG)

    def write_block(title: str | None, body: str, *, size: float = 11, bold_title: bool = True, required: bool = False) -> None:
        nonlocal cy
        if not body:
            return
        title_h = 200_000 if title else 0
        # Prefer complete text; do not ellipsis-clip sidebar
        fitted = body
        needed = measure_text_height(fitted, w - 2 * pad, size, line_spacing=1.2)
        avail = max_bottom - cy - 160_000 - title_h
        if needed > avail:
            # Keep complete sentences only (no ellipsis). When even the first
            # sentence cannot fit: the mandatory headline is a QA failure, but
            # optional trailing blocks are simply dropped whole (title included)
            # — a shorter sidebar beats failing the entire render.
            sentences = re.split(r"(?<=[.!?…])\s+", fitted)
            kept: list[str] = []
            for sent in sentences:
                trial = " ".join(kept + [sent]).strip()
                if measure_text_height(trial, w - 2 * pad, size, line_spacing=1.2) <= avail:
                    kept.append(sent)
                else:
                    break
            if not kept:
                if required:
                    fail("sidebar overflow without complete sentence")
                return
            fitted = " ".join(kept)
        if title:
            box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(w - 2 * pad), Emu(220_000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = title
            r.font.name = FONT
            r.font.bold = bold_title
            r.font.size = Pt(10.5)
            r.font.color.rgb = NAVY
            cy += 200_000
        bh = measure_text_height(fitted, w - 2 * pad, size, line_spacing=1.2)
        box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(cy), Emu(w - 2 * pad), Emu(max(bh, 120_000)))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = fitted
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = BODY_COLOR
        cy += bh + gap

    write_block(None, headline, size=12, required=True)
    write_block(mid_title, mid_body, size=11)
    if meaning and meaning != mid_body and meaning != headline:
        write_block("Что это значит", meaning, size=11)
    if action:
        write_block("Что сделать", action, size=11)
    if provenance:
        # Fine print, no frame
        if cy < max_bottom - 80_000:
            box = ctx.slide.shapes.add_textbox(Emu(x + pad), Emu(min(cy, max_bottom - 120_000)), Emu(w - 2 * pad), Emu(140_000))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = provenance
            r.font.name = FONT
            r.font.size = Pt(8.5)
            r.font.color.rgb = MUTED_COLOR

    if qa_failures:
        raise RuntimeError("ORION sidebar QA failed: " + "; ".join(qa_failures))



def _tone_fill(tone: str) -> RGBColor:
    return {"risk": RISK_BG, "warn": WARN_BG, "good": GOOD_BG, "accent": ACCENT_SOFT}.get(tone, CARD_BG)


def _tone_value_color(tone: str) -> RGBColor:
    return {"risk": TONE_RISK, "warn": TONE_WARN, "good": TONE_GOOD}.get(tone, NAVY)


def _render_kpi_cards(ctx: _Ctx, metrics: list[dict[str, Any]], x: int, y: int, width: int, cols: int = 2) -> int:
    items = [m for m in metrics if isinstance(m, dict) and _safe(m.get("value"))][:6]
    if not items:
        return y
    gap = 80_000
    card_w = (width - gap * (cols - 1)) // cols
    card_h = 780_000
    row_y = y
    for idx, m in enumerate(items):
        col = idx % cols
        if idx > 0 and col == 0:
            row_y += card_h + gap
        cx = x + col * (card_w + gap)
        tone = str(m.get("tone") or "neutral")
        # Keep room for Russian status phrases like «Данные не собраны» / «0 / 10».
        value = _clip_words(_safe(m.get("value")), 36)
        label = _clip_words(_safe(m.get("label")), 28)
        ctx.card(row_y, h=card_h, x=cx, w=card_w, fill=_tone_fill(tone))
        box = ctx.slide.shapes.add_textbox(
            Emu(cx + 70_000), Emu(row_y + 100_000), Emu(card_w - 140_000), Emu(card_h - 180_000)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = value
        r0.font.name = FONT
        r0.font.bold = True
        r0.font.size = Pt(18 if len(value) <= 10 else 12 if len(value) <= 22 else 10)
        r0.font.color.rgb = _tone_value_color(tone)
        p1 = tf.add_paragraph()
        p1.space_before = Pt(6)
        r1 = p1.add_run()
        r1.text = label
        r1.font.name = FONT
        r1.font.size = Pt(11)
        r1.font.color.rgb = MUTED_COLOR
    rows = (len(items) + cols - 1) // cols
    return y + rows * card_h + max(0, rows - 1) * gap



def _render_status_badge(ctx: _Ctx, badge: dict[str, Any] | None, x: int, y: int, width: int) -> int:
    if not isinstance(badge, dict) or not _safe(badge.get("label")):
        return y
    tone = str(badge.get("tone") or "neutral")
    h = 360_000
    ctx.card(y, h=h, x=x, w=width, fill=_tone_fill(tone))
    box = ctx.slide.shapes.add_textbox(Emu(x + 90_000), Emu(y + 90_000), Emu(width - 180_000), Emu(200_000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = _clip_words(_safe(badge.get("label")), 48)
    r.font.name = FONT
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = _tone_value_color(tone)
    return y + h


def _render_executive_dashboard(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title, 240000, NAVY, FS_SECTION)
    raw_narrative = str(slide.get("narrative") or "")
    paras = [p.strip() for p in re.split(r"\n+", raw_narrative) if p.strip()][:3]
    if not paras and raw_narrative:
        # Split long single paragraph into ~3 chunks on sentence boundaries.
        parts = re.split(r"(?<=[.!?])\s+", _safe(raw_narrative))
        paras = []
        buf = ""
        for part in parts:
            trial = f"{buf} {part}".strip() if buf else part
            if len(trial) > 280 and buf:
                paras.append(buf)
                buf = part
            else:
                buf = trial
            if len(paras) >= 3:
                break
        if buf and len(paras) < 3:
            paras.append(buf)
        paras = paras[:3] or [_clip_words(_safe(raw_narrative), 420)]
    left_w = int(CONTENT_W * 0.62)
    right_w = CONTENT_W - left_w - 120_000
    left_x = MARGIN_X
    right_x = MARGIN_X + left_w + 120_000
    cy = y
    for para in paras:
        cy = ctx.content_card(
            title=None,
            text=_clip_words(_safe(para), 420),
            x=left_x,
            y=cy,
            width=left_w,
            min_h=240_000,
            max_h=900_000,
            tone="neutral",
            body_size=11,
        )
        cy += 50_000
    metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
    _render_kpi_cards(ctx, metrics[:4], right_x, y, right_w, cols=2)
    bottom_y = max(cy, y + 1_600_000) + 40_000
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)][:2]
    actions = [a for a in (slide.get("actions") or []) if isinstance(a, dict)][:1]
    cards: list[tuple[str, str, str]] = []
    for finding in findings:
        tone = str(finding.get("tone") or "warn")
        detail = _safe(finding.get("detail") or "")
        headline = _safe(finding.get("headline") or "")
        # Prefer detail; avoid duplicating headline when detail already starts with it.
        if detail and headline and detail.lower().startswith(headline.lower()[:24].lower()):
            text = detail
        elif detail:
            text = detail
        else:
            text = headline
        # Guard dangling tails from upstream clips.
        text = _trim_dangling_tail(_safe(text))
        cards.append(("Риск", text, tone))
    if actions:
        act = actions[0]
        label = _safe(act.get("label"))
        text = _clip_words(label, 220)
        cards.append(("Следующий шаг", text, "accent"))
    if cards:
        gap = 80_000
        col_w = (CONTENT_W - gap * (len(cards) - 1)) // max(1, len(cards))
        fx = MARGIN_X
        card_max = min(1_500_000, max(520_000, CONTENT_BOTTOM - bottom_y - 40_000))
        for card_title, text, tone in cards:
            ctx.content_card(
                title=card_title,
                text=text,
                x=fx,
                y=bottom_y,
                width=col_w,
                min_h=420_000,
                max_h=card_max,
                tone=tone,
                title_size=11,
                body_size=11,
            )
            fx += col_w + gap


def _render_risk_matrix_grid(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title or "Матрица рисков", 240000, NAVY, FS_SECTION)
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)]
    if not findings:
        bullets = [_safe(b) for b in slide.get("bullets") or [] if _safe(b)]
        findings = [{"headline": b.split("—")[0].strip()[:60], "detail": b, "tone": "warn"} for b in bullets[:6]]
    badge_w = 1_900_000
    dangling = re.compile(
        r"(?:\bв\s+т\.?\s*ч\.?|\bс\s+[А-ЯA-Z]\.?|\b[А-ЯA-Z]\.?|,|;|—|–|-)\s*$",
        re.I,
    )
    for finding in findings[:6]:
        tone = str(finding.get("tone") or "warn")
        pill = _safe(finding.get("status") or finding.get("severity") or "")
        if len(pill) > 28:
            pill = _clip_words(pill, 28)
        headline = _safe(finding.get("headline") or "Тема")
        if len(headline) > 72:
            headline = _clip_words(headline, 72)
        detail = _safe(finding.get("detail") or "")
        # Prefer complete source text; only sentence-fit if height is tight.
        marker = _safe(finding.get("manualReview") or "")
        # Embed "requires review" into status instead of a cramped footer.
        if marker and re.search(r"требует", pill or "", re.I) is None and "проверк" in marker.lower():
            if pill and "проверк" not in pill.lower():
                pill = pill  # keep level; marker dropped from footer
        text_w = CONTENT_W - badge_w - 220_000 if pill else int(CONTENT_W * 0.92)
        left = MARGIN_X + 100_000
        pad_y = 80_000
        headline_h = measure_text_height(headline, text_w, 13, line_spacing=1.15)
        detail_budget = 2_200_000
        if detail:
            # Grow card to fit complete detail when possible.
            detail_h = measure_text_height(detail, text_w, 11, line_spacing=1.2)
            needed = pad_y + headline_h + 40_000 + detail_h + pad_y
            max_h = max(480_000, CONTENT_BOTTOM - y - 40_000)
            if needed > max_h:
                fitted = _fit_text_to_height(detail, text_w, 11, max(180_000, max_h - pad_y - headline_h - pad_y - 40_000))
                if dangling.search(fitted) or fitted != detail and not fitted.endswith((".", "!", "?")):
                    # Fall back to first complete sentence only.
                    sentences = re.split(r"(?<=[.!?])\s+", detail)
                    fitted = sentences[0].rstrip(".,;: ") + ("." if sentences and not sentences[0].endswith((".", "!", "?")) else "")
                    if dangling.search(fitted):
                        raise RuntimeError(f"ORION risk-matrix dangling detail on p{ctx.page}: {fitted[-40:]}")
                detail = fitted
                detail_h = measure_text_height(detail, text_w, 11, line_spacing=1.2)
            h = min(max_h, pad_y + headline_h + 40_000 + detail_h + pad_y)
        else:
            h = max(420_000, pad_y + headline_h + pad_y)
        h = max(420_000, min(h, CONTENT_BOTTOM - y - 40_000))
        ctx.card(y, h=h, fill=_tone_fill(tone))
        # Headline
        box = ctx.slide.shapes.add_textbox(Emu(left), Emu(y + pad_y), Emu(text_w), Emu(max(headline_h + 20_000, 160_000)))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = headline
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = NAVY
        text_y = y + pad_y + headline_h + 30_000
        if detail:
            rem = max(140_000, y + h - text_y - pad_y)
            box = ctx.slide.shapes.add_textbox(Emu(left), Emu(text_y), Emu(text_w), Emu(rem))
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = detail
            r.font.name = FONT
            r.font.size = Pt(11)
            r.font.color.rgb = BODY_COLOR
        if pill:
            bx = MARGIN_X + CONTENT_W - badge_w - 80_000
            by = y + pad_y
            bh = 280_000
            ctx.card(by, h=bh, x=bx, w=badge_w, fill=WHITE)
            b = ctx.slide.shapes.add_textbox(Emu(bx + 50_000), Emu(by + 70_000), Emu(badge_w - 100_000), Emu(160_000))
            btf = b.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = pill
            br.font.name = FONT
            br.font.bold = True
            br.font.size = Pt(11)
            br.font.color.rgb = _tone_value_color(tone)
        y += h + 50_000
        if y > CONTENT_BOTTOM - 360_000:
            break



def _render_profile_overview(ctx: _Ctx, slide: dict[str, Any], title: str) -> None:
    ctx.light_bg()
    y = ctx.title(title, 240000, NAVY, FS_SECTION)
    y = _render_status_badge(
        ctx, slide.get("statusBadge") if isinstance(slide.get("statusBadge"), dict) else None, MARGIN_X, y, CONTENT_W
    )
    y += 80_000
    metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
    ru = [m for m in metrics if "RU" in _safe(m.get("label")).upper() or "Росс" in _safe(m.get("label"))]
    uae = [m for m in metrics if "UAE" in _safe(m.get("label")).upper() or "ОАЭ" in _safe(m.get("label"))]
    other = [m for m in metrics if m not in ru and m not in uae]
    half = (CONTENT_W - 100_000) // 2
    if ru or uae:
        left_metrics = (ru or metrics[:4])[:4]
        right_metrics = (uae or metrics[4:8])[:4]
        hdr_h = 320_000
        for label, xx in (("Россия", MARGIN_X), ("ОАЭ", MARGIN_X + half + 100_000)):
            ctx.card(y, h=hdr_h, x=xx, w=half, fill=ACCENT_SOFT)
            box = ctx.slide.shapes.add_textbox(Emu(xx + 90_000), Emu(y + 90_000), Emu(half - 180_000), Emu(160_000))
            rr = box.text_frame.paragraphs[0].add_run()
            rr.text = label
            rr.font.name = FONT
            rr.font.bold = True
            rr.font.size = Pt(14)
            rr.font.color.rgb = NAVY
        y += hdr_h + 80_000
        left_bottom = _render_kpi_cards(ctx, left_metrics, MARGIN_X, y, half, cols=2)
        right_bottom = _render_kpi_cards(ctx, right_metrics, MARGIN_X + half + 100_000, y, half, cols=2)
        y = max(left_bottom, right_bottom) + 100_000
    else:
        y = _render_kpi_cards(ctx, metrics[:8], MARGIN_X, y, CONTENT_W, cols=4) + 100_000
    if other:
        y = _render_kpi_cards(ctx, other[:4], MARGIN_X, y, CONTENT_W, cols=4) + 80_000
    findings = [f for f in (slide.get("keyFindings") or []) if isinstance(f, dict)][:3]
    for finding in findings:
        if y > CONTENT_BOTTOM - 360_000:
            break
        y = ctx.content_card(
            title=_clip_words(_safe(finding.get("headline")), 48),
            text=_safe(finding.get("detail") or ""),
            x=MARGIN_X,
            y=y,
            width=CONTENT_W,
            min_h=280_000,
            max_h=700_000,
            tone=str(finding.get("tone") or "neutral"),
            title_size=11,
            body_size=11,
        )
        y += 60_000



def _render_visual_with_sidebar(
    ctx: _Ctx,
    slide: dict[str, Any],
    assets: dict[str, dict[str, Any]],
    title: str,
) -> None:
    """Title + left visual (contain) + right analytical sidebar."""
    ctx.light_bg()
    y = ctx.title(title, 280000, NAVY)
    refs = slide.get("assetRefs") or []
    visual = _first_visual_asset(refs, assets)
    has_sidebar = bool(slide.get("visualAnalysis") or slide.get("clientTakeaway") or slide.get("bullets"))
    img_w = int(CONTENT_W * 0.62) if has_sidebar else CONTENT_W
    side_w = CONTENT_W - img_w - 120000
    img_h = CONTENT_BOTTOM - y - 80000
    if visual:
        # Temporarily embed into a narrower box by using contain math inline
        raw = _resolve_image_bytes(visual)
        if raw:
            iw, ih = img_w, img_h
            if Image is not None:
                try:
                    with Image.open(io.BytesIO(raw)) as im:
                        iw, ih = im.size
                except Exception:  # noqa: BLE001
                    pass
            # Top-align visual with sidebar (do not vertically center in full column).
            scale = min(img_w / max(iw, 1), (img_h - 60000) / max(ih, 1))
            dw, dh = int(iw * scale), int(ih * scale)
            left = MARGIN_X + (img_w - dw) // 2
            top = y + 60000
            ctx.slide.shapes.add_picture(io.BytesIO(raw), Emu(left), Emu(top), width=Emu(dw), height=Emu(dh))
        else:
            ctx.body("Визуальный материал недоступен.", y + 80000, max_h=600000, color=MUTED_COLOR)
    else:
        reason = _safe(slide.get("blockedReason") or "Визуальный материал недоступен для данного раздела.")
        ctx.body(reason, y + 80000, max_h=800000, color=MUTED_COLOR)
    if has_sidebar and side_w > 400000:
        _sidebar_analysis(ctx, slide, MARGIN_X + img_w + 120000, y + 60000, side_w, img_h - 60000)


def _title_line_estimate(text: str, col_width_emu: int, font_pt: float, max_lines: int = 2) -> int:
    """Word-aware line estimate mirroring TS search-results-pagination.ts."""
    text = (text or "").strip()
    if not text:
        return 1
    char_w = font_pt * EMU_PER_PT * 0.52
    max_chars = max(1, int(col_width_emu / char_w))
    words = text.split()
    lines = 1
    cur = 0
    for w in words:
        add = len(w) if cur == 0 else cur + 1 + len(w)
        if add <= max_chars:
            cur = add
        else:
            lines += 1
            cur = min(len(w), max_chars)
            if lines >= max_lines:
                return max_lines
    return min(lines, max_lines)


def _status_tone(status: str) -> tuple[str, "RGBColor"]:
    s = (status or "").strip().lower()
    if "нежелат" in s:
        return "●", RGBColor(0xB9, 0x1C, 0x1C)
    if "проверк" in s or "требует" in s:
        return "●", RGBColor(0xC2, 0x41, 0x0C)
    return "●", RGBColor(0x04, 0x78, 0x57)


def _add_search_table(
    ctx: _Ctx,
    y: int,
    headers: list[str],
    rows: list[list[str]],
    groups: list[dict[str, Any]] | None = None,
) -> None:
    """
    Grouped SERP position table. Renders EVERY row the slide carries (no cap) —
    TS pagination already guaranteed geometric fit. Query is shown as a compact
    group-header band (spec §4), status as a colored badge (spec §5).
    """
    # Body layout is 4 cols: Позиция | Домен | Заголовок | Статус.
    # If TS sends a leading «Запрос» column, drop it — query lives in group bands.
    hdr = [str(h) for h in headers]
    data_rows = [list(r) for r in rows]
    if len(hdr) >= 5 and re.search(r"запрос|query", hdr[0], re.I):
        hdr = hdr[1:]
        data_rows = [r[1:] if len(r) > 1 else r for r in data_rows]
    cols = max(1, min(4, len(hdr)))
    headers = hdr
    groups = groups or []

    # Row plan: header + interleaved group bands + data rows.
    plan: list[tuple[str, Any]] = [("header", headers)]
    if groups:
        for g in groups:
            start = int(g.get("rowStart", 0))
            count = int(g.get("rowCount", 0))
            label = str(g.get("queryDisplay") or "")
            qtag = g.get("qTag")
            band = f"Запрос: {label}" if not qtag else f"{qtag} — {label}"
            plan.append(("group", band))
            for r in data_rows[start : start + count]:
                plan.append(("data", r))
    else:
        for r in data_rows:
            plan.append(("data", r))

    # Column widths (Позиция | Домен | Заголовок | Статус) — spec §4 proportions.
    # Two-column tables (Параметр | Значение) need a readable label column, and
    # a textual first column (e.g. «База данных») needs more than the numeric
    # position width.
    if cols == 2:
        prop = [0.24, 0.76]
    elif headers and len(str(headers[0]).strip()) > 3:
        prop = [0.14, 0.26, 0.42, 0.18][:cols]
    else:
        prop = [0.07, 0.22, 0.53, 0.18][:cols]
    widths = [max(500_000, int(CONTENT_W * p)) for p in prop]
    leftover = CONTENT_W - sum(widths)
    if leftover != 0 and widths:
        widths[2 if cols > 2 else len(widths) - 1] += leftover
    title_col_w = widths[2] if cols > 2 else widths[-1]

    # Per-row heights.
    body_pt = 10.0
    line_h = int(body_pt * EMU_PER_PT * 1.2)
    pad = int(6 * EMU_PER_PT)
    header_h = int(26 * EMU_PER_PT)
    group_h = int(18 * EMU_PER_PT)
    heights: list[int] = []
    for kind, payload in plan:
        if kind == "header":
            heights.append(header_h)
        elif kind == "group":
            heights.append(group_h)
        else:
            lines = _title_line_estimate(str(payload[2]) if len(payload) > 2 else "", title_col_w, body_pt)
            heights.append(lines * line_h + pad)

    table_rows = len(plan)
    table_h = sum(heights)
    shape = ctx.slide.shapes.add_table(table_rows, cols, Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(table_h))
    tbl = shape.table
    for i, w in enumerate(widths):
        tbl.columns[i].width = Emu(w)
    for i, h in enumerate(heights):
        tbl.rows[i].height = Emu(h)

    def paint(cell: Any, text: str, *, bold: bool = False, color: Any = BODY_COLOR, bg: Any = WHITE, size: float = 10.0, clip: bool = True) -> None:
        # Status badges ("● Нежелательный") are complete labels, not clipped
        # prose — the dangling-tail trimmer would strip the word after the dot.
        cell.text = _clip_words(text, 200) if clip else _safe(text)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
        cell.text_frame.word_wrap = True
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = bg

    for r_idx, (kind, payload) in enumerate(plan):
        if kind == "header":
            for c in range(cols):
                label = str(payload[c]) if c < len(payload) else ""
                paint(tbl.cell(r_idx, c), label, bold=True, color=WHITE, bg=NAVY, size=10.0)
        elif kind == "group":
            merged = tbl.cell(r_idx, 0)
            merged.merge(tbl.cell(r_idx, cols - 1))
            paint(merged, str(payload), bold=True, color=NAVY, bg=ACCENT_SOFT, size=10.0)
        else:
            row = payload
            status = str(row[cols - 1] if len(row) >= cols else "").strip()
            adverse = "нежелат" in status.lower()
            row_bg = RGBColor(0xFE, 0xF2, 0xF2) if adverse else WHITE
            for c in range(cols):
                val = str(row[c]) if c < len(row) else ""
                if c == cols - 1:
                    dot, tone = _status_tone(val)
                    paint(tbl.cell(r_idx, c), f"{dot} {val}", color=tone, bg=row_bg, size=9.5, clip=False)
                else:
                    paint(tbl.cell(r_idx, c), val, bg=row_bg, size=10.0)


def _render_slide(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    template = str(slide.get("template") or "")
    title = _safe(slide.get("title") or "ORION")
    narrative = _safe(slide.get("narrative") or "")
    bullets = [_safe(b) for b in slide.get("bullets") or [] if _safe(b)]
    refs = slide.get("assetRefs") or []
    primary = assets.get(str(refs[0])) if refs else None

    if template == "orion_golden_cover":
        ctx.dark_bg()
        y = ctx.title("ORION Digital Profile", 1800000, WHITE, 34)
        ctx.body(narrative or title, y, max_h=700000, color=RGBColor(0xBF, 0xDB, 0xFE))
        ctx.body("Клиентский аудит · предварительная оценка", y + 900000, max_h=400000, color=MUTED_COLOR)
        return

    if template == "orion_golden_toc":
        ctx.dark_bg()
        y = ctx.title("Содержание отчёта", 400000, WHITE, FS_SECTION)
        ctx.bullets(
            bullets or ["Резюме", "Россия", "ОАЭ", "Compliance", "LexisNexis", "Рекомендации"],
            y,
            color=WHITE,
            max_items=22,
            max_chars=120,
        )
        return

    if template == "orion_golden_executive_dashboard":
        _render_executive_dashboard(ctx, slide, title)
        return

    if template == "orion_golden_risk_matrix_grid":
        _render_risk_matrix_grid(ctx, slide, title)
        return

    if template == "orion_golden_profile_overview":
        _render_profile_overview(ctx, slide, title)
        return

    if template == "orion_golden_executive_card":
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
        if metrics:
            y = _render_kpi_cards(ctx, metrics[:4], MARGIN_X, y, CONTENT_W, cols=4) + 100_000
        narr = narrative.strip()
        if narr and not bullets:
            y = ctx.content_card(
                title=None,
                text=_clip_words(narr, 1800),
                x=MARGIN_X,
                y=y,
                width=CONTENT_W,
                min_h=400_000,
                max_h=min(3_200_000, CONTENT_BOTTOM - y - 100_000),
                tone="neutral",
                body_size=11,
            )
            return
        if narr:
            y = ctx.content_card(
                title=None,
                text=_clip_words(narr, 1200),
                x=MARGIN_X,
                y=y,
                width=CONTENT_W,
                min_h=320_000,
                max_h=min(2_000_000, CONTENT_BOTTOM - y - (1_200_000 if bullets else 100_000)),
                tone="neutral",
                body_size=11,
            )
            y += 100_000
        if bullets:
            ctx.bullets(bullets, y, max_items=7, max_chars=280)
        return

    if template == "orion_golden_risk_matrix":
        _render_risk_matrix_grid(ctx, slide, title or "Матрица рисков")
        return

    if template == "orion_golden_region_divider":
        ctx.dark_bg()
        ctx.title(title, 2800000, WHITE, 34)
        return

    if template == "orion_golden_metrics_dashboard":
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        badge = slide.get("statusBadge") if isinstance(slide.get("statusBadge"), dict) else None
        if badge:
            y = _render_status_badge(ctx, badge, MARGIN_X, y, CONTENT_W) + 80_000
        if narrative:
            y = ctx.content_card(
                title=None,
                text=_clip_words(narrative, 420),
                x=MARGIN_X,
                y=y,
                width=CONTENT_W,
                min_h=260_000,
                max_h=800_000,
                tone="neutral",
                body_size=11,
            )
            y += 80_000
        metrics = [m for m in (slide.get("metrics") or []) if isinstance(m, dict)]
        if metrics:
            y = _render_kpi_cards(ctx, metrics[:6], MARGIN_X, y, CONTENT_W, cols=3) + 80_000
        actions = [a for a in (slide.get("actions") or []) if isinstance(a, dict)]
        if actions:
            y = ctx.content_card(
                title="Действие",
                text=_clip_words(_safe(actions[0].get("label")), 160),
                x=MARGIN_X,
                y=y,
                width=CONTENT_W,
                min_h=260_000,
                max_h=600_000,
                tone="warn",
                title_size=11,
                body_size=11,
            )
            y += 60_000
        if bullets:
            ctx.bullets(bullets, y, max_items=5, max_chars=260)
        return

    if template == "orion_golden_serp_screenshot":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_image_grid":
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        if len(refs) == 1:
            primary_grid = assets.get(str(refs[0])) if refs else None
            if primary_grid and _resolve_image_bytes(primary_grid):
                _embed_image(ctx, primary_grid, y + 60000, h=5_200_000)
                cap = _safe(primary_grid.get("caption") or "")
                if cap:
                    ctx.body(cap, CONTENT_BOTTOM - 380000, max_h=320000, color=MUTED_COLOR)
                return
        cols = 3
        cell_w = CONTENT_W // 3 - 80_000
        cell_h = 1_600_000
        gap = 120000
        max_rows = max(1, int((CONTENT_BOTTOM - y + gap) // (cell_h + gap)))
        max_cells = max_rows * cols
        for idx, ref in enumerate(refs):
            if idx >= max_cells:
                break
            row = idx // cols
            col = idx % cols
            cx = MARGIN_X + col * (cell_w + gap)
            cy = y + row * (cell_h + gap)
            asset = assets.get(str(ref))
            raw = _resolve_image_bytes(asset) if asset else None
            if raw:
                stream = io.BytesIO(raw)
                iw, ih = cell_w, cell_h
                if Image is not None:
                    try:
                        with Image.open(io.BytesIO(raw)) as im:
                            iw, ih = im.size
                    except Exception:
                        pass
                scale = min(cell_w / max(iw, 1), cell_h / max(ih, 1))
                dw, dh = int(iw * scale), int(ih * scale)
                left = cx + (cell_w - dw) // 2
                top = cy + (cell_h - dh) // 2
                ctx.slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(dw), height=Emu(dh))
            else:
                shape = ctx.slide.shapes.add_shape(1, Emu(cx), Emu(cy), Emu(cell_w), Emu(cell_h))
                shape.fill.solid()
                shape.fill.fore_color.rgb = CARD_BG
                shape.line.color.rgb = CARD_BORDER
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = _safe((asset or {}).get("title") or "Недоступно")
                r.font.size = Pt(FS_CAPTION)
        return

    if template == "orion_golden_video_cards":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_knowledge_panel":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_surface_panel":
        _render_visual_with_sidebar(ctx, slide, assets, title)
        return

    if template == "orion_golden_lexis_visual_page":
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        _embed_image(ctx, primary, y + 60000, h=5_200_000)
        return

    if template == "orion_golden_compliance_visual_page":
        # Dow Jones / World-Check approved screenshots — same layout as Lexis visual.
        if slide.get("visualAnalysis") or slide.get("clientTakeaway"):
            _render_visual_with_sidebar(ctx, slide, assets, title)
            return
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY)
        _embed_image(ctx, primary, y + 60000, h=5_200_000)
        return

    if template == "orion_golden_search_table":
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        if narrative:
            # Keep 1–2 complete sentences above the table; never end on «как/и/с».
            intro = _safe(narrative)
            sentences = [s.strip() for s in re.split(r"(?<=[.!?…])\s+", intro) if s.strip()]
            complete = [
                s
                for s in sentences
                if s.endswith((".", "!", "?", "…"))
                and not re.search(r"(?:\bкак|\bи|\bс|\bв|\bпо|,|;|—)\s*$", s, re.I)
            ]
            if complete:
                intro = " ".join(complete[:2])
            else:
                intro = "Таблица фиксирует сохранённые позиции поисковой выдачи."
            intro = _trim_dangling_tail(intro)
            y = ctx.body(intro, y, max_h=900000, color=MUTED_COLOR)
            y = y + 40000
        table = slide.get("table") if isinstance(slide.get("table"), dict) else None
        headers = list((table or {}).get("headers") or [])
        rows = list((table or {}).get("rows") or [])
        groups = list((table or {}).get("groups") or [])
        if not rows and bullets:
            # Fallback: parse bullet lines into a compact table
            headers = ["Поз.", "Домен", "Заголовок", "Риск"]
            parsed: list[list[str]] = []
            for bullet in bullets[:10]:
                raw = _safe(bullet)
                m = re.match(
                    r"^(?:\[([Н·N.])\]\s*)?#?\s*(\d+)\s+([^\s—\-]+)\s*[—\-–]\s*(.+)$",
                    raw,
                )
                if m:
                    mark = "Н" if m.group(1) in ("Н", "N") else "·"
                    parsed.append([m.group(2), m.group(3), _clip_words(m.group(4), 70), mark])
                else:
                    parsed.append(["—", "—", _clip_words(raw, 80), "·"])
            rows = parsed
        if headers and rows:
            # Render every row the (paginated) slide carries — no hidden cap.
            # Keep up to 5 headers so Запрос can be stripped inside the helper
            # without also dropping Статус.
            _add_search_table(ctx, y, headers[:5], rows, groups)
        elif bullets:
            avail = max(400000, CONTENT_BOTTOM - y)
            box = ctx.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(avail))
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for bullet in bullets[:18]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_before = Pt(2)
                p.space_after = Pt(5)
                r = p.add_run()
                clipped = _clip_words(bullet, 160)
                r.text = f"• {clipped}"
                r.font.name = FONT
                r.font.size = Pt(11)
                r.font.color.rgb = (
                    RGBColor(0xB9, 0x1C, 0x1C) if clipped.startswith("[Н]") else BODY_COLOR
                )
        return

    if template == "orion_golden_no_data_compact":
        ctx.light_bg()
        y = ctx.title(title, 320000, NAVY)
        # Coverage empty states carry a full client explanation (what the
        # surface is, why it matters, recommendation) — allow multi-paragraph.
        ctx.body(
            narrative or "Для данного раздела недостаточно подтверждённых данных.",
            y,
            2600000,
        )
        return

    if template == "orion_golden_audit_dashboard":
        # ORION regional résumé: themes left-ish via bullets top, KPI counters below.
        ctx.light_bg()
        y = ctx.title(title, 280000, NAVY, FS_SECTION)
        if narrative:
            y = ctx.body(_clip_words(narrative, 520), y, max_h=1000000)
            y = y + 80000
        if bullets:
            ctx.bullets(bullets, y, max_items=14, max_chars=220)
        return

    # default section summary / appendix
    ctx.light_bg()
    y = ctx.title(title, 280000, NAVY, FS_SECTION)
    # Prefer bullets for dense content; keep narrative short to avoid overlap
    short_narrative = _clip_words(narrative, 480) if narrative else ""
    if short_narrative and not bullets:
        ctx.body(short_narrative, y, max_h=CONTENT_BOTTOM - y - 100000)
        return
    if short_narrative:
        y = ctx.body(short_narrative, y, max_h=900000)
        y = y + 80000
    if bullets:
        ctx.bullets(bullets, y, max_items=8)


def _pdf_cyrillic_fontfile() -> str | None:
    """Prefer a system font that can render Russian in fitz insert_textbox."""
    candidates = [
        Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts" / "arial.ttf",
        Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts" / "arialuni.ttf",
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
    ]
    for p in candidates:
        try:
            if p.is_file():
                return str(p)
        except OSError:
            continue
    return None


def _write_pdf_fallback(
    slides: list[dict[str, Any]],
    pdf_path: Path,
    subject: str,
    assets: dict[str, dict[str, Any]] | None = None,
) -> None:
    """Text+image PDF when LibreOffice is unavailable. Must embed SERP/Lexis imageData."""
    doc = fitz.open()
    asset_map = assets or {}
    # Match PPTX master 16:10 (do not invent an extra cover page — slides already include cover).
    all_slides = list(slides)
    total = len(all_slides)
    # 16:10 page geometry (px @ 100dpi of 12.8"×8.0")
    page_w, page_h = 1280, 800
    margin_x, title_bottom, content_bottom, footer_y = 48, 72, 740, 770
    visual_templates = {
        "orion_golden_serp_screenshot",
        "orion_golden_lexis_visual_page",
        "orion_golden_image_grid",
        "orion_golden_video_cards",
        "orion_golden_knowledge_panel",
        "orion_golden_surface_panel",
        "orion_golden_lexis_visual_page",
        "orion_golden_compliance_visual_page",
    }
    fontfile = _pdf_cyrillic_fontfile()
    cyr_font = "ArialCyr"

    def esc(t: str) -> str:
        return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    def textbox(page: fitz.Page, rect: fitz.Rect, text: str, *, fontsize: float, color: tuple[float, float, float]) -> None:
        if fontfile:
            try:
                page.insert_font(fontname=cyr_font, fontfile=fontfile)
                page.insert_textbox(rect, text, fontsize=fontsize, fontname=cyr_font, color=color)
                return
            except Exception:  # noqa: BLE001
                pass
        page.insert_textbox(rect, text, fontsize=fontsize, fontname="helv", color=color)

    for idx, slide in enumerate(all_slides, start=1):
        page = doc.new_page(width=page_w, height=page_h)
        title = _safe(slide.get("title") or "ORION")
        template = str(slide.get("template") or "")
        refs = slide.get("assetRefs") or []
        primary = None
        for ref in refs:
            cand = asset_map.get(str(ref))
            if cand and _resolve_image_bytes(cand):
                primary = cand
                break
        if primary is None and refs:
            primary = asset_map.get(str(refs[0]))
        img_bytes: bytes | None = None
        if template in visual_templates and primary:
            img_bytes = _resolve_image_bytes(primary)

        if img_bytes and len(img_bytes) > 500:
            # Title strip + embedded visual (same data PPTX path uses).
            textbox(
                page,
                fitz.Rect(margin_x, 28, page_w - margin_x, title_bottom),
                title,
                fontsize=18,
                color=(0.04, 0.10, 0.20),
            )
            analysis = slide.get("visualAnalysis") if isinstance(slide.get("visualAnalysis"), dict) else {}
            has_side = bool(analysis) or bool(slide.get("clientTakeaway"))
            img_right = int(page_w * 0.62) if has_side else (page_w - margin_x)
            try:
                page.insert_image(
                    fitz.Rect(margin_x, 80, img_right, content_bottom),
                    stream=img_bytes,
                    keep_proportion=True,
                )
            except Exception:  # noqa: BLE001
                textbox(
                    page,
                    fitz.Rect(margin_x, 100, img_right, 200),
                    "Визуальный материал недоступен для данного раздела.",
                    fontsize=12,
                    color=(0.2, 0.25, 0.33),
                )
            if has_side:
                side_bits = [
                    _safe(analysis.get("headlineConclusion") or slide.get("clientTakeaway") or ""),
                    _safe(analysis.get("whatIsVisible") or ""),
                    _safe(analysis.get("whyItMatters") or ""),
                    _safe((analysis.get("limitations") or [None])[0] or ""),
                    _safe(analysis.get("provenanceLabel") or ""),
                ]
                side_text = "\n\n".join([b for b in side_bits if b])
                textbox(
                    page,
                    fitz.Rect(img_right + 16, 80, page_w - margin_x, content_bottom),
                    side_text[:1200],
                    fontsize=10,
                    color=(0.2, 0.25, 0.33),
                )
            textbox(
                page,
                fitz.Rect(page_w - 180, footer_y, page_w - margin_x, footer_y + 20),
                f"{idx}/{total}",
                fontsize=10,
                color=(0.58, 0.64, 0.72),
            )
            continue

        body = esc(_safe(slide.get("body") or slide.get("narrative") or ""))
        bullets = slide.get("bullets") or []
        bullet_html = "".join(f"<li>{esc(_safe(b))}</li>" for b in bullets[:8])
        html = (
            "<div style='font-family:Arial,sans-serif;color:#0b1a33;padding:8px;'>"
            f"<h1 style='font-size:22px;margin:0;'>{esc(title)}</h1>"
            f"<p style='margin-top:12px;font-size:12px;color:#334155;'>{body}</p>"
            f"<ul style='margin-top:12px;font-size:11px;color:#334155;'>{bullet_html}</ul>"
            f"<p style='position:absolute;bottom:16px;right:24px;color:#94a3b8;font-size:10px;'>{idx}/{total}</p>"
            "</div>"
        )
        page.insert_htmlbox(fitz.Rect(margin_x, 40, page_w - margin_x, content_bottom), html)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(pdf_path))
    doc.close()


def _export_png_pages(pdf_path: Path) -> list[dict[str, Any]]:
    doc = fitz.open(str(pdf_path))
    pages: list[dict[str, Any]] = []
    try:
        for i in range(len(doc)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
            pages.append(
                {
                    "pageNumber": i + 1,
                    "width": pix.width,
                    "height": pix.height,
                    "contentBase64": base64.b64encode(pix.tobytes("png")).decode("ascii"),
                }
            )
    finally:
        doc.close()
    return pages


def render_orion_golden(payload: dict[str, Any]) -> dict[str, Any]:
    assert_render_font_family()
    reset_layout_telemetry()
    deck = payload.get("deckManifest") or {}
    report_spec = payload.get("reportSpec") or {}
    slides = list(deck.get("finalSlides") or [])
    if not slides:
        raise ValueError("deckManifest.finalSlides is empty")

    assets = _asset_map(payload)
    # Diagnostics for blank SERP slides (lengths only — never log base64).
    asset_diag = []
    for ref, asset in list(assets.items())[:20]:
        raw = _resolve_image_bytes(asset)
        asset_diag.append(
            {
                "assetRef": ref,
                "kind": asset.get("kind"),
                "hasImageData": bool(asset.get("imageData")),
                "imageDataChars": len(str(asset.get("imageData") or "")),
                "hasStorageKey": bool(asset.get("storageKey")),
                "resolvedBytes": len(raw) if raw else 0,
            }
        )
    print(
        "[orion-golden-render] assets",
        json.dumps(
            {
                "assetCount": len(assets),
                "serpSlides": sum(
                    1 for s in slides if str(s.get("template") or "") == "orion_golden_serp_screenshot"
                ),
                "sample": asset_diag[:8],
            },
            ensure_ascii=False,
        ),
        flush=True,
    )
    subject = (report_spec.get("subject") or {}).get("displayName") or "Цифровой профиль"
    total = len(slides)
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for idx, slide in enumerate(slides, start=1):
        ctx = _Ctx(prs, idx, total)
        _render_slide(ctx, slide, assets)
        ctx.footer()

    warnings: list[str] = []
    with tempfile.TemporaryDirectory(prefix="orion-golden-") as tmp:
        tmp_path = Path(tmp)
        pptx_path = tmp_path / "report.pptx"
        prs.save(str(pptx_path))
        pdf_path = tmp_path / "report.pdf"
        pdf_ok = False
        pdf_mode = "fitz-fallback"
        try:
            from convert_pdf import convert_to_pdf

            convert_to_pdf(str(pptx_path), str(pdf_path))
            pdf_ok = pdf_path.exists() and pdf_path.stat().st_size > 0
            if pdf_ok:
                pdf_mode = "libreoffice"
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"libreoffice-failed:{exc}")

        if not pdf_ok:
            _write_pdf_fallback(slides, pdf_path, str(subject), assets)
            pdf_mode = "fitz-fallback"

        print(f"[orion-golden-render] pdfExportMode={pdf_mode} warnings={warnings}", flush=True)
        pages = _export_png_pages(pdf_path)
        telemetry = get_layout_telemetry()
        return {
            "slideCount": len(prs.slides),
            "pptxBase64": base64.b64encode(pptx_path.read_bytes()).decode("ascii"),
            "pdfBase64": base64.b64encode(pdf_path.read_bytes()).decode("ascii") if pdf_path.exists() else "",
            "pages": pages,
            "pdfExportMode": pdf_mode,
            "warnings": warnings,
            "layoutTelemetry": {
                "version": "orion-layout-telemetry-v1",
                "entries": telemetry,
            },
        }


if __name__ == "__main__":
    import sys

    data = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    out = render_orion_golden(data)
    Path(sys.argv[2]).write_bytes(base64.b64decode(out["pptxBase64"]))
    if out.get("pdfBase64"):
        Path(sys.argv[3]).write_bytes(base64.b64decode(out["pdfBase64"]))
    pages_dir = Path(sys.argv[4])
    pages_dir.mkdir(parents=True, exist_ok=True)
    for page in out.get("pages") or []:
        Path(pages_dir / f"page-{page['pageNumber']:02d}.png").write_bytes(
            base64.b64decode(page["contentBase64"])
        )
    meta = {
        "slideCount": out["slideCount"],
        "pages": len(out.get("pages") or []),
        "pdfExportMode": out.get("pdfExportMode"),
    }
    Path(pages_dir.parent / "golden-render-meta.json").write_text(json.dumps(meta), encoding="utf-8")
    telemetry = out.get("layoutTelemetry")
    if telemetry:
        Path(pages_dir.parent / "layout-telemetry.json").write_text(
            json.dumps(telemetry, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    print(json.dumps(meta))
