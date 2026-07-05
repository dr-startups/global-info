"""Render ORION v2 manifest report JSON into PPTX, PDF and PNG page previews."""

from __future__ import annotations

import base64
import json
import os
import re
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

try:
    from render_pptx import build_pptx as _renderer_build_pptx
except Exception:  # pragma: no cover - optional import
    _renderer_build_pptx = None

ALLOW_BRANDS = {
    "orion",
    "google",
    "yandex",
    "dow jones",
    "world-check",
    "lexisnexis",
    "pep",
    "rca",
    "kyc",
}


def _safe_text(value: object) -> str:
    text = str(value or "").strip()
    text = re.sub(r"(c:\\\\|/mnt/|storage/digital-profile|https?://[^\s]+)", "", text, flags=re.I)
    text = re.sub(r"(openai[_-]?api[_-]?key|sk-[a-z0-9]{10,})", "", text, flags=re.I)
    return text.strip()


def _is_russian_report(report_json: dict[str, Any]) -> bool:
    lang = str((report_json.get("meta") or {}).get("language") or "").lower()
    return lang in {"ru", "russian", ""}


def _strip_english_leakage(text: str) -> str:
    if not text:
        return text
    out = text
    for token in re.findall(r"[A-Za-z][A-Za-z\-]{3,}", out):
        if token.lower() in ALLOW_BRANDS:
            continue
        out = out.replace(token, "")
    return re.sub(r"\s{2,}", " ", out).strip()


def _collect_slides(report_json: dict[str, Any]) -> list[dict[str, Any]]:
    manifest = (
        report_json.get("finalDeckManifest")
        or report_json.get("orionFinalDeckManifest")
        or {}
    )
    slides: list[dict[str, Any]] = [
        {"title": "Цифровой профиль", "subtitle": "ORION", "slideType": "cover_orion"},
        {"title": "Содержание", "subtitle": "Глобальная структура", "slideType": "toc_orion"},
    ]
    for section in manifest.get("sections") or []:
        for slide in section.get("slides") or []:
            slides.append(slide)
    return slides


def _write_pptx_fallback(report_json: dict[str, Any], pptx_path: Path) -> int:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    slides = _collect_slides(report_json)
    total = max(1, len(slides))
    ru_mode = _is_russian_report(report_json)

    for idx, src in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        title = _safe_text(src.get("title") or src.get("slideType") or f"Слайд {idx}")
        subtitle = _safe_text(src.get("subtitle") or "")
        if ru_mode:
            title = _strip_english_leakage(title)
            subtitle = _strip_english_leakage(subtitle)

        box = slide.shapes.add_textbox(Emu(500000), Emu(300000), Emu(8200000), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title or f"Слайд {idx}"
        r.font.bold = True
        r.font.size = Pt(26)
        r.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(13)
            r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        bullets_box = slide.shapes.add_textbox(Emu(500000), Emu(1400000), Emu(8200000), Emu(4300000))
        btf = bullets_box.text_frame
        btf.word_wrap = True
        bullets: list[str] = []
        for item in src.get("narrativeBlocks") or []:
            if isinstance(item, dict):
                text = _safe_text(item.get("text") or item.get("title") or "")
            else:
                text = _safe_text(item)
            if ru_mode:
                text = _strip_english_leakage(text)
            if text:
                bullets.append(text[:420])
        if not bullets:
            bullets = ["Раздел сформирован из структуры слайдов по этапам анализа."]
        for b_i, line in enumerate(bullets[:6]):
            pp = btf.paragraphs[0] if b_i == 0 else btf.add_paragraph()
            rr = pp.add_run()
            rr.text = f"• {line}"
            rr.font.size = Pt(14)

        footer = slide.shapes.add_textbox(Emu(500000), Emu(6200000), Emu(8200000), Emu(300000))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.add_run()
        fr.text = f"{idx}/{total}"
        fr.font.size = Pt(10)
        fr.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return total


def _write_pptx(report_json: dict[str, Any], pptx_path: Path, audience: str) -> int:
    if _renderer_build_pptx is not None:
        warnings, slide_count = _renderer_build_pptx(
            report_json,
            str(pptx_path),
            os.getcwd(),
            "report-template-v3",
            audience,
            "draft",
        )
        if int(slide_count or 0) > 0:
            return int(slide_count or 0)
        _ = warnings
        return _write_pptx_fallback(report_json, pptx_path)
    return _write_pptx_fallback(report_json, pptx_path)


def _escape_html(value: str) -> str:
    return (
        value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _slide_bullets(src: dict[str, Any], ru_mode: bool) -> list[str]:
    bullets: list[str] = []
    for item in src.get("narrativeBlocks") or []:
        if isinstance(item, dict):
            text = _safe_text(item.get("text") or item.get("title") or "")
        else:
            text = _safe_text(item)
        if ru_mode:
            text = _strip_english_leakage(text)
        if text:
            bullets.append(text)
    if not bullets:
        bullets = ["Раздел сформирован из структуры слайдов по этапам анализа."]
    return bullets[:8]


def _write_pdf(slide_count: int, report_json: dict[str, Any], pdf_path: Path) -> None:
    """Fallback PDF when LibreOffice conversion is unavailable (Unicode-safe HTML layout)."""
    slides = _collect_slides(report_json)
    doc = fitz.open()
    total = max(1, slide_count)
    ru_mode = _is_russian_report(report_json)
    for idx in range(total):
        page = doc.new_page(width=1280, height=720)
        src = slides[idx] if idx < len(slides) else {}
        title = _safe_text(src.get("title") or src.get("slideType") or f"Слайд {idx + 1}")
        subtitle = _safe_text(src.get("subtitle") or "")
        if ru_mode:
            title = _strip_english_leakage(title)
            subtitle = _strip_english_leakage(subtitle)
        bullets = _slide_bullets(src, ru_mode)
        bullet_html = "".join(f"<li>{_escape_html(line)}</li>" for line in bullets)
        subtitle_html = (
            f"<p style='margin:8px 0 0;color:#555;font-size:13px;'>{_escape_html(subtitle)}</p>"
            if subtitle
            else ""
        )
        html = (
            "<div style='font-family:sans-serif;color:#1f3a5f;'>"
            f"<h1 style='font-size:24px;margin:0;'>{_escape_html(title or f'Слайд {idx + 1}')}</h1>"
            f"{subtitle_html}"
            f"<ul style='margin-top:18px;font-size:14px;color:#222;line-height:1.45;'>{bullet_html}</ul>"
            f"<p style='position:absolute;bottom:12px;right:0;color:#777;font-size:10px;'>{idx + 1}/{total}</p>"
            "</div>"
        )
        page.insert_htmlbox(fitz.Rect(50, 40, 1230, 680), html)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(pdf_path))
    doc.close()


def _write_pdf_from_pptx(pptx_path: Path, pdf_path: Path) -> bool:
    try:
        from convert_pdf import convert_to_pdf

        convert_to_pdf(str(pptx_path), str(pdf_path))
        return pdf_path.exists() and pdf_path.stat().st_size > 0
    except Exception:
        return False


def _export_png_pages(pdf_path: Path) -> list[dict[str, Any]]:
    doc = fitz.open(str(pdf_path))
    pages: list[dict[str, Any]] = []
    try:
        for i in range(len(doc)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
            pages.append(
                {
                    "pageNumber": i + 1,
                    "width": pix.width,
                    "height": pix.height,
                    "contentBase64": base64.b64encode(pix.tobytes("png")).decode("ascii"),
                }
            )
    finally:
        doc.close()
    return pages


def render_orion_manifest(report_json: dict[str, Any], audience: str = "internal") -> dict[str, Any]:
    import tempfile

    with tempfile.TemporaryDirectory(prefix="orion-manifest-") as tmp:
        work = Path(tmp)
        pptx_path = work / "report.pptx"
        pdf_path = work / "report.pdf"
        slide_count = _write_pptx(report_json, pptx_path, audience)
        if not _write_pdf_from_pptx(pptx_path, pdf_path):
            _write_pdf(slide_count, report_json, pdf_path)
        pages = _export_png_pages(pdf_path)
        pptx_bytes = pptx_path.read_bytes()
        pdf_bytes = pdf_path.read_bytes()
    return {
        "slideCount": slide_count,
        "pptxBase64": base64.b64encode(pptx_bytes).decode("ascii"),
        "pdfBase64": base64.b64encode(pdf_bytes).decode("ascii"),
        "pages": pages,
    }
