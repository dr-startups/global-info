"""Render ORION ReportSpec v1 target sections — R9.7b visual fidelity renderer."""

from __future__ import annotations

import base64
import json
import re
import tempfile
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ORION type scale — single font family
FONT = "Arial"
FS_TITLE = 24
FS_SUBTITLE = 14
FS_BODY = 12
FS_CAPTION = 10
FS_METRIC_VALUE = 20
FS_METRIC_LABEL = 11

MARGIN_X = 450000
CONTENT_W = 8200000
SLIDE_H = 6858000
FOOTER_Y = 6400000

TITLE_COLOR = RGBColor(0x1F, 0x3A, 0x5F)
BODY_COLOR = RGBColor(0x33, 0x41, 0x55)
MUTED_COLOR = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


def _safe(text: object) -> str:
    val = re.sub(r"\s+", " ", str(text or "")).strip()
    val = re.sub(r"(storage/|C:\\\\|/mnt/|openai[_-]?api[_-]?key)", "", val, flags=re.I)
    return val


class _SlideCtx:
    def __init__(self, prs: Presentation, page_num: int, total: int):
        self.prs = prs
        self.page_num = page_num
        self.total = total
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        self.slide = prs.slides.add_slide(layout)

    def footer(self) -> None:
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(FOOTER_Y), Emu(CONTENT_W), Emu(250000))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{self.page_num} / {self.total}"
        r.font.name = FONT
        r.font.size = Pt(FS_CAPTION)
        r.font.color.rgb = MUTED_COLOR

    def title_block(self, title: str, subtitle: str = "") -> int:
        y = 280000
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(800000))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(title)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_TITLE)
        r.font.color.rgb = TITLE_COLOR
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = _safe(subtitle)
            r2.font.name = FONT
            r2.font.size = Pt(FS_SUBTITLE)
            r2.font.color.rgb = MUTED_COLOR
        return y + 850000

    def body(self, text: str, y: int, max_h: int = 900000) -> int:
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(max_h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(text)
        r.font.name = FONT
        r.font.size = Pt(FS_BODY)
        r.font.color.rgb = BODY_COLOR
        return y + max_h

    def bullets(self, items: list[str], y: int, max_items: int = 6) -> int:
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(2800000))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for bullet in items[:max_items]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = f"• {_safe(bullet)}"
            r.font.name = FONT
            r.font.size = Pt(FS_BODY)
            r.font.color.rgb = BODY_COLOR
        return y + 2800000

    def metrics_row(self, metrics: list[dict[str, Any]], y: int) -> int:
        card_w = 1900000
        gap = 180000
        for idx, metric in enumerate(metrics[:4]):
            cx = MARGIN_X + idx * (card_w + gap)
            shape = self.slide.shapes.add_shape(1, Emu(cx), Emu(y), Emu(card_w), Emu(700000))
            shape.fill.solid()
            shape.fill.fore_color.rgb = CARD_BG
            shape.line.color.rgb = CARD_BORDER
            tf = shape.text_frame
            tf.margin_left = Emu(80000)
            tf.margin_top = Emu(60000)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = _safe(metric.get("label"))
            r.font.name = FONT
            r.font.size = Pt(FS_METRIC_LABEL)
            r.font.color.rgb = MUTED_COLOR
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = _safe(metric.get("value"))
            r2.font.name = FONT
            r2.font.bold = True
            r2.font.size = Pt(FS_METRIC_VALUE)
            r2.font.color.rgb = TITLE_COLOR
        return y + 800000


def _asset_map(report_spec: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {str(a.get("assetRef")): a for a in report_spec.get("assets") or []}


def _collect_render_slides(report_spec: dict[str, Any]) -> list[tuple[dict[str, Any], dict[str, Any]]]:
    out: list[tuple[dict[str, Any], dict[str, Any]]] = []
    for section in report_spec.get("sections") or []:
        for slide in section.get("slides") or []:
            out.append((slide, section))
    return out


def _render_executive(ctx: _SlideCtx, slide: dict[str, Any], section: dict[str, Any]) -> None:
    narrative = section.get("clientNarrative") or {}
    y = ctx.title_block(slide.get("title") or "Executive Summary", slide.get("subtitle") or narrative.get("headline"))
    y = ctx.body(slide.get("narrative") or narrative.get("summary") or "", y, max_h=700000)
    y = ctx.metrics_row(section.get("metrics") or [], y + 100000)
    ctx.bullets(slide.get("bullets") or narrative.get("whatWasFound") or [], y + 100000)


def _render_section_summary(ctx: _SlideCtx, slide: dict[str, Any], section: dict[str, Any]) -> None:
    narrative = section.get("clientNarrative") or {}
    y = ctx.title_block(slide.get("title") or section.get("title"), slide.get("subtitle") or section.get("subtitle"))
    y = ctx.body(slide.get("narrative") or narrative.get("summary") or "", y, max_h=750000)
    y = ctx.metrics_row(section.get("metrics") or [], y + 80000)
    ctx.bullets(slide.get("bullets") or narrative.get("whatWasFound") or [], y + 80000)


def _render_serp(ctx: _SlideCtx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    y = ctx.title_block(slide.get("title") or "Поисковая выдача", slide.get("subtitle") or "")
    asset_ref = (slide.get("assetRefs") or [None])[0]
    asset = assets.get(str(asset_ref)) if asset_ref else None
    img_y = y + 80000
    if asset and asset.get("imageData"):
        img_path = Path(tempfile.gettempdir()) / f"orion-reportspec-{asset_ref}.png"
        img_path.write_bytes(base64.b64decode(str(asset.get("imageData"))))
        ctx.slide.shapes.add_picture(str(img_path), Emu(MARGIN_X), Emu(img_y), width=Emu(CONTENT_W), height=Emu(5000000))
        try:
            img_path.unlink(missing_ok=True)
        except OSError:
            pass
    else:
        ctx.body("Снимок поисковой выдачи недоступен для этого кейса.", img_y)


def _render_evidence(ctx: _SlideCtx, slide: dict[str, Any], section: dict[str, Any]) -> None:
    narrative = section.get("clientNarrative") or {}
    y = ctx.title_block(slide.get("title") or "Пояснение доказательной базы", "")
    y = ctx.body(slide.get("narrative") or narrative.get("whyItMatters") or "", y, max_h=650000)
    ctx.bullets(slide.get("bullets") or narrative.get("manualReviewQueue") or [], y + 80000)


def _write_pdf_fallback(report_spec: dict[str, Any], pdf_path: Path) -> None:
    doc = fitz.open()
    subject = report_spec.get("subject") or {}
    pairs = _collect_render_slides(report_spec)
    slides: list[tuple[str, str, list[str]]] = [
        ("ORION Digital Profile", str(subject.get("displayName") or ""), ["Executive Summary • Россия 2.1 • Россия 2.2"]),
    ]
    for slide, section in pairs:
        narrative = section.get("clientNarrative") or {}
        bullets = slide.get("bullets") or narrative.get("whatWasFound") or []
        slides.append(
            (
                _safe(slide.get("title") or section.get("title")),
                _safe(slide.get("subtitle") or ""),
                [_safe(b) for b in bullets if _safe(b)],
            )
        )
    total = max(1, len(slides))

    def esc(t: str) -> str:
        return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    for idx, (title, subtitle, bullets) in enumerate(slides, start=1):
        page = doc.new_page(width=1280, height=720)
        bullet_html = "".join(f"<li>{esc(line)}</li>" for line in bullets[:8])
        subtitle_html = f"<p style='margin:6px 0 0;color:#64748b;font-size:13px;'>{esc(subtitle)}</p>" if subtitle else ""
        html = (
            "<div style='font-family:Arial,sans-serif;color:#1f3a5f;padding:8px;'>"
            f"<h1 style='font-size:22px;margin:0;'>{esc(title)}</h1>"
            f"{subtitle_html}"
            f"<ul style='margin-top:16px;font-size:12px;color:#334155;line-height:1.5;'>{bullet_html}</ul>"
            f"<p style='position:absolute;bottom:16px;right:24px;color:#94a3b8;font-size:10px;'>{idx}/{total}</p>"
            "</div>"
        )
        page.insert_htmlbox(fitz.Rect(48, 40, 1232, 680), html)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(pdf_path))
    doc.close()


def _export_png_pages(pdf_path: Path) -> list[dict[str, Any]]:
    doc = fitz.open(str(pdf_path))
    pages: list[dict[str, Any]] = []
    try:
        for i in range(len(doc)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
            pages.append(
                {
                    "pageNumber": i + 1,
                    "width": pix.width,
                    "height": pix.height,
                    "contentBase64": base64.b64encode(pix.tobytes("png")).decode("ascii"),
                }
            )
    finally:
        doc.close()
    return pages


def render_report_spec(report_spec: dict[str, Any]) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(SLIDE_H)
    assets = _asset_map(report_spec)
    pairs = _collect_render_slides(report_spec)
    total = 1 + len(pairs)

    subject = report_spec.get("subject") or {}
    cover = _SlideCtx(prs, 1, total)
    cover.title_block("ORION Digital Profile", str(subject.get("displayName") or "Цифровой профиль"))
    cover.body("Executive Summary • Россия 2.1 • Россия 2.2", 1200000, max_h=400000)
    cover.footer()

    for idx, (slide, section) in enumerate(pairs, start=2):
        ctx = _SlideCtx(prs, idx, total)
        template = str(slide.get("template") or "")
        if template == "orion_executive_summary":
            _render_executive(ctx, slide, section)
        elif template == "orion_section_summary":
            _render_section_summary(ctx, slide, section)
        elif template == "orion_serp_screenshot":
            _render_serp(ctx, slide, assets)
        elif template == "orion_evidence_explanation":
            _render_evidence(ctx, slide, section)
        else:
            _render_section_summary(ctx, slide, section)
        ctx.footer()

    with tempfile.TemporaryDirectory(prefix="orion-reportspec-") as tmp:
        tmp_path = Path(tmp)
        pptx_path = tmp_path / "report.pptx"
        prs.save(str(pptx_path))
        pdf_path = tmp_path / "report.pdf"
        pdf_ok = False
        try:
            from convert_pdf import convert_to_pdf

            convert_to_pdf(str(pptx_path), str(pdf_path))
            pdf_ok = pdf_path.exists() and pdf_path.stat().st_size > 0
        except Exception:
            pdf_ok = False
        if not pdf_ok:
            _write_pdf_fallback(report_spec, pdf_path)
        pages = _export_png_pages(pdf_path)
        return {
            "slideCount": len(prs.slides),
            "pptxBase64": base64.b64encode(pptx_path.read_bytes()).decode("ascii"),
            "pdfBase64": base64.b64encode(pdf_path.read_bytes()).decode("ascii") if pdf_path.exists() else "",
            "pages": pages,
        }


if __name__ == "__main__":
    import sys

    data = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))
    out = render_report_spec(data)
    Path(sys.argv[2]).write_bytes(base64.b64decode(out["pptxBase64"]))
    if out.get("pdfBase64"):
        Path(sys.argv[3]).write_bytes(base64.b64decode(out["pdfBase64"]))
    pages_dir = Path(sys.argv[4])
    pages_dir.mkdir(parents=True, exist_ok=True)
    for page in out.get("pages") or []:
        Path(pages_dir / f"page-{page['pageNumber']:02d}.png").write_bytes(
            base64.b64decode(page["contentBase64"])
        )
    print(json.dumps({"slideCount": out["slideCount"], "pages": len(out.get("pages") or [])}))
