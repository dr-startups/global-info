"""ORION Client Storyboard Visual Composer — R9.9/R9.11 deterministic slide renderer."""

from __future__ import annotations

import base64
import io
import re
import tempfile
from pathlib import Path
from typing import Any

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None  # type: ignore

FONT = "Arial"
FS_TITLE = 24
FS_SECTION = 22
FS_SUBTITLE = 13
FS_BODY = 11
FS_CAPTION = 9
FS_TAKEAWAY = 12
FS_LABEL = 10
FS_METRIC_VALUE = 18
FS_METRIC_LABEL = 10
FS_BADGE = 10

MARGIN_X = 420000
CONTENT_W = 8300000
SLIDE_H = 6858000
FOOTER_Y = 6420000
HEADER_H = 180000
REPORT_TITLE = "ORION Digital Profile"

TITLE_COLOR = RGBColor(0x0F, 0x17, 0x2A)
BODY_COLOR = RGBColor(0x33, 0x41, 0x55)
MUTED_COLOR = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
WARN_BG = RGBColor(0xFF, 0xFB, 0xEB)
LOW_BG = RGBColor(0xEC, 0xFD, 0xF5)
MED_BG = RGBColor(0xFF, 0xF7, 0xED)
HIGH_BG = RGBColor(0xFE, 0xF2, 0xF2)

FORBIDDEN = re.compile(
    r"\b(PRESENT|UNKNOWN|adverse_media|pep|mock|fallback|provider|runtime|debug|cmr[a-z0-9]{10,}|executive_summary-rf-|ru_audit_summary-rf-|-sr-cmr)\b",
    re.I,
)


def _safe(text: object) -> str:
    val = re.sub(r"\s+", " ", str(text or "")).strip()
    val = FORBIDDEN.sub("", val)
    val = re.sub(r"(storage/|C:\\\\|/mnt/|openai[_-]?api[_-]?key)", "", val, flags=re.I)
    val = re.sub(r"\+\s*\d+\s*more items.*", "", val, flags=re.I)
    return val.strip()


def _truncate_at_word(text: str, max_len: int) -> str:
    safe = _safe(text)
    if len(safe) <= max_len:
        return safe
    clipped = safe[:max_len]
    last_space = clipped.rfind(" ")
    if last_space > int(max_len * 0.55):
        return clipped[:last_space].rstrip() + "…"
    return clipped.rstrip() + "…"


def _risk_palette(level: str) -> tuple[RGBColor, RGBColor, str]:
    lvl = (level or "unknown").lower()
    if lvl == "high":
        return HIGH_BG, RGBColor(0xB9, 0x1C, 0x1C), "Повышенное внимание"
    if lvl == "medium":
        return MED_BG, RGBColor(0xC2, 0x41, 0x0C), "Умеренный риск"
    if lvl == "low":
        return LOW_BG, RGBColor(0x04, 0x78, 0x57), "Низкий риск"
    return CARD_BG, MUTED_COLOR, "Требует проверки"


class _Ctx:
    def __init__(self, prs: Presentation, page: int, total: int):
        self.prs = prs
        self.page = page
        self.total = total
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        self.slide = prs.slides.add_slide(layout)

    def brand_band(self) -> None:
        band = self.slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(9144000), Emu(HEADER_H))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(40000), Emu(CONTENT_W), Emu(120000))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = REPORT_TITLE
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_CAPTION)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    def footer(self) -> None:
        line = self.slide.shapes.add_shape(1, Emu(MARGIN_X), Emu(FOOTER_Y - 60000), Emu(CONTENT_W), Emu(12000))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.fill.background()
        left = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(FOOTER_Y), Emu(4200000), Emu(220000))
        lp = left.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = REPORT_TITLE
        lr.font.name = FONT
        lr.font.size = Pt(FS_CAPTION)
        lr.font.color.rgb = MUTED_COLOR
        right = self.slide.shapes.add_textbox(Emu(MARGIN_X + CONTENT_W - 900000), Emu(FOOTER_Y), Emu(900000), Emu(220000))
        rp = right.text_frame.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        rr = rp.add_run()
        rr.text = f"{self.page} / {self.total}"
        rr.font.name = FONT
        rr.font.size = Pt(FS_CAPTION)
        rr.font.color.rgb = MUTED_COLOR

    def section_header(self, title: str, subtitle: str = "") -> int:
        y = HEADER_H + 120000
        accent = self.slide.shapes.add_shape(1, Emu(MARGIN_X), Emu(y), Emu(140000), Emu(520000))
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X + 180000), Emu(y), Emu(CONTENT_W - 180000), Emu(900000))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(title)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_SECTION)
        r.font.color.rgb = TITLE_COLOR
        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = _safe(subtitle)
            r2.font.name = FONT
            r2.font.size = Pt(FS_SUBTITLE)
            r2.font.color.rgb = MUTED_COLOR
        return y + 980000

    def badge(self, text: str, risk_level: str, x: int, y: int) -> None:
        bg, fg, label = _risk_palette(risk_level)
        shape = self.slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(2100000), Emu(340000))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = fg
        tf = shape.text_frame
        tf.margin_left = Emu(80000)
        tf.margin_top = Emu(50000)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(text or label)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_BADGE)
        r.font.color.rgb = fg

    def labeled_block(self, label: str, text: str, y: int, height: int = 620000) -> int:
        safe_text = _truncate_at_word(text, 280)
        line_estimate = max(1, len(safe_text) // 55 + 1)
        dynamic_h = min(max(height, 380000 + line_estimate * 90000), 900000)
        shape = self.slide.shapes.add_shape(1, Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(dynamic_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame
        tf.margin_left = Emu(100000)
        tf.margin_top = Emu(70000)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(label)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_LABEL)
        r.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = safe_text
        r2.font.name = FONT
        r2.font.size = Pt(FS_TAKEAWAY)
        r2.font.color.rgb = BODY_COLOR
        return y + dynamic_h + 100000

    def callout(self, title: str, text: str, y: int, tone: str = "info") -> int:
        bg = WARN_BG if tone == "warn" else CARD_BG
        shape = self.slide.shapes.add_shape(1, Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(520000))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = CARD_BORDER
        tf = shape.text_frame
        tf.margin_left = Emu(90000)
        tf.margin_top = Emu(60000)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(title)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_LABEL)
        r.font.color.rgb = TITLE_COLOR
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = _truncate_at_word(text, 260)
        r2.font.name = FONT
        r2.font.size = Pt(FS_BODY)
        r2.font.color.rgb = BODY_COLOR
        return y + 560000

    def card(self, title: str, body: str, x: int, y: int, w: int, h: int) -> None:
        shape = self.slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_BORDER
        tf = shape.text_frame
        tf.margin_left = Emu(70000)
        tf.margin_top = Emu(60000)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _safe(title)
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(FS_LABEL)
        r.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = _truncate_at_word(body, 220)
        r2.font.name = FONT
        r2.font.size = Pt(FS_BODY)
        r2.font.color.rgb = BODY_COLOR

    def metrics(self, metrics: list[dict[str, Any]], y: int) -> int:
        if not metrics:
            return y
        card_w = 1900000
        gap = 160000
        for idx, metric in enumerate(metrics[:4]):
            cx = MARGIN_X + idx * (card_w + gap)
            self.card(
                _safe(metric.get("label")),
                _safe(metric.get("value")),
                cx,
                y,
                card_w,
                720000,
            )
        return y + 820000

    def bullets(self, items: list[str], y: int, max_items: int = 5) -> int:
        cleaned = [_safe(x) for x in items if _safe(x)][:max_items]
        if not cleaned:
            return y
        box = self.slide.shapes.add_textbox(Emu(MARGIN_X), Emu(y), Emu(CONTENT_W), Emu(2200000))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for bullet in cleaned:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = f"• {_truncate_at_word(bullet, 180)}"
            r.font.name = FONT
            r.font.size = Pt(FS_BODY)
            r.font.color.rgb = BODY_COLOR
        return y + min(2200000, 420000 * len(cleaned) + 180000)

    def image_contain(
        self,
        img_data: str,
        x: int,
        y: int,
        box_w: int,
        box_h: int,
        asset_ref: str,
    ) -> bool:
        if not img_data:
            return False
        raw = base64.b64decode(str(img_data))
        tmp = Path(tempfile.gettempdir()) / f"orion-img-{asset_ref}.png"
        tmp.write_bytes(raw)
        iw, ih = box_w, box_h
        if Image is not None:
            try:
                with Image.open(io.BytesIO(raw)) as im:
                    iw, ih = im.size
            except Exception:  # noqa: BLE001
                pass
        scale = min(box_w / max(iw, 1), box_h / max(ih, 1))
        draw_w = int(iw * scale)
        draw_h = int(ih * scale)
        left = x + (box_w - draw_w) // 2
        top = y + (box_h - draw_h) // 2
        frame = self.slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(box_w), Emu(box_h))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        frame.line.color.rgb = CARD_BORDER
        self.slide.shapes.add_picture(str(tmp), Emu(left), Emu(top), width=Emu(draw_w), height=Emu(draw_h))
        try:
            tmp.unlink(missing_ok=True)
        except OSError:
            pass
        return True

    def unavailable_card(self, message: str, y: int) -> int:
        return self.callout("Материал недоступен", message, y, tone="warn")


def _assets_map(payload: dict[str, Any]) -> dict[str, dict[str, Any]]:
    return {str(a.get("assetRef")): a for a in payload.get("assets") or []}


def _first_asset(slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> dict[str, Any] | None:
    refs = slide.get("assetRefs") or []
    asset_ref = refs[0].get("assetRef") if refs and isinstance(refs[0], dict) else None
    if not asset_ref and refs:
        asset_ref = str(refs[0])
    return assets.get(str(asset_ref)) if asset_ref else None


def render_cover(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "ORION Digital Profile", slide.get("subtitle") or "")
    ctx.badge("Статус проверки", str(slide.get("riskLevel") or "unknown"), MARGIN_X + CONTENT_W - 2200000, HEADER_H + 140000)
    ctx.labeled_block("Главный вывод", slide.get("clientTakeaway") or "", y + 20000, height=780000)


def render_global_toc(ctx: _Ctx, slide: dict[str, Any], all_slides: list[dict[str, Any]]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Содержание", "Структура отчёта")
    items = []
    for s in all_slides:
        t = _safe(s.get("title"))
        if t and s.get("slideType") not in ("cover", "global_toc"):
            items.append(t)
    ctx.bullets(items[:8], y, max_items=8)


def render_executive_summary(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Executive Summary", slide.get("subtitle") or "")
    ctx.badge("Уровень риска", str(slide.get("riskLevel") or "unknown"), MARGIN_X + CONTENT_W - 2200000, HEADER_H + 140000)
    y = ctx.labeled_block("Главный вывод", slide.get("clientTakeaway") or "", y, height=480000)
    y = ctx.metrics(slide.get("metrics") or [], y)
    actions = [_safe(a.get("label")) for a in (slide.get("recommendedActions") or []) if isinstance(a, dict)][:2]
    if actions:
        ctx.callout("Следующий шаг", actions[0], y, tone="info")


def render_region_summary(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Сводка региона", slide.get("subtitle") or "")
    y = ctx.labeled_block("Краткий вывод", slide.get("clientTakeaway") or "", y, height=480000)
    findings = [_safe(f.get("summary")) for f in (slide.get("findings") or []) if isinstance(f, dict)][:3]
    col_w = (CONTENT_W - 240000) // 3
    gap = 120000
    card_y = y + 40000
    if card_y + 1050000 > FOOTER_Y - 200000:
        card_y = y
    ctx.card("Что найдено", findings[0] if findings else "Подтверждённые материалы ограничены.", MARGIN_X, card_y, col_w, 980000)
    ctx.card(
        "Почему это важно",
        findings[1] if len(findings) > 1 else "Сигналы влияют на общую картину риска по региону.",
        MARGIN_X + col_w + gap,
        card_y,
        col_w,
        980000,
    )
    ctx.card(
        "Что требует проверки",
        findings[2] if len(findings) > 2 else "Ручная верификация спорных совпадений.",
        MARGIN_X + 2 * (col_w + gap),
        card_y,
        col_w,
        980000,
    )


def render_search_overview(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Поисковая выдача", slide.get("subtitle") or "")
    y = ctx.labeled_block("Обзор", slide.get("clientTakeaway") or "", y, height=520000)
    y = ctx.metrics(slide.get("metrics") or [], y)
    y = ctx.callout(
        "Риск и релевантность",
        "Отдельные результаты могут быть косвенными; каждый существенный сигнал требует ручной проверки связи с субъектом.",
        y,
        tone="warn",
    )
    evidence = slide.get("evidenceRefs") or []
    ctx.bullets(
        [_safe(f"{e.get('label')}: {e.get('statusLabel')}") for e in evidence if isinstance(e, dict)],
        y + 80000,
        max_items=5,
    )


def render_serp_screenshot(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Снимок выдачи", slide.get("subtitle") or "")
    y = ctx.callout(
        "Релевантность и риск",
        _safe(slide.get("clientTakeaway") or "Визуализация поисковой выдачи для аналитической проверки."),
        y,
        tone="warn",
    )
    asset = _first_asset(slide, assets)
    img_y = y + 70000
    box_h = 4300000
    if asset and asset.get("imageData"):
        ctx.image_contain(str(asset.get("imageData")), MARGIN_X, img_y, CONTENT_W, box_h, str(asset.get("assetRef")))
    else:
        ctx.unavailable_card("Данные снимка поисковой выдачи не обнаружены.", img_y)


def render_lexis_summary(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "LexisNexis — сводка", slide.get("subtitle") or "Compliance")
    ctx.badge("Статус импорта", str(slide.get("riskLevel") or "medium"), MARGIN_X + CONTENT_W - 2200000, HEADER_H + 140000)
    y = ctx.metrics(slide.get("metrics") or [], y)
    y = ctx.labeled_block("Аналитическая сводка", slide.get("clientTakeaway") or "", y)
    ctx.callout(
        "Ручная проверка",
        "Материалы LexisNexis не являются юридическим заключением и требуют аналитического подтверждения перед использованием в решении.",
        y + 80000,
        tone="warn",
    )


def render_lexis_visual_page(
    ctx: _Ctx,
    slide: dict[str, Any],
    assets: dict[str, dict[str, Any]],
    page_idx: int,
    page_total: int,
) -> None:
    ctx.brand_band()
    subtitle = f"Приложение — страница {page_idx} из {page_total}"
    y = ctx.section_header(slide.get("title") or "LexisNexis", subtitle)
    ctx.callout(
        "Примечание",
        _safe(
            slide.get("clientTakeaway")
            or "Оригинальная страница включена как визуальное подтверждение импорта; детальная проверка выполняется по исходному DOCX/PDF."
        ),
        y,
        tone="info",
    )
    asset = _first_asset(slide, assets)
    img_y = y + 620000
    box_h = min(5200000, FOOTER_Y - img_y - 120000)
    if asset and asset.get("imageData"):
        ctx.image_contain(str(asset.get("imageData")), MARGIN_X, img_y, CONTENT_W, box_h, str(asset.get("assetRef")))
    else:
        ctx.unavailable_card("Визуальная страница недоступна.", img_y)


def render_search_results_table(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Результаты поиска", "Ключевые источники")
    evidence = slide.get("evidenceRefs") or []
    rows = []
    for e in evidence[:5]:
        if not isinstance(e, dict):
            continue
        rows.append(f"{_safe(e.get('label'))} — {_safe(e.get('statusLabel'))}")
    ctx.bullets(rows, y, max_items=5)


def render_adverse_media_summary(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Негативные публикации", "")
    y = ctx.labeled_block("Интерпретация", slide.get("clientTakeaway") or "", y, height=480000)
    findings = slide.get("findings") or []
    ctx.bullets([_truncate_at_word(_safe(f.get("summary")), 200) for f in findings if isinstance(f, dict)], y + 60000, max_items=4)


def render_scope_overview(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Что проверялось", "")
    y = ctx.labeled_block("Область проверки", slide.get("clientTakeaway") or "", y, height=440000)
    findings = [_safe(f.get("summary")) for f in (slide.get("findings") or []) if isinstance(f, dict)]
    ctx.bullets(findings, y, max_items=5)


def render_risk_conclusion(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Главные выводы", "")
    ctx.badge("Уровень риска", str(slide.get("riskLevel") or "unknown"), MARGIN_X + CONTENT_W - 2200000, HEADER_H + 140000)
    y = ctx.labeled_block("Интерпретация риска", slide.get("clientTakeaway") or "", y, height=480000)
    findings = [_safe(f.get("summary")) for f in (slide.get("findings") or []) if isinstance(f, dict)][:3]
    ctx.bullets(findings, y, max_items=3)
    actions = [_safe(a.get("label")) for a in (slide.get("recommendedActions") or []) if isinstance(a, dict)][:2]
    if actions:
        ctx.callout("Что проверить дальше", actions[0], y + 1400000, tone="warn")


def render_relevant_sources(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Релевантные источники", slide.get("subtitle") or "")
    y = ctx.labeled_block("Сводка", slide.get("clientTakeaway") or "", y, height=420000)
    evidence = slide.get("evidenceRefs") or []
    rows = []
    for e in evidence[:5]:
        if not isinstance(e, dict):
            continue
        rows.append(f"{_safe(e.get('label'))}: {_truncate_at_word(_safe(e.get('summary')), 120)}")
    ctx.bullets(rows, y, max_items=5)


def render_excluded_matches(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Исключённые совпадения", "")
    y = ctx.callout("Пояснение", _safe(slide.get("clientTakeaway") or ""), y, tone="info")
    findings = [_safe(f.get("summary")) for f in (slide.get("findings") or []) if isinstance(f, dict)]
    ctx.bullets(findings, y + 80000, max_items=5)


def render_lexis_signals(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "LexisNexis — сигналы", "")
    y = ctx.labeled_block("Ключевые сигналы", slide.get("clientTakeaway") or "", y, height=420000)
    findings = slide.get("findings") or []
    col_w = (CONTENT_W - 120000) // 2
    for idx, f in enumerate(findings[:4]):
        if not isinstance(f, dict):
            continue
        col = idx % 2
        row = idx // 2
        cx = MARGIN_X + col * (col_w + 120000)
        cy = y + row * 950000
        ctx.card(_safe(f.get("headline")), _truncate_at_word(_safe(f.get("summary")), 180), cx, cy, col_w, 880000)


def render_image_grid(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Изображения", "")
    y = ctx.labeled_block("Контекст", slide.get("clientTakeaway") or "", y, height=480000)
    refs = slide.get("assetRefs") or []
    placed = 0
    for idx, ref in enumerate(refs[:4]):
        asset_ref = ref.get("assetRef") if isinstance(ref, dict) else str(ref)
        asset = assets.get(str(asset_ref))
        if not asset or not asset.get("imageData"):
            continue
        col = idx % 2
        row = idx // 2
        cx = MARGIN_X + col * (CONTENT_W // 2 + 80000)
        cy = y + row * 2300000
        ctx.image_contain(str(asset.get("imageData")), cx, cy, CONTENT_W // 2 - 100000, 2100000, str(asset_ref))
        placed += 1
    if placed == 0:
        ctx.unavailable_card("Изображения по данному региону не обнаружены.", y + 80000)


def render_video_cards(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Видеоматериалы", "")
    asset = _first_asset(slide, assets)
    if asset and asset.get("imageData"):
        ctx.image_contain(str(asset.get("imageData")), MARGIN_X, y + 80000, CONTENT_W, 4200000, str(asset.get("assetRef")))
    else:
        ctx.unavailable_card("Видеоматериалы не обнаружены или недоступны для предпросмотра.", y + 80000)


def render_knowledge_panel(ctx: _Ctx, slide: dict[str, Any], assets: dict[str, dict[str, Any]]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Справочная карточка", "")
    asset = _first_asset(slide, assets)
    if asset and asset.get("imageData"):
        ctx.image_contain(str(asset.get("imageData")), MARGIN_X, y + 80000, CONTENT_W // 2 + 2000000, 3800000, str(asset.get("assetRef")))
    else:
        ctx.labeled_block("Справочные данные", slide.get("clientTakeaway") or "Справочные данные ограничены.", y)


def render_recommended_actions(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Рекомендуемые действия", "Практические шаги для клиента")
    y = ctx.labeled_block("Контекст", slide.get("clientTakeaway") or "", y, height=480000)
    actions = slide.get("recommendedActions") or []
    labels = [_safe(a.get("label")) for a in actions if isinstance(a, dict) and _safe(a.get("label"))][:5]
    defaults = [
        "Проверить совпадения идентификационных данных субъекта",
        "Просмотреть отмеченные источники и подтвердить связь",
        "Подтвердить совпадения LexisNexis вручную",
        "Зафиксировать клиентский вывод после проверки",
    ]
    merged: list[str] = []
    for item in labels + defaults:
        if item and item not in merged:
            merged.append(item)
        if len(merged) >= 5:
            break
    col_w = (CONTENT_W - 120000) // 2
    for idx, label in enumerate(merged[:5]):
        col = idx % 2
        row = idx // 2
        cx = MARGIN_X + col * (col_w + 120000)
        cy = y + row * 900000
        ctx.card(f"Шаг {idx + 1}", label, cx, cy, col_w, 820000)


def render_no_data_compact(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Недостаточно данных", "")
    ctx.unavailable_card(slide.get("clientTakeaway") or "Данные не обнаружены / не применимо", y)


def render_generic(ctx: _Ctx, slide: dict[str, Any]) -> None:
    ctx.brand_band()
    y = ctx.section_header(slide.get("title") or "Раздел", slide.get("subtitle") or "")
    ctx.labeled_block("Сводка", slide.get("clientTakeaway") or "", y)


def render_client_storyboard(payload: dict[str, Any]) -> dict[str, Any]:
    storyboard = payload.get("storyboard") or payload
    slides = list(storyboard.get("slides") or [])
    assets = _assets_map(payload)
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(SLIDE_H)
    total = max(len(slides), 1)

    lexis_visual_slides = [s for s in slides if str(s.get("slideType")) == "lexisnexis_visual_page"]
    lexis_total = len(lexis_visual_slides)
    lexis_idx = 0

    for idx, slide in enumerate(slides, start=1):
        stype = str(slide.get("slideType") or "region_summary")
        ctx = _Ctx(prs, idx, total)
        if stype == "cover":
            render_cover(ctx, slide)
        elif stype == "global_toc":
            render_global_toc(ctx, slide, slides)
        elif stype == "executive_summary":
            render_executive_summary(ctx, slide)
        elif stype == "scope_overview":
            render_scope_overview(ctx, slide)
        elif stype == "risk_conclusion":
            render_risk_conclusion(ctx, slide)
        elif stype == "region_summary":
            render_region_summary(ctx, slide)
        elif stype == "search_overview":
            render_search_overview(ctx, slide)
        elif stype == "relevant_sources":
            render_relevant_sources(ctx, slide)
        elif stype == "excluded_matches":
            render_excluded_matches(ctx, slide)
        elif stype == "serp_screenshot":
            render_serp_screenshot(ctx, slide, assets)
        elif stype == "search_results_table":
            render_search_results_table(ctx, slide)
        elif stype == "adverse_media_summary":
            render_adverse_media_summary(ctx, slide)
        elif stype == "image_grid":
            render_image_grid(ctx, slide, assets)
        elif stype == "video_cards":
            render_video_cards(ctx, slide, assets)
        elif stype == "knowledge_panel":
            render_knowledge_panel(ctx, slide, assets)
        elif stype == "lexisnexis_summary":
            render_lexis_summary(ctx, slide)
        elif stype == "lexisnexis_signals":
            render_lexis_signals(ctx, slide)
        elif stype == "lexisnexis_visual_page":
            lexis_idx += 1
            render_lexis_visual_page(ctx, slide, assets, lexis_idx, lexis_total)
        elif stype == "recommended_actions":
            render_recommended_actions(ctx, slide)
        elif stype == "no_data_compact":
            render_no_data_compact(ctx, slide)
        else:
            render_generic(ctx, slide)
        ctx.footer()

    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "storyboard.pptx"
        pdf_path = Path(tmp) / "storyboard.pdf"
        prs.save(str(pptx_path))
        pdf_mode = "libreoffice"
        try:
            from convert_pdf import convert_to_pdf

            convert_to_pdf(str(pptx_path), str(pdf_path))
        except Exception:
            pdf_mode = "fitz-fallback"
            doc = fitz.open()
            for slide in slides:
                page = doc.new_page(width=595, height=842)
                page.insert_text((50, 60), _safe(slide.get("title")), fontsize=14)
                page.insert_text((50, 90), _safe(slide.get("clientTakeaway"))[:500], fontsize=10)
            doc.save(str(pdf_path))
            doc.close()

        pptx_b64 = base64.b64encode(pptx_path.read_bytes()).decode("ascii")
        pdf_b64 = base64.b64encode(pdf_path.read_bytes()).decode("ascii")
        pages = []
        doc = fitz.open(str(pdf_path))
        for i in range(len(doc)):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
            pages.append(
                {
                    "pageNumber": i + 1,
                    "contentBase64": base64.b64encode(pix.tobytes("png")).decode("ascii"),
                    "width": pix.width,
                    "height": pix.height,
                }
            )
        doc.close()

    return {
        "slideCount": len(slides),
        "pptxBase64": pptx_b64,
        "pdfBase64": pdf_b64,
        "pages": pages,
        "pdfExportMode": pdf_mode,
        "warnings": [],
    }
