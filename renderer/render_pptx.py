"""Build a PPTX from a Digital Profile `report_json`.

Approach:
  - If `template.pptx` exists next to this file, it is used as the base so a
    designer can control master slides / branding. Otherwise a clean default
    presentation is created.
  - Dynamic pages (person-specific) are rendered first, then static commercial
    pages. Each page maps to one slide.
  - A faint "DRAFT" watermark is stamped on every slide while the report is not
    final (controlled by report_json.meta.watermark).

This module never invents facts; it only lays out the data passed in report_json.
"""

from __future__ import annotations

import os
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from report_template_v1 import build_report_v1
from report_template_v2 import build_report_v2
from report_template_v3 import build_report_v3
from report_i18n import normalize_lang, watermark_text

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")
TEMPLATE_V1_PATH = os.path.join(
    os.path.dirname(__file__), "templates", "report-template-v1.pptx"
)
TEMPLATE_V2_PATH = os.path.join(
    os.path.dirname(__file__), "templates", "report-template-v2.pptx"
)
TEMPLATE_V3_PATH = os.path.join(
    os.path.dirname(__file__), "templates", "report-template-v3.pptx"
)
DEFAULT_TEMPLATE_VERSION = "report-template-v3"

SLIDE_W = Emu(9144000)  # 10 in
SLIDE_H = Emu(6858000)  # 7.5 in

ACCENT = RGBColor(0x1F, 0x3A, 0x5F)
MUTED = RGBColor(0x66, 0x66, 0x66)
WATERMARK_COLOR = RGBColor(0xD9, 0xD9, 0xD9)


def _new_presentation(base_path: str | None = None) -> Presentation:
    path = base_path if base_path and os.path.exists(base_path) else TEMPLATE_PATH
    if os.path.exists(path):
        prs = Presentation(path)
        # Reuse the base only for its master/branding; drop any starter slides.
        xml_slides = prs.slides._sldIdLst  # noqa: SLF001
        for sld in list(xml_slides):
            xml_slides.remove(sld)
        return prs
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_layout(prs: Presentation):
    # Prefer a truly blank layout (usually index 6 in the default template).
    layouts = prs.slide_layouts
    return layouts[6] if len(layouts) > 6 else layouts[-1]


def _add_title(slide, title: str, subtitle: str | None = None) -> Emu:
    box = slide.shapes.add_textbox(Emu(457200), Emu(304800), Emu(8229600), Emu(914400))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title or ""
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    if subtitle:
        sp = tf.add_paragraph()
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(14)
        srun.font.color.rgb = MUTED
    return Emu(1371600)  # y offset where content can start


def _add_bullets(slide, top: Emu, lines: list[str]) -> Emu:
    if not lines:
        return top
    box = slide.shapes.add_textbox(Emu(457200), top, Emu(8229600), Emu(3200400))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022 {line}"
        run.font.size = Pt(14)
    return Emu(int(top) + 457200 * max(1, len(lines)))


def _add_table(slide, top: Emu, columns: list[str], rows: list[list[Any]]):
    n_rows = len(rows) + 1
    n_cols = max(1, len(columns))
    width = Emu(8229600)
    height = Emu(min(457200 * n_rows, 4114800))
    graphic = slide.shapes.add_table(n_rows, n_cols, Emu(457200), top, width, height)
    table = graphic.table
    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = "" if val is None else str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)


def _add_images(slide, top: Emu, data_root: str, images: list[dict]):
    x = Emu(457200)
    y = int(top)
    max_w = Emu(2743200)  # 3 in
    for img in images[:6]:
        key = img.get("storageKey")
        if not key:
            continue
        path = os.path.join(data_root, key)
        if not os.path.exists(path):
            continue
        try:
            slide.shapes.add_picture(path, x, Emu(y), width=max_w)
        except Exception:
            continue
        x = Emu(int(x) + int(max_w) + 228600)
        if int(x) + int(max_w) > int(SLIDE_W):
            x = Emu(457200)
            y += 2057400


def _add_watermark(slide, text: str):
    box = slide.shapes.add_textbox(Emu(1371600), Emu(2743200), Emu(6400800), Emu(1371600))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = WATERMARK_COLOR


def _render_page(prs: Presentation, page: dict, data_root: str, watermark: str | None):
    slide = prs.slides.add_slide(_blank_layout(prs))
    if watermark:
        _add_watermark(slide, watermark)
    top = _add_title(slide, page.get("title", ""), page.get("subtitle"))
    if page.get("body"):
        top = _add_bullets(slide, top, page["body"])
    table = page.get("table")
    if table and table.get("columns"):
        _add_table(slide, top, table["columns"], table.get("rows", []))
        top = Emu(int(top) + 2286000)
    if page.get("images"):
        _add_images(slide, top, data_root, page["images"])


def _build_simple(report_json: dict, prs: Presentation, data_root: str) -> None:
    """The original generic renderer: one slide per report_json page."""
    meta = report_json.get("meta", {})
    lang = normalize_lang(report_json.get("reportLanguage") or meta.get("language"))
    watermark = watermark_text(lang, meta.get("watermark"))
    for page in report_json.get("dynamicPages", []):
        _render_page(prs, page, data_root, watermark)
    for page in report_json.get("staticPages", []):
        _render_page(prs, page, data_root, watermark)


def build_pptx(
    report_json: dict,
    out_path: str,
    data_root: str,
    template_version: str | None = None,
    audience: str = "internal",
    watermark_mode: str = "draft",
) -> tuple[list[str], int]:
    """Render report_json into a PPTX saved at out_path.

    template_version:
      - "simple"            -> original generic page-per-slide renderer
      - "report-template-v1"-> corporate audit template (default)
      - "report-template-v2"-> full 36-page dynamic audit template
      - "report-template-v3"-> polished audit + final commercial block

    audience / watermark_mode only affect v3 (others keep prior behaviour).

    Returns (warnings, slide_count). If a template fails entirely it falls back
    to the simple renderer so a deck is always produced.
    """
    version = (template_version or DEFAULT_TEMPLATE_VERSION).strip()
    warnings: list[str] = []

    if version == "simple":
        prs = _new_presentation()
        _build_simple(report_json, prs, data_root)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)
        return warnings, len(prs.slides)

    if version == "report-template-v3":
        prs = _new_presentation(TEMPLATE_V3_PATH)
        try:
            build_report_v3(report_json, prs, data_root, warnings, audience, watermark_mode)
            if len(prs.slides) == 0:
                raise RuntimeError("template v3 produced no slides")
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"template v3 failed ({exc}); fell back to simple renderer")
            prs = _new_presentation()
            _build_simple(report_json, prs, data_root)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)
        return warnings, len(prs.slides)

    if version == "report-template-v2":
        prs = _new_presentation(TEMPLATE_V2_PATH)
        try:
            build_report_v2(report_json, prs, data_root, warnings)
            if len(prs.slides) == 0:
                raise RuntimeError("template v2 produced no slides")
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"template v2 failed ({exc}); fell back to simple renderer")
            prs = _new_presentation()
            _build_simple(report_json, prs, data_root)
        os.makedirs(os.path.dirname(out_path), exist_ok=True)
        prs.save(out_path)
        return warnings, len(prs.slides)

    # report-template-v1 (default)
    prs = _new_presentation(TEMPLATE_V1_PATH)
    try:
        build_report_v1(report_json, prs, data_root, warnings)
        if len(prs.slides) == 0:
            raise RuntimeError("template v1 produced no slides")
    except Exception as exc:  # noqa: BLE001
        warnings.append(f"template v1 failed ({exc}); fell back to simple renderer")
        prs = _new_presentation()
        _build_simple(report_json, prs, data_root)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return warnings, len(prs.slides)
