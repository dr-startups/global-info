"""Corporate audit report template v1 (Stage K1).

Builds a multi-page PPTX in a compliance-audit structure: dynamic analytical
pages (from the audit summary) followed by static commercial offer pages. The
visual style is intentionally simple-but-presentational (dark cover, blue
accent, light tables, risk badges) and is meant to be refined in K2.

Robustness: each slide is built in isolation. If one slide fails, the error is
captured as a warning and the rest of the deck still renders.

No LLM, no network — only lays out the data passed in report_json.
"""

from __future__ import annotations

from typing import Any, Callable

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from report_mapper import build_view_model

# --- palette -----------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x1F, 0x3A)
ACCENT = RGBColor(0x1C, 0x6F, 0xD6)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_HEAD = RGBColor(0x1F, 0x3A, 0x5F)
TABLE_ALT = RGBColor(0xEF, 0xF3, 0xF9)
WATERMARK_COLOR = RGBColor(0xDD, 0xDD, 0xDD)

RISK_COLORS = {
    "LOW": RGBColor(0x2E, 0x7D, 0x32),
    "MEDIUM": RGBColor(0xB8, 0x86, 0x00),
    "HIGH": RGBColor(0xC6, 0x4A, 0x00),
    "CRITICAL": RGBColor(0xB0, 0x1E, 0x1E),
    "UNKNOWN": RGBColor(0x6B, 0x6B, 0x6B),
    "NONE": RGBColor(0x6B, 0x6B, 0x6B),
}

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(457200)
CONTENT_W = Emu(8229600)
MAX_TABLE_ROWS = 12

# Localizable table footnote ("Showing top N of M."), shared with v2. Set per
# render via ``set_table_strings`` so v1/v2 tables honour the report language.
_SHOWING_TOP = "Showing top {n} of {total}."


def set_table_strings(showing_top: str | None) -> None:
    global _SHOWING_TOP
    if showing_top:
        _SHOWING_TOP = showing_top


def _blank(prs):
    layouts = prs.slide_layouts
    return layouts[6] if len(layouts) > 6 else layouts[-1]


def _bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    return box


def _watermark(slide, text: str) -> None:
    box = _textbox(slide, Emu(1371600), Emu(2743200), Emu(6400800), Emu(1371600))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = WATERMARK_COLOR


def _header(slide, title: str, subtitle: str | None = None) -> Emu:
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Emu(64000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    box = _textbox(slide, MARGIN, Emu(228600), CONTENT_W, Emu(914400))
    tf = box.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = title or ""
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = NAVY
    if subtitle:
        sp = tf.add_paragraph()
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.size = Pt(13)
        srun.font.color.rgb = MUTED
    return Emu(1280160)


def _bullets(slide, top: Emu, lines: list[str], size: int = 14) -> Emu:
    if not lines:
        return top
    box = _textbox(slide, MARGIN, top, CONTENT_W, Emu(3600000))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"\u2022 {line}"
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        p.space_after = Pt(4)
    return Emu(int(top) + 320000 * max(1, len(lines)))


def _risk_badge(slide, x: Emu, y: Emu, level: str) -> None:
    lvl = str(level or "UNKNOWN").upper()
    shape = slide.shapes.add_shape(5, x, y, Emu(1600000), Emu(520000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RISK_COLORS.get(lvl, RISK_COLORS["UNKNOWN"])
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = lvl
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _table(slide, top: Emu, columns: list[str], rows: list[list[Any]], note: str | None = None) -> Emu:
    total = len(rows)
    rows = rows[:MAX_TABLE_ROWS]
    n_rows = len(rows) + 1
    n_cols = max(1, len(columns))
    height = Emu(min(360000 * n_rows, 4400000))
    graphic = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, CONTENT_W, height)
    table = graphic.table
    for c, col in enumerate(columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEAD
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(11)
        para.font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            val = row[c] if c < len(row) else ""
            cell.text = "" if val is None else str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(10)
            para.font.color.rgb = INK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
    bottom = Emu(int(top) + int(height) + 120000)
    label = note
    if total > MAX_TABLE_ROWS:
        label = _SHOWING_TOP.format(n=MAX_TABLE_ROWS, total=total) + (f" {note}" if note else "")
    if label:
        box = _textbox(slide, MARGIN, bottom, CONTENT_W, Emu(360000))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(10)
        run.font.italic = True
        run.font.color.rgb = MUTED
        bottom = Emu(int(bottom) + 360000)
    return bottom


def _empty_note(slide, top: Emu, text: str) -> Emu:
    box = _textbox(slide, MARGIN, top, CONTENT_W, Emu(700000))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = MUTED
    return Emu(int(top) + 600000)


# --- slide builders ----------------------------------------------------------

def _slide_cover(prs, vm, _dr):
    slide = prs.slides.add_slide(_blank(prs))
    _bg(slide, NAVY)
    cover = vm["cover"]
    box = _textbox(slide, MARGIN, Emu(2057400), CONTENT_W, Emu(2743200))
    tf = box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = cover["reportTitle"]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = cover["subjectFullName"]
    r2.font.size = Pt(24)
    r2.font.color.rgb = RGBColor(0x9E, 0xC2, 0xF0)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = f"{vm['labels']['audit_date']}: {cover['auditDate']}   ·   {cover['brand']}"
    r3.font.size = Pt(14)
    r3.font.color.rgb = RGBColor(0xC9, 0xD6, 0xEA)
    _risk_badge(slide, MARGIN, Emu(5000000), cover["overallRiskLevel"])
    if vm["meta"].get("watermark"):
        box2 = _textbox(slide, Emu(2057400), Emu(6000000), Emu(5029200), Emu(600000))
        run = box2.text_frame.paragraphs[0].add_run()
        run.text = str(vm["meta"]["watermark"])
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x3A, 0x4F, 0x73)


def _section(prs, vm, watermark):
    slide = prs.slides.add_slide(_blank(prs))
    _bg(slide, WHITE)
    if watermark:
        _watermark(slide, watermark)
    return slide


def _slide_contents(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    top = _header(slide, L["contents"])
    _bullets(slide, top, list(L["contents_list_v1"]))


def _slide_executive(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    top = _header(slide, L["executive_summary"], f"{L['overall_risk']}: {vm['executiveSummary']['overallRiskLevel']}")
    _risk_badge(slide, Emu(7000000), Emu(228600), vm["executiveSummary"]["overallRiskLevel"])
    _bullets(slide, top, vm["executiveSummary"]["bullets"])


def _slide_risk_matrix(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    rm = vm["riskMatrix"]
    top = _header(slide, L["compliance_risk_matrix"], rm["subject"])
    _risk_badge(slide, Emu(7000000), Emu(228600), rm["overallRiskLevel"])
    rows = [[k, str(v)] for k, v in (rm.get("byLevel") or {}).items()]
    if not rows:
        rows = [[L["rm_no_findings"], "0"]]
    top = _table(slide, top, [L["th_risk_level"], L["th_findings"]], rows)
    themes = ", ".join(f"{t['theme']} ({t['count']})" for t in rm.get("topThemes", [])) or "—"
    _bullets(slide, top, [
        L["rm_highest_risk"].format(level=rm["highestRiskLevel"]),
        L["rm_total_findings"].format(n=rm["totalFindings"]),
        L["rm_top_themes"].format(themes=themes),
        L["rm_consequences"].format(items="; ".join(rm.get("consequences", []))),
    ])


def _slide_overview(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    o = vm["digitalProfileOverview"]
    top = _header(slide, L["digital_profile_overview"])
    _bullets(slide, top, [
        L["ov_shares"].format(ru=o["negativeShareRu"], uae=o["negativeShareUae"]),
        L["ov_results"].format(total=o["searchTotal"], neg=o["searchNegative"], share=o["searchNegativeShare"]),
        L["ov_wikipedia"].format(status=o["wikipediaStatus"]),
        L["ov_compliance"].format(summary=o["complianceSummary"]),
    ])


def _region_summary_slide(code: str) -> Callable:
    def builder(prs, vm, _dr):
        L = vm["labels"]
        slide = _section(prs, vm, vm["meta"].get("watermark"))
        r = vm["regions"][code]
        title = L["pg_ru_search_audit"] if code == "RU" else L["pg_uae_search_audit"]
        top = _header(slide, title, f"{L['region_risk']}: {r['riskLevel']}")
        _risk_badge(slide, Emu(7000000), Emu(228600), r["riskLevel"])
        if not r["present"]:
            _empty_note(slide, top, L["no_evidence_region"].format(label=code))
            return
        top = _table(slide, top, [L["th_metric"], L["th_value"]], [
            [L["m_organic_total"], r["organicTotal"]],
            [L["m_organic_negative"], r["organicNegative"]],
            [L["m_negative_share"], r["organicNegativeShare"]],
            [L["m_unique_neg_urls"], r["uniqueNegativeUrls"]],
            [L["m_suggestions_nt"], r["suggestions"]],
            [L["m_images_nt"], r["images"]],
            [L["m_videos_nt"], r["videos"]],
            [L["m_knowledge"], r["knowledgeBlockStatus"]],
        ])
        _bullets(slide, top, [r["conclusion"]] if r["conclusion"] else [])
    return builder


def _region_results_slide(code: str) -> Callable:
    def builder(prs, vm, _dr):
        L = vm["labels"]
        slide = _section(prs, vm, vm["meta"].get("watermark"))
        r = vm["regions"][code]
        title = L["pg_top_search_results"].format(label=code)
        top = _header(slide, title)
        rows = [[x["provider"], x["rank"], x["domain"], x["title"], x["classification"]] for x in r["topResults"]]
        if not rows:
            _empty_note(slide, top, L["nd_no_organic_region"])
            return
        _table(slide, top, [L["th_provider"], L["th_rank"], L["th_domain"], L["th_title"], L["th_class"]], rows)
    return builder


def _slide_ru_themes(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    s = vm["search"]
    top = _header(slide, L["pg_ru_themes"])
    themes = [f"{t['theme']} ({t['count']})" for t in s["topNegativeThemes"]]
    top = _bullets(slide, top, [L["top_themes"] + " " + (", ".join(themes) or "—"),
                                L["negative_domains"] + " " + (", ".join(s["negativeDomains"]) or "—")])
    rows = [[u["title"], u["url"]] for u in s["topNegativeUrls"]]
    if rows:
        _table(slide, top, [L["th_title"], L["th_domain"]], rows)
    else:
        _empty_note(slide, top, L["nd_no_negative_urls"])


def _slide_ru_suggestions(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    r = vm["regions"]["RU"]
    top = _header(slide, L["pg_ru_suggestions"], f"{L['m_suggestions_nt']}: {r['suggestions']}")
    if r["topSuggestions"]:
        _bullets(slide, top, r["topSuggestions"])
    else:
        _empty_note(slide, top, L["nd_no_suggestions"])


def _slide_ru_media(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    r = vm["regions"]["RU"]
    top = _header(slide, L["pg_ru_media"], f"{L['m_images_nt']}: {r['images']}  ·  {L['m_videos_nt']}: {r['videos']}")
    rows = [["Image", i["title"], i["url"]] for i in r["topImages"]]
    rows += [["Video", v["title"], v["url"]] for v in r["topVideos"]]
    if rows:
        _table(slide, top, [L["th_type"], L["th_title"], L["th_source"]], rows)
    else:
        _empty_note(slide, top, L["nd_no_media"])


def _slide_ru_knowledge(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    r = vm["regions"]["RU"]
    su = vm["surfaces"]
    top = _header(slide, L["pg_ru_knowledge"])
    _bullets(slide, top, [
        L["knowledge_block_status"].format(status=r["knowledgeBlockStatus"]),
        f"{L['m_knowledge']}: {su['knowledgeBlocks']} ({su['knowledgeMismatches']}).",
        f"{L['th_source']}: {su['screenshots']}.",
        *L["snapshot_lines"][:1],
    ])


def _slide_wikipedia(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    w = vm["wikipedia"]
    top = _header(slide, L["pg_wikipedia_summary"], w["status"])
    top = _table(slide, top, [L["th_field"], L["th_value"]], [
        [L["m_exists"], L["yes"] if w["exists"] else L["no"]],
        ["URL", w["pageUrl"] or "—"],
        [L["m_language"], w["language"] or "—"],
        [L["m_notability"], w["notabilityScore"]],
    ])
    _bullets(slide, top, [w["conclusion"]] if w["conclusion"] else [])


def _slide_compliance(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    c = vm["complianceDatabases"]
    top = _header(slide, L["pg_compliance_databases"])
    top = _table(slide, top, [L["th_metric"], L["th_value"]], [
        [L["providers_checked"], ", ".join(c["providersChecked"]) or "—"],
        [L["m_active_matches"], c["activeMatches"]],
        ["PEP", c["pepMatches"]],
        ["RCA", c["rcaMatches"]],
        [L["m_sanctions"], c["sanctionsMatches"]],
        [L["m_adverse_media"], c["adverseMediaMatches"]],
    ])
    _bullets(slide, top, [c["conclusion"]] if c["conclusion"] else [])


def _slide_risk_findings(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    rf = vm["riskFindings"]
    top = _header(slide, L["pg_risk_findings_plain"], L["rf_active_findings"].format(n=rf["totalFindings"]))
    rows = [[f["severity"], f["theme"], f["title"], f["reviewStatus"], f["evidenceCount"]] for f in rf["topFindings"]]
    if rows:
        _table(slide, top, [L["th_severity"], L["th_theme"], L["th_finding"], L["th_review"], L["th_evidence"]], rows)
    else:
        _empty_note(slide, top, L["nd_no_findings_global"])


def _slide_data_quality(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    d = vm["dataQuality"]
    top = _header(slide, L["pg_data_quality_plain"])
    top = _table(slide, top, [L["th_metric"], L["th_value"]], [
        [L["dq_evidence_items"], d["evidenceCount"]],
        [L["dq_reviewed"], d["reviewedFindings"]],
        [L["dq_pending"], d["pendingFindings"]],
        [L["dq_dismissed"], d["dismissedFindings"]],
        [L["dq_missing_sections"], ", ".join(d["missingSections"]) or L["dq_none"]],
    ])
    _bullets(slide, top, d["warnings"] or [L["dq_coverage_adequate"]])


def _slide_recommended(prs, vm, _dr):
    L = vm["labels"]
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    top = _header(slide, L["pg_recommended_plain"])
    _bullets(slide, top, vm["recommendedActions"])


def _offer_slide(page: dict) -> Callable:
    def builder(prs, vm, _dr):
        slide = prs.slides.add_slide(_blank(prs))
        _bg(slide, WHITE)
        top = _header(slide, page.get("title", ""), page.get("subtitle"))
        if page.get("bullets"):
            top = _bullets(slide, top, page["bullets"])
        tbl = page.get("table")
        if tbl and tbl.get("columns"):
            _table(slide, top, tbl["columns"], tbl.get("rows", []))
    return builder


def build_report_v1(report_json: dict, prs, data_root: str, warnings: list[str]) -> None:
    vm, vm_warnings = build_view_model(report_json)
    warnings.extend(vm_warnings)
    set_table_strings(vm["labels"]["showing_top"])

    builders: list[tuple[str, Callable]] = [
        ("cover", _slide_cover),
        ("contents", _slide_contents),
        ("executive", _slide_executive),
        ("risk_matrix", _slide_risk_matrix),
        ("overview", _slide_overview),
        ("ru_summary", _region_summary_slide("RU")),
        ("ru_results", _region_results_slide("RU")),
        ("ru_themes", _slide_ru_themes),
        ("ru_suggestions", _slide_ru_suggestions),
        ("ru_media", _slide_ru_media),
        ("ru_knowledge", _slide_ru_knowledge),
        ("wikipedia", _slide_wikipedia),
        ("uae_summary", _region_summary_slide("UAE")),
        ("uae_results", _region_results_slide("UAE")),
        ("compliance", _slide_compliance),
        ("risk_findings", _slide_risk_findings),
        ("data_quality", _slide_data_quality),
        ("recommended", _slide_recommended),
    ]
    for page in vm["offerPages"]:
        builders.append((f"offer:{page.get('title', '')}", _offer_slide(page)))

    for name, fn in builders:
        try:
            fn(prs, vm, data_root)
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"Slide '{name}' failed: {exc}")
