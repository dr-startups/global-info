"""Full dynamic audit report template v2 (Stage K2).

Extends the corporate template to the full 36-page dynamic structure:
  1-5    Cover / Contents / Executive / Risk Matrix / Overview
  6-21   Russia (RU) digital-profile audit block
  22-31  UAE / International audit block
  32-36  Compliance databases + final dynamic conclusion
Followed by static offer pages (37+).

Reuses the v1 styling/builders. Every page is built through a safe wrapper: a
failing slide produces a fallback note + a warning instead of breaking the deck.
Robust against empty cases and missing regions/sections.

No LLM, no network, no live SERP screenshots — only lays out report_json data.
"""

from __future__ import annotations

from typing import Any, Callable

from report_mapper import build_view_model_v2
from report_template_v1 import (
    ACCENT,
    INK,
    MUTED,
    NAVY,
    WHITE,
    _bg,
    _blank,
    _bullets,
    _empty_note,
    _header,
    _risk_badge,
    _section,
    _table,
    _textbox,
)
from pptx.util import Emu, Pt


# --- shared mini helpers -----------------------------------------------------

def _kv_table(slide, top, pairs: list[tuple[str, Any]]):
    return _table(slide, top, ["Metric", "Value"], [[k, v] for k, v in pairs])


def _note(slide, top, text: str):
    return _empty_note(slide, top, text)


def _price(offer: dict, value: Any) -> str:
    currency = offer.get("currency", "EUR")
    try:
        return f"{int(value):,} {currency}"
    except (TypeError, ValueError):
        return f"0 {currency}"


# ===========================================================================
# 1-5: front matter
# ===========================================================================

def _p_cover(prs, vm):
    slide = prs.slides.add_slide(_blank(prs))
    _bg(slide, NAVY)
    c = vm["cover"]
    box = _textbox(slide, Emu(457200), Emu(1828800), Emu(8229600), Emu(3000000))
    tf = box.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = "Digital Profile"
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = WHITE
    for text, size, color in [
        (c["subjectFullName"], 26, (0x9E, 0xC2, 0xF0)),
        (f"Audit date: {c['auditDate']}", 14, (0xC9, 0xD6, 0xEA)),
        (c["brand"], 14, (0xC9, 0xD6, 0xEA)),
        (" · ".join(x for x in [c.get("website", ""), c.get("contact", "")] if x), 12, (0x8A, 0x9C, 0xBC)),
    ]:
        if not text:
            continue
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        from pptx.dml.color import RGBColor

        run.font.color.rgb = RGBColor(*color)
    _risk_badge(slide, Emu(457200), Emu(5200000), c["overallRiskLevel"])
    if vm["meta"].get("watermark"):
        wm = _textbox(slide, Emu(2057400), Emu(6050000), Emu(5029200), Emu(560000))
        run = wm.text_frame.paragraphs[0].add_run()
        run.text = str(vm["meta"]["watermark"])
        run.font.size = Pt(26)
        run.font.bold = True
        from pptx.dml.color import RGBColor

        run.font.color.rgb = RGBColor(0x3A, 0x4F, 0x73)


def _p_contents(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    top = _header(slide, "Contents")
    _bullets(slide, top, vm["contents"]["sections"])


def _p_executive(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    e = vm["executiveSummary"]
    top = _header(slide, "Executive summary", f"Overall risk: {e['overallRiskLevel']}")
    _risk_badge(slide, Emu(7000000), Emu(228600), e["overallRiskLevel"])
    lines = list(e.get("bullets", []))
    for g in e.get("keyFindings", [])[:5]:
        if g.get("title"):
            lines.append(g["title"])
    top = _bullets(slide, top, lines[:10])
    if e.get("dataQualityWarning"):
        _note(slide, top, f"Data quality: {e['dataQualityWarning']}")


def _p_risk_matrix(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    rm = vm["riskMatrix"]
    top = _header(slide, "Compliance risk matrix", rm["subject"])
    _risk_badge(slide, Emu(7000000), Emu(228600), rm["overallRiskLevel"])
    rows = [[r["area"], r["problems"], r["level"], r["consequences"]] for r in rm["rows"]]
    _table(slide, top, ["Compliance area", "Problems & risks", "Risk", "Possible consequences"], rows)


def _p_overview(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    o = vm["overview"]
    top = _header(slide, "Digital profile overview")
    _risk_badge(slide, Emu(7000000), Emu(228600), o.get("overallRiskLevel", "UNKNOWN"))
    _kv_table(
        slide,
        top,
        [
            ("RU negative share", o.get("negativeShareRu", "0%")),
            ("UAE negative share", o.get("negativeShareUae", "0%")),
            ("Search results (total / negative)", f"{o.get('searchTotal', 0)} / {o.get('searchNegative', 0)}"),
            ("Wikipedia", o.get("wikipediaStatus", "")),
            ("Compliance", o.get("complianceSummary", "")),
        ],
    )


# ===========================================================================
# region pages (used for RU and UAE/International)
# ===========================================================================

def _rb(block_key: str, builder: Callable[[Any, dict, dict], None], title: str, subtitle: str | None = None) -> Callable:
    def page(prs, vm):
        slide = _section(prs, vm, vm["meta"].get("watermark"))
        blk = vm[block_key]
        sub = subtitle
        full_title = title.replace("{label}", blk["label"])
        top = _header(slide, full_title, sub)
        _risk_badge(slide, Emu(7000000), Emu(228600), blk["riskLevel"])
        if not blk["present"]:
            _note(slide, top, blk["noDataText"] or "No evidence collected for this region.")
            return
        builder(slide, blk, vm)

    return page


def _b_summary(slide, blk, vm):
    s = blk["summary"]
    top = Emu(1280160)
    top = _kv_table(
        slide,
        top,
        [
            ("Organic total", s.get("organicTotal", 0)),
            ("Organic negative", s.get("organicNegative", 0)),
            ("Negative share", s.get("organicNegativeShare", "0%")),
            ("Suggestions (neg/total)", s.get("suggestions", "0/0")),
            ("Images (neg/total)", s.get("images", "0/0")),
            ("Videos (neg/total)", s.get("videos", "0/0")),
            ("Knowledge block", s.get("knowledgeBlockStatus", "ABSENT")),
        ],
    )
    _bullets(slide, top, [blk["conclusion"]] if blk["conclusion"] else [])


def _b_organic_overview(slide, blk, vm):
    o = blk["organicOverview"]
    top = Emu(1280160)
    top = _kv_table(
        slide,
        top,
        [
            ("Organic total", o.get("organicTotal", 0)),
            ("Organic negative", o.get("organicNegative", 0)),
            ("Unique negative URLs", o.get("uniqueNegativeUrls", 0)),
            ("Total unique URLs", o.get("totalUniqueUrls", 0)),
            ("Negative share", o.get("negativeShare", "0%")),
        ],
    )
    if o.get("observedQueries"):
        _bullets(slide, top, ["Observed suggestions/queries:"] + o["observedQueries"])


def _b_top_results(slide, blk, vm):
    rows = [[x["provider"], x["rank"], x["domain"], x["title"], x["classification"]] for x in blk["topResults"]]
    if rows:
        _table(slide, Emu(1280160), ["Provider", "Rank", "Domain", "Title", "Class"], rows)
    else:
        _note(slide, Emu(1280160), "No organic results collected for this region.")


def _b_themes(slide, blk, vm):
    t = blk["themes"]
    top = Emu(1280160)
    themes = [f"{x['theme']} ({x['count']})" for x in t["topThemes"]]
    top = _bullets(
        slide,
        top,
        [
            "Top themes: " + (", ".join(themes) or "—"),
            "Negative domains: " + (", ".join(t["negativeDomains"]) or "—"),
        ],
    )
    rows = [[u["title"], u["domain"], u["classification"]] for u in t["negativeUrls"]]
    if rows:
        _table(slide, top, ["Title", "Domain", "Class"], rows)
    else:
        _note(slide, top, "No negative URLs detected.")


def _b_snapshots(slide, blk, vm):
    _bullets(
        slide,
        Emu(1280160),
        [
            f"Knowledge block status: {blk['summary'].get('knowledgeBlockStatus', 'ABSENT')}.",
            "Screenshot evidence is shown where captured and stored privately.",
            "Synthetic snapshots are generated from API data — they are NOT live SERP screenshots.",
            "Live browser automation / scraping is not used by this audit.",
        ],
    )


def _b_suggestions(slide, blk, vm):
    sg = blk["suggestions"]
    top = _header_metrics(slide, f"Total: {sg['total']}  ·  Negative: {sg['negative']}")
    if sg["list"]:
        _bullets(slide, top, sg["list"])
    else:
        _note(slide, top, "No suggestions collected for this region.")


def _b_related(slide, blk, vm):
    rq = blk["relatedQueries"]
    top = _header_metrics(slide, f"Total: {rq['total']}  ·  Negative: {rq['negative']}")
    if rq["list"]:
        _bullets(slide, top, rq["list"])
    else:
        _note(slide, top, "No related queries collected for this region.")


def _b_images(slide, blk, vm):
    im = blk["images"]
    top = _header_metrics(slide, f"Images total: {im['total']}  ·  Negative: {im['negative']}")
    rows = [[i["title"], i["source"]] for i in im["items"]]
    if rows:
        _table(slide, top, ["Image title", "Source"], rows)
    else:
        _note(slide, top, "No image results collected for this region.")


def _b_videos(slide, blk, vm):
    vi = blk["videos"]
    top = _header_metrics(slide, f"Videos total: {vi['total']}  ·  Negative: {vi['negative']}")
    rows = [[v["title"], v["source"]] for v in vi["items"]]
    if rows:
        _table(slide, top, ["Video title", "Source"], rows)
    else:
        _note(slide, top, "No video results collected for this region.")


def _b_images_videos(slide, blk, vm):
    im, vi = blk["images"], blk["videos"]
    top = _header_metrics(slide, f"Images: {im['negative']}/{im['total']}  ·  Videos: {vi['negative']}/{vi['total']}")
    rows = [["Image", i["title"], i["source"]] for i in im["items"]]
    rows += [["Video", v["title"], v["source"]] for v in vi["items"]]
    if rows:
        _table(slide, top, ["Type", "Title", "Source"], rows)
    else:
        _note(slide, top, "No image/video results collected for this region.")


def _b_knowledge(slide, blk, vm):
    kb = blk["knowledgeBlock"]
    top = Emu(1280160)
    if kb and kb.get("title"):
        top = _kv_table(
            slide,
            top,
            [
                ("Status", kb.get("status", "ABSENT")),
                ("Title", kb.get("title", "")),
                ("Source", kb.get("source", "") or "—"),
            ],
        )
        if kb.get("snippet"):
            _bullets(slide, top, [kb["snippet"]])
    else:
        _note(slide, top, f"Knowledge block status: {(kb or {}).get('status', 'ABSENT')}. No knowledge block content collected.")


def _b_wikipedia(slide, blk, vm):
    w = blk["wikipedia"]
    top = Emu(1280160)
    if w.get("exists"):
        top = _kv_table(
            slide,
            top,
            [
                ("Page URL", w.get("pageUrl") or "—"),
                ("Language", w.get("language") or "—"),
                ("Notability score", w.get("notabilityScore", 0)),
            ],
        )
        _bullets(slide, top, [w.get("conclusion", "")])
    else:
        _bullets(
            slide,
            top,
            [
                "Wikipedia article not found.",
                "This indicates the absence of a controlled authoritative profile — it is a low/medium profile risk, NOT an adverse fact.",
            ],
        )


def _b_wiki_knowledge(slide, blk, vm):
    w = blk["wikipedia"]
    kb = blk["knowledgeBlock"]
    top = Emu(1280160)
    _bullets(
        slide,
        top,
        [
            f"Wikipedia: {'page exists' if w.get('exists') else 'no page found'} (language: {w.get('language') or '—'}).",
            f"Knowledge block status: {(kb or {}).get('status', 'ABSENT')}.",
            "Relevant language versions are reviewed where available; no data is treated as neutral.",
        ],
    )


def _b_findings(slide, blk, vm):
    rows = [[f["severity"], f["theme"], f["title"], f["reviewStatus"], f["evidenceCount"]] for f in blk["riskFindings"]]
    if rows:
        _table(slide, Emu(1280160), ["Severity", "Theme", "Finding", "Review", "Evidence"], rows)
    else:
        _note(slide, Emu(1280160), "No risk findings for this region. Run the Risk Classifier to populate.")


def _b_data_quality(slide, blk, vm):
    dq = blk["dataQuality"]
    top = _kv_table(
        slide,
        Emu(1280160),
        [
            ("Organic evidence", dq.get("organic", 0)),
            ("Surface evidence", dq.get("surfaces", 0)),
        ],
    )
    _bullets(slide, top, dq.get("warnings") or ["Evidence coverage is adequate for a preliminary regional assessment."])


def _b_recommended(slide, blk, vm):
    _bullets(slide, Emu(1280160), blk.get("recommendedActions") or ["Expand data collection for this region."])


def _b_evidence(slide, blk, vm):
    rows = [[e["title"], e["domain"], e["provider"], e["classification"]] for e in blk["evidenceAppendix"]]
    if rows:
        _table(slide, Emu(1280160), ["Title", "Domain", "Provider", "Class"], rows)
    else:
        _note(slide, Emu(1280160), "No evidence to list for this region.")


def _b_conclusion(slide, blk, vm):
    top = Emu(1280160)
    _bullets(
        slide,
        top,
        [
            f"Region risk level: {blk['riskLevel']}.",
            blk["conclusion"] or "Region assessment is preliminary and requires manual confirmation.",
        ],
    )


def _header_metrics(slide, subtitle: str):
    # Region pages already drew the header; add a metrics subtitle line.
    box = _textbox(slide, Emu(457200), Emu(1097280), Emu(8229600), Emu(360000))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = subtitle
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED
    return Emu(1500000)


# ===========================================================================
# 32-36: compliance + final
# ===========================================================================

def _p_compliance_overview(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    c = vm["compliance"]
    top = _header(slide, "Compliance databases — overview")
    top = _kv_table(
        slide,
        top,
        [
            ("Providers checked", ", ".join(c["providersChecked"]) or "—"),
            ("Active matches", c["activeMatches"]),
            ("PEP / RCA", f"{c['pepMatches']} / {c['rcaMatches']}"),
            ("Sanctions", c["sanctionsMatches"]),
            ("Adverse media", c["adverseMediaMatches"]),
        ],
    )
    _bullets(slide, top, [c["conclusion"]] if c["conclusion"] else [])


def _p_compliance_provider(prs, vm, key: str, title: str):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    c = vm["compliance"]
    top = _header(slide, title)
    rows = [[r["provider"], r["importMethod"], r["matchType"], r["score"]] for r in c[key]]
    if rows:
        _table(slide, top, ["Provider", "Source method", "Category", "Score"], rows)
    else:
        _note(slide, top, "No screening records for these providers (manual import or official API required).")


def _p_compliance_findings(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    c = vm["compliance"]
    top = _header(slide, "Compliance database — risk findings")
    rows = [[f["severity"], f["theme"], f["title"], f["reviewStatus"], f["evidenceCount"]] for f in c["findings"]]
    if rows:
        _table(slide, top, ["Severity", "Theme", "Finding", "Review", "Evidence"], rows)
    else:
        _note(slide, top, "No compliance-database risk findings recorded.")


def _p_final(prs, vm):
    slide = _section(prs, vm, vm["meta"].get("watermark"))
    f = vm["finalConclusion"]
    top = _header(slide, "Final dynamic audit conclusion", f"Overall risk: {f['overallRiskLevel']}")
    _risk_badge(slide, Emu(7000000), Emu(228600), f["overallRiskLevel"])
    themes = ", ".join(f"{t['theme']} ({t['count']})" for t in f.get("topThemes", [])) or "—"
    lines = [f"Highest risk themes: {themes}"]
    lines += list(f.get("recommendedActions", []))[:5]
    if f.get("missingSections"):
        lines.append("Missing sections: " + ", ".join(f["missingSections"]))
    top = _bullets(slide, top, lines)
    _note(slide, top, "The following pages present our services and proposed remediation plan.")


# ===========================================================================
# static offer pages (37+)
# ===========================================================================

def _p_offer(prs, vm):
    offer = vm["offer"]

    def slide_with(title, subtitle, bullets=None, table=None):
        slide = prs.slides.add_slide(_blank(prs))
        _bg(slide, WHITE)
        top = _header(slide, title, subtitle)
        if bullets:
            top = _bullets(slide, top, bullets)
        if table:
            _table(slide, top, table[0], table[1])

    name = offer.get("productName", "Digital Profile Audit")
    company = offer.get("companyName", name)
    slide_with(name, "Services & solutions", ["Evidence-first digital profile and compliance audits."])
    slide_with("Product overview", company, [
        "Open-source search audit across regions.",
        "Compliance database screening via official API or manual import.",
        "Deterministic risk classification with human review.",
    ])
    slide_with("Solution 1 — Digital Profile", offer.get("solution1Title", "Basic"), [
        "RU + international search audit.", "Search surfaces and negative-link analysis.",
    ])
    slide_with("Solution 1 — Work plan", offer.get("solution1Title", "Basic"), [
        "Scope & lawful basis.", "Search + surfaces collection.", "Risk findings & review.",
    ])
    slide_with("Solution 1 — Expected results", offer.get("solution1Title", "Basic"), [
        "Clear map of the digital footprint.", "Prioritised negative items to address.",
    ])
    slide_with("Solution 1 — Pricing", _price(offer, offer.get("solution1Price")), [offer.get("pricingNotes", "")])
    slide_with("Solution 2 — Compliance Databases", offer.get("solution2Title", "Standard"), [
        "Dow Jones / LexisNexis / World-Check screening.", "PEP / RCA / sanctions / adverse-media categorization.",
    ])
    slide_with("Solution 2 — Work plan", offer.get("solution2Title", "Standard"), [
        "Provider screening (official API / manual import).", "Match verification & documentation.",
    ])
    slide_with("Solution 2 — Pricing", _price(offer, offer.get("solution2Price")), [offer.get("pricingNotes", "")])
    slide_with("Solution 3 — Wikipedia & Authority", offer.get("solution3Title", "Enterprise"), [
        "Authoritative profile assessment.", "Knowledge-panel consistency checks.",
    ])
    slide_with("Solution 3 — Work plan", offer.get("solution3Title", "Enterprise"), [
        "Notability review.", "Authoritative source strategy.", "Ongoing monitoring.",
    ])
    slide_with("Solution 3 — Pricing", _price(offer, offer.get("solution3Price")), [offer.get("pricingNotes", "")])
    slide_with("Process / Timeline", "How an engagement runs", [
        "1. Scope & lawful basis.", "2. Evidence collection.", "3. Risk classification.",
        "4. Analyst review.", "5. Report delivery (PPTX/PDF).",
    ])
    slide_with("About / Contacts", company, [
        f"Contact: {offer.get('contactEmail', '')}",
        f"Website: {offer.get('website', '')}",
        "Reports are advisory; all findings require manual verification.",
    ])


# ===========================================================================
# orchestration
# ===========================================================================

def build_report_v2(report_json: dict, prs, data_root: str, warnings: list[str]) -> None:
    vm, vm_warnings = build_view_model_v2(report_json)
    warnings.extend(vm_warnings)

    ru_pages = [
        ("ru_summary", _rb("ru", _b_summary, "{label} — audit summary")),
        ("ru_organic", _rb("ru", _b_organic_overview, "{label} — organic search overview")),
        ("ru_results", _rb("ru", _b_top_results, "{label} — top search results")),
        ("ru_themes", _rb("ru", _b_themes, "{label} — negative publications & themes")),
        ("ru_snapshots", _rb("ru", _b_snapshots, "{label} — search screens / snapshots")),
        ("ru_suggestions", _rb("ru", _b_suggestions, "{label} — search suggestions")),
        ("ru_related", _rb("ru", _b_related, "{label} — related queries")),
        ("ru_images", _rb("ru", _b_images, "{label} — images")),
        ("ru_videos", _rb("ru", _b_videos, "{label} — videos")),
        ("ru_knowledge", _rb("ru", _b_knowledge, "{label} — knowledge block")),
        ("ru_wikipedia", _rb("ru", _b_wikipedia, "{label} — Wikipedia")),
        ("ru_findings", _rb("ru", _b_findings, "{label} — risk findings")),
        ("ru_quality", _rb("ru", _b_data_quality, "{label} — data quality")),
        ("ru_recommended", _rb("ru", _b_recommended, "{label} — recommended actions")),
        ("ru_evidence", _rb("ru", _b_evidence, "{label} — evidence appendix")),
        ("ru_conclusion", _rb("ru", _b_conclusion, "{label} — interim conclusion")),
    ]
    intl_pages = [
        ("intl_summary", _rb("intl", _b_summary, "{label} — audit summary")),
        ("intl_organic", _rb("intl", _b_organic_overview, "{label} — organic search overview")),
        ("intl_results", _rb("intl", _b_top_results, "{label} — top search results")),
        ("intl_themes", _rb("intl", _b_themes, "{label} — negative themes")),
        ("intl_suggestions", _rb("intl", _b_suggestions, "{label} — suggestions")),
        ("intl_media", _rb("intl", _b_images_videos, "{label} — images & videos")),
        ("intl_wiki", _rb("intl", _b_wiki_knowledge, "{label} — Wikipedia / knowledge context")),
        ("intl_findings", _rb("intl", _b_findings, "{label} — risk findings")),
        ("intl_quality", _rb("intl", _b_data_quality, "{label} — data quality")),
        ("intl_conclusion", _rb("intl", _b_conclusion, "{label} — conclusion")),
    ]

    builders: list[tuple[str, Callable]] = [
        ("cover", _p_cover),
        ("contents", _p_contents),
        ("executive", _p_executive),
        ("risk_matrix", _p_risk_matrix),
        ("overview", _p_overview),
        *ru_pages,
        *intl_pages,
        ("compliance_overview", _p_compliance_overview),
        ("compliance_dow_world", lambda prs_, vm_: _p_compliance_provider(prs_, vm_, "dowWorldRows", "Dow Jones / World-Check summary")),
        ("compliance_lexis", lambda prs_, vm_: _p_compliance_provider(prs_, vm_, "lexisRows", "LexisNexis summary")),
        ("compliance_findings", _p_compliance_findings),
        ("final_conclusion", _p_final),
        ("offer_block", _p_offer),
    ]

    for name, fn in builders:
        try:
            fn(prs, vm)
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"Slide '{name}' failed: {exc}")
            try:
                slide = prs.slides.add_slide(_blank(prs))
                _bg(slide, WHITE)
                _header(slide, name.replace("_", " ").title())
                _empty_note(slide, Emu(1280160), f"This section could not be rendered: {exc}")
            except Exception:  # noqa: BLE001
                pass
