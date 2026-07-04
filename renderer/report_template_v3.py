"""Polished report template v3 (Stage K3).

Same 36-page dynamic structure as v2 plus a full final commercial block, but
rendered through the shared design system in `theme.py`: consistent page frame
(accent bar + footer + page number), metric/risk cards, polished zebra tables
with risk-coloured cells, and a near-final client-facing offer block.

Render options:
  - audience: "internal" (default) shows technical notes; "client" softens them.
  - watermark_mode: "draft" (default) stamps the watermark; "none" hides it.

Every slide is built through a safe wrapper: a failing slide adds a warning +
fallback card instead of breaking the deck. Robust to empty / missing data.

No LLM, no scraping, no live SERP screenshots — only lays out report_json data.
"""

from __future__ import annotations

import io
import re
from typing import Any, Callable

from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

import theme as T
from report_mapper import build_view_model_v3


class Ctx:
    def __init__(self, brand: str, watermark: str | None, internal: bool):
        self.brand = brand
        self.watermark = watermark
        self.internal = internal
        self.page = 1
        self.total = 0
        self.layout_warnings: list[str] = []


# ---------------------------------------------------------------------------
# small helpers
# ---------------------------------------------------------------------------

def _frame(slide, ctx: Ctx, title: str, subtitle: str | None = None, title_width: Emu | None = None) -> Emu:
    return T.page_frame(
        slide,
        title,
        subtitle,
        brand=ctx.brand,
        page_no=ctx.page,
        total=ctx.total,
        watermark=ctx.watermark,
        title_width=title_width,
    )


def _section(prs, ctx: Ctx, title: str, subtitle: str | None = None, title_width: Emu | None = None):
    slide = T.blank_slide(prs)
    top = _frame(slide, ctx, title, subtitle, title_width=title_width)
    return slide, top


def _risk_card_value(level: str, L: dict) -> Any:
    return {"label": L["m_risk_level"], "value": str(level or "UNKNOWN"), "tone": T.RISK_COLORS.get(str(level or "UNKNOWN").upper(), T.NEUTRAL_GRAY)}


def _norm(text: Any) -> str:
    """Normalize text for duplicate detection (case/space-insensitive)."""
    return " ".join(str(text or "").split()).strip().lower()


def _region_saved_for_review_note(hidden: int, L: dict) -> str:
    if hidden <= 0:
        return ""
    if "подсказ" in str(L.get("metric_suggestions", "")).lower():
        return f"+ ещё {hidden} {T.r2_ru_plural_suggestions(hidden)} сохранены для проверки."
    return L.get("region_saved_for_review", "+ ещё {n} записей сохранены для проверки.").format(n=hidden)


# ===========================================================================
# 1-5 front matter
# ===========================================================================

def _p_cover(prs, vm, ctx):
    L = vm["labels"]
    slide = T.blank_slide(prs)
    T.set_bg(slide, T.BRAND_PRIMARY)
    c = vm["cover"]
    ob = vm["offerBlock"]
    # R6.2: keep the cover premium, but still follow the global page-marker contract.
    T.footer(slide, c.get("brand", ctx.brand), ctx.page, ctx.total)
    # accent rule
    rule = slide.shapes.add_shape(T.RECT, T.MARGIN, Emu(1700000), Emu(2200000), Emu(54864))
    rule.fill.solid()
    rule.fill.fore_color.rgb = T.ACCENT
    rule.line.fill.background()

    box = T.textbox(slide, T.MARGIN, Emu(1900000), T.CONTENT_W, Emu(2900000))
    tf = box.text_frame
    T._run(tf.paragraphs[0], L["cover_heading"], T.FS_COVER_TITLE, T.WHITE, bold=True)
    T._run(tf.add_paragraph(), ob["cover"].get("subtitle", ""), T.FS_SUBTITLE + 3, T.ACCENT_SOFT)
    p = tf.add_paragraph()
    T._run(p, c.get("subjectFullName", ""), 26, T.WHITE, bold=True)
    p.space_before = Pt(14)
    T._run(tf.add_paragraph(), f"{L['audit_date']}: {c.get('auditDate', '')}", T.FS_SUBTITLE, T.ACCENT_SOFT)
    T._run(tf.add_paragraph(), c.get("brand", ""), T.FS_SUBTITLE, T.ACCENT_SOFT)
    safe_website = re.sub(r"^https?://", "", str(c.get("website", "")).strip(), flags=re.I)
    contact = " · ".join(x for x in [safe_website, c.get("contact", "")] if x)
    if contact:
        T._run(tf.add_paragraph(), contact, T.FS_NOTE + 1, T.NEUTRAL_GRAY)

    T.risk_badge(slide, T.MARGIN, Emu(5050000), c.get("overallRiskLevel", "UNKNOWN"), w=Emu(1800000), h=Emu(520000))
    if ctx.watermark:
        wm = T.textbox(slide, Emu(5200000), Emu(5100000), Emu(3400000), Emu(520000))
        wp = wm.text_frame.paragraphs[0]
        wp.alignment = 2  # right
        T._run(wp, str(ctx.watermark), 24, T.RGBColor(0x3A, 0x4F, 0x73), bold=True)


def _p_contents(prs, vm, ctx):
    L = vm["labels"]
    slide, top = _section(prs, ctx, L["contents"], L["contents_subtitle"])
    sec = L["section"]
    cards = [
        {"label": f"{sec} 1", "value": L["sec_executive"]},
        {"label": f"{sec} 2", "value": L["sec_russia"]},
        {"label": f"{sec} 3", "value": L.get("sec_international_short", L["sec_international"])},
        {"label": f"{sec} 4", "value": L["sec_compliance"]},
        {"label": f"{sec} 5", "value": L["sec_solutions"]},
        {"label": f"{sec} 6", "value": L["sec_about"]},
    ]
    top = T.metric_cards(slide, top, cards, per_row=3)
    T.bullets(slide, top, vm["contents"]["sections"])


def _p_executive(prs, vm, ctx):
    L = vm["labels"]
    e = vm["executiveSummary"]
    slide, top = _section(prs, ctx, L["executive_summary"], f"{L['overall_risk']}: {e['overallRiskLevel']}")
    o = vm["overview"]
    c = vm["compliance"]
    cards = [
        {"label": L["overall_risk"], "value": e["overallRiskLevel"], "tone": T.RISK_COLORS.get(str(e["overallRiskLevel"]).upper(), T.NEUTRAL_GRAY)},
        {"label": L["m_ru_negative"], "value": o.get("negativeShareRu", "0%")},
        {"label": L["m_uae_negative"], "value": o.get("negativeShareUae", "0%")},
        {"label": L["m_compliance_matches"], "value": c.get("activeMatches", 0), "tone": T.DANGER if c.get("activeMatches") else T.NEUTRAL_GRAY},
    ]
    top = T.metric_cards(slide, top, cards, per_row=4)
    bullet_lines = [b for b in e.get("bullets", [])[:6] if b]
    warning = e.get("dataQualityWarning") if ctx.internal else None
    # Drop the warning if it merely repeats the last bullet (no duplication).
    if warning and bullet_lines and _norm(warning) == _norm(bullet_lines[-1]):
        warning = None
    # Stack sequentially by computed text height; bullets() returns the bottom Y
    # (incl. a safe gap), so the warning always sits below the list.
    top = T.bullets(slide, top, bullet_lines)
    if warning:
        T.note(slide, top, warning, "warning")


def _p_risk_matrix(prs, vm, ctx):
    L = vm["labels"]
    rm = vm["riskMatrix"]
    slide, top = _section(prs, ctx, L["compliance_risk_matrix"], rm["subject"])
    rows = [[r["area"], T.truncate(r["problems"], 60), r["level"], T.truncate(r["consequences"], 46)] for r in rm["rows"]]
    T.table(slide, top, [L["th_compliance_area"], L["th_problems_risks"], L["th_risk"], L["th_consequences"]], rows,
            col_widths=[0.26, 0.34, 0.12, 0.28])


def _p_overview(prs, vm, ctx):
    L = vm["labels"]
    o = vm["overview"]
    slide, top = _section(prs, ctx, L["digital_profile_overview"])
    cards = [
        {"label": L["m_ru_negative_share"], "value": o.get("negativeShareRu", "0%")},
        {"label": L["m_uae_negative_share"], "value": o.get("negativeShareUae", "0%")},
        {"label": L["m_search_neg_total"], "value": f"{o.get('searchNegative', 0)}/{o.get('searchTotal', 0)}"},
        _risk_card_value(o.get("overallRiskLevel", "UNKNOWN"), L),
    ]
    top = T.metric_cards(slide, top, cards, per_row=4)
    T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["profile_summary"],
        [
            f"{L['wikipedia_label']} {o.get('wikipediaStatus', '')}",
            T.truncate(o.get("complianceSummary", ""), 120),
        ],
        min_h=560000,
    )


def _b_results_r2(slide, top, blk, vm, ctx):
    """R2 pilot for RU top-search-results table (slide 8 only)."""
    L = vm["labels"]
    rows = [[x["provider"], x["rank"], x["domain"], T.truncate(x["title"], 46), x["classification"]] for x in blk["topResults"]]
    if not rows:
        T.r2_no_data_state(slide, top, L["nd_no_organic_region"])
        return
    T.table(
        slide,
        top,
        [L["th_provider"], L["th_rank"], L["th_domain"], L["th_title"], L["th_class"]],
        rows,
        col_widths=[0.13, 0.08, 0.22, 0.37, 0.20],
        max_rows=10,
        layout_warnings=ctx.layout_warnings,
    )


def _p_ru_results_r2(prs, vm, ctx):
    """R2 pilot page wrapper for RU top search results (slide 8)."""
    L = vm["labels"]
    blk = vm["ru"]
    title = L["pg_top_results"].replace("{label}", blk["label"])
    slide = T.blank_slide(prs)
    top = T.r2_page_header(slide, title=title, section_marker=L.get("section"))
    T.r2_page_footer(slide, brand=ctx.brand, page_no=ctx.page, total=ctx.total)
    if ctx.watermark:
        T._watermark(slide, str(ctx.watermark))
    T.risk_badge(
        slide,
        Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1500000),
        Emu(228600),
        blk["riskLevel"],
    )
    if not blk["present"]:
        T.r2_no_data_state(
            slide, top,
            blk.get("noDataText") or L["no_evidence_region"].format(label=blk["label"]),
        )
        return
    _b_results_r2(slide, top, blk, vm, ctx)


# ===========================================================================
# region pages (generic across RU / international)
# ===========================================================================

def _rb(block_key: str, builder: Callable, title: str, subtitle: str | None = None) -> Callable:
    def page(prs, vm, ctx):
        blk = vm[block_key]
        page_title = title.replace("{label}", blk["label"])
        title_w = Emu(int(T.CONTENT_W) - 1700000)
        slide, top = _section(prs, ctx, page_title, subtitle, title_width=title_w)
        T.risk_badge(
            slide,
            Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1500000),
            Emu(228600),
            blk["riskLevel"],
        )
        if not blk["present"]:
            T.no_data_card(
                slide, top,
                blk["noDataText"] or vm["labels"]["no_evidence_region"].format(label=blk["label"]),
            )
            return
        builder(slide, top, blk, vm, ctx)

    return page


def _rb_r24(block_key: str, builder: Callable, title: str, subtitle: str | None = None) -> Callable:
    """R2.4 pilot wrapper: region header rhythm + premium no-data state."""

    def page(prs, vm, ctx):
        blk = vm[block_key]
        L = vm["labels"]
        page_title = title.replace("{label}", blk["label"])
        page_title = page_title.replace("{label}", "").strip(" —-")
        if block_key == "intl" and title == L.get("pg_suggestions"):
            page_title = L.get("r24_intl_suggestions_title", "Международные поисковые подсказки")
        slide = T.blank_slide(prs)
        top = T.r2_region_header(
            slide,
            title=page_title,
            subtitle=subtitle,
            section_marker=L.get("section"),
            risk_level=blk.get("riskLevel") or "UNKNOWN",
        )
        T.r2_page_footer(slide, brand=ctx.brand, page_no=ctx.page, total=ctx.total)
        if ctx.watermark:
            T._watermark(slide, str(ctx.watermark))
        if not blk["present"]:
            body = blk.get("noDataText") or L.get("r24_region_no_data_body", L["no_evidence_region"].format(label=blk["label"]))
            T.r2_region_no_data_state(
                slide,
                top,
                headline=L.get("r24_region_no_data_title", "Данные по этому региону не зафиксированы."),
                body=T.r2_region_safe_text(body, labels=L, max_len=170),
                width_ratio=0.70 if str(blk.get("code", "")).upper() == "RU" else 0.72,
            )
            return
        builder(slide, top, blk, vm, ctx)

    return page


def _b_summary(slide, top, blk, vm, ctx):
    L = vm["labels"]
    s = blk["summary"]
    cards = [
        {"label": L["m_organic_total"], "value": s.get("organicTotal", 0)},
        {"label": L["m_organic_negative"], "value": s.get("organicNegative", 0), "tone": T.DANGER},
        {"label": L["m_negative_share"], "value": s.get("organicNegativeShare", "0%")},
        _risk_card_value(blk["riskLevel"], L),
    ]
    top = T.metric_cards(slide, top, cards, per_row=4)
    cards2 = [
        {"label": L["m_suggestions_nt"], "value": s.get("suggestions", "0/0")},
        {"label": L["m_images_nt"], "value": s.get("images", "0/0")},
        {"label": L["m_videos_nt"], "value": s.get("videos", "0/0")},
        {"label": L["m_knowledge"], "value": T.r2_region_metric_value(s.get("knowledgeBlockStatus", "ABSENT"), labels=L)},
    ]
    top = T.metric_cards(slide, top, cards2, per_row=4)
    subs = blk.get("subregions") or {}
    if subs:
        parts = []
        for key in ("uae", "international"):
            sub = subs.get(key) or {}
            if sub.get("present"):
                ssub = sub.get("summary") or {}
                parts.append(
                    f"{sub.get('label', key)}: "
                    f"{L['m_organic_total']} {ssub.get('organicTotal', 0)}, "
                    f"{L['m_suggestions_nt']} {ssub.get('suggestions', '0/0')}, "
                    f"related {(sub.get('relatedQueries') or {}).get('total', 0)}"
                )
        if parts:
            top = T.bullets(slide, top, parts)
    if blk["conclusion"]:
        T.note(slide, top, blk["conclusion"], "info")


def _b_summary_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    s = blk["summary"]
    cards = [
        {"label": L["m_organic_total"], "value": s.get("organicTotal", 0)},
        {"label": L["m_organic_negative"], "value": s.get("organicNegative", 0), "tone": T.DANGER},
        {"label": L["m_negative_share"], "value": s.get("organicNegativeShare", "0%")},
        _risk_card_value(blk["riskLevel"], L),
    ]
    top = T.r2_region_metric_row(slide, top, cards, per_row=4)
    cards2 = [
        {"label": L.get("metric_suggestions", L["m_suggestions_nt"]), "value": s.get("suggestions", "0/0")},
        {"label": L.get("metric_images", L["m_images_nt"]), "value": s.get("images", "0/0")},
        {"label": L.get("metric_videos", L["m_videos_nt"]), "value": s.get("videos", "0/0")},
        {"label": L.get("metric_knowledge", L["m_knowledge"]), "value": T.r2_region_metric_value(s.get("knowledgeBlockStatus", "ABSENT"), labels=L)},
    ]
    top = T.r2_region_metric_row(slide, top, cards2, per_row=4)
    is_ru = str(blk.get("code", "")).upper() == "RU"
    if not is_ru:
        top = T.r2_region_bullet_sections(
            slide,
            top,
            [
                {
                    "label": L.get("region_international_segment", "Международный сегмент"),
                    "items": [
                        f"Органических материалов: {s.get('organicTotal', 0)}",
                        f"Подсказок для проверки: {s.get('suggestions', '0/0')}",
                        L.get("region_international_no_subject_results", "Подтверждённых международных материалов по субъекту не выявлено."),
                    ],
                }
            ],
            labels=L,
            max_items_per_section=4,
            layout_warnings=ctx.layout_warnings,
        )
        return
    T.r2_region_note(
        slide,
        top,
        L.get("region_no_adverse_organic", "Негативных органических материалов по выбранным релевантным результатам не выявлено."),
        labels=L,
        kind="info",
    )


def _b_organic(slide, top, blk, vm, ctx):
    L = vm["labels"]
    o = blk["organicOverview"]
    cards = [
        {"label": L["m_organic_total"], "value": o.get("organicTotal", 0)},
        {"label": L["m_negative"], "value": o.get("organicNegative", 0), "tone": T.DANGER},
        {"label": L["m_unique_neg_urls"], "value": o.get("uniqueNegativeUrls", 0)},
        {"label": L["m_negative_share"], "value": o.get("negativeShare", "0%")},
    ]
    top = T.metric_cards(slide, top, cards, per_row=4)
    if o.get("observedQueries"):
        T.bullets(slide, top, [L["observed_queries"]] + o["observedQueries"])


def _b_organic_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    o = blk["organicOverview"]
    cards = [
        {"label": L["m_organic_total"], "value": o.get("organicTotal", 0)},
        {"label": L["m_negative"], "value": o.get("organicNegative", 0), "tone": T.DANGER},
        {"label": L["m_unique_neg_urls"], "value": o.get("uniqueNegativeUrls", 0)},
        {"label": L["m_negative_share"], "value": o.get("negativeShare", "0%")},
    ]
    top = T.r2_region_metric_row(slide, top, cards, per_row=4)
    observed = [T.r2_region_safe_text(x, labels=L, max_len=86) for x in (o.get("observedQueries") or []) if x]
    if observed:
        top = T.note(slide, top, L["observed_queries"].rstrip(" :"), "section")
        shown_out: list[int] = []
        top = T.bullets(
            slide,
            top,
            observed,
            max_items=7,
            overflow_note="+ {n} hidden",
            emit_overflow_note=False,
            shown_out=shown_out,
            layout_warnings=ctx.layout_warnings,
        )
        hidden = max(0, 1 + len(observed) - (shown_out[0] if shown_out else 0))
        note = _region_saved_for_review_note(hidden, L)
        if note:
            T.r2_region_note(slide, top, note, labels=L, kind="info")
    else:
        T.r2_region_no_data_state(
            slide,
            top,
            headline=L.get("r24_organic_empty_title", "Органические наблюдения не зафиксированы"),
            body=L.get("r24_organic_empty_body", "Подтверждённых материалов для отдельного вывода не выявлено."),
            width_ratio=0.70 if str(blk.get("code", "")).upper() == "RU" else 0.72,
        )


def _b_results(slide, top, blk, vm, ctx):
    L = vm["labels"]
    rows = [[x["provider"], x["rank"], x["domain"], T.truncate(x["title"], 46), x["classification"]] for x in blk["topResults"]]
    if rows:
        T.table(slide, top, [L["th_provider"], L["th_rank"], L["th_domain"], L["th_title"], L["th_class"]], rows,
                col_widths=[0.13, 0.08, 0.22, 0.37, 0.20], max_rows=10,
                layout_warnings=ctx.layout_warnings)
    else:
        T.no_data_card(slide, top, L["nd_no_organic_region"])


def _b_results_r23(slide, top, blk, vm, ctx):
    """R2.3c polish: compact client-safe top-results table (slides 8 and 24)."""
    L = vm["labels"]
    rows = list(blk.get("topResults") or [])
    if not rows:
        T.no_data_card(slide, top, L["nd_no_organic_region"])
        return
    T.r2_top_results_table(
        slide,
        top,
        rows=rows,
        region=str(blk.get("label") or ""),
        max_rows=8,
        labels=L,
        layout_warnings=ctx.layout_warnings,
    )


def _b_themes(slide, top, blk, vm, ctx):
    L = vm["labels"]
    t = blk["themes"]
    themes = [f"{x['theme']} ({x['count']})" for x in t["topThemes"]]
    top = T.bullets(slide, top, [
        L["top_themes"] + " " + (", ".join(themes) or "—"),
        L["negative_domains"] + " " + (", ".join(t["negativeDomains"]) or "—"),
    ])
    rows = [[T.truncate(u["title"], 56), u["domain"], u["classification"]] for u in t["negativeUrls"]]
    if rows:
        T.table(slide, top, [L["th_title"], L["th_domain"], L["th_class"]], rows, col_widths=[0.5, 0.3, 0.2],
                max_rows=11, layout_warnings=ctx.layout_warnings)
    else:
        has_signals = bool(t.get("topThemes") or t.get("negativeDomains"))
        T.no_data_card(
            slide,
            top,
            L.get("nd_no_confirmed_negative_urls_with_signals", L["nd_no_negative_urls"])
            if has_signals
            else L["nd_no_negative_urls"],
        )


def _serp_caption(slide, top: Emu, ss: dict, L: dict) -> None:
    box = T.textbox(slide, T.MARGIN, top, T.CONTENT_W, Emu(540000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    T._run(p, ss.get("caption") or L["serp_snapshot_caption"], T.FS_NOTE, T.NEUTRAL_GRAY, italic=True)
    # Stage N1.2 — provenance line (real / mixed / mock / empty).
    source_note = ss.get("source_note")
    if source_note:
        snp = tf.add_paragraph()
        snp.alignment = PP_ALIGN.CENTER
        T._run(snp, source_note, T.FS_NOTE, T.NEUTRAL_GRAY)
    detail = " · ".join(x for x in [ss.get("query", ""), ss.get("generatedAt", "")] if x)
    if detail:
        sp = tf.add_paragraph()
        sp.alignment = PP_ALIGN.CENTER
        T._run(sp, detail, T.FS_FOOTER, T.NEUTRAL_GRAY)


def _place_serp_image(slide, top: Emu, ss: dict, L: dict) -> None:
    """Embed the ORION-style PNG: contain-scaled, centred, above the footer."""
    caption_reserve = 600000  # room for the caption + provenance line below the image
    gap = 90000
    avail_w = int(T.CONTENT_W)
    avail_top = int(top)
    avail_h = max(1, int(T.FOOTER_Y) - avail_top - caption_reserve - gap)

    stream = io.BytesIO(ss["image_bytes"])
    img_w = int(ss.get("width") or 0)
    img_h = int(ss.get("height") or 0)

    if img_w > 0 and img_h > 0:
        scale = min(avail_w / img_w, avail_h / img_h)
        draw_w = max(1, int(img_w * scale))
        draw_h = max(1, int(img_h * scale))
        left = Emu(int(T.MARGIN) + (avail_w - draw_w) // 2)
        slide.shapes.add_picture(stream, left, top, width=Emu(draw_w), height=Emu(draw_h))
        cap_top = Emu(avail_top + draw_h + gap)
    else:
        # Unknown native size: add by width, then clamp height if needed.
        pic = slide.shapes.add_picture(stream, T.MARGIN, top, width=Emu(avail_w))
        if int(pic.height) > avail_h:
            ratio = avail_h / int(pic.height)
            pic.width = Emu(int(int(pic.width) * ratio))
            pic.height = Emu(avail_h)
        pic.left = Emu(int(T.MARGIN) + (avail_w - int(pic.width)) // 2)
        pic.top = top
        cap_top = Emu(int(pic.top) + int(pic.height) + gap)

    _serp_caption(slide, cap_top, ss, L)


def _p_snapshots(prs, vm, ctx):
    """Search screens / snapshots page (Stage S1.5).

    If a synthetic ORION-style SERP snapshot exists, render the image almost
    full-width (contain-scaled, centred) under a localized title/subtitle, with a
    small synthetic-source caption. Otherwise fall back to the original
    informational card so the page (and the 50-slide count) never changes.
    The watermark is drawn by the page frame, so it stays consistent (draft shows
    through around the image; none never appears).
    """
    L = vm["labels"]
    ss = vm.get("serp_snapshot") or {}
    blk = vm.get("ru") or {}

    if ss.get("exists") and ss.get("image_bytes"):
        slide, top = _section(prs, ctx, ss.get("title") or L["search_screens_title"], ss.get("subtitle"))
        try:
            _place_serp_image(slide, top, ss, L)
            return
        except Exception:  # noqa: BLE001 - never crash the deck on an image issue
            T.no_data_card(slide, top, L["serp_snapshot_missing"])
            return

    # Fallback: original search-screens page (title + informational card).
    title = L["pg_search_screens"].replace("{label}", blk.get("label", ""))
    slide, top = _section(prs, ctx, title)
    summary = blk.get("summary") or {}
    T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["search_screens_title"],
        [
            L["knowledge_block_status"].format(
                status=T.r2_region_metric_value(summary.get("knowledgeBlockStatus", "ABSENT"), labels=L)
            ),
            *L["snapshot_lines"],
        ],
        min_h=720000,
    )


def _b_suggestions(slide, top, blk, vm, ctx):
    L = vm["labels"]
    sg = blk["suggestions"]
    internal = vm.get("audience") == "internal"
    lw = ctx.layout_warnings
    top = T.metric_cards(slide, top, [
        {"label": L["m_total"], "value": sg["total"]},
        {"label": L["m_negative"], "value": sg["negative"], "tone": T.DANGER},
    ], per_row=2)
    disclaimer = sg.get("exposureDisclaimer") or L.get("autocomplete_disclaimer", "")
    if disclaimer:
        top = T.note(slide, top, disclaimer, "source")
    overflow = L.get("bullets_overflow_note", "+ {n} more suggestions preserved in evidence.")
    groups = sg.get("groups") or []
    if groups:
        sections = [{"label": g.get("label"), "items": g.get("items") or []} for g in groups if g.get("items") or g.get("label")]
        top = T.bounded_bullet_sections(
            slide, top, sections,
            max_items_per_section=8,
            overflow_note=overflow,
            layout_warnings=lw,
        )
    elif sg["list"]:
        top = T.bullets(slide, top, sg["list"], max_items=8, overflow_note=overflow, layout_warnings=lw)
    else:
        T.no_data_card(slide, top, L["nd_no_suggestions"])


def _b_suggestions_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    sg = blk["suggestions"]
    top = T.r2_region_metric_row(
        slide,
        top,
        [
            {"label": L["m_total"], "value": sg["total"]},
            {"label": L["m_negative"], "value": sg["negative"], "tone": T.DANGER},
        ],
        per_row=2,
    )
    disclaimer = L.get("region_search_suggestions_note") or sg.get("exposureDisclaimer") or L.get("autocomplete_disclaimer", "")
    if disclaimer:
        top = T.r2_region_note(slide, top, disclaimer, labels=L, kind="disclaimer")
    groups = sg.get("groups") or []
    if groups:
        sections = [{"label": g.get("label"), "items": g.get("items") or []} for g in groups if g.get("items") or g.get("label")]
        flat: list[str] = []
        for sec in sections:
            if sec.get("label"):
                flat.append(str(sec["label"]))
            flat.extend([str(x) for x in (sec.get("items") or []) if x])
        shown_out: list[int] = []
        top = T.bullets(
            slide,
            top,
            [T.r2_region_safe_text(x, labels=L, max_len=120) for x in flat if str(x).strip()],
            max_items=8,
            overflow_note="+ {n} hidden",
            emit_overflow_note=False,
            shown_out=shown_out,
            layout_warnings=ctx.layout_warnings,
        )
        hidden = max(0, len(flat) - (shown_out[0] if shown_out else 0))
        note = _region_saved_for_review_note(hidden, L)
        if note:
            T.r2_region_note(slide, top, note, labels=L, kind="info")
    elif sg["list"]:
        shown_out: list[int] = []
        top = T.bullets(
            slide,
            top,
            [T.r2_region_safe_text(x, labels=L, max_len=120) for x in sg["list"]],
            max_items=7,
            overflow_note="+ {n} hidden",
            emit_overflow_note=False,
            shown_out=shown_out,
            layout_warnings=ctx.layout_warnings,
        )
        hidden = max(0, len(sg["list"]) - (shown_out[0] if shown_out else 0))
        note = _region_saved_for_review_note(hidden, L)
        if note:
            T.r2_region_note(slide, top, note, labels=L, kind="info")
    else:
        T.r2_region_no_data_state(
            slide,
            top,
            headline=L.get("r24_suggestions_empty_title", "Рекомендации по подсказкам не зафиксированы"),
            body=L.get("r24_suggestions_empty_body", "Подтверждённых материалов для отдельного вывода не выявлено."),
            width_ratio=0.70 if str(blk.get("code", "")).upper() == "RU" else 0.72,
        )


def _b_related(slide, top, blk, vm, ctx):
    L = vm["labels"]
    rq = blk["relatedQueries"]
    top = T.metric_cards(slide, top, [
        {"label": L["m_total"], "value": rq["total"]},
        {"label": L["m_negative"], "value": rq["negative"], "tone": T.DANGER},
    ], per_row=2)
    subs = blk.get("subregions") or {}
    if subs and rq.get("total", 0) > 0:
        parts = []
        for key in ("uae", "international"):
            sub = subs.get(key) or {}
            sr = sub.get("relatedQueries") or {}
            if (sr.get("total") or 0) > 0:
                parts.append(f"{sub.get('label', key)} ({sr.get('total', 0)}):")
                parts.extend(sr.get("list") or [])
        if parts:
            top = T.bullets(
                slide, top, parts, max_items=8,
                overflow_note=L.get("bullets_overflow_note", "+ {n} more items preserved in evidence."),
                layout_warnings=ctx.layout_warnings,
            )
            return
    if rq["list"]:
        top = T.bullets(
            slide, top, rq["list"], max_items=8,
            overflow_note=L.get("bullets_overflow_note", "+ {n} more items preserved in evidence."),
            layout_warnings=ctx.layout_warnings,
        )
    elif str(rq.get("collectionStatus", "")).upper() == "COLLECTED":
        T.no_data_card(slide, top, L.get("nd_none_found_related", L["nd_no_related"]))
    elif str(rq.get("collectionStatus", "")).upper() in ("NOT_QUERIED", "NOT_CONFIGURED", "NOT_SUPPORTED"):
        T.no_data_card(slide, top, rq.get("statusMessage") or L.get("nd_not_queried_surface", L["nd_no_related"]))
    else:
        T.no_data_card(slide, top, L["nd_no_related"])


def _b_related_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    rq = blk["relatedQueries"]
    top = T.r2_region_metric_row(
        slide,
        top,
        [
            {"label": L["m_total"], "value": rq["total"]},
            {"label": L["m_negative"], "value": rq["negative"], "tone": T.DANGER},
        ],
        per_row=2,
    )
    if rq["list"]:
        top = T.r2_region_bullet_sections(
            slide,
            top,
            [{"label": L.get("pg_related_queries", "Related queries"), "items": rq["list"]}],
            labels=L,
            max_items_per_section=7,
            layout_warnings=ctx.layout_warnings,
        )
    elif str(rq.get("collectionStatus", "")).upper() == "COLLECTED":
        T.r2_region_no_data_state(
            slide,
            top,
            headline=L.get("r24_related_empty_title", "Связанные запросы не зафиксированы"),
            body=L.get("r24_related_empty_body", "Подтверждённых материалов для отдельного вывода не выявлено."),
            width_ratio=0.70,
        )
    else:
        body = rq.get("statusMessage") or L.get("r24_region_no_data_body", L["nd_no_related"])
        T.r2_region_no_data_state(
            slide,
            top,
            headline=L.get("r24_region_no_data_title", "Данные по этому региону не зафиксированы."),
            body=T.r2_region_safe_text(body, labels=L, max_len=170),
            width_ratio=0.70,
        )


def _p_ru_images_orion(prs, vm, ctx):
    """Slide 13 — ORION-style images page (replaces gallery card layout)."""
    L = vm["labels"]
    blk = vm["ru"]
    slide = T.blank_slide(prs)
    if not blk.get("present"):
        top = T.page_frame(
            slide,
            L["pg_images"].replace("{label}", blk["label"]),
            brand=ctx.brand,
            page_no=ctx.page,
            total=ctx.total,
            watermark=ctx.watermark,
        )
        T.no_data_card(
            slide, top,
            blk.get("noDataText") or L["no_evidence_region"].format(label=blk["label"]),
        )
        return
    T.orion_images_slide(slide, blk, vm, ctx, layout_warnings=ctx.layout_warnings)


def _b_images(slide, top, blk, vm, ctx):
    L = vm["labels"]
    im = blk["images"]
    collected = im.get("total", 0)
    items = im.get("items") or []
    selected = im.get("selected", len(items))
    excluded = max(0, collected - selected)
    top = T.metric_cards(slide, top, [
        {"label": L.get("m_images_collected", L["m_images_total"]), "value": collected},
        {"label": L.get("m_images_subject_matched", "Subject-matched"), "value": selected, "tone": T.SUCCESS if selected else T.NEUTRAL_GRAY},
        {"label": L.get("m_images_excluded", "Excluded"), "value": excluded, "tone": T.NEUTRAL_GRAY},
    ], per_row=3)
    if items:
        lw = ctx.layout_warnings
        top = T.image_grid(
            slide, top, items,
            show_identity=False,
            labels=L, layout_warnings=lw,
            orion_gallery=True,
        )
    else:
        T.no_data_card(slide, top, L.get("nd_no_relevant_images", L["nd_no_images"]))


def _b_videos(slide, top, blk, vm, ctx):
    L = vm["labels"]
    vi = blk["videos"]
    collected = vi.get("total", 0)
    selected = vi.get("selected", len(vi.get("items") or []))
    excluded = max(0, collected - selected)
    top = T.metric_cards(slide, top, [
        {"label": L.get("m_videos_collected", L["m_videos_total"]), "value": collected},
        {"label": L.get("m_videos_subject_matched", "Subject-matched"), "value": selected, "tone": T.SUCCESS if selected else T.NEUTRAL_GRAY},
        {"label": L.get("m_videos_excluded", "Excluded"), "value": excluded, "tone": T.NEUTRAL_GRAY},
    ], per_row=3)
    items = vi.get("items") or []
    if items:
        note_key = "media_videos_note" if ctx.internal else "media_videos_note_client"
        top = T.video_cards(
            slide,
            top,
            items,
            L.get("video_open_source", "Open source"),
            labels=L,
            layout_warnings=ctx.layout_warnings,
            note_template=L.get(note_key),
        )
    else:
        T.no_data_card(slide, top, L.get("nd_no_relevant_videos", L["nd_no_videos"]))


def _b_media(slide, top, blk, vm, ctx):
    L = vm["labels"]
    im, vi = blk["images"], blk["videos"]
    internal = vm.get("audience") == "internal"
    no_intl = blk.get("noIntlSubjectResults") or (
        not (im.get("items") or []) and not (vi.get("items") or []) and blk.get("code") == "INTL"
    )
    top = T.metric_cards(slide, top, [
        {"label": L.get("m_images_collected", L["m_images_nt"]), "value": f"{im.get('selected', 0)}/{im.get('total', 0)}"},
        {"label": L.get("m_videos_collected", L["m_videos_nt"]), "value": f"{vi.get('selected', 0)}/{vi.get('total', 0)}"},
    ], per_row=2)
    img_items = im.get("items") or []
    vid_items = vi.get("items") or []
    if no_intl and not img_items and not vid_items:
        T.no_data_card(slide, top, L.get("nd_no_intl_subject_results", L.get("nd_no_relevant_media", L["nd_no_media"])))
        return
    if img_items:
        top = T.image_grid(
            slide,
            top,
            img_items,
            show_identity=internal,
            labels=L,
            layout_warnings=ctx.layout_warnings,
            intl_compact=True,
            allow_cover=False,
            max_items=2,
        )
    if vid_items:
        top = T.video_cards(
            slide,
            top,
            vid_items,
            L.get("video_open_source", "Open source"),
            labels=L,
            layout_warnings=ctx.layout_warnings,
        )
    elif not img_items:
        T.no_data_card(slide, top, L.get("nd_no_relevant_media", L["nd_no_media"]))


def _b_intl_media_r2(slide, top, blk, vm, ctx):
    """R2.2c: intl media cleanup with polished empty-state."""
    L = vm["labels"]
    im, vi = blk["images"], blk["videos"]
    img_items = list(im.get("items") or [])
    vid_items = list(vi.get("items") or [])
    raw_rel_img = int(im.get("selected", len(img_items)) or 0)
    raw_rel_vid = int(vi.get("selected", len(vid_items)) or 0)
    no_media = blk.get("noIntlSubjectResults") or (not img_items and not vid_items)
    rel_img = 0 if no_media else raw_rel_img
    rel_vid = 0 if no_media else raw_rel_vid
    total = int(im.get("total", 0) or 0) + int(vi.get("total", 0) or 0)
    relevant = rel_img + rel_vid
    excluded = max(0, total - relevant)
    top = T.r2_metric_cards(slide, top, [
        {"label": L.get("m_intl_images_relevant", "Relevant images"), "value": rel_img, "tone": T.SUCCESS if rel_img else T.NEUTRAL_GRAY},
        {"label": L.get("m_intl_videos_relevant", "Relevant videos"), "value": rel_vid, "tone": T.SUCCESS if rel_vid else T.NEUTRAL_GRAY},
        {"label": L.get("m_intl_excluded", "Excluded"), "value": excluded, "tone": T.NEUTRAL_GRAY},
    ], per_row=3)

    if no_media:
        T.r2_media_empty_state(
            slide,
            top,
            message=L.get(
                "nd_no_relevant_intl_media_subject",
                "No relevant international images or videos were found for the subject.",
            ),
            detail=L.get(
                "nd_intl_media_partial_saved",
                "Partial-name or patronymic-only materials are saved in internal evidence and excluded from the client report.",
            ),
        )
        return

    shown_img = 0
    shown_vid = 0
    if img_items:
        shown_img = min(2, len(img_items))
        top = T.image_grid(
            slide,
            top,
            img_items,
            show_identity=False,
            labels=L,
            layout_warnings=ctx.layout_warnings,
            intl_compact=True,
            allow_cover=False,
            max_items=2,
        )
    if vid_items:
        shown_vid = min(2, len(vid_items))
        top = T.video_cards(
            slide,
            top,
            vid_items,
            L.get("video_open_source", "Open source"),
            labels=L,
            layout_warnings=ctx.layout_warnings,
            max_items=2,
        )
    shown_total = shown_img + shown_vid
    if relevant > shown_total and shown_total > 0:
        tpl = L.get("media_showing_intl_compact", "Showing {shown} of {total} relevant media sources.")
        T.r2_overflow_note(slide, top, tpl.format(shown=shown_total, total=relevant))


def _p_intl_media_r2(prs, vm, ctx):
    """R2.2c: isolated intl media page path (slide 27)."""
    blk = vm["intl"]
    L = vm["labels"]
    page_title = L["pg_images_videos"].replace("{label}", blk["label"])
    title_w = Emu(int(T.CONTENT_W) - 1700000)
    slide, top = _section(prs, ctx, page_title, title_width=title_w)
    T.risk_badge(
        slide,
        Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1500000),
        Emu(228600),
        blk["riskLevel"],
    )
    if not blk["present"]:
        T.no_data_card(
            slide, top,
            blk["noDataText"] or vm["labels"]["no_evidence_region"].format(label=blk["label"]),
        )
        return
    _b_intl_media_r2(slide, top, blk, vm, ctx)


def _b_knowledge(slide, top, blk, vm, ctx):
    L = vm["labels"]
    kb = blk["knowledgeBlock"]
    if kb and kb.get("title"):
        T.card_auto(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            kb.get("title", L["knowledge_block_default"]),
            [
                f"{L['m_status']}: {T.r2_region_metric_value(kb.get('status', 'ABSENT'), labels=L)}",
                f"{L['th_source']}: {kb.get('source', '') or L.get('domain_unavailable', 'домен не указан')}",
                T.truncate(kb.get("snippet", ""), 160),
            ],
            min_h=620000,
        )
    else:
        T.no_data_card(
            slide,
            top,
            L["no_knowledge_content"].format(
                status=T.r2_region_metric_value((kb or {}).get("status", "ABSENT"), labels=L)
            ),
        )


def _b_wikipedia(slide, top, blk, vm, ctx):
    L = vm["labels"]
    w = blk["wikipedia"]
    if w.get("exists"):
        cards = [
            {"label": L["m_status"], "value": L["m_exists"], "tone": T.SUCCESS},
            {"label": L["m_language"], "value": w.get("language") or "—"},
            {"label": L["m_notability"], "value": w.get("notabilityScore", 0)},
        ]
        top = T.metric_cards(slide, top, cards, per_row=3)
        top = T.bullets(slide, top, [w.get("pageUrl", ""), w.get("conclusion", "")])
    else:
        T.card_auto(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            L["wiki_not_found_title"],
            list(L["wiki_not_found_lines"]),
            tone=T.WARNING,
            min_h=600000,
        )


def _b_wiki_knowledge(slide, top, blk, vm, ctx):
    L = vm["labels"]
    w = blk["wikipedia"]
    kb = blk["knowledgeBlock"]
    state = L["wiki_page_exists"] if w.get("exists") else L["wiki_no_page"]
    T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["wiki_knowledge_title"],
        [
            L["wiki_context_line"].format(state=state, lang=w.get("language") or "—"),
            L["wiki_kb_line"].format(
                status=T.r2_region_metric_value((kb or {}).get("status", "ABSENT"), labels=L)
            ),
            L["wiki_review_line"],
        ],
        min_h=620000,
    )


def _b_findings(slide, top, blk, vm, ctx):
    L = vm["labels"]
    findings = list(blk.get("riskFindings") or [])
    if findings:
        T.r2_risk_findings_table(
            slide,
            top,
            rows=findings,
            labels=L,
            max_rows=6,
            col_widths=[0.46, 0.22, 0.14, 0.18],
            layout_warnings=ctx.layout_warnings,
        )
        return

    is_ru = str(blk.get("code", "")).upper() == "RU"
    T.r2_risk_findings_empty_state(
        slide,
        top,
        headline=L.get(
            "rf_empty_ru_title" if is_ru else "rf_empty_intl_title",
            "Риск-находки по российскому сегменту не зафиксированы" if is_ru else "Международные риск-находки не зафиксированы",
        ),
        body=L.get(
            "rf_empty_ru_body" if is_ru else "rf_empty_intl_body",
            (
                "В доступных источниках нет подтверждённых находок, требующих отдельного вывода."
                if is_ru
                else "По международному сегменту нет подтверждённых находок, требующих отдельного вывода."
            ),
        ),
        width_ratio=0.70 if is_ru else 0.72,
    )


def _b_data_quality(slide, top, blk, vm, ctx):
    L = vm["labels"]
    dq = blk["dataQuality"]
    top = T.metric_cards(slide, top, [
        {"label": L["m_organic_evidence"], "value": dq.get("organic", 0)},
        {"label": L["m_surface_evidence"], "value": dq.get("surfaces", 0)},
    ], per_row=2)
    if ctx.internal:
        T.bullets(slide, top, dq.get("warnings") or [L["coverage_adequate_region"]])
    else:
        T.note(slide, top, L["coverage_on_request"], "disclaimer")


def _b_data_quality_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    dq = blk["dataQuality"]
    top = T.r2_region_metric_row(
        slide,
        top,
        [
            {"label": L["m_organic_evidence"], "value": dq.get("organic", 0)},
            {"label": L["m_surface_evidence"], "value": dq.get("surfaces", 0)},
        ],
        per_row=2,
    )
    T.r2_region_summary_card(
        slide,
        Emu(int(top) + 160000),
        title=L.get("region_data_quality_title", "Качество данных"),
        lines=[
            L.get(
                "r24_data_quality_body",
                "Данные собраны из доступных источников; спорные совпадения требуют ручной проверки. Сводка покрытия доступна для аналитической проверки.",
            )
        ],
        labels=L,
        card_h=Emu(1440000),
    )


def _b_recommended(slide, top, blk, vm, ctx):
    L = vm["labels"]
    T.bullets(slide, top, blk.get("recommendedActions") or [L["expand_region_collection"]])


def _b_evidence(slide, top, blk, vm, ctx):
    L = vm["labels"]
    audience = vm.get("audience") or ("internal" if getattr(ctx, "internal", False) else "client")
    internal = audience == "internal"
    evidence_rows = list(blk.get("evidenceAppendix") or [])
    excluded_rows = list(blk.get("excludedAppendix") or [])

    def _is_review_entry(e: dict) -> bool:
        v = " ".join(
            str(e.get(k, "") or "")
            for k in ("review", "classification", "identity")
        ).lower()
        return any(tok in v for tok in ("review", "pending", "needs", "провер"))

    confirmed_rows = [e for e in evidence_rows if not _is_review_entry(e)]
    review_rows = [e for e in evidence_rows if _is_review_entry(e)]

    top = T.metric_cards(slide, top, [
        {"label": L.get("m_evidence_confirmed", "Подтверждено"), "value": len(confirmed_rows), "tone": T.SUCCESS if confirmed_rows else T.NEUTRAL_GRAY},
        {"label": L.get("m_evidence_review", "На проверке"), "value": len(review_rows), "tone": T.WARNING if review_rows else T.NEUTRAL_GRAY},
        {"label": L.get("m_evidence_excluded", "Исключено"), "value": len(excluded_rows), "tone": T.NEUTRAL_GRAY},
    ], per_row=3)

    def _to_table_row(e: dict) -> list[str]:
        identity_raw = str(e.get("identity") or e.get("classification") or "")
        status_raw = str(e.get("review") or e.get("classification") or "")
        identity = T.r2_status_pill(identity_raw, L)
        status = T.r2_status_pill(status_raw, L)
        return [
            T.r2_truncate_cell_text(e.get("title"), 58),
            T.r2_domain_text(str(e.get("domain") or e.get("link") or ""), 34),
            identity,
            status,
        ]

    groups: list[dict] = []
    if confirmed_rows:
        groups.append(
            {
                "kind": "confirmed",
                "title": L.get("appendix_section_confirmed", L.get("appendix_confirmed_title", "Подтверждённые материалы")),
                "columns": [
                    L.get("th_evidence_material", "Материал"),
                    L["th_domain"],
                    L.get("th_identity", "Идентичность"),
                    L.get("th_status", "Статус"),
                ],
                "rows": [_to_table_row(e) for e in confirmed_rows],
                "col_widths": [0.40, 0.20, 0.20, 0.20],
                "max_rows": 4,
            }
        )
    if internal and review_rows:
        groups.append(
            {
                "kind": "review",
                "title": L.get("appendix_section_review", "Требуют проверки"),
                "columns": [
                    L.get("th_evidence_material", "Материал"),
                    L["th_domain"],
                    L.get("th_identity", "Идентичность"),
                    L.get("th_status", "Статус"),
                ],
                "rows": [_to_table_row(e) for e in review_rows],
                "col_widths": [0.40, 0.20, 0.20, 0.20],
                "max_rows": 2,
            }
        )
    if internal and excluded_rows:
        groups.append(
            {
                "kind": "excluded",
                "title": L.get("appendix_section_excluded", L.get("appendix_excluded_title", "Исключено / шум")),
                "columns": [
                    L.get("th_evidence_material", "Материал"),
                    L["th_domain"],
                    L.get("th_reason", "Причина"),
                    L.get("th_status", "Статус"),
                ],
                "rows": [
                    [
                        T.r2_truncate_cell_text(e.get("title"), 52),
                        T.r2_domain_text(str(e.get("domain") or ""), 34),
                        T.r2_truncate_cell_text(e.get("reason"), 26),
                        L.get("status_excluded_noise", "Исключено / шум"),
                    ]
                    for e in excluded_rows
                ],
                "col_widths": [0.40, 0.20, 0.22, 0.18],
                "max_rows": 2,
            }
        )

    if groups:
        top = T.r2_grouped_evidence_sections(
            slide,
            top,
            groups=groups,
            audience=audience,
            labels=L,
            layout_warnings=ctx.layout_warnings,
        )
    else:
        top = T.no_data_card(slide, top, L.get("nd_no_confirmed_evidence", L["nd_no_evidence_region"]))

    if not internal and (review_rows or excluded_rows):
        T.note(
            slide,
            top,
            L.get(
                "appendix_client_hidden_note",
                "Исключённые и спорные совпадения сохранены во внутреннем evidence и не включены в клиентский отчёт.",
            ),
            "disclaimer",
        )


def _b_conclusion(slide, top, blk, vm, ctx):
    L = vm["labels"]
    T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        f"{L['region_risk']}: {blk['riskLevel']}",
        [blk["conclusion"] or L["interim_conclusion_fallback"]],
        tone=T.RISK_COLORS.get(str(blk["riskLevel"]).upper(), T.NEUTRAL_GRAY),
    )


def _b_conclusion_r24(slide, top, blk, vm, ctx):
    L = vm["labels"]
    summary_line = blk["conclusion"] or L["interim_conclusion_fallback"]
    if str(blk.get("code", "")).upper() == "INTL":
        if re.search(
            r"No international subject-matched results (in collected data|found)\.?",
            str(summary_line),
            flags=re.I,
        ):
            summary_line = L.get("region_international_no_subject_results", summary_line)
    top = T.r2_region_summary_card(
        slide,
        top,
        title=f"{L['region_risk']}: {blk['riskLevel']}",
        lines=[summary_line],
        labels=L,
    )
    T.r2_region_note(
        slide,
        top,
        L.get("r24_conclusion_note", "Дополнительные материалы сохранены для проверки."),
        labels=L,
        kind="disclaimer",
    )


# ===========================================================================
# 32-36 compliance (R1 enterprise due diligence layout)
# ===========================================================================

def _p_compliance_overview(prs, vm, ctx):
    L = vm["labels"]
    c = vm["compliance"]
    slide, top = _section(prs, ctx, L["compliance_overview_title"])
    T.risk_badge(slide, Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1700000), Emu(250000), "LOW", w=Emu(1700000))
    cards = [
        {"label": L["m_total_hits"], "value": c.get("totalHits", 0), "tone": T.ACCENT if c.get("totalHits") else T.NEUTRAL_GRAY},
        {"label": L["m_pending_review"], "value": c.get("pendingHits", 0), "tone": T.WARNING if c.get("pendingHits") else T.NEUTRAL_GRAY},
        {"label": L["m_confirmed_matches"], "value": c.get("confirmedHits", 0), "tone": T.DANGER if c.get("confirmedHits") else T.NEUTRAL_GRAY},
    ]
    top = T.metric_cards(slide, top, cards, per_row=3)
    rows = [
        [
            T.r2_provider_source_text(p.get("provider"), labels=L),
            T.r2_compliance_status_pill(p.get("status"), L),
            T.r2_provider_source_text(p.get("sourceType"), labels=L),
            L.get("comp_comment_manual_review", "Часть материалов требует ручной проверки перед финальной интерпретацией.")
            if "ручн" in str(p.get("sourceType", "")).lower() or "manual" in str(p.get("sourceType", "")).lower()
            else "—",
        ]
        for p in c.get("providerTable", [])
    ]
    if rows:
        top = T.r2_compliance_table(
            slide,
            top,
            columns=[
                L.get("th_source_compact", L["th_source"]),
                L.get("th_status_compact", L["th_status"]),
                L.get("th_materials", "Материалы"),
                L.get("th_comment", "Комментарий"),
            ],
            rows=rows,
            max_rows=6,
            labels=L,
            col_widths=[0.27, 0.17, 0.18, 0.38],
            layout_warnings=ctx.layout_warnings,
        )
    else:
        top = T.no_data_card(slide, top, L["nd_no_compliance_hits"])
    T.warning_card(slide, top, L.get("comp_manual_review_note", "Часть материалов требует ручной проверки перед финальной интерпретацией."))


def _p_compliance_risk_types(prs, vm, ctx):
    L = vm["labels"]
    c = vm["compliance"]
    slide, top = _section(prs, ctx, L["compliance_risk_types_title"])
    T.risk_badge(slide, Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1700000), Emu(250000), "LOW", w=Emu(1700000))
    breakdown = c.get("riskTypeBreakdown") or []
    rows = []
    for b in breakdown:
        found = int(b.get("total", 0) or 0)
        pending = int(b.get("pending", 0) or 0)
        confirmed = int(b.get("confirmed", 0) or 0)
        if confirmed > 0:
            level = L.get("comp_level_high", "Высокий")
        elif pending > 0:
            level = L.get("comp_level_medium", "Средний")
        else:
            level = L.get("comp_level_low", "Низкий")
        rows.append([
            T.r2_compliance_type_label(b.get("riskType"), L),
            found,
            confirmed,
            pending,
            T.r2_compliance_status_pill(level, L, compact=True),
        ])
    if rows:
        T.r2_compliance_table(
            slide,
            top,
            columns=[
                L.get("th_risk_type", "Тип риска"),
                L.get("th_found", L.get("th_total", "Найдено")),
                L.get("th_confirmed_short", L.get("th_confirmed", "Подтв.")),
                L.get("th_review_short", L.get("th_pending", "На проверке")),
                L.get("th_level", "Уровень"),
            ],
            rows=rows,
            max_rows=8,
            labels=L,
            col_widths=[0.33, 0.14, 0.14, 0.17, 0.22],
            layout_warnings=ctx.layout_warnings,
        )
    else:
        T.r2_compliance_empty_state(
            slide,
            top,
            headline=L.get("comp_empty_risk_types_title", "По типам риска совпадений не найдено"),
            body=L.get(
                "comp_empty_risk_types_body",
                "В доступных комплаенс-источниках не зафиксированы категории риска, требующие отдельной группировки.",
            ),
        )


def _p_compliance_top_matches(prs, vm, ctx):
    L = vm["labels"]
    c = vm["compliance"]
    title = L["compliance_top_matches_title"]
    if "Ключевые комплаенс-совпадения" in title:
        title = "Ключевые комплаенс-\nсовпадения"
    slide, top = _section(prs, ctx, title, title_width=Emu(int(T.CONTENT_W) - 2300000))
    T.risk_badge(slide, Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1700000), Emu(250000), "LOW", w=Emu(1700000))
    hits = c.get("topHits") or []
    rows = []
    for h in hits:
        score = float(h.get("score", 0) or 0)
        if score >= 80:
            level = L.get("comp_level_high", "Высокий")
        elif score >= 50:
            level = L.get("comp_level_medium", "Средний")
        else:
            level = L.get("comp_level_low", "Низкий")
        src_primary = h.get("provider") or h.get("source") or "—"
        rows.append([
            T.r2_provider_source_text(src_primary, labels=L),
            T.r2_truncate_cell_text(h.get("matchedName"), 52),
            T.r2_compliance_type_label(h.get("riskTypes"), L),
            T.r2_compliance_status_pill(level, L, compact=False),
            T.r2_compliance_status_pill(h.get("reviewStatus"), L, compact=False),
        ])
    if rows:
        top = T.r2_compliance_table(
            slide,
            top,
            columns=[
                L.get("th_source_compact", L["th_source"]),
                L.get("th_match", L.get("th_matched_name", "Совпадение")),
                L.get("th_type", L.get("th_risk_type", "Тип")),
                L.get("th_level", "Уровень"),
                L.get("th_check", L.get("th_review", "Проверка")),
            ],
            rows=rows,
            max_rows=5,
            labels=L,
            col_widths=[0.14, 0.33, 0.19, 0.14, 0.20],
            note=L.get(
                "comp_top_matches_note",
                "Показаны ключевые совпадения, требующие аналитической проверки.",
            ),
            layout_warnings=ctx.layout_warnings,
        )
    else:
        T.r2_compliance_empty_state(
            slide,
            top,
            headline=L.get("comp_empty_top_matches_title", "Ключевые комплаенс-совпадения не найдены"),
            body=L.get(
                "comp_empty_top_matches_body",
                "По проверенным источникам нет совпадений, которые требуют вынесения в таблицу ключевых результатов.",
            ),
        )


def _p_compliance_review_quality(prs, vm, ctx):
    L = vm["labels"]
    c = vm["compliance"]
    slide, top = _section(prs, ctx, L["compliance_review_quality_title"])
    intel = vm.get("complianceRiskIntel") or {}
    counts = intel.get("counts") or {}
    manual = intel.get("manualImport") or {}
    # R3.5 — make review/data-quality analytical: show buckets, not just warnings.
    if counts:
        top = T.metric_cards(
            slide,
            top,
            [
                {"label": L.get("r35_bucket_confirmed", "Подтверждено"), "value": int(counts.get("confirmed", 0) or 0), "tone": T.DANGER if counts.get("confirmed") else T.NEUTRAL_GRAY},
                {"label": L.get("r35_bucket_review", "На проверке"), "value": int(counts.get("review", 0) or 0), "tone": T.WARNING if counts.get("review") else T.NEUTRAL_GRAY},
                {"label": L.get("r35_bucket_excluded", "Исключено"), "value": int(counts.get("excluded", 0) or 0), "tone": T.NEUTRAL_GRAY},
            ],
            per_row=3,
        )
    warnings = [w for w in (c.get("dataQualityWarnings") or []) if w]
    if c.get("pendingHits", 0) > 0:
        warnings.insert(0, L["lang_requires_review"])
    not_configured = [p["provider"] for p in c.get("providerTable", []) if L["src_not_configured"] in str(p.get("sourceType", ""))]
    if not_configured:
        warnings.append(f"{L['warn_provider_not_queried']} ({', '.join(not_configured[:4])})")
    if int(manual.get("review", 0) or 0) > 0:
        warnings.append(L.get("r35_manual_import_note", "Материалы ручного импорта требуют проверки."))
    if not warnings:
        warnings = [L["dq_coverage_adequate"], L["warn_not_legal"]]
    top = T.bullets(slide, top, warnings[:6], max_items=6, layout_warnings=ctx.layout_warnings)
    T.warning_card(slide, top, L["warn_not_legal"])


def _p_compliance_findings(prs, vm, ctx):
    L = vm["labels"]
    c = vm["compliance"]
    f = vm["finalConclusion"]
    lw = ctx.layout_warnings
    slide, top = _section(prs, ctx, L["compliance_risk_findings_title"], f"{L['overall_risk']}: {f.get('overallRiskLevel', 'UNKNOWN')}")
    T.risk_badge(slide, Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1700000), Emu(250000), "LOW", w=Emu(1700000))
    rows = [
        [
            T.r2_truncate_cell_text(fnd.get("title"), 54),
            T.r2_provider_source_text(fnd.get("source") or fnd.get("riskType"), labels=L),
            T.r2_compliance_status_pill(fnd.get("severity"), L, compact=True),
            T.r2_compliance_status_pill(fnd.get("reviewStatus"), L, compact=False),
        ]
        for fnd in c.get("findings", [])
    ]
    if rows:
        top = T.r2_compliance_table(
            slide,
            top,
            columns=[
                L.get("th_finding", "Находка"),
                L.get("th_source_compact", L["th_source"]),
                L.get("th_level", "Уровень"),
                L.get("th_action", "Действие"),
            ],
            rows=rows,
            max_rows=7,
            labels=L,
            col_widths=[0.48, 0.24, 0.12, 0.16],
            layout_warnings=lw,
        )
    else:
        top = T.r2_compliance_empty_state(
            slide,
            top,
            headline=L.get("comp_empty_findings_title", "Комплаенс-риск-находки не зафиксированы"),
            body=L.get(
                "comp_empty_findings_body",
                "Материалы не содержат подтверждённых риск-находок по комплаенс-базам на момент формирования отчёта.",
            ),
        )
    T.note(
        slide,
        Emu(int(top) + 60000),
        L.get(
            "comp_analytical_summary_note",
            "Материалы являются аналитической сводкой и требуют проверки перед использованием в юридически значимых решениях.",
        ),
        "disclaimer",
    )


def _p_final(prs, vm, ctx):
    """Legacy final slide — kept for v2 parity; v3 folds conclusion into compliance findings."""
    L = vm["labels"]
    f = vm["finalConclusion"]
    slide, top = _section(prs, ctx, L["final_title"], f"{L['overall_risk']}: {f['overallRiskLevel']}")
    T.risk_badge(slide, Emu(int(T.SLIDE_W) - int(T.MARGIN) - 1700000), Emu(250000), f["overallRiskLevel"], w=Emu(1700000))
    themes = ", ".join(f"{t['theme']} ({t['count']})" for t in f.get("topThemes", [])) or "—"
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["highest_risk_themes"],
        [themes],
        min_h=460000,
    ) or top
    top = T.bullets(slide, top, list(f.get("recommendedActions", []))[:5])
    if ctx.internal and f.get("missingSections"):
        T.note(slide, top, L["missing_sections_inline"].format(items=", ".join(f["missingSections"])), "warning")


# ===========================================================================
# final commercial block (37+)
# ===========================================================================

def _p_offer_cover(prs, vm, ctx):
    ob = vm["offerBlock"]
    slide = T.blank_slide(prs)
    T.set_bg(slide, T.BRAND_PRIMARY)
    T.footer(slide, ctx.brand, ctx.page, ctx.total)
    title_box = T.textbox(slide, T.MARGIN, Emu(380000), T.CONTENT_W, Emu(560000))
    ttf = title_box.text_frame
    T._run(ttf.paragraphs[0], ob["cover"]["title"], T.FS_MAIN_TITLE + 2, T.WHITE, bold=True)
    T._run(ttf.add_paragraph(), ob["cover"]["subtitle"], T.FS_SUBTITLE, T.ACCENT_SOFT)
    rule = slide.shapes.add_shape(T.RECT, T.MARGIN, Emu(2300000), Emu(2200000), Emu(54864))
    rule.fill.solid()
    rule.fill.fore_color.rgb = T.ACCENT
    rule.line.fill.background()
    box = T.textbox(slide, T.MARGIN, Emu(2500000), T.CONTENT_W, Emu(2200000))
    tf = box.text_frame
    T._run(tf.paragraphs[0], ob["cover"]["title"], T.FS_COVER_TITLE - 6, T.WHITE, bold=True)
    T._run(tf.add_paragraph(), ob["cover"]["subtitle"], T.FS_SUBTITLE + 2, T.ACCENT_SOFT)
    T._run(tf.add_paragraph(), ob["cover"]["brand"], T.FS_SUBTITLE, T.NEUTRAL_GRAY)
    T._run(tf.add_paragraph(), T.truncate(ob["cover"].get("valueProposition", ""), 180), T.FS_BODY - 1, T.WHITE)
    outcomes = ob["cover"].get("outcomeCards", [])[:3]
    if outcomes:
        cards = [
            {
                "label": T.truncate(str(c.get("title", "")), 26),
                "value": T.truncate(str(c.get("text", "")), 44),
                "tone": T.ACCENT,
            }
            for c in outcomes
        ]
        T.metric_cards(slide, Emu(4960000), cards, per_row=3)


def _p_product_overview(prs, vm, ctx):
    L = vm["labels"]
    ob = vm["offerBlock"]["productOverview"]
    slide, top = _section(prs, ctx, L["offer_product_overview"], vm["offerBlock"]["cover"]["brand"])
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["offer_what_we_do"],
        [ob.get("description", "")],
    ) or top
    caps = ob.get("capabilities", [])
    if caps:
        cards = [
            {
                "label": T.truncate(str(c.get("title", "")), 28),
                "value": T.truncate(str(c.get("text", "")), 48),
                "tone": T.ACCENT,
            }
            for c in caps[:4]
        ]
        top = T.metric_cards(slide, top, cards, per_row=2)
    inc = ob.get("includedItems", [])[:3]
    if inc:
        top = T.card_auto(slide, T.MARGIN, top, T.CONTENT_W, L["offer_includes"], [x for x in inc[:2]], min_h=420000) or top
    T.note(slide, top, ob.get("audienceNote", ""), "info")


def _p_pricing_summary(prs, vm, ctx):
    L = vm["labels"]
    sols = vm["offerBlock"]["solutions"]
    slide, top = _section(prs, ctx, L["offer_solutions_pricing"], L["offer_indicative"])
    cards = [
        {
            "label": T.truncate(s["title"], 26),
            "value": f"{s['price']} · {T.truncate(s.get('duration', '—'), 16)}",
            "tone": T.ACCENT,
        }
        for s in sols[:3]
    ]
    top = T.metric_cards(slide, top, cards, per_row=3)
    rows = [[T.truncate(s["title"], 34), T.truncate(s["duration"], 18), s["price"]] for s in sols]
    if rows:
        top = T.polished_table(slide, top, [L["th_solution"], L["th_duration"], L["th_price"]], rows, col_widths=[0.5, 0.25, 0.25])
    T.note(slide, top, vm["offerBlock"]["solutions"][0].get("pricingNotes", "") if sols else "", "disclaimer")


def _p_solution_objective(prs, vm, ctx, idx: int):
    L = vm["labels"]
    s = _solution(vm, idx)
    slide, top = _section(prs, ctx, s["title"], s["subtitle"])
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["offer_objective"],
        [s["objective"]],
    ) or top
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L.get("offer_business_value", "Business value"),
        [s.get("businessValue", s.get("expectedResults", [""])[0] if s.get("expectedResults") else "")],
        min_h=500000,
    ) or top
    col_w = Emu((int(T.CONTENT_W) - int(T.GUTTER)) // 2)
    left = T.card_auto(
        slide,
        T.MARGIN,
        top,
        col_w,
        L["offer_deliverables"],
        [x for x in (s["deliverables"] or ["—"])[:3]],
        min_h=620000,
    ) or top
    T.card_auto(
        slide,
        Emu(int(T.MARGIN) + int(col_w) + int(T.GUTTER)),
        top,
        col_w,
        L["offer_expected_results"],
        [x for x in (s["expectedResults"] or ["—"])[:3]],
        min_h=620000,
    )
    top = left


def _p_solution_workplan(prs, vm, ctx, idx: int):
    L = vm["labels"]
    s = _solution(vm, idx)
    slide, top = _section(prs, ctx, s["title"] + L["offer_workplan_suffix"], L["offer_duration"].format(d=s["duration"]))
    steps = s["workPlan"] or [L["offer_workplan_default"]]
    top = T.step_cards(slide, top, [T.truncate(step, 84) for step in steps[:5]], per_row=2)
    if s["expectedResults"]:
        top = T.card_auto(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            L["offer_expected_results"],
            [x for x in s["expectedResults"][:3]],
            min_h=500000,
        ) or top


def _p_solution_pricing(prs, vm, ctx, idx: int):
    L = vm["labels"]
    s = _solution(vm, idx)
    slide, top = _section(prs, ctx, s["title"] + L["offer_pricing_suffix"], s["subtitle"])
    top = T.metric_cards(
        slide,
        top,
        [
            {"label": L["m_package_price"], "value": s["price"], "tone": T.ACCENT},
            {"label": L["th_duration"], "value": T.truncate(s.get("duration", "—"), 22), "tone": T.NEUTRAL_DARK},
        ],
        per_row=2,
    )
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["offer_included"],
        [x for x in (s["includedItems"] or ["—"])[:4]],
        min_h=560000,
    ) or top
    T.note(slide, top, s.get("pricingNotes", ""), "disclaimer")


def _p_process(prs, vm, ctx):
    L = vm["labels"]
    ob = vm["offerBlock"]["process"]
    slide, top = _section(prs, ctx, L["offer_process_title"], L["offer_process_subtitle"])
    steps = ob.get("steps", []) or list(L["op_process_bullets"])
    top = T.step_cards(slide, top, [T.truncate(step, 82) for step in steps[:5]], per_row=2)
    outcomes = [x for x in (ob.get("outcomes", []) or []) if str(x).strip()]
    if outcomes:
        top = T.card_auto(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            L.get("offer_expected_results", "Expected results"),
            [x for x in outcomes[:3]],
            min_h=500000,
        ) or top
    T.note(slide, top, L["offer_value"], "info")


def _p_about(prs, vm, ctx):
    L = vm["labels"]
    ob = vm["offerBlock"]["contact"]
    slide, top = _section(prs, ctx, L["offer_about_title"], ob.get("company", ""))
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        L["offer_next_step"],
        [ob.get("cta", L["offer_contact_default"])],
        tone=T.ACCENT,
    ) or top
    safe_website = re.sub(r"^https?://", "", str(ob.get("website", "")).strip(), flags=re.I)
    top = T.metric_cards(
        slide,
        top,
        [
            {"label": L.get("offer_email_label", "Email"), "value": T.truncate(ob.get("email", ""), 34), "tone": T.ACCENT},
            {"label": L.get("offer_website_label", "Website"), "value": T.truncate(safe_website, 34), "tone": T.NEUTRAL_DARK},
        ],
        per_row=2,
    )
    if ob.get("deliveryNote") or ob.get("responseSla"):
        top = T.card_auto(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            L.get("offer_delivery_title", "Delivery"),
            [ob.get("deliveryNote", ""), ob.get("responseSla", "")],
            min_h=460000,
        ) or top
    for d in ob.get("disclaimers", []):
        top = T.note(slide, top, d, "disclaimer")


def _r31_evidence_rows(entries: list[dict], L: dict, *, include_reason: bool = False) -> list[list[str]]:
    def _safe_identity(raw: Any) -> str:
        s = str(raw or "").lower()
        if any(tok in s for tok in ("exact", "confirmed", "match confirmed", "подтверж")):
            return L.get("status_confirmed", "Подтверждено")
        if any(tok in s for tok in ("likely", "возможно", "вероят")):
            return L.get("status_likely_subject", "Вероятно относится к субъекту")
        if any(tok in s for tok in ("review", "pending", "needs", "провер")):
            return L.get("status_needs_review", "Требует проверки")
        return "—"

    def _safe_status(raw: Any) -> str:
        s = str(raw or "").lower()
        if any(tok in s for tok in ("excluded", "noise", "false", "dismiss", "исключ")):
            return L.get("status_excluded_noise", "Исключено / шум")
        if any(tok in s for tok in ("review", "pending", "needs", "провер")):
            return L.get("status_needs_review", "Требует проверки")
        return L.get("status_confirmed", "Подтверждено")

    rows: list[list[str]] = []
    for e in entries:
        title = T.r2_truncate_cell_text(e.get("title"), 54)
        domain = T.r2_domain_text(str(e.get("domain") or e.get("link") or ""), 30)
        if include_reason:
            reason = T.r2_truncate_cell_text(e.get("reason") or L.get("r31_reason_filtered", "Отфильтровано"), 28)
            status = L.get("status_excluded_noise", "Исключено / шум")
            rows.append([title, domain, reason, status])
        else:
            identity = _safe_identity(e.get("identity") or e.get("classification"))
            status = _safe_status(e.get("review") or e.get("classification"))
            rows.append([title, domain, identity, status])
    return rows


def _p_r31_appendix_cover(prs, vm, ctx):
    L = vm["labels"]
    slide, top = _section(
        prs,
        ctx,
        L.get("r31_appendix_cover_title", "Расширенное приложение с доказательствами"),
        L.get("r31_appendix_cover_subtitle", "Структурированная навигация по подтверждённым и спорным материалам"),
    )
    top = T.card(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        Emu(1650000),
        L.get("r31_appendix_cover_card_title", "Назначение раздела"),
        [
            L.get("r31_appendix_cover_line_1", "Раздел объединяет подтверждённые материалы, очередь аналитической проверки и исключённые совпадения."),
            L.get("r31_appendix_cover_line_2", "Все блоки показываются в клиент-safe формате без технических внутренних обозначений."),
        ],
        tone=T.ACCENT,
    )
    T.note(slide, top, L.get("r31_appendix_cover_note", "Подробные материалы сохраняются для последующей аналитической трассируемости."), "disclaimer")


def _p_r31_appendix_overview(prs, vm, ctx):
    L = vm["labels"]
    ov = vm.get("appendixOverview") or {}
    slide, top = _section(prs, ctx, L.get("r31_appendix_overview_title", "Структура расширенного приложения"))
    cards = list(ov.get("cards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    lines = list(ov.get("lines") or [])
    if lines:
        top = T.bullets(slide, top, lines[:6], max_items=6, layout_warnings=ctx.layout_warnings)
    else:
        top = T.no_data_card(slide, top, L.get("r31_no_data_overview", "Достаточные структурированные материалы для обзора не зафиксированы."))
    T.note(slide, top, L.get("r31_appendix_overview_note", "Каждый блок раздела готов к отдельному экспертному разбору."), "info")


def _p_r31_region_evidence(prs, vm, ctx, *, title_key: str, subtitle_key: str, entries: list[dict], include_reason: bool = False):
    L = vm["labels"]
    slide, top = _section(prs, ctx, L.get(title_key, "Раздел"), L.get(subtitle_key, ""))
    rows = _r31_evidence_rows(entries, L, include_reason=include_reason)
    if rows:
        if include_reason:
            top = T.table(
                slide,
                top,
                [L.get("th_evidence_material", "Материал"), L["th_domain"], L.get("th_reason", "Причина"), L.get("th_status", "Статус")],
                rows,
                max_rows=8,
                col_widths=[0.40, 0.20, 0.22, 0.18],
                note_text=L.get("r31_table_note", "Показаны приоритетные записи; полный перечень сохранён в материалах проверки."),
                layout_warnings=ctx.layout_warnings,
            )
        else:
            top = T.table(
                slide,
                top,
                [L.get("th_evidence_material", "Материал"), L["th_domain"], L.get("th_identity", "Идентичность"), L.get("th_status", "Статус")],
                rows,
                max_rows=8,
                col_widths=[0.40, 0.20, 0.20, 0.20],
                note_text=L.get("r31_table_note", "Показаны приоритетные записи; полный перечень сохранён в материалах проверки."),
                layout_warnings=ctx.layout_warnings,
            )
    else:
        top = T.no_data_card(slide, top, L.get("r31_no_data_region_evidence", "По текущему сегменту релевантные записи не зафиксированы."))
    T.note(
        slide,
        top,
        L.get(
            "r31_region_evidence_note",
            "Материалы представлены в клиент-safe формате и сохраняют трассируемость для экспертной проверки.",
        ),
        "disclaimer",
    )


def _p_r31_media_overview(prs, vm, ctx):
    L = vm["labels"]
    md = vm.get("mediaEvidenceOverview") or {}
    slide, top = _section(prs, ctx, L.get("r31_media_overview_title", "Сводка по медиа-доказательствам"))
    cards = list(md.get("cards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    narrative = list(md.get("lines") or [])
    if narrative:
        top = T.card(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            Emu(1650000),
            L.get("r31_media_narrative_title", "Ключевые наблюдения"),
            narrative[:3],
            tone=T.ACCENT,
        )
    else:
        top = T.no_data_card(slide, top, L.get("r31_no_data_media", "Подтверждённые медиа-материалы для отдельной сводки не зафиксированы."))
    T.note(slide, top, L.get("r31_media_note", "Источник и контекст каждого медиа-материала доступны в приложении."), "info")


def _p_r31_risk_reasoning_overview(prs, vm, ctx):
    L = vm["labels"]
    rr = vm.get("riskReasoningOverview") or {}
    slide, top = _section(prs, ctx, L.get("r31_risk_reasoning_overview_title", "Обоснование итогового уровня риска"))
    # R3.5 — connect the overall risk level to evidence buckets (client-safe).
    buckets = rr.get("evidenceBuckets") or {}
    reasoning_summary = str(rr.get("reasoningSummary") or "")
    if reasoning_summary or buckets:
        top = T.metric_cards(
            slide,
            top,
            [
                _risk_card_value(rr.get("overallRiskLevel", "UNKNOWN"), L),
                {"label": L.get("r35_bucket_confirmed", "Подтверждено"), "value": int(buckets.get("confirmed", 0) or 0), "tone": T.DANGER if buckets.get("confirmed") else T.NEUTRAL_GRAY},
                {"label": L.get("r35_bucket_review", "На проверке"), "value": int(buckets.get("review", 0) or 0), "tone": T.WARNING if buckets.get("review") else T.NEUTRAL_GRAY},
                {"label": L.get("r35_bucket_excluded", "Исключено"), "value": int(buckets.get("excluded", 0) or 0), "tone": T.NEUTRAL_GRAY},
            ],
            per_row=4,
        )
        summary_lines = [reasoning_summary] if reasoning_summary else list(rr.get("supportingSignals") or [])[:2]
        if summary_lines:
            top = T.card(
                slide,
                T.MARGIN,
                top,
                T.CONTENT_W,
                Emu(1150000),
                L.get("r35_reasoning_summary_title", "Обоснование"),
                summary_lines,
                tone=T.ACCENT,
            )
            top = Emu(int(top) + 1230000)
        limiting = [x for x in (rr.get("limitingFactors") or []) if x]
        if limiting:
            top = T.bullets(slide, top, limiting[:3], max_items=3, layout_warnings=ctx.layout_warnings)
        T.note(
            slide,
            top,
            str(rr.get("legalSafeDisclaimer") or L.get("r35_legal_safe_note", "")),
            "disclaimer",
        )
        return
    # Fallback (pre-R3.5 data): original supporting-signals layout.
    top = T.metric_cards(
        slide,
        top,
        [
            _risk_card_value(rr.get("overallRiskLevel", "UNKNOWN"), L),
            {"label": L.get("r31_risk_reasoning_themes", "Ключевые темы"), "value": len(rr.get("topThemes") or [])},
            {"label": L.get("r31_risk_reasoning_signals", "Опорные сигналы"), "value": len(rr.get("rows") or [])},
            {"label": L.get("r31_risk_reasoning_actions", "Рекомендуемые действия"), "value": len(rr.get("recommendedActions") or [])},
        ],
        per_row=4,
    )
    lines = list(rr.get("supportingSignals") or [])
    if lines:
        top = T.bullets(slide, top, lines[:6], max_items=6, layout_warnings=ctx.layout_warnings)
    themes = [str(t.get("theme", "")) for t in (rr.get("topThemes") or []) if str(t.get("theme", ""))]
    if themes:
        T.note(slide, top, L.get("r31_risk_reasoning_top_themes", "Приоритетные темы: {themes}.").format(themes=", ".join(themes[:4])), "disclaimer")


def _p_r31_risk_reasoning_by_region(prs, vm, ctx):
    L = vm["labels"]
    rr = vm.get("riskReasoningByRegion") or {}
    slide, top = _section(prs, ctx, L.get("r31_risk_reasoning_region_title", "Региональная детализация обоснования"))
    ru = rr.get("ru") or {}
    intl = rr.get("intl") or {}
    top = T.metric_cards(
        slide,
        top,
        [
            {"label": f"{ru.get('label', 'RU')} — {L['region_risk']}", "value": ru.get("riskLevel", "UNKNOWN"), "tone": T.RISK_COLORS.get(str(ru.get("riskLevel", "UNKNOWN")).upper(), T.NEUTRAL_GRAY)},
            {"label": f"{intl.get('label', 'INTL')} — {L['region_risk']}", "value": intl.get("riskLevel", "UNKNOWN"), "tone": T.RISK_COLORS.get(str(intl.get("riskLevel", "UNKNOWN")).upper(), T.NEUTRAL_GRAY)},
        ],
        per_row=2,
    )
    ru_lines = [ru.get("conclusion", "")] + list(ru.get("signals") or [])[:2]
    intl_lines = [intl.get("conclusion", "")] + list(intl.get("signals") or [])[:2]
    top = T.card(slide, T.MARGIN, top, Emu(int(T.CONTENT_W * 0.49)), Emu(1700000), ru.get("label", "RU"), [x for x in ru_lines if x], tone=T.ACCENT)
    T.card(
        slide,
        Emu(int(T.MARGIN) + int(T.CONTENT_W * 0.51)),
        Emu(int(top) - 1780000),
        Emu(int(T.CONTENT_W * 0.49)),
        Emu(1700000),
        intl.get("label", "INTL"),
        [x for x in intl_lines if x],
        tone=T.ACCENT_SOFT,
    )
    T.note(slide, Emu(int(top) + 80000), L.get("r31_risk_reasoning_region_note", "Региональные выводы отражают текущий объём подтверждённых сигналов и требуют периодической актуализации."), "disclaimer")


def _p_r31_appendix_conclusion(prs, vm, ctx):
    L = vm["labels"]
    ac = vm.get("appendixConclusion") or {}
    slide, top = _section(prs, ctx, L.get("r31_appendix_conclusion_title", "Навигация по расширенному приложению"))
    lines = list(ac.get("lines") or [])
    if lines:
        top = T.card(
            slide,
            T.MARGIN,
            top,
            T.CONTENT_W,
            Emu(1650000),
            L.get("r31_appendix_conclusion_card_title", "Итог раздела"),
            lines[:3],
            tone=T.ACCENT,
        )
    else:
        top = T.no_data_card(slide, top, L.get("r31_no_data_conclusion", "Итоговые материалы приложения не сформированы."))
    T.note(
        slide,
        top,
        L.get("r31_appendix_conclusion_note", "Детальные материалы сохранены для аналитической проверки и доказательной трассируемости."),
        "disclaimer",
    )


def _p_r32_provider_diagnostics(prs, vm, ctx):
    if not ctx.internal:
        return
    L = vm["labels"]
    pd = vm.get("providerDiagnostics") or {}
    slide, top = _section(
        prs,
        ctx,
        L.get("r32_provider_diag_title", "Диагностика источников"),
        L.get(
            "r32_provider_diag_subtitle",
            "Матрица доступности источников по текущей конфигурации и режиму выполнения.",
        ),
    )
    cards = list(pd.get("cards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    rows = []
    for r in list(pd.get("rows") or []):
        rows.append(
            [
                T.r2_truncate_cell_text(r.get("source"), 28),
                T.r2_truncate_cell_text(r.get("category"), 16),
                T.r2_truncate_cell_text(r.get("mode"), 14),
                T.r2_truncate_cell_text(r.get("status"), 16),
                T.r2_truncate_cell_text(r.get("risk"), 10),
                T.r2_truncate_cell_text(r.get("note"), 48),
            ]
        )
    if rows:
        top = T.table(
            slide,
            top,
            [
                L.get("r32_diag_th_source", "Источник"),
                L.get("r32_diag_th_category", "Категория"),
                L.get("r32_diag_th_mode", "Режим"),
                L.get("r32_diag_th_status", "Статус"),
                L.get("r32_diag_th_risk", "Риск"),
                L.get("r32_diag_th_note", "Комментарий"),
            ],
            rows,
            max_rows=8,
            col_widths=[0.18, 0.12, 0.11, 0.14, 0.10, 0.35],
            note_text=L.get("r32_provider_diag_table_note", "Показаны ключевые источники и текущая готовность."),
            layout_warnings=ctx.layout_warnings,
        )
    else:
        top = T.no_data_card(slide, top, L.get("r32_provider_diag_no_data", "Диагностические данные недоступны."))
    notes = [str(x) for x in list(pd.get("auditNotes") or []) if str(x).strip()]
    if notes:
        top = T.bullets(slide, top, notes[:3], max_items=3, layout_warnings=ctx.layout_warnings)
    T.note(
        slide,
        top,
        L.get(
            "r32_provider_diag_note",
            "Технические детали ограничены безопасным описанием, без URL, ключей и секретов.",
        ),
        "disclaimer",
    )


# ===========================================================================
# R3.4 — ORION-like evidence appendix expansion (additive card-style pages)
# ===========================================================================

_R34_CARD_H = Emu(1180000)
_R34_NOTE_RESERVE = 820000


def _r34_source_card_lines(raw: dict, L: dict, kind: str) -> tuple[str, list[str], Any]:
    """Return (card_title, body_lines, accent_tone) for one client-safe source card."""
    title = T.r2_truncate_cell_text(raw.get("title") or L.get("th_evidence_material", "Материал"), 64)
    domain = T.r2_domain_text(str(raw.get("domain") or ""), 40)
    lines: list[str] = []
    if domain and domain not in ("—", ""):
        lines.append(f"{L.get('r34_source_label', 'Источник')}: {domain}")
    if kind == "confirmed":
        status = L.get("r34_status_confirmed", "Подтверждённые материалы")
        tone = T.SUCCESS
    elif kind == "review":
        status = L.get("r34_status_review", "Требует проверки")
        tone = T.WARNING
    else:
        status = L.get("r34_status_excluded", "Исключено как шум")
        tone = T.NEUTRAL_GRAY
    lines.append(f"{L.get('th_status', 'Статус')}: {status}")
    if kind == "review":
        lines.append(L.get("r34_requires_analyst", "Требуется аналитическая проверка"))
    return title, lines[:3], tone


def _r34_source_cards_page(prs, vm, ctx, *, title_key, subtitle_key, entries, kind, empty_key):
    L = vm["labels"]
    slide, top = _section(prs, ctx, L.get(title_key, "Раздел"), L.get(subtitle_key, ""))
    cards = [_r34_source_card_lines(e, L, kind) for e in (entries or [])]
    shown = 0
    for ctitle, lines, tone in cards:
        if int(top) + int(_R34_CARD_H) + _R34_NOTE_RESERVE > int(T.CONTENT_SAFE_BOTTOM):
            break
        top = T.card(slide, T.MARGIN, top, T.CONTENT_W, _R34_CARD_H, ctitle, lines, tone=tone)
        shown += 1
    if shown == 0:
        top = T.no_data_card(slide, top, L.get(empty_key, L.get("r34_no_confirmed", "Материалы отсутствуют.")))
    hidden = len(cards) - shown
    if hidden > 0:
        top = T.note(slide, top, L.get("r34_more_saved", "+ ещё {n} материалов сохранены для проверки.").format(n=hidden), "info")
    T.note(slide, top, L.get("r34_traceability_note", "Материалы показаны в клиент-safe формате и сохраняют трассируемость."), "disclaimer")


def _p_r34_map(prs, vm, ctx):
    L = vm["labels"]
    r34 = vm.get("r34Appendix") or {}
    slide, top = _section(prs, ctx, L.get("r34_map_title", "Карта раздела доказательств"), L.get("r34_map_subtitle", ""))
    cards = list(r34.get("navCards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    sections = [s for s in (r34.get("sections") or []) if str(s).strip()]
    if sections:
        top = T.bullets(slide, top, sections[:6], max_items=6, layout_warnings=ctx.layout_warnings)
    else:
        top = T.no_data_card(slide, top, L.get("r34_no_confirmed", "Материалы отсутствуют."))
    T.note(slide, top, L.get("r34_map_note", "Каждый слой можно рассматривать отдельно."), "info")


def _p_r34_confirmed_ru(prs, vm, ctx):
    r34 = vm.get("r34Appendix") or {}
    _r34_source_cards_page(
        prs, vm, ctx,
        title_key="r34_confirmed_ru_title", subtitle_key="r34_confirmed_ru_subtitle",
        entries=list(r34.get("confirmedRu") or []), kind="confirmed", empty_key="r34_no_confirmed",
    )


def _p_r34_review_ru(prs, vm, ctx):
    r34 = vm.get("r34Appendix") or {}
    _r34_source_cards_page(
        prs, vm, ctx,
        title_key="r34_review_ru_title", subtitle_key="r34_review_ru_subtitle",
        entries=list(r34.get("reviewRu") or []), kind="review", empty_key="r34_no_review",
    )


def _p_r34_confirmed_intl(prs, vm, ctx):
    r34 = vm.get("r34Appendix") or {}
    _r34_source_cards_page(
        prs, vm, ctx,
        title_key="r34_confirmed_intl_title", subtitle_key="r34_confirmed_intl_subtitle",
        entries=list(r34.get("confirmedIntl") or []), kind="confirmed", empty_key="r34_no_confirmed",
    )


def _p_r34_review_intl(prs, vm, ctx):
    r34 = vm.get("r34Appendix") or {}
    _r34_source_cards_page(
        prs, vm, ctx,
        title_key="r34_review_intl_title", subtitle_key="r34_review_intl_subtitle",
        entries=list(r34.get("reviewIntl") or []), kind="review", empty_key="r34_no_review",
    )


def _p_r34_excluded(prs, vm, ctx):
    L = vm["labels"]
    r34 = vm.get("r34Appendix") or {}
    exc = r34.get("excluded") or {}
    slide, top = _section(prs, ctx, L.get("r34_excluded_title", "Сводка по исключённым материалам"), L.get("r34_excluded_subtitle", ""))
    cards = list(exc.get("cards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    top = T.card(
        slide, T.MARGIN, top, T.CONTENT_W, Emu(1500000),
        L.get("r34_status_excluded", "Исключено как шум"),
        [
            L.get("r34_excluded_line_1", ""),
            L.get("r34_excluded_line_2", ""),
        ],
        tone=T.NEUTRAL_GRAY,
    )
    reasons = [str(x) for x in (exc.get("reasons") or []) if str(x).strip()]
    if ctx.internal and reasons:
        top = T.bullets(slide, top, reasons[:4], max_items=4, layout_warnings=ctx.layout_warnings)
    T.note(slide, top, L.get("r34_traceability_note", "Материалы сохраняют трассируемость."), "disclaimer")


def _p_r34_media(prs, vm, ctx):
    L = vm["labels"]
    r34 = vm.get("r34Appendix") or {}
    media = r34.get("media") or {}
    slide, top = _section(prs, ctx, L.get("r34_media_title", "Карточки медиа-доказательств"), L.get("r34_media_subtitle", ""))
    cards = list(media.get("cards") or [])
    if cards:
        top = T.metric_cards(slide, top, cards[:4], per_row=4)
    lines = [x for x in [L.get("r34_media_line_1", ""), L.get("r34_media_line_2", "")] if x]
    extra = [str(x) for x in (media.get("lines") or []) if str(x).strip()]
    body = (lines + extra)[:3]
    if body:
        top = T.card(slide, T.MARGIN, top, T.CONTENT_W, Emu(1550000), L.get("r34_media_card_title", "Логика отбора медиа"), body, tone=T.ACCENT)
    else:
        top = T.no_data_card(slide, top, L.get("r34_no_media", "Медиа-доказательства отсутствуют."))
    T.note(slide, top, L.get("r34_traceability_note", "Материалы сохраняют трассируемость."), "info")


def _p_r34_provenance(prs, vm, ctx):
    L = vm["labels"]
    slide, top = _section(prs, ctx, L.get("r34_provenance_title", "Происхождение источников и скриншотов"), L.get("r34_provenance_subtitle", ""))
    top = T.card(
        slide, T.MARGIN, top, T.CONTENT_W, Emu(2050000),
        L.get("r34_provenance_card_title", "Сводка о происхождении"),
        [
            L.get("r34_provenance_line_1", ""),
            L.get("r34_provenance_line_2", ""),
            L.get("r34_provenance_line_3", ""),
        ],
        tone=T.ACCENT,
    )
    T.note(slide, top, L.get("r34_traceability_note", "Материалы сохраняют трассируемость."), "disclaimer")


def _p_r34_risk(prs, vm, ctx):
    L = vm["labels"]
    r34 = vm.get("r34Appendix") or {}
    risk = r34.get("risk") or {}
    slide, top = _section(prs, ctx, L.get("r34_risk_title", "Обоснование риска по доказательствам"), L.get("r34_risk_subtitle", ""))
    top = T.metric_cards(
        slide, top,
        [
            _risk_card_value(risk.get("overallRiskLevel", "UNKNOWN"), L),
            {"label": L.get("r31_risk_reasoning_themes", "Ключевые темы"), "value": len(risk.get("topThemes") or [])},
            {"label": L.get("r31_risk_reasoning_actions", "Рекомендуемые действия"), "value": len(risk.get("recommendedActions") or [])},
        ],
        per_row=3,
    )
    top = T.card(
        slide, T.MARGIN, top, T.CONTENT_W, Emu(1550000),
        L.get("r34_risk_card_title", "Сводка обоснования"),
        [
            L.get("r34_risk_line_1", ""),
            L.get("r34_risk_line_2", ""),
        ],
        tone=T.ACCENT,
    )
    T.note(slide, top, L.get("r34_traceability_note", "Материалы сохраняют трассируемость."), "disclaimer")


def _p_r34_conclusion(prs, vm, ctx):
    L = vm["labels"]
    slide, top = _section(prs, ctx, L.get("r34_conclusion_title", "Итог приложения и трассируемость"), L.get("r34_conclusion_subtitle", ""))
    top = T.card(
        slide, T.MARGIN, top, T.CONTENT_W, Emu(2050000),
        L.get("r34_conclusion_card_title", "Заметка о трассируемости"),
        [
            L.get("r34_conclusion_line_1", ""),
            L.get("r34_conclusion_line_2", ""),
            L.get("r34_conclusion_line_3", ""),
        ],
        tone=T.ACCENT,
    )
    T.note(slide, top, L.get("r34_stored_internal", "Материалы сохранены для внутренней проверки."), "disclaimer")


def _solution(vm, idx: int) -> dict:
    sols = vm["offerBlock"]["solutions"]
    if idx < len(sols):
        return sols[idx]
    return {
        "title": f"Solution {idx + 1}", "subtitle": "", "objective": "", "price": "—",
        "duration": "—", "includedItems": [], "deliverables": [], "expectedResults": [],
        "workPlan": [], "pricingNotes": "",
    }


# ===========================================================================
# R7.4 — LexisNexis hybrid appendix (visual + parsed analytics)
# ===========================================================================

def _p_r74_lexis_intro(prs, vm, ctx):
    lx = vm.get("lexisHybrid") or {}
    if not lx.get("present"):
        return
    L = vm["labels"]
    labels = lx.get("labels") or {}
    slide, top = _section(prs, ctx, labels.get("introTitle", L.get("r74_intro_title", "Импортированный отчёт LexisNexis")))
    summary = lx.get("summary") or {}
    top = T.metric_cards(
        slide,
        top,
        [
            {"label": L.get("r74_docs", "Документы"), "value": int(summary.get("totalDocuments") or 0)},
            {"label": L.get("r74_signals", "Сигналы"), "value": int(summary.get("totalSignals") or 0)},
            {"label": L.get("r74_review_required", "Требуют проверки"), "value": int(summary.get("reviewRequired") or 0), "tone": T.WARNING},
            {"label": L.get("r74_parser_status", "Статус парсера"), "value": str(summary.get("parserStatus") or "unknown")},
        ],
        per_row=4,
    )
    top = T.card_auto(
        slide,
        T.MARGIN,
        top,
        T.CONTENT_W,
        labels.get("introTitle", L.get("r74_intro_title", "Импортированный отчёт LexisNexis")),
        [
            labels.get("introBody", L.get("r74_intro_body", "Импортированный отчёт LexisNexis добавлен как доказательный материал. Автоматический разбор выделил сигналы для проверки; выводы не являются юридическим заключением.")),
            str(summary.get("executiveSummaryClient") or ""),
        ],
        tone=T.ACCENT,
        min_h=620000,
    )
    T.note(slide, top, lx.get("legalSafeDisclaimer") or L.get("r74_legal_disclaimer", ""), "disclaimer")


def _p_r74_lexis_analytics(prs, vm, ctx):
    lx = vm.get("lexisHybrid") or {}
    if not lx.get("present"):
        return
    L = vm["labels"]
    labels = lx.get("labels") or {}
    slide, top = _section(prs, ctx, labels.get("analyticsTitle", "Аналитика импортированного отчёта"))
    docs = list(lx.get("documents") or [])
    if not docs:
        T.no_data_card(slide, top, "Данные отсутствуют.")
        return
    doc = docs[0]
    pa = doc.get("parsedAnalytics") or {}
    counts = pa.get("signalCounts") or {}
    top = T.metric_cards(
        slide,
        top,
        [
            {"label": L.get("r74_parser_status", "Статус парсера"), "value": str(pa.get("parserStatus") or "unknown")},
            {"label": L.get("r74_signals", "Сигналы"), "value": int(counts.get("totalSignals") or 0)},
            {"label": L.get("r74_review_required", "Требуют проверки"), "value": int(counts.get("reviewRequired") or 0), "tone": T.WARNING},
            {"label": L.get("r74_pages", "Страницы"), "value": int(doc.get("pageCount") or 0)},
        ],
        per_row=4,
    )
    rows = []
    for s in list(pa.get("signals") or [])[:10]:
        rows.append(
            [
                T.r2_truncate_cell_text(s.get("categoryLabelRu") or s.get("categoryLabelEn") or "—", 26),
                T.r2_truncate_cell_text(s.get("clientSafeFinding") or "—", 40),
                T.r2_truncate_cell_text(s.get("reviewStatus") or "review_required", 18),
                T.r2_truncate_cell_text(s.get("confidenceLabel") or "unknown", 12),
                f"p.{s.get('pageRef')}" if s.get("pageRef") else "—",
            ]
        )
    if rows:
        top = T.table(
            slide,
            top,
            ["Категория", "Сигнал", "Статус", "Уверенность", "Стр."],
            rows,
            col_widths=[0.20, 0.40, 0.16, 0.12, 0.12],
            max_rows=10,
            layout_warnings=ctx.layout_warnings,
        )
    else:
        top = T.no_data_card(slide, top, "Сигналы не выделены.")
    if pa.get("executiveSummaryClient"):
        T.note(slide, top, str(pa.get("executiveSummaryClient")), "info")


def _p_r74_lexis_visual_page(prs, vm, ctx, doc: dict, page: dict):
    labels = (vm.get("lexisHybrid") or {}).get("labels") or {}
    title = labels.get("visualTitle", "Страница импортированного документа")
    slide, top = _section(prs, ctx, title, f"{doc.get('fileName', 'LexisNexis')} · page {page.get('pageNumber', 0)}")
    raw = page.get("imageBytes")
    if raw:
        box_h = Emu(int(T.CONTENT_SAFE_BOTTOM) - int(top) - 250000)
        try:
            T._fit_picture_contain(slide, io.BytesIO(raw), int(T.MARGIN), int(top), int(T.CONTENT_W), int(box_h))
        except Exception:  # noqa: BLE001
            T.no_data_card(slide, top, "Не удалось визуализировать страницу.")
    else:
        T.no_data_card(slide, top, "Визуальная страница недоступна.")
    T.note(
        slide,
        Emu(int(T.CONTENT_SAFE_BOTTOM) - 220000),
        f"LexisNexis · page {page.get('pageNumber', 0)}",
        "disclaimer",
    )


# ===========================================================================
# orchestration
# ===========================================================================

def build_report_v3(
    report_json: dict,
    prs,
    data_root: str,
    warnings: list[str],
    audience: str = "internal",
    watermark_mode: str = "draft",
) -> None:
    vm, vm_warnings = build_view_model_v3(report_json, audience)
    warnings.extend(vm_warnings)
    L = vm["labels"]
    T.set_table_strings(L["showing_top"])
    T.set_note_strings(L.get("source_prefix", "Source: "))

    brand = vm["offerBlock"]["cover"]["brand"]
    meta_wm = vm["meta"].get("watermark")
    effective_wm = None if str(watermark_mode).lower() == "none" else meta_wm
    ctx = Ctx(brand=brand, watermark=effective_wm, internal=(str(audience).lower() != "client"))

    ru_pages = [
        _rb_r24("ru", _b_summary_r24, L["pg_audit_summary"]),
        _rb_r24("ru", _b_organic_r24, L["pg_organic_overview"]),
        _rb("ru", _b_results_r23, L["pg_top_results"]),
        _rb("ru", _b_themes, L["pg_neg_publications"]),
        _p_snapshots,
        _rb_r24("ru", _b_suggestions_r24, L["pg_suggestions"]),
        _rb_r24("ru", _b_related_r24, L["pg_related_queries"]),
        _p_ru_images_orion,
        _rb("ru", _b_videos, L["pg_videos"]),
        _rb("ru", _b_knowledge, L["pg_knowledge_block"]),
        _rb("ru", _b_wikipedia, L["pg_wikipedia"]),
        _rb("ru", _b_findings, L["pg_risk_findings"]),
        _rb("ru", _b_data_quality, L["pg_data_quality"]),
        _rb("ru", _b_recommended, L["pg_recommended"]),
        _rb("ru", _b_evidence, L["pg_evidence_appendix"]),
        _rb("ru", _b_conclusion, L["pg_interim_conclusion"]),
    ]
    intl_pages = [
        _rb_r24("intl", _b_summary_r24, L["pg_audit_summary"]),
        _rb_r24("intl", _b_organic_r24, L["pg_organic_overview"]),
        _rb("intl", _b_results_r23, L["pg_top_results"]),
        _rb("intl", _b_themes, L["pg_neg_themes"]),
        _rb_r24("intl", _b_suggestions_r24, L["pg_suggestions"]),
        _p_intl_media_r2,
        _rb("intl", _b_wiki_knowledge, L["pg_wiki_knowledge"]),
        _rb("intl", _b_findings, L["pg_risk_findings"]),
        _rb_r24("intl", _b_data_quality_r24, L["pg_data_quality"]),
        _rb_r24("intl", _b_conclusion_r24, L["pg_conclusion"]),
    ]
    offer_pages = [
        _p_offer_cover,
        _p_product_overview,
        _p_pricing_summary,
        lambda prs_, vm_, ctx_: _p_solution_objective(prs_, vm_, ctx_, 0),
        lambda prs_, vm_, ctx_: _p_solution_workplan(prs_, vm_, ctx_, 0),
        lambda prs_, vm_, ctx_: _p_solution_pricing(prs_, vm_, ctx_, 0),
        lambda prs_, vm_, ctx_: _p_solution_objective(prs_, vm_, ctx_, 1),
        lambda prs_, vm_, ctx_: _p_solution_workplan(prs_, vm_, ctx_, 1),
        lambda prs_, vm_, ctx_: _p_solution_pricing(prs_, vm_, ctx_, 1),
        lambda prs_, vm_, ctx_: _p_solution_objective(prs_, vm_, ctx_, 2),
        lambda prs_, vm_, ctx_: _p_solution_workplan(prs_, vm_, ctx_, 2),
        lambda prs_, vm_, ctx_: _p_solution_pricing(prs_, vm_, ctx_, 2),
        _p_process,
        _p_about,
    ]
    r31_pages = [
        _p_r31_appendix_cover,
        _p_r31_appendix_overview,
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_ru_confirmed_title",
            subtitle_key="r31_ru_confirmed_subtitle",
            entries=list((vm_.get("evidenceConfirmedRu") or {}).get("confirmed") or []),
        ),
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_ru_review_title",
            subtitle_key="r31_ru_review_subtitle",
            entries=list((vm_.get("evidenceReviewRu") or {}).get("rows") or []),
        ),
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_ru_excluded_title",
            subtitle_key="r31_ru_excluded_subtitle",
            entries=list((vm_.get("evidenceExcludedRu") or {}).get("rows") or []),
            include_reason=True,
        ),
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_intl_confirmed_title",
            subtitle_key="r31_intl_confirmed_subtitle",
            entries=list((vm_.get("evidenceConfirmedIntl") or {}).get("confirmed") or []),
        ),
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_intl_review_title",
            subtitle_key="r31_intl_review_subtitle",
            entries=list((vm_.get("evidenceReviewIntl") or {}).get("rows") or []),
        ),
        lambda prs_, vm_, ctx_: _p_r31_region_evidence(
            prs_,
            vm_,
            ctx_,
            title_key="r31_intl_excluded_title",
            subtitle_key="r31_intl_excluded_subtitle",
            entries=list((vm_.get("evidenceExcludedIntl") or {}).get("rows") or []),
            include_reason=True,
        ),
        _p_r31_media_overview,
        _p_r31_risk_reasoning_overview,
        _p_r31_risk_reasoning_by_region,
        _p_r31_appendix_conclusion,
    ]
    r34_pages = [
        _p_r34_map,
        _p_r34_confirmed_ru,
        _p_r34_review_ru,
        _p_r34_excluded,
        _p_r34_confirmed_intl,
        _p_r34_review_intl,
        _p_r34_media,
        _p_r34_provenance,
        _p_r34_risk,
        _p_r34_conclusion,
    ]
    diagnostics_pages = [_p_r32_provider_diagnostics] if ctx.internal else []
    lexis_pages: list[Callable] = []
    lx = vm.get("lexisHybrid") or {}
    if lx.get("present"):
        lexis_pages.extend([_p_r74_lexis_intro, _p_r74_lexis_analytics])
        for doc in list(lx.get("documents") or []):
            for page in list(doc.get("pages") or []):
                lexis_pages.append(
                    lambda prs_, vm_, ctx_, d=doc, p=page: _p_r74_lexis_visual_page(prs_, vm_, ctx_, d, p)
                )

    builders: list[Callable] = [
        _p_cover, _p_contents, _p_executive, _p_risk_matrix, _p_overview,
        *ru_pages,
        *intl_pages,
        _p_compliance_overview,
        _p_compliance_risk_types,
        _p_compliance_top_matches,
        _p_compliance_review_quality,
        _p_compliance_findings,
        *lexis_pages,
        *offer_pages,
        *r31_pages,
        *r34_pages,
        *diagnostics_pages,
    ]

    ctx.total = len(builders)
    for i, fn in enumerate(builders):
        ctx.page = i + 1
        ctx.layout_warnings = []
        try:
            fn(prs, vm, ctx)
            for w in ctx.layout_warnings:
                if w not in warnings:
                    warnings.append(w)
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"Slide {ctx.page} failed: {exc}")
            try:
                slide, top = _section(prs, ctx, L["section"])
                T.no_data_card(slide, top, L["section_render_error"].format(error=exc))
            except Exception:  # noqa: BLE001
                pass
