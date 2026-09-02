#!/usr/bin/env python3
"""Смок: рендерер не печатает пустоту — ни вместо буллета, ни вместо карточки.

На прогоне 76 страница «Россия — Википедия» пришла с двумя буллетами из трёх:
нарратив в 1387 знаков исчез без следа. Причина — `_clip_structured_bullet`
искал последнюю границу предложения `rfind`-ом по **всей** строке, а не в
пределах бюджета: у длинного текста она лежит за 900-м знаком, условие
«граница помещается» не выполнялось, и функция возвращала пустую строку.
`ctx.bullets` пропускал такой буллет через `continue` — без записи в
телеметрию, то есть правило §8 «выброшенное поднимает
CONTENT_DROPPED_BY_RENDERER» исполнять было нечем.

Проверяется оба конца:
  * длинный прозаический буллет сохраняет голову до границы предложения
    **внутри** бюджета и рисуется;
  * любая потеря — выброшенный, опустевший или укороченный буллет — попадает
    в телеметрию, и инспектор геометрии её видит. Проверка, не умеющая дать
    не-ноль, — не проверка.

Здесь же — карточный макет страницы проверки: без нарратива (страница-
продолжение) карточка статуса не рисуется вовсе, иначе на листе остаётся
озаглавленный пустой блок, а это дефект вёрстки по мерке приёмки.

Сеть, база и LibreOffice не нужны: презентация строится в памяти, инспектор
геометрии читает файл из временного каталога.

Запуск: python3 renderer/smoke_bullet_clip_keeps_head.py (нужен python-pptx)
"""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
from orion_golden_render.slides import _render_status_cards  # noqa: E402
from orion_golden_render.common import (  # noqa: E402
    CONTENT_BOTTOM,
    SLIDE_H,
    SLIDE_W,
    _Ctx,
    _clip_structured_bullet,
    get_layout_telemetry,
    reset_layout_telemetry,
)

REPO_ROOT = Path(__file__).resolve().parent.parent
INSPECTOR = REPO_ROOT / "scripts" / "inspect-first36-pptx-geometry.py"

BUDGET = 900

#: Нарратив страницы Википедии живого прогона: предложения по ~180 знаков,
#: последняя граница предложения — далеко за бюджетом.
LIVE_NARRATIVE = " ".join(
    [
        "Наличие статьи о проверяемом лице проверено через официальный поисковый API "
        "Википедии в русскоязычном разделе (ru.wikipedia.org); запрос: «Рашников Виктор "
        "Филиппович», проверка выполнена 17.08.2026.",
        "Результат проверки: статья о проверяемом лице найдена — «Рашников, Виктор "
        "Филиппович», https://ru.wikipedia.org/wiki/Рашников.",
        "В поисковой выдаче по контуру Россия зафиксировано 7 энциклопедических строк "
        "(ru.wikipedia.org, ru.ruwiki.ru, cyclowiki.org); из них о субъекте — 7.",
        "Существенных негативных или спорных формулировок в них не выявлено, а материалов "
        "об одноимённых лицах в контуре Россия не зафиксировано вовсе.",
        "Содержимое статьи попадает в панели знаний и ответы ИИ и воспринимается читателем "
        "как проверенная биография субъекта.",
        "Править статью может любой участник Википедии, поэтому искажение или "
        "недоброжелательное изменение быстро тиражируется поисковыми системами.",
        "Отдельные карточки выдачи содержат чувствительные формулировки, и их содержание "
        "отражено в темах повышенного внимания настоящего отчёта.",
        "Мы предлагаем контролировать корректность содержимого статьи: периодически сверять "
        "изложенные в ней факты и отслеживать историю правок, чтобы своевременно фиксировать "
        "искажения.",
    ]
)

#: Одно предложение длиннее бюджета: границы внутри бюджета нет вовсе.
ONE_LONG_SENTENCE = (
    "В открытых источниках описан длительный корпоративный конфликт вокруг актива, "
    "стороны публично обвиняют друг друга в давлении на суд и в использовании "
    "административного ресурса, публикации выходят в изданиях нескольких стран и "
    "повторяются с интервалом в несколько недель, дополняясь новыми подробностями сделки, "
    "составом участников, ссылками на решения судов, комментариями юристов и цитатами "
    "представителей сторон, что делает предложение бесконечным и лишает ножницы точки опоры "
    "внутри отведённого бюджета знаков, а также перечислением площадок, на которых материал "
    "повторяется: отраслевые издания, деловые агентства, региональные сайты и агрегаторы "
    "новостей, каждый со своей формулировкой заголовка и своим набором подробностей, "
    "которые в сумме и составляют этот единственный, ничем не прерываемый период, "
    "продолжающийся ровно столько, сколько нужно, чтобы превысить отведённый буллету "
    "бюджет знаков и не дать ножницам ни одной законной точки внутри него"
)

#: Буллет с заглушкой модели: «Что делать:.» выбрасывается чисткой внутри
#: ножниц. Это не потеря содержимого — обрубок и печатать было нечем.
STUB_BULLET = (
    "«Судебные споры» Найдены публикации по теме: 4. "
    "«Суд отказал в иске» — источник audit-it.ru Что делать:."
)

#: Структурный буллет темы: короткий, многострочный — его трогать нельзя.
THEME_BULLET = (
    "«Судебные споры»\n"
    "Найдены публикации по теме: 4.\n"
    "«Суд отказал в иске» — источник audit-it.ru\n"
    "Источники: audit-it.ru, x.com."
)

#: Самая длинная рекомендация страницы проверки — её карточка меряется 845 202.
LONG_ADVICE = (
    "Мы предлагаем сначала подтвердить принадлежность найденной статьи проверяемому лицу; "
    "если статья о нём — контролировать корректность её содержимого, если о другом лице — "
    "рассмотреть создание отдельной нейтральной биографической статьи."
)

#: Каркасный футнот шаблона.
METHODOLOGY_NOTE = (
    "Наличие статьи проверяется прямым запросом к официальному API Википедии в каждом "
    "языковом разделе; строки поисковой выдачи приводятся отдельно от результата проверки."
)

failures: list[str] = []

#: Счётчик выполненных проверок — раннер обязан видеть, сколько их было.
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)


def draw_bullets(items: list[str], y: int = 1_200_000, page: int = 30) -> tuple[Any, list[dict[str, Any]]]:
    """Нарисовать буллеты на пустой странице и вернуть презентацию с телеметрией."""
    reset_layout_telemetry()
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, page, 48)
    ctx.bullets(items, y, max_items=9, max_chars=BUDGET)
    return prs, get_layout_telemetry()


def drawn_text(prs: Any) -> str:
    return "\n".join(
        sh.text_frame.text
        for sh in prs.slides[0].shapes
        if getattr(sh, "has_text_frame", False)
    )


def bullet_entries(telemetry: list[dict[str, Any]], page: int = 30) -> list[dict[str, Any]]:
    return [
        e
        for e in telemetry
        if e.get("page") == page and str(e.get("name") or "").startswith("orion_bullets")
    ]


def inspector_codes(prs: Any, telemetry: list[dict[str, Any]]) -> list[str]:
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "page.pptx"
        prs.save(str(pptx_path))
        (Path(tmp) / "layout-telemetry.json").write_text(
            json.dumps({"version": "orion-layout-telemetry-v1", "entries": telemetry}, ensure_ascii=False),
            encoding="utf-8",
        )
        out = subprocess.run(
            [sys.executable, "-X", "utf8", str(INSPECTOR), str(pptx_path)],
            capture_output=True,
            encoding="utf-8",
            check=False,
        )
        report = json.loads(out.stdout or "{}")
    return [str(c.get("code")) for c in report.get("clipping") or []]


def main() -> int:
    # --- К1. Живое состояние: граница предложения лежит за бюджетом ----------
    clipped = _clip_structured_bullet(LIVE_NARRATIVE, BUDGET)
    check(
        "К1а: непустой вход не даёт пустого выхода",
        bool(clipped.strip()),
        f"вход {len(LIVE_NARRATIVE)} знаков, выход {len(clipped)}",
    )
    check(
        "К1б: сохранена голова текста, а не его середина",
        bool(clipped.strip()) and LIVE_NARRATIVE.startswith(clipped.strip()[:200]),
        f"начало выхода: {clipped[:60]!r}",
    )
    check(
        "К1в: выход держится в бюджете и кончается целым предложением",
        len(clipped) <= BUDGET and clipped.rstrip().endswith((".", "!", "?", "…")),
        f"длина {len(clipped)}, хвост {clipped.rstrip()[-24:]!r}",
    )

    # --- К2. Границы предложения в бюджете нет вовсе -------------------------
    one = _clip_structured_bullet(ONE_LONG_SENTENCE, BUDGET)
    check(
        "К2: предложение длиннее бюджета тоже не исчезает",
        bool(one.strip()) and len(one) <= BUDGET,
        f"вход {len(ONE_LONG_SENTENCE)}, выход {len(one)}",
    )

    # --- К3. Короткий структурный буллет не тронут ---------------------------
    theme = _clip_structured_bullet(THEME_BULLET, BUDGET)
    check(
        "К3: структурный буллет в пределах бюджета проходит без изменений",
        theme == THEME_BULLET,
        f"строк на входе {THEME_BULLET.count(chr(10)) + 1}, на выходе {theme.count(chr(10)) + 1}",
    )

    # --- К4. Буллет рисуется, а укороченный текст слышен --------------------
    prs, telemetry = draw_bullets([LIVE_NARRATIVE])
    text = drawn_text(prs)
    check(
        "К4а: длинный буллет нарисован на странице",
        "официальный поисковый API" in text,
        f"нарисовано {len(text)} знаков",
    )
    entries = bullet_entries(telemetry)
    check(
        "К4б: потеря хвоста записана в телеметрию",
        len(entries) == 1 and entries[0].get("clipped") is True,
        f"записей {len(entries)}: {entries[0] if entries else '—'}",
    )
    check(
        "К4в: инспектор геометрии видит потерю",
        bool(inspector_codes(prs, telemetry)),
        f"коды: {', '.join(sorted(set(inspector_codes(prs, telemetry)))) or '—'}",
    )

    # --- К5. Буллет, влезающий целиком, телеметрию не поднимает -------------
    prs_ok, telemetry_ok = draw_bullets([THEME_BULLET])
    check(
        "К5а: короткий буллет нарисован",
        "«Судебные споры»" in drawn_text(prs_ok),
        "",
    )
    check(
        "К5б: без потери телеметрия молчит",
        bullet_entries(telemetry_ok) == [],
        f"записей {len(bullet_entries(telemetry_ok))}",
    )

    # --- К6. Контрольный дефект: по высоте буллет не влез -------------------
    # Место под содержимое исчерпано: буллет не рисуется вовсе, и это обязано
    # быть слышно — страница без буллета при чистой телеметрии и есть тихая
    # потеря, из-за которой дефект жил при зелёном контуре.
    prs_bad, telemetry_bad = draw_bullets([THEME_BULLET], y=CONTENT_BOTTOM - 200_000)
    entries_bad = bullet_entries(telemetry_bad)
    dropped = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_bad
    )
    check(
        "К6а: не влезший буллет не нарисован",
        "«Судебные споры»" not in drawn_text(prs_bad),
        "",
    )
    check(
        "К6б: телеметрия сообщает о выброшенном буллете",
        len(entries_bad) == 1 and dropped > 0,
        f"записей {len(entries_bad)}, потеряно {dropped}",
    )
    check(
        "К6в: инспектор геометрии даёт CONTENT_DROPPED_BY_RENDERER",
        "CONTENT_DROPPED_BY_RENDERER" in inspector_codes(prs_bad, telemetry_bad),
        f"коды: {', '.join(sorted(set(inspector_codes(prs_bad, telemetry_bad)))) or '—'}",
    )

    # --- К7. Страница-продолжение не рисует карточку без текста ------------
    # Нарратив принадлежит первой странице блока, у продолжения его нет. Пока
    # карточка рисовалась всегда, на листе оставался озаглавленный пустой блок —
    # ровно то, что приёмка считает дефектом вёрстки.
    reset_layout_telemetry()
    prs_cont = Presentation()
    prs_cont.slide_width = Emu(SLIDE_W)
    prs_cont.slide_height = Emu(SLIDE_H)
    ctx_cont = _Ctx(prs_cont, 26, 48)
    _render_status_cards(
        ctx_cont,
        {"methodologyNote": "Методология страницы."},
        "Россия — Википедия (продолжение 2/2)",
        "",
        ["Строка выдачи — ru.wikipedia.org"],
        status_title="Результат проверки",
        bullets_as_card=False,
    )
    cont_text = drawn_text(prs_cont)
    check(
        "К7а: без нарратива карточка статуса не рисуется",
        "Результат проверки" not in cont_text,
        f"нарисовано: {cont_text[:80]!r}",
    )
    check(
        "К7б: строки выдачи и методология на продолжении печатаются",
        "Строка выдачи — ru.wikipedia.org" in cont_text and "Методология страницы." in cont_text,
        "",
    )

    # --- К8. Чистка заглушек модели — не потеря содержимого ----------------
    # Ножницы выбрасывают обрубки вроде «Что делать:.» ещё до счёта бюджета.
    # Пока разница длин считалась от сырого входа, косметическая чистка
    # поднимала `text-clipping` CRITICAL — на живых GPT-страницах приёмка
    # краснела бы там, где ничего не потеряно.
    prs_stub, telemetry_stub = draw_bullets([STUB_BULLET])
    check(
        "К8а: заглушка модели не напечатана",
        "Что делать:." not in drawn_text(prs_stub),
        "",
    )
    check(
        "К8б: чистка заглушки не считается потерей содержимого",
        bullet_entries(telemetry_stub) == [],
        f"записей {len(bullet_entries(telemetry_stub))}: {bullet_entries(telemetry_stub)}",
    )

    # --- К9. Футнот методологии печатается или слышно, что не напечатан -----
    # Резерв под рекомендацию и футнот держит место на полном листе; если
    # страница всё же не оставила места, пропуск обязан попасть в телеметрию —
    # молча терять методологию нельзя, ради этого шаг и делался.
    reset_layout_telemetry()
    prs_full = Presentation()
    prs_full.slide_width = Emu(SLIDE_W)
    prs_full.slide_height = Emu(SLIDE_H)
    ctx_full = _Ctx(prs_full, 25, 48)
    _render_status_cards(
        ctx_full,
        {
            "actions": [{"label": LONG_ADVICE}],
            "methodologyNote": METHODOLOGY_NOTE,
            "sourceNote": "Источники — ru.wikipedia.org, en.wikipedia.org.",
        },
        "Россия — Википедия",
        LIVE_NARRATIVE[:820],
        [f"Строка выдачи {i} — ru.wikipedia.org" for i in range(4)],
        status_title="Результат проверки",
        bullets_as_card=False,
    )
    full_text = drawn_text(prs_full)
    check(
        "К9а: на полном листе футнот методологии напечатан",
        METHODOLOGY_NOTE[:40] in full_text,
        f"нарисовано {len(full_text)} знаков",
    )
    check(
        "К9б: рекомендация на полном листе тоже напечатана",
        LONG_ADVICE[:40] in full_text,
        "",
    )
    footnote_entries = [
        e for e in get_layout_telemetry() if str(e.get("name") or "").startswith("orion_footnote")
    ]
    check(
        "К9в: напечатанный футнот телеметрию не поднимает",
        footnote_entries == [],
        f"записей {len(footnote_entries)}",
    )

    # Контрольный дефект: карточки съели лист, футноту места нет.
    reset_layout_telemetry()
    prs_squeezed = Presentation()
    prs_squeezed.slide_width = Emu(SLIDE_W)
    prs_squeezed.slide_height = Emu(SLIDE_H)
    ctx_squeezed = _Ctx(prs_squeezed, 25, 48)
    _render_status_cards(
        ctx_squeezed,
        {
            "actions": [{"label": LONG_ADVICE}],
            "methodologyNote": METHODOLOGY_NOTE,
        },
        "ОАЭ — Википедия",
        LIVE_NARRATIVE,
        [ONE_LONG_SENTENCE[:700], ONE_LONG_SENTENCE[:700], ONE_LONG_SENTENCE[:700], ONE_LONG_SENTENCE[:700]],
        status_title="Статус сбора",
        bullets_as_card=True,
    )
    squeezed_text = drawn_text(prs_squeezed)
    squeezed_entries = [
        e for e in get_layout_telemetry() if str(e.get("name") or "").startswith("orion_footnote")
    ]
    check(
        "К9г: не напечатанный футнот записан в телеметрию",
        (METHODOLOGY_NOTE[:40] in squeezed_text) or len(squeezed_entries) == 1,
        f"футнот на листе: {METHODOLOGY_NOTE[:40] in squeezed_text}, записей {len(squeezed_entries)}",
    )

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
