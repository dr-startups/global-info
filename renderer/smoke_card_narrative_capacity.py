#!/usr/bin/env python3
"""Смок: объявленная ёмкость абзаца карточной страницы подтверждается отрисовкой.

`CARD_NARRATIVE_CHAR_BUDGET` в `template-registry.ts` — число померенное, а не
назначенное: столько знаков построитель отдаёт карточке «Результат проверки».
Пока замер существовал только записью в комментарии, его никто не повторял, а
`content_card` режет невлезшее **молча**: в её теле нет ни одного вызова
`record_text_layout` и ни одного выставления `droppedLines` (проверено). Значит,
правка рендерера, снижающая ёмкость, сделала бы объявленный бюджет
несуществующим, не уронив ни телеметрию, ни приёмку.

Проверка зовёт **сам** `_render_status_cards` — тот код, который рисует
карточную страницу в деке, — и сравнивает длину написанного в текстовое поле с
длиной поданного. Аргументы карточки (`min_h`, `max_h`, `padding`, кегли) сюда
не копируются намеренно: копия — это второй ответ, и правка `max_h` в рендерере
оставила бы проверку зелёной при упавшей ёмкости.

Число читается из TypeScript: второго ответа о ёмкости здесь заводить нельзя,
а ненайденная константа — отказ, а не пропуск.

**Слабость, которую надо знать:** ёмкость меряется тем же инструментом, который
её создаёт. Смок ловит расхождение объявленного бюджета с рендерером, но не
ошибку самого рендерера; второе мнение о карточной странице даёт растровая
проверка и не заменяется этой.

Сеть, база и LibreOffice не нужны: презентация строится в памяти.

Запуск: python3 renderer/smoke_card_narrative_capacity.py (нужен python-pptx)
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
from smoke_ts_constants import ts_int  # noqa: E402
from orion_golden_render.common import SLIDE_H, SLIDE_W, _Ctx, _safe  # noqa: E402
from orion_golden_render.slides import _render_status_cards  # noqa: E402

REPO_ROOT = Path(__file__).resolve().parent.parent
REGISTRY_TS = (
    REPO_ROOT
    / "src/modules/digital-profile/orion-golden/deck-sections/template-registry.ts"
)

failures: list[str] = []
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)


CARD_NARRATIVE_CHAR_BUDGET = ts_int(
    REGISTRY_TS, r"export const CARD_NARRATIVE_CHAR_BUDGET = (\d+);"
)

#: Абзац карточной страницы эталона-72 (`p13_ru_wikipedia`) — проза, которую
#: построитель действительно отдаёт этой карточке.
PROSE_CLIENT = (
    "Результат отдельной проверки наличия статьи в Википедии для этого контура в "
    "отчёте отсутствует, поэтому вывод о наличии или отсутствии статьи не делается. "
    "В поисковой выдаче по контуру Россия зафиксирована 1 энциклопедическая строка "
    "(ru.wikipedia.org); принадлежность материала проверяемому лицу не подтверждена. "
    "В заголовках зафиксированных строк выдачи существенных негативных или спорных "
    "формулировок не выявлено."
)

#: Та же проза без единой латинской буквы: кириллица шире, и ёмкость на ней ниже.
PROSE_CYRILLIC = (
    "Наличие статьи о субъекте в энциклопедии и связанных справочных материалов "
    "проверено по обоим контурам сбора, и результат проверки описан словами. "
    "Открытые источники повторяют один сюжет несколько месяцев подряд и добавляют "
    "подробности состава участников, сроков рассмотрения и позиции сторон спора в "
    "судах нескольких инстанций."
)

#: Неразрывный токен длиной 38 знаков — транслитерация из адреса материала.
#: Такие слова рендерер переносить не умеет и тратит на них остаток строки.
UNBREAKABLE_TOKEN = "byvshijpartneroligarhovnahodkatranspor"
PROSE_TOKENS = " ".join([UNBREAKABLE_TOKEN] * 40)

SAMPLES = [
    ("клиентская проза страницы", PROSE_CLIENT),
    ("проза без латиницы", PROSE_CYRILLIC),
    (f"неразрывные токены по {len(UNBREAKABLE_TOKEN)} знаков", PROSE_TOKENS),
]


def material(sample: str, length: int) -> str:
    """Ровно `length` знаков материала из повторов образца.

    Повторы склеиваются через `strip()`: `_safe` схлопывает двойной пробел, и
    тогда сравнение длин падает не на обрезке, а на уборке пробела.
    """
    one = sample.strip()
    return " ".join([one] * (length // len(one) + 2))[:length].strip()


def drawn_body(text: str) -> str:
    """Текст, который карточная страница написала в тело своей первой карточки."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, 1, 1)
    _render_status_cards(
        ctx,
        {},
        "Россия — Википедия",
        text,
        [],
        status_title="Результат проверки",
        bullets_as_card=False,
    )
    # Надписи страницы идут в порядке отрисовки: заголовок страницы, заголовок
    # карточки, тело. Брать «последнюю фигуру» нельзя: когда тело выброшено
    # целиком, последним остаётся заголовок карточки, и проверка прочитала бы
    # его как нарисованный текст. Позиция, а не `id()`: python-pptx создаёт
    # обёртку фигуры заново на каждом обходе, и множество идентификаторов
    # «до отрисовки» на следующем обходе не совпадает само с собой.
    written = [
        sh.text_frame.text
        for sh in prs.slides[0].shapes
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text
    ]
    return written[2] if len(written) >= 3 else ""


def fits(text: str) -> bool:
    """Нарисовано столько же, сколько подано.

    Сравниваются **длины**: рендерер разбивает абзац на параграфы, и в
    прочитанном тексте появляются переводы строк, которых во входе не было.
    Посимвольное равенство находило бы «границу» на первом же переносе.
    """
    return len(drawn_body(text)) >= len(_safe(text)) - 2


def capacity(sample: str, ceiling: int = 3000) -> int:
    """Наибольшая длина материала, которая доезжает до листа целиком."""
    lo, hi, best = 100, ceiling, 0
    while lo <= hi:
        mid = (lo + hi) // 2
        if fits(material(sample, mid)):
            best = mid
            lo = mid + 1
        else:
            hi = mid - 1
    return best


def main() -> int:
    print(f"# CARD_NARRATIVE_CHAR_BUDGET (из template-registry.ts) = {CARD_NARRATIVE_CHAR_BUDGET}")
    for name, sample in SAMPLES:
        text = material(sample, CARD_NARRATIVE_CHAR_BUDGET)
        # Материал обязан пережить уборку текста. Идентификаторы и служебные
        # слова рендерер вырезает (`FORBIDDEN`), и на вырезанном материале
        # проверка «нарисовано столько же» зеленеет вакуумно: нарисовано ноль
        # из нуля. На этом замер ёмкости однажды и соврал.
        check(
            f"материал переживает уборку текста — {name}",
            len(_safe(text)) >= CARD_NARRATIVE_CHAR_BUDGET - 2,
            f"подано {len(text)}, после уборки {len(_safe(text))}",
        )
        measured = capacity(sample)
        check(
            f"объявленный бюджет доезжает до листа — {name}",
            fits(text),
            f"нарисовано {len(drawn_body(text))} из {len(_safe(text))}; "
            f"померенная ёмкость {measured} знаков",
        )
        # Число печатается, чтобы следующий читал ёмкость, а не верил зелёному
        # коду возврата.
        print(f"#   ёмкость на образце «{name}»: {measured} знаков")

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
