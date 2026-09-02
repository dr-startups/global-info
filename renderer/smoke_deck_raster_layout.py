#!/usr/bin/env python3
"""Смок: проверка вёрстки по растру (ADR-0007).

Инспектор геометрии зелёный на деке, которую человек называет сломанной:
он считает переполнением только выход рамки за границы слайда, а текст,
не влезший в свой блок, рамку не двигает. Эта проверка смотрит на
отрисованную страницу и нашими метриками не пользуется вовсе.

Контрольный дефект обязателен. Проверка, которая на исправной деке даёт ноль,
но не умеет дать не-ноль, — не гейт: ровно так `overflow: 0` и держался.
Поэтому сначала синтетические страницы (заведомо чистая и заведомо сломанная),
и только потом настоящая дека, если она отрендерена.

Сеть и БД не нужны. Настоящие страницы — только если есть артефакты прогона.

Запуск: python3 renderer/smoke_deck_raster_layout.py
"""

from __future__ import annotations

import base64
import io
import json
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from PIL import Image, ImageDraw  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
from smoke_ts_constants import ts_int  # noqa: E402
# Пределы ячеек и слагаемые ёмкости второй таблицы выдачи берутся у смока
# ширин, а не объявляются здесь заново: вопрос «какая строка худшая законная»
# обязан иметь один ответ на оба прибора, иначе один из них однажды нарисует
# бывшую худшую строку и промолчит.
from smoke_search_table_layout import (  # noqa: E402
    SERP_ADDRESS_MAX_CHARS,
    HDR_SERP_EXTRA,
    SERP_FOUND_BY_MAX_CHARS,
    SERP_TITLE_MAX_CHARS,
    serp_capacity_terms,
)
sys.path.insert(0, str(Path(__file__).resolve().parent))
from orion_golden_render.common import FONT, TYPE_SCALE_PT  # noqa: E402
from deck_raster_layout import (  # noqa: E402
    INK_BOTTOM,
    SLIDE_H,
    STAGE_BOTTOM,
    _ink_row_counts,
    check_pages,
    detect_furniture_top,
)

W, H = 1844, 1152
FURNITURE_RULE_Y = int(H * 0.934)
FURNITURE_TEXT_Y = int(H * 0.950)
INK_BOTTOM_PX = int(INK_BOTTOM / SLIDE_H * H)

#: Полоса, в которой лежит нижняя кромка белой сцены вместе со своей тенью:
#: от низа сцены до последнего ряда тени. Числами не задаётся — обе границы
#: выведены из той же геометрии, что и `INK_BOTTOM`.
#: Замер на эталоне 25.08: у здоровых страниц ряды 962…968 несут 823–850
#: отсчётов чернил, ниже 969 — ноль; у испорченной стр. 17 там 46–71.
STAGE_EDGE_FROM_PX = int(STAGE_BOTTOM / SLIDE_H * H)
STAGE_EDGE_TO_PX = INK_BOTTOM_PX - 1
#: Полноширинная линия: при шаге выборки 2 это ~920 отсчётов на всю ширину.
STAGE_EDGE_INK = 600

failures: list[str] = []

#: Счётчик выполненных проверок — раннер обязан видеть, сколько их было.
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)



def count_first_level(runs_by_page: dict) -> dict:
    """Фигуры, несущие крупнейший кегль страницы, — её первый уровень."""
    out: dict = {}
    for page, runs in runs_by_page.items():
        if not runs:
            continue
        top = max(r[0] for r in runs)
        out[page] = {r[1] for r in runs if r[0] == top}
    return out


#: Цвета шкалы степени в отрисованной карточке: закрашенное деление и пустое.
BAR_COLORS = {"red": (0xE1, 0x3D, 0x3D), "amber": (0xE3, 0x8A, 0x24), "cold": (0xD8, 0xE3, 0xDA)}
#: Деление шириной 120 000 EMU — около 18 px на растре; глиф уже вдвое.
BAR_MIN_RUN_PX = 12
#: Бейджи стоят в правой четверти листа — там же и слово ступени, и её шкала.
BADGE_COLUMN_FROM = 0.75


def render_matrix_page(cards: list[dict]) -> bytes | None:
    """Одна страница матрицы, отрисованная настоящим рендерером; None — нечего смотреть.

    Шкалу делений смотрим на карточках, которых в эталоне 72 нет вовсе
    («Требует подтверждения», «Нет данных»), поэтому страница строится здесь.

    Растр принимается только от LibreOffice: при её отсутствии рендерер рисует
    PDF запасным кодом, и проверять вёрстку по нему — значит проверять не ту
    вёрстку. Такой прогон объявляется пропуском, а не считается пройденным.
    """
    from orion_golden_render import render_orion_golden  # noqa: PLC0415

    payload = {
        "subjectName": "Тест",
        "deckManifest": {
            "toc": [],
            "sectionPageRanges": [],
            "finalSlides": [
                {
                    "slideKey": "p04_risk_dashboard",
                    "template": "orion_golden_risk_matrix_grid",
                    "title": "Матрица комплаенс-рисков",
                    "pageNumber": 1,
                    "totalPageCount": 1,
                    "keyFindings": cards,
                }
            ],
        },
    }
    out = render_orion_golden(payload)
    pages = out.get("pages") or []
    if not pages or out.get("pdfExportMode") != "libreoffice":
        return None
    return base64.b64decode(pages[0]["contentBase64"])


#: Сколько знаков построитель отдаёт одному объяснению панели
#: (`SIDEBAR_HIGHLIGHT_BUDGET` реестра шаблонов). Число читается из TypeScript:
#: второго ответа о бюджете здесь заводить нельзя.
SIDEBAR_HIGHLIGHT_BUDGET = ts_int(
    Path(__file__).resolve().parent.parent
    / "src/modules/digital-profile/orion-golden/deck-sections/template-registry.ts",
    r"export const SIDEBAR_HIGHLIGHT_BUDGET = (\d+);",
)


def _padded(text: str, length: int) -> str:
    """Фраза ровно объявленной длины — худший законный случай, а не типичный."""
    tail = " Формулировка приведена дословно по строке выдачи и не пересказывается."
    while len(text) < length:
        text = f"{text}{tail}"
    return text[:length].rstrip()


#: Два объяснения негативной подсказки — то, что построитель кладёт на панель
#: поверхности при первом же прогоне с негативной строкой.
PANEL_REASONS = [
    _padded(
        "Негативная формулировка в подсказке: запрос связывает субъекта с уголовным делом "
        "и хищением средств; строка показана поисковой системой в первой десятке и "
        "повторяется на обоих контурах сбора. Страница не читалась в этом прогоне.",
        SIDEBAR_HIGHLIGHT_BUDGET,
    ),
    _padded(
        "Негативная формулировка в связанном запросе: субъект упоминается вместе со "
        "словами о банкротстве и претензиях кредиторов; запрос показан над органической "
        "выдачей. Оценка сделана по формулировке запроса.",
        SIDEBAR_HIGHLIGHT_BUDGET,
    ),
]
PANEL_MORE_SIGNALS = 3


def render_surface_panel_page() -> tuple[bytes, bytes] | None:
    """Панель поверхности с блоком «Почему выделено»; None — растр не от LibreOffice.

    Блок рисует `_sidebar_analysis` при `sidebarMode == "adverse_explanation"`,
    и на панели поверхности (подсказки, связанные запросы) он **не рисовался
    ни разу**: у всех восьми таких слайдов обоих эталонов объяснений нет —
    непустые `highlightExplanations` есть только у снимка выдачи и сетки
    изображений. Значит, первым зрителем этой геометрии оказался бы клиент
    первого же живого прогона с негативной подсказкой, а прогон этот платный.

    Панель без картинки вырождается в полностраничные карточки
    (`_render_visual_with_sidebar` → `_render_analysis_cards_full_width`), и
    сайдбара на такой странице нет вовсе. Поэтому ассет обязателен: без него
    проверка смотрела бы на другой макет.
    """
    from orion_golden_render import render_orion_golden  # noqa: PLC0415

    buf = io.BytesIO()
    Image.new("RGB", (1200, 760), (0xF2, 0xF5, 0xF3)).save(buf, format="PNG")
    payload = {
        "subjectName": "Тест",
        "assets": [
            {
                "assetRef": "surface_panel_adverse",
                "kind": "visual",
                "title": "Подсказки",
                "imageData": base64.b64encode(buf.getvalue()).decode("ascii"),
            }
        ],
        "deckManifest": {
            "toc": [],
            "sectionPageRanges": [],
            "finalSlides": [
                {
                    "slideKey": "p11_ru_suggestions_yandex",
                    "sectionKey": "RU_PROFILE",
                    "template": "orion_golden_surface_panel",
                    "title": "Россия — подсказки Яндекса: есть негативные формулировки",
                    "pageNumber": 1,
                    "totalPageCount": 1,
                    "assetRefs": ["surface_panel_adverse"],
                    "visualAnalysis": {
                        "sidebarMode": "adverse_explanation",
                        "headlineConclusion": (
                            "На панели — 10 подсказок: 3 несут негативную формулировку, "
                            "остальные нейтральны."
                        ),
                        "whatIsVisible": "Подсказки показывают, что чаще всего ищут о субъекте.",
                        "clientMeaning": (
                            "Негативная подсказка видна каждому, кто набирает имя в строке поиска."
                        ),
                        "highlightExplanations": [
                            {"clientReason": PANEL_REASONS[0], "frameTone": "red"},
                            {"clientReason": PANEL_REASONS[1], "frameTone": "amber"},
                        ],
                        "moreSignalsCount": PANEL_MORE_SIGNALS,
                        "recommendedActions": [
                            "Сверить формулировки подсказок с профилем субъекта."
                        ],
                        "provenanceLabel": (
                            "Источник — поисковая выдача: у показанных элементов "
                            "нет отдельных адресов."
                        ),
                    },
                }
            ],
        },
    }
    out = render_orion_golden(payload)
    pages = out.get("pages") or []
    if not pages or out.get("pdfExportMode") != "libreoffice" or not out.get("pdfBase64"):
        return None
    return base64.b64decode(pages[0]["contentBase64"]), base64.b64decode(out["pdfBase64"])


#: Вводный абзац листа второй таблицы, заведомо длиннее своего потолка.
#: Ёмкость листа выведена от **объявленного** верха таблицы (заголовок плюс
#: потолок абзаца), поэтому лист с коротким абзацем начинался бы выше — и
#: проверялся бы не худший лист, а удобный.
SERP_EXTRA_INTRO_OVERFLOW = " ".join(
    f"Предложение номер {n} этого абзаца написано нарочно длинным, чтобы вводный "
    "текст страницы упёрся в объявленный потолок мерой, а не поместился в него."
    for n in range(1, 13)
)


def render_serp_extra_table_page(row_count: int) -> tuple[bytes, bytes] | None:
    """Лист второй таблицы выдачи; None — растр не от LibreOffice.

    Страниц второй таблицы выдачи («Найдено по дополнительным запросам») не
    рисовал **ни один эталон**: у корпуса отчёта 72 дополнительных запросов нет
    вовсе, а золотой кейс до рендерера не доходит. Её ёмкость выведена
    арифметикой — бюджет листа, поделённый на худшую законную строку, — и с
    нарисованной страницей эту арифметику до сих пор не сверял никто. Первым
    зрителем геометрии оказался бы клиент первого же прогона с дополнительным
    запросом, а прогон этот платный.

    Строка берётся предельная, а не правдоподобная: адрес, заголовок и запрос
    на своих пределах, написанные самым широким знаком 9 pt. Именно из такой
    строки ёмкость и выведена — правдоподобная строка проверяла бы не то.
    """
    from orion_golden_render import render_orion_golden  # noqa: PLC0415

    worst_row = [
        "Ю" * SERP_ADDRESS_MAX_CHARS,
        "Ю" * SERP_TITLE_MAX_CHARS,
        "Ю" * SERP_FOUND_BY_MAX_CHARS,
        "Официальный сайт / госресурс",
        "Нежелательный",
    ]
    payload = {
        "subjectName": "Тест",
        "deckManifest": {
            "toc": [],
            "sectionPageRanges": [],
            "finalSlides": [
                {
                    "slideKey": "p12_ru_serp_extra_queries",
                    "sectionKey": "RU_PROFILE",
                    "template": "orion_golden_search_table",
                    "title": "Россия — найдено по дополнительным запросам",
                    "pageNumber": 1,
                    "totalPageCount": 1,
                    "narrative": SERP_EXTRA_INTRO_OVERFLOW,
                    "table": {
                        "headers": list(HDR_SERP_EXTRA),
                        "rows": [list(worst_row) for _ in range(row_count)],
                    },
                }
            ],
        },
    }
    out = render_orion_golden(payload)
    pages = out.get("pages") or []
    if not pages or out.get("pdfExportMode") != "libreoffice" or not out.get("pdfBase64"):
        return None
    return base64.b64decode(pages[0]["contentBase64"]), base64.b64decode(out["pdfBase64"])


def stage_edge_ink(png: bytes) -> tuple[int, int, tuple[int, int]]:
    """Чернила в полосе нижней кромки сцены, самый низкий ряд чернил и размер растра.

    Полоса кромки задана в пикселях модуля, то есть привязана к размеру
    эталонной страницы; размер возвращается, чтобы проверка могла его назвать,
    а не считать совпадение само собой разумеющимся.
    """
    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "page-01.png"
        path.write_bytes(png)
        rows, w, h = _ink_row_counts(path)
    return (
        max(rows[STAGE_EDGE_FROM_PX : STAGE_EDGE_TO_PX + 1], default=0),
        max((y for y, n in enumerate(rows) if n > 0), default=0),
        (w, h),
    )


def header_left_edges(pdf: bytes, headers: list[str]) -> dict[str, float]:
    """Левые края нарисованных заголовков таблицы — из текстового слоя PDF.

    Ширина колонки меряется по тому, что нарисовано, а не по долям, которые мы
    сами и объявили: разница левых краёв соседних заголовков и есть ширина
    колонки — отступ ячейки у обоих одинаков и сокращается.
    """
    with tempfile.TemporaryDirectory() as tmp:
        path = Path(tmp) / "page.pdf"
        path.write_bytes(pdf)
        import fitz  # noqa: PLC0415

        words = fitz.open(str(path))[0].get_text("words")
    edges: dict[str, float] = {}
    for header in headers:
        # Заголовок из двух слов ищется по первому: слова в текстовом слое
        # лежат по одному, а левый край строки задаёт именно первое.
        hits = [w[0] for w in words if w[4] == header.split()[0]]
        if hits:
            edges[header] = min(hits)
    return edges


def scan_badges(png: bytes) -> dict:
    """Что нарисовано в колонке бейджей: шкалы делений и палитра пикселей.

    Мерами рендерера не пользуемся вовсе — ищем то, что нарисовано: подряд
    идущие ряды, где стоят четыре и более широких пятна цвета делений. Слово
    ступени тоже цвет, поэтому пиксели считаются в том же проходе: по ним видно,
    что бейдж нарисован и без шкалы.
    """
    with Image.open(io.BytesIO(png)) as raw:
        im = raw.convert("RGB")
    w, h = im.size
    px = im.load()
    x_from = int(w * BADGE_COLUMN_FROM)
    pixels = {name: 0 for name in BAR_COLORS}

    def color_of(c) -> str | None:
        for name, t in BAR_COLORS.items():
            if all(abs(c[i] - t[i]) <= 26 for i in range(3)):
                return name
        return None

    groups: list[dict] = []
    for y in range(h):
        runs: list[str] = []
        x = x_from
        while x < w:
            kind = color_of(px[x, y])
            if kind is None:
                x += 1
                continue
            start = x
            while kind is not None:
                pixels[kind] += 1
                x += 1
                kind = color_of(px[x, y]) if x < w else None
            if x - start >= BAR_MIN_RUN_PX:
                # Цвет деления берётся по середине пятна: край размыт конвертацией.
                runs.append(color_of(px[(start + x) // 2, y]) or "cold")
        if len(runs) < 4:
            continue
        # Сама шкала и её тень идут одной группой: тень отстоит от делений на
        # несколько рядов и без склейки читалась бы второй шкалой.
        if groups and y - groups[-1]["lastY"] <= 12:
            group = groups[-1]
            group["lastY"] = y
            if len(runs) > len(group["widest"]):
                group["widest"] = runs
        else:
            groups.append({"y": y, "lastY": y, "widest": runs})

    return {
        "pixels": pixels,
        "bars": [
            {name: sum(1 for kind in g["widest"] if kind == name) for name in BAR_COLORS}
            for g in groups
        ],
    }


def make_page(path: Path, *, spill: bool = False, outside: bool = False) -> None:
    """Страница с обычной мебелью; при желании — с дефектом."""
    im = Image.new("RGB", (W, H), (255, 255, 255))
    d = ImageDraw.Draw(im)
    # Содержимое в отведённой области.
    d.rectangle([120, 120, W - 120, 400], fill=(51, 65, 85))
    # Мебель: линейка колонтитула и его текст.
    d.rectangle([90, FURNITURE_RULE_Y, W - 90, FURNITURE_RULE_Y + 2], fill=(120, 130, 150))
    d.rectangle([90, FURNITURE_TEXT_Y, 400, FURNITURE_TEXT_Y + 14], fill=(120, 130, 150))
    if spill:
        # Текст, вышедший за нижнюю границу своего блока.
        d.rectangle([120, INK_BOTTOM_PX + 8, 900, INK_BOTTOM_PX + 40], fill=(51, 65, 85))
    if outside:
        # Текст, вылезший за боковое поле.
        d.rectangle([2, 300, 40, 420], fill=(51, 65, 85))
    im.save(path)


def main() -> int:
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        clean = [d / f"page-{i:02d}.png" for i in range(1, 5)]
        for p in clean:
            make_page(p)

        report = check_pages(clean)
        check(
            "чистые страницы проходят",
            report.passed and not report.findings,
            f"проверено {report.pages_checked}, дефектов {len(report.findings)}",
        )
        check(
            "мебель найдена по самим страницам, а не задана числом",
            report.furniture_top_y is not None
            and abs(report.furniture_top_y - FURNITURE_RULE_Y) <= 2,
            f"верх мебели y={report.furniture_top_y}, ожидалось ~{FURNITURE_RULE_Y}",
        )

        # Контрольный дефект — тот самый случай, который сегодня проходит мимо
        # инспектора геометрии.
        broken = d / "page-05.png"
        make_page(broken, spill=True)
        r2 = check_pages(clean + [broken])
        spilled = [f for f in r2.findings if f.code == "TEXT_BELOW_CONTENT_AREA"]
        check(
            "текст ниже границы контента пойман",
            len(spilled) == 1 and spilled[0].page == "page-05.png",
            spilled[0].detail if spilled else "не пойман",
        )
        check("страница с дефектом не считается пройденной", not r2.passed)

        wide = d / "page-06.png"
        make_page(wide, outside=True)
        r3 = check_pages(clean + [wide])
        outside = [f for f in r3.findings if f.code == "TEXT_OUTSIDE_MARGINS"]
        check(
            "текст за боковым полем пойман",
            len(outside) == 1 and outside[0].page == "page-06.png",
            outside[0].detail if outside else "не пойман",
        )

        # Ноль не может быть ответом, когда проверять было нечего.
        empty = check_pages([])
        check(
            "пустой набор страниц не выдаётся за успех",
            not empty.passed and empty.pages_checked == 0,
            "0 дефектов на 0 страниц — это «не проверяли», а не «всё хорошо»",
        )
        check(
            "мебель не ищется по заголовкам в верхней части листа",
            detect_furniture_top([[1] * H]) is not None,
        )

    # Настоящая дека — если отрендерена.
    pages_dir = (
        Path(__file__).resolve().parent.parent
        / "baselines/report-72/artifacts/deck-sections/pages-png"
    )
    real = sorted(pages_dir.glob("page-*.png")) if pages_dir.is_dir() else []
    if not real:
        print("# SKIP растровая проверка эталонной деки — нет отрендеренных страниц")
    else:
        rep = check_pages(real)
        for f in rep.findings:
            print(f"    {f.page}: {f.code} — {f.detail}")
        check(
            "эталонная дека проходит растровую проверку",
            rep.passed,
            f"проверено страниц: {rep.pages_checked}, дефектов: {len(rep.findings)}",
        )

        # --- Нижняя кромка сцены видна на страницах таблицы выдачи -----------
        #
        # Таблица лежит на белой сцене. Если она объявила высоту меньше
        # нарисованной, LibreOffice тянет строку по содержимому, таблица едет
        # вниз и закрашивает кромку сцены собой — сцена «пропадает», а
        # содержимое уходит ниже неё. Прежний порог растровых ворот лежал
        # почти на десяток пикселей ниже кромки и такую страницу пропускал:
        # на эталоне 25.08 стр. 17 шла до y=972 при кромке 962, и ворота
        # молчали.
        #
        # Какая страница является таблицей выдачи — берётся из манифеста
        # (указатель, а не измерение); сама кромка ищется по растру.
        manifest_path = (
            Path(__file__).resolve().parent.parent
            / "baselines/report-72/artifacts/deck-sections/report-deck-manifest.json"
        )
        if not manifest_path.is_file():
            print("# SKIP кромка сцены на страницах выдачи — нет манифеста деки")
        else:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            serp_pages = {
                int(s["pageNumber"])
                for s in manifest.get("slides", [])
                if s.get("templateId") == "serp-table"
            }
            missing_edge: list[str] = []
            # Считаются **осмотренные** страницы, а не объявленные манифестом:
            # наборы расходятся (страницы не отрисованы, файл переименован), и
            # «ноль дефектов на ноль осмотренных» — это «не проверяли».
            inspected = 0
            for path in real:
                number = int(path.stem.split("-")[-1])
                if number not in serp_pages:
                    continue
                inspected += 1
                rows, _w, _h = _ink_row_counts(path)
                if max(rows[STAGE_EDGE_FROM_PX : STAGE_EDGE_TO_PX + 1], default=0) < STAGE_EDGE_INK:
                    lowest = max((y for y, n in enumerate(rows) if n > 0), default=0)
                    missing_edge.append(f"{path.name} (самый низкий ряд чернил y={lowest})")
            check(
                "нижняя кромка сцены видна на каждой странице таблицы выдачи",
                not missing_edge and inspected == len(serp_pages) and inspected > 0,
                f"осмотрено {inspected} из {len(serp_pages)} объявленных; "
                f"без кромки: {missing_edge or 'нет'}",
            )

    # --- Шкала делений принадлежит ступени, а не статусу ------------------
    #
    # «Требует подтверждения» — статус идентификации, а не степень риска: на
    # шкале ему места нет, в том числе визуального. Пока деления рисовались по
    # тону карточки, непроверенный материал стоял на середине шкалы (3 из 5).
    detail = (
        "7 свидетельств (5 негативных) в источниках audit-it.ru, x.com, m.sledst.org.\n"
        "В чём проблема: материалы связывают субъекта с судебными разбирательствами.\n"
        "Что делать: Проверить статусы дел и первоисточники до принятия решений."
    )
    steps_png = render_matrix_page(
        [
            {"headline": "Тема ступени", "status": "Высокий", "detail": detail, "tone": "danger"},
            {"headline": "Тема ступени", "status": "Средний", "detail": detail, "tone": "warn"},
        ]
    )
    if steps_png is None:
        print("# SKIP шкала делений — растр не получен от LibreOffice")
    else:
        bars = scan_badges(steps_png)["bars"]
        check(
            "карточка ступени рисует шкалу делений",
            len(bars) == 2,
            f"найдено шкал: {len(bars)}, ожидалось 2",
        )
        high, medium = (bars + [{}, {}])[:2]
        check(
            "«Высокий» — пять красных делений",
            high.get("red") == 5 and high.get("cold") == 0,
            f"красных {high.get('red')}, серых {high.get('cold')}",
        )
        check(
            "«Средний» — три янтарных из пяти",
            medium.get("amber") == 3 and medium.get("cold") == 2,
            f"янтарных {medium.get('amber')}, серых {medium.get('cold')}",
        )

        statuses_png = render_matrix_page(
            [
                {
                    "headline": "Тема статуса",
                    "status": "Требует подтверждения",
                    "detail": detail,
                    "tone": "warn",
                },
                {
                    "headline": "Нет подтверждённых тем",
                    "status": "Нет данных",
                    "detail": detail,
                    "tone": "warn",
                },
            ]
        )
        if statuses_png is None:
            print("# SKIP шкала делений у статусов — растр не получен от LibreOffice")
        else:
            statuses = scan_badges(statuses_png)
            check(
                "статус делений не рисует вовсе",
                not statuses["bars"],
                f"найдено шкал: {len(statuses['bars'])}",
            )
            check(
                "слово статуса при этом нарисовано янтарным",
                statuses["pixels"]["amber"] > 100,
                f"янтарных пикселей: {statuses['pixels']['amber']}",
            )

    # --- Панель поверхности рисует «Почему выделено» ---------------------
    #
    # Первая отрисовка этой геометрии. Блок объяснений на панели поверхности
    # необязателен для рендерера: `write_block` выбрасывает его **целиком**,
    # если он не влезает в колонку, и делает это молча — ни телеметрии, ни
    # предупреждения. Пока такой страницы не было ни в одном эталоне, ответ на
    # вопрос «влезает ли» никто не получал.
    #
    # Свойства проверяются двумя приборами, потому что вопросы разные.
    # Геометрию — по растру (`check_pages`, наших мер не касается). Наличие
    # именно этих фраз — по текстовому слою PDF: растр умеет сказать «чернила
    # есть», но не «нарисована вторая фраза», а выброшен бывает как раз
    # блок целиком, и колонка при этом остаётся заполненной заголовком.
    panel = render_surface_panel_page()
    if panel is None:
        print("# SKIP панель поверхности — растр не получен от LibreOffice")
    else:
        panel_png, panel_pdf = panel
        with tempfile.TemporaryDirectory() as tmp:
            panel_path = Path(tmp) / "page-01.png"
            panel_path.write_bytes(panel_png)
            panel_report = check_pages([panel_path])
            check(
                "панель поверхности с объяснениями не выходит за края листа",
                panel_report.passed and not panel_report.findings,
                "; ".join(f"{f.code} — {f.detail}" for f in panel_report.findings)
                or f"проверено страниц: {panel_report.pages_checked}",
            )
            pdf_path = Path(tmp) / "panel.pdf"
            pdf_path.write_bytes(panel_pdf)
            import fitz  # noqa: PLC0415

            drawn = " ".join(fitz.open(str(pdf_path))[0].get_text().split())
        check(
            "заголовок блока «Почему выделено» нарисован",
            "Почему выделено" in drawn,
        )
        missing_reasons = [
            i + 1
            for i, reason in enumerate(PANEL_REASONS)
            if " ".join(reason.split()) not in drawn
        ]
        check(
            "обе фразы объяснений нарисованы целиком",
            not missing_reasons,
            f"не нарисованы объяснения: {missing_reasons}"
            if missing_reasons
            else f"по {SIDEBAR_HIGHLIGHT_BUDGET} знаков каждое",
        )
        check(
            "остаток сигналов назван числом",
            f"Ещё {PANEL_MORE_SIGNALS} похожих сигнала" in drawn,
        )

    # --- Лист второй таблицы выдачи нарисован впервые ---------------------
    #
    # Растровая проверка выше смотрит на страницы эталона 72, а листов второй
    # таблицы выдачи там нет ни одного: в её корпусе нет дополнительных
    # запросов. Значит, про эту страницу растр не говорил ничего — её вёрстку
    # держала одна арифметика.
    #
    # Главный вопрос — тот же, что у первой таблицы: не уехал ли лист за низ
    # белой сцены. Общая растровая проверка на него **не отвечает**: замер
    # ниже показывает, что `check_pages` проходит и на листе, который сцену
    # собой закрыл, — она считает переполнением выход за край слайда, а лист
    # уезжает за край сцены, оставаясь внутри слайда.
    #
    # Число строк берётся не литералом, а ёмкостью из реестра: занизив худшую
    # строку, разработчик поднимет ёмкость, и проверка нарисует лист, которого
    # не бывает, — то есть увидит подмену числа рисованием.
    terms = serp_capacity_terms()
    extra_capacity = (
        terms["SERP_TABLE_ROW_BUDGET_EMU"] // terms["SERP_EXTRA_TABLE_WORST_ROW_EMU"]
    )
    extra = render_serp_extra_table_page(extra_capacity)
    if extra is None:
        print("# SKIP лист второй таблицы выдачи — растр не получен от LibreOffice")
    else:
        extra_png, extra_pdf = extra
        with tempfile.TemporaryDirectory() as tmp:
            extra_path = Path(tmp) / "page-01.png"
            extra_path.write_bytes(extra_png)
            extra_report = check_pages([extra_path])
        check(
            "лист второй таблицы выдачи не выходит за края листа",
            extra_report.passed and not extra_report.findings,
            "; ".join(f"{f.code} — {f.detail}" for f in extra_report.findings)
            or f"строк на листе: {extra_capacity}",
        )
        edge_ink, lowest, size = stage_edge_ink(extra_png)
        check(
            f"нижняя кромка сцены видна под листом из {extra_capacity} худших строк",
            edge_ink >= STAGE_EDGE_INK and size == (W, H),
            f"чернил в полосе кромки {edge_ink} при пороге {STAGE_EDGE_INK}, "
            f"самый низкий ряд чернил y={lowest}, растр {size[0]}×{size[1]}",
        )

        # Контрольный дефект: лист на одну строку длиннее ёмкости обязан кромку
        # закрыть. Без него проверка выше умеет только зелёный — а такая
        # проверка не гейт.
        over = render_serp_extra_table_page(extra_capacity + 1)
        if over is None:
            print("# SKIP контрольный дефект второй таблицы — растр не от LibreOffice")
        else:
            over_ink, over_lowest, _size = stage_edge_ink(over[0])
            check(
                f"контрольный дефект: лист из {extra_capacity + 1} строк кромку закрывает",
                over_ink < STAGE_EDGE_INK,
                f"чернил в полосе кромки {over_ink} при пороге {STAGE_EDGE_INK}, "
                f"самый низкий ряд чернил y={over_lowest}",
            )

        # Лист попал в свою ветку рендерера, а не в ветку первой таблицы.
        # Ветку выбирает **отсутствие** колонки «№», и ошибка здесь не
        # громкая: доли первой таблицы отдали бы адресу 5 % листа вместо 30 %,
        # страница нарисовалась бы без отказа, а телеметрия молчала бы.
        # Меряется по нарисованному: левые края заголовков в текстовом слое.
        edges = header_left_edges(extra_pdf, list(HDR_SERP_EXTRA))
        missing = [h for h in HDR_SERP_EXTRA if h not in edges]
        check(
            "все пять заголовков второй таблицы нарисованы",
            not missing,
            f"не нарисованы: {missing}" if missing else ", ".join(HDR_SERP_EXTRA),
        )
        if not missing:
            span = edges["Оценка"] - edges["Ссылка"]
            share = (edges["Заголовок"] - edges["Ссылка"]) / span if span else 0.0
            # Полоса, а не точное число: этот прибор отвечает на вопрос «та ли
            # ветка», а доли колонок по отдельности держит Т2в смока ширин.
            # 0.349 — своя ветка (0.30 из 0.86), 0.058 — ветка первой таблицы.
            check(
                "лист нарисован долями своей ветки, а не долями первой таблицы",
                0.30 <= share <= 0.40,
                f"колонка «Ссылка» — {share:.3f} ширины первых четырёх колонок "
                f"(своя ветка 0.30/0.86 = 0.349, ветка первой таблицы 0.05/0.86 = 0.058)",
            )

    # --- ADR-0008: типографическая шкала закрыта -------------------------
    #
    # Замер до шага: шестнадцать разных кеглей в деке, до восьми на странице.
    # Глаз читает такой набор не как иерархию, а как шум.
    pptx_path = (
        Path(__file__).resolve().parent.parent
        / "baselines/report-72/artifacts/deck-sections/rendered-client.pptx"
    )
    if not pptx_path.is_file():
        print("# SKIP шкала кеглей эталонной деки — нет отрисованного PPTX")
    else:
        from pptx import Presentation

        used: dict[float, int] = {}
        per_page: dict[int, set] = {}
        runs_by_page: dict[int, list] = {}
        for idx, slide in enumerate(Presentation(str(pptx_path)).slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is None:
                            continue
                        pt = round(run.font.size.pt, 2)
                        used[pt] = used.get(pt, 0) + 1
                        per_page.setdefault(idx, set()).add(pt)
                        if (run.text or "").strip():
                            runs_by_page.setdefault(idx, []).append(
                                (pt, shape.name or "?")
                            )
        first_level_shapes = count_first_level(runs_by_page)
        off = sorted(p for p in used if p not in TYPE_SCALE_PT)
        check(
            "в деке нет кеглей вне шкалы",
            not off,
            f"вне шкалы: {off}" if off else f"использовано ступеней: {sorted(used)}",
        )
        # Иерархия: на странице ровно один элемент первого уровня — лид или
        # ключевая цифра, остальное заведомо тише (ADR-0008, п.3). Свойство
        # соблюдается на всех 45 страницах эталона; проверка держит его,
        # чтобы следующая правка шаблона не вернула «всё одинаково важно».
        offenders = [pg for pg, shapes in first_level_shapes.items() if len(shapes) != 1]
        check(
            "на каждой странице ровно один элемент первого уровня",
            not offenders,
            f"страниц с нарушением: {offenders[:5]}" if offenders else f"страниц: {len(first_level_shapes)}",
        )
        # Отрицательный контроль: две фигуры одного крупнейшего кегля — это
        # нарушение, и проверка обязана его увидеть.
        synthetic = count_first_level({1: [(22.0, "Заголовок A"), (22.0, "Заголовок Б"), (11.0, "текст")]})
        check(
            "две фигуры первого уровня считаются нарушением",
            len(synthetic[1]) == 2,
            f"найдено фигур: {sorted(synthetic[1])}",
        )

        widest = max(per_page.items(), key=lambda kv: len(kv[1])) if per_page else (0, set())
        check(
            "на странице не больше четырёх ступеней",
            len(widest[1]) <= 4,
            f"страница {widest[0]}: {sorted(widest[1])}",
        )

    # --- Гарнитура в PDF: та же, что мы меряем, и встроена ---------------
    #
    # PPTX несёт только имя семейства; рисует его LibreOffice при конвертации.
    # Если гарнитуры нет в образе, он молча подставит запасную — и клиент
    # получит документ, набранный не тем, чем мы его мерили. Кириллица при
    # такой подмене обычно и уезжает.
    pdf_path = (
        Path(__file__).resolve().parent.parent
        / "baselines/report-72/artifacts/deck-sections/rendered-client.pdf"
    )
    if not pdf_path.is_file():
        print("# SKIP гарнитура в PDF — нет отрендеренного PDF")
    else:
        import fitz

        doc = fitz.open(str(pdf_path))
        families, not_embedded = set(), set()
        for page_no in range(len(doc)):
            for font in doc.get_page_fonts(page_no):
                base = str(font[3]).split("+")[-1]
                families.add(base)
                if font[1] == "n/a":
                    not_embedded.add(base)
        family = FONT.replace(" ", "").lower()
        foreign = sorted(f for f in families if family not in f.replace(" ", "").lower())
        check(
            "в PDF только объявленная гарнитура",
            not foreign,
            f"чужие: {foreign}" if foreign else f"найдено: {sorted(families)}",
        )
        check("все начертания встроены в PDF", not not_embedded, f"невстроенные: {sorted(not_embedded)}")

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
