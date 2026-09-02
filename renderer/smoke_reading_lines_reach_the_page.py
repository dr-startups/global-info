#!/usr/bin/env python3
"""Смок: слова о чтении и выводы страниц доезжают до листа, а не до payload.

На прогоне 76 отчёт запросил 120 страниц, прочитал 74 — и не сказал об этом
нигде. Оба носителя честности убила вёрстка, а не построители:

  * `orion_golden_search_table` печатал над таблицей два первых законченных
    предложения, и строка покрытия, стоявшая третьим, исчезала без записи;
  * страница региона (`orion_golden_metrics_dashboard`) поля `statusNote` не
    получала вовсе — фраза с процентом и базой не имела печатного носителя.

Счётчик предложений снят: сколько абзаца влезет, решает мера в пределах
объявленного потолка. Отсюда Т8в и Т8г — «почему важно», «что проверить» и
правило ручной верификации обязаны быть на листе, а не только в нагрузке.

Проверяется именно лист: текст шейпов после отрисовки. Числа здесь
синтетические — только счётчики, без имён и материалов реального дела.

Сеть, база и LibreOffice не нужны: презентация строится в памяти.

Запуск: python3 renderer/smoke_reading_lines_reach_the_page.py (нужен python-pptx)
"""

from __future__ import annotations

import sys
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
from orion_golden_render.common import (  # noqa: E402
    CONTENT_BOTTOM,
    CONTENT_W,
    MARGIN_X,
    SLIDE_H,
    SLIDE_W,
    _Ctx,
    get_layout_telemetry,
    reset_layout_telemetry,
)
from orion_golden_render.slides import _render_slide  # noqa: E402
from orion_golden_render.visual import _sidebar_analysis  # noqa: E402

#: Интро страницы тем в том виде, в каком его собирает построитель секций.
THEMES_INTRO = (
    "Из 120 отобранных по отчёту страниц прочитано 74; каждая прочитанная отнесена "
    "к теме по её содержанию, нежелательных публикаций: 15. "
    "Из непрочитанных: 23 закрыли доступ, 17 без читаемого текста, 6 не ответили."
)

#: Статусная строка страницы региона: доля негатива и база, по которой она посчитана.
READ_SHARE = (
    "Негатив среди прочитанных страниц региона: 15 из 50 (30%); "
    "прочитано 50 из 86 отобранных."
)

failures: list[str] = []

#: Счётчик выполненных проверок — раннер обязан видеть, сколько их было.
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)


def render(slide: dict[str, Any], page: int = 9) -> Any:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, page, 48)
    _render_slide(ctx, slide, {})
    return prs


def page_text(prs: Any) -> str:
    """Весь текст листа: текстовые короба плюс ячейки таблиц."""
    parts: list[str] = []
    for sh in prs.slides[0].shapes:
        if getattr(sh, "has_text_frame", False):
            parts.append(sh.text_frame.text)
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    parts.append(cell.text)
    return "\n".join(parts)


#: Блоки панели: вывод, «что показывает экран», «что это значит», «что сделать».
PANEL_ANALYSIS = {
    "sidebarMode": "context",
    "headlineConclusion": "Собрано 50 подсказок. На панели показаны 10.",
    "whatIsVisible": (
        "Подсказки описывают общий интерес к личности, биографии и написанию имени. "
        "Негативных формулировок среди показанных строк нет: прямых запросов про суды, "
        "долги и санкции не видно. Это не снимает репутационный риск по профилю в целом."
    ),
    "clientMeaning": (
        "Подсказки влияют на первое впечатление при поиске: здесь они не подталкивают "
        "проверяющего к негативным темам, но оставляют интерес к идентификации персоны."
    ),
    "recommendedActions": [
        "Проверить строки со статусом «вероятно» и регулярно отслеживать появление "
        "риск-формулировок в подсказках обоих поисковиков."
    ],
    "provenanceLabel": "Источник — поисковая выдача: у показанных элементов нет отдельных адресов.",
    "highlightExplanations": [],
    "moreSignalsCount": 0,
}


#: Короб панели — тот самый, что даёт `_render_visual_with_sidebar`: картинка
#: занимает 62 % ширины, панели остаётся остальное за вычетом промежутка.
PANEL_IMG_W = int(CONTENT_W * 0.62)
PANEL_X = MARGIN_X + PANEL_IMG_W + 120_000
PANEL_W = CONTENT_W - PANEL_IMG_W - 120_000
PANEL_Y = 1_210_000

#: Высота панели на странице без второго заголовка: до низа сцены.
PANEL_H_WHOLE = CONTENT_BOTTOM - 1_150_000 - 80_000 - 60_000

#: Высота, на которой места хватает не всем блокам, и видно, **кто уступает**.
#:
#: Выведена из меры самого рендерера: вывод занимает 197 815 EMU, рекомендация
#: с заголовком — 991 260, и под средний блок с «Что это значит» остаётся
#: меньше, чем им нужно. До пола рекомендации на этой высоте печатались вывод и
#: средний блок целиком, а «Что сделать» и «Что это значит» выбрасывались; с
#: полом печатается рекомендация, средний блок обрезается по предложениям, а
#: «Что это значит» уступает целиком.
PANEL_H_TIGHT = 2_200_000

#: Высота, при которой не помещается даже обязательный вывод: он заменяется
#: запасной фразой, а подпись источников уходит за низ панели.
PANEL_H_CRAMPED = 400_000


def sidebar_panel(height: int) -> tuple[Any, str]:
    """Панель заданной высоты: контекст отрисовки и текст, который на ней есть."""
    reset_layout_telemetry()
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, 46, 94)
    _sidebar_analysis(
        ctx,
        {"visualAnalysis": PANEL_ANALYSIS},
        x=PANEL_X,
        y=PANEL_Y,
        w=PANEL_W,
        h=height,
    )
    return ctx, page_text(prs)


def loss_kinds(ctx: Any) -> dict[str, str]:
    """Поле панели → что с ним случилось: `dropped` или `truncated`."""
    kinds: dict[str, str] = {}
    for w in ctx.warnings:
        parts = w.split(":")
        if len(parts) >= 4 and parts[0] == "sidebar-qa" and parts[3] in {"dropped", "truncated"}:
            kinds[parts[2]] = parts[3]
    return kinds


def sidebar_drop_checks() -> None:
    """Проверки Т8д: рекомендация переживает тесную панель, потеря слышна."""
    # Тесная панель: место уступают средние блоки, а не рекомендация — она
    # то, ради чего клиент читает страницу.
    ctx, drawn = sidebar_panel(PANEL_H_TIGHT)
    entries = get_layout_telemetry()
    kinds = loss_kinds(ctx)
    action = PANEL_ANALYSIS["recommendedActions"][0]
    visible_tail = "Это не снимает репутационный риск"
    check(
        "Т8д1: тесная панель печатает рекомендацию, а местом уступает средний блок",
        action[:40] in drawn and visible_tail not in drawn,
        drawn[:200],
    )
    check(
        "Т8д2: уступивший блок назван предупреждением",
        any("sidebar-qa" in w and "whatIsVisible" in w for w in ctx.warnings),
        f"предупреждения: {ctx.warnings}",
    )
    check(
        "Т8д2б: о рекомендации предупреждений нет — она не уступает",
        not any("recommendedActions" in w for w in ctx.warnings),
        f"предупреждения: {ctx.warnings}",
    )
    check(
        "Т8д3: потеря записана разметкой, и ни одна запись не о рекомендации",
        len(entries) > 0
        and all(e.get("role") == "sidebar" and e.get("clipped") is True for e in entries)
        and not any("recommendedActions" in str(e.get("name") or "") for e in entries),
        f"записей {len(entries)}: {[e.get('name') for e in entries]}",
    )
    dropped_names = {
        f"orion_sidebar_{field}_p46" for field, kind in kinds.items() if kind == "dropped"
    }
    check(
        "Т8д4: выброшенный блок объявлен потерей содержимого",
        len(dropped_names) > 0
        and all(
            int(e.get("droppedLines") or 0) >= 1
            for e in entries
            if e.get("name") in dropped_names
        ),
        f"исходы {kinds}; записи {[(e.get('name'), e.get('droppedLines')) for e in entries]}",
    )
    check(
        "Т8д5: запись называет свою страницу",
        all(e.get("page") == 46 for e in entries) and len(entries) > 0,
        f"страницы: {[e.get('page') for e in entries]}",
    )

    # Обрезка — не потеря целого блока: часть предложений напечатана, и
    # блокером выдачи такая запись не становится.
    check(
        "Т8д4б: обрезанный блок остаётся клипом, а не потерей",
        kinds.get("whatIsVisible") == "truncated"
        and all(
            not int(e.get("droppedLines") or 0)
            for e in entries
            if e.get("name") == "orion_sidebar_whatIsVisible_p46"
        ),
        f"исходы {kinds}; записи {[(e.get('name'), e.get('droppedLines')) for e in entries]}",
    )
    truncated = next(
        (e for e in entries if str(e.get("name", "")).endswith("whatIsVisible_p46")), {}
    )
    check(
        "Т8д8: запись обрезанного блока меряет потерянное, а не весь блок",
        0 < int(truncated.get("textLength") or 0) < len(str(PANEL_ANALYSIS["whatIsVisible"])),
        f"textLength={truncated.get('textLength')} при длине блока "
        f"{len(str(PANEL_ANALYSIS['whatIsVisible']))}",
    )

    # Подпись источников — третья молчавшая ветка. Она срабатывает только на
    # панели, где даже обязательный вывод не поместился и заменён запасной
    # фразой: у прочих блоков в запасе остаётся 160 000 EMU, и подпись после
    # них всегда влезает (проверено арифметикой `write_block`).
    cramped_ctx, cramped = sidebar_panel(PANEL_H_CRAMPED)
    check(
        "Т8д9: подпись источников на тесной панели действительно не напечатана",
        str(PANEL_ANALYSIS["provenanceLabel"])[:40] not in cramped,
        cramped[:200],
    )
    cramped_entries = get_layout_telemetry()
    reset_layout_telemetry()
    check(
        "Т8д10: невлезшая подпись источников объявлена",
        any("sidebar-qa" in w and "provenanceLabel" in w for w in cramped_ctx.warnings),
        f"предупреждения: {cramped_ctx.warnings}",
    )
    # Подпись источников — мелкая строка происхождения, а не блок, ради
    # которого заводился блокер выдачи («клиент видит риск и не видит, что
    # делать»). Её невлезание остаётся клипом: отчёт из-за неё не пропадает.
    check(
        "Т8д12: потеря подписи источников блокером выдачи не становится",
        all(
            not int(e.get("droppedLines") or 0)
            for e in cramped_entries
            if str(e.get("name") or "").endswith("provenanceLabel_p46")
        )
        and any(
            str(e.get("name") or "").endswith("provenanceLabel_p46") for e in cramped_entries
        ),
        f"записи: {[(e.get('name'), e.get('droppedLines')) for e in cramped_entries]}",
    )

    # Вторая рекомендация: ворот приёмки проверяет каждый элемент
    # `recommendedActions`, а панель рисовала только первый — приёмка краснела
    # бы на законной ветке рендерера.
    two_actions = dict(PANEL_ANALYSIS)
    two_actions["recommendedActions"] = [
        "Проверить строки со статусом «вероятно».",
        "Отслеживать появление риск-формулировок ежемесячно.",
    ]
    reset_layout_telemetry()
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    both_ctx = _Ctx(prs, 47, 94)
    _sidebar_analysis(
        both_ctx,
        {"visualAnalysis": two_actions},
        x=PANEL_X,
        y=PANEL_Y,
        w=PANEL_W,
        h=PANEL_H_WHOLE,
    )
    both = page_text(prs)
    reset_layout_telemetry()
    check(
        "Т8д11: на просторной панели рисуются все рекомендации, а не первая",
        all(a[:40] in both for a in two_actions["recommendedActions"]),
        both[-260:],
    )

    # Просторная панель: все четыре блока на листе, и объявлять нечего.
    whole_ctx, whole = sidebar_panel(PANEL_H_WHOLE)
    whole_entries = get_layout_telemetry()
    reset_layout_telemetry()
    printed = all(
        str(text)[:40] in whole
        for text in (
            PANEL_ANALYSIS["headlineConclusion"],
            PANEL_ANALYSIS["clientMeaning"],
            action,
        )
    )
    check("Т8д6: просторная панель печатает свои блоки", printed, whole[:200])
    check(
        "Т8д7: целая панель ни предупреждений, ни записей не даёт",
        not whole_ctx.warnings and whole_entries == [],
        f"предупреждения {whole_ctx.warnings}, записей {len(whole_entries)}",
    )


def main() -> int:
    # --- Т8а. Двухпредложное интро страницы тем напечатано целиком -----------
    themes = render(
        {
            "template": "orion_golden_search_table",
            "title": "Россия — о чём публикации в ТОП-20",
            "narrative": THEMES_INTRO,
            "table": {
                "headers": ["Тема", "Публикаций", "Из них нежелательных"],
                "rows": [
                    ["Криминальные / судебные материалы", "20", "10"],
                    ["Деловая репутация", "12", "5"],
                ],
            },
        }
    )
    themes_text = page_text(themes)
    check("Т8а1: страница тем называет прочитанное", "74" in themes_text, themes_text[:120])
    check("Т8а2: страница тем называет базу", "120 отобранных" in themes_text, themes_text[:120])
    check(
        "Т8а3: причины отказов напечатаны вторым предложением",
        "23 закрыли доступ" in themes_text and "6 не ответили" in themes_text,
        themes_text[:200],
    )
    check(
        "Т8а4: абсолюта «прочитаны» без базы на листе нет",
        "прочитаны" not in themes_text,
        themes_text[:120],
    )

    # --- Т8в. Страница выдачи печатает весь свой абзац ------------------------
    #
    # Шаблон брал из нарратива первые два законченных предложения и остальное
    # выбрасывал без записи. «Почему важно» и «что проверить» страницы выдачи
    # до клиента не доезжали вовсе: на эталоне 72 так терялись 13 × 2
    # предложения. Проверяется текст фигур готового листа, а не мера: мерить
    # лист тем же прибором, который его резал, бессмысленно.
    serp = render(
        {
            "template": "orion_golden_search_table",
            "title": "Россия — Яндекс, ТОП-20 (1/4)",
            "narrative": "\n".join(
                [
                    "Показана выдача Яндекса по запросу «Глинка Сергей Михайлович». "
                    "«PEP / RCA / watchlist-сигналы» — средний уровень внимания "
                    "(строки 1 и 2).",
                    "На странице 1 тема повышенного внимания — эти материалы видны при "
                    "первой же проверке субъекта.",
                    "Запросить первичные карточки баз и подтвердить принадлежность совпадений.",
                ]
            ),
            "table": {
                "headers": ["№", "Заголовок", "Тип источника", "Оценка"],
                "rows": [
                    ["1", "Материал о проверке", "СМИ", "Нежелательный"],
                    ["2", "Деловая публикация", "СМИ", "Нейтральный"],
                ],
                "rowAddresses": ["kommersant.ru/doc/1", "forbes.ru/2"],
            },
        }
    )
    serp_text = page_text(serp)
    check(
        "Т8в1: страница выдачи называет запрос своей таблицы",
        "по запросу «Глинка Сергей Михайлович»" in serp_text,
        serp_text[:200],
    )
    check(
        "Т8в2: вывод страницы напечатан",
        "средний уровень внимания" in serp_text,
        serp_text[:200],
    )
    check(
        "Т8в3: «почему важно» доезжает до листа",
        "видны при первой же проверке субъекта" in serp_text,
        serp_text[:400],
    )
    check(
        "Т8в4: «что проверить» доезжает до листа",
        "подтвердить принадлежность совпадений" in serp_text,
        serp_text[:400],
    )

    # --- Т8г. Сводка комплаенса печатает правило ручной верификации -----------
    #
    # «Совпадение по комплаенсу не подтверждается автоматически» — правило
    # продукта, и до этого шага у него не было печатного носителя: рекомендация
    # стояла третьим предложением и отбрасывалась.
    compliance = render(
        {
            "template": "orion_golden_search_table",
            "title": "Комплаенс-базы: что проверялось",
            "narrative": "\n".join(
                [
                    "Записей, отобранных по имени субъекта в комплаенс-базах: 3 (Dow "
                    "Jones, LexisNexis и World-Check). Совпадение по базе не "
                    "подтверждается автоматически: "
                    "подтверждено аналитиком — 0, требует ручной проверки — 0, статус не "
                    "зафиксирован — 3.",
                    "Верифицировать каждое потенциальное совпадение вручную: сопоставить "
                    "идентификаторы субъекта с записью базы.",
                ]
            ),
            "table": {
                "headers": [
                    "База данных",
                    "Тип совпадения",
                    "Совпадение по имени",
                    "Статус проверки",
                ],
                "rows": [
                    ["Dow Jones", "PEP", "Глинка Сергей Михайлович", "Не подтверждено"],
                    ["LexisNexis", "Публикации", "Sergey Glinka", "Не подтверждено"],
                ],
            },
        },
        page=33,
    )
    compliance_text = page_text(compliance)
    check(
        "Т8г1: сводка комплаенса печатает предложение о ручной верификации",
        "Верифицировать каждое потенциальное совпадение вручную" in compliance_text,
        compliance_text[:400],
    )
    check(
        "Т8г2: разбивка статусов на листе осталась",
        "статус не зафиксирован — 3" in compliance_text,
        compliance_text[:400],
    )

    # --- Т8б. Статусная строка страницы региона ------------------------------
    dashboard = render(
        {
            "template": "orion_golden_metrics_dashboard",
            "title": "Россия: в выдаче есть материалы повышенного внимания",
            "narrative": "Предмет аудита по региону «Россия» — ТОП-20 выдачи: 20 материалов.",
            "statusNote": READ_SHARE,
            "metrics": [
                {"label": "Собрано по региону", "value": "404", "tone": "neutral"},
                {"label": "Тем повышенного внимания", "value": "2", "tone": "risk"},
            ],
            "actions": [{"label": "Проверить актуальные статусы дел."}],
            "bullets": ["«Криминальные / судебные материалы». Всего по теме: 7 материалов."],
        },
        page=7,
    )
    dash_text = page_text(dashboard)
    check("Т8б1: доля негатива напечатана", "15 из 50 (30%)" in dash_text, dash_text[:200])
    check("Т8б2: база доли напечатана", "прочитано 50 из 86 отобранных" in dash_text, dash_text[:200])
    check(
        "Т8б3: нарратив и карточка действия на месте",
        "ТОП-20 выдачи" in dash_text and "статусы дел" in dash_text,
        dash_text[:200],
    )

    without = render(
        {
            "template": "orion_golden_metrics_dashboard",
            "title": "Обзор цифрового профиля",
            "narrative": "Предмет аудита по региону «Россия» — ТОП-20 выдачи: 20 материалов.",
            "metrics": [{"label": "Собрано по региону", "value": "404", "tone": "neutral"}],
        },
        page=5,
    )
    check(
        "Т8б4: без статусной строки страница её не выдумывает",
        "прочитано" not in page_text(without),
        page_text(without)[:200],
    )

    # --- Т8д. Панель объявляет блок, который не поместился ------------------
    #
    # Прогон 91, стр. 46: панель подсказок напечатала вывод и «что показывает
    # экран», а «Что сделать» не напечатала вовсе — необязательный блок, у
    # которого не влезло даже первое предложение, `write_block` выбрасывал
    # целиком (`else: return`). Ни предупреждения, ни записи разметки при этом
    # не появлялось: в `layout-telemetry.json` прогона 62 записи и ни одной про
    # эту страницу. Потерянный клиентский текст обязан быть слышен.
    sidebar_drop_checks()

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
