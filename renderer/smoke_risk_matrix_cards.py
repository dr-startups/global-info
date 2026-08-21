#!/usr/bin/env python3
"""Смок: карточная сетка матрицы рисков рисует всё, что ей дали, и говорит о потере.

На прогоне 72 страница 6 несла три карточки в модели и две на бумаге: третья
(«PEP / RCA / watchlist-сигналы») не была нарисована нигде в деке. Причина —
`_safe` на входе detail: он схлопывал переносы, строчная логика карточки
умирала, карточки раздувались и последняя не помещалась. Симптом никто не
поймал: карточная сетка не писала телеметрию, а без входа правило
`CONTENT_DROPPED_BY_RENDERER` исполнять нечему.

Проверяется оба конца:
  * здоровая страница — нарисованы все карточки, переносы дожили до абзацев,
    «Что делать» стоит своей строкой, потерь по телеметрии ноль;
  * контрольные дефекты — лишняя карточка сверх ёмкости листа и страница
    карточек-переростков: они не рисуются, телеметрия сообщает потерю, и
    инспектор геометрии даёт по ней `CONTENT_DROPPED_BY_RENDERER`. Проверка,
    не умеющая дать не-ноль, — не проверка.

Отдельно разведены три события, которые в записи выглядят похоже: потеря
содержимого (`dropped*`, блокирует выдачу), сработавшая защита — понижение
кегля (`clipped`, предупреждение) и уборка повисшего двоеточия, которая не
является ни тем, ни другим. Сама мера сверена с растровым замером страницы
эталона: без внешней опоры «нарисовано столько, сколько обещано» остаётся
верным при любом запасе.

Сеть, база и LibreOffice не нужны: презентация строится в памяти, инспектор
геометрии читает файл из временного каталога.

Запуск: python3 renderer/smoke_risk_matrix_cards.py (нужен python-pptx)
"""

from __future__ import annotations

import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
import orion_golden_render.executive as executive_mod  # noqa: E402
from orion_golden_render.common import (  # noqa: E402
    CONTENT_BOTTOM,
    CONTENT_W,
    EMU_PER_PT,
    FS_BODY,
    RISK_BG,
    SLIDE_H,
    SLIDE_W,
    _Ctx,
    _split_structured_bullet,
    _wrapped_line_count,
    get_layout_telemetry,
    measure_text_height,
    reset_layout_telemetry,
)
from orion_golden_render.executive import (  # noqa: E402
    _card_body_height,
    _card_line_style,
    _card_paragraph_spacing,
    _render_risk_matrix_grid,
)

REPO_ROOT = Path(__file__).resolve().parent.parent
INSPECTOR = REPO_ROOT / "scripts" / "inspect-first36-pptx-geometry.py"
SECTIONS_TS = REPO_ROOT / "src" / "modules" / "digital-profile" / "orion-golden" / "deck-sections"


def ts_int(path: Path, pattern: str) -> int:
    """Число из исходника на той стороне языковой границы.

    Бюджеты знаков карточки живут в TS (реестр шаблонов и построитель), а
    худшую легальную карточку строит этот смок. Пока числа были здесь
    литералами, поднятие бюджета на той стороне оставляло проверку зелёной, а
    реальная худшая карточка росла — и ёмкость листа переставала быть
    гарантией. Питон в TS не ходит, поэтому читается исходник; не нашли
    число — это отказ, а не повод считать по умолчанию.
    """
    match = re.search(pattern, path.read_text(encoding="utf-8"), re.S)
    if not match:
        raise RuntimeError(f"бюджет по шаблону {pattern!r} не найден в {path.name}")
    return int(match.group(1))


#: Бюджет тела карточки матрицы целиком (`itemCharBudget` шаблона `risk-matrix`).
ITEM_CHAR_BUDGET = ts_int(
    SECTIONS_TS / "template-registry.ts",
    r'"risk-matrix": \{.*?itemCharBudget:\s*(\d+)',
)
#: Бюджет строки «в чём проблема».
PROBLEM_CHARS = ts_int(
    SECTIONS_TS / "fragment-builders" / "executive.ts",
    r"const RISK_MATRIX_PROBLEM_CHARS = (\d+)",
)
#: Столько знаков оставляет от заголовка сам рендерер (`_clip_words` ниже по
#: `_render_risk_matrix_grid`); длиннее клиент не увидит при любом построителе.
HEADLINE_CHARS = 72
ACTION_PREFIX = "Что делать: "
#: Бюджет знаков, которым рендерер режет тело карточки перед отрисовкой
#: (`_clip_structured_bullet(detail, 900)` в `_render_risk_matrix_grid`).
RENDERER_CHAR_BUDGET = 900

#: Сводная карточка в том виде, в каком её собирает построитель секций.
SUMMARY_DETAIL = (
    "7 свидетельств (5 негативных) в источниках audit-it.ru, x.com, m.sledst.org.\n"
    "Что делать: Проверить статусы дел и первоисточники до принятия решений."
)

#: Проза-наполнитель: слова обычной длины, чтобы перенос считался как на прозе.
_FILLER = (
    "Опубликованы материалы о корпоративном конфликте вокруг профильного актива "
    "и о давлении на суд, издания нескольких стран повторяют сюжет третий месяц "
    "подряд и добавляют подробности сделки, состава участников и позиции сторон "
    "спора, а также сроков разбирательства в судах нескольких инстанций подряд. "
)


def prose(chars: int) -> str:
    """Ровно `chars` знаков прозы — столько, сколько отдаёт бюджет построителя."""
    return (_FILLER * (chars // len(_FILLER) + 1))[:chars].strip()


#: Худшая легальная карточка: столько текста способен выпустить построитель.
#:
#: Строка «в чём проблема» — бюджет `RISK_MATRIX_PROBLEM_CHARS`, «Что делать» —
#: остаток бюджета карточки (минус префикс и перенос), заголовок — то, что
#: оставляет рендерер. Обе части тела при этом рисуются в две строки: четыре
#: строки тела — потолок построителя. Ёмкость листа в реестре выведена из неё,
#: а не из типичной карточки: на типичной ёмкости живой прогон терял карточки
#: дважды.
WORST_PROBLEM = prose(PROBLEM_CHARS)
WORST_ACTION = ACTION_PREFIX + prose(
    ITEM_CHAR_BUDGET - PROBLEM_CHARS - len(ACTION_PREFIX) - len("\n")
)
WORST_HEADLINE = prose(HEADLINE_CHARS)
#: Карточка-переросток: столько текста на странице помещается не всё. Втрое
#: длиннее бюджета знаков рендерера — значит, упирается и в него, и в лист.
OVERSIZED_DETAIL = prose(3 * RENDERER_CHAR_BUDGET)

failures: list[str] = []

#: Счётчик выполненных проверок — раннер обязан видеть, сколько их было.
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)


def card(idx: int, detail: str) -> dict[str, str]:
    return {
        "headline": f"Тема риска {idx}",
        "status": "Высокий",
        "detail": detail,
        "tone": "risk",
    }


#: Геометрия карточной сетки — те же числа, которыми рендерер её раскладывает.
BADGE_W = 1_900_000
PAD_Y = 100_000
CARD_GAP = 50_000
#: Полоса под сеткой: лист без заголовка страницы (240 000 + 950 000) и без
#: нижнего зазора рендерера (100 000).
PAGE_AVAILABLE = CONTENT_BOTTOM - (240_000 + 950_000) - 100_000


def render_page(cards: list[dict[str, str]], page: int = 6) -> tuple[Any, list[dict[str, Any]]]:
    """Отрисовать одну страницу матрицы и вернуть презентацию с телеметрией."""
    reset_layout_telemetry()
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, page, 48)
    _render_risk_matrix_grid(ctx, {"keyFindings": cards}, "Матрица комплаенс-рисков")
    return prs, get_layout_telemetry()


def render_with_fitter(cards: list[dict[str, str]], fit: Any) -> tuple[Any, list[dict[str, Any]]]:
    """Отрисовать страницу с подменённой страховкой `_fit_lines_to_height`.

    На легальном входе фиттер не срабатывает по построению: его бюджет
    производен от той же меры, которой посчитана высота карточки. Значит,
    вызвать срабатывание можно только подменой — а проверка «потеря слышна»,
    неспособная дать не-ноль, ничего не проверяет.
    """
    original = executive_mod._fit_lines_to_height
    executive_mod._fit_lines_to_height = fit
    try:
        return render_page(cards)
    finally:
        executive_mod._fit_lines_to_height = original


def card_fills(prs: Any) -> list[str]:
    """Заливки карточек страницы в порядке отрисовки — шестнадцатеричными кодами."""
    out: list[str] = []
    for sh in prs.slides[0].shapes:
        if not str(getattr(sh, "name", "")).startswith("orion_card_p"):
            continue
        try:
            out.append(str(sh.fill.fore_color.rgb))
        except Exception:  # noqa: BLE001 — форма без сплошной заливки нам не интересна
            continue
    return out


def rects(prs: Any, width: int) -> list[tuple[int, int]]:
    """(верх, высота) прямоугольников карточек либо плашек — по их ширине."""
    return [
        (int(sh.top), int(sh.height))
        for sh in prs.slides[0].shapes
        if str(getattr(sh, "name", "")).startswith("orion_card_p") and int(sh.width) == width
    ]


def drawn_headlines(prs: Any, cards: list[dict[str, str]]) -> list[str]:
    wanted = {c["headline"] for c in cards}
    slide = prs.slides[0]
    return [
        sh.text_frame.text
        for sh in slide.shapes
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text in wanted
    ]


def detail_paragraphs(prs: Any) -> list[list[str]]:
    """Абзацы текстовых коробов, несущих тело карточки."""
    out: list[list[str]] = []
    for sh in prs.slides[0].shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        paras = [p.text for p in sh.text_frame.paragraphs if p.text.strip()]
        if any(p.startswith("Что делать:") for p in paras):
            out.append(paras)
    return out


def risk_entries(telemetry: list[dict[str, Any]], page: int) -> list[dict[str, Any]]:
    return [
        e
        for e in telemetry
        if e.get("page") == page and str(e.get("name") or "").startswith("orion_risk_matrix")
    ]


def inspector_codes(prs: Any, telemetry: list[dict[str, Any]]) -> list[str]:
    """Прогнать инспектор геометрии по странице вместе с её телеметрией."""
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "page.pptx"
        prs.save(str(pptx_path))
        (Path(tmp) / "layout-telemetry.json").write_text(
            json.dumps({"version": "orion-layout-telemetry-v1", "entries": telemetry}, ensure_ascii=False),
            encoding="utf-8",
        )
        out = subprocess.run(
            [sys.executable, "-X", "utf8", str(INSPECTOR), str(pptx_path)],
            capture_output=True,
            encoding="utf-8",
            check=False,
        )
        report = json.loads(out.stdout or "{}")
    return [str(c.get("code")) for c in report.get("clipping") or []]


def card_height(headline: str, detail: str, pill: str) -> int:
    """Высота карточки по мере рендерера — той же, которой он решает «влезает».

    Литерала ёмкости здесь больше нет: он был вторым ответом на вопрос реестра
    и держался комментарием, а не замером. Проверка «нарисованы все» имеет
    смысл только при выведенном числе — иначе она сверяет мнение с мнением.
    """
    text_w = CONTENT_W - BADGE_W - 220_000 if pill else int(CONTENT_W * 0.92)
    headline_h = measure_text_height(headline, text_w, 13, line_spacing=1.15, bold=True)
    if not detail:
        return max(420_000, PAD_Y + headline_h + PAD_Y)
    return PAD_Y + headline_h + 40_000 + _card_body_height(detail, text_w, FS_BODY) + PAD_Y


def page_capacity(headline: str, detail: str, pill: str) -> int:
    """Сколько таких карточек физически помещается на лист.

    Рендерер рисует карточку, пока её низ не ниже `CONTENT_BOTTOM − 100 000`,
    и добавляет зазор 50 000 между карточками; сетка начинается под полосой
    заголовка страницы. То же неравенство, выраженное числом карточек.
    """
    h = card_height(headline, detail, pill)
    n = 0
    while (n + 1) * h + n * CARD_GAP <= PAGE_AVAILABLE:
        n += 1
    return n


#: Шаг 11-pt строки тела, замеренный по растру эталона: страница 6
#: (`baselines/report-72/artifacts/deck-sections/pages-png/page-06.png`, 144 dpi)
#: даёт 34 px между началами соседних абзацев — 2 pt межабзацной отбивки и 30 px
#: самой строки. Это опора вне кода меры: `smoke_text_measurement.py` сверяет с
#: ней `font_line_step_emu`, а здесь ею проверяется запас карточной меры.
RASTER_LINE_STEP = 190_500


def raster_body_height(detail: str, text_w: int, detail_font: float) -> int:
    """Сколько высоты телу карточки нужно физически — по растровому шагу строки.

    Прибор нарочно другой: `_card_body_height` считает шаг из метрик шрифта и
    умножает на именованный запас, и проверять её тем же способом значит
    проверять мнение мнением. Здесь шаг — замеренный на бумаге, запаса нет.
    """
    lines = _split_structured_bullet(detail) or ([detail] if detail else [])
    total = 0
    for index, line in enumerate(lines):
        bold, _color, size_pt = _card_line_style(line, index, detail_font)
        step = int(RASTER_LINE_STEP * size_pt / 11)
        total += _wrapped_line_count(line, text_w, size_pt, bold) * step
        total += sum(_card_paragraph_spacing(index)) * EMU_PER_PT
    return total


def main() -> int:
    # --- К1. Здоровая страница: полный лист сводных карточек -----------------
    PAGE_CARDS = page_capacity("Тема риска 1", SUMMARY_DETAIL, "Высокий")
    healthy = [card(i, SUMMARY_DETAIL) for i in range(1, PAGE_CARDS + 1)]
    prs, telemetry = render_page(healthy)
    drawn = drawn_headlines(prs, healthy)
    check(
        f"К1а: нарисованы все {PAGE_CARDS} сводные карточки листа",
        len(drawn) == PAGE_CARDS,
        f"нарисовано {len(drawn)}: {', '.join(drawn) or '—'}",
    )
    bodies = detail_paragraphs(prs)
    check(
        "К1б: тело каждой карточки — не меньше двух абзацев (переносы дожили)",
        len(bodies) == PAGE_CARDS and all(len(b) >= 2 for b in bodies),
        f"абзацев по карточкам: {[len(b) for b in bodies]}",
    )
    check(
        "К1в: «Что делать» стоит отдельной строкой, а не вклейкой",
        len(bodies) == PAGE_CARDS and all(b[-1].startswith("Что делать:") for b in bodies),
        f"последние строки: {[b[-1][:32] for b in bodies]}",
    )
    entries = risk_entries(telemetry, 6)
    check(
        "К1г: страница матрицы написала запись телеметрии",
        len(entries) == 1,
        f"записей {len(entries)} из {len(telemetry)}",
    )
    dropped = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries
    )
    check("К1д: потерь по телеметрии ноль", dropped == 0, f"потеряно {dropped}")

    # BX: клиентская шкала метит верхнюю ступень тоном `danger`, а рендерер знал
    # только `risk|warn|good|accent` — незнакомый тон красился дефолтом, то есть
    # **белым**. Карточка «Высокий» выходила единственной без тревожного фона.
    danger_prs, _ = render_page(
        [
            {"headline": "Тема ступени", "status": "Высокий", "detail": prose(120), "tone": "danger"},
            {"headline": "Соседняя тема", "status": "Средний", "detail": prose(120), "tone": "warn"},
        ]
    )
    fills = card_fills(danger_prs)
    check(
        "BX: верхняя ступень не красится белым",
        len(fills) >= 1 and fills[0].upper() != "FFFFFF",
        f"заливки: {fills}",
    )
    check(
        "BX: `danger` красится так же, как `risk`",
        len(fills) >= 1 and fills[0].upper() == str(RISK_BG).upper(),
        f"верхняя карточка {fills[0] if fills else '—'} против RISK_BG {RISK_BG}",
    )
    check(
        "К1е: инспектор геометрии молчит о потере содержимого",
        "CONTENT_DROPPED_BY_RENDERER" not in inspector_codes(prs, telemetry),
        "",
    )

    # --- К1ж. Лишняя карточка на лист не влезает, и это слышно ---------------
    # Контроль ёмкости: без него «нарисованы все N» вырождается в тавтологию
    # «сколько мера обещала, столько и нарисовали» — обещание проверяет само
    # себя. Красной проверку делает ровно одна карточка сверх выведенного N.
    over = [card(i, SUMMARY_DETAIL) for i in range(1, PAGE_CARDS + 2)]
    prs_over, telemetry_over = render_page(over)
    drawn_over = drawn_headlines(prs_over, over)
    check(
        f"К1ж: {PAGE_CARDS + 1}-я карточка на лист не попала",
        len(drawn_over) == PAGE_CARDS,
        f"нарисовано {len(drawn_over)} из {len(over)}",
    )
    entries_over = risk_entries(telemetry_over, 6)
    check(
        "К1з: невлезшая карточка посчитана в droppedBullets",
        len(entries_over) == 1 and int(entries_over[0].get("droppedBullets") or 0) == 1,
        f"запись: {entries_over[0] if entries_over else '—'}",
    )
    check(
        "К1и: инспектор геометрии называет это потерей содержимого",
        "CONTENT_DROPPED_BY_RENDERER" in inspector_codes(prs_over, telemetry_over),
        "",
    )
    check(
        # `clipped` и `dropped_*` отвечают на разные вопросы: первый — «защита
        # сработала, содержимое доехало», второй — «содержимое потеряно». Пока
        # признак клипа выставлялся по факту потери, он значил ровно
        # противоположное определению, и ветка «клип без потери» у обоих
        # классификаторов была недостижима для карточной страницы.
        "К1к: выброшенная карточка клипом не объявлена",
        len(entries_over) == 1 and entries_over[0].get("clipped") is False,
        f"clipped={entries_over[0].get('clipped') if entries_over else '—'}",
    )

    # --- К2. Список карточек рендерер не режет -------------------------------
    # Ёмкость страницы — ответ построителя; рендерер отвечает только за то,
    # влезает ли карточка целиком. Здесь проверяется именно отсутствие своего
    # среза списка, поэтому карточки взяты без тела: семь голых заголовков
    # (424 040 EMU каждый) по высоте помещаются, а срез на шесть — нет.
    seven = [card(i, "") for i in range(1, 8)]
    prs7, telemetry7 = render_page(seven)
    drawn7 = drawn_headlines(prs7, seven)
    check("К2а: семь коротких карточек нарисованы все", len(drawn7) == 7, f"нарисовано {len(drawn7)}")
    entries7 = risk_entries(telemetry7, 6)
    check(
        "К2б: телеметрия насчитала семь нарисованных карточек",
        len(entries7) == 1 and entries7[0].get("drawnBlocks") == 7,
        f"запись: {entries7[0] if entries7 else '—'}",
    )
    check(
        "К2в: строк тела у голых заголовков — ноль, и поле говорит именно о строках",
        len(entries7) == 1 and entries7[0].get("measuredLines") == 0,
        f"measuredLines={entries7[0].get('measuredLines') if entries7 else '—'}",
    )

    # --- К3. Контрольный дефект: переростки не рисуются и об этом слышно -----
    monsters = [card(i, OVERSIZED_DETAIL) for i in range(1, 7)]
    prs_bad, telemetry_bad = render_page(monsters)
    drawn_bad = drawn_headlines(prs_bad, monsters)
    check(
        "К3а: карточка, не влезающая целиком, не нарисована",
        len(drawn_bad) < len(monsters),
        f"нарисовано {len(drawn_bad)} из {len(monsters)}",
    )
    entries_bad = risk_entries(telemetry_bad, 6)
    dropped_bad = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_bad
    )
    check(
        # Невлезшие карточки считаются блоками, а срез по бюджету знаков (тело
        # переростка длиннее 900) — строками: складывать их в одно число
        # значит терять, что именно пропало.
        "К3б: телеметрия сообщает потерю",
        len(entries_bad) == 1
        and int(entries_bad[0].get("droppedBullets") or 0) == len(monsters) - len(drawn_bad)
        and dropped_bad > 0,
        f"потеряно по телеметрии {dropped_bad} (блоков "
        f"{entries_bad[0].get('droppedBullets') if entries_bad else '—'}), "
        f"не нарисовано {len(monsters) - len(drawn_bad)}",
    )
    codes = inspector_codes(prs_bad, telemetry_bad)
    check(
        "К3в: инспектор геометрии даёт CONTENT_DROPPED_BY_RENDERER",
        "CONTENT_DROPPED_BY_RENDERER" in codes,
        f"коды: {', '.join(sorted(set(codes))) or '—'}",
    )

    # --- К4. Строка, снятая страховкой внутри карточки, слышна ---------------
    # Предрисовочные «ножницы» съели «Что делать» на живой странице 8 при
    # чистой телеметрии. Ножницы убраны, страховка фиттера осталась последним
    # рубежом — и о её срабатывании обязан узнать контур.
    trimmed = [card(i, SUMMARY_DETAIL) for i in range(1, 4)]
    prs_fit, telemetry_fit = render_with_fitter(
        trimmed, lambda text, *a, **kw: "\n".join(str(text).split("\n")[:-1])
    )
    entries_fit = risk_entries(telemetry_fit, 6)
    check(
        "К4а: снятая фиттером строка посчитана в droppedLines",
        len(entries_fit) == 1 and int(entries_fit[0].get("droppedLines") or 0) == len(trimmed),
        f"снято по строке с каждой из {len(trimmed)} карточек, в записи "
        f"droppedLines={entries_fit[0].get('droppedLines') if entries_fit else '—'}",
    )
    codes_fit = inspector_codes(prs_fit, telemetry_fit)
    check(
        "К4б: инспектор геометрии называет это потерей содержимого",
        "CONTENT_DROPPED_BY_RENDERER" in codes_fit,
        f"коды: {', '.join(sorted(set(codes_fit))) or '—'}",
    )
    check(
        "К4в: карточки при этом нарисованы все",
        len(drawn_headlines(prs_fit, trimmed)) == len(trimmed),
        f"нарисовано {len(drawn_headlines(prs_fit, trimmed))} из {len(trimmed)}",
    )

    # --- К5. Карточка с опустевшим телом занимает свою высоту -----------------
    # `continue` пропускал и бейдж, и приращение y: следующая карточка ложилась
    # поверх текущей.
    pair = [card(1, SUMMARY_DETAIL), card(2, SUMMARY_DETAIL)]
    prs_empty, _telemetry_empty = render_with_fitter(pair, lambda text, *a, **kw: "")
    card_rects = rects(prs_empty, CONTENT_W)
    check(
        "К5а: следующая карточка начинается ниже низа предыдущей",
        len(card_rects) == 2 and card_rects[1][0] >= card_rects[0][0] + card_rects[0][1],
        f"прямоугольники: {card_rects}",
    )
    check(
        "К5б: бейдж карточки с пустым телом нарисован",
        len(rects(prs_empty, BADGE_W)) == 2,
        f"плашек {len(rects(prs_empty, BADGE_W))} из 2",
    )

    # --- К6. Три худших легальных карточки помещаются целиком ----------------
    # Закрепление ёмкости реестра: она выведена из этой карточки, а не из
    # типичной. Число «3» — зеркало реестра через языковую границу (питон в TS
    # не ходит); при регрессе меры или бюджетов текста проверка краснеет.
    REGISTRY_CAPACITY = 3
    worst = [
        {
            "headline": WORST_HEADLINE,
            "status": "Требует подтверждения",
            "detail": f"{WORST_PROBLEM}\n{WORST_ACTION}",
            "tone": "warn",
        }
        for _ in range(REGISTRY_CAPACITY)
    ]
    prs_worst, telemetry_worst = render_page(worst)
    drawn_worst = drawn_headlines(prs_worst, worst)
    check(
        f"К6а: {REGISTRY_CAPACITY} худших легальных карточки нарисованы все",
        len(drawn_worst) == REGISTRY_CAPACITY,
        f"нарисовано {len(drawn_worst)} из {REGISTRY_CAPACITY}; "
        f"ёмкость по мере {page_capacity(WORST_HEADLINE, worst[0]['detail'], worst[0]['status'])}",
    )
    entries_worst = risk_entries(telemetry_worst, 6)
    lost_worst = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_worst
    )
    check("К6б: потерь на худшем легальном листе нет", lost_worst == 0, f"потеряно {lost_worst}")
    worst_bodies = detail_paragraphs(prs_worst)
    check(
        "К6в: «Что делать» стоит последней строкой каждой карточки",
        len(worst_bodies) == REGISTRY_CAPACITY
        and all(b[-1] == WORST_ACTION for b in worst_bodies),
        f"последние строки: {[b[-1][:28] for b in worst_bodies]}",
    )
    # Поле схемы называется «сколько строк намерено» и у всех остальных
    # производителей означает строки переноса. У худшей карточки структурных
    # строк две, а нарисованных — четыре: считать структурные значит отвечать
    # не на тот вопрос, ради которого поле читают глазами.
    structural = sum(len(b) for b in worst_bodies)
    check(
        "К6г: measuredLines считает нарисованные строки переноса, а не структурные",
        len(entries_worst) == 1
        and int(entries_worst[0].get("measuredLines") or 0) == 4 * REGISTRY_CAPACITY,
        f"measuredLines={entries_worst[0].get('measuredLines') if entries_worst else '—'}, "
        f"структурных строк {structural}",
    )

    # --- К7. Снятое повисшее двоеточие потерей не считается -------------------
    # Страховка фиттера заканчивается уборкой оборванного ввода: строка,
    # обещающая продолжение двоеточием, снимается, даже когда по высоте всё
    # влезло. Уборка обещания без продолжения — не потеря содержимого (тем же
    # правилом живёт учёт в `ctx.bullets`), а посчитанная как потеря она даёт
    # CRITICAL обоих инспекторов и останавливает выдачу оплаченного отчёта.
    dangling = [
        card(
            1,
            "7 свидетельств (5 негативных) в источниках audit-it.ru, x.com.\n"
            "Что делать: Проверить следующее:",
        )
    ]
    prs_dang, telemetry_dang = render_page(dangling)
    entries_dang = risk_entries(telemetry_dang, 6)
    drawn_dang = drawn_headlines(prs_dang, dangling)
    check(
        "К7а: карточка с повисшим двоеточием нарисована",
        len(drawn_dang) == 1,
        f"нарисовано {len(drawn_dang)} из 1",
    )
    lost_dang = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_dang
    )
    check(
        "К7б: потери при этом нет",
        len(entries_dang) == 1 and lost_dang == 0,
        f"запись: {entries_dang[0] if entries_dang else '—'}",
    )
    codes_dang = inspector_codes(prs_dang, telemetry_dang)
    check(
        "К7в: инспектор геометрии молчит",
        codes_dang == [],
        f"коды: {', '.join(codes_dang) or '—'}",
    )

    # --- К8. Рез по бюджету знаков рендерера слышен ---------------------------
    # Тело карточки режется бюджетом `_clip_structured_bullet(detail, 900)` до
    # отрисовки. Живой построитель в него не упирается (его бюджет — 320), а
    # замороженный пак старого прогона и чужой пак упираются: с июльских
    # карточек уезжало до трети текста, и приёмка об этом не знала.
    huge = [card(1, prose(RENDERER_CHAR_BUDGET + 350))]
    prs_huge, telemetry_huge = render_page(huge)
    entries_huge = risk_entries(telemetry_huge, 6)
    drawn_huge = drawn_headlines(prs_huge, huge)
    check(
        "К8а: карточка нарисована (лист её вмещает)",
        len(drawn_huge) == 1,
        f"нарисовано {len(drawn_huge)} из 1",
    )
    lost_huge = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_huge
    )
    check(
        "К8б: срезанный бюджетом текст посчитан потерей",
        len(entries_huge) == 1 and lost_huge > 0,
        f"запись: {entries_huge[0] if entries_huge else '—'}",
    )
    codes_huge = inspector_codes(prs_huge, telemetry_huge)
    check(
        "К8в: инспектор геометрии даёт CONTENT_DROPPED_BY_RENDERER",
        "CONTENT_DROPPED_BY_RENDERER" in codes_huge,
        f"коды: {', '.join(sorted(set(codes_huge))) or '—'}",
    )

    # --- К9. Понижение кегля — сработавшая защита, и она слышна ---------------
    # Карточке, которой не хватило места на 11 pt, рендерер понижает кегль до
    # подписи и рисует целиком. Содержимое доехало, поэтому это не потеря, — но
    # клиент видит карточку мельче соседних, и приёмка обязана об этом знать:
    # `clipped` по §8 и означает «защита сработала».
    stepped = [
        {
            "headline": f"{WORST_HEADLINE} {i}",
            "status": "Требует подтверждения",
            "detail": f"{WORST_PROBLEM}\n{WORST_ACTION}",
            "tone": "warn",
        }
        for i in range(REGISTRY_CAPACITY)
    ] + [card(4, SUMMARY_DETAIL)]
    prs_step, telemetry_step = render_page(stepped)
    sizes = [
        {float(r.font.size.pt) for p in sh.text_frame.paragraphs for r in p.runs if r.font.size}
        for sh in prs_step.slides[0].shapes
        if getattr(sh, "has_text_frame", False)
    ]
    check(
        "К9а: последней карточке понижен кегль до подписи",
        {9.0} in sizes,
        f"кегли коробов: {[sorted(s) for s in sizes if s]}",
    )
    entries_step = risk_entries(telemetry_step, 6)
    lost_step = sum(
        int(e.get("droppedBullets") or 0) + int(e.get("droppedLines") or 0) for e in entries_step
    )
    check(
        "К9б: содержимое при этом не потеряно",
        len(entries_step) == 1 and lost_step == 0,
        f"потеряно {lost_step}",
    )
    check(
        "К9в: сработавшая защита записана как клип",
        len(entries_step) == 1 and entries_step[0].get("clipped") is True,
        f"запись: {entries_step[0] if entries_step else '—'}",
    )
    codes_step = inspector_codes(prs_step, telemetry_step)
    check(
        "К9г: инспектор зовёт это клипом, а не потерей содержимого",
        "text-clipping" in codes_step and "CONTENT_DROPPED_BY_RENDERER" not in codes_step,
        f"коды: {', '.join(sorted(set(codes_step))) or '—'}",
    )

    # --- К10. Мера карточки проверена прибором, которым не считает ------------
    # Внешняя опора для запаса меры: без неё «сколько обещано, столько и
    # нарисовано» остаётся верным при любом множителе, а лишний запас — это
    # снова треть пустоты на листе и недостача ёмкости. Растровый шаг строки
    # замерен по странице 6 эталона и живёт вне кода меры.
    for label, detail, pill in (
        ("сводной", SUMMARY_DETAIL, "Высокий"),
        ("худшей легальной", f"{WORST_PROBLEM}\n{WORST_ACTION}", "Требует подтверждения"),
    ):
        text_w = CONTENT_W - BADGE_W - 220_000 if pill else int(CONTENT_W * 0.92)
        measured = _card_body_height(detail, text_w, FS_BODY)
        physical = raster_body_height(detail, text_w, FS_BODY)
        ratio = measured / physical
        check(
            f"К10: мера тела {label} карточки не ниже растровой и не выше её на 12 %",
            1.0 <= ratio <= 1.12,
            f"мера {measured}, растровая потребность {physical}, отношение {ratio:.3f}",
        )
    summary_h = card_height("Тема риска 1", SUMMARY_DETAIL, "Высокий")
    check(
        # Число из комментария реестра и из ENGINEERING.md: сводная карточка
        # эталона. Допуск ±5 % — меньше одной строки тела (строка стоит ~21 %
        # высоты такой карточки) и больше любого дрейфа метрик гарнитуры.
        "К10в: сводная карточка стоит столько, сколько записано в реестре",
        abs(summary_h - 906_117) <= 906_117 * 0.05,
        f"высота {summary_h} против 906 117",
    )

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
