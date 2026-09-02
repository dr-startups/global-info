#!/usr/bin/env python3
"""Смок: ширины колонок и высоты строк клиентских таблиц (`_add_search_table`).

`_add_search_table` — единственный отрисовщик таблиц шаблона
`orion_golden_search_table`, то есть всех таблиц деки, кроме матрицы рисков.
Пропорции колонок выбирались по длине первого заголовка, а не по смыслу:
трёхколоночная таблица тем («Тема | Публикаций | Из них нежелательных»)
отдавала тексту 14 % ширины, а двум счётчикам — 86 %. Высота строки мерилась
по одной колонке, поэтому строка объявлялась в 2–4 раза ниже своего
содержимого.

Растровый смок этого не видит по построению: текст, зажатый в узкую колонку,
переносится и остаётся внутри ячейки, за поля не выходит. Поэтому ширины
проверяются здесь — прямо по геометрии таблицы, которую построил рендерер.

Сеть, база и LibreOffice не нужны: строится презентация в памяти, читаются
`table.columns[i].width` и `table.rows[i].height`.

Запуск: python3 renderer/smoke_search_table_layout.py (нужен python-pptx)
"""

from __future__ import annotations

import re
import sys
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from smoke_counters import print_tap_counters  # noqa: E402
from smoke_ts_constants import ts_int  # noqa: E402
from orion_golden_render.common import (  # noqa: E402
    CONTENT_W,
    EMU_PER_INCH,
    EMU_PER_PT,
    SLIDE_H,
    SLIDE_W,
    _Ctx,
)
from orion_golden_render.common import (  # noqa: E402
    _wrapped_line_count,
    get_bullet_measure,
    get_layout_telemetry,
    reset_bullet_measure,
    reset_layout_telemetry,
    text_width_px,
)
from orion_golden_render.visual import (  # noqa: E402
    BADGE_PT,
    CELL_MARGINS_EMU,
    TABLE_MEASURE_KEY_SUFFIX,
    _add_search_table,
    _status_tone,
)
from orion_golden_render.slides import (  # noqa: E402
    SEARCH_TABLE_INTRO_GAP as INTRO_GAP,
    SEARCH_TABLE_INTRO_MAX_H as INTRO_MAX_H,
    _render_slide,
)

# Наборы заголовков, которые сегодня посылают построители секций.
HDR_SERP = ["№", "Ссылка", "Заголовок", "Тип источника", "Оценка"]
HDR_SERP_EXTRA = ["Ссылка", "Заголовок", "Найдено по запросу", "Тип источника", "Оценка"]
HDR_THEMES = ["Тема", "Публикаций", "Из них нежелательных"]
HDR_METRICS = ["Система", "Показатель", "Объём", "Комментарий"]
HDR_COMPLIANCE = ["База данных", "Тип совпадения", "Совпадение по имени", "Статус проверки"]
HDR_FALLBACK = ["Поз.", "Домен", "Заголовок", "Риск"]
HDR_PAIR = ["Параметр", "Значение"]

# Тема на предельной длине схемы (`LinkVerdict.theme`, max 120) — худший
# законный случай текстовой колонки таблицы тем.
THEME_120 = (
    "Совместный бизнес с действующим чиновником регионального уровня "
    "и участие в закупках подведомственных учреждений области"
)
TITLE_SERP = "Как устроена сделка: партнёры, доли и подряды регионального оператора"

#: Опора темы в том виде, в каком её печатает `themeAttentionLine` (`shared.ts`):
#: смежные номера — диапазоном. Литерал повторён здесь потому, что смок
#: питоновский и до TS не дотягивается; формат меняется вместе с построителем.
#: На самом длинном листе номеров четыре, по ёмкости таблицы, и все двузначные.
ROW_ANCHOR = " (строки 17–20)"

#: Опора, которую та же ёмкость даёт в худшем случае: ни один номер не смежен
#: соседнему, и сворачивать нечего. Диапазон вместо неё печатать нельзя — он
#: назвал бы строки 14, 16 и 18, которые тему не несут.
ROW_ANCHOR_WORST = " (строки 13, 15, 17, 19)"


#: Имена субъектов разной длины. Имя входит в абзац страницы выдачи шесть раз —
#: лид плюс пять запросов прогона, — поэтому каждая лишняя буква стоит листу
#: шести знаков, и проверка на единственном имени сторожит одну точку кромки,
#: а не кромку.
INTRO_NAMES = [
    "Глинка Сергей Михайлович",
    "Глинкин Сергей Михайлович",
    "Глинкина Светлана Михайловна",
]

#: Темы каталога `config/finding-themes.ts` целиком — от самой длинной (48
#: знаков) до самой короткой. Тема входит в вывод один раз, но стоит рядом с
#: опорой, и вместе они и решают, каким кеглем нарисован лист.
CATALOGUE_THEMES = [
    "Внимание по линии безопасности / оборонный контур",
    "Политические связи / публичная экспозиция",
    "Финансовые претензии / долговые споры",
    "Криминальные / судебные материалы",
    "Офшоры / корпоративное владение",
    "PEP / RCA / watchlist-сигналы",
    "Семья и деловые связи",
    "Деловой профиль",
]

#: Кегль абзаца страницы выдачи на решётке «тема × имя» — **замер**, а не
#: пожелание.
#:
#: Живой состав абзаца (лид с запросом, пропущенные позиции, вывод по теме с
#: опорой, справка о наборе запросов, «почему важно» и рекомендация) упирается в
#: потолок короба `ctx.body` = 1 000 000 EMU. Пока справка печатала каждый
#: запрос целиком, имя субъекта входило в абзац шесть раз, и 11 pt держали
#: только 8 ячеек из 24 при сворачиваемой опоре и 6 при худшей. Справка стала
#: называть общую часть один раз («…«имя» и он же с добавлением «отзывы», …»),
#: имя входит дважды — и решётка держится целиком.
#:
#: **Замер снят на худшей законной опоре** (`ROW_ANCHOR_WORST`): номера строк
#: выбирает не отчёт, а данные, и ячейка, которую держит только сворачиваемый
#: набор, безопасной не является.
#:
#: **Запас у худшей ячейки — около 32 знаков широкой кириллицы**, то есть
#: кромка не исчезла, а отодвинулась за пределы решётки: любое новое
#: предложение в абзаце снова уронит кегль. Ячейка, изменившая кегль, — повод
#: прочитать замер заново: в худшую сторону значит, что формулировка
#: подорожала, в лучшую — что запас вырос и таблицу пора обновить.
INTRO_FONT_BY_THEME = {
    "Внимание по линии безопасности / оборонный контур": [11.0, 11.0, 11.0],
    "Политические связи / публичная экспозиция": [11.0, 11.0, 11.0],
    "Финансовые претензии / долговые споры": [11.0, 11.0, 11.0],
    "Криминальные / судебные материалы": [11.0, 11.0, 11.0],
    "Офшоры / корпоративное владение": [11.0, 11.0, 11.0],
    "PEP / RCA / watchlist-сигналы": [11.0, 11.0, 11.0],
    "Семья и деловые связи": [11.0, 11.0, 11.0],
    "Деловой профиль": [11.0, 11.0, 11.0],
}

#: Пределы ячеек — **из самих построителей**, а не копией чисел здесь.
#:
#: Пока смок держал свою копию, он мерил худшую строку по одному числу, а
#: построитель резал по другому: подъём предела запроса с 80 до 100 не краснил
#: ни одной проверки, хотя лист при этом выходит за поле на 752 120 EMU. Из
#: расхождения двух копий и растёт ёмкость, выведенная из строки, которой не
#: бывает, — история числа 12 в третий раз. Ёмкость листа читается из реестра
#: тем же приёмом (см. `serp_capacity`).
_BUILDERS = (
    Path(__file__).resolve().parent.parent
    / "src/modules/digital-profile/orion-golden/deck-sections/fragment-builders"
)
SERP_TITLE_MAX_CHARS = ts_int(_BUILDERS / "serp.ts", r"const SERP_TITLE_MAX_CHARS = (\d+);")
SERP_FOUND_BY_MAX_CHARS = ts_int(
    _BUILDERS / "serp.ts", r"export const SERP_FOUND_BY_MAX_CHARS = (\d+);"
)
SERP_ADDRESS_MAX_CHARS = ts_int(
    _BUILDERS / "shared.ts", r"const SERP_ADDRESS_MAX_CHARS = (\d+);"
)

#: Самый широкий знак 9 pt среди тех, что встречаются в адресах корпуса и в
#: русском тексте («W» 12,10 px, «Ю» 12,15 px). Из него и выведен предел полосы.
WIDEST_GLYPH_PX = text_width_px("Ю" * 20, 9) / 20

#: Худший правдоподобный заголовок строки: 95 знаков русской прозы.
TITLE_95 = "Расследование о деловых связях и зарубежных активах предпринимателя в юрисдикциях Европы сегодня"[
    :SERP_TITLE_MAX_CHARS
]
#: Тот же предел, но с патологическим переносом: три слова, каждое чуть шире
#: половины колонки. Больше двух строк — законно, и объявить их обязаны все три.
TITLE_95_WORST_WRAP = ("Ю" * 31 + " ") * 3
TITLE_95_WORST_WRAP = TITLE_95_WORST_WRAP[:SERP_TITLE_MAX_CHARS]
#: Адрес обычным письмом на предельной длине полосы — две строки.
ADDRESS_BAND_PLAIN = (
    "kompromat1.online/articles/364300-byvshij-partner-oligarhov-usmanova-i-ananeva-stal-"
    "figurantom-dela-o-moshennichestve-v-osobo-krupnom-razmere-podrobnosti-materiala-"
    "prodolzhenie-chast-vtoraya-i-tretya-arhivnaya-kopiya-stranicy-2026-goda-dopolnenie"
)
ADDRESS_BAND_PLAIN = ADDRESS_BAND_PLAIN.ljust(SERP_ADDRESS_MAX_CHARS, "x")[
    :SERP_ADDRESS_MAX_CHARS
]

#: Худший законный адрес полосы — **из корпуса прогона**, а не придуманный:
#: адрес инстаграма с процентно закодированными иероглифами в строке
#: параметров, 280 знаков в корпусе, обрезанный печатником по пределу полосы.
#: Три нарисованные строки, и это рядовой случай: длиннее 163 знаков в корпусе
#: 16 адресов из 243 (правило: уникальный host+path без схемы, www. и якоря).
ADDRESS_BAND_WORST = (
    "instagram.com/p/DWnYkpiinGF?__d=1Bug%E8%BA%AB%E4%BB%BD%E8%AF%81%E6%9F%A5%E6%89%8B%E6%9C%BA"
    "%E5%8F%B7%E7%A4%BE%E4%BF%9D%E5%8D%A1%E5%8F%B7%E6%9F%A5%E8%AF%A2%E5%A7%93%E5%90%8D%E6%9F%A5"
    "%E8%AF%A2%E8%BA%AB%E4%BB%BD%E8%AF%81%E5%8F%B7%E7%A0%81%E6%9F%A5%E8%AF%A2%E4%B8%AA%E4%BA%BA"
)
ADDRESS_BAND_WORST = ADDRESS_BAND_WORST[: SERP_ADDRESS_MAX_CHARS - 1] + "…"
#: Самый длинный полный адрес корпуса прогона 72 — 153 знака после разбора.
ADDRESS_CORPUS_MAX = (
    "kompromat1.online/articles/364300-byvshij_partner_oligarhov_usmanova_i_ananeva_stal_"
    "figurantom_dela_o_moshennichestve_v_osobo_krupnom_razmere_podrobnosti"
)

# Модель высоты строки таблицы: межстрочный 1.2, отбивка 6pt.
#
# Кегль — тот, которым красится ячейка (кегль подписи). Пока высота считалась
# при 10pt, а рисовалась при 9, мера и рисование расходились на ступень шкалы.
SERP_PT = 9.0
SERP_LINE_H = int(SERP_PT * EMU_PER_PT * 1.2)
PAD = int(6 * EMU_PER_PT)

def client_name_limit() -> int:
    """Предел длины имени записи — из фильтра инвентаря, а не вторым числом здесь.

    `isNarrativeOrPlaceholderMatchName` объявляет нарративом всё, что длиннее
    этого числа, поэтому в ячейку «Совпадение по имени» доезжают только имена не
    длиннее. Читается сам исходник: вопрос «какое имя бывает самым длинным»
    обязан иметь один ответ в любой момент, а на этом ответе стоят оба потолка
    страниц комплаенса.
    """
    adapter = (
        Path(__file__).resolve().parent.parent
        / "src/modules/digital-profile/services/compliance-inventory-adapter.ts"
    ).read_text(encoding="utf-8")
    found = re.search(r"if \(n\.length > (\d+)\) return true;", adapter)
    if not found:
        raise RuntimeError("в фильтре имён инвентаря не найден предел длины")
    return int(found.group(1))


# Самая длинная законная ячейка статуса: у прогона, чей статус проверки в
# артефактах не зафиксирован, печатается именно она.
STATUS_UNRECORDED = "Не подтверждено (статус в артефактах прогона не зафиксирован)"
#: Имя записи в колонке совпадения — трёхчастное ФИО заглавными: так его
#: печатает OpenSanctions.
COMPLIANCE_NAME = "КИРИЛЛ СЕРГЕЕВИЧ КУЛЕБАКИН"

#: Худшее **законное** значение той же колонки: заглавная кириллица — самая
#: широкая форма записи имени, а длина берётся из клиентского фильтра имён.
#: Остальные колонки этой таблицы смок проверяет именно худшим законным
#: значением, и у имени не должно быть исключения.
#:
#: Число читается из `isNarrativeOrPlaceholderMatchName`, а не хранится копией:
#: **на нём посчитаны оба потолка страниц комплаенса**. Со своей копией смок
#: оставался зелёным при поднятом пороге фильтра — а замер говорит, что уже при
#: 120 знаках пять строк сводки пробивают сцену, то есть ослабление фильтра
#: молча возвращало бы отказ без выхода, ради устранения которого делался шаг.
CLIENT_NAME_LIMIT = client_name_limit()
_NAME_UNIT = "КУЛЕБАКИН КИРИЛЛ СЕРГЕЕВИЧ "
COMPLIANCE_NAME_MAX = (_NAME_UNIT * (CLIENT_NAME_LIMIT // len(_NAME_UNIT) + 1))[:CLIENT_NAME_LIMIT]
COMPLIANCE_ROW = [
    "OpenSanctions",
    "PEP (политически значимое лицо)",
    COMPLIANCE_NAME,
    STATUS_UNRECORDED,
]

#: Карточка записи предельного размера: восемь строк, каждая на своём клампе
#: построителя (алиасы и страны — 200 знаков, сводка — 300).
COMPLIANCE_CARD_MAX = [
    ["Совпадение по имени", COMPLIANCE_NAME_MAX],
    ["Категория", "Связь с санкционным лицом"],
    ["Статус проверки", STATUS_UNRECORDED],
    ["Также числится как", ("Кулебакин К. С.; Kulebakin Kirill; Кулебакін Кирило Сергійович; " * 4)[:200]],
    ["Страны в записи", ("Российская Федерация, Швейцария, Кипр, Соединённое Королевство, " * 4)[:200]],
    ["Даты рождения в записи", "1965-04-12, 1965-04-13, 12.04.1965"],
    ["Сводка записи", ("Запись базы связывает субъекта с санкционным лицом через владение компанией; " * 5)[:300]],
    ["Карточка записи", "https://www.opensanctions.org/entities/ru-inn-504309044808/?utm_source=orion"],
]

#: Карточка из одних обязательных строк — самая дешёвая законная запись.
COMPLIANCE_CARD_MIN = COMPLIANCE_CARD_MAX[:3]

#: Карточка на четыре строки: обязательные плюс самая длинная ячейка (сводка).
#: Ею набирается самый плотный законный лист — четыре слота на запись дешевле
#: девяти, а строки при этом остаются предельными.
COMPLIANCE_CARD_MID = COMPLIANCE_CARD_MAX[:3] + [COMPLIANCE_CARD_MAX[6]]

#: Справка раздела: печатается на последней странице базы, своей полосой.
#: У LexisNexis она **на строку шире**, чем у Dow Jones, — третья строка
#: «Визуальный экспорт». Резерв слотов считается по своей базе, а худший
#: законный лист меряется по самой широкой справке.
COMPLIANCE_CARD_INFO_MAX = [
    [
        "Почему важно",
        "Негативные публикации в базе увеличивают репутационный риск и требуют проверки первоисточников.",
    ],
    [
        "Что сделать",
        "Запросить полную карточку записи LexisNexis, включая связанных лиц (RCA), и проверить первоисточники публикаций.",
    ],
    [
        "Визуальный экспорт",
        "Недоступен в текущем наборе; данные приведены в текстовом виде без потерь.",
    ],
]

#: Вводный абзац страницы карточек — тот, который получает рендерер.
#:
#: Абзац склеивает `composeFindingProse` (`deck-sections/run-deck-build.ts`):
#: вводная фраза страницы, «что обнаружено», а на **не последнем** листе ещё и
#: «почему важно» с «что сделать» — там справки в таблице нет, и дедупликация по
#: предложениям их не снимает. Разница существенная: 351 знак против 560, то
#: есть верх таблицы 1 903 445 против 2 270 000. Мерить ёмкость на короткой
#: фикстуре значило бы мерить не ту страницу, которую бюджет защищает.
COMPLIANCE_CARD_NARRATIVE_LAST = (
    "Страница профиля LexisNexis. Визуальный экспорт страницы в текущем наборе недоступен; "
    "содержимое записи приведено в текстовом виде без потерь. Вторая страница профиля из "
    "отчёта v72 объединена с этой: отдельного содержимого у неё нет.\n"
    "Потенциальное совпадение категории «Связь с санкционным лицом»; совпадение не подтверждено "
    "и требует ручной проверки."
)
COMPLIANCE_CARD_NARRATIVE_CONT = (
    COMPLIANCE_CARD_NARRATIVE_LAST
    + " Негативные публикации в базе увеличивают репутационный риск и требуют проверки "
    "первоисточников.\nЗапросить полную карточку записи LexisNexis, включая связанных лиц (RCA), "
    "и проверить первоисточники публикаций."
)

#: Подпись источника страниц комплаенса — та же строка, что и у построителя.
COMPLIANCE_SOURCE_NOTE = "Источник: комплаенс-базы (существующий контур, без расширения источников)."

#: Худшая законная строка сводной таблицы: предельное имя и самый длинный статус.
COMPLIANCE_SUMMARY_ROW_MAX = [
    "OpenSanctions",
    "Связь с санкционным лицом",
    COMPLIANCE_NAME_MAX,
    STATUS_UNRECORDED,
]

#: Худший законный вводный абзац сводной страницы: четыре базы в перечислении,
#: все три статуса и оговорка про объединённые дубли — плюс рекомендация, которую
#: `composeFindingProse` вклеивает в тот же абзац.
COMPLIANCE_SUMMARY_NARRATIVE_MAX = (
    "Записей, отобранных по имени субъекта в комплаенс-базах: {n} "
    "(Dow Jones, LexisNexis, OpenSanctions и World-Check). "
    "Совпадение по базе не подтверждается автоматически: подтверждено аналитиком — 10, "
    "требует ручной проверки — 25, статус не зафиксирован — 5; повторные записи объединены: 12.\n"
    "Верифицировать каждое потенциальное совпадение вручную: сопоставить идентификаторы "
    "субъекта с записью базы."
)

#: Цвета `_status_tone` по ступеням клиентской шкалы: danger / warn / neutral.
RED_RISK = "B91C1C"
AMBER_OPEN = "C2410C"
SLATE_NEUTRAL = "64748B"

#: Зелёный `_status_tone` по умолчанию — тот самый «всё в порядке».
#: Сравнивается строкой: RGBColor — подкласс tuple, и сравнение с числом
#: всегда ложно, то есть проверка «не зелёный» была бы тождеством.
GREEN_OK = "047857"

#: Поля ячейки python-pptx берутся у самого отрисовщика: второе такое число
#: разошлось бы с ним молча.
EMU_PER_PX = 9525

failures: list[str] = []

#: Счётчик выполненных проверок — раннер обязан видеть, сколько их было.
passed_checks = 0


def check(name: str, ok: bool, detail: str = "") -> None:
    global passed_checks
    print(f"[{'PASS' if ok else 'FAIL'}] {name}" + (f" — {detail}" if detail else ""))
    if ok:
        passed_checks += 1
    else:
        failures.append(name)


def build_table(headers: list[str], rows: list[list[str]]) -> Any:
    """Отрисовать таблицу так, как её рисует рендерер, и отдать её геометрию.

    Полос адреса у отрисовщика больше нет: адрес печатается колонкой «Ссылка».
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, 1, 1)
    _add_search_table(ctx, 1_500_000, headers, rows)
    tables = [sh.table for sh in ctx.slide.shapes if getattr(sh, "has_table", False)]
    if len(tables) != 1:
        raise RuntimeError(f"ожидалась одна таблица, получено {len(tables)}")
    return tables[0]


def widths_of(headers: list[str], rows: list[list[str]]) -> list[int]:
    return [int(c.width) for c in build_table(headers, rows).columns]


def serp_capacity() -> int:
    """Ёмкость листа выдачи — из реестра шаблонов, а не вторым числом здесь.

    Реестр объявлен в TypeScript, поэтому читается сам файл: артефакт приёмки
    отстаёт от кода ровно до следующего прогона ворот, а вопрос «сколько строк
    на листе» обязан иметь один ответ в любой момент.
    """
    registry = (
        Path(__file__).resolve().parent.parent
        / "src/modules/digital-profile/orion-golden/deck-sections/template-registry.ts"
    )
    src = registry.read_text(encoding="utf-8")
    block = re.search(r'"serp-table":\s*\{(.*?)\n  \},', src, re.S)
    if not block:
        raise RuntimeError(f"в {registry.name} не найден блок шаблона serp-table")
    # Ёмкость в реестре не записана числом, а поделена: бюджет строк на высоту
    # худшей законной строки. Читаются оба слагаемых — иначе смок знал бы
    # ёмкость, но не знал бы, из чего она получена, и подмена любого из чисел
    # прошла бы мимо него.
    if not re.search(
        r"maxTableRowsPerSlide:\s*Math\.floor\(\s*SERP_TABLE_ROW_BUDGET_EMU\s*/\s*SERP_TABLE_WORST_ROW_EMU\s*\)",
        block.group(1),
    ):
        raise RuntimeError(
            "в блоке serp-table ёмкость больше не выводится делением бюджета строк "
            "на худшую законную строку — замер ниже проверял бы не то число"
        )
    terms = serp_capacity_terms()
    return terms["SERP_TABLE_ROW_BUDGET_EMU"] // terms["SERP_TABLE_WORST_ROW_EMU"]


def serp_capacity_terms() -> dict[str, int]:
    """Слагаемые вывода ёмкости — из самого реестра, а не числами здесь."""
    registry = (
        Path(__file__).resolve().parent.parent
        / "src/modules/digital-profile/orion-golden/deck-sections/template-registry.ts"
    )
    src = registry.read_text(encoding="utf-8")
    terms: dict[str, int] = {}
    for name in (
        "SERP_TABLE_ROW_BUDGET_EMU",
        "SERP_TABLE_WORST_ROW_EMU",
        "SERP_EXTRA_TABLE_WORST_ROW_EMU",
    ):
        found = re.search(rf"export const {name} = ([0-9_]+);", src)
        if not found:
            raise RuntimeError(f"в {registry.name} не найдено слагаемое вывода ёмкости {name}")
        terms[name] = int(found.group(1).replace("_", ""))
    return terms


def render_page(payload: dict[str, Any]) -> Any:
    """Отрисовать одну страницу рендерером и отдать её ctx.

    Ключ страницы отдаётся контексту так же, как его отдаёт `_draw_deck`: по
    нему построитель находит свою страницу в вердикте меры, и без него мерные
    записи страницы были бы безымянными.
    """
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, 1, 1, slide_key=str(payload.get("slideKey") or ""))
    _render_slide(ctx, payload, {})
    return ctx


def page_shapes(ctx: Any, where: str) -> tuple[Any, Any]:
    """Таблица и белая сцена нарисованной страницы.

    Бюджет листа — низ сцены, а не низ слайда, и меряется он на настоящей
    странице, а не арифметикой по константам.
    """
    table = next((sh for sh in ctx.slide.shapes if getattr(sh, "has_table", False)), None)
    stage = next((sh for sh in ctx.slide.shapes if (sh.name or "").startswith("orion_card_p")), None)
    if table is None or stage is None:
        raise RuntimeError(f"{where}: нет таблицы или белой сцены")
    return table, stage


def render_search_table_page(rows: list[list[str]], intro: str) -> Any:
    """Отрисовать страницу шаблона `orion_golden_search_table` и отдать её ctx.

    Полосы адреса под строкой у страницы больше нет: адрес стоит колонкой, а
    поле `rowAddresses` сборка деки теперь отвергает. Подавать его сюда значило
    бы рисовать страницу, которой не бывает.
    """
    return render_page(
        {
            "slideKey": "p09_ru_serp_table",
            "template": "orion_golden_search_table",
            "title": "Россия — Яндекс: собранная выдача (1/4)",
            "narrative": intro,
            "table": {"headers": HDR_SERP, "rows": rows},
        }
    )


def search_table_page(rows: list[list[str]], intro: str) -> tuple[Any, Any]:
    """Страница шаблона `orion_golden_search_table` целиком.

    Возвращает фигуру таблицы и фигуру белой сцены: бюджет листа — низ сцены,
    а не низ слайда, и проверяется он на настоящей странице, а не арифметикой
    по константам.
    """
    return page_shapes(render_search_table_page(rows, intro), "страница выдачи")


def compliance_summary_page(rows: list[list[str]]) -> tuple[Any, Any]:
    """Страница сводки комплаенса целиком: таблица и белая сцена.

    Бюджет листа меряется на настоящей странице, а не арифметикой по
    константам: высота строки зависит от того, во сколько строк ляжет имя
    записи, а высота вводного абзаца — от того, сколько баз и статусов он
    называет. Абзац здесь худший законный: на нём страница и переполняется.
    """
    ctx = render_page(
        {
            "slideKey": "p33_compliance_toc",
            "template": "orion_golden_search_table",
            "title": "Комплаенс — сводка баз данных",
            "narrative": COMPLIANCE_SUMMARY_NARRATIVE_MAX.format(n=len(rows)),
            "table": {"headers": HDR_COMPLIANCE, "rows": rows},
            "sourceNote": COMPLIANCE_SOURCE_NOTE,
        }
    )
    return page_shapes(ctx, "сводная страница комплаенса")


def compliance_card_page(
    records: list[list[list[str]]],
    info: list[list[str]] | None,
    narrative: str,
) -> tuple[Any, Any]:
    """Страница карточек комплаенса целиком: таблица и белая сцена.

    `records` — строки карточек по записям; каждая запись получает свою
    полосу-заголовок, справка (`info`) — свою. Бюджет листа меряется на
    настоящей странице: мера таблиц у рендерера отсутствует, и ёмкость
    построителя держится этим замером.
    """
    rows: list[list[str]] = []
    groups: list[dict[str, Any]] = []
    for i, rec in enumerate(records):
        groups.append(
            {
                "rowStart": len(rows),
                "rowCount": len(rec),
                "qTag": f"Запись {i + 1} из {len(records)}",
                "queryDisplay": COMPLIANCE_NAME_MAX[:60],
            }
        )
        rows.extend(rec)
    if info:
        groups.append(
            {
                "rowStart": len(rows),
                "rowCount": len(info),
                "qTag": "Справка",
                "queryDisplay": "значение раздела и рекомендации",
            }
        )
        rows.extend(info)
    ctx = render_page(
        {
            "slideKey": "p35_lexis_visual",
            "template": "orion_golden_search_table",
            "title": "Комплаенс — LexisNexis: карточки записей",
            "narrative": narrative,
            "table": {"headers": ["Параметр", "Значение"], "rows": rows, "groups": groups},
            "sourceNote": COMPLIANCE_SOURCE_NOTE,
        }
    )
    return page_shapes(ctx, "страница карточек комплаенса")


def compliance_budget(name: str) -> int:
    """Потолок страницы комплаенса — из построителя, а не вторым числом здесь.

    Построитель объявлен в TypeScript, поэтому читается сам файл: вопрос
    «сколько строк на листе» обязан иметь один ответ в любой момент. Оба листа
    комплаенса объявляют свой потолок там: `CARD_PAGE_SLOTS` — карточки,
    `SUMMARY_PAGE_ROWS` — сводка.
    """
    builder = (
        Path(__file__).resolve().parent.parent
        / "src/modules/digital-profile/orion-golden/deck-sections/fragment-builders/compliance.ts"
    ).read_text(encoding="utf-8")
    found = re.search(rf"const {name} = (\d+);", builder)
    if not found:
        raise RuntimeError(f"в построителе комплаенса нет {name}")
    return int(found.group(1))


def compliance_card_fill(
    record_budget: int, shapes: list[list[list[str]]]
) -> list[list[list[str]]]:
    """Самый плотный законный набор карточек на бюджет из данных форм.

    Формы перебираются от дорогой к дешёвой, как набирает построитель: запись
    едет на лист целиком или не едет вовсе, поэтому остаток меньше самой дешёвой
    карточки просто пропадает. Набор строится **из** бюджета, а не задаётся
    числом: фикстура с фиксированным числом слотов не растёт вместе с потолком и
    его подъёма не замечает.
    """
    recs: list[list[list[str]]] = []
    used = 0
    for shape in shapes:
        cost = len(shape) + 1
        while used + cost <= record_budget:
            recs.append(shape)
            used += cost
    return recs


def intro_box(ctx: Any) -> Any:
    """Короб вводного абзаца страницы выдачи — по имени, которое даёт `ctx.body`."""
    box = next(
        (sh for sh in ctx.slide.shapes if (sh.name or "").startswith("orion_text_body_p")),
        None,
    )
    if box is None:
        raise RuntimeError("на странице выдачи нет вводного абзаца")
    return box


def intro_font_pt(box: Any) -> float | None:
    """Кегль, которым нарисован абзац: `ctx.body` понижает его, когда не влезло."""
    for para in box.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return None


def shares_of(widths: list[int]) -> list[float]:
    return [w / CONTENT_W for w in widths]


def fmt(widths: list[int]) -> str:
    return ", ".join(f"{w} ({w / CONTENT_W:.3f})" for w in widths)


def usable_px(width_emu: int) -> float:
    """Полезная ширина ячейки в пикселях — так же, как её считает перенос.

    Поля ячейки вычитаются, и берётся тот же запас 0.90, что у
    `_wrapped_line_count`: проверка ширины обязана мерить ту ширину, по которой
    рендерер потом объявит высоту.
    """
    return max(1, width_emu - CELL_MARGINS_EMU) / EMU_PER_INCH * 96 * 0.90


def lines_needed(text: str, width_emu: int, pt: float, bold: bool = False) -> int:
    """Во сколько строк ляжет текст в колонке такой ширины.

    Меряется настоящими метриками шрифта, а не моделью `_title_line_estimate`:
    та упирается в потолок «две строки» и на длинной ячейке отвечает «две»
    всегда — то есть проверка через неё была бы тождеством.
    """
    usable = max(1, width_emu - CELL_MARGINS_EMU)
    lines = 1
    cur = ""
    for word in text.split():
        trial = f"{cur} {word}".strip()
        if text_width_px(trial, pt, bold) * EMU_PER_PX <= usable:
            cur = trial
        else:
            lines += 1
            cur = word
    return lines


def main() -> int:
    # --- Т1. Текстовая колонка ведёт таблицу тем ----------------------------
    themes_rows = [[THEME_120, "7", "3"]]
    w_themes = widths_of(HDR_THEMES, themes_rows)
    s_themes = shares_of(w_themes)
    check(
        "Т1а: колонка «Тема» занимает не меньше 55 % ширины",
        s_themes[0] >= 0.55,
        f"доли: {fmt(w_themes)}",
    )
    check(
        "Т1б: колонки-счётчики занимают не больше 22 % каждая",
        max(s_themes[1], s_themes[2]) <= 0.22,
        f"«Публикаций» {s_themes[1]:.3f}, «Из них нежелательных» {s_themes[2]:.3f}",
    )

    # --- Т2. Таблица выдачи: пять колонок, адрес своей колонкой --------------
    #
    # Полос адреса больше нет: адрес вернулся в колонку «Ссылка», и ветка
    # рендерера узнаётся по колонке «№» — у второй таблицы выдачи её нет вовсе
    # по решению владельца, и это самый устойчивый признак разворота.
    serp_rows = [
        ["1", ADDRESS_BAND_PLAIN, TITLE_95, "Официальный сайт / госресурс", "Нежелательный"]
    ]
    w_serp = widths_of(HDR_SERP, serp_rows)
    s_serp = shares_of(w_serp)
    check(
        "Т2а: колонка «Ссылка» занимает не меньше трети ширины",
        s_serp[1] >= 0.33,
        f"доли: {fmt(w_serp)}",
    )
    check(
        "Т2б: номерная колонка «№» не шире 6 % ширины",
        s_serp[0] <= 0.06,
        f"«№» {w_serp[0]} ({s_serp[0]:.3f})",
    )
    # Вторая таблица выдачи — тоже пять колонок, но без «№»: она обязана
    # получить свои доли, иначе адрес молча уедет на 5 % ширины первой колонки.
    extra_rows = [[ADDRESS_BAND_PLAIN, TITLE_95, "запрос", "Официальный сайт / госресурс", "Нежелательный"]]
    s_extra = shares_of(widths_of(HDR_SERP_EXTRA, extra_rows))
    check(
        "Т2в: у второй таблицы выдачи свои доли, а не доли первой",
        abs(s_extra[0] - 0.30) <= 0.005 and abs(s_extra[2] - 0.16) <= 0.005,
        f"вторая таблица: {fmt(widths_of(HDR_SERP_EXTRA, extra_rows))}",
    )
    # --- Т2г. Границы колонки адреса: обычное письмо и худший законный случай -
    #
    # Предел выведен от **узкой** из двух колонок адреса (0.30 листа, 287 px):
    # 7 × 287 / 12,15 = 165 знаков. Семь строк, а не меньше, потому что из
    # семистрочной строки и выведена ёмкость листа. Ниже закреплены обе
    # границы — иначе вывод ёмкости стоял бы на слове.
    usable_band_px = usable_px(widths_of(HDR_SERP_EXTRA, extra_rows)[0])
    band_px = text_width_px(ADDRESS_BAND_PLAIN, 9)
    check(
        "Т2г: адрес обычным письмом на пределе колонки — не больше пяти строк",
        band_px <= 5 * usable_band_px,
        f"{band_px}px при полезной ширине колонки {usable_band_px:.0f}px",
    )
    worst_band_px = text_width_px(ADDRESS_BAND_WORST, 9)
    check(
        "Т2г2: худший законный адрес колонки — не больше семи строк",
        worst_band_px <= 7 * usable_band_px,
        f"{worst_band_px}px при полезной ширине колонки {usable_band_px:.0f}px",
    )
    check(
        "Т2г3: предел адреса гарантирует семь строк самым широким знаком",
        SERP_ADDRESS_MAX_CHARS * WIDEST_GLYPH_PX <= 7 * usable_band_px,
        f"{SERP_ADDRESS_MAX_CHARS} × {WIDEST_GLYPH_PX:.2f}px = "
        f"{SERP_ADDRESS_MAX_CHARS * WIDEST_GLYPH_PX:.0f}px при семи строках "
        f"{7 * usable_band_px:.0f}px",
    )
    title_px = text_width_px(TITLE_95, 9)
    usable_title_px = usable_px(w_serp[2])
    check(
        "Т2д: худший правдоподобный заголовок укладывается в три строки",
        title_px <= 3 * usable_title_px,
        f"{title_px}px при полезной ширине колонки {usable_title_px:.0f}px",
    )
    widest_source_type = "Официальный сайт / госресурс"
    source_px = text_width_px(widest_source_type, 9)
    usable_source_px = usable_px(w_serp[3])
    check(
        "Т2е: самое длинное значение «Тип источника» — одна строка",
        source_px <= usable_source_px,
        f"{source_px}px при полезной ширине {usable_source_px:.0f}px",
    )
    # --- Т2ж. Адрес печатается целиком в своей ячейке -------------------------
    serp_table = build_table(HDR_SERP, serp_rows)
    address_cell = serp_table.cell(1, HDR_SERP.index("Ссылка"))
    check(
        "Т2ж: адрес доезжает до ячейки без реза",
        address_cell.text == ADDRESS_BAND_PLAIN,
        f"в ячейке {len(address_cell.text)} знаков из {len(ADDRESS_BAND_PLAIN)}",
    )
    check(
        "Т2з: адрес стоит колонкой, а не объединённой полосой под строкой",
        not address_cell.is_merge_origin,
        f"ячейка объединяет {getattr(address_cell, 'span_width', 1)} колонок",
    )
    # --- Т2и. Высота строки без потолка «две строки» ---------------------------
    #
    # `_title_line_estimate` объявляла две строки при любом содержимом: строка
    # в три нарисованных строки объявлялась двумя, LibreOffice тянул её по
    # содержимому, и таблица уезжала ниже поля при «чистой» геометрии.
    worst_wrap = build_table(
        HDR_SERP,
        [["1", ADDRESS_CORPUS_MAX, TITLE_95_WORST_WRAP, "СМИ", "Нейтральный"]],
    )
    worst_row_h = int(worst_wrap.rows[1].height)
    check(
        "Т2и: ячейка в три нарисованных строки объявлена тремя",
        worst_row_h >= 3 * SERP_LINE_H + PAD,
        f"высота {worst_row_h}, нужно ≥ {3 * SERP_LINE_H + PAD}",
    )

    # --- Т3. Высота строки — по самой высокой ячейке ------------------------
    themes_table = build_table(HDR_THEMES, themes_rows)
    data_row_h = int(themes_table.rows[1].height)
    check(
        "Т3: строка со 120-символьной темой объявлена не ниже двух строк текста",
        data_row_h >= 2 * SERP_LINE_H + PAD,
        f"высота {data_row_h}, нужно ≥ {2 * SERP_LINE_H + PAD} (тема {len(THEME_120)} симв.)",
    )

    # --- Т4. Таблицы вне задачи не тронуты ----------------------------------
    # Метрики региона сохраняют сегодняшние доли: их пропорции тоже кривые, но
    # это отдельный пункт бэклога, а не эта правка.
    expected_four = [0.14, 0.26, 0.42, 0.18]
    for label, headers, rows in (
        (
            "метрики региона",
            HDR_METRICS,
            [["Яндекс", "Найдено страниц", "312", "Проверено вручную 20 первых"]],
        ),
    ):
        w = widths_of(headers, rows)
        s = shares_of(w)
        ok = len(s) == 4 and all(abs(s[i] - expected_four[i]) <= 0.005 for i in range(4))
        check(
            f"Т4: доли таблицы «{label}» не изменились",
            ok,
            f"{fmt(w)}",
        )

    # --- Т5. Двухколоночная таблица ------------------------------------------
    w_pair = widths_of(HDR_PAIR, [["Регион выдачи", "Российская Федерация"]])
    s_pair = shares_of(w_pair)
    check(
        "Т5: две колонки делятся как 0.24 / 0.76",
        len(s_pair) == 2
        and abs(s_pair[0] - 0.24) <= 0.005
        and abs(s_pair[1] - 0.76) <= 0.005,
        f"{fmt(w_pair)}",
    )

    # --- Т6. Номерная колонка узнаётся по смыслу ------------------------------
    w_fallback = widths_of(HDR_FALLBACK, [["1", "kommersant.ru", TITLE_SERP, "Высокий"]])
    check(
        "Т6: номерная колонка «Поз.» не шире 8 % ширины",
        w_fallback[0] <= 0.08 * CONTENT_W,
        f"«Поз.» {w_fallback[0]} ({w_fallback[0] / CONTENT_W:.3f})",
    )

    # --- Т7. Остаток не теряется и не утекает ---------------------------------
    for label, widths in (
        ("выдача ТОП-20", w_serp),
        ("темы публикаций", w_themes),
        ("метрики региона", widths_of(HDR_METRICS, [["Яндекс", "Найдено страниц", "312", "—"]])),
        (
            "комплаенс",
            widths_of(HDR_COMPLIANCE, [["OpenSanctions", "Имя", COMPLIANCE_NAME, "Требует проверки"]]),
        ),
        ("запасной разбор", w_fallback),
        ("две колонки", w_pair),
    ):
        total = sum(widths)
        check(
            f"Т7: ширины таблицы «{label}» дают ровно ширину контента",
            total == CONTENT_W,
            f"сумма {total}, ширина контента {CONTENT_W}",
        )

    # --- Т8. Комплаенс-сводка: своя ветка долей -------------------------------
    # Общая четырёхколоночная ветка отдавала третьей колонке 42 %, а «Статусу
    # проверки» — 18 % под строку из шести слов. Доли выверены настоящими
    # метриками шрифта: самая длинная законная ячейка статуса — «Не
    # подтверждено (статус в артефактах прогона не зафиксирован)», а в третьей
    # колонке теперь стоит имя записи, а не «71/100», и ему нужна ширина под
    # трёхчастное ФИО в одну строку.
    w_comp = widths_of(HDR_COMPLIANCE, [COMPLIANCE_ROW])
    s_comp = shares_of(w_comp)
    expected_compliance = [0.14, 0.26, 0.26, 0.34]
    check(
        "Т8а: доли комплаенс-сводки — своя ветка 0.14 / 0.26 / 0.26 / 0.34",
        len(s_comp) == 4
        and all(abs(s_comp[i] - expected_compliance[i]) <= 0.005 for i in range(4)),
        f"{fmt(w_comp)}",
    )
    check(
        "Т8б: «Статус проверки» шире колонки имени",
        w_comp[3] > w_comp[2],
        f"статус {w_comp[3]}, имя {w_comp[2]}",
    )
    # Строка статуса рисуется бейджем: «● » впереди и кегль 9.5.
    status_lines = lines_needed(f"● {STATUS_UNRECORDED}", w_comp[3], 9.5)
    check(
        "Т8в: самая длинная законная ячейка статуса укладывается в две строки",
        status_lines <= 2,
        f"{status_lines} строк(и) при ширине {w_comp[3]}",
    )
    for i, (label, text, bold) in enumerate(
        (
            ("База данных", "OpenSanctions", False),
            ("Тип совпадения", "Импортированный отчёт LexisNexis", False),
            ("Совпадение по имени", HDR_COMPLIANCE[2], True),
        )
    ):
        n = lines_needed(text, w_comp[i], 10.0, bold)
        check(
            f"Т8г: колонка «{label}» держит «{text}» в одну строку",
            n == 1,
            f"{n} строк(и) при ширине {w_comp[i]}",
        )

    # Колонка имени затем и заведена, чтобы читатель увидел, на кого запись:
    # перенос ФИО на вторую строку её не ломает, но одна строка — тот запас,
    # ради которого доли и пересчитаны.
    name_lines = lines_needed(COMPLIANCE_NAME, w_comp[2], 10.0)
    check(
        "Т8д: трёхчастное ФИО заглавными укладывается в одну строку",
        name_lines == 1,
        f"{name_lines} строк(и) при ширине {w_comp[2]}",
    )
    # Имя предельной длины **не режется**: обрезанное имя — это другое имя, то
    # есть ложное утверждение о том, на кого запись. Оно занимает столько
    # строк, сколько занимает, и цена этого — высота строки таблицы; предел
    # ниже и проверяет, что цена посчитана, а не забыта.
    #
    # Прежний Т8д1 («имя ровно 90 знаков») снят: имя строится из предела фильтра,
    # и проверка стала бы утверждением о собственной фикстуре. Ослабление фильтра
    # ловится по существу — здесь и в Т8е/Т8ж, где то же имя стоит в измеряемой
    # строке.
    max_name_lines = lines_needed(COMPLIANCE_NAME_MAX, w_comp[2], 10.0)
    check(
        f"Т8д2: имя предельной длины ({CLIENT_NAME_LIMIT} знаков — предел фильтра "
        f"инвентаря) занимает не больше четырёх строк",
        max_name_lines <= 4,
        f"{max_name_lines} строк(и) при ширине {w_comp[2]}",
    )

    # --- Т8е. Бюджет листа сводки ---------------------------------------------
    #
    # У сводной таблицы комплаенса свой потолок строк (`SUMMARY_PAGE_ROWS` в
    # построителе): клип на странице комплаенса останавливает выдачу целиком,
    # поэтому страница без потолка — это отказ, из которого оплаченный прогон не
    # выходит ни пересборкой, ни повтором рендера. Лист строится **из** числа
    # построителя и меряется настоящей страницей: низ таблицы против низа белой
    # сцены, худшая законная строка и худший законный вводный абзац.
    summary_rows = compliance_budget("SUMMARY_PAGE_ROWS")
    table_max, stage_max = compliance_summary_page([COMPLIANCE_SUMMARY_ROW_MAX] * summary_rows)
    bottom_max = int(table_max.top) + int(table_max.height)
    stage_bottom = int(stage_max.top) + int(stage_max.height)
    margin_summary = stage_bottom - bottom_max
    check(
        f"Т8е1: худший законный лист сводки ({summary_rows} строк) влезает в сцену с запасом",
        bottom_max <= stage_bottom and margin_summary >= 400_000,
        f"низ таблицы {bottom_max}, низ сцены {stage_bottom}, запас {margin_summary} EMU",
    )
    # Число стоит там, где надо: строкой больше — и лист пробивает сцену.
    # Без этой половины проверка ловила бы только заведомо огромный потолок.
    over_table, over_stage = compliance_summary_page(
        [COMPLIANCE_SUMMARY_ROW_MAX] * (summary_rows + 1)
    )
    over_bottom = int(over_table.top) + int(over_table.height)
    over_stage_bottom = int(over_stage.top) + int(over_stage.height)
    check(
        f"Т8е2: {summary_rows + 1}-я строка сводки на лист уже не помещается",
        over_bottom > over_stage_bottom,
        f"низ таблицы {over_bottom}, низ сцены {over_stage_bottom}",
    )

    # --- Т8ж. Бюджет листа карточек комплаенса --------------------------------
    #
    # Ёмкость страницы карточек считается в **слотах** (строка таблицы или
    # полоса-заголовок), и число живёт в построителе. Здесь оно сверяется с
    # геометрией настоящей страницы: обрезанная строка карточки останавливает
    # выдачу целиком, поэтому промах числа стоит не косметики, а отказа прогона.
    #
    # Лист набирается **из** числа построителя тремя законными способами: слот
    # не различает высоту строки, и самый плотный лист — не всегда самый
    # высокий. Резерв под справку берётся по самой широкой из них (LexisNexis,
    # три строки плюс полоса), потому что худший лист бывает именно там.
    slots = compliance_budget("CARD_PAGE_SLOTS")
    card_budget = slots - (len(COMPLIANCE_CARD_INFO_MAX) + 1)
    for label, shapes in (
        ("предельные, затем короткие", [COMPLIANCE_CARD_MAX, COMPLIANCE_CARD_MID, COMPLIANCE_CARD_MIN]),
        ("четырёхстрочные", [COMPLIANCE_CARD_MID]),
        ("минимальные", [COMPLIANCE_CARD_MIN]),
    ):
        recs = compliance_card_fill(card_budget, shapes)
        # Последний лист базы: справка внизу таблицы, из абзаца она вычищена.
        last_table, last_stage = compliance_card_page(
            recs, COMPLIANCE_CARD_INFO_MAX, COMPLIANCE_CARD_NARRATIVE_LAST
        )
        last_bottom = int(last_table.top) + int(last_table.height)
        last_stage_bottom = int(last_stage.top) + int(last_stage.height)
        last_slots = sum(len(r) + 1 for r in recs) + len(COMPLIANCE_CARD_INFO_MAX) + 1
        last_margin = last_stage_bottom - last_bottom
        check(
            f"Т8ж1 [{label}]: последний лист базы ({last_slots} слотов) влезает в сцену с запасом",
            last_bottom <= last_stage_bottom and last_margin >= 400_000,
            f"низ таблицы {last_bottom}, низ сцены {last_stage_bottom}, "
            f"запас {last_margin} EMU при бюджете {slots} слотов",
        )
        # Лист-продолжение: справки в таблице нет, поэтому «Почему важно» и
        # «Что сделать» остаются во вводном абзаце и он вырастает вдвое.
        cont_table, cont_stage = compliance_card_page(
            recs, None, COMPLIANCE_CARD_NARRATIVE_CONT
        )
        cont_bottom = int(cont_table.top) + int(cont_table.height)
        cont_stage_bottom = int(cont_stage.top) + int(cont_stage.height)
        cont_margin = cont_stage_bottom - cont_bottom
        check(
            f"Т8ж2 [{label}]: лист-продолжение с длинным абзацем влезает в сцену с запасом",
            cont_bottom <= cont_stage_bottom and cont_margin >= 400_000,
            f"низ таблицы {cont_bottom}, низ сцены {cont_stage_bottom}, "
            f"запас {cont_margin} EMU при бюджете {slots} слотов",
        )
    two_table, two_stage = compliance_card_page(
        [COMPLIANCE_CARD_MAX, COMPLIANCE_CARD_MAX], None, COMPLIANCE_CARD_NARRATIVE_CONT
    )
    check(
        "Т8ж3: две записи предельного размера на одном листе не помещаются",
        int(two_table.top) + int(two_table.height)
        > int(two_stage.top) + int(two_stage.height),
        f"низ таблицы {int(two_table.top) + int(two_table.height)}, "
        f"низ сцены {int(two_stage.top) + int(two_stage.height)}",
    )

    # --- Т9. Бейдж статуса ----------------------------------------------------
    # Незнакомое слово в `_status_tone` зелёное по умолчанию, и подтверждённый
    # комплаенс-риск рисовался цветом «всё в порядке».
    _, tone_confirmed = _status_tone("Подтверждено аналитиком")
    check(
        "Т9а: «Подтверждено аналитиком» — не зелёный бейдж",
        str(tone_confirmed) != GREEN_OK,
        f"цвет {tone_confirmed}",
    )
    _, tone_unrecorded = _status_tone(STATUS_UNRECORDED)
    check(
        "Т9б: незафиксированный статус — не зелёный бейдж",
        str(tone_unrecorded) != GREEN_OK,
        f"цвет {tone_unrecorded}",
    )
    # Клиентская шкала печатает три слова, и бейдж обязан их знать: незнакомое
    # слово здесь зелёное по умолчанию, то есть «Высокий» читался бы как «всё в
    # порядке».
    steps = {word: str(_status_tone(word)[1]) for word in ("Высокий", "Средний", "Низкий")}
    check(
        "Т9в: ступени шкалы не зелёные",
        all(colour != GREEN_OK for colour in steps.values()),
        f"цвета {steps}",
    )
    check(
        "Т9г: цвет ступени растёт вместе со ступенью",
        steps["Высокий"] == RED_RISK and steps["Средний"] == AMBER_OPEN and steps["Низкий"] == SLATE_NEUTRAL,
        f"цвета {steps}",
    )

    # Зелёный — цвет явного «всё в порядке», и достаться он должен только
    # слову, которое это значит. Незнакомое слово зелёным быть не может:
    # статус, которого функция не знает, — не благополучие, а неизвестность, и
    # в документе для банка он читался бы как одобрение. Живой прогон 22.08:
    # все шесть встреченных значений попали в названные ветки, то есть зелёный
    # достижим сегодня **только** незнакомым словом (пункт V).
    unknown = {
        word: str(_status_tone(word)[1])
        for word in ("Проверено", "Совпадений не найдено", "Ожидает выгрузки", "Clean")
    }
    check(
        "Т9д: незнакомый статус не зелёный",
        all(colour != GREEN_OK for colour in unknown.values()),
        f"цвета {unknown}",
    )
    # «Не проверено» — третье значение колонки оценки, и читаться как
    # «Нейтральный» оно не должно: серым красятся оба, а это разные вещи —
    # проверенная страница и незакрытый вопрос. Янтарь тут по той же причине,
    # что у «Не подтверждено».
    _, tone_unverified = _status_tone("Не проверено")
    check(
        "Т9ж: «Не проверено» — янтарный, а не серый «Нейтрального»",
        str(tone_unverified) == AMBER_OPEN,
        f"цвет {tone_unverified}, у «Нейтрального» {_status_tone('Нейтральный')[1]}",
    )
    # У зелёного есть законный владелец: «Позитивный» стоит в легенде таблицы.
    # Пока он доставался по умолчанию, легенда обещала маркер, которого функция
    # не знала по имени.
    _, tone_positive = _status_tone("Позитивный")
    check(
        "Т9е: «Позитивный» — зелёный по имени, а не по умолчанию",
        str(tone_positive) == GREEN_OK,
        f"цвет {tone_positive}",
    )

    # --- Т10. Подсветка строки — только от статусной колонки ------------------
    # В карточке записи последняя колонка называется «Значение». Пока подсветка
    # смотрела на её текст, янтарь доставался строке «Категория: Требует ручной
    # классификации», а строка настоящего статуса оставалась белой.
    card = build_table(
        HDR_PAIR,
        [
            ["Совпадение по имени", "Глинка Сергей Михайлович"],
            ["Категория", "Требует ручной классификации"],
            ["Статус проверки", STATUS_UNRECORDED],
        ],
    )
    card_fills = {str(card.cell(r, 1).fill.fore_color.rgb) for r in range(1, 4)}
    check(
        "Т10а: строки карточки записи не подсвечиваются по колонке «Значение»",
        card_fills == {"FFFFFF"},
        f"фоны строк: {sorted(card_fills)}",
    )
    serp = build_table(
        HDR_SERP, [["1", ADDRESS_CORPUS_MAX, TITLE_SERP, "СМИ", "Нежелательный"]]
    )
    check(
        "Т10б: нежелательная строка выдачи подсветку сохраняет",
        str(serp.cell(1, 1).fill.fore_color.rgb) != "FFFFFF",
        f"фон строки: {serp.cell(1, 1).fill.fore_color.rgb}",
    )

    # --- Т11. Ёмкость листа выдачи и контрольный дефект ----------------------
    #
    # Меры таблиц у рендерера нет, поэтому ёмкость выведена из бюджета листа и
    # худшей законной пары «строка результата + полоса адреса». Арифметика
    # считается по **объявленному** верху таблицы: заголовок страницы плюс
    # потолок вводного абзаца (`SEARCH_TABLE_INTRO_MAX_H`, читается из
    # `slides.py`) плюс отбивка. Считать от фактического верха нельзя — он
    # зависит от длины абзаца и замолчит ровно тогда, когда нужен.
    #
    # Все слагаемые меряются самим рендерером, а не повторяются здесь числами:
    # низ заголовка — верх таблицы на странице без абзаца; высота шапки и пары
    # — строки настоящей таблицы; низ бюджета — низ белой сцены.
    capacity = serp_capacity()
    terms = serp_capacity_terms()
    # Худшая законная строка — не правдоподобная, а предельная: адрес на своём
    # пределе, написанный самым широким знаком 9 pt. Полосы под строкой больше
    # нет, поэтому и пары больше нет: высоту строки задаёт её самая высокая
    # ячейка, и это колонка адреса.
    worst_row = [
        str(1),
        "Ю" * SERP_ADDRESS_MAX_CHARS,
        TITLE_95,
        "Официальный сайт / госресурс",
        "Нежелательный",
    ]
    bare_table, stage_shape = search_table_page([worst_row], "")
    stage_bottom = int(stage_shape.top) + int(stage_shape.height)
    title_bottom = int(bare_table.top)
    declared_top = title_bottom + INTRO_MAX_H + INTRO_GAP
    header_h = int(bare_table.table.rows[0].height)
    row_h = int(bare_table.table.rows[1].height)
    budget = stage_bottom - declared_top - header_h
    # Главная сверка: нарисованная худшая строка равна **объявленной** в
    # реестре. Разойдясь, эти два числа дали бы ёмкость, выведенную из высоты,
    # которой не бывает, — ровно историю числа 12.
    check(
        "Т11а: нарисованная худшая законная строка равна объявленной в реестре",
        row_h == terms["SERP_TABLE_WORST_ROW_EMU"],
        f"нарисовано {row_h}, объявлено {terms['SERP_TABLE_WORST_ROW_EMU']}",
    )
    check(
        f"Т11б: {capacity} худших строк влезают в бюджет объявленного листа",
        capacity * row_h <= budget,
        f"нужно {capacity * row_h}, бюджет {budget} (низ сцены {stage_bottom}, "
        f"объявленный верх таблицы {declared_top}, шапка {header_h})",
    )
    check(
        f"Т11в: контрольный дефект — {capacity + 1} строк в бюджет не влезают",
        (capacity + 1) * row_h > budget,
        f"нужно {(capacity + 1) * row_h}, бюджет {budget}",
    )
    # Вторая таблица выдачи выводит свою ёмкость тем же делением и из своей
    # худшей строки. Без этой сверки её число не держало **ничто**: юнит вывода
    # читал только первую таблицу, а `serp_capacity()` искал её блок регуляркой.
    extra_worst = build_table(
        HDR_SERP_EXTRA,
        [
            [
                "Ю" * SERP_ADDRESS_MAX_CHARS,
                "Ю" * SERP_TITLE_MAX_CHARS,
                "Ю" * SERP_FOUND_BY_MAX_CHARS,
                "Официальный сайт / госресурс",
                "Нежелательный",
            ]
        ],
    )
    extra_row_h = int(extra_worst.rows[1].height)
    extra_capacity = terms["SERP_TABLE_ROW_BUDGET_EMU"] // terms["SERP_EXTRA_TABLE_WORST_ROW_EMU"]
    check(
        "Т11в2: нарисованная худшая строка второй таблицы равна объявленной",
        extra_row_h == terms["SERP_EXTRA_TABLE_WORST_ROW_EMU"],
        f"нарисовано {extra_row_h}, объявлено {terms['SERP_EXTRA_TABLE_WORST_ROW_EMU']}",
    )
    check(
        f"Т11в3: {extra_capacity} худших строк второй таблицы влезают, {extra_capacity + 1} — нет",
        extra_capacity * extra_row_h <= budget and (extra_capacity + 1) * extra_row_h > budget,
        f"нужно {extra_capacity * extra_row_h} и {(extra_capacity + 1) * extra_row_h}, бюджет {budget}",
    )
    # Второй свидетель: настоящая страница с заведомо переполняющим абзацем.
    # Её верх обязан быть не ниже объявленного потолка — иначе это не потолок,
    # а ёмкость листа, выведенная от него, ничем не держится. Счётчика
    # предложений в шаблоне больше нет, поэтому абзац подаётся заведомо
    # длиннее любого законного: за `max_h` не может выйти и он.
    intro = " ".join(
        [
            f"Предложение номер {n} этого абзаца написано нарочно длинным, чтобы "
            "вводный текст страницы заведомо не поместился в объявленный потолок "
            "и упёрся в него мерой, а не счётчиком предложений."
            for n in range(1, 13)
        ]
    )
    real_table, real_stage = search_table_page([worst_row] * capacity, intro)
    check(
        "Т11г: фактический верх таблицы не ниже объявленного потолка",
        int(real_table.top) <= declared_top,
        f"фактический {real_table.top}, объявленный {declared_top}",
    )
    check(
        "Т11д: настоящая страница из худших пар кончается не ниже низа сцены",
        int(real_table.top) + int(real_table.height)
        <= int(real_stage.top) + int(real_stage.height),
        f"низ таблицы {int(real_table.top) + int(real_table.height)}, "
        f"низ сцены {int(real_stage.top) + int(real_stage.height)}",
    )

    # --- Т11е/Т11ж. Абзац страницы выдачи нарисован целиком ------------------
    #
    # Составы собраны из того, что печатают сами построители, и подаются одной
    # строкой — ровно так их видит `ctx.body` (`_safe` схлопывает переводы
    # строк). Проверяется текст листа, а не мера: вопрос «напечаталось ли»
    # задан не тому, кто резал.
    MISSING = (
        "Позиции 1–3, 5, 11–20 в собранных данных отсутствуют: эти строки потеряны "
        "при сборе, а не пусты в выдаче."
    )
    WHY = (
        "На странице 1 тема повышенного внимания — эти материалы видны при первой же "
        "проверке субъекта."
    )
    TAIL = "Запросить первичные карточки баз и подтвердить принадлежность совпадений."

    def live_intro_for(
        name: str, theme: str, anchor: str = ROW_ANCHOR, tail: str = TAIL
    ) -> str:
        """Живой состав абзаца страницы выдачи для этого субъекта и этой темы.

        Корпусный состав плюс два предложения, которых нет на эталоне (там
        запрос выдачи в артефактах не записан) — о пропущенных позициях и о
        наборе запросов прогона. Справка о наборе повторена здесь литералом:
        смок питоновский и до TS не дотягивается, формат меняется вместе с
        `subjectQueriesLine`. Имя субъекта входит в абзац **дважды** — лид и
        база набора запросов; до сокращения справки входило шесть раз.
        """
        low = name.lower()
        return " ".join(
            [
                f"Показана выдача Яндекса по запросу «{name}».",
                MISSING,
                f"«{theme}» — средний уровень внимания{anchor}.",
                f"Выдача проверена по 5 запросам: «{low}» и он же с добавлением "
                "«бизнес», «компромат», «отзывы», «суд».",
                WHY,
                tail,
            ]
        )

    # Живой состав абзаца — на решётке «тема каталога × имя субъекта», а не на
    # одном имени: прежняя проверка стояла ровно на кромке (эталонное имя, тема
    # в 29 знаков) и о том, что имя на букву длиннее роняет лист на девятку, не
    # знала. Ожидания по кеглю — замер, объявленный в `INTRO_FONT_BY_THEME`;
    # нарисован абзац обязан быть целиком в каждой ячейке.
    grid_bad: list[str] = []
    grid_cut: list[str] = []
    for theme, expected in INTRO_FONT_BY_THEME.items():
        for name, want in zip(INTRO_NAMES, expected):
            text = live_intro_for(name, theme, ROW_ANCHOR_WORST)
            drawn = intro_box(
                render_search_table_page([worst_row] * 2, text)
            )
            got = intro_font_pt(drawn)
            cell = f"«{theme[:24]}…» + «{name}»"
            if got != want:
                grid_bad.append(f"{cell}: {got} вместо {want} ({len(text)} знаков)")
            if TAIL[-40:] not in drawn.text_frame.text:
                grid_cut.append(f"{cell}: {len(drawn.text_frame.text)} из {len(text)}")
    check(
        "Т11е1: живой состав абзаца нарисован целиком на всей решётке «тема × имя»",
        not grid_cut,
        "; ".join(grid_cut),
    )
    check(
        "Т11е2: кегль абзаца на решётке «тема × имя» совпадает с объявленным замером",
        not grid_bad,
        "; ".join(grid_bad),
    )

    # Самая дорогая ячейка решётки — прицельно и обеими записями опоры.
    #
    # Решётка выше уже гоняет худшую опору, но проверка, названная своим
    # именем, нужна отдельно: она называет **ту самую** ячейку, на которой лист
    # ломался раньше всех, — самая длинная тема каталога и самое длинное имя, —
    # и печатает её кегли обеими записями опоры. Пока справка о наборе запросов
    # повторяла имя субъекта пять раз, эта ячейка шла девяткой при любой опоре.
    WORST_CELL_THEME = CATALOGUE_THEMES[0]
    WORST_CELL_NAME = INTRO_NAMES[-1]

    def worst_cell_font(anchor: str) -> float | None:
        text = live_intro_for(WORST_CELL_NAME, WORST_CELL_THEME, anchor)
        return intro_font_pt(
            intro_box(
                render_search_table_page([worst_row] * 2, text)
            )
        )

    folded, unfolded = worst_cell_font(ROW_ANCHOR), worst_cell_font(ROW_ANCHOR_WORST)
    check(
        "Т11е3: самая дорогая ячейка решётки держит 11 pt обеими записями опоры",
        (folded, unfolded) == (11, 11),
        f"диапазоном {folded}, перечнем {unfolded} "
        f"(тема «{WORST_CELL_THEME}», имя «{WORST_CELL_NAME}»)",
    )

    # Худший состав по **схемным пределам**, а не по корпусу: тема на пределе
    # `LinkVerdict.theme` (120 знаков) и рекомендация на клампе `whatToCheck`
    # (220). Кегль здесь уже девятка — это известно и объявлено; терять
    # предложения абзац по-прежнему не вправе.
    worst_intro = live_intro_for(
        INTRO_NAMES[0],
        THEME_120,
        tail=(
            "Запросить первичные карточки баз, подтвердить принадлежность совпадений "
            "по идентификаторам субъекта, сверить заголовки публикаций с "
            "первоисточниками и зафиксировать результат сверки в карточке проверки."
        ),
    )
    worst_text = intro_box(
        render_search_table_page([worst_row] * 2, worst_intro)
    ).text_frame.text
    check(
        "Т11ж: состав по схемным пределам нарисован целиком",
        "зафиксировать результат сверки в карточке проверки" in worst_text,
        f"нарисовано {len(worst_text)} из {len(worst_intro)} знаков: {worst_text[-100:]!r}",
    )

    # --- Т12 снята вместе с полосой адреса ------------------------------------
    #
    # Проверяла, что полоса под строкой наследует её фон: белая полоса
    # разрезала красную плашку пополам, и адрес нежелательного материала
    # оказывался вне красного. Полосы больше нет — адрес стоит ячейкой той же
    # строки и её фон получает сам, — поэтому у проверки не осталось предмета.
    # Подсветку самой строки держит Т10б.

    # --- Т13. Мера ячейки совпадает с тем, что в ней нарисовано --------------
    #
    # Бейдж красится как «● Нежелательный», а мерился по «Нежелательный» —
    # на два знака короче. Та же щель «мерим одно, рисуем другое», которую
    # закрыли в остальных колонках.
    # Контрольная строка подобрана по метрикам шрифта: без точки 130 px (одна
    # строка колонки оценки), с точкой 146 px при полезных 145 — две. Всё
    # остальное на странице заведомо ниже, чтобы высоту решал именно бейдж.
    BOUNDARY_STATUS = "проверки аналитика"
    badge_table = build_table(
        HDR_SERP, [["1", "a.example.org/1", "Кратко", "СМИ", BOUNDARY_STATUS]]
    )
    badge_cell = badge_table.cell(1, 3)
    # Меряем **нарисованное**: тем же переносом, которым рендерер объявляет
    # высоту, и по тексту ячейки, а не по значению из строки.
    badge_lines = _wrapped_line_count(
        badge_cell.text, int(badge_table.columns[3].width) - CELL_MARGINS_EMU, BADGE_PT
    )
    check(
        "Т13: строка объявлена не ниже, чем требует нарисованный бейдж",
        int(badge_table.rows[1].height) >= badge_lines * int(BADGE_PT * EMU_PER_PT * 1.2) + PAD,
        f"высота {badge_table.rows[1].height}, бейджу «{badge_cell.text}» нужно "
        f"{badge_lines} строк(и) по {int(BADGE_PT * EMU_PER_PT * 1.2)}",
    )

    # --- Т14. Обе пятиколоночные таблицы рисуются без отказа -----------------
    #
    # `slides.py` пропускает до пяти заголовков, и у обеих таблиц выдачи их
    # ровно пять, но составы разные. Ветка первой узнаётся по «№», ветка второй
    # — общая на пять колонок; пока список долей был четырёхэлементным, пятая
    # колонка роняла страницу `IndexError`ом при подсчёте высот.
    for label, headers, row in (
        ("первая", HDR_SERP, ["1", "example.org/1", TITLE_SERP, "СМИ", "Нейтральный"]),
        ("вторая", HDR_SERP_EXTRA, ["example.org/1", TITLE_SERP, "запрос", "СМИ", "Нейтральный"]),
    ):
        try:
            five = build_table(headers, [row])
            five_ok = len(list(five.columns)) == 5
            five_detail = f"долей {len(list(five.columns))}"
        except Exception as exc:  # noqa: BLE001
            five_ok = False
            five_detail = f"{type(exc).__name__}: {exc}"
        check(f"Т14 [{label} таблица]: пять колонок рисуются без отказа", five_ok, five_detail)

    # --- Т15. Бюджет листа объявляется только тогда, когда он известен -------
    #
    # `content_stage` возвращает низ сцены, а при отказе рисовать — свой вход.
    # Одно значение на два вопроса: поданное как бюджет, сентинельное значение
    # даёт ложный CRITICAL с нулевым запасом.
    reset_layout_telemetry()
    build_table(HDR_SERP, [worst_row])
    check(
        "Т15а: без объявленного бюджета таблица о переполнении не заявляет",
        get_layout_telemetry() == [],
        f"записей телеметрии {len(get_layout_telemetry())}",
    )
    reset_layout_telemetry()
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    ctx = _Ctx(prs, 1, 1)
    _add_search_table(ctx, 1_500_000, HDR_SERP, [worst_row] * 20, bottom=2_000_000)
    entries = get_layout_telemetry()
    check(
        "Т15б: превышение объявленного бюджета слышно записью разметки",
        len(entries) == 1 and entries[0].get("clipped") is True and entries[0]["role"] == "table",
        f"записей {len(entries)}: {entries[:1]}",
    )
    reset_layout_telemetry()

    # --- Т16. Мера таблицы пишется всегда, а не только при переполнении -----
    #
    # `render/layout-telemetry.json` прогона 91: 62 записи, из них `role:
    # "table"` — ноль, при 39 листах с таблицами. То есть о таблицах рендерер
    # не сообщал ничего, и ёмкость листа приходилось выводить из худшего
    # случая: 3 510 000 / 1 036 320 = три строки при медианной строке
    # 350 520 EMU. Мера снимает этот вывод — но только если она есть у **всех**
    # листов, включая те, что влезли.
    reset_bullet_measure()
    reset_layout_telemetry()
    fitting_rows = [
        [str(i + 1), f"example-{i + 1}.org/material", TITLE_SERP, "СМИ", "Нейтральный"]
        for i in range(3)
    ]
    ctx_fit = render_search_table_page(fitting_rows, "Показана выдача Яндекса по запросу «Кремлев Умар Назарович».")
    table_fit, _stage_fit = page_shapes(ctx_fit, "страница выдачи")
    measured = get_bullet_measure()
    entry = measured[0] if len(measured) == 1 else {}
    check(
        "Т16а: лист выдачи, который влез, пишет меру своим ключом",
        len(measured) == 1 and entry.get("slideKey") == f"p09_ru_serp_table{TABLE_MEASURE_KEY_SUFFIX}",
        f"записей {len(measured)}: {[m.get('slideKey') for m in measured]}",
    )
    # Суффикс ключа объявлен по обе стороны провода, и его равенство — не
    # формальность: пустой суффикс отдал бы ключ страницы мере таблицы, а по
    # этому ключу перекладка буллетов ищет свою запись.
    ts_suffix = re.search(
        r'export const TABLE_MEASURE_KEY_SUFFIX = "([^"]*)";',
        (
            Path(__file__).resolve().parent.parent
            / "src/modules/digital-profile/orion-golden/deck-sections/measured-table-fit.ts"
        ).read_text(encoding="utf-8"),
    )
    check(
        "Т16а2: суффикс ключа мерной записи совпадает с объявленным в TypeScript и непуст",
        bool(ts_suffix)
        and ts_suffix.group(1) == TABLE_MEASURE_KEY_SUFFIX
        and TABLE_MEASURE_KEY_SUFFIX != ""
        and entry.get("slideKey") != "p09_ru_serp_table",
        f"рендерер {TABLE_MEASURE_KEY_SUFFIX!r}, TypeScript "
        f"{ts_suffix.group(1) if ts_suffix else 'не найден'!r}, ключ записи {entry.get('slideKey')!r}",
    )
    heights = list(entry.get("itemHeights") or [])
    drawn = [int(r.height) for r in table_fit.table.rows]
    check(
        "Т16б: мера перечисляет шапку и строки теми же высотами, какие нарисованы",
        heights == drawn and len(heights) == len(fitting_rows) + 1,
        f"мера {heights} против нарисованного {drawn}",
    )
    check(
        "Т16в: сумма меры равна высоте нарисованной таблицы",
        sum(heights) == int(table_fit.height),
        f"сумма {sum(heights)}, таблица {int(table_fit.height)}",
    )
    row_budget = serp_capacity_terms()["SERP_TABLE_ROW_BUDGET_EMU"]
    check(
        "Т16г: бюджет меры за вычетом шапки равен бюджету строк реестра",
        heights and int(entry.get("availableHeight") or 0) - heights[0] == row_budget,
        f"available {entry.get('availableHeight')} − шапка {heights[0] if heights else '?'} "
        f"против SERP_TABLE_ROW_BUDGET_EMU={row_budget}",
    )
    reset_bullet_measure()
    reset_layout_telemetry()
    long_intro = (
        "Показана выдача Яндекса по запросу «Кремлев Умар Назарович». "
        + "Позиции 17–20 в собранных данных отсутствуют. " * 6
    )
    ctx_long = render_search_table_page(fitting_rows, long_intro)
    long_entry = (get_bullet_measure() or [{}])[0]
    check(
        "Т16д: длина вводного абзаца бюджет меры не двигает — он объявленный, а не от факта",
        long_entry.get("availableHeight") == entry.get("availableHeight"),
        f"короткий абзац {entry.get('availableHeight')}, длинный {long_entry.get('availableHeight')}",
    )
    table_long, _stage_long = page_shapes(ctx_long, "страница выдачи с длинным абзацем")
    check(
        "Т16е: таблица с длинным абзацем начинается ниже, чем с коротким",
        int(table_long.top) > int(table_fit.top),
        f"верх таблицы {int(table_fit.top)} против {int(table_long.top)}",
    )
    # Переполнение таблицы **не** объявляется потерей вердикта: строки таблицы
    # циклу буллетов не подвластны, и `droppedLines > 0` объявил бы ему
    # несходимость — то есть уронил бы оплаченный прогон из-за одной высокой
    # строки. О переполнении говорит запись разметки (Т15б), а не мера.
    reset_bullet_measure()
    reset_layout_telemetry()
    over_rows = [
        [str(i + 1), ADDRESS_BAND_PLAIN, TITLE_95, "Официальный сайт / госресурс", "Нежелательный"]
        for i in range(12)
    ]
    render_search_table_page(over_rows, "Показана выдача Яндекса.")
    over = (get_bullet_measure() or [{}])[0]
    clipped = [e for e in get_layout_telemetry() if e.get("role") == "table"]
    check(
        "Т16ж: переполнение остаётся громким записью разметки, а не потерей в мере",
        int(over.get("droppedBullets") or 0) == 0
        and int(over.get("droppedLines") or 0) == 0
        and len(clipped) == 1
        and clipped[0].get("clipped") is True,
        f"мера {over.get('droppedBullets')}/{over.get('droppedLines')}, записей разметки {len(clipped)}",
    )
    reset_bullet_measure()
    reset_layout_telemetry()

    print(f"\n{'FAILED (' + str(len(failures)) + ')' if failures else 'PASSED (0 failures)'}")
    print_tap_counters(passed=passed_checks, failed=len(failures))
    return 1 if failures else 0


if __name__ == "__main__":
    raise SystemExit(main())
