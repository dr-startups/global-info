"""Visual design tokens + reusable slide components for template v3 (Stage K3).

A single source of truth for colours, typography, spacing and the building
blocks (page frame, cards, badges, polished tables, notes) used by the polished
report template. Keeping this isolated means v1/v2/simple renderers are
untouched and the look can be tuned in one place.

No LLM, no network — pure layout helpers over python-pptx.
"""

from __future__ import annotations

import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# Palette (neutral "Digital Profile Audit" brand — no third-party branding)
# ---------------------------------------------------------------------------
BRAND_PRIMARY = RGBColor(0x0E, 0x1F, 0x3A)   # deep navy
ACCENT = RGBColor(0x1C, 0x6F, 0xD6)          # blue accent
ACCENT_SOFT = RGBColor(0x9E, 0xC2, 0xF0)     # light blue
NEUTRAL_DARK = RGBColor(0x1A, 0x1F, 0x29)    # near-black ink
NEUTRAL_GRAY = RGBColor(0x68, 0x71, 0x80)    # muted text
NEUTRAL_LINE = RGBColor(0xD8, 0xDE, 0xE6)    # hairline / card border
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)        # page background
BG_PANEL = RGBColor(0xF3, 0xF6, 0xFB)        # subtle panel fill
TABLE_HEAD = RGBColor(0x16, 0x2C, 0x4A)
TABLE_ZEBRA = RGBColor(0xEF, 0xF3, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WATERMARK_COLOR = RGBColor(0xDD, 0xE2, 0xEA)

SUCCESS = RGBColor(0x2E, 0x7D, 0x32)
WARNING = RGBColor(0xB8, 0x86, 0x00)
DANGER = RGBColor(0xB0, 0x1E, 0x1E)

RISK_COLORS = {
    "LOW": SUCCESS,
    "MEDIUM": WARNING,
    "HIGH": RGBColor(0xC6, 0x4A, 0x00),
    "CRITICAL": DANGER,
    "UNKNOWN": NEUTRAL_GRAY,
    "NONE": NEUTRAL_GRAY,
}

# Review-status badge tones (match enum keys after uppercasing).
REVIEW_STATUS_COLORS = {
    "PENDING": WARNING,
    "NEEDS_REVIEW": WARNING,
    "MATCH_CONFIRMED": DANGER,
    "FALSE_POSITIVE": SUCCESS,
    "DISMISSED": NEUTRAL_GRAY,
}

# Source-type badge tones for compliance tables.
SOURCE_TYPE_COLORS = {
    "REAL API": ACCENT,
    "MANUAL IMPORT": RGBColor(0x15, 0x65, 0xC0),
    "MOCK/DEMO": NEUTRAL_GRAY,
    "NOT CONFIGURED": NEUTRAL_GRAY,
    "STUB": NEUTRAL_GRAY,
}

# ---------------------------------------------------------------------------
# Typography (pt) and spacing (EMU; 914400 = 1 inch)
# ---------------------------------------------------------------------------
FS_COVER_TITLE = 46
FS_MAIN_TITLE = 26
FS_SECTION_TITLE = 18
FS_SUBTITLE = 13
FS_TABLE_HEAD = 11
FS_TABLE_BODY = 10
FS_METRIC_LABEL = 11
FS_METRIC_VALUE = 24
FS_BODY = 14
FS_NOTE = 10
FS_FOOTER = 9

SLIDE_W = Emu(9144000)
SLIDE_H = Emu(6858000)
MARGIN = Emu(548640)            # ~0.6"
CONTENT_W = Emu(int(SLIDE_W) - 2 * int(MARGIN))
CONTENT_TOP = Emu(1310640)      # where content starts under the header
FOOTER_Y = Emu(6492240)
GUTTER = Emu(137160)            # ~0.15" gap between cards
# Content must stay above footer band (0.2" above footer rule).
CONTENT_SAFE_BOTTOM = Emu(int(FOOTER_Y) - 182880)

ROUNDED_RECT = 5
RECT = 1

# ---------------------------------------------------------------------------
# R2 Design System Foundation (tokens only; backward-compatible, unused by default)
# ---------------------------------------------------------------------------

# R2 typography scale (pt)
R2_TYPO_SECTION_MARKER = 9
R2_TYPO_PAGE_TITLE = 28
R2_TYPO_SUBTITLE = 13
R2_TYPO_BODY = 11
R2_TYPO_METRIC_LABEL = 10
R2_TYPO_METRIC_VALUE = 28
R2_TYPO_TABLE_HEADER = 11
R2_TYPO_TABLE_BODY = 10
R2_TYPO_NOTE = 10
R2_TYPO_BUTTON = 10

# R2 spacing scale (EMU)
R2_SPACE_XS = Emu(45720)     # 0.05"
R2_SPACE_SM = Emu(91440)     # 0.10"
R2_SPACE_MD = Emu(137160)    # 0.15"
R2_SPACE_LG = Emu(182880)    # 0.20"
R2_SPACE_XL = Emu(274320)    # 0.30"
R2_CARD_PADDING = Emu(114300)
R2_SECTION_GAP = Emu(160000)
R2_TABLE_CELL_PAD = Emu(80000)
R2_FOOTER_GAP = Emu(182880)
R2_MEDIA_CARD_GAP = Emu(137160)

# R2 safe-area zones
R2_SLIDE_LEFT = MARGIN
R2_SLIDE_RIGHT = Emu(int(SLIDE_W) - int(MARGIN))
R2_CONTENT_LEFT = MARGIN
R2_CONTENT_RIGHT = Emu(int(MARGIN) + int(CONTENT_W))
R2_HEADER_TOP = Emu(131064)
R2_HEADER_BOTTOM = Emu(1219200)
R2_BODY_TOP = CONTENT_TOP
R2_BODY_BOTTOM = CONTENT_SAFE_BOTTOM
R2_FOOTER_TOP = FOOTER_Y
R2_FOOTER_BOTTOM = Emu(int(SLIDE_H) - 36000)
R2_FOOTER_SAFE_BOTTOM = CONTENT_SAFE_BOTTOM
R2_MIN_GAP_BEFORE_FOOTER = Emu(182880)

# R2 visual style tokens
R2_CARD_RADIUS = ROUNDED_RECT
R2_CARD_BORDER_WIDTH_PT = 0.75
R2_CARD_SHADOW_DX = Emu(8000)
R2_CARD_SHADOW_DY = Emu(12000)
R2_BG_LIGHT = BG_LIGHT
R2_TEXT_MUTED = NEUTRAL_GRAY
R2_ACCENT = ACCENT
R2_WARNING = WARNING
R2_SUCCESS = SUCCESS
R2_DANGER = DANGER

# Localizable table footnote ("Showing top N of M."). Set per-render via
# ``set_table_strings`` so v3 tables honour the report language.
_SHOWING_TOP = "Showing top {n} of {total}."


def set_table_strings(showing_top: str | None) -> None:
    global _SHOWING_TOP
    if showing_top:
        _SHOWING_TOP = showing_top


# ---------------------------------------------------------------------------
# Base helpers
# ---------------------------------------------------------------------------

def blank_slide(prs):
    layouts = prs.slide_layouts
    layout = layouts[6] if len(layouts) > 6 else layouts[-1]
    return prs.slides.add_slide(layout)


def set_bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.text_frame.word_wrap = True
    return box


def _run(p, text: str, size: int, color: RGBColor, bold: bool = False, italic: bool = False):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def truncate(text: Any, length: int) -> str:
    s = "" if text is None else str(text)
    return s if len(s) <= length else s[: length - 1] + "\u2026"


# ---------------------------------------------------------------------------
# Page frame: accent bar + title + subtitle + footer + page number
# ---------------------------------------------------------------------------

def page_frame(
    slide,
    title: str,
    subtitle: str | None = None,
    brand: str = "Digital Profile Audit",
    page_no: int | None = None,
    total: int | None = None,
    watermark: str | None = None,
    title_width: Emu | None = None,
) -> Emu:
    set_bg(slide, BG_LIGHT)
    # top accent bar
    bar = slide.shapes.add_shape(RECT, Emu(0), Emu(0), SLIDE_W, Emu(73152))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    if watermark:
        _watermark(slide, watermark)

    tw = title_width or CONTENT_W
    box = textbox(slide, MARGIN, Emu(228600), tw, Emu(960120))
    tf = box.text_frame
    _run(tf.paragraphs[0], title or "", FS_MAIN_TITLE, BRAND_PRIMARY, bold=True)
    if subtitle:
        sp = tf.add_paragraph()
        _run(sp, subtitle, FS_SUBTITLE, NEUTRAL_GRAY)

    footer(slide, brand, page_no, total)
    return CONTENT_TOP


def footer(slide, brand: str, page_no: int | None, total: int | None) -> None:
    line = slide.shapes.add_shape(RECT, MARGIN, FOOTER_Y, CONTENT_W, Emu(12700))
    line.fill.solid()
    line.fill.fore_color.rgb = NEUTRAL_LINE
    line.line.fill.background()

    box = textbox(slide, MARGIN, Emu(int(FOOTER_Y) + 36000), CONTENT_W, Emu(260000))
    tf = box.text_frame
    p = tf.paragraphs[0]
    _run(p, brand, FS_FOOTER, NEUTRAL_GRAY)
    if page_no is not None:
        rbox = textbox(slide, MARGIN, Emu(int(FOOTER_Y) + 36000), CONTENT_W, Emu(260000))
        rp = rbox.text_frame.paragraphs[0]
        rp.alignment = PP_ALIGN.RIGHT
        label = f"{page_no}" + (f" / {total}" if total else "")
        _run(rp, label, FS_FOOTER, NEUTRAL_GRAY)


def _watermark(slide, text: str) -> None:
    box = textbox(slide, Emu(1371600), Emu(2900000), Emu(6400800), Emu(1200000))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, 66, WATERMARK_COLOR, bold=True)


# ---------------------------------------------------------------------------
# R2 primitives (foundation only; intentionally not wired into slide builders yet)
# ---------------------------------------------------------------------------

def r2_truncate_lines(text: Any, max_lines: int = 2, line_len: int = 64) -> str:
    """Conservative line-cap truncation for stable card/table text zones."""
    s = " ".join(str(text or "").split())
    if not s:
        return ""
    words = s.split(" ")
    lines: list[str] = []
    current = ""
    for w in words:
        candidate = (current + " " + w).strip()
        if len(candidate) <= line_len:
            current = candidate
            continue
        if current:
            lines.append(current)
        current = w
        if len(lines) >= max_lines:
            break
    if current and len(lines) < max_lines:
        lines.append(current)
    joined = "\n".join(lines[:max_lines])
    if len(" ".join(words)) > len(" ".join(" ".join(lines).split())):
        return truncate(joined.replace("\n", " "), max_lines * line_len)
    return joined


def r2_text_box(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    text: str,
    *,
    size: int = R2_TYPO_BODY,
    color: RGBColor = NEUTRAL_DARK,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN | None = None,
):
    """Simple typed text-box primitive with R2 defaults."""
    box = textbox(slide, x, y, w, h)
    tf = box.text_frame
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    _run(p, text, size, color, bold=bold, italic=italic)
    return box


def r2_card(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    *,
    fill: RGBColor = R2_BG_LIGHT,
    border: RGBColor = NEUTRAL_LINE,
    radius: int = R2_CARD_RADIUS,
):
    """Base rounded card primitive for future slide migration."""
    shp = slide.shapes.add_shape(radius, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(R2_CARD_BORDER_WIDTH_PT)
    return shp


def r2_overflow_note(slide, top: Emu, text: str) -> Emu:
    """R2 overflow note primitive; draw only if safe area allows."""
    if not text:
        return top
    h = int(text_block_height([text], R2_TYPO_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0))
    if int(top) + h > int(R2_FOOTER_SAFE_BOTTOM):
        return top
    return note(slide, top, text, "info")


def r2_page_header(
    slide,
    *,
    title: str,
    subtitle: str | None = None,
    section_marker: str | None = None,
    right_meta: str | None = None,
) -> Emu:
    """R2 header primitive; returns body top anchor."""
    set_bg(slide, BG_LIGHT)
    bar = slide.shapes.add_shape(RECT, Emu(0), Emu(0), SLIDE_W, Emu(73152))
    bar.fill.solid()
    bar.fill.fore_color.rgb = R2_ACCENT
    bar.line.fill.background()
    if section_marker:
        m = r2_text_box(
            slide,
            R2_CONTENT_LEFT,
            Emu(165000),
            Emu(2200000),
            Emu(140000),
            section_marker,
            size=R2_TYPO_SECTION_MARKER,
            color=R2_ACCENT,
            bold=True,
        )
        del m
    title_w = Emu(int(CONTENT_W) - 1700000 if right_meta else int(CONTENT_W))
    r2_text_box(
        slide,
        R2_CONTENT_LEFT,
        Emu(250000),
        title_w,
        Emu(900000),
        r2_truncate_lines(title, max_lines=2, line_len=56),
        size=R2_TYPO_PAGE_TITLE,
        color=BRAND_PRIMARY,
        bold=True,
    )
    if subtitle:
        r2_text_box(
            slide,
            R2_CONTENT_LEFT,
            Emu(760000),
            title_w,
            Emu(260000),
            truncate(subtitle, 90),
            size=R2_TYPO_SUBTITLE,
            color=R2_TEXT_MUTED,
        )
    if right_meta:
        r2_text_box(
            slide,
            Emu(int(SLIDE_W) - int(MARGIN) - 1700000),
            Emu(250000),
            Emu(1700000),
            Emu(260000),
            truncate(right_meta, 42),
            size=R2_TYPO_NOTE,
            color=R2_TEXT_MUTED,
            align=PP_ALIGN.RIGHT,
        )
    return R2_BODY_TOP


def r2_page_footer(slide, *, brand: str, page_no: int | None, total: int | None) -> None:
    """R2 footer primitive with fixed safe-area behavior."""
    footer(slide, brand, page_no, total)


def r2_metric_cards(slide, top: Emu, cards: list[dict], per_row: int = 4) -> Emu:
    """R2 metric cards primitive (wrapper over stable existing renderer)."""
    return metric_cards(slide, top, cards, per_row=per_row)


def r2_no_data_state(slide, top: Emu, text: str, *, tone: str = "info") -> Emu:
    """R2 no-data primitive with client-safe default copy."""
    _ = tone  # reserved for future palette variants
    return no_data_card(slide, top, text or "No relevant data available for this section.")


R2_REGION_INTERNAL_RE = re.compile(
    r"internal|debug|classifier|providerAdapter|sourceMode|rawMetadata|reviewQueue|UNCLASSIFIED"
    r"|not collected|unavailable|undefined|null|none",
    re.I,
)
R2_REGION_URL_RE = re.compile(r"https?://\S+|www\.\S+", re.I)


def r2_region_safe_text(value: Any, *, labels: dict | None = None, max_len: int = 220) -> str:
    """Sanitize region-page copy to client-safe visible wording."""
    labels = labels or {}
    raw = " ".join(str(value or "").split())
    if not raw:
        return ""
    safe = raw
    if R2_REGION_URL_RE.search(safe):
        safe = R2_REGION_URL_RE.sub(lambda m: r2_domain_text(m.group(0), 36), safe)
    safe = re.sub(r"\{label\}", "", safe, flags=re.I)
    safe = re.sub(r"\bSource:\s*", labels.get("region_source_prefix", "Примечание: "), safe, flags=re.I)
    safe = re.sub(r"\bevidence\b", labels.get("region_review_word", "проверки"), safe, flags=re.I)
    safe = re.sub(r"\bABSENT\b", labels.get("metric_not_found", "НЕ НАЙДЕНО"), safe, flags=re.I)
    safe = re.sub(r"\bn/ ?t\b|н/в", "", safe, flags=re.I)
    safe = re.sub(r"\brelated:\s*\d+\b", "", safe, flags=re.I)
    safe = re.sub(
        r"No international subject-matched results in collected data\.",
        labels.get("region_international_no_subject_results", "Подтверждённых международных материалов по субъекту не выявлено."),
        safe,
        flags=re.I,
    )
    safe = re.sub(
        r"No adverse organic content in selected subject-matched results\.",
        labels.get("region_no_adverse_organic", "Негативных органических материалов по выбранным релевантным результатам не выявлено."),
        safe,
        flags=re.I,
    )
    safe = R2_REGION_INTERNAL_RE.sub("", safe)
    safe = re.sub(r"\s+", " ", safe).strip(" -,:;")
    if not safe:
        safe = labels.get("r24_fallback_no_data", "Insufficient data for a confident conclusion.")
    return truncate(safe, max_len)


def r2_ru_plural_suggestions(n: int) -> str:
    """Russian plural form for 'подсказка'."""
    n10 = n % 10
    n100 = n % 100
    if n10 == 1 and n100 != 11:
        return "подсказка"
    if 2 <= n10 <= 4 and not (12 <= n100 <= 14):
        return "подсказки"
    return "подсказок"


def r2_region_metric_value(value: Any, *, labels: dict | None = None) -> str:
    """Normalize metric values for client-safe region cards."""
    labels = labels or {}
    raw = str(value or "").strip()
    up = raw.upper()
    if up in {"ABSENT", "NONE", "NO_DATA", "UNKNOWN"}:
        return labels.get("metric_no_data", "НЕТ")
    return r2_region_safe_text(raw, labels=labels, max_len=26) or labels.get("metric_no_data", "НЕТ ДАННЫХ")


def r2_region_header(
    slide,
    *,
    title: str,
    risk_level: str,
    section_marker: str | None = None,
    subtitle: str | None = None,
) -> Emu:
    """Consistent region-page header with protected title/badge spacing."""
    top = r2_page_header(
        slide,
        title=title,
        subtitle=subtitle,
        section_marker=section_marker,
        right_meta=" ",
    )
    risk_badge(
        slide,
        Emu(int(SLIDE_W) - int(MARGIN) - 1500000),
        Emu(228600),
        risk_level or "UNKNOWN",
    )
    return top


def r2_region_metric_row(slide, top: Emu, cards: list[dict], *, per_row: int = 3) -> Emu:
    """Region KPI row primitive with stable equal-height cards."""
    return r2_metric_cards(slide, top, cards, per_row=per_row)


def r2_region_summary_card(
    slide,
    top: Emu,
    *,
    title: str,
    lines: list[str],
    labels: dict | None = None,
    card_h: Emu = Emu(1320000),
) -> Emu:
    """Short analytical summary card (max 4 concise lines)."""
    labels = labels or {}
    body = [r2_region_safe_text(x, labels=labels, max_len=140) for x in (lines or []) if str(x or "").strip()]
    return card(slide, MARGIN, top, CONTENT_W, card_h, title, body[:4])


def r2_region_bullet_sections(
    slide,
    top: Emu,
    sections: list[dict[str, Any]],
    *,
    labels: dict | None = None,
    max_items_per_section: int = 5,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Bounded grouped bullets with overflow continuation note."""
    labels = labels or {}
    cleaned: list[dict[str, Any]] = []
    for sec in sections:
        raw_items = sec.get("items") or []
        items = [r2_region_safe_text(x, labels=labels, max_len=120) for x in raw_items if str(x or "").strip()]
        cleaned.append(
            {
                "label": r2_region_safe_text(sec.get("label"), labels=labels, max_len=56),
                "items": items,
            }
        )
    return bounded_bullet_sections(
        slide,
        top,
        cleaned,
        max_items_per_section=max_items_per_section,
        overflow_note=labels.get("bullets_overflow_note", "+ {n} more items preserved in evidence."),
        layout_warnings=layout_warnings,
    )


def r2_region_no_data_state(
    slide,
    top: Emu,
    *,
    headline: str,
    body: str,
    width_ratio: float = 0.70,
) -> Emu:
    """Premium region no-data card with separated title/body zones."""
    return r2_risk_findings_empty_state(
        slide,
        top,
        headline=headline,
        body=body,
        width_ratio=width_ratio,
    )


def r2_region_note(slide, top: Emu, text: str, *, labels: dict | None = None, kind: str = "disclaimer") -> Emu:
    """Client-safe region note style wrapper."""
    safe = r2_region_safe_text(text, labels=labels, max_len=240)
    if not safe:
        return top
    return note(slide, top, safe, kind)


def r2_region_two_column(slide, top: Emu) -> tuple[Emu, Emu, Emu]:
    """Two-column layout helper: left/right anchors plus shared bottom."""
    gutter = Emu(220000)
    left_w = Emu((int(CONTENT_W) - int(gutter)) // 2)
    right_x = Emu(int(MARGIN) + int(left_w) + int(gutter))
    right_w = Emu(int(CONTENT_W) - int(left_w) - int(gutter))
    r2_card(slide, MARGIN, top, left_w, Emu(2160000), fill=WHITE)
    r2_card(slide, right_x, top, right_w, Emu(2160000), fill=WHITE)
    return Emu(int(top) + 180000), Emu(int(top) + 180000), Emu(int(top) + 2240000)


def r2_warning_card(slide, top: Emu, text: str) -> Emu:
    """R2 warning/recommendation primitive."""
    return warning_card(slide, top, text)


def r2_media_empty_state(
    slide,
    top: Emu,
    *,
    message: str,
    detail: str | None = None,
) -> Emu:
    """Centered premium no-media state card for client-facing pages."""
    card_h = Emu(1280000 if detail else 980000)
    if int(top) + int(card_h) > int(R2_FOOTER_SAFE_BOTTOM):
        card_h = Emu(max(760000, int(R2_FOOTER_SAFE_BOTTOM) - int(top) - 60000))
    r2_card(slide, MARGIN, top, CONTENT_W, card_h, fill=BG_LIGHT, border=NEUTRAL_LINE)
    msg = r2_truncate_lines(message, max_lines=3, line_len=72)
    r2_text_box(
        slide,
        Emu(int(MARGIN) + 180000),
        Emu(int(top) + 150000),
        Emu(int(CONTENT_W) - 360000),
        Emu(360000),
        msg,
        size=R2_TYPO_BODY,
        color=BRAND_PRIMARY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if detail:
        det = r2_truncate_lines(detail, max_lines=3, line_len=96)
        r2_text_box(
            slide,
            Emu(int(MARGIN) + 180000),
            Emu(int(top) + 540000),
            Emu(int(CONTENT_W) - 360000),
            Emu(max(200000, int(card_h) - 640000)),
            det,
            size=R2_TYPO_NOTE,
            color=R2_TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )
    return Emu(int(top) + int(card_h) + int(_BLOCK_GAP))


def r2_safe_table(
    slide,
    top: Emu,
    columns: list[str],
    rows: list[list[Any]],
    *,
    col_widths: list[float] | None = None,
    max_rows: int | None = None,
    overflow_note: str | None = None,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """R2 table primitive with optional overflow note below the table."""
    bottom = table(
        slide,
        top,
        columns,
        rows,
        col_widths=col_widths,
        max_rows=max_rows,
        layout_warnings=layout_warnings,
    )
    if max_rows is not None and len(rows) > max_rows:
        line = overflow_note or _SHOWING_TOP.format(n=max_rows, total=len(rows))
        bottom = r2_overflow_note(slide, bottom, line)
    return bottom


def r2_truncate_cell_text(text: Any, max_len: int = 56) -> str:
    """Truncate table-cell text with stable ellipsis behavior."""
    raw = str(text or "").strip()
    if len(raw) <= max_len:
        return raw
    if max_len <= 1:
        return raw[:max_len]
    return raw[: max_len - 1].rstrip() + "…"


def r2_status_pill(value: str, labels: dict | None = None) -> str:
    """Return client-safe localized status label for evidence rows."""
    labels = labels or {}
    v = (value or "").strip().upper()
    if v in ("EXACT_SUBJECT", "MATCH_CONFIRMED", "CONFIRMED", "ПОДТВЕРЖДЕНО"):
        return labels.get("status_confirmed", "Подтверждено")
    if v in ("LIKELY_SUBJECT", "LIKELY", "VERIFIED_LIKELY"):
        return labels.get("status_likely_subject", "Вероятно относится к субъекту")
    if v in ("PENDING", "NEEDS_REVIEW", "REVIEW_REQUIRED", "ON_REVIEW"):
        return labels.get("status_needs_review", "Требует проверки")
    txt = (value or "").strip()
    if not txt:
        return labels.get("status_confirmed", "Подтверждено")
    if "review" in txt.lower() or "провер" in txt.lower():
        return labels.get("status_needs_review", "Требует проверки")
    return r2_truncate_cell_text(txt, 28)


def r2_result_class_label(value: Any, labels: dict | None = None) -> str:
    """Convert raw top-results class/status enum to client-safe label."""
    labels = labels or {}
    raw = str(value or "").strip()
    key = raw.upper().strip()
    mapping = {
        "UNCLASSIFIED": labels.get("class_unclassified_compact", "Neutral"),
        "RISK": labels.get("class_risk", "Risk"),
        "NEGATIVE": labels.get("class_negative", "Negative"),
        "EXACT_SUBJECT": labels.get("class_exact_subject_compact", "Exact"),
        "LIKELY_SUBJECT": labels.get("class_likely_subject_compact", "Likely"),
        "NAMESAKE": labels.get("class_namesake_compact", "Namesake"),
        "INSUFFICIENT": labels.get("class_insufficient_compact", "Insuff."),
    }
    if key in mapping:
        return mapping[key]
    if not raw:
        return labels.get("class_unclassified_compact", "Neutral")
    # Keep unknown values readable but non-technical.
    if "_" in raw:
        raw = raw.replace("_", " ").title()
    return r2_truncate_cell_text(raw, 10)


def r2_domain_text(value: Any, max_len: int = 34) -> str:
    """Normalize URL/domain text for compact appendix tables."""
    from urllib.parse import urlparse

    raw = str(value or "").strip()
    if not raw:
        return ""
    try:
        parsed = urlparse(raw if "://" in raw else f"https://{raw}")
        host = (parsed.netloc or parsed.path or "").strip().lower()
        host = host.replace("www.", "").split("/")[0]
    except Exception:
        host = raw.replace("http://", "").replace("https://", "").split("/")[0].lower()
    return r2_truncate_cell_text(host, max_len)


def r2_top_results_table(
    slide,
    top: Emu,
    *,
    rows: list[dict],
    region: str,
    max_rows: int,
    labels: dict,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Standardized compact top-results table for RU/INTL pages."""
    _ = region  # reserved for future region-specific variants
    internal_label_re = re.compile(
        r"CLIENT_INCLUDE|REVIEW_REQUIRED|EXCLUDE|RELATED_QUERY|SEARCH_SUGGESTION"
        r"|sourceMode|rawMetadata|reviewQueue|providerAdapter|contentClass",
        re.I,
    )

    def _sanitize_text(text: Any) -> str:
        s = str(text or "")
        s = re.sub(r"https?://\S+", "", s, flags=re.I)
        s = internal_label_re.sub("", s)
        s = re.sub(r"\s+", " ", s).strip()
        return s

    def _extract_domain_fallback(text: str) -> str:
        m = re.search(r"[A-Za-zА-Яа-я0-9-]+\.[A-Za-zА-Яа-я0-9-]+", text or "")
        if not m:
            return ""
        return r2_domain_text(m.group(0), 28)

    safe_rows: list[list[str]] = []
    for item in rows:
        rank = r2_truncate_cell_text(_sanitize_text(item.get("rank")), 8)
        title_src = _sanitize_text(item.get("title"))
        domain_txt = r2_domain_text(item.get("domain") or item.get("url") or "", 28)
        if not domain_txt:
            domain_txt = _extract_domain_fallback(title_src)
        if not domain_txt:
            domain_txt = "—"
        title_txt = r2_truncate_cell_text(title_src, 56)
        cls_txt = r2_result_class_label(_sanitize_text(item.get("classification")), labels)
        safe_rows.append([rank, domain_txt, title_txt, cls_txt])

    if not safe_rows:
        return top

    total = len(safe_rows)
    shown = min(total, max_rows)
    note_tpl = labels.get("table_showing_first", labels.get("showing_top", "Showing first {n} of {total}."))

    while shown > 0:
        foot = note_tpl.format(shown=shown, total=total, n=shown)
        if shown <= _max_table_data_rows(top, foot):
            break
        shown -= 1
    shown = max(1, shown)

    bottom = table(
        slide,
        top,
        [labels.get("th_rank_compact", "#"), labels["th_domain"], labels["th_title"], labels["th_class"]],
        safe_rows[:shown],
        col_widths=[0.07, 0.21, 0.54, 0.18],
        max_rows=shown,
        layout_warnings=layout_warnings,
    )
    return r2_table_overflow_note(slide, bottom, shown=shown, total=total, template=note_tpl)


def r2_compliance_status_label(value: Any, labels: dict | None = None) -> str:
    """Normalize compliance status/severity/review to client-safe label."""
    labels = labels or {}
    raw = str(value or "").strip()
    key = raw.lower().replace("_", " ")
    if not key:
        return labels.get("class_unclassified_compact", "Unclassified")
    if any(tok in key for tok in ("confirm", "ok", "pass", "подтвержд")):
        return labels.get("comp_status_confirmed", "Подтверждено")
    if any(tok in key for tok in ("pending", "review", "на проверке", "needs", "требует", "ожида")):
        return labels.get("comp_status_review", "На проверке")
    if any(tok in key for tok in ("exclude", "noise", "исключ", "false", "ложн", "dismiss", "отклон")):
        return labels.get("comp_status_excluded", "Исключено")
    if "high" in key or "высок" in key:
        return labels.get("comp_level_high", "Высокий")
    if "medium" in key or "сред" in key:
        return labels.get("comp_level_medium", "Средний")
    if "low" in key or "низк" in key:
        return labels.get("comp_level_low", "Низкий")
    if any(tok in key for tok in ("unknown", "unclassified", "не классиф")):
        return labels.get("comp_level_unclassified", "Не классиф.")
    return r2_truncate_cell_text(raw, 22)


def r2_compliance_status_pill(value: Any, labels: dict | None = None, *, compact: bool = False) -> str:
    """Readable status token for compliance table cells."""
    txt = r2_compliance_status_label(value, labels)
    level_short_map = {
        labels.get("comp_level_high", "Высокий"): labels.get("comp_level_high_short", "Выс."),
        labels.get("comp_level_medium", "Средний"): labels.get("comp_level_medium_short", "Ср."),
        labels.get("comp_level_low", "Низкий"): labels.get("comp_level_low_short", "Низ."),
        labels.get("comp_level_unclassified", "Не классиф."): labels.get("comp_level_unclassified_short", "Не класс."),
    }
    if compact:
        if txt in level_short_map:
            return r2_truncate_cell_text(level_short_map[txt], 10)
        short_map = {
            labels.get("comp_status_confirmed", "Подтверждено"): labels.get("comp_status_confirmed_short", "Подтв."),
            labels.get("comp_status_review", "На проверке"): labels.get("comp_status_review_short", "Проверка"),
            labels.get("comp_status_excluded", "Исключено"): labels.get("comp_status_excluded_short", "Искл."),
        }
        return r2_truncate_cell_text(short_map.get(txt, txt), 14)
    return r2_truncate_cell_text(txt, 28)


def r2_compliance_type_label(value: Any, labels: dict | None = None) -> str:
    """Normalize compliance type enums to client-safe labels."""
    labels = labels or {}
    text = str(value or "").strip()
    if not text:
        return labels.get("comp_level_unclassified", "Не классиф.")

    mapping = {
        "SANCTIONS": labels.get("comp_type_sanctions", "Санкции"),
        "PEP": labels.get("comp_type_pep", "Политически значимое лицо"),
        "PEP_RCA": labels.get("comp_type_pep_rca", "Связанные лица PEP"),
        "ADVERSE_MEDIA": labels.get("comp_type_adverse_media", "Негативные публикации"),
        "LEGAL": labels.get("comp_type_legal", "Правовые материалы"),
        "WATCHLIST": labels.get("comp_type_watchlist", "Списки наблюдения"),
        "COMPLIANCE": labels.get("comp_type_compliance", "Комплаенс"),
        "BUSINESS": labels.get("comp_type_business", "Деловые связи"),
    }
    split_tokens = re.split(r"[,/|;]+", text)
    out: list[str] = []
    for token in split_tokens:
        tok = token.strip().upper().replace("-", "_").replace(" ", "_")
        if not tok:
            continue
        out.append(mapping.get(tok, ""))
    out = [x for x in out if x]
    if not out:
        return labels.get("comp_level_unclassified", "Не классиф.")

    # Compact special-case for combined SANCTIONS + PEP_RCA.
    up = {t.strip().upper().replace("-", "_").replace(" ", "_") for t in split_tokens if t.strip()}
    if "SANCTIONS" in up and "PEP_RCA" in up:
        return labels.get("comp_type_sanctions_pep_compact", "Санкции / PEP-связи")
    return r2_truncate_cell_text(" / ".join(dict.fromkeys(out)), 28)


def r2_provider_source_text(value: Any, *, labels: dict | None = None) -> str:
    """Normalize provider/source strings to compact readable text."""
    labels = labels or {}
    raw = str(value or "").strip()
    if not raw:
        return "—"
    raw = re.sub(r"https?://\S+", "", raw, flags=re.I)
    raw = re.sub(r"providerAdapter|sourceMode|rawMetadata|reviewQueue|warn_potential_review|contentClass", "", raw, flags=re.I)
    raw = re.sub(r"\s+", " ", raw).strip(" -,:;")
    if not raw:
        return "—"
    if re.fullmatch(r"manual[_\s-]*import", raw, flags=re.I):
        return labels.get("src_manual_import", "Manual import")
    # Keep known provider labels concise.
    raw = raw.replace("Не настроен / не запрашивался", "Не настроен")
    return r2_truncate_cell_text(raw, 28)


def r2_risk_level_label(value: Any, labels: dict | None = None) -> str:
    """Normalize risk level value to client-safe localized label."""
    labels = labels or {}
    raw = str(value or "").strip()
    key = raw.lower().replace("_", " ")
    if not key:
        return labels.get("risk_level_unknown", "Unknown")
    if any(tok in key for tok in ("critical", "крит")):
        return labels.get("risk_level_critical", "Critical")
    if any(tok in key for tok in ("high", "высок")):
        return labels.get("risk_level_high", "High")
    if any(tok in key for tok in ("medium", "сред")):
        return labels.get("risk_level_medium", "Medium")
    if any(tok in key for tok in ("low", "низк")):
        return labels.get("risk_level_low", "Low")
    if any(tok in key for tok in ("unknown", "undefined", "unclassified", "не определ", "не классиф")):
        return labels.get("risk_level_unknown", "Unknown")
    return labels.get("risk_level_unknown", "Unknown")


def r2_source_label(value: Any, *, labels: dict | None = None) -> str:
    """Normalize source text: no raw URLs, no technical/debug tokens."""
    labels = labels or {}
    raw = str(value or "").strip()
    if not raw:
        return labels.get("source_not_specified", "Source not specified")
    clean = re.sub(
        r"providerAdapter|sourceMode|rawMetadata|reviewQueue|classifier|internal|debug|contentClass",
        "",
        raw,
        flags=re.I,
    )
    clean = re.sub(r"\s+", " ", clean).strip(" -,:;")
    if not clean:
        return labels.get("source_not_specified", "Source not specified")
    if "http://" in clean or "https://" in clean or "www." in clean.lower():
        host = r2_domain_text(clean, 28)
        return host or labels.get("source_not_specified", "Source not specified")
    up = clean.upper().replace("-", "_").replace(" ", "_")
    mapping = {
        "MANUAL_IMPORT": labels.get("src_manual_import", "Manual import"),
        "GOOGLE_SEARCH_PROVIDER": labels.get("src_google", "Google"),
        "GOOGLE": labels.get("src_google", "Google"),
        "WIKIPEDIA": "Wikipedia",
        "YOUTUBE": "YouTube",
        "ADVERSE_MEDIA": labels.get("comp_type_adverse_media", "Adverse media"),
        "ADVERSE_PUBLICITY": labels.get("comp_type_adverse_media", "Adverse media"),
        "LEGAL": labels.get("comp_type_legal", "Legal"),
        "WATCHLIST": labels.get("comp_type_watchlist", "Watchlist"),
        "SANCTIONS": labels.get("comp_type_sanctions", "Sanctions"),
        "REPUTATION": labels.get("rf_source_reputation", "Reputational sources"),
        "LITIGATION": labels.get("rf_source_litigation", "Litigation sources"),
        "OWNERSHIP": labels.get("rf_source_corporate", "Corporate records"),
        "CONTROL": labels.get("rf_source_corporate", "Corporate records"),
        "REGULATORY": labels.get("rf_source_regulatory", "Regulatory sources"),
    }
    if up in mapping:
        return mapping[up]
    if "_" in clean and " " not in clean:
        clean = clean.replace("_", " ")
    if re.fullmatch(r"[A-Z0-9_]{3,}", clean):
        clean = clean.replace("_", " ").title()
    return r2_truncate_cell_text(clean, 30)


def r2_risk_action_label(value: Any, labels: dict | None = None) -> str:
    """Normalize review/internal status into client-safe action wording."""
    labels = labels or {}
    raw = str(value or "").strip()
    key = raw.lower().replace("_", " ")
    if any(tok in key for tok in ("confirm", "подтвержд")):
        return labels.get("rf_action_no_action", "No action needed")
    if any(tok in key for tok in ("false", "dismiss", "exclude", "исключ", "ложн", "отклон")):
        return labels.get("rf_action_exclude_noise", "Exclude as noise")
    if any(tok in key for tok in ("pending", "review", "needs", "требует", "ожида")):
        return labels.get("rf_action_review_source", "Review source")
    if not key:
        return labels.get("rf_action_assess_relevance", "Assess relevance")
    return labels.get("rf_action_assess_relevance", "Assess relevance")


def r2_risk_findings_table(
    slide,
    top: Emu,
    *,
    rows: list[dict[str, Any]],
    labels: dict,
    max_rows: int = 6,
    col_widths: list[float] | None = None,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Compact report-grade risk findings table used on slides 17/29."""
    internal_re = re.compile(
        r"reviewStatus|sourceMode|rawMetadata|providerAdapter|classifier|UNCLASSIFIED|internal|debug",
        re.I,
    )

    def _clean_text(value: Any, limit: int) -> str:
        s = str(value or "").strip()
        s = internal_re.sub("", s)
        s = re.sub(r"https?://\S+", "", s, flags=re.I)
        s = re.sub(r"\s+", " ", s).strip(" -,:;")
        return r2_truncate_cell_text(s, limit)

    safe_rows: list[list[str]] = []
    for item in rows:
        finding = _clean_text(item.get("title"), 74)
        source_raw = item.get("source") or item.get("theme")
        source = r2_source_label(source_raw, labels=labels)
        level = r2_risk_level_label(item.get("severity"), labels=labels)
        action = r2_risk_action_label(item.get("reviewStatus"), labels=labels)
        safe_rows.append(
            [
                finding or "—",
                source,
                r2_truncate_cell_text(level, 22),
                r2_truncate_cell_text(action, 24),
            ]
        )
    if not safe_rows:
        return top

    widths = col_widths or [0.46, 0.22, 0.14, 0.18]
    total = len(safe_rows)
    shown = min(max_rows, total)
    note_tpl = labels.get("table_showing_first", labels.get("showing_top", "Showing first {n} of {total}."))
    while shown > 0:
        note = note_tpl.format(shown=shown, total=total, n=shown)
        if shown <= _max_table_data_rows(top, note):
            break
        shown -= 1
    shown = max(1, shown)

    bottom = table(
        slide,
        top,
        [
            labels.get("th_finding", "Finding"),
            labels.get("th_source", "Source"),
            labels.get("th_level", "Level"),
            labels.get("th_action", "Action"),
        ],
        safe_rows[:shown],
        col_widths=widths,
        max_rows=shown,
        semantic_colors=False,
        layout_warnings=layout_warnings,
    )
    return r2_table_overflow_note(slide, bottom, shown=shown, total=total, template=note_tpl)


def r2_compliance_table(
    slide,
    top: Emu,
    *,
    columns: list[str],
    rows: list[list[Any]],
    max_rows: int,
    labels: dict,
    col_widths: list[float] | None = None,
    note: str | None = None,
    compact: bool = True,
    semantic_colors: bool = False,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Compliance-safe table wrapper with footer-safe continuation note."""
    internal_re = re.compile(
        r"CLIENT_INCLUDE|REVIEW_REQUIRED|EXCLUDE|RELATED_QUERY|SEARCH_SUGGESTION"
        r"|sourceMode|rawMetadata|reviewQueue|providerAdapter|warn_potential_review|contentClass",
        re.I,
    )

    def _clean_cell(val: Any, col_idx: int) -> str:
        s = str(val or "").strip()
        s = internal_re.sub("", s)
        s = re.sub(r"\s+", " ", s).strip()
        if "http://" in s or "https://" in s:
            s = r2_domain_text(s, 28)
        if col_idx == 0:
            return r2_provider_source_text(s, labels=labels)
        return r2_truncate_cell_text(s, 44 if compact else 56)

    safe_rows: list[list[str]] = []
    for row in rows:
        safe_rows.append([_clean_cell(v, i) for i, v in enumerate(row)])
    if not safe_rows:
        return top

    shown = min(max_rows, len(safe_rows))
    bottom = table(
        slide,
        top,
        columns,
        safe_rows[:shown],
        max_rows=shown,
        col_widths=col_widths,
        semantic_colors=semantic_colors,
        layout_warnings=layout_warnings,
    )
    note_tpl = labels.get("table_showing_first", labels.get("showing_top", "Showing first {n} of {total}."))
    bottom = r2_table_overflow_note(slide, bottom, shown=shown, total=len(safe_rows), template=note_tpl)
    if note:
        bottom = _safe_content_note(slide, bottom, note, "disclaimer")
    return bottom


def r2_compliance_empty_state(
    slide,
    top: Emu,
    *,
    headline: str,
    body: str,
    labels: dict | None = None,
    checked: str | None = None,
    result: str | None = None,
    width_ratio: float = 0.78,
) -> Emu:
    """Premium no-data card for compliance slides."""
    _ = labels
    ratio = max(0.72, min(0.82, width_ratio))
    card_w = Emu(int(CONTENT_W) * ratio)
    card_h = Emu(1180000 if not (checked and result) else 1420000)
    card_x = Emu(int(MARGIN) + (int(CONTENT_W) - int(card_w)) // 2)
    card_y = top

    card = slide.shapes.add_shape(ROUNDED_RECT, card_x, card_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_PANEL
    card.line.color.rgb = NEUTRAL_LINE
    card.line.width = Pt(0.9)

    marker_d = Emu(160000)
    marker_x = Emu(int(card_x) + 170000)
    marker_y = Emu(int(card_y) + 170000)
    marker = slide.shapes.add_shape(OVAL, marker_x, marker_y, marker_d, marker_d)
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT_SOFT
    marker.line.color.rgb = ACCENT
    marker.line.width = Pt(0.8)

    text_x = Emu(int(card_x) + 420000)
    text_w = Emu(int(card_w) - 520000)
    r2_text_box(
        slide,
        text_x,
        Emu(int(card_y) + 145000),
        text_w,
        Emu(340000),
        r2_truncate_lines(headline, max_lines=1, line_len=72),
        size=FS_SECTION_TITLE - 2,
        color=BRAND_PRIMARY,
        bold=True,
    )
    r2_text_box(
        slide,
        text_x,
        Emu(int(card_y) + 490000),
        text_w,
        Emu(420000),
        r2_truncate_lines(body, max_lines=2, line_len=88),
        size=FS_BODY - 2,
        color=NEUTRAL_GRAY,
    )

    if checked and result:
        meta = f"{truncate(checked, 24)}: {truncate(result, 32)}"
        r2_text_box(
            slide,
            text_x,
            Emu(int(card_y) + 910000),
            text_w,
            Emu(260000),
            meta,
            size=FS_NOTE,
            color=NEUTRAL_GRAY,
            italic=True,
        )
    return Emu(int(card_y) + int(card_h) + 120000)


def r2_risk_findings_empty_state(
    slide,
    top: Emu,
    *,
    headline: str,
    body: str,
    width_ratio: float = 0.70,
) -> Emu:
    """Premium no-data card with strict non-overlapping title/body zones."""
    ratio = max(0.66, min(0.72, width_ratio))
    card_w = Emu(int(CONTENT_W) * ratio)
    card_x = Emu(int(MARGIN) + (int(CONTENT_W) - int(card_w)) // 2)
    card_y = Emu(max(int(top), 1650000))

    CARD_PAD_TOP = 190000
    CARD_PAD_BOTTOM = 240000
    CARD_PAD_LEFT = 260000
    CARD_PAD_RIGHT = 260000
    ICON_SIZE = 128000
    ICON_TO_TEXT_GAP = 240000
    TITLE_H = 430000
    TITLE_BODY_GAP = 90000
    BODY_H = 430000

    card_h = Emu(CARD_PAD_TOP + TITLE_H + TITLE_BODY_GAP + BODY_H + CARD_PAD_BOTTOM)

    card = slide.shapes.add_shape(ROUNDED_RECT, card_x, card_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_PANEL
    card.line.color.rgb = NEUTRAL_LINE
    card.line.width = Pt(0.9)

    icon_d = Emu(ICON_SIZE)
    icon_x = Emu(int(card_x) + 180000)
    icon_y = Emu(int(card_y) + 190000)
    icon = slide.shapes.add_shape(OVAL, icon_x, icon_y, icon_d, icon_d)
    icon.fill.solid()
    icon.fill.fore_color.rgb = ACCENT_SOFT
    icon.line.color.rgb = ACCENT
    icon.line.width = Pt(0.8)

    title_x = Emu(int(card_x) + CARD_PAD_LEFT + ICON_SIZE + ICON_TO_TEXT_GAP)
    title_y = Emu(int(card_y) + CARD_PAD_TOP)
    title_h = Emu(TITLE_H)
    text_w = Emu(int(card_w) - (int(title_x) - int(card_x)) - CARD_PAD_RIGHT)

    r2_text_box(
        slide,
        title_x,
        title_y,
        text_w,
        title_h,
        r2_truncate_lines(headline, max_lines=2, line_len=60),
        size=FS_SECTION_TITLE - 1,
        color=BRAND_PRIMARY,
        bold=True,
    )

    body_y = Emu(int(title_y) + TITLE_H + TITLE_BODY_GAP)
    body_h = Emu(BODY_H)
    body_box = textbox(slide, title_x, body_y, text_w, body_h)
    body_tf = body_box.text_frame
    body_p = body_tf.paragraphs[0]
    body_p.line_spacing = 1.15
    _run(
        body_p,
        r2_truncate_lines(body, max_lines=2, line_len=76),
        FS_BODY - 2,
        NEUTRAL_GRAY,
    )
    return Emu(int(card_y) + int(card_h) + 120000)


def r2_table_overflow_note(
    slide,
    top: Emu,
    *,
    shown: int,
    total: int,
    template: str | None = None,
) -> Emu:
    """Place continuation note below table without footer collision."""
    if total <= shown:
        return top
    tpl = template or _SHOWING_TOP
    msg = tpl.format(shown=shown, total=total, n=shown)
    return r2_overflow_note(slide, top, msg)


def r2_evidence_appendix_table(
    slide,
    top: Emu,
    *,
    columns: list[str],
    rows: list[list[Any]],
    col_widths: list[float],
    max_rows: int,
    audience: str,
    section_kind: str,
    labels: dict,
    layout_warnings: list[str] | None = None,
) -> tuple[Emu, int]:
    """Appendix-specific safe table with per-section overflow note."""
    def _sanitize_cell(text: Any) -> str:
        raw = str(text or "")
        replacements = {
            "CLIENT_INCLUDE": labels.get("status_confirmed", "Подтверждено"),
            "REVIEW_REQUIRED": labels.get("status_needs_review", "Требует проверки"),
            "EXCLUDE": labels.get("status_excluded_noise", "Исключено / шум"),
            "RELATED_QUERY": "",
            "SEARCH_SUGGESTION": "",
            "sourceMode": "",
            "rawMetadata": "",
            "reviewQueue": "",
            "providerAdapter": "",
            "contentClass": "",
        }
        out = raw
        for key, val in replacements.items():
            out = re.sub(key, val, out, flags=re.I)
        out = re.sub(r"\s+", " ", out).strip()
        return out

    safe_rows: list[list[str]] = []
    for row in rows:
        cells = [r2_truncate_cell_text(_sanitize_cell(c), 64) for c in row]
        if len(cells) >= 2:
            cells[1] = r2_domain_text(cells[1], 34) or cells[1]
        if audience == "client":
            cells = [c.replace("http://", "").replace("https://", "") for c in cells]
        safe_rows.append(cells)

    shown = min(len(safe_rows), max_rows)
    table_bottom = table(
        slide,
        top,
        columns,
        safe_rows,
        col_widths=col_widths,
        max_rows=max_rows,
        layout_warnings=layout_warnings,
    )
    return table_bottom, shown


def r2_grouped_evidence_sections(
    slide,
    top: Emu,
    *,
    groups: list[dict],
    audience: str,
    labels: dict,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Render grouped evidence sections with audience-aware gating."""
    for group in groups:
        rows = list(group.get("rows") or [])
        if not rows:
            continue
        kind = str(group.get("kind") or "")
        if audience == "client" and kind == "excluded":
            continue
        header = r2_truncate_lines(str(group.get("title") or ""), max_lines=1, line_len=72)
        # Skip section when it cannot fit safely above footer.
        if int(top) + 1000000 > int(CONTENT_SAFE_BOTTOM):
            msg_tpl = labels.get(
                "appendix_section_hidden_overflow",
                "Section \"{title}\" is preserved in internal evidence.",
            )
            top = _safe_content_note(slide, top, msg_tpl.format(title=header), "disclaimer")
            continue
        top = note(slide, top, header, "info" if kind != "excluded" else "warning")
        top, _ = r2_evidence_appendix_table(
            slide,
            top,
            columns=list(group.get("columns") or []),
            rows=rows,
            col_widths=list(group.get("col_widths") or [0.4, 0.2, 0.2, 0.2]),
            max_rows=int(group.get("max_rows") or 6),
            audience=audience,
            section_kind=kind,
            labels=labels,
            layout_warnings=layout_warnings,
        )
    return top


# ---------------------------------------------------------------------------
# Text-height estimation (python-pptx cannot measure rendered text, so we
# approximate wrapped line counts to stack elements without overlap).
# ---------------------------------------------------------------------------

EMU_PER_PT = 12700
# Default text-frame left+right insets (~0.1" each).
_TEXT_INSET_PT = 14.0
# Conservative average glyph advance as a fraction of the font size. Slightly
# high on purpose so we over- rather than under-estimate height (no overlap).
_CHAR_W_FACTOR = 0.55
# Vertical gap kept below stacked blocks.
_BLOCK_GAP = Emu(140000)


def _emu_to_pt(emu: Emu) -> float:
    return int(emu) / EMU_PER_PT


def _wrapped_lines(text: str, size: int, width: Emu) -> int:
    """Estimate how many visual lines `text` occupies at `size` within `width`."""
    usable_pt = max(1.0, _emu_to_pt(width) - _TEXT_INSET_PT)
    char_w = max(1.0, size * _CHAR_W_FACTOR)
    chars_per_line = max(1, int(usable_pt / char_w))
    length = max(1, len(text or ""))
    return max(1, -(-length // chars_per_line))  # ceil division


def text_block_height(lines: list[str], size: int, width: Emu,
                      space_after_pt: float = 5.0, pad_pt: float = 12.0) -> Emu:
    """Estimated rendered height of a stacked set of paragraphs."""
    line_h_pt = size * 1.2
    total_pt = pad_pt
    for line in lines:
        vis = _wrapped_lines(line, size, width)
        total_pt += vis * line_h_pt + space_after_pt
    return Emu(int(total_pt * EMU_PER_PT))


def assert_shape_within_safe_area(
    shape,
    context: str,
    warnings: list[str] | None = None,
    *,
    strict: bool = False,
) -> bool:
    """Layout guard: shape bottom must not intrude on footer/page-number band."""
    try:
        bottom = int(shape.top) + int(shape.height)
    except Exception:
        return True
    limit = int(CONTENT_SAFE_BOTTOM)
    ok = bottom <= limit
    if not ok:
        msg = f"Layout safe-area overflow ({context}): bottom={bottom} > {limit}"
        if strict:
            raise ValueError(msg)
        if warnings is not None:
            warnings.append(msg)
    return ok


def _shape_bottom(shape) -> int:
    return int(shape.top) + int(shape.height)


def assert_no_vertical_overlap(zones: list[tuple[int, int]], min_gap: int | None = None) -> bool:
    """zones: list of (top, bottom) in EMU; must be sorted top-to-bottom."""
    gap = int(min_gap if min_gap is not None else CARD_ZONE_GAP)
    for i in range(len(zones) - 1):
        if zones[i][1] + gap > zones[i + 1][0]:
            return False
    return True


# ---------------------------------------------------------------------------
# Bullets
# ---------------------------------------------------------------------------

def bullets(
    slide,
    top: Emu,
    lines: list[str],
    size: int = FS_BODY,
    width: Emu | None = None,
    *,
    max_items: int = 8,
    bottom_limit: Emu | None = None,
    overflow_note: str | None = None,
    emit_overflow_note: bool = True,
    shown_out: list[int] | None = None,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Bounded bullet list — never draws below footer safe line."""
    lines = [l for l in lines if l]
    if not lines:
        if shown_out is not None:
            shown_out.clear()
            shown_out.append(0)
        return top
    w = width or CONTENT_W
    limit = int(bottom_limit or CONTENT_SAFE_BOTTOM)
    cap = max(1, max_items)
    overflow_tpl = overflow_note or "+ {n} more items preserved in evidence."
    overflow_reserve = int(
        text_block_height([overflow_tpl.format(n=99)], FS_NOTE, w, space_after_pt=0.0, pad_pt=8.0)
    ) + int(_BLOCK_GAP)

    shown: list[str] = []
    for line in lines[:cap]:
        trial = shown + [line]
        rendered = [f"\u2022 {x}" for x in trial]
        h = text_block_height(rendered, size, w)
        need_overflow = len(lines) > len(trial)
        reserve = overflow_reserve if (need_overflow and emit_overflow_note) else 0
        if int(top) + int(h) + reserve > limit:
            break
        shown.append(line)
    if not shown:
        shown = [lines[0]]

    remaining = len(lines) - len(shown)
    if remaining > 0 and emit_overflow_note:
        while remaining > 0:
            rendered = [f"\u2022 {line}" for line in shown]
            h = text_block_height(rendered, size, w)
            if int(top) + int(h) + overflow_reserve <= limit:
                break
            if len(shown) <= 1:
                break
            shown.pop()
            remaining = len(lines) - len(shown)

    if shown_out is not None:
        shown_out.clear()
        shown_out.append(len(shown))

    rendered = [f"\u2022 {line}" for line in shown]
    h = text_block_height(rendered, size, w)
    box = textbox(slide, MARGIN, top, w, h)
    tf = box.text_frame
    for i, line in enumerate(rendered):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, line, size, NEUTRAL_DARK)
        p.space_after = Pt(5)
    bottom = Emu(int(top) + int(h) + int(_BLOCK_GAP))
    assert_shape_within_safe_area(box, "bullets", layout_warnings)
    remaining = len(lines) - len(shown)
    if remaining > 0 and emit_overflow_note:
        bottom = _safe_content_note(slide, bottom, overflow_tpl.format(n=remaining), "info")
    return bottom


def bounded_bullet_sections(
    slide,
    top: Emu,
    sections: list[dict[str, Any]],
    *,
    max_items_per_section: int = 8,
    overflow_note: str | None = None,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Render labelled bullet groups; one overflow note total; stay above footer safe line."""
    if not sections:
        return top
    limit = int(CONTENT_SAFE_BOTTOM)
    overflow_tpl = overflow_note or "+ {n} more items preserved in evidence."
    overflow_reserve = int(
        text_block_height([overflow_tpl.format(n=99)], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0)
    ) + int(_BLOCK_GAP)
    section_label_h = int(text_block_height(["Section"], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0))
    min_bullet_h = int(text_block_height(["\u2022 item"], FS_BODY, CONTENT_W))

    total_items = sum(len(s.get("items") or []) for s in sections)
    shown_total = 0

    for sec in sections:
        label = str(sec.get("label") or "")
        items = [str(x) for x in (sec.get("items") or []) if x]
        if not items and not label:
            continue
        need_label = bool(label)
        label_reserve = section_label_h + int(_BLOCK_GAP) if need_label else 0
        will_hide = shown_total < total_items
        reserve = overflow_reserve if will_hide else 0
        if int(top) + label_reserve + min_bullet_h + reserve > limit:
            break
        if need_label:
            top = note(slide, top, label, "section")
        if items:
            count_box: list[int] = []
            bullet_limit = limit - reserve if will_hide else limit
            top = bullets(
                slide,
                top,
                items,
                max_items=max_items_per_section,
                bottom_limit=Emu(bullet_limit),
                emit_overflow_note=False,
                shown_out=count_box,
                layout_warnings=layout_warnings,
            )
            shown_total += count_box[0] if count_box else 0

    hidden = total_items - shown_total
    if hidden > 0:
        top = _safe_content_note(slide, top, overflow_tpl.format(n=hidden), "info")
    return top


# ---------------------------------------------------------------------------
# Badges
# ---------------------------------------------------------------------------

def risk_badge(slide, x: Emu, y: Emu, level: str, w: Emu = Emu(1500000), h: Emu = Emu(470000)) -> None:
    lvl = str(level or "UNKNOWN").upper()
    shape = slide.shapes.add_shape(ROUNDED_RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RISK_COLORS.get(lvl, NEUTRAL_GRAY)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, lvl, 13, WHITE, bold=True)


def _pill_badge(slide, x: Emu, y: Emu, text: str, tone: RGBColor, w: Emu = Emu(1700000), h: Emu = Emu(420000)) -> None:
    shape = slide.shapes.add_shape(ROUNDED_RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = tone
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, truncate(text, 22), 10, WHITE, bold=True)


def review_status_badge(slide, x: Emu, y: Emu, status: str, label: str | None = None) -> None:
    key = str(status or "PENDING").upper()
    _pill_badge(slide, x, y, label or key, REVIEW_STATUS_COLORS.get(key, NEUTRAL_GRAY))


def source_badge(slide, x: Emu, y: Emu, source_type: str, label: str | None = None) -> None:
    key = str(source_type or "").upper()
    tone = SOURCE_TYPE_COLORS.get(key, NEUTRAL_GRAY)
    for k, v in SOURCE_TYPE_COLORS.items():
        if k in key:
            tone = v
            break
    _pill_badge(slide, x, y, label or source_type, tone, w=Emu(1900000))


def metric_card(slide, top: Emu, label: str, value: Any, tone: RGBColor = ACCENT) -> Emu:
    return metric_cards(slide, top, [{"label": label, "value": value, "tone": tone}], per_row=1)


def warning_card(slide, top: Emu, text: str) -> Emu:
    if not text:
        return top
    h = 820000
    shape = slide.shapes.add_shape(ROUNDED_RECT, MARGIN, top, CONTENT_W, Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_PANEL
    shape.line.color.rgb = WARNING
    shape.line.width = Pt(1.25)
    strip = slide.shapes.add_shape(RECT, MARGIN, top, Emu(64008), Emu(h))
    strip.fill.solid()
    strip.fill.fore_color.rgb = WARNING
    strip.line.fill.background()
    box = textbox(slide, Emu(int(MARGIN) + 160000), Emu(int(top) + 90000), Emu(int(CONTENT_W) - 240000), Emu(h - 160000))
    tf = box.text_frame
    _run(tf.paragraphs[0], text, FS_NOTE + 1, NEUTRAL_DARK)
    return Emu(int(top) + h + 80000)


def source_note(slide, top: Emu, text: str) -> Emu:
    return note(slide, top, text, "source")


def polished_table(slide, top: Emu, columns: list[str], rows: list[list[Any]], **kwargs) -> Emu:
    return table(slide, top, columns, rows, **kwargs)


# ---------------------------------------------------------------------------
# Cards / metrics
# ---------------------------------------------------------------------------

def card(slide, x: Emu, y: Emu, w: Emu, h: Emu, title: str, lines: list[str], tone: RGBColor = ACCENT) -> Emu:
    shape = slide.shapes.add_shape(ROUNDED_RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_PANEL
    shape.line.color.rgb = NEUTRAL_LINE
    shape.line.width = Pt(0.75)
    # accent strip
    strip = slide.shapes.add_shape(RECT, x, y, Emu(64008), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = tone
    strip.line.fill.background()

    box = textbox(slide, Emu(int(x) + 160000), Emu(int(y) + 90000), Emu(int(w) - 240000), Emu(int(h) - 160000))
    tf = box.text_frame
    if title:
        _run(tf.paragraphs[0], title, FS_SECTION_TITLE - 2, BRAND_PRIMARY, bold=True)
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        _run(p, line, FS_BODY - 2, NEUTRAL_DARK)
        p.space_after = Pt(2)
    return Emu(int(y) + int(h) + int(_BLOCK_GAP))


def metric_cards(slide, top: Emu, cards: list[dict], per_row: int = 4) -> Emu:
    """cards: [{label, value, tone?}]. Lays a responsive grid of metric cards."""
    if not cards:
        return top
    per_row = max(1, min(per_row, len(cards)))
    gap = int(GUTTER)
    card_w = (int(CONTENT_W) - gap * (per_row - 1)) // per_row
    card_h = 980000
    y = int(top)
    for idx, c in enumerate(cards):
        col = idx % per_row
        if col == 0 and idx > 0:
            y += card_h + gap
        x = int(MARGIN) + col * (card_w + gap)
        shape = slide.shapes.add_shape(ROUNDED_RECT, Emu(x), Emu(y), Emu(card_w), Emu(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_PANEL
        shape.line.color.rgb = NEUTRAL_LINE
        shape.line.width = Pt(0.75)
        box = textbox(slide, Emu(x + 110000), Emu(y + 90000), Emu(card_w - 180000), Emu(card_h - 160000))
        tf = box.text_frame
        _run(tf.paragraphs[0], str(c.get("label", "")), FS_METRIC_LABEL, NEUTRAL_GRAY)
        vp = tf.add_paragraph()
        val = str(c.get("value", ""))
        val_size = FS_METRIC_VALUE - (8 if len(val) > 12 else 0) - (4 if len(val) > 8 else 0)
        _run(vp, val, max(14, val_size), c.get("tone", ACCENT), bold=True)
    rows = (len(cards) + per_row - 1) // per_row
    return Emu(int(top) + rows * (card_h + gap) + 60000)


def no_data_card(slide, top: Emu, text: str) -> Emu:
    h = 760000
    shape = slide.shapes.add_shape(ROUNDED_RECT, MARGIN, top, CONTENT_W, Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_PANEL
    shape.line.color.rgb = NEUTRAL_LINE
    shape.line.width = Pt(0.75)
    icon = slide.shapes.add_shape(ROUNDED_RECT, Emu(int(MARGIN) + 180000), Emu(int(top) + 280000), Emu(90000), Emu(90000))
    icon.fill.solid()
    icon.fill.fore_color.rgb = NEUTRAL_LINE
    icon.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, text, FS_BODY - 1, NEUTRAL_GRAY, italic=True)
    return Emu(int(top) + h + 80000)


def step_cards(slide, top: Emu, steps: list[str], per_row: int = 1) -> Emu:
    """Numbered process/timeline blocks for commercial pages."""
    if not steps:
        return top
    per_row = max(1, min(per_row, 2))
    gap = int(GUTTER)
    card_w = (int(CONTENT_W) - gap * (per_row - 1)) // per_row
    card_h = 720000
    y = int(top)
    for idx, step in enumerate(steps):
        col = idx % per_row
        if col == 0 and idx > 0:
            y += card_h + gap
        x = int(MARGIN) + col * (card_w + gap)
        shape = slide.shapes.add_shape(ROUNDED_RECT, Emu(x), Emu(y), Emu(card_w), Emu(card_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_PANEL
        shape.line.color.rgb = NEUTRAL_LINE
        shape.line.width = Pt(0.75)
        num = slide.shapes.add_shape(ROUNDED_RECT, Emu(x + 90000), Emu(y + 90000), Emu(340000), Emu(340000))
        num.fill.solid()
        num.fill.fore_color.rgb = ACCENT
        num.line.fill.background()
        ntf = num.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.alignment = PP_ALIGN.CENTER
        _run(np, str(idx + 1), 12, WHITE, bold=True)
        box = textbox(slide, Emu(x + 480000), Emu(y + 80000), Emu(card_w - 560000), Emu(card_h - 140000))
        tf = box.text_frame
        _run(tf.paragraphs[0], truncate(step, 120), FS_BODY - 2, NEUTRAL_DARK)
    rows = (len(steps) + per_row - 1) // per_row
    return Emu(int(top) + rows * (card_h + gap) + 60000)


# kind -> (border tone, label)
NOTE_TONES = {
    "warning": WARNING,
    "disclaimer": NEUTRAL_GRAY,
    "source": ACCENT,
    "info": ACCENT,
}


def note(slide, top: Emu, text: str, kind: str = "info") -> Emu:
    if not text:
        return top
    tone = NOTE_TONES.get(kind, ACCENT)
    prefix = {"warning": "\u26a0 ", "disclaimer": "", "source": "Source: ", "info": ""}.get(kind, "")
    full = f"{prefix}{text}"
    h = text_block_height([full], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=10.0)
    box = textbox(slide, MARGIN, top, CONTENT_W, h)
    tf = box.text_frame
    p = tf.paragraphs[0]
    _run(p, full, FS_NOTE, tone, italic=(kind != "warning"))
    return Emu(int(top) + int(h) + int(_BLOCK_GAP))


def _safe_content_note(slide, top: Emu, text: str, kind: str = "info") -> Emu:
    """Draw note only if it fits above CONTENT_SAFE_BOTTOM."""
    if not text:
        return top
    tone = NOTE_TONES.get(kind, ACCENT)
    prefix = {"warning": "\u26a0 ", "disclaimer": "", "source": "Source: ", "info": ""}.get(kind, "")
    full = f"{prefix}{text}"
    h = int(text_block_height([full], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=10.0))
    if int(top) + h > int(CONTENT_SAFE_BOTTOM):
        return top
    return note(slide, top, text, kind)


# ---------------------------------------------------------------------------
# Polished table (zebra + header contrast + risk-coloured cells)
# ---------------------------------------------------------------------------

_RISK_WORDS = {"LOW", "MEDIUM", "HIGH", "CRITICAL"}
_REVIEW_WORDS = {"PENDING", "NEEDS REVIEW", "MATCH CONFIRMED", "FALSE POSITIVE", "DISMISSED"}
_SOURCE_WORDS = {
    "REAL API", "MANUAL IMPORT", "MOCK/DEMO", "NOT CONFIGURED", "STUB",
    "РЕАЛЬНЫЙ API", "РУЧНОЙ ИМПОРТ", "MOCK", "DEMO",
}


def _cell_color(val: str) -> RGBColor | None:
    up = str(val).upper().strip()
    if up in _RISK_WORDS:
        return RISK_COLORS.get(up, NEUTRAL_DARK)
    for key, color in REVIEW_STATUS_COLORS.items():
        if key.replace("_", " ") in up or up == key:
            return color
    for src, color in SOURCE_TYPE_COLORS.items():
        if src in up:
            return color
    if up in _REVIEW_WORDS:
        return WARNING
    if any(s in up for s in _SOURCE_WORDS):
        return ACCENT
    return None


TABLE_ROW_H = 350000
TABLE_NOTE_GAP = 320000
TABLE_SAFE_MARGIN = 220000
TABLE_PDF_BLEED = 140000


def _table_footnote(total: int, shown: int, note_text: str | None) -> str | None:
    if total > shown:
        base = _SHOWING_TOP.format(n=shown, total=total)
        return f"{base} {note_text}".strip() if note_text else base
    return note_text


def _table_footnote_height(label: str | None) -> int:
    if not label:
        return 0
    return (
        int(text_block_height([label], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=10.0))
        + TABLE_NOTE_GAP
        + TABLE_SAFE_MARGIN
        + TABLE_PDF_BLEED
    )


def _max_table_data_rows(top: Emu, footnote: str | None) -> int:
    """Conservative row cap — reserves footnote + PDF export bleed before drawing rows."""
    reserve = _table_footnote_height(footnote)
    available_bottom = max(0, int(CONTENT_SAFE_BOTTOM) - reserve)
    avail_h = max(0, available_bottom - int(top))
    max_rows_including_header = avail_h // TABLE_ROW_H
    return max(1, max_rows_including_header - 1)


def table(
    slide,
    top: Emu,
    columns: list[str],
    rows: list[list[Any]],
    max_rows: int = 12,
    col_widths: list[float] | None = None,
    note_text: str | None = None,
    *,
    semantic_colors: bool = True,
    layout_warnings: list[str] | None = None,
) -> Emu:
    """Safe-area table: rows paginated; footnote always below table, never overlapping."""
    total = len(rows)
    shown_count = min(total, max_rows)
    while shown_count > 0:
        footnote = _table_footnote(total, shown_count, note_text)
        space_rows = _max_table_data_rows(top, footnote)
        if shown_count <= space_rows:
            break
        shown_count -= 1
    shown_count = max(0, shown_count)
    footnote = _table_footnote(total, shown_count, note_text)
    rows = rows[:shown_count]
    if not rows:
        if footnote:
            return note(slide, top, footnote, "info")
        return top
    n_rows = len(rows) + 1
    n_cols = max(1, len(columns))
    height = Emu(TABLE_ROW_H * n_rows)
    graphic = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, CONTENT_W, height)
    tbl = graphic.table

    if col_widths and len(col_widths) == n_cols:
        total_w = int(CONTENT_W)
        s = sum(col_widths)
        for c, frac in enumerate(col_widths):
            tbl.columns[c].width = Emu(int(total_w * frac / s))

    for c, col in enumerate(columns):
        cell = tbl.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEAD
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(FS_TABLE_HEAD)
        para.font.color.rgb = WHITE
        para.word_wrap = False

    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            val = row[c] if c < len(row) else ""
            cell.text = "" if val is None else str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(FS_TABLE_BODY)
            para.word_wrap = False
            tone = _cell_color(str(val)) if semantic_colors else None
            if semantic_colors and tone:
                para.font.color.rgb = tone
                para.font.bold = True
            else:
                para.font.color.rgb = NEUTRAL_DARK
                if not semantic_colors:
                    para.font.bold = False
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ZEBRA

    assert_shape_within_safe_area(graphic, "table", layout_warnings)
    tbl_bottom = int(top) + int(height)
    note_top = tbl_bottom + TABLE_NOTE_GAP + TABLE_PDF_BLEED
    bottom = Emu(note_top)
    if footnote:
        bottom = note(slide, Emu(note_top), footnote, "info")
    if note_text:
        bottom = note(slide, bottom, note_text, "source")
    return bottom


# O5.4.2 / v16 — media card layout constants.
CARD_PAD = Emu(91440)  # ~0.1"
# Vertical gaps between fixed caption zones inside media cards (~8 px).
CARD_ZONE_GAP = Emu(101600)
IMG_TITLE_ZONE_H = Emu(266700)   # max ~2 title lines
IMG_DOMAIN_ZONE_H = Emu(139700)  # 1 domain line
IMG_BADGE_ZONE_H = Emu(139700)   # 1 identity badge line (extra line height)
IMAGE_AREA_FRAC = 0.50
MIN_IMAGE_VISUAL_EMU = 1333350
MIN_GALLERY_IMG_H = 1143000
MIN_GALLERY_IMG_W = 1524000
MIN_GALLERY_CARD_IMG_FRAC = 0.45
MAX_GALLERY_ITEMS = 4
MAX_VIDEO_ITEMS = 4
MAX_UPSCALE_RATIO = 2.0
MAX_INTL_SLIDE_FRAC = 0.55
IMG_SLOT_FRAC = 0.50
GAL_ORION_IMG_FRAC = 0.65
GAL_ORION_IMG_FRAC_TALL = 0.72  # single-row 2-up — image dominates card
GAL_ORION_MIN_IMG_FRAC = 0.55
GAL_ORION_MIN_FRAME_H = 1333500  # ~140 px at 96 dpi
GAL_ORION_MAX_FRAME_ASPECT = 2.35  # reject banner-like image frames (w/h)
GAL_ASPECT_CONTAIN_LO = 0.55
GAL_ASPECT_CONTAIN_HI = 2.2
CARD_BOTTOM_PAD = 91440
GAL_TITLE_ZONE_H = 228600   # max 2 title lines
GAL_DOMAIN_ZONE_H = 101600  # 1 domain line
GAL_BADGE_ZONE_H = 101600
GAL_TWO_LINE_INNER_H = 2400000


def _gallery_title_zone_h(inner_h: int) -> int:
    return int(GAL_TITLE_ZONE_H) if inner_h >= GAL_TWO_LINE_INNER_H else int(GAL_DOMAIN_ZONE_H)


VID_PLAY_DOMAIN_ZONE_H = 101600  # legacy: play icon + domain row (text-only cards)
VID_TITLE_ZONE_H = 228600        # max 2 title lines
VID_DOMAIN_ZONE_H = 101600         # domain row under preview
VID_BADGE_ZONE_H = 101600
VID_BUTTON_ZONE_H = 190000
VID_THUMB_MIN_H = 280000
VID_TWO_LINE_INNER_H = 2000000
MIN_THUMB_B64_LEN = 2000

# ORION-like video evidence card — preview band + domain + title + button
VID_PREVIEW_MIN_H = 380000       # min designed preview/placeholder band height
VID_PREVIEW_MAX_H = 900000       # cap so preview stays balanced vs. text zones
VID_PLAY_CIRCLE_D = 260000       # play-icon circle diameter inside preview
VID_BUTTON_W = 1650000           # fixed "open source" button width
OVAL = 9


def _video_title_zone_h(inner_h: int) -> int:
    return int(VID_TITLE_ZONE_H) if inner_h >= VID_TWO_LINE_INNER_H else int(VID_DOMAIN_ZONE_H)


def short_display_url(url: str, max_len: int = 38) -> str:
    from urllib.parse import urlparse

    try:
        p = urlparse(str(url or "").strip())
        host = (p.netloc or "").replace("www.", "")
        path = (p.path or "").strip("/")
        display = f"{host}/{path}" if path else host
        return truncate(display or url, max_len)
    except Exception:
        return truncate(url, max_len)


def _fit_picture_cover(slide, stream, box_left: int, box_top: int, box_w: int, box_h: int):
    """Cover-crop image into box; preserve aspect ratio; center; clip to box."""
    if hasattr(stream, "seek"):
        stream.seek(0)
    pic = slide.shapes.add_picture(stream, box_left, box_top)
    scale = max(box_w / max(pic.width, 1), box_h / max(pic.height, 1))
    sw = int(pic.width * scale)
    sh = int(pic.height * scale)
    pic.width = sw
    pic.height = sh
    pic.left = box_left + (box_w - sw) // 2
    pic.top = box_top + (box_h - sh) // 2
    if sw > box_w:
        crop = (sw - box_w) / max(sw, 1)
        pic.crop_left = crop / 2
        pic.crop_right = crop / 2
    if sh > box_h:
        crop = (sh - box_h) / max(sh, 1)
        pic.crop_top = crop / 2
        pic.crop_bottom = crop / 2
    pic.left = box_left
    pic.top = box_top
    pic.width = box_w
    pic.height = box_h
    return pic


def _fit_picture_contain(slide, stream, box_left: int, box_top: int, box_w: int, box_h: int):
    """Place image scaled to fit box; preserve aspect ratio; center in box."""
    if hasattr(stream, "seek"):
        stream.seek(0)
    pic = slide.shapes.add_picture(stream, box_left, box_top)
    scale = min(box_w / max(pic.width, 1), box_h / max(pic.height, 1), 1.0)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = box_left + (box_w - pic.width) // 2
    pic.top = box_top + (box_h - pic.height) // 2
    return pic


def _card_text_zone(
    slide,
    x: int,
    y: int,
    w: int,
    h: int,
    text: str,
    size: int,
    color: RGBColor,
    *,
    bold: bool = False,
    italic: bool = False,
    hyperlink: str | None = None,
    layout_warnings: list[str] | None = None,
    context: str = "card_zone",
    word_wrap: bool = False,
) -> int:
    """Fixed-height text zone; returns bottom y (EMU)."""
    box = textbox(slide, Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.line_spacing = 1.0
    p.space_after = Pt(0)
    if hyperlink and text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = True
        try:
            run.hyperlink.address = hyperlink
        except Exception:
            pass
    else:
        _run(p, text, size, color, bold=bold, italic=italic)
    assert_shape_within_safe_area(box, context, layout_warnings)
    return y + h


def _image_bytes_from_item(item: dict[str, Any]) -> bytes | None:
    import base64

    b64 = item.get("thumbnailBytesBase64") or item.get("thumbnailBase64")
    if not b64 or len(str(b64)) < MIN_THUMB_B64_LEN:
        return None
    try:
        return base64.b64decode(b64)
    except Exception:
        return None


def _image_item_key(item: dict[str, Any]) -> str:
    import hashlib

    key = str(item.get("thumbnailStorageKey") or item.get("imageUrl") or item.get("url") or "")
    if key:
        return key
    raw = _image_bytes_from_item(item)
    if raw:
        return hashlib.sha256(raw).hexdigest()
    return str(item.get("title") or "")


def _dedupe_gallery_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[str] = set()
    out: list[dict[str, Any]] = []
    for it in items:
        k = _image_item_key(it)
        if k and k in seen:
            continue
        if k:
            seen.add(k)
        out.append(it)
    return out


def _layout_image_card_zones(
    ih: int,
    *,
    show_identity: bool,
    intl_compact: bool = False,
    orion_tile: bool = False,
    orion_tall: bool = False,
) -> dict[str, tuple[int, int]] | None:
    """Return fixed y-zones (top, bottom) relative to card top for image + captions."""
    pad = int(CARD_PAD)
    gap = int(CARD_ZONE_GAP)
    inner_h = max(1, ih - 2 * pad)
    title_h = _gallery_title_zone_h(inner_h)
    domain_h = int(GAL_DOMAIN_ZONE_H)
    badge_h = int(GAL_BADGE_ZONE_H) if show_identity and not orion_tile else 0
    bottom_pad = int(CARD_BOTTOM_PAD)
    text_stack = title_h + gap + domain_h + bottom_pad
    if show_identity and not orion_tile:
        text_stack += gap + badge_h
    if orion_tile:
        slot_frac = GAL_ORION_IMG_FRAC_TALL if orion_tall else GAL_ORION_IMG_FRAC
        min_img = max(int(inner_h * GAL_ORION_MIN_IMG_FRAC), 450000)
    elif intl_compact:
        slot_frac = 0.42
        min_img = min(int(MIN_GALLERY_IMG_H), max(450000, int(inner_h * 0.38)))
    else:
        slot_frac = IMG_SLOT_FRAC
        min_img = min(int(MIN_GALLERY_IMG_H), max(450000, int(inner_h * 0.38)))
    max_img_h = inner_h - text_stack - gap
    if max_img_h < 400000:
        return None
    img_h = min(int(inner_h * slot_frac), max_img_h)
    if img_h < min_img:
        if max_img_h >= min_img:
            img_h = min(min_img, max_img_h)
        else:
            img_h = max_img_h
    if img_h < 350000:
        return None
    y = pad
    zones: dict[str, tuple[int, int]] = {"img": (y, y + img_h)}
    y = zones["img"][1] + gap
    zones["title"] = (y, y + title_h)
    y = zones["title"][1] + gap
    zones["domain"] = (y, y + domain_h)
    if show_identity and not orion_tile:
        y = zones["domain"][1] + gap
        zones["badge"] = (y, y + badge_h)
    if not assert_no_vertical_overlap(list(zones.values())):
        return None
    if zones[list(zones.keys())[-1]][1] > ih - bottom_pad:
        return None
    return zones


def _gallery_image_acceptable(raw: bytes, slot_w: int, slot_h: int) -> bool:
    """Reject thumbnails that would upscale into blurry gallery tiles."""
    size = _image_native_size(raw)
    if not size:
        return len(raw) >= MIN_THUMB_B64_LEN * 3 // 4
    nw, nh = size
    if nw <= 0 or nh <= 0:
        return False
    return max(slot_w / nw, slot_h / nh) <= MAX_UPSCALE_RATIO


def _gallery_image_fit_mode(nw: int, nh: int, box_w: int, box_h: int) -> str:
    """Cover when safe; contain for banners, upscale, or portrait-in-wide-box (no eye-strip crop)."""
    if nw <= 0 or nh <= 0:
        return "contain"
    img_aspect = nw / nh
    box_aspect = box_w / max(box_h, 1)
    if max(box_w / nw, box_h / nh) > MAX_UPSCALE_RATIO:
        return "contain"
    if img_aspect > GAL_ASPECT_CONTAIN_HI or img_aspect < GAL_ASPECT_CONTAIN_LO:
        return "contain"
    # Portrait headshot in a wide short frame — letterbox instead of horizontal strip crop.
    if img_aspect < 1.05 and box_aspect > 1.08:
        return "contain"
    if img_aspect < 0.90 and box_aspect > img_aspect * 1.35:
        return "contain"
    return "cover"


def _fit_gallery_image(
    slide,
    stream,
    box_left: int,
    box_top: int,
    box_w: int,
    box_h: int,
):
    """Smart gallery fit — cover when safe, else letterboxed contain."""
    if hasattr(stream, "seek"):
        stream.seek(0)
    pic = slide.shapes.add_picture(stream, box_left, box_top)
    nw, nh = int(pic.width), int(pic.height)
    mode = _gallery_image_fit_mode(nw, nh, box_w, box_h)
    try:
        slide.shapes._spTree.remove(pic._element)
    except Exception:
        pass
    if hasattr(stream, "seek"):
        stream.seek(0)
    if mode == "cover":
        return _fit_picture_cover(slide, stream, box_left, box_top, box_w, box_h)
    return _fit_picture_contain(slide, stream, box_left, box_top, box_w, box_h)


def _fit_picture_for_card(
    slide,
    stream,
    box_left: int,
    box_top: int,
    box_w: int,
    box_h: int,
    *,
    allow_cover: bool = True,
):
    """Cover or contain; never upscale beyond MAX_UPSCALE_RATIO."""
    if hasattr(stream, "seek"):
        stream.seek(0)
    pic = slide.shapes.add_picture(stream, box_left, box_top)
    nw, nh = int(pic.width), int(pic.height)
    if nw <= 0 or nh <= 0:
        return pic
    upscale = max(box_w / nw, box_h / nh)
    if upscale > MAX_UPSCALE_RATIO or not allow_cover:
        scale = min(box_w / nw, box_h / nh, 1.0)
        pic.width = int(nw * scale)
        pic.height = int(nh * scale)
        pic.left = box_left + (box_w - pic.width) // 2
        pic.top = box_top + (box_h - pic.height) // 2
        return pic
    try:
        slide.shapes._spTree.remove(pic._element)
    except Exception:
        pass
    if hasattr(stream, "seek"):
        stream.seek(0)
    return _fit_picture_cover(slide, stream, box_left, box_top, box_w, box_h)


def _image_card_caption_heights(show_identity: bool) -> tuple[int, int, int, int]:
    """Return title_h, domain_h, badge_h, caption_stack_h (incl. gaps)."""
    gap = int(CARD_ZONE_GAP)
    title_h = int(IMG_TITLE_ZONE_H)
    domain_h = int(IMG_DOMAIN_ZONE_H)
    badge_h = int(IMG_BADGE_ZONE_H) if show_identity else 0
    stack = title_h + gap + domain_h
    if show_identity:
        stack += gap + badge_h
    return title_h, domain_h, badge_h, stack


def _gallery_caption_heights(show_identity: bool, inner_h: int | None = None) -> tuple[int, int, int, int]:
    """Compact caption stack for evidence gallery cards."""
    gap = int(CARD_ZONE_GAP)
    ih_ref = inner_h if inner_h is not None else GAL_TWO_LINE_INNER_H
    title_h = _gallery_title_zone_h(ih_ref)
    domain_h = int(GAL_DOMAIN_ZONE_H)
    badge_h = int(GAL_BADGE_ZONE_H) if show_identity else 0
    stack = title_h + gap + domain_h
    if show_identity:
        stack += gap + badge_h
    return title_h, domain_h, badge_h, stack


def _image_card_min_height(show_identity: bool, *, orion_tile: bool = False) -> int:
    _, _, _, cap_stack = _gallery_caption_heights(show_identity and not orion_tile, inner_h=0)
    pad = int(CARD_PAD)
    gap = int(CARD_ZONE_GAP)
    if orion_tile:
        min_img = 550000
    else:
        min_img = MIN_GALLERY_IMG_H
    return 2 * pad + min_img + cap_stack + gap


def _gallery_notes_reserve_height(
    labels: dict[str, str],
    total: int,
    *,
    max_shown: int = MAX_GALLERY_ITEMS,
    orion_gallery: bool = False,
    planned_shown: int | None = None,
) -> int:
    """Conservative vertical reserve for overflow notes below the gallery grid."""
    if total <= 0:
        return 0
    if orion_gallery:
        shown = planned_shown if planned_shown is not None else min(total, max_shown)
        if shown >= total:
            return 0
        tpl = labels.get(
            "media_gallery_footer_note",
            "Showing {shown} of {total} relevant images. Others saved in evidence.",
        )
        line = tpl.format(shown=shown, total=total)
        return int(text_block_height([line], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0)) + int(_BLOCK_GAP) + 100000
    candidates: list[str] = []
    if total > max_shown:
        candidates.append(
            labels.get("media_showing_images", "Showing {shown} of {total} subject-matched images.").format(
                shown=max_shown, total=total,
            )
        )
        extra = labels.get("media_saved_in_evidence", "+ {n} saved in evidence.")
        if extra:
            candidates.append(extra.format(n=max(1, total - max_shown)))
    else:
        candidates.append(
            labels.get("media_gallery_skipped", "{n} images skipped (thumbnail too small or unavailable).").format(n=1)
        )
    h = 0
    for line in candidates[:2]:
        h += int(text_block_height([line], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0))
        h += int(_BLOCK_GAP)
    return h + 100000


def _orion_image_frame_ok(cell_w: int, row_h: int, *, orion_tall: bool = False) -> bool:
    """True when the ORION image zone is tall enough and not a banner strip."""
    zones = _layout_image_card_zones(
        row_h, show_identity=False, orion_tile=True, orion_tall=orion_tall,
    )
    if not zones:
        return False
    pad = int(CARD_PAD)
    inner_w = max(1, cell_w - 2 * pad)
    img_h = zones["img"][1] - zones["img"][0]
    if img_h < GAL_ORION_MIN_FRAME_H:
        return False
    if inner_w / max(img_h, 1) > GAL_ORION_MAX_FRAME_ASPECT:
        return False
    return True


def _orion_gallery_layout_candidates(usable: int) -> list[tuple[int, int, int]]:
    """ORION preference: 2x2 only when frames qualify; else 2-up large cards."""
    if usable >= 4:
        return [(2, 2, 4), (2, 1, 2)]
    if usable == 3:
        return [(2, 2, 3), (2, 1, 2)]
    if usable == 2:
        return [(2, 1, 2)]
    return [(1, 1, min(usable, 1))]


def _gallery_layout_candidates(try_n: int, *, orion_gallery: bool = False) -> list[tuple[int, int]]:
    if orion_gallery:
        return [(c, r) for c, r, _ in _orion_gallery_layout_candidates(try_n)]
    if try_n <= 1:
        return [(1, 1)]
    if try_n == 2:
        return [(2, 1), (1, 1)]
    if try_n == 3:
        return [(3, 1), (2, 2)]
    return [(2, 2)]


def _gallery_grid_bottom(top: int, rows: int, row_h: int, gap: int) -> int:
    if rows <= 0:
        return int(top)
    return int(top) + rows * row_h + max(0, rows - 1) * gap


def _plan_image_gallery(
    count: int,
    top: Emu,
    *,
    show_identity: bool = False,
    max_shown: int = MAX_GALLERY_ITEMS,
    labels: dict[str, str] | None = None,
    orion_gallery: bool = False,
) -> tuple[int, int, int, int, int, bool]:
    """Pick cols/rows/count so the full grid + notes fit within CONTENT_SAFE_BOTTOM."""
    labels = labels or {}
    usable = min(count, max_shown)
    gap = int(GUTTER)
    note_reserve = _gallery_notes_reserve_height(
        labels, count, max_shown=max_shown, orion_gallery=orion_gallery,
        planned_shown=2 if orion_gallery and count >= 4 else None,
    )
    safe_bottom = int(CONTENT_SAFE_BOTTOM)
    top_i = int(top)
    min_card_h = _image_card_min_height(show_identity, orion_tile=orion_gallery)
    inner_min_w = MIN_GALLERY_IMG_W

    if orion_gallery:
        for cols, rows, plan_n in _orion_gallery_layout_candidates(usable):
            plan_n = min(plan_n, usable)
            orion_tall = rows == 1 and cols == 2
            if cols == 2 and rows == 2:
                probe_note = _gallery_notes_reserve_height(
                    labels, count, max_shown=plan_n, orion_gallery=True,
                    planned_shown=plan_n,
                )
            else:
                probe_note = _gallery_notes_reserve_height(
                    labels, count, orion_gallery=True, planned_shown=plan_n,
                )
            avail_h = max(1, safe_bottom - top_i - probe_note)
            max_row = (avail_h - gap * max(0, rows - 1)) // max(rows, 1)
            if max_row < min_card_h:
                continue
            cell_w = (int(CONTENT_W) - gap * (cols - 1)) // cols
            if cell_w - 2 * int(CARD_PAD) < inner_min_w:
                continue
            if not _orion_image_frame_ok(cell_w, max_row, orion_tall=orion_tall):
                continue
            grid_bottom = _gallery_grid_bottom(top_i, rows, max_row, gap)
            if grid_bottom > safe_bottom - probe_note:
                continue
            if _layout_image_card_zones(
                max_row,
                show_identity=show_identity,
                orion_tile=True,
                orion_tall=orion_tall,
            ) is None:
                continue
            return plan_n, cols, rows, cell_w, max_row, orion_tall
        return 0, 0, 0, 0, 0, False

    for try_n in range(usable, 0, -1):
        for cols, rows in _gallery_layout_candidates(try_n):
            avail_h = max(1, safe_bottom - top_i - note_reserve)
            max_row = (avail_h - gap * max(0, rows - 1)) // max(rows, 1)
            if max_row < min_card_h:
                continue
            cell_w = (int(CONTENT_W) - gap * (cols - 1)) // cols
            if cell_w - 2 * int(CARD_PAD) < inner_min_w:
                continue
            grid_bottom = _gallery_grid_bottom(top_i, rows, max_row, gap)
            if grid_bottom > safe_bottom - note_reserve:
                continue
            if _layout_image_card_zones(
                max_row,
                show_identity=show_identity,
                orion_tile=orion_gallery,
            ) is None:
                continue
            return try_n, cols, rows, cell_w, max_row, False
    return 0, 0, 0, 0, 0, False


def _image_gallery_geometry(
    count: int,
    top: Emu,
    *,
    show_identity: bool = False,
    max_shown: int = 6,
) -> tuple[int, int, int, int]:
    """Backward-compatible wrapper — returns geometry for the best-fit plan."""
    n, cols, rows, cell_w, row_h, _ = _plan_image_gallery(
        count, top, show_identity=show_identity, max_shown=max_shown,
    )
    if n <= 0:
        return 2, 1, int(CONTENT_W), _image_card_min_height(show_identity)
    return cols, rows, cell_w, row_h


def _layout_video_card_zones(
    ih: int,
    *,
    show_badge: bool = False,
    has_thumb: bool = False,
) -> dict[str, tuple[int, int]] | None:
    """ORION video card zones — preview band, domain row, title, button (top→bottom)."""
    pad = int(CARD_PAD)
    gap = int(CARD_ZONE_GAP)
    bottom_pad = int(CARD_BOTTOM_PAD)
    btn_h = int(VID_BUTTON_ZONE_H)
    title_h = int(VID_TITLE_ZONE_H)
    domain_h = int(VID_DOMAIN_ZONE_H)

    y = ih - pad - bottom_pad
    button = (y - btn_h, y)
    y = button[0] - gap
    title = (y - title_h, y)
    y = title[0] - gap
    domain = (y - domain_h, y)
    y = domain[0] - gap

    preview_top = pad
    preview_bottom = y
    preview_h = preview_bottom - preview_top
    if preview_h < int(VID_PREVIEW_MIN_H):
        return None
    if preview_h > int(VID_PREVIEW_MAX_H):
        preview_top = preview_bottom - int(VID_PREVIEW_MAX_H)

    zones: dict[str, tuple[int, int]] = {
        "preview": (preview_top, preview_bottom),
        "domain": domain,
        "title": title,
        "button": button,
    }
    if not assert_no_vertical_overlap(sorted(zones.values(), key=lambda z: z[0])):
        return None
    return zones


def _video_card_min_height(*, show_badge: bool = False, has_thumb: bool = False) -> int:
    pad = 2 * int(CARD_PAD)
    gap = int(CARD_ZONE_GAP)
    bottom = int(CARD_BOTTOM_PAD)
    stack = (
        int(VID_PREVIEW_MIN_H) + gap
        + int(VID_DOMAIN_ZONE_H) + gap
        + int(VID_TITLE_ZONE_H) + gap
        + int(VID_BUTTON_ZONE_H) + bottom
    )
    return pad + stack


def _video_notes_reserve_height(labels: dict[str, str], total: int, *, max_shown: int) -> int:
    if total <= max_shown:
        return 0
    line = labels.get(
        "media_videos_note", "Showing {shown} of {total} relevant video sources. + {n} saved in evidence.",
    ).format(shown=max_shown, total=total, n=max(1, total - max_shown))
    h = int(text_block_height([line], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=8.0))
    return h + int(_BLOCK_GAP) + 100000


def _plan_video_grid(
    count: int,
    top: Emu,
    *,
    show_badge: bool = False,
    max_shown: int = MAX_VIDEO_ITEMS,
    labels: dict[str, str] | None = None,
) -> tuple[int, int, int]:
    """Return (plan_n, row_h, cols) so grid + notes fit within CONTENT_SAFE_BOTTOM."""
    labels = labels or {}
    usable = min(count, max_shown)
    cols = 2
    gap = int(GUTTER)
    note_reserve = _video_notes_reserve_height(labels, count, max_shown=max_shown)
    safe_bottom = int(CONTENT_SAFE_BOTTOM)
    top_i = int(top)
    min_row = _video_card_min_height(show_badge=show_badge, has_thumb=False)

    for try_n in range(usable, 0, -1):
        rows = max(1, (try_n + cols - 1) // cols)
        avail_h = max(1, safe_bottom - top_i - note_reserve)
        max_row = (avail_h - gap * max(0, rows - 1)) // max(rows, 1)
        if max_row < min_row:
            continue
        grid_bottom = top_i + rows * max_row + max(0, rows - 1) * gap
        if grid_bottom > safe_bottom - note_reserve:
            continue
        if _layout_video_card_zones(max_row, show_badge=show_badge, has_thumb=False) is None:
            continue
        return try_n, max_row, cols
    return 0, 0, cols


def _image_native_size(raw: bytes) -> tuple[int, int] | None:
    """Read width/height from PNG/JPEG headers without Pillow."""
    if len(raw) >= 24 and raw[:8] == b"\x89PNG\r\n\x1a\n":
        return int.from_bytes(raw[16:20], "big"), int.from_bytes(raw[20:24], "big")
    if len(raw) >= 4 and raw[:2] == b"\xff\xd8":
        i = 2
        while i + 9 < len(raw):
            if raw[i] != 0xFF:
                i += 1
                continue
            marker = raw[i + 1]
            if marker in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF):
                h = int.from_bytes(raw[i + 5 : i + 7], "big")
                w = int.from_bytes(raw[i + 7 : i + 9], "big")
                return w, h
            if marker in (0xD8, 0xD9):
                break
            seg_len = int.from_bytes(raw[i + 2 : i + 4], "big")
            i += 2 + max(seg_len, 2)
    return None


def _intl_image_acceptable(raw: bytes, slot_w: int, slot_h: int) -> bool:
    """Reject thumbnails that would upscale beyond quality gate."""
    size = _image_native_size(raw)
    if not size:
        return len(raw) >= MIN_THUMB_B64_LEN * 3 // 4
    nw, nh = size
    if nw <= 0 or nh <= 0:
        return False
    upscale = max(slot_w / nw, slot_h / nh)
    if upscale > MAX_UPSCALE_RATIO:
        return False
    slide_w = int(SLIDE_W) * MAX_INTL_SLIDE_FRAC
    slide_h = int(SLIDE_H) * MAX_INTL_SLIDE_FRAC
    disp_w = nw * min(slot_w / nw, slot_h / nh, 1.0)
    disp_h = nh * min(slot_w / nw, slot_h / nh, 1.0)
    return disp_w <= slide_w and disp_h <= slide_h


def add_image_card(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    item: dict[str, Any],
    *,
    show_identity: bool = False,
    link_label: str = "Source",
    layout_warnings: list[str] | None = None,
    intl_compact: bool = False,
    allow_cover: bool = True,
    orion_tile: bool = False,
    orion_tall: bool = False,
) -> bool:
    """Image evidence tile — ORION single-frame gallery or legacy nested slot."""
    import io

    ix, iy, iw, ih = int(x), int(y), int(w), int(h)
    pad = int(CARD_PAD)
    gap = int(CARD_ZONE_GAP)
    inner_w = max(1, iw - 2 * pad)
    inner_x = ix + pad
    if orion_tile:
        show_identity = False
    zones = _layout_image_card_zones(
        ih, show_identity=show_identity, intl_compact=intl_compact,
        orion_tile=orion_tile, orion_tall=orion_tall,
    )
    if not zones and show_identity:
        show_identity = False
        zones = _layout_image_card_zones(
            ih, show_identity=False, intl_compact=intl_compact,
            orion_tile=orion_tile, orion_tall=orion_tall,
        )
    if not zones:
        return False
    if iy + ih > int(CONTENT_SAFE_BOTTOM):
        return False

    raw = _image_bytes_from_item(item)
    if not raw:
        return False

    img_top, img_bottom = zones["img"]
    img_box_h = img_bottom - img_top
    if intl_compact and not _intl_image_acceptable(raw, inner_w, img_box_h):
        return False

    frame = slide.shapes.add_shape(ROUNDED_RECT, Emu(ix), Emu(iy), Emu(iw), Emu(ih))
    frame.fill.solid()
    frame.fill.fore_color.rgb = BG_LIGHT
    frame.line.color.rgb = NEUTRAL_LINE
    frame.line.width = Pt(0.75)

    img_slot = None
    if not orion_tile:
        img_slot = slide.shapes.add_shape(
            ROUNDED_RECT, Emu(inner_x), Emu(iy + img_top), Emu(inner_w), Emu(img_box_h),
        )
        img_slot.fill.solid()
        img_slot.fill.fore_color.rgb = BG_PANEL
        img_slot.line.color.rgb = NEUTRAL_LINE

    pic_ok = False
    try:
        if orion_tile:
            pic = _fit_gallery_image(
                slide, io.BytesIO(raw), inner_x, iy + img_top, inner_w, img_box_h,
            )
        else:
            pic = _fit_picture_for_card(
                slide, io.BytesIO(raw), inner_x, iy + img_top, inner_w, img_box_h,
                allow_cover=allow_cover and not intl_compact,
            )
        pic_ok = int(pic.width) >= 1 and int(pic.height) >= 1
        pic_bottom = int(pic.top) + int(pic.height)
        pic_right = int(pic.left) + int(pic.width)
        if pic_ok and (
            int(pic.top) < iy + img_top - 5000
            or pic_bottom > iy + img_bottom + 5000
            or int(pic.left) < inner_x - 5000
            or pic_right > inner_x + inner_w + 5000
        ):
            pic_ok = False
        if pic_ok and orion_tile:
            pic.width = min(int(pic.width), inner_w)
            pic.height = min(int(pic.height), img_box_h)
            pic.left = max(inner_x, min(int(pic.left), inner_x + inner_w - int(pic.width)))
            pic.top = max(iy + img_top, min(int(pic.top), iy + img_bottom - int(pic.height)))
            if int(pic.height) < img_box_h * 0.22 or int(pic.width) < inner_w * 0.18:
                pic_ok = False
        if pic_ok and not orion_tile and not intl_compact and allow_cover:
            pic_ok = int(pic.width) >= MIN_GALLERY_IMG_W // 3 and int(pic.height) >= MIN_GALLERY_IMG_H // 3
        if pic_ok:
            assert_shape_within_safe_area(pic, "image_card_picture", layout_warnings)
        else:
            try:
                slide.shapes._spTree.remove(pic._element)
            except Exception:
                pass
    except Exception:
        pic_ok = False
    if not pic_ok:
        for shape in (frame, img_slot):
            if shape is not None:
                try:
                    slide.shapes._spTree.remove(shape._element)
                except Exception:
                    pass
        return False

    source_url = str(item.get("sourcePageUrl") or item.get("url") or item.get("imageUrl") or "")
    title_top, title_bottom = zones["title"]
    title_zone_h = title_bottom - title_top
    two_line_title = title_zone_h >= int(GAL_TITLE_ZONE_H) - 5000
    domain = truncate(item.get("source") or short_display_url(source_url), 32 if two_line_title else 28)
    title = truncate(item.get("title"), 44 if two_line_title else 26)

    _card_text_zone(
        slide, inner_x, iy + title_top, inner_w, title_zone_h,
        title, FS_NOTE, BRAND_PRIMARY, bold=True,
        word_wrap=two_line_title,
        layout_warnings=layout_warnings, context="image_card_title",
    )
    dom_top, dom_bottom = zones["domain"]
    _card_text_zone(
        slide, inner_x, iy + dom_top, inner_w, dom_bottom - dom_top,
        domain, FS_NOTE, ACCENT if source_url.startswith("http") else NEUTRAL_GRAY,
        hyperlink=source_url if source_url.startswith("http") else None,
        word_wrap=False,
        layout_warnings=layout_warnings, context="image_card_domain",
    )
    if show_identity and item.get("identityDecision") and "badge" in zones:
        b_top, b_bottom = zones["badge"]
        _card_text_zone(
            slide, inner_x, iy + b_top, inner_w, b_bottom - b_top,
            truncate(str(item.get("identityDecision")), 22), FS_NOTE - 1, NEUTRAL_GRAY,
            italic=True, word_wrap=False,
            layout_warnings=layout_warnings, context="image_card_badge",
        )

    zone_list = [zones[k] for k in ("img", "title", "domain") + (("badge",) if "badge" in zones else ())]
    if not assert_no_vertical_overlap([(iy + z[0], iy + z[1]) for z in zone_list]):
        try:
            slide.shapes._spTree.remove(frame._element)
        except Exception:
            pass
        return False

    assert_shape_within_safe_area(frame, "image_card_frame", layout_warnings)
    return True


def _video_preview_placeholder(slide, x: int, y: int, w: int, h: int) -> None:
    """Designed neutral preview — light band with a centered play circle."""
    circle_d = min(int(VID_PLAY_CIRCLE_D), max(1, h - 60000))
    cx = x + (w - circle_d) // 2
    cy = y + (h - circle_d) // 2
    circle = slide.shapes.add_shape(OVAL, Emu(cx), Emu(cy), Emu(circle_d), Emu(circle_d))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.color.rgb = ACCENT_SOFT
    circle.line.width = Pt(1.0)
    ctf = circle.text_frame
    ctf.margin_left = 0
    ctf.margin_right = 0
    ctf.margin_top = 0
    ctf.margin_bottom = 0
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    _run(cp, "\u25b6", 9, ACCENT)


def add_video_card(
    slide,
    x: Emu,
    y: Emu,
    w: Emu,
    h: Emu,
    item: dict[str, Any],
    link_label: str = "Open source",
    layout_warnings: list[str] | None = None,
) -> bool:
    """ORION video evidence card — preview band, domain, 2-line title, link button."""
    import io

    ix, iy, iw, ih = int(x), int(y), int(w), int(h)
    pad = int(CARD_PAD)
    inner_w = max(1, iw - 2 * pad)
    inner_x = ix + pad
    url = str(item.get("url") or item.get("sourcePageUrl") or "")
    domain = truncate(item.get("source") or short_display_url(url), 32)
    raw = _image_bytes_from_item(item)

    zones = _layout_video_card_zones(ih)
    if not zones:
        return False
    if iy + ih > int(CONTENT_SAFE_BOTTOM):
        return False

    title = truncate(item.get("title"), 44)

    frame = slide.shapes.add_shape(ROUNDED_RECT, Emu(ix), Emu(iy), Emu(iw), Emu(ih))
    frame.fill.solid()
    frame.fill.fore_color.rgb = BG_LIGHT
    frame.line.color.rgb = NEUTRAL_LINE
    frame.line.width = Pt(0.75)

    p_top, p_bottom = zones["preview"]
    preview_h = p_bottom - p_top
    preview = slide.shapes.add_shape(
        ROUNDED_RECT, Emu(inner_x), Emu(iy + p_top), Emu(inner_w), Emu(preview_h),
    )
    preview.fill.solid()
    preview.fill.fore_color.rgb = BG_PANEL
    preview.line.color.rgb = NEUTRAL_LINE
    preview.line.width = Pt(0.5)

    thumb_ok = False
    if raw:
        try:
            pic = _fit_picture_for_card(
                slide, io.BytesIO(raw), inner_x, iy + p_top, inner_w, preview_h,
                allow_cover=False,
            )
            thumb_ok = int(pic.width) >= 1 and int(pic.height) >= 1
            if not thumb_ok:
                try:
                    slide.shapes._spTree.remove(pic._element)
                except Exception:
                    pass
        except Exception:
            thumb_ok = False
    if not thumb_ok:
        _video_preview_placeholder(slide, inner_x, iy + p_top, inner_w, preview_h)

    d_top, d_bottom = zones["domain"]
    _card_text_zone(
        slide, inner_x, iy + d_top, inner_w, d_bottom - d_top,
        domain, FS_NOTE - 1, NEUTRAL_GRAY, word_wrap=False,
        layout_warnings=layout_warnings, context="video_card_domain",
    )

    title_top, title_bottom = zones["title"]
    _card_text_zone(
        slide, inner_x, iy + title_top, inner_w, title_bottom - title_top,
        title, FS_NOTE, BRAND_PRIMARY, bold=True, word_wrap=True,
        layout_warnings=layout_warnings, context="video_card_title",
    )

    btn_top, btn_bottom = zones["button"]
    btn_h = btn_bottom - btn_top
    btn_w = min(int(VID_BUTTON_W), inner_w)
    has_link = url.startswith("http")
    button = slide.shapes.add_shape(ROUNDED_RECT, Emu(inner_x), Emu(iy + btn_top), Emu(btn_w), Emu(btn_h))
    button.fill.solid()
    button.fill.fore_color.rgb = RGBColor(0xEA, 0xF1, 0xFB)
    button.line.color.rgb = ACCENT if has_link else NEUTRAL_LINE
    button.line.width = Pt(0.75)
    btf = button.text_frame
    btf.margin_top = 0
    btf.margin_bottom = 0
    btf.word_wrap = False
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    brun = _run(bp, link_label, FS_NOTE - 1, ACCENT if has_link else NEUTRAL_GRAY, bold=True)
    if has_link:
        try:
            brun.hyperlink.address = url
        except Exception:
            pass

    zone_pts = sorted([(iy + z[0], iy + z[1]) for z in zones.values()], key=lambda z: z[0])
    if not assert_no_vertical_overlap(zone_pts):
        try:
            slide.shapes._spTree.remove(frame._element)
        except Exception:
            pass
        return False

    assert_shape_within_safe_area(frame, "video_card", layout_warnings)
    return True


def _grid_geometry(
    count: int,
    top: Emu,
    *,
    max_shown: int = 6,
    min_row_h: int = 560000,
    max_row_h: int = 1150000,
) -> tuple[int, int, int, int]:
    """Return cols, rows, cell_w, row_h fitting within footer safe area."""
    shown = min(count, max_shown)
    cols = 2 if shown <= 4 else 3
    rows = max(1, (shown + cols - 1) // cols)
    gap = int(GUTTER)
    note_reserve = 260000 if count > max_shown else 0  # overflow note below grid only
    avail_h = max(1, int(CONTENT_SAFE_BOTTOM) - int(top) - note_reserve)
    row_h = (avail_h - gap * max(0, rows - 1)) // max(rows, 1)
    row_h = max(min_row_h, min(row_h, max_row_h))
    cell_w = (int(CONTENT_W) - gap * (cols - 1)) // cols
    return cols, rows, cell_w, row_h


def _safe_gallery_note(slide, top: Emu, text: str) -> Emu:
    """Draw note only if it fits above footer safe line."""
    if not text:
        return top
    h = int(text_block_height([text], FS_NOTE, CONTENT_W, space_after_pt=0.0, pad_pt=10.0))
    if int(top) + h > int(CONTENT_SAFE_BOTTOM):
        return top
    return note(slide, top, text, "info")


def image_grid(
    slide,
    top: Emu,
    items: list[dict[str, Any]],
    *,
    cols: int = 3,
    max_items: int = MAX_GALLERY_ITEMS,
    show_identity: bool = False,
    labels: dict[str, str] | None = None,
    layout_warnings: list[str] | None = None,
    intl_compact: bool = False,
    allow_cover: bool = False,
    orion_gallery: bool = False,
) -> Emu:
    """Image evidence gallery — max 4 validated cards, deduped; grid + notes stay in safe area."""
    labels = labels or {}
    deduped = _dedupe_gallery_items(items)
    total = len(items)
    usable = [it for it in deduped if _image_bytes_from_item(it)]
    skipped_bytes = len(deduped) - len(usable)
    dup_skipped = len(items) - len(deduped)
    if not usable:
        if total > 0:
            msg = labels.get("nd_gallery_no_usable_images", "Selected images unavailable for gallery display.")
            return note(slide, top, msg, "info")
        return top

    plan_n, cols, rows, cell_w, row_h, orion_tall = _plan_image_gallery(
        len(usable), top, show_identity=show_identity, max_shown=max_items, labels=labels,
        orion_gallery=orion_gallery,
    )
    if plan_n <= 0:
        msg = labels.get("nd_gallery_no_usable_images", "Selected images unavailable for gallery display.")
        return note(slide, top, msg, "info")

    gap = int(GUTTER)
    y0 = int(top)
    rendered = 0
    skipped_layout = 0
    candidate_idx = 0
    max_slots = plan_n

    while rendered < max_slots and candidate_idx < len(usable):
        item = usable[candidate_idx]
        candidate_idx += 1
        row, col = divmod(rendered, cols)
        left = int(MARGIN) + col * (cell_w + gap)
        cell_top = y0 + row * (row_h + gap)
        if cell_top + row_h > int(CONTENT_SAFE_BOTTOM):
            break
        if add_image_card(
            slide, Emu(left), Emu(cell_top), Emu(cell_w), Emu(row_h),
            item, show_identity=show_identity,
            link_label=labels.get("media_source_link", "Source"),
            layout_warnings=layout_warnings,
            intl_compact=intl_compact,
            allow_cover=allow_cover,
            orion_tile=orion_gallery,
            orion_tall=orion_tall,
        ):
            rendered += 1
        else:
            skipped_layout += 1

    if rendered == 0:
        msg = labels.get("nd_gallery_no_usable_images", "Selected images unavailable for gallery display.")
        return note(slide, top, msg, "info")

    actual_rows = max(1, (rendered + cols - 1) // cols)
    bottom = Emu(_gallery_grid_bottom(y0, actual_rows, row_h, gap) + int(_BLOCK_GAP))
    skipped = skipped_bytes + skipped_layout + dup_skipped + max(0, len(usable) - candidate_idx)
    hidden = max(0, total - rendered)
    if orion_gallery:
        if hidden > 0:
            tpl = labels.get(
                "media_gallery_footer_note",
                "Showing {shown} of {total} relevant images. Others saved in evidence.",
            )
            bottom = _safe_gallery_note(slide, bottom, tpl.format(shown=rendered, total=total))
    else:
        if skipped > 0:
            tpl = labels.get("media_gallery_skipped", "{n} images skipped (thumbnail too small or unavailable).")
            bottom = _safe_gallery_note(slide, bottom, tpl.format(n=skipped))
        if total > rendered:
            tpl = labels.get("media_showing_images", "Showing {shown} of {total} subject-matched images.")
            bottom = _safe_gallery_note(slide, bottom, tpl.format(shown=rendered, total=total))
        if hidden > 0:
            extra = labels.get("media_saved_in_evidence", "+ {n} saved in evidence.")
            if extra:
                bottom = _safe_gallery_note(slide, bottom, extra.format(n=hidden))
    return bottom


def video_cards(
    slide,
    top: Emu,
    items: list[dict[str, Any]],
    link_label: str = "Open source",
    *,
    max_items: int = MAX_VIDEO_ITEMS,
    labels: dict[str, str] | None = None,
    layout_warnings: list[str] | None = None,
    note_template: str | None = None,
) -> Emu:
    """ORION 2-column video evidence grid — max 4 cards with preview bands."""
    labels = labels or {}
    total = len(items)
    if not items:
        return top

    plan_n, row_h, cols = _plan_video_grid(
        len(items), top, max_shown=max_items, labels=labels,
    )
    if plan_n <= 0:
        return top

    gap = int(GUTTER)
    cell_w = (int(CONTENT_W) - gap * (cols - 1)) // cols
    y0 = int(top)
    rendered = 0
    skipped = 0

    for item in items:
        if rendered >= plan_n:
            break
        row, col = divmod(rendered, cols)
        left = int(MARGIN) + col * (cell_w + gap)
        cell_top = y0 + row * (row_h + gap)
        if cell_top + row_h > int(CONTENT_SAFE_BOTTOM):
            break
        if add_video_card(
            slide, Emu(left), Emu(cell_top), Emu(cell_w), Emu(row_h),
            item, link_label=link_label, layout_warnings=layout_warnings,
        ):
            rendered += 1
        else:
            skipped += 1

    if rendered == 0:
        return top
    actual_rows = max(1, (rendered + cols - 1) // cols)
    bottom = Emu(_gallery_grid_bottom(y0, actual_rows, row_h, gap) + int(_BLOCK_GAP))
    hidden = max(0, total - rendered)
    if hidden > 0:
        tpl = note_template or labels.get(
            "media_videos_note", "Showing {shown} of {total} relevant video sources. + {n} saved in evidence.",
        )
        bottom = _safe_gallery_note(slide, bottom, tpl.format(shown=rendered, total=total, n=hidden))
    return bottom


# ---------------------------------------------------------------------------
# ORION-style slide 13 — compact fixed-grid layout (4:3)
# ---------------------------------------------------------------------------

YANDEX_RED = RGBColor(0xFC, 0x3F, 0x1C)

# Zone fractions (content width)
ORION_LEFT_FRAC = 0.45
ORION_RIGHT_FRAC = 0.48
ORION_GUTTER_FRAC = 0.035

# Header / content zones (EMU)
ORION_MARKER_TOP = 140000
ORION_MARKER_H = 85000
ORION_HEADLINE_Y = 335000
ORION_HEADLINE_H = 300000
ORION_HEADLINE_W_FRAC = 0.70
ORION_CONTENT_Y = 1500000
ORION_HEADER_GAP = 180000

# Left column block geometry (fixed Y offsets from CONTENT_Y)
ORION_SUMMARY_H = 880000
ORION_SUMMARY_GAP = 230000
ORION_QUERY_TITLE_H = 140000
ORION_TITLE_CHIP_GAP = 80000
ORION_CHIP_H = 285000
ORION_CHIP_GAP = 70000
ORION_QUERY_GAP = 250000
ORION_EXPLAINER_H = 1050000

# Left column internal padding (ORION premium spacing)
ORION_SUMMARY_PAD_X = 95000
ORION_SUMMARY_PAD_TOP = 75000
ORION_SUMMARY_METRIC_H = 340000
ORION_SUMMARY_METRIC_BODY_GAP = 60000
ORION_SUMMARY_PAD_BOTTOM = 120000
ORION_CHIP_PAD_X = 70000
ORION_CHIP_ICON_W = 68000
ORION_CHIP_TEXT_GAP = 55000
ORION_EXPLAINER_PAD_X = 100000
ORION_EXPLAINER_PAD_TOP = 90000
ORION_EXPLAINER_PAD_BOTTOM = 95000
ORION_EXPLAINER_BODY_LINE_SPACING = 1.2
ORION_CHIP_LINE_SPACING = 1.05

# Right panel — content-driven outer frame (generous inset)
ORION_PANEL_Y_OFFSET = 50000
ORION_PANEL_SIDE_PAD = 110000
ORION_PANEL_TOP_PAD = 125000
ORION_PANEL_BOTTOM_PAD = 140000
ORION_SEARCH_BAND_H = 180000
ORION_TABS_GAP = 45000
ORION_TABS_BAND_H = 70000
ORION_TITLE_GAP = 60000
ORION_TITLE_BAND_H = 90000
ORION_GRID_TOP_GAP = 80000
ORION_BADGE_SIZE = 76000
ORION_HIGHLIGHT_BADGE = 44000
ORION_HIGHLIGHT_RING_PAD = 12000
ORION_GRID_GAP = 36000
ORION_FRAME_CORNER_ADJ = 0.035

# Typography (pt) — compact ORION
FS_ORION_MARKER = 7
FS_ORION_DATE = 7
FS_ORION_BRAND = 10
FS_ORION_HEADLINE = 16
FS_ORION_METRIC = 26
FS_ORION_BODY = 8
FS_ORION_CHIP_TITLE = 9
FS_ORION_CHIP = 8
FS_ORION_TAB = 5
FS_ORION_LABEL = 9
FS_ORION_INPUT = 7
FS_ORION_WHY_TITLE = 11

ORION_MIN_CELL = 210000
ORION_MAX_GRID = 9
ORION_MAX_HIGHLIGHTS = 3


def _orion_compact_layout() -> dict[str, int]:
    """Fixed zone coordinates — no dynamic stacking."""
    cw = int(CONTENT_W)
    left_w = int(cw * ORION_LEFT_FRAC)
    gutter = max(int(cw * ORION_GUTTER_FRAC), int(GUTTER))
    right_w = int(cw * ORION_RIGHT_FRAC)
    left_x = int(MARGIN)
    right_x = left_x + left_w + gutter
    content_y = ORION_CONTENT_Y
    panel_y = content_y + ORION_PANEL_Y_OFFSET
    content_bottom = int(CONTENT_SAFE_BOTTOM)

    summary_y = content_y
    summary_bottom = summary_y + ORION_SUMMARY_H
    query_title_y = summary_bottom + ORION_SUMMARY_GAP
    chips_y = query_title_y + ORION_QUERY_TITLE_H + ORION_TITLE_CHIP_GAP
    chips_bottom = chips_y + 2 * ORION_CHIP_H + ORION_CHIP_GAP
    explainer_y = chips_bottom + ORION_QUERY_GAP
    explainer_h = min(ORION_EXPLAINER_H, content_bottom - explainer_y)
    explainer_bottom = explainer_y + explainer_h

    headline_w = int(int(SLIDE_W) * ORION_HEADLINE_W_FRAC)
    headline_bottom = ORION_HEADLINE_Y + ORION_HEADLINE_H

    return {
        "left_x": left_x,
        "left_w": left_w,
        "right_x": right_x,
        "right_w": right_w,
        "gutter": gutter,
        "content_y": content_y,
        "panel_y": panel_y,
        "content_bottom": content_bottom,
        "panel_max_bottom": content_bottom - 70000,
        "headline_w": headline_w,
        "headline_bottom": headline_bottom,
        "summary_y": summary_y,
        "summary_bottom": summary_bottom,
        "query_title_y": query_title_y,
        "chips_y": chips_y,
        "chips_bottom": chips_bottom,
        "explainer_y": explainer_y,
        "explainer_h": explainer_h,
        "explainer_bottom": explainer_bottom,
    }


def _orion_screenshot_frame(slide, x: int, y: int, w: int, h: int):
    """Visible outer screenshot card — shadow (rect) then rounded frame behind content."""
    shadow = slide.shapes.add_shape(RECT, Emu(x + 8000), Emu(y + 12000), Emu(w), Emu(h))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0xE4, 0xE9, 0xF0)
    shadow.line.fill.background()
    card = slide.shapes.add_shape(ROUNDED_RECT, Emu(x), Emu(y), Emu(w), Emu(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = NEUTRAL_LINE
    card.line.width = Pt(1.0)
    try:
        if card.adjustments:
            card.adjustments[0] = ORION_FRAME_CORNER_ADJ
    except Exception:
        pass
    return card


def _orion_shadow_card(slide, x: int, y: int, w: int, h: int, *, radius: int = ROUNDED_RECT):
    return _orion_screenshot_frame(slide, x, y, w, h)


def _orion_yandex_badge(slide, x: int, y: int, size: int = 85000) -> None:
    badge = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(size), Emu(size))
    badge.fill.solid()
    badge.fill.fore_color.rgb = YANDEX_RED
    badge.line.fill.background()
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Я", max(7, int(size / 12700 * 0.5)), WHITE, bold=True)


def _orion_header(
    slide,
    oi: dict,
    brand: str,
    layout: dict[str, int],
    layout_warnings: list[str] | None,
) -> None:
    """Compact ORION header — fixed zones only."""
    x0 = int(MARGIN)
    y0 = ORION_MARKER_TOP
    right_x = layout["right_x"]
    right_w = layout["right_w"]

    marker = textbox(slide, Emu(x0), Emu(y0), Emu(layout["left_w"]), Emu(ORION_MARKER_H))
    sec = str(oi.get("section") or "04  Images")
    num, _, label = sec.partition("  ")
    if not label:
        label, num = sec, "04"
    _run(marker.text_frame.paragraphs[0], f"{num.strip()}  {label.strip()}", FS_ORION_MARKER, ACCENT, bold=True)

    headline = textbox(
        slide, Emu(x0), Emu(ORION_HEADLINE_Y), Emu(layout["headline_w"]), Emu(ORION_HEADLINE_H),
    )
    htf = headline.text_frame
    htf.word_wrap = False
    headline_text = str(oi.get("headline") or "")
    headline_lines = [ln.strip() for ln in headline_text.split("\n") if ln.strip()]
    if not headline_lines:
        headline_lines = [""]
    _run(htf.paragraphs[0], headline_lines[0], FS_ORION_HEADLINE, NEUTRAL_DARK, bold=True)
    for ln in headline_lines[1:]:
        _run(htf.add_paragraph(), ln, FS_ORION_HEADLINE, NEUTRAL_DARK, bold=True)

    date_box = textbox(slide, Emu(right_x), Emu(y0), Emu(right_w), Emu(90000))
    dp = date_box.text_frame.paragraphs[0]
    dp.alignment = PP_ALIGN.RIGHT
    _run(dp, str(oi.get("asOf") or ""), FS_ORION_DATE, NEUTRAL_GRAY)

    brand_label = str(oi.get("brandDisplay") or brand or "Digital Profile Audit")
    brand_box = textbox(slide, Emu(right_x), Emu(y0 + 82000), Emu(right_w), Emu(105000))
    btf = brand_box.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.RIGHT
    _run(bp, truncate(brand_label, 28), FS_ORION_BRAND, NEUTRAL_DARK, bold=True)


def _orion_summary_box(slide, layout: dict[str, int], oi: dict) -> None:
    x, y, w = layout["left_x"], layout["summary_y"], layout["left_w"]
    h = ORION_SUMMARY_H
    frame = slide.shapes.add_shape(ROUNDED_RECT, Emu(x), Emu(y), Emu(w), Emu(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = BG_LIGHT
    frame.line.color.rgb = ACCENT_SOFT
    frame.line.width = Pt(0.75)

    pad_x = ORION_SUMMARY_PAD_X
    inner_w = w - 2 * pad_x

    metric = str(oi.get("metricLabel") or f"{oi.get('metricX', 0)} / {oi.get('metricY', 0)}")
    metric_box = textbox(slide, Emu(x + pad_x), Emu(y + ORION_SUMMARY_PAD_TOP), Emu(inner_w), Emu(ORION_SUMMARY_METRIC_H))
    mtf = metric_box.text_frame
    mtf.word_wrap = True
    mtf.margin_top = 0
    mtf.margin_bottom = 0
    _run(mtf.paragraphs[0], metric, FS_ORION_METRIC, RGBColor(0xE6, 0x5C, 0x00), bold=True)

    body_y = y + ORION_SUMMARY_PAD_TOP + ORION_SUMMARY_METRIC_H + ORION_SUMMARY_METRIC_BODY_GAP
    body_h = h - (ORION_SUMMARY_PAD_TOP + ORION_SUMMARY_METRIC_H + ORION_SUMMARY_METRIC_BODY_GAP) - ORION_SUMMARY_PAD_BOTTOM
    body_box = textbox(slide, Emu(x + pad_x), Emu(body_y), Emu(inner_w), Emu(body_h))
    btf = body_box.text_frame
    btf.word_wrap = True
    btf.margin_top = 0
    btf.margin_bottom = 0
    bp = btf.paragraphs[0]
    bp.line_spacing = 1.15
    _run(bp, truncate(str(oi.get("summaryLine") or ""), 72), FS_ORION_BODY, NEUTRAL_DARK)


def _orion_query_chips(slide, layout: dict[str, int], oi: dict) -> None:
    x, w = layout["left_x"], layout["left_w"]
    y = layout["query_title_y"]
    tb = textbox(slide, Emu(x), Emu(y), Emu(w), Emu(ORION_QUERY_TITLE_H))
    _run(tb.text_frame.paragraphs[0], str(oi.get("queriesTitle") or ""), FS_ORION_CHIP_TITLE, NEUTRAL_DARK, bold=True)

    queries = list(oi.get("queries") or [])[:4]
    if not queries:
        return
    chip_y = layout["chips_y"]
    gap = ORION_CHIP_GAP
    chip_h = ORION_CHIP_H
    cols = 2
    chip_w = (w - gap) // cols
    pad_x = ORION_CHIP_PAD_X
    icon_w = ORION_CHIP_ICON_W
    for idx, q in enumerate(queries):
        row, col = divmod(idx, cols)
        cx = x + col * (chip_w + gap)
        cy = chip_y + row * (chip_h + gap)
        chip = slide.shapes.add_shape(ROUNDED_RECT, Emu(cx), Emu(cy), Emu(chip_w), Emu(chip_h))
        chip.fill.solid()
        chip.fill.fore_color.rgb = WHITE
        chip.line.color.rgb = NEUTRAL_LINE
        chip.line.width = Pt(0.5)
        icon = textbox(slide, Emu(cx + pad_x), Emu(cy + (chip_h - icon_w) // 2), Emu(icon_w), Emu(icon_w))
        itf = icon.text_frame
        itf.margin_top = 0
        itf.margin_bottom = 0
        ip = itf.paragraphs[0]
        ip.alignment = PP_ALIGN.CENTER
        _run(ip, "⌕", FS_ORION_CHIP, ACCENT)
        text_x = cx + pad_x + icon_w + ORION_CHIP_TEXT_GAP
        text_w = cx + chip_w - pad_x - text_x
        qbox = textbox(slide, Emu(text_x), Emu(cy), Emu(text_w), Emu(chip_h))
        qtf = qbox.text_frame
        qtf.word_wrap = True
        qtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        qtf.margin_left = 0
        qtf.margin_right = 0
        qtf.margin_top = 0
        qtf.margin_bottom = 0
        qp = qtf.paragraphs[0]
        qp.line_spacing = ORION_CHIP_LINE_SPACING
        _run(qp, truncate(q, 42), FS_ORION_CHIP, NEUTRAL_DARK)


def _orion_explainer_box(slide, layout: dict[str, int], oi: dict) -> None:
    x, y, w = layout["left_x"], layout["explainer_y"], layout["left_w"]
    h = layout["explainer_h"]
    box = slide.shapes.add_shape(ROUNDED_RECT, Emu(x), Emu(y), Emu(w), Emu(h))
    box.fill.solid()
    box.fill.fore_color.rgb = BG_PANEL
    box.line.fill.background()
    pad_x = ORION_EXPLAINER_PAD_X
    inner = textbox(
        slide,
        Emu(x + pad_x),
        Emu(y + ORION_EXPLAINER_PAD_TOP),
        Emu(w - 2 * pad_x),
        Emu(h - ORION_EXPLAINER_PAD_TOP - ORION_EXPLAINER_PAD_BOTTOM),
    )
    tf = inner.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    _run(tf.paragraphs[0], str(oi.get("whyTitle") or ""), FS_ORION_WHY_TITLE, NEUTRAL_DARK, bold=True)
    body_p = tf.add_paragraph()
    body_p.line_spacing = ORION_EXPLAINER_BODY_LINE_SPACING
    body_p.space_before = Pt(4)
    _run(body_p, truncate(str(oi.get("whyBody") or ""), 150), FS_ORION_BODY, NEUTRAL_GRAY)


def _orion_screenshot_plan(layout: dict[str, int], oi: dict) -> dict[str, int | list]:
    """Content-driven outer frame — union of internal bands + explicit outer padding."""
    panel_x = layout["right_x"]
    panel_w = layout["right_w"]
    panel_y = layout["panel_y"]
    footer_limit = layout["panel_max_bottom"]

    side = ORION_PANEL_SIDE_PAD
    top = ORION_PANEL_TOP_PAD
    inner_x = panel_x + side
    inner_w = max(1, panel_w - 2 * side)

    search_y = panel_y + top
    search_h = ORION_SEARCH_BAND_H
    tabs_y = search_y + search_h + ORION_TABS_GAP
    tabs_h = ORION_TABS_BAND_H
    title_y = tabs_y + tabs_h + ORION_TITLE_GAP
    title_h = ORION_TITLE_BAND_H
    grid_y = title_y + title_h + ORION_GRID_TOP_GAP

    items = list(oi.get("gridItems") or [])[:ORION_MAX_GRID]
    grid_limit_y = footer_limit - ORION_PANEL_BOTTOM_PAD - 110000
    avail_h = max(1, grid_limit_y - grid_y)
    cols, rows, cell_w, cell_h, gap = _orion_grid_geometry(inner_w, avail_h, len(items))
    slots = min(len(items), cols * rows, ORION_MAX_GRID)
    rows_used = 0 if slots <= 0 else (slots - 1) // cols + 1

    badge_size = ORION_BADGE_SIZE
    badge_y = search_y + (search_h - badge_size) // 2
    search_x = inner_x + badge_size + 50000
    search_w = inner_w - badge_size - 50000

    tab_widths = (210000, 380000, 210000, 210000)
    tab_gap = 35000
    tabs_right = inner_x + sum(tab_widths) + tab_gap * (len(tab_widths) - 1)
    tabs_right = min(inner_x + inner_w, tabs_right)

    ring_out = 3500 + ORION_HIGHLIGHT_RING_PAD
    content_boxes: list[dict[str, int]] = [
        {"x": inner_x, "y": badge_y, "right": inner_x + badge_size, "bottom": badge_y + badge_size},
        {"x": search_x, "y": search_y, "right": search_x + search_w, "bottom": search_y + search_h},
        {"x": inner_x, "y": tabs_y, "right": tabs_right, "bottom": tabs_y + tabs_h},
        {"x": inner_x, "y": title_y, "right": inner_x + inner_w, "bottom": title_y + title_h},
    ]
    for idx in range(slots):
        row, col = divmod(idx, cols)
        cx = inner_x + col * (cell_w + gap)
        cy = grid_y + row * (cell_h + gap)
        item = items[idx]
        highlight = bool(item.get("highlight"))
        badge_extra = ORION_HIGHLIGHT_BADGE + 3000 if highlight else 0
        content_boxes.append(
            {
                "x": cx - ring_out,
                "y": cy - ring_out,
                "right": cx + cell_w + ring_out,
                "bottom": cy + cell_h + ring_out + badge_extra,
            }
        )

    min_x = min(b["x"] for b in content_boxes)
    min_y = min(b["y"] for b in content_boxes)
    max_x = max(b["right"] for b in content_boxes)
    max_y = max(b["bottom"] for b in content_boxes)

    out_l = out_t = out_r = 100000
    out_b = 120000
    frame_x = min_x - out_l
    frame_y = min_y - out_t
    frame_right = max_x + out_r
    frame_bottom = max_y + out_b
    if frame_bottom > footer_limit:
        frame_bottom = footer_limit
    frame_w = frame_right - frame_x
    frame_h = frame_bottom - frame_y

    grid_bottom = max_y if slots > 0 else grid_y

    return {
        "frame_x": frame_x,
        "frame_y": frame_y,
        "frame_w": frame_w,
        "frame_h": frame_h,
        "panel_x": frame_x,
        "panel_y": frame_y,
        "panel_w": frame_w,
        "panel_h": frame_h,
        "panel_bottom": frame_bottom,
        "inner_x": inner_x,
        "inner_w": inner_w,
        "search_y": search_y,
        "search_h": search_h,
        "search_x": search_x,
        "search_w": search_w,
        "tabs_y": tabs_y,
        "tabs_h": tabs_h,
        "title_y": title_y,
        "title_h": title_h,
        "grid_y": grid_y,
        "grid_bottom": grid_bottom,
        "cols": cols,
        "rows": rows,
        "rows_used": rows_used,
        "cell_w": cell_w,
        "cell_h": cell_h,
        "gap": gap,
        "slots": slots,
        "items": items,
        "content_left": min_x,
        "content_right": max_x,
        "content_top": min_y,
        "content_bottom": max_y,
    }


def _orion_centered_input_text(slide, x: int, y: int, w: int, h: int, text: str, size: int, color: RGBColor) -> None:
    box = textbox(slide, Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0
    tf.margin_bottom = 0
    _run(tf.paragraphs[0], text, size, color)


def _orion_search_tabs(slide, x: int, y: int, w: int, h: int, tabs: list[str], active_idx: int = 1) -> None:
    shown = [tabs[0], tabs[1], tabs[2], tabs[3]] if len(tabs) >= 4 else tabs[:4]
    widths = (210000, 380000, 210000, 210000)
    gap = 35000
    cx = x
    for i, tab in enumerate(shown):
        tw = widths[i] if i < len(widths) else 210000
        tb = textbox(slide, Emu(cx), Emu(y), Emu(tw), Emu(h))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        p = tf.paragraphs[0]
        idx = tabs.index(tab) if tab in tabs else i
        color = NEUTRAL_DARK if idx == active_idx else NEUTRAL_GRAY
        _run(p, tab, FS_ORION_TAB, color, bold=(idx == active_idx))
        cx += tw + gap
        if cx > x + w - 60000:
            break


def _orion_grid_geometry(inner_w: int, inner_h: int, item_count: int) -> tuple[int, int, int, int, int]:
    """Prefer 3×3, degrade to 2×3 → 2×2."""
    n = min(max(item_count, 1), ORION_MAX_GRID)
    for cols, rows in ((3, 3), (2, 3), (2, 2)):
        if cols * rows < min(n, 4):
            continue
        gap = ORION_GRID_GAP
        cell_w = (inner_w - gap * (cols - 1)) // cols
        cell_h = (inner_h - gap * (rows - 1)) // rows
        if cell_w >= ORION_MIN_CELL and cell_h >= ORION_MIN_CELL:
            return cols, rows, cell_w, cell_h, gap
    gap = ORION_GRID_GAP
    cols, rows = 2, 2
    cell_w = (inner_w - gap) // cols
    cell_h = (inner_h - gap) // rows
    return cols, rows, cell_w, cell_h, gap


def _orion_thumb_in_cell(
    slide,
    item: dict[str, Any],
    cell_x: int,
    cell_y: int,
    cell_w: int,
    cell_h: int,
    *,
    highlight: bool = False,
) -> bool:
    import io

    pad = 4000
    ix, iy, iw, ih = cell_x + pad, cell_y + pad, cell_w - 2 * pad, cell_h - 2 * pad
    raw = _image_bytes_from_item(item)
    if not raw:
        return False
    try:
        pic = _fit_picture_contain(slide, io.BytesIO(raw), ix, iy, iw, ih)
        if int(pic.width) <= 0 or int(pic.height) <= 0:
            return False
    except Exception:
        return False

    if highlight:
        ring = slide.shapes.add_shape(ROUNDED_RECT, Emu(ix - 3500), Emu(iy - 3500), Emu(iw + 7000), Emu(ih + 7000))
        ring.fill.background()
        ring.line.color.rgb = DANGER
        ring.line.width = Pt(0.75)
        bs = ORION_HIGHLIGHT_BADGE
        badge = slide.shapes.add_shape(1, Emu(ix + 3000), Emu(iy + 3000), Emu(bs), Emu(bs))
        badge.fill.solid()
        badge.fill.fore_color.rgb = DANGER
        badge.line.fill.background()
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        _run(bp, "×", 7, WHITE, bold=True)
    return True


def _orion_synthetic_image_panel(
    slide,
    layout: dict[str, int],
    oi: dict,
    *,
    layout_warnings: list[str] | None = None,
) -> None:
    plan = _orion_screenshot_plan(layout, oi)
    fx = int(plan["frame_x"])
    fy = int(plan["frame_y"])
    fw = int(plan["frame_w"])
    fh = int(plan["frame_h"])
    _orion_screenshot_frame(slide, fx, fy, fw, fh)

    inner_x = int(plan["inner_x"])
    inner_w = int(plan["inner_w"])
    search_y = int(plan["search_y"])
    search_h = int(plan["search_h"])
    search_x = int(plan["search_x"])
    search_w = int(plan["search_w"])

    badge_size = ORION_BADGE_SIZE
    badge_y = search_y + (search_h - badge_size) // 2
    _orion_yandex_badge(slide, inner_x, badge_y, badge_size)
    bar = slide.shapes.add_shape(ROUNDED_RECT, Emu(search_x), Emu(search_y), Emu(search_w), Emu(search_h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE
    bar.line.color.rgb = NEUTRAL_LINE
    bar.line.width = Pt(0.5)

    close_w = 55000
    inp_pad = 60000
    query_text = truncate(str(oi.get("primaryQuery") or ""), 28)
    _orion_centered_input_text(
        slide,
        search_x + inp_pad,
        search_y,
        search_w - inp_pad - close_w - 20000,
        search_h,
        query_text,
        FS_ORION_INPUT,
        NEUTRAL_DARK,
    )
    _orion_centered_input_text(
        slide,
        search_x + search_w - close_w,
        search_y,
        close_w,
        search_h,
        "×",
        FS_ORION_TAB + 1,
        NEUTRAL_GRAY,
    )

    _orion_search_tabs(
        slide, inner_x, int(plan["tabs_y"]), inner_w, int(plan["tabs_h"]),
        list(oi.get("tabs") or []), active_idx=1,
    )
    lbl = textbox(slide, Emu(inner_x), Emu(int(plan["title_y"])), Emu(inner_w), Emu(int(plan["title_h"])))
    ltf = lbl.text_frame
    ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ltf.margin_top = 0
    ltf.margin_bottom = 0
    _run(ltf.paragraphs[0], str(oi.get("gridTitle") or ""), FS_ORION_LABEL, NEUTRAL_DARK, bold=True)

    grid_top = int(plan["grid_y"])
    cols = int(plan["cols"])
    cell_w = int(plan["cell_w"])
    cell_h = int(plan["cell_h"])
    gap = int(plan["gap"])
    slots = int(plan["slots"])
    items = list(plan["items"])
    hi_used = 0
    for idx in range(slots):
        row, col = divmod(idx, cols)
        cx = inner_x + col * (cell_w + gap)
        cy = grid_top + row * (cell_h + gap)
        item = items[idx]
        highlight = bool(item.get("highlight")) and hi_used < ORION_MAX_HIGHLIGHTS
        if highlight:
            hi_used += 1
        _orion_thumb_in_cell(slide, item, cx, cy, cell_w, cell_h, highlight=highlight)


def orion_images_slide(
    slide,
    blk: dict[str, Any],
    vm: dict[str, Any],
    ctx,
    *,
    layout_warnings: list[str] | None = None,
) -> None:
    """Full ORION-style slide 13 — compact fixed-grid analytical layout."""
    L = vm.get("labels") or {}
    oi = dict(blk.get("orionImages") or {})
    brand = str(getattr(ctx, "brand", None) or vm.get("meta", {}).get("brand", ""))
    layout = _orion_compact_layout()

    set_bg(slide, BG_LIGHT)
    footer(slide, brand, getattr(ctx, "page", None), getattr(ctx, "total", None))

    _orion_header(slide, oi, brand, layout, layout_warnings)

    if not oi.get("gridItems"):
        nd = textbox(slide, MARGIN, Emu(layout["content_y"] + 120000), CONTENT_W, Emu(400000))
        _run(nd.text_frame.paragraphs[0], str(oi.get("noData") or L.get("nd_no_relevant_images", "")), FS_BODY, NEUTRAL_GRAY, italic=True)
        return

    _orion_summary_box(slide, layout, oi)
    _orion_query_chips(slide, layout, oi)
    _orion_explainer_box(slide, layout, oi)
    _orion_synthetic_image_panel(slide, layout, oi, layout_warnings=layout_warnings)
