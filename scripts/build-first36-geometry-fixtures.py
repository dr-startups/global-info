#!/usr/bin/env python3
"""Build intentional First36 geometry PPTX fixtures (background/container overlaps)."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

SLIDE_W = 11_704_320
SLIDE_H = 7_315_200
OUT_DIR = Path(__file__).resolve().parents[1] / "src/modules/digital-profile/orion-golden/classic/fixtures/geometry"


def _blank_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    return prs


def _add_slide(prs: Presentation):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    return prs.slides.add_slide(layout)


def build_bg_plus_text(path: Path) -> None:
    """Full-slide background + textbox → expected PASS."""
    prs = _blank_prs()
    slide = _add_slide(prs)
    bg = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H))
    bg.name = "orion_bg_full"
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0B, 0x1A, 0x33)
    bg.line.fill.background()
    box = slide.shapes.add_textbox(Emu(480_000), Emu(800_000), Emu(8_000_000), Emu(1_200_000))
    box.name = "orion_text_title"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Background plus textbox should pass geometry v2"
    r.font.size = Pt(22)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    prs.save(str(path))


def build_textbox_in_card(path: Path) -> None:
    """Textbox contained in card → PASS."""
    prs = _blank_prs()
    slide = _add_slide(prs)
    card = slide.shapes.add_shape(1, Emu(480_000), Emu(700_000), Emu(10_000_000), Emu(4_000_000))
    card.name = "orion_card_main"
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    box = slide.shapes.add_textbox(Emu(700_000), Emu(1_000_000), Emu(9_000_000), Emu(2_000_000))
    box.name = "orion_text_body"
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Text inside card is intentional containment"
    r.font.size = Pt(14)
    prs.save(str(path))


def build_text_collision(path: Path) -> None:
    """Two overlapping text boxes → CRITICAL."""
    prs = _blank_prs()
    slide = _add_slide(prs)
    a = slide.shapes.add_textbox(Emu(500_000), Emu(800_000), Emu(5_000_000), Emu(1_500_000))
    a.name = "orion_text_a"
    a.text_frame.paragraphs[0].add_run().text = "First colliding text box content here"
    b = slide.shapes.add_textbox(Emu(1_500_000), Emu(1_000_000), Emu(5_000_000), Emu(1_500_000))
    b.name = "orion_text_b"
    b.text_frame.paragraphs[0].add_run().text = "Second colliding text box content here"
    prs.save(str(path))


def build_out_of_bounds(path: Path) -> None:
    """Text shape past slide bounds → CRITICAL."""
    prs = _blank_prs()
    slide = _add_slide(prs)
    box = slide.shapes.add_textbox(Emu(500_000), Emu(6_800_000), Emu(4_000_000), Emu(1_200_000))
    box.name = "orion_text_oob"
    box.text_frame.paragraphs[0].add_run().text = "This text box extends past the slide bottom edge"
    prs.save(str(path))


def build_intentional_overlaps_deck(path: Path) -> None:
    """Multi-slide fixture with intentional bg/card overlaps (integration)."""
    prs = _blank_prs()
    # slide 1: bg + text
    s1 = _add_slide(prs)
    bg = s1.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H))
    bg.name = "orion_bg_full"
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0B, 0x1A, 0x33)
    t = s1.shapes.add_textbox(Emu(480_000), Emu(900_000), Emu(9_000_000), Emu(1_000_000))
    t.name = "orion_text_title"
    t.text_frame.paragraphs[0].add_run().text = "Intentional background overlap page"
    # slide 2: card + text
    s2 = _add_slide(prs)
    card = s2.shapes.add_shape(1, Emu(400_000), Emu(600_000), Emu(10_500_000), Emu(5_000_000))
    card.name = "orion_card_main"
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    body = s2.shapes.add_textbox(Emu(700_000), Emu(900_000), Emu(9_500_000), Emu(3_000_000))
    body.name = "orion_text_body"
    body.text_frame.paragraphs[0].add_run().text = "Card containment is allowed by geometry v2"
    prs.save(str(path))


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_bg_plus_text(OUT_DIR / "fixture-bg-plus-text.pptx")
    build_textbox_in_card(OUT_DIR / "fixture-textbox-in-card.pptx")
    build_text_collision(OUT_DIR / "fixture-text-collision.pptx")
    build_out_of_bounds(OUT_DIR / "fixture-out-of-bounds.pptx")
    build_intentional_overlaps_deck(OUT_DIR / "fixture-intentional-overlaps.pptx")
    print(f"wrote fixtures under {OUT_DIR}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
