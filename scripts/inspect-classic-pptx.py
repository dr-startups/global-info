from pptx import Presentation
from pathlib import Path
import json
import re

p = Path(r"c:\Users\Константин\Downloads\orion-classic-audit.pptx")
out = Path(r"c:\Global Info\storage\digital-profile\qa-r10-11-classic-orion-audit\pptx-inspect.json")
out.parent.mkdir(parents=True, exist_ok=True)
prs = Presentation(str(p))
slides = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    slides.append({"n": i, "texts": texts, "chars": sum(len(t) for t in texts)})

titles = [(s["texts"][0].split("\n")[0] if s["texts"] else "") for s in slides]
commercial = sum(
    1
    for t in titles
    if any(x in t for x in ["Наше предложение", "обзор продукта", "Решение", "О нас"])
)
joined_all = "\n".join("\n".join(s["texts"]) for s in slides)
needed = {
    "serp_matrix": bool(re.search(r"позиц|TOP-?20|доля нежелат|нежелательн.*ссыл", joined_all, re.I)),
    "themes_with_counts": bool(re.search(r"\d+\s*(публикац|ссыл|из\s+\d+)", joined_all, re.I)),
    "suggestions_audit": bool(re.search(r"подсказ", joined_all, re.I)),
    "related_queries": bool(re.search(r"похож|смежн", joined_all, re.I)),
    "wikipedia": bool(re.search(r"википед|wikipedia", joined_all, re.I)),
    "dow_jones": bool(re.search(r"Dow Jones|RCA", joined_all, re.I)),
    "lexis": bool(re.search(r"LexisNexis|PEP", joined_all, re.I)),
    "concrete_facts": bool(
        re.search(r"Трансмаш|Махмудов|Бокарев|Ликсутов|Молдав|санкц|офшор", joined_all, re.I)
    ),
    "risk_level_banner": bool(
        re.search(r"Крайне высокий|Высокий уровень риска|Compliance риски", joined_all, re.I)
    ),
    "url_examples": bool(re.search(r"https?://|rucompromat|forbes\.ru|tadviser", joined_all, re.I)),
}
payload = {
    "slideCount": len(slides),
    "commercialApprox": commercial,
    "auditApprox": len(slides) - commercial,
    "needed": needed,
    "titles": titles,
    "slides": slides,
}
out.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
print("slides", len(slides), "commercialApprox", commercial)
print(json.dumps(needed, ensure_ascii=False, indent=2))
for s in slides:
    t0 = s["texts"][0].replace("\n", " | ")[:180] if s["texts"] else ""
    print(f"{s['n']:02d} [{s['chars']}] {t0}")
