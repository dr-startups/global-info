# -*- coding: utf-8 -*-
"""Manual-quality visual acceptance for the assembled ORION deck.

Fail-closed checks on the RENDERED PPTX (not just the model):
  * emptySidebarCount=0        — no visually empty titled sidebar panels;
  * emptyTitledContainerCount=0 — no titled content container without body;
  * materiallyEmptyPageCount=0  — no near-blank full page unless it is a
                                  divider or an explicit valid empty state;
  * blankVisualPageCount=0      — every visual page renders an asset or an
                                  explicit client-safe fallback label.

Usage:
  python inspect-deck-manual-visual.py <pptx> <report-deck-manifest.json> <out.json>
"""

import json
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Canonical templates allowed to be structurally sparse.
STRUCTURAL_TEMPLATES = {"cover", "toc", "section-divider", "coverage-empty-state", "continuation"}
# Canonical templates that are expected to carry a visual asset (or fallback).
VISUAL_TEMPLATES = {
    "serp-screenshot-analysis",
    "suggestions",
    "image-grid",
    "wikipedia-knowledge",
    "ai-overview",
    "related-queries",
}
FALLBACK_MARKERS = ("недоступ", "не зафиксирован", "нет данных", "VISUAL_ASSET_UNAVAILABLE")


def iter_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(sh.shapes)
        else:
            yield sh


def shape_text(sh) -> str:
    if not getattr(sh, "has_text_frame", False):
        return ""
    return (sh.text_frame.text or "").strip()


def analyze(pptx_path: str, manifest_path: str) -> dict:
    prs = Presentation(pptx_path)
    with open(manifest_path, encoding="utf-8") as fh:
        manifest = json.load(fh)
    template_by_page = {s["pageNumber"]: s["templateId"] for s in manifest.get("slides", [])}

    sw, sh_h = prs.slide_width, prs.slide_height

    empty_sidebars = []
    empty_containers = []
    materially_empty = []
    blank_visual = []

    for idx, slide in enumerate(prs.slides, start=1):
        template = template_by_page.get(idx, "")
        shapes = list(iter_shapes(slide.shapes))
        pictures = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        tables = [s for s in shapes if getattr(s, "has_table", False)]
        texts = [(s, shape_text(s)) for s in shapes if shape_text(s)]

        cards = [
            s
            for s in shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.width and s.height
            and s.width > sw * 0.12
            and s.height > sh_h * 0.10
        ]

        def inside(inner, outer) -> bool:
            cx = inner.left + inner.width / 2
            cy = inner.top + inner.height / 2
            return (
                outer.left <= cx <= outer.left + outer.width
                and outer.top <= cy <= outer.top + outer.height
            )

        for card in cards:
            inner_texts = [t for s, t in texts if s is not card and inside(s, card)]
            inner_pics = [p for p in pictures if inside(p, card)]
            inner_tables = [t for t in tables if inside(t, card)]
            joined = " ".join(inner_texts).strip()
            own = shape_text(card)
            total = (own + " " + joined).strip()
            if inner_pics or inner_tables:
                continue
            # Big right-hand analytical sidebar with (almost) no body text.
            is_sidebar = card.left > sw * 0.55 and card.height > sh_h * 0.35
            if is_sidebar and len(total) < 40:
                empty_sidebars.append({"page": idx, "template": template, "text": total})
                continue
            # Any titled container whose entire content is a bare short title.
            if len(total) > 0 and len(total) <= 8:
                empty_containers.append({"page": idx, "template": template, "text": total})

        # Materially empty page: exclude title (top) and footer (bottom) text.
        body_chars = 0
        for s, t in texts:
            top_frac = s.top / sh_h if s.top is not None else 0.5
            if top_frac < 0.12 or top_frac > 0.90:
                continue
            body_chars += len(t)
        if (
            template not in STRUCTURAL_TEMPLATES
            and not pictures
            and not tables
            and body_chars < 40
        ):
            materially_empty.append({"page": idx, "template": template, "bodyChars": body_chars})

        # Visual pages must render an asset, an explicit labeled fallback, or a
        # substantive text layout (downgraded template carrying real content).
        if template in VISUAL_TEMPLATES and not pictures:
            page_text = " ".join(t for _, t in texts)
            has_fallback = any(m.lower() in page_text.lower() for m in FALLBACK_MARKERS)
            if not has_fallback and body_chars < 60:
                blank_visual.append({"page": idx, "template": template, "bodyChars": body_chars})

    report = {
        "version": "deck-manual-visual-report-v1",
        "pageCount": len(list(prs.slides)),
        "emptySidebarCount": len(empty_sidebars),
        "emptySidebars": empty_sidebars,
        "emptyTitledContainerCount": len(empty_containers),
        "emptyTitledContainers": empty_containers,
        "materiallyEmptyPageCount": len(materially_empty),
        "materiallyEmptyPages": materially_empty,
        "blankVisualPageCount": len(blank_visual),
        "blankVisualPages": blank_visual,
    }
    report["passed"] = (
        report["emptySidebarCount"] == 0
        and report["emptyTitledContainerCount"] == 0
        and report["materiallyEmptyPageCount"] == 0
        and report["blankVisualPageCount"] == 0
    )
    return report


def main() -> int:
    pptx_path, manifest_path, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    report = analyze(pptx_path, manifest_path)
    with open(out_path, "w", encoding="utf-8") as fh:
        json.dump(report, fh, ensure_ascii=False, indent=2)
    print(json.dumps({k: v for k, v in report.items() if isinstance(v, (int, bool))}, ensure_ascii=False))
    return 0 if report["passed"] else 1


if __name__ == "__main__":
    sys.exit(main())
