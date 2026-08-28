#!/usr/bin/env python3
"""First36 PPTX geometry inspector v2 — role-aware layout QA (no naive shape-intersection = CRITICAL)."""

from __future__ import annotations

import json
import sys
import zipfile
from pathlib import Path
from typing import Any

INSPECTOR_VERSION = "first36-geometry-v2"

# Match renderer/orion_golden_renderer.py 16:10 master
SLIDE_W = 11_704_320
SLIDE_H = 7_315_200
FOOTER_Y = SLIDE_H - 440_000
CONTENT_BOTTOM = SLIDE_H - 700_000
FOOTER_ZONE_TOP = FOOTER_Y - 40_000

# Intersection area / min(area_a, area_b) above this → text-text collision
TEXT_COLLISION_RATIO = 0.12
# Shape must be this fraction inside container to count as intentional containment
CONTAINMENT_RATIO = 0.85
# Near full-slide → background
BACKGROUND_AREA_RATIO = 0.88
# Thin shapes → decorative / border
DECORATIVE_MAX_THICKNESS = 80_000


def _area(b: dict[str, int]) -> int:
    return max(0, b["cx"]) * max(0, b["cy"])


def _intersect_area(a: dict[str, int], b: dict[str, int]) -> int:
    x0 = max(a["x"], b["x"])
    y0 = max(a["y"], b["y"])
    x1 = min(a["x"] + a["cx"], b["x"] + b["cx"])
    y1 = min(a["y"] + a["cy"], b["y"] + b["cy"])
    if x1 <= x0 or y1 <= y0:
        return 0
    return (x1 - x0) * (y1 - y0)


def _contained(inner: dict[str, int], outer: dict[str, int]) -> bool:
    ia = _area(inner)
    if ia <= 0:
        return False
    return (_intersect_area(inner, outer) / ia) >= CONTAINMENT_RATIO


def _out_of_bounds(b: dict[str, int], margin: int = 20_000) -> bool:
    return (
        b["x"] < -margin
        or b["y"] < -margin
        or b["x"] + b["cx"] > SLIDE_W + margin
        or b["y"] + b["cy"] > SLIDE_H + margin
    )


def _classify_role(shape: Any, z: int) -> str:
    name = (getattr(shape, "name", None) or "").lower()
    if "footer" in name or "orion_footer" in name:
        return "footer"
    if "background" in name or "orion_bg" in name:
        return "background"
    if "border" in name or "decor" in name or "orion_decor" in name:
        return "decorative"
    if "card" in name or "container" in name or "orion_card" in name:
        return "container"
    if "table" in name or "orion_table" in name:
        return "table"
    if "image" in name or "picture" in name or "orion_img" in name:
        return "image"

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        if st == MSO_SHAPE_TYPE.TABLE:
            return "table"
        if st == MSO_SHAPE_TYPE.LINE:
            return "decorative"
    except Exception:  # noqa: BLE001
        pass

    has_text = False
    text_len = 0
    try:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            t = (shape.text_frame.text or "").strip()
            text_len = len(t)
            has_text = text_len > 0
    except Exception:  # noqa: BLE001
        pass

    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    area = max(0, width) * max(0, height)
    slide_area = SLIDE_W * SLIDE_H

    if top >= FOOTER_ZONE_TOP and height < 500_000:
        return "footer"

    if width <= DECORATIVE_MAX_THICKNESS or height <= DECORATIVE_MAX_THICKNESS:
        return "decorative"

    if area >= slide_area * BACKGROUND_AREA_RATIO and not has_text:
        return "background"

    if has_text:
        return "text"

    if area >= slide_area * 0.05:
        return "container"

    return "decorative"


def _drawn_text(shape: Any) -> str:
    """Нарисованный текст фигуры: рамка текста или ячейки таблицы.

    Геометрия рядом читает только рамку — ей нужна длина текста для роли
    фигуры, а таблицу она узнаёт по типу. Ворот следа поля обязан видеть
    больше: у страницы выдачи почти весь текст лежит в ячейках таблицы, и без
    них страница выглядела бы пустой.
    """
    try:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            return (shape.text_frame.text or "").strip()
        if getattr(shape, "has_table", False) and shape.has_table:
            cells = [(cell.text or "").strip() for row in shape.table.rows for cell in row.cells]
            return "\n".join(c for c in cells if c)
    except Exception:  # noqa: BLE001
        return ""
    return ""


def page_texts(pptx: Path) -> list[str]:
    """Текст каждой страницы презентации — по одной строке на страницу."""
    from pptx import Presentation

    return [
        "\n".join(t for t in (_drawn_text(shape) for shape in slide.shapes) if t)
        for slide in Presentation(str(pptx)).slides
    ]


def _extract_shapes(slide: Any, page: int) -> list[dict[str, Any]]:
    shapes: list[dict[str, Any]] = []
    for z, shape in enumerate(slide.shapes):
        try:
            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
        except Exception:  # noqa: BLE001
            continue
        text = ""
        has_text = False
        try:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = (shape.text_frame.text or "").strip()
                has_text = bool(text)
        except Exception:  # noqa: BLE001
            pass
        role = _classify_role(shape, z)
        if role == "text" and top >= FOOTER_ZONE_TOP and len(text) < 24:
            role = "footer"
        shape_type = "unknown"
        try:
            shape_type = str(shape.shape_type)
        except Exception:  # noqa: BLE001
            pass
        shapes.append(
            {
                "page": page,
                "id": z,
                "name": getattr(shape, "name", None) or f"shape_{z}",
                "type": shape_type,
                "role": role,
                "bbox": {"x": left, "y": top, "cx": width, "cy": height},
                "zOrder": z,
                "hasText": has_text,
                "textLength": len(text),
            }
        )
    return shapes


def _ignore_pair(a: dict[str, Any], b: dict[str, Any]) -> bool:
    ignore_roles = {"background", "border", "decorative", "footer"}
    if a["role"] in ignore_roles or b["role"] in ignore_roles:
        return True
    for inner, outer in ((a, b), (b, a)):
        if outer["role"] in {"container", "card"} and inner["role"] in {"text", "image", "table"}:
            if _contained(inner["bbox"], outer["bbox"]):
                return True
        if outer["role"] == "background" and inner["role"] in {"text", "image", "table", "container"}:
            return True
    return False


def inspect_presentation(pptx: Path, *, expect_pages: int | None = None) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(pptx))

    all_shapes: list[dict[str, Any]] = []
    overlaps: list[dict[str, Any]] = []
    overflow: list[dict[str, Any]] = []
    clipping: list[dict[str, Any]] = []
    empty_pages: list[dict[str, Any]] = []
    pages_meta: list[dict[str, Any]] = []

    slide_count = len(prs.slides)
    scanned = 0
    for page_idx, slide in enumerate(prs.slides):
        page = page_idx + 1
        if expect_pages and page > expect_pages:
            break
        scanned = page
        shapes = _extract_shapes(slide, page)
        all_shapes.extend(shapes)

        meaningful = [
            s
            for s in shapes
            if s["role"] in {"text", "image", "table"} and (s["role"] != "text" or s["textLength"] >= 8)
        ]
        if not meaningful:
            empty_pages.append(
                {
                    "page": page,
                    "code": "empty-page",
                    "severity": "CRITICAL",
                    "detail": "page has no meaningful text/image/table content",
                }
            )

        for s in shapes:
            if s["role"] in {"background", "decorative", "border", "footer"}:
                continue
            if _out_of_bounds(s["bbox"]):
                overflow.append(
                    {
                        "page": page,
                        "code": "out-of-bounds",
                        "severity": "CRITICAL",
                        "detail": (
                            f"content shape '{s['name']}' role={s['role']} "
                            f"bbox={s['bbox']} exceeds slide bounds"
                        ),
                        "shapeId": s["id"],
                        "role": s["role"],
                    }
                )

        texts = [s for s in shapes if s["role"] == "text"]
        images = [s for s in shapes if s["role"] in {"image", "table"}]
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                a, b = texts[i], texts[j]
                if _ignore_pair(a, b):
                    continue
                ia = _intersect_area(a["bbox"], b["bbox"])
                if ia <= 0:
                    continue
                min_a = min(_area(a["bbox"]), _area(b["bbox"]))
                ratio = ia / min_a if min_a else 0
                if ratio >= TEXT_COLLISION_RATIO:
                    overlaps.append(
                        {
                            "page": page,
                            "code": "text-text-collision",
                            "severity": "CRITICAL",
                            "detail": (
                                f"text-text collision ratio={ratio:.2f} "
                                f"'{a['name']}' vs '{b['name']}'"
                            ),
                            "shapeIds": [a["id"], b["id"]],
                        }
                    )

        for t in texts:
            for img in images:
                if _ignore_pair(t, img):
                    continue
                ia = _intersect_area(t["bbox"], img["bbox"])
                if ia <= 0:
                    continue
                t_mid = t["bbox"]["y"] + t["bbox"]["cy"] / 2
                img_bottom = img["bbox"]["y"] + img["bbox"]["cy"]
                if t_mid >= img_bottom - 80_000:
                    continue
                min_a = min(_area(t["bbox"]), _area(img["bbox"]))
                ratio = ia / min_a if min_a else 0
                if ratio >= 0.08:
                    overlaps.append(
                        {
                            "page": page,
                            "code": "text-over-image",
                            "severity": "CRITICAL",
                            "detail": (
                                f"text '{t['name']}' overlaps {img['role']} '{img['name']}' "
                                f"ratio={ratio:.2f}"
                            ),
                            "shapeIds": [t["id"], img["id"]],
                        }
                    )

        pages_meta.append(
            {
                "page": page,
                "shapeCount": len(shapes),
                "roles": {
                    r: sum(1 for s in shapes if s["role"] == r) for r in sorted({s["role"] for s in shapes})
                },
                "shapes": shapes,
            }
        )

    telemetry_path = pptx.parent / "layout-telemetry.json"
    if telemetry_path.is_file():
        try:
            telemetry = json.loads(telemetry_path.read_text(encoding="utf-8"))
            for row in telemetry.get("entries") or telemetry.get("textBoxes") or []:
                dropped = int(row.get("droppedBullets") or 0) + int(row.get("droppedLines") or 0)
                if dropped > 0:
                    # Потеря содержимого — не то же самое, что вылезший за рамку
                    # текст: до читателя блок не дошёл вовсе. Правило записано в
                    # ENGINEERING.md и реализовано в TS-инспекторе
                    # (`inspectLayoutTelemetry`), но прогон зовёт этот файл — и
                    # здесь правила не было. Один вопрос, два ответа, причём
                    # строгий ответ не исполнялся никогда: на эталонной деке
                    # выброшенные буллеты страниц 11 и 29 докладывались как
                    # `text-clipping`.
                    clipping.append(
                        {
                            "page": int(row.get("page") or 0),
                            "code": "CONTENT_DROPPED_BY_RENDERER",
                            "severity": "CRITICAL",
                            "detail": (
                                f"рендерер выбросил содержимое: "
                                f"блоков={row.get('droppedBullets') or 0} "
                                f"строк={row.get('droppedLines') or 0} "
                                f"requiredHeight={row.get('requiredHeight')} "
                                f"availableHeight={row.get('availableHeight')} "
                                f"name={row.get('name') or row.get('role')} "
                                f"— страницу должна была разбить пагинация"
                            ),
                        }
                    )
                elif row.get("clipped") is True:
                    clipping.append(
                        {
                            "page": int(row.get("page") or 0),
                            "code": "text-clipping",
                            "severity": "CRITICAL",
                            "detail": (
                                f"text clipping requiredHeight={row.get('requiredHeight')} "
                                f"availableHeight={row.get('availableHeight')} "
                                f"name={row.get('name') or row.get('role')}"
                            ),
                        }
                    )
                elif row.get("measurementUncertain") is True:
                    clipping.append(
                        {
                            "page": int(row.get("page") or 0),
                            "code": "text-measurement-uncertain",
                            "severity": "WARNING",
                            "detail": "font measurement unavailable; clipping not proven",
                        }
                    )
        except Exception as exc:  # noqa: BLE001
            clipping.append(
                {
                    "page": 0,
                    "code": "telemetry-read-error",
                    "severity": "WARNING",
                    "detail": str(exc),
                }
            )

    if expect_pages and slide_count < expect_pages:
        for page in range(slide_count + 1, expect_pages + 1):
            overflow.append(
                {
                    "page": page,
                    "code": "missing-slide",
                    "severity": "CRITICAL",
                    "detail": "missing slide xml/page",
                }
            )

    return {
        "inspectorVersion": INSPECTOR_VERSION,
        "overlaps": overlaps,
        "overflow": overflow,
        "clipping": clipping,
        "emptyPages": empty_pages,
        "shapes": all_shapes,
        "pages": pages_meta,
        "slideCount": scanned,
        "constants": {
            "SLIDE_W": SLIDE_W,
            "SLIDE_H": SLIDE_H,
            "FOOTER_Y": FOOTER_Y,
            "CONTENT_BOTTOM": CONTENT_BOTTOM,
            "TEXT_COLLISION_RATIO": TEXT_COLLISION_RATIO,
        },
    }


def inspect_zip_fallback(pptx: Path, *, expect_pages: int | None = None) -> dict[str, Any]:
    overlaps: list[dict[str, Any]] = []
    overflow: list[dict[str, Any]] = []
    with zipfile.ZipFile(pptx) as z:
        pages = expect_pages or 0
        for page in range(1, (pages or 1) + 1):
            name = f"ppt/slides/slide{page}.xml"
            if name not in z.namelist():
                if expect_pages:
                    overflow.append(
                        {
                            "page": page,
                            "code": "missing-slide",
                            "severity": "CRITICAL",
                            "detail": "missing slide xml",
                        }
                    )
                continue
            if page == 1:
                overlaps.append(
                    {
                        "page": 0,
                        "code": "inspector-fallback",
                        "severity": "WARNING",
                        "detail": "python-pptx unavailable; role-aware overlap checks skipped",
                    }
                )
    return {
        "inspectorVersion": f"{INSPECTOR_VERSION}-fallback",
        "overlaps": overlaps,
        "overflow": overflow,
        "clipping": [],
        "emptyPages": [],
        "shapes": [],
        "pages": [],
        "slideCount": 0,
    }


def inspect(pptx: Path, *, expect_pages: int | None = None) -> dict[str, Any]:
    try:
        return inspect_presentation(pptx, expect_pages=expect_pages)
    except ImportError:
        return inspect_zip_fallback(pptx, expect_pages=expect_pages)


def main() -> int:
    if len(sys.argv) < 2:
        print(
            json.dumps(
                {"error": "usage: inspect-first36-pptx-geometry.py <pptx> [--expect-pages=36] [--texts]"}
            )
        )
        return 2
    pptx = Path(sys.argv[1])
    expect_pages: int | None = None
    texts_only = False
    for arg in sys.argv[2:]:
        if arg.startswith("--expect-pages="):
            expect_pages = int(arg.split("=", 1)[1])
        elif arg == "--expect-pages":
            expect_pages = 36
        elif arg == "--texts":
            texts_only = True
    if not pptx.exists():
        print(json.dumps({"error": f"missing {pptx}"}))
        return 2
    if texts_only:
        # Отдельным режимом, а не полем отчёта геометрии: `geometry-report.json`
        # лежит в эталоне, и текст всех страниц раздул бы его вдвое, меняясь от
        # каждой правки формулировки.
        try:
            print(json.dumps(page_texts(pptx), ensure_ascii=False))
            return 0
        except Exception as exc:  # noqa: BLE001
            print(json.dumps({"error": str(exc)}))
            return 1
    try:
        print(json.dumps(inspect(pptx, expect_pages=expect_pages), ensure_ascii=False))
        return 0
    except Exception as exc:  # noqa: BLE001
        print(json.dumps({"error": str(exc)}))
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
