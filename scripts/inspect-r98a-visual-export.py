#!/usr/bin/env python3
"""R9.8a visual export inspection — PDF/PPTX image presence + PNG page rasterization."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    import fitz  # type: ignore
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None


def _pdf_page_image_stats(pdf_path: Path) -> list[dict]:
    if fitz is None or not pdf_path.exists():
        return []
    doc = fitz.open(str(pdf_path))
    stats: list[dict] = []
    for i in range(len(doc)):
        page = doc[i]
        images = page.get_images(full=True)
        stats.append({"page": i + 1, "imageCount": len(images), "hasImages": len(images) > 0})
    return stats


def _pptx_picture_stats(pptx_path: Path) -> dict:
    slide_count = 0
    slides_with_pictures = 0
    total_pictures = 0
    serp_slide_pictures = 0

    if pptx_path.exists():
        import zipfile
        import re

        with zipfile.ZipFile(str(pptx_path), "r") as zf:
            slide_names = sorted(n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
            slide_count = len(slide_names)
            for idx, name in enumerate(slide_names, start=1):
                xml = zf.read(name).decode("utf-8", errors="ignore")
                pics = len(re.findall(r"<p:pic\b", xml))
                blips = len(re.findall(r"<a:blip\b", xml))
                count = max(pics, blips)
                if count:
                    slides_with_pictures += 1
                total_pictures += count
                text = re.sub(r"<[^>]+>", " ", xml).lower()
                if count and ("serp" in text or "поиск" in text or "снимок" in text or "выдач" in text):
                    serp_slide_pictures += count

    if Presentation is not None and pptx_path.exists() and total_pictures == 0:
        prs = Presentation(str(pptx_path))
        slide_count = len(prs.slides)
        for idx, slide in enumerate(prs.slides, start=1):
            pics = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pics += 1
            if pics:
                slides_with_pictures += 1
            total_pictures += pics
            text = " ".join(
                shape.text_frame.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame
            ).lower()
            if pics and ("serp" in text or "поиск" in text or "снимок" in text or "выдач" in text):
                serp_slide_pictures += pics

    return {
        "slideCount": slide_count,
        "slidesWithPictures": slides_with_pictures,
        "totalPictures": total_pictures,
        "serpSlidePictures": serp_slide_pictures,
    }


def _rasterize_pdf(pdf_path: Path, pages_out: Path) -> int:
    if fitz is None or not pdf_path.exists():
        return 0
    pages_out.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf_path))
    for i in range(len(doc)):
        pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(str(pages_out / f"page-{i + 1:02d}.png"))
    return len(doc)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--pdf", required=True)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--pages-out", required=True)
    args = parser.parse_args()

    pdf_path = Path(args.pdf)
    pptx_path = Path(args.pptx)
    pages_out = Path(args.pages_out)

    page_count = _rasterize_pdf(pdf_path, pages_out)
    pdf_pages = _pdf_page_image_stats(pdf_path)
    pptx_stats = _pptx_picture_stats(pptx_path)

    serp_pdf_pages = [p for p in pdf_pages if p["hasImages"] and p["page"] >= 8 and p["page"] <= 14]
    pdf_serp_has_images = len(serp_pdf_pages) > 0 or any(p["hasImages"] for p in pdf_pages[7:14] if pdf_pages)

    result = {
        "pageCount": page_count,
        "pdfSizeBytes": pdf_path.stat().st_size if pdf_path.exists() else 0,
        "pptxSizeBytes": pptx_path.stat().st_size if pptx_path.exists() else 0,
        "pdfPages": pdf_pages,
        "pptx": pptx_stats,
        "pdfSerpHasImages": pdf_serp_has_images,
        "pdfAnyImages": any(p["hasImages"] for p in pdf_pages),
        "pptxHasPictures": pptx_stats.get("totalPictures", 0) > 0,
        "serpPptxPictures": pptx_stats.get("serpSlidePictures", 0),
    }
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    sys.exit(main())
