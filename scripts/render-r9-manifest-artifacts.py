from __future__ import annotations

import json
import os
import re
import sys
from pathlib import Path

import fitz  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.util import Emu, Pt  # type: ignore

try:
    _RENDERER_PATH = Path(__file__).resolve().parents[1] / "renderer"
    if str(_RENDERER_PATH) not in sys.path:
        sys.path.insert(0, str(_RENDERER_PATH))
    from render_pptx import build_pptx as _renderer_build_pptx  # type: ignore
except Exception:  # pragma: no cover - fallback path
    _renderer_build_pptx = None

ALLOW_BRANDS = {
    "orion",
    "google",
    "yandex",
    "dow jones",
    "world-check",
    "lexisnexis",
    "pep",
    "rca",
    "kyc",
}


def _safe_text(value: object) -> str:
    text = str(value or "").strip()
    text = re.sub(r"(c:\\\\|/mnt/|storage/digital-profile|https?://[^\s]+)", "", text, flags=re.I)
    text = re.sub(r"(openai[_-]?api[_-]?key|sk-[a-z0-9]{10,})", "", text, flags=re.I)
    return text.strip()


def _is_russian_report(report_json: dict) -> bool:
    lang = str((report_json.get("meta") or {}).get("language") or "").lower()
    if lang in {"ru", "russian"}:
        return True
    return True


def _strip_english_leakage(text: str) -> str:
    if not text:
        return text
    out = text
    for token in re.findall(r"[A-Za-z][A-Za-z\-]{3,}", out):
        low = token.lower()
        if low in ALLOW_BRANDS:
            continue
        out = out.replace(token, "")
    out = re.sub(r"\s{2,}", " ", out).strip()
    return out


def _collect_slides(report_json: dict) -> list[dict]:
    manifest = (
        report_json.get("finalDeckManifest")
        or report_json.get("orionFinalDeckManifest")
        or {}
    )
    slides: list[dict] = []
    slides.append({"title": "Цифровой профиль", "subtitle": "ORION", "slideType": "cover_orion"})
    slides.append({"title": "Содержание", "subtitle": "Глобальная структура", "slideType": "toc_orion"})
    for section in manifest.get("sections") or []:
        for slide in section.get("slides") or []:
            slides.append(slide)
    return slides


def _write_pptx_fallback(report_json: dict, pptx_path: Path) -> int:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(6858000)
    slides = _collect_slides(report_json)
    total = max(1, len(slides))
    ru_mode = _is_russian_report(report_json)

    for idx, src in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0])
        title = _safe_text(src.get("title") or src.get("slideType") or f"Слайд {idx}")
        subtitle = _safe_text(src.get("subtitle") or "")
        if ru_mode:
            title = _strip_english_leakage(title)
            subtitle = _strip_english_leakage(subtitle)

        box = slide.shapes.add_textbox(Emu(500000), Emu(300000), Emu(8200000), Emu(900000))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title or f"Слайд {idx}"
        r.font.bold = True
        r.font.size = Pt(26)
        r.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

        if subtitle:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.size = Pt(13)
            r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        bullets_box = slide.shapes.add_textbox(Emu(500000), Emu(1400000), Emu(8200000), Emu(4300000))
        btf = bullets_box.text_frame
        bullets: list[str] = []
        for item in src.get("narrativeBlocks") or []:
            if isinstance(item, dict):
                text = _safe_text(item.get("text") or item.get("title") or "")
            else:
                text = _safe_text(item)
            if ru_mode:
                text = _strip_english_leakage(text)
            if text:
                bullets.append(text)
        if not bullets:
            bullets = ["Раздел сформирован из структуры слайдов по этапам анализа."]
        for b_i, line in enumerate(bullets[:8]):
            pp = btf.paragraphs[0] if b_i == 0 else btf.add_paragraph()
            rr = pp.add_run()
            rr.text = f"• {line}"
            rr.font.size = Pt(14)

        footer = slide.shapes.add_textbox(Emu(500000), Emu(6200000), Emu(8200000), Emu(300000))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.RIGHT
        fr = fp.add_run()
        fr.text = f"{idx}/{total}"
        fr.font.size = Pt(10)
        fr.font.color.rgb = RGBColor(0x77, 0x77, 0x77)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return total


def _write_pptx(report_json: dict, pptx_path: Path) -> int:
    if _renderer_build_pptx is not None:
        warnings, slide_count = _renderer_build_pptx(
            report_json,
            str(pptx_path),
            os.getcwd(),
            "report-template-v3",
            "client" if "client" in str(pptx_path).lower() else "internal",
            "draft",
        )
        if int(slide_count or 0) > 0:
            return int(slide_count or 0)
        _ = warnings
        return _write_pptx_fallback(report_json, pptx_path)
    return _write_pptx_fallback(report_json, pptx_path)


def _write_pdf(slide_count: int, report_json: dict, pdf_path: Path) -> None:
    slides = _collect_slides(report_json)
    doc = fitz.open()
    total = max(1, slide_count)
    ru_mode = _is_russian_report(report_json)
    for idx in range(total):
        page = doc.new_page(width=1280, height=720)
        src = slides[idx] if idx < len(slides) else {}
        title = _safe_text(src.get("title") or src.get("slideType") or f"Слайд {idx+1}")
        if ru_mode:
            title = _strip_english_leakage(title)
        page.insert_text((50, 60), title, fontsize=24)
        page.insert_text((50, 110), "ORION Section Pipeline v1", fontsize=12)
        page.insert_text((1160, 700), f"{idx+1}/{total}", fontsize=10)
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(pdf_path))
    doc.close()


def _export_png_pages(pdf_path: Path, out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    for child in out_dir.glob("page-*.png"):
        child.unlink(missing_ok=True)
    doc = fitz.open(str(pdf_path))
    for i in range(len(doc)):
        pix = doc[i].get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(str(out_dir / f"page-{i+1:02d}.png"))
    doc.close()


def main() -> int:
    if len(sys.argv) != 5:
        print(
            "Usage: render-r9-manifest-artifacts.py <report-json> <pptx-out> <pdf-out> <pages-dir>",
            file=sys.stderr,
        )
        return 1
    report_json_path = Path(sys.argv[1])
    pptx_out = Path(sys.argv[2])
    pdf_out = Path(sys.argv[3])
    pages_dir = Path(sys.argv[4])
    report_json = json.loads(report_json_path.read_text(encoding="utf-8"))
    slide_count = _write_pptx(report_json, pptx_out)
    _write_pdf(slide_count, report_json, pdf_out)
    _export_png_pages(pdf_out, pages_dir)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

